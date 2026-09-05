from __future__ import annotations

import sys
from pathlib import Path
from types import SimpleNamespace

import pytest
from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

import material_parser
from material_models import MaterialAsset, ParsedDocument
from material_parser import (
    LegacyOfficeConversionParser,
    ScannedPdfOcrParser,
    _blocks_from_pdf_backend,
    _enrich_docx_structure,
    _pptx_visual_evidence,
)
from material_storage import MaterialRepository


class FakeUpload:
    def __init__(self, filename: str, content: bytes, content_type: str):
        self.filename = filename
        self.content_type = content_type
        self._content = content
        self._read = False

    async def read(self, _size: int) -> bytes:
        if self._read:
            return b""
        self._read = True
        return self._content


def asset(extension: str) -> MaterialAsset:
    return MaterialAsset(
        asset_id="mat-test",
        filename=f"source{extension}",
        extension=extension,
        mime_type="application/octet-stream",
        detected_mime="application/octet-stream",
        size_bytes=10,
        sha256="abc",
        source_name=f"source{extension}",
        uploaded_at="2026-08-27T00:00:00Z",
        updated_at="2026-08-27T00:00:00Z",
    )


@pytest.mark.asyncio
async def test_legacy_office_binary_is_accepted_but_fake_binary_is_rejected(tmp_path):
    repository = MaterialRepository(tmp_path / "materials")
    signature = b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1"

    stored = await repository.save_upload(FakeUpload(
        "old-lesson.doc",
        signature + b"legacy-office-payload",
        "application/msword",
    ))

    assert stored.extension == ".doc"
    assert stored.detected_mime == "application/msword"
    with pytest.raises(ValueError, match="旧版 Office 文件结构无效"):
        await repository.save_upload(FakeUpload("fake.ppt", b"not-office", "application/vnd.ms-powerpoint"))


def test_pptx_visual_evidence_keeps_layout_chart_notes_and_picture_boundary(tmp_path):
    deck_path = tmp_path / "lesson.pptx"
    image_path = tmp_path / "figure.png"
    Image.new("RGB", (32, 24), "white").save(image_path)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1)).text = "导数图像"
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(2), Inches(2), Inches(1.5))
    chart_data = CategoryChartData()
    chart_data.categories = ["一", "二"]
    chart_data.add_series("变化率", (1, 3))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(4),
        Inches(1),
        Inches(4),
        Inches(3),
        chart_data,
    )
    slide.notes_slide.notes_text_frame.text = "提醒学生区分平均变化率与瞬时变化率。"
    presentation.save(deck_path)

    blocks, quality, warnings = _pptx_visual_evidence(deck_path, starting_order=3)

    visual = next(item for item in blocks if item.metadata.get("evidence_kind") == "pptx_visual_structure")
    notes = next(item for item in blocks if item.metadata.get("evidence_kind") == "speaker_notes")
    assert visual.locator.slide == 1
    assert visual.metadata["shape_count"] >= 3
    assert quality["chart_count"] == 1
    assert quality["picture_count"] == 1
    assert quality["unread_picture_count"] == 1
    assert "瞬时变化率" in notes.text
    assert any("未将图片内容冒充" in warning for warning in warnings)


def test_scanned_pdf_ocr_keeps_page_locator_and_confidence(monkeypatch, tmp_path):
    class FakeBitmap:
        @staticmethod
        def to_pil():
            return Image.new("RGB", (40, 30), "white")

    class FakePage:
        @staticmethod
        def render(scale: float):
            assert scale == 2.0
            return FakeBitmap()

        @staticmethod
        def close():
            return None

    class FakePdf:
        def __len__(self):
            return 2

        def __getitem__(self, _index: int):
            return FakePage()

        @staticmethod
        def close():
            return None

    monkeypatch.setitem(sys.modules, "pypdfium2", SimpleNamespace(PdfDocument=lambda _path: FakePdf()))
    monkeypatch.setitem(sys.modules, "rapidocr_onnxruntime", SimpleNamespace(RapidOCR=lambda: object()))
    monkeypatch.setattr(
        material_parser,
        "_ocr_image",
        lambda _path, *, engine, page_number: [{
            "text": f"第 {page_number} 页导数定义",
            "confidence": .91,
            "bbox": {"x": .1, "y": .2, "width": .5, "height": .1},
        }],
    )
    source = tmp_path / "scan.pdf"
    source.write_bytes(b"%PDF-fake")

    parsed = ScannedPdfOcrParser().parse(asset(".pdf"), source)

    assert parsed.parse_status == "degraded"
    assert [item.locator.page for item in parsed.blocks] == [1, 2]
    assert parsed.quality["ocr_confidence"] == .91
    assert parsed.quality["source_page_count"] == 2


def test_legacy_conversion_is_temporary_and_reuses_existing_parser(monkeypatch, tmp_path):
    source = tmp_path / "source.doc"
    original = b"legacy-source"
    source.write_bytes(original)
    monkeypatch.setattr(material_parser.shutil, "which", lambda _name: "/usr/bin/soffice")

    def fake_run(arguments, **_kwargs):
        output_dir = Path(arguments[arguments.index("--outdir") + 1])
        (output_dir / "source.docx").write_bytes(b"converted")
        return SimpleNamespace(returncode=0, stdout="", stderr="")

    monkeypatch.setattr(material_parser.subprocess, "run", fake_run)
    monkeypatch.setattr(
        material_parser.DoclingDocumentParser,
        "parse",
        lambda self, converted_asset, _path: ParsedDocument(
            document_id="doc-1",
            asset_id=converted_asset.asset_id,
            source_sha256=converted_asset.sha256,
            parse_status="parsed",
            parser_name=self.name,
            parser_version="test",
            parse_options_hash="test",
            created_at="2026-08-27T00:00:00Z",
        ),
    )

    parsed = LegacyOfficeConversionParser().parse(asset(".doc"), source)

    assert source.read_bytes() == original
    assert parsed.parser_name == "libreoffice_legacy+docling"
    assert "临时目录转换" in parsed.warnings[0]


def test_docx_structure_keeps_heading_list_table_and_review_boundaries(tmp_path):
    from docx import Document

    source = tmp_path / "lesson.docx"
    document = Document()
    document.add_heading("第一讲 导数", level=1)
    document.add_paragraph("能解释导数的几何意义。")
    document.add_paragraph("使用切线图像判断变化率。", style="List Bullet")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "教学环节"
    table.cell(0, 1).text = "时长"
    table.cell(1, 0).text = "导入"
    table.cell(1, 1).text = "5"
    document.sections[0].header.paragraphs[0].text = "高等数学"
    document.save(source)
    blocks = [
        material_parser.DocumentBlock(block_id="b1", kind="paragraph", text="第一讲 导数", order=0),
        material_parser.DocumentBlock(block_id="b2", kind="paragraph", text="能解释导数的几何意义。", order=1),
        material_parser.DocumentBlock(block_id="b3", kind="paragraph", text="使用切线图像判断变化率。", order=2),
        material_parser.DocumentBlock(block_id="b4", kind="table", text="教学环节 | 时长 | 导入 | 5", order=3),
    ]

    enriched, quality, warnings = _enrich_docx_structure(source, blocks)

    assert enriched[0].kind == "title"
    assert enriched[0].metadata["heading_level"] == 1
    assert enriched[2].kind == "list_item"
    assert enriched[3].metadata["rows"][1] == ["导入", "5"]
    assert any(item.metadata.get("evidence_kind") == "docx_header" for item in enriched)
    assert quality["heading_count"] == 1
    assert quality["table_count"] == 1
    assert warnings == []


def test_pdf_text_cells_keep_page_and_source_locator() -> None:
    class Cell:
        def __init__(self, text: str, left: float):
            self.text = text
            self.rect = SimpleNamespace(left=left, top=1.0, right=left + 4.0, bottom=2.0)

    class Page:
        def get_text_cells(self):
            return [Cell("教学目标", 1.0), Cell("核心概念", 6.0)]

        def unload(self):
            return None

    class Backend:
        def iter_pages(self):
            return [Page()]

        def unload(self):
            return None

    blocks = _blocks_from_pdf_backend(Backend())

    assert [item.text for item in blocks] == ["教学目标", "核心概念"]
    assert all(item.locator.page == 1 for item in blocks)
    assert blocks[0].locator.bbox == {"l": 1.0, "t": 1.0, "r": 5.0, "b": 2.0}
    assert blocks[1].metadata["page_cell_index"] == 2
