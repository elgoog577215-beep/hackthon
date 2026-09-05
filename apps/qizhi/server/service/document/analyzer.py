"""文档分析服务：提取文本 + LLM 多维度 map-reduce 分析。"""

from __future__ import annotations

import json
import math
import random
import re
import time
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from typing import Any

from openai import OpenAI
from pydantic import TypeAdapter

from common.config import settings
from common.utils.logger import get_logger

from .models import (
    DimensionEvaluation,
    DimensionScore,
    DocumentAnalysisResult,
    SegmentAnalysisResult,
    SegmentDetail,
)

logger = get_logger(__name__)

# 6 分析维度
ANALYSIS_DIMENSIONS = [
    "教学目标清晰度",
    "内容组织逻辑",
    "教学方法适切性",
    "学生互动设计",
    "思政融合",
    "规范性",
]

# 分段配置
_DOCX_MAX_CHARS_PER_SEGMENT = 3000  # docx 按字符数分段上限
_PPTX_SLIDES_PER_SEGMENT = 8       # pptx 每段最多 8 张幻灯片


# ---------------------------------------------------------------------------
# 文本提取
# ---------------------------------------------------------------------------

def extract_text_from_docx(file_path: str) -> list[dict]:
    """从 .docx 提取文本段落，按标题层级分组。

    返回: [{"title": "章节标题", "content": "正文..."}]
    """
    from docx import Document

    doc = Document(file_path)
    sections: list[dict] = []
    current_title = ""
    current_lines: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        style_name = (para.style.name or "").lower()
        is_heading = style_name.startswith("heading") or style_name.startswith("title")

        if is_heading:
            # 遇到新标题 -> 存储前一节
            if current_lines:
                sections.append({"title": current_title, "content": "\n".join(current_lines)})
                current_lines = []
            current_title = text
        else:
            current_lines.append(text)

    # 尾部
    if current_lines:
        sections.append({"title": current_title, "content": "\n".join(current_lines)})

    # 如果整个文档无标题结构，将所有内容作为一段
    if not sections:
        full_text = "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())
        if full_text:
            sections.append({"title": "全文", "content": full_text})

    logger.info("[docx] 提取到 %d 个文本段落", len(sections))
    return sections


def extract_text_from_pptx(file_path: str) -> list[dict]:
    """从 .pptx 提取幻灯片文本。

    返回: [{"title": "幻灯片 N", "content": "文本..."}]
    """
    from pptx import Presentation

    prs = Presentation(file_path)
    slides: list[dict] = []

    for idx, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        slide_title = ""

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)

            # 尝试从标题占位符获取标题
            if hasattr(shape, "placeholder_format") and shape.placeholder_format is not None:
                if shape.placeholder_format.idx == 0:  # 标题占位符
                    slide_title = shape.text.strip()

        content = "\n".join(texts)
        if content:
            title = slide_title or f"幻灯片 {idx}"
            slides.append({"title": title, "content": content})

    logger.info("[pptx] 提取到 %d 张幻灯片文本", len(slides))
    return slides


# ---------------------------------------------------------------------------
# 分段
# ---------------------------------------------------------------------------

def _segment_sections(sections: list[dict], file_type: str) -> list[dict]:
    """将提取的 sections 聚合为分析段。

    每段: {"segment_index": 1, "segment_title": "...", "text": "..."}
    """
    if file_type == "pptx":
        return _segment_by_count(sections, _PPTX_SLIDES_PER_SEGMENT)
    return _segment_by_chars(sections, _DOCX_MAX_CHARS_PER_SEGMENT)


def _segment_by_chars(sections: list[dict], max_chars: int) -> list[dict]:
    """按字符数合并相邻 section 到同一段。"""
    segments: list[dict] = []
    buf_titles: list[str] = []
    buf_text: list[str] = []
    buf_len = 0

    for sec in sections:
        sec_len = len(sec["content"])
        if buf_len > 0 and buf_len + sec_len > max_chars:
            segments.append(_make_segment(len(segments) + 1, buf_titles, buf_text))
            buf_titles, buf_text, buf_len = [], [], 0
        buf_titles.append(sec["title"])
        buf_text.append(sec["content"])
        buf_len += sec_len

    if buf_text:
        segments.append(_make_segment(len(segments) + 1, buf_titles, buf_text))

    return segments


def _segment_by_count(sections: list[dict], group_size: int) -> list[dict]:
    """按固定个数分组。"""
    segments: list[dict] = []
    for i in range(0, len(sections), group_size):
        group = sections[i : i + group_size]
        titles = [s["title"] for s in group]
        texts = [s["content"] for s in group]
        segments.append(_make_segment(len(segments) + 1, titles, texts))
    return segments


def _make_segment(index: int, titles: list[str], texts: list[str]) -> dict:
    title = " / ".join(t for t in titles if t)
    return {
        "segment_index": index,
        "segment_title": title or f"第 {index} 段",
        "text": "\n\n".join(texts),
    }


# ---------------------------------------------------------------------------
# LLM 调用
# ---------------------------------------------------------------------------

_TRANSIENT_KEYWORDS = (
    "503", "429", "500", "unavailable", "resource_exhausted",
    "timeout", "timed out", "connection", "reset by peer",
    "overloaded", "try again",
)


def _is_transient(e: Exception) -> bool:
    s = str(e).lower()
    return any(k in s for k in _TRANSIENT_KEYWORDS)


def _strip_fences(text: str) -> str:
    t = text.strip()
    if t.startswith("```"):
        nl = t.find("\n")
        if nl != -1:
            t = t[nl + 1:]
        if t.rstrip().endswith("```"):
            t = t.rstrip()[:-3]
    return t.strip()


def _make_clients() -> list[OpenAI]:
    """从配置创建 OpenAI 客户端（支持逗号分隔多实例）。"""
    urls = [u.strip() for u in settings.LOCAL_ANALYSIS_TEXT_BASE_URL.split(",") if u.strip()]
    api_key = settings.LOCAL_ANALYSIS_API_KEY or "EMPTY"
    return [OpenAI(api_key=api_key, base_url=u, timeout=180) for u in urls]


def _llm_chat(
    clients: list[OpenAI],
    system: str,
    user: str,
    model: str,
    *,
    temperature: float = 0.3,
    max_tokens: int | None = None,
    max_tries: int = 5,
) -> str:
    """调用 vLLM 聊天补全（JSON 模式 + 重试）。"""
    messages = [{"role": "system", "content": system}, {"role": "user", "content": user}]
    delay = 1.0
    last: Exception | None = None

    for i in range(max_tries):
        try:
            client = random.choice(clients)
            resp = client.chat.completions.create(
                model=model,
                messages=messages,
                temperature=temperature,
                response_format={"type": "json_object"},
                extra_body={"chat_template_kwargs": {"enable_thinking": False}},
                max_tokens=max_tokens,
            )
            return resp.choices[0].message.content or ""
        except Exception as e:
            last = e
            if _is_transient(e) and i < max_tries - 1:
                wait = delay + random.uniform(0, delay * 0.5)
                logger.warning("[llm] 瞬时错误，%.1fs 后重试（%d/%d）：%s", wait, i + 1, max_tries, str(e)[:160])
                time.sleep(wait)
                delay = min(delay * 2, 8.0)
                continue
            raise

    raise last  # type: ignore[misc]


def _analyze_segment(
    clients: list[OpenAI],
    model: str,
    segment: dict,
    retries: int = 2,
) -> SegmentAnalysisResult:
    """对单段文本调用 LLM 进行 6 维度分析（map 阶段）。"""
    system_prompt = (
        "你是一名资深的高校教学督导与课堂分析专家。请对以下教学文档内容进行多维度分析评价。\n"
        "请只输出 JSON，不要任何多余文字。"
    )

    dimensions_desc = "、".join(ANALYSIS_DIMENSIONS)
    schema_json = json.dumps(SegmentAnalysisResult.model_json_schema(), ensure_ascii=False)

    user_prompt = (
        f"以下是教学文档的一个段落（标题：{segment['segment_title']}）：\n\n"
        f"{segment['text']}\n\n"
        f"请从以下 6 个维度进行分析评价：{dimensions_desc}。\n"
        f"对每个维度，给出：\n"
        f"- dimension：维度名称\n"
        f"- score：得分 (0-100)\n"
        f"- strengths：优点（字符串列表）\n"
        f"- weaknesses：不足（字符串列表）\n"
        f"- suggestions：改进建议（字符串列表）\n\n"
        f"请输出一个 JSON 对象，必须符合以下 JSON Schema：\n{schema_json}"
    )

    temp = 0.3
    last: Exception | None = None

    for attempt in range(retries + 1):
        try:
            raw = _llm_chat(clients, system_prompt, user_prompt, model, temperature=temp)
            obj = json.loads(_strip_fences(raw))
            return TypeAdapter(SegmentAnalysisResult).validate_python(obj)
        except Exception as e:
            last = e
            logger.warning(
                "[segment-%d] 解析失败（第 %d 次）：%s",
                segment["segment_index"], attempt + 1, str(e)[:160],
            )
            # 追加纠错提示
            user_prompt = (
                f"以下是教学文档的一个段落（标题：{segment['segment_title']}）：\n\n"
                f"{segment['text']}\n\n"
                f"请从以下 6 个维度进行分析评价：{dimensions_desc}。\n"
                f"对每个维度，给出：\n"
                f"- dimension：维度名称\n"
                f"- score：得分 (0-100)\n"
                f"- strengths：优点（字符串列表）\n"
                f"- weaknesses：不足（字符串列表）\n"
                f"- suggestions：改进建议（字符串列表）\n\n"
                f"请输出一个 JSON 对象，必须符合以下 JSON Schema：\n{schema_json}\n\n"
                f"上次输出无法解析或不符合结构（{str(e)[:100]}）。请严格只输出符合上述结构的 JSON。"
            )
            temp = max(0.0, temp - 0.15)

    raise RuntimeError(f"[segment-{segment['segment_index']}] 多次重试仍失败：{last}")


# ---------------------------------------------------------------------------
# Reduce：合并分段结果
# ---------------------------------------------------------------------------

def _merge_segment_results(
    segments: list[dict],
    segment_results: list[SegmentAnalysisResult],
) -> DocumentAnalysisResult:
    """将各段分析结果合并为整体评价（reduce 阶段）。"""
    # 按维度聚合
    dim_scores: dict[str, list[int]] = {d: [] for d in ANALYSIS_DIMENSIONS}
    dim_strengths: dict[str, list[str]] = {d: [] for d in ANALYSIS_DIMENSIONS}
    dim_weaknesses: dict[str, list[str]] = {d: [] for d in ANALYSIS_DIMENSIONS}
    dim_suggestions: dict[str, list[str]] = {d: [] for d in ANALYSIS_DIMENSIONS}

    segment_details: list[SegmentDetail] = []

    for seg, result in zip(segments, segment_results):
        seg_dims: list[DimensionEvaluation] = []
        for dim_eval in result.dimensions:
            name = dim_eval.dimension
            if name not in dim_scores:
                # 容错：LLM 可能输出略有偏差的维度名，尝试模糊匹配
                matched = _fuzzy_match_dimension(name)
                if matched:
                    name = matched
                else:
                    continue
            dim_scores[name].append(dim_eval.score)
            dim_strengths[name].extend(dim_eval.strengths)
            dim_weaknesses[name].extend(dim_eval.weaknesses)
            dim_suggestions[name].extend(dim_eval.suggestions)
            seg_dims.append(dim_eval)

        segment_details.append(SegmentDetail(
            segment_index=seg["segment_index"],
            segment_title=seg["segment_title"],
            dimensions=seg_dims,
        ))

    # 综合得分
    dimension_score_list: list[DimensionScore] = []
    for dim_name in ANALYSIS_DIMENSIONS:
        scores = dim_scores[dim_name]
        avg = round(sum(scores) / len(scores)) if scores else 0

        # 汇总评语
        strengths = list(dict.fromkeys(dim_strengths[dim_name]))[:5]
        weaknesses = list(dict.fromkeys(dim_weaknesses[dim_name]))[:5]
        comments_parts: list[str] = []
        if strengths:
            comments_parts.append("优点：" + "；".join(strengths))
        if weaknesses:
            comments_parts.append("不足：" + "；".join(weaknesses))
        comments = "。".join(comments_parts)

        suggestions = list(dict.fromkeys(dim_suggestions[dim_name]))[:5]

        dimension_score_list.append(DimensionScore(
            dimension=dim_name,
            score=avg,
            comments=comments,
            suggestions=suggestions,
        ))

    all_scores = [ds.score for ds in dimension_score_list]
    overall_score = round(sum(all_scores) / len(all_scores)) if all_scores else 0

    # 雷达图数据
    radar_data = [
        {"dimension": ds.dimension, "score": ds.score, "fullScore": 100}
        for ds in dimension_score_list
    ]

    # 能力概述
    strong = [ds.dimension for ds in dimension_score_list if ds.score >= 80]
    weak = [ds.dimension for ds in dimension_score_list if ds.score < 60]
    summary_parts: list[str] = [f"文档总体得分 {overall_score} 分。"]
    if strong:
        summary_parts.append(f"在{'、'.join(strong)}方面表现较好。")
    if weak:
        summary_parts.append(f"在{'、'.join(weak)}方面有较大提升空间。")
    capability_summary = "".join(summary_parts)

    # 整体改进建议（从各维度中挑选得分最低的维度建议）
    sorted_dims = sorted(dimension_score_list, key=lambda d: d.score)
    improvement_suggestions: list[str] = []
    for ds in sorted_dims[:3]:
        for sug in ds.suggestions[:2]:
            improvement_suggestions.append(f"【{ds.dimension}】{sug}")

    return DocumentAnalysisResult(
        file_name="",   # 由调用方填充
        file_type="",   # 由调用方填充
        total_segments=len(segments),
        dimension_scores=dimension_score_list,
        overall_score=overall_score,
        radar_data=radar_data,
        capability_summary=capability_summary,
        improvement_suggestions=improvement_suggestions,
        segment_details=segment_details,
    )


def _fuzzy_match_dimension(name: str) -> str | None:
    """模糊匹配维度名称（容错 LLM 输出偏差）。"""
    for dim in ANALYSIS_DIMENSIONS:
        # 完全包含即匹配
        if dim in name or name in dim:
            return dim
    return None


# ---------------------------------------------------------------------------
# 主入口
# ---------------------------------------------------------------------------

def analyze_document(
    file_path: str,
    file_type: str,
    *,
    progress_callback: Any | None = None,
) -> DocumentAnalysisResult:
    """分析教学文档（同步，需由调用方用 asyncio.to_thread 包装）。

    Args:
        file_path: 文档文件路径
        file_type: 文件类型 (docx / pptx)
        progress_callback: 可选的进度回调 fn(current, total, message)
    """
    file_name = Path(file_path).name

    # 1. 提取文本
    logger.info("[doc] 开始提取文本：%s (类型=%s)", file_name, file_type)
    if progress_callback:
        progress_callback(0, 100, "正在提取文档内容...")

    if file_type == "pptx":
        sections = extract_text_from_pptx(file_path)
    else:
        sections = extract_text_from_docx(file_path)

    if not sections:
        raise ValueError("文档内容为空，无法分析")

    # 2. 分段
    segments = _segment_sections(sections, file_type)
    total_segments = len(segments)
    logger.info("[doc] 文档分为 %d 个分析段", total_segments)

    if progress_callback:
        progress_callback(10, 100, f"文档已分为 {total_segments} 段，开始 LLM 分析...")

    # 3. Map 阶段：并行分析各段
    clients = _make_clients()
    model = settings.LOCAL_ANALYSIS_MODEL
    workers = min(4, total_segments)

    def analyze_one(seg: dict) -> SegmentAnalysisResult:
        idx = seg["segment_index"]
        logger.info("[doc] 分析第 %d/%d 段：%s", idx, total_segments, seg["segment_title"])
        result = _analyze_segment(clients, model, seg)
        if progress_callback:
            pct = 10 + int(80 * idx / total_segments)
            progress_callback(pct, 100, f"已完成第 {idx}/{total_segments} 段分析")
        return result

    if workers > 1 and total_segments > 1:
        with ThreadPoolExecutor(max_workers=workers) as pool:
            segment_results = list(pool.map(analyze_one, segments))
    else:
        segment_results = [analyze_one(seg) for seg in segments]

    # 4. Reduce 阶段：合并结果
    if progress_callback:
        progress_callback(90, 100, "正在汇总分析结果...")

    result = _merge_segment_results(segments, segment_results)
    result.file_name = file_name
    result.file_type = file_type

    if progress_callback:
        progress_callback(100, 100, "分析完成")

    logger.info("[doc] 分析完成：%s，总分=%d", file_name, result.overall_score)
    return result
