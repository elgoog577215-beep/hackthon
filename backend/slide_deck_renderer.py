"""Editable PPTX renderer for the shared structured slide deck contract."""

from __future__ import annotations

import os
import re
import shutil
import subprocess
import tempfile
from collections import Counter
from copy import deepcopy
from difflib import SequenceMatcher
from functools import lru_cache
from pathlib import Path
from typing import Any

from slide_asset_repository import SlideAssetRepository, slide_asset_repository
from slide_deck import SlideBlockSpec, SlideDeckContent, SlideSpec, validate_slide_deck
from slide_layout_geometry import (
    BALANCED_TWO_COLUMN_BODY_FONT_WIDTH_SAFETY_FACTOR,
    BALANCED_TWO_COLUMN_BODY_V1,
    CLASSIFICATION_THREE_CARDS_V1,
    HORIZONTAL_PROCESS_CARDS_V1,
    balanced_two_column_body_metrics,
    classification_three_card_metrics,
    diagram_node_layout_metrics,
    horizontal_process_card_metrics,
    wrapped_line_count,
)
from slide_theme import load_slide_theme_pack, slide_theme_asset_path

THEMES: dict[str, dict[str, Any]] = {
    "qingfeng-classroom": {
        "surface": "F7FAFC",
        "canvas": "EBF4FF",
        "chart_bg": "E2E8F0",
        "title": "1A365D",
        "ink": "4A5568",
        "muted": "718096",
        "accent": "2B6CB0",
        "accent_soft": "BEE3F8",
        "green": "2F855A",
        "green_soft": "F0FFF4",
        "amber": "ED8936",
        "amber_soft": "FFFAF0",
        "red": "B54735",
        "red_soft": "FFF1EE",
        "code": "1A365D",
        "title_font": "Noto Sans SC",
        "title_east_asian_font": "Microsoft YaHei",
        "body_font": "Noto Sans SC",
        "body_east_asian_font": "Microsoft YaHei",
        "math_font": "Times New Roman",
    },
    "academic-bluegray": {
        "surface": "FCFCFD",
        "canvas": "E8EBEE",
        "chart_bg": "E8EBEE",
        "title": "2C3E50",
        "ink": "5D6D7E",
        "muted": "7F8C8D",
        "accent": "2E86C1",
        "accent_soft": "D6EAF8",
        "green": "2874A6",
        "green_soft": "EAF2F8",
        "amber": "B9770E",
        "amber_soft": "FDF2E9",
        "red": "922B21",
        "red_soft": "FDEDEC",
        "code": "2C3E50",
        "title_font": "Noto Serif SC",
        "title_east_asian_font": "SimSun",
        "body_font": "Noto Sans SC",
        "body_east_asian_font": "Microsoft YaHei",
        "math_font": "Times New Roman",
    },
    "qizhi-classroom": {
        "surface": "FFFDF7",
        "canvas": "FFF3D6",
        "chart_bg": "E6F2FF",
        "title": "17365D",
        "ink": "34465C",
        "muted": "6D7D91",
        "accent": "2F6FE4",
        "accent_soft": "DCE9FF",
        "green": "16856B",
        "green_soft": "E2F7F0",
        "amber": "F29D38",
        "amber_soft": "FFF0D9",
        "red": "C45443",
        "red_soft": "FCE7E3",
        "code": "14243B",
        "title_font": "Noto Sans SC",
        "title_east_asian_font": "Microsoft YaHei",
        "body_font": "Noto Sans SC",
        "body_east_asian_font": "Microsoft YaHei",
        "math_font": "Times New Roman",
    },
    "academic-editorial": {
        "surface": "FBFAF7",
        "canvas": "EFEEE9",
        "chart_bg": "E4E7E9",
        "title": "273340",
        "ink": "45515D",
        "muted": "727A82",
        "accent": "315E7D",
        "accent_soft": "DCE6EB",
        "green": "4F6D64",
        "green_soft": "E6EDEA",
        "amber": "8B6B3E",
        "amber_soft": "F2EBDD",
        "red": "824C45",
        "red_soft": "F3E4E1",
        "code": "252E35",
        "title_font": "Noto Serif SC",
        "title_east_asian_font": "SimSun",
        "body_font": "Noto Sans SC",
        "body_east_asian_font": "Microsoft YaHei",
        "math_font": "Times New Roman",
    },
    "grid-notebook": {
        "surface": "FAF8F0",
        "canvas": "F2EEDC",
        "chart_bg": "DFE8E3",
        "title": "283B36",
        "ink": "40524D",
        "muted": "73817C",
        "accent": "2D7464",
        "accent_soft": "D9EAE3",
        "green": "648B57",
        "green_soft": "E8F0E2",
        "amber": "D18A32",
        "amber_soft": "F8E9CF",
        "red": "B75A48",
        "red_soft": "F4DFDA",
        "code": "253D38",
        "title_font": "Noto Sans SC",
        "title_east_asian_font": "Microsoft YaHei",
        "body_font": "Noto Sans SC",
        "body_east_asian_font": "Microsoft YaHei",
        "math_font": "Times New Roman",
    },
    "modern-geometric": {
        "surface": "F6F3FF",
        "canvas": "E9E1FF",
        "chart_bg": "DAD2F2",
        "title": "231A4A",
        "ink": "463D62",
        "muted": "746C8B",
        "accent": "6548E8",
        "accent_soft": "DDD4FF",
        "green": "138D85",
        "green_soft": "DDF5F1",
        "amber": "F08B3E",
        "amber_soft": "FFE6D5",
        "red": "D45168",
        "red_soft": "FADCE3",
        "code": "211B3A",
        "title_font": "Noto Sans SC",
        "title_east_asian_font": "Microsoft YaHei",
        "body_font": "Noto Sans SC",
        "body_east_asian_font": "Microsoft YaHei",
        "math_font": "Times New Roman",
    },
    "dark-tech": {
        "surface": "0C1321",
        "canvas": "16243A",
        "chart_bg": "22334B",
        "title": "F3F8FF",
        "ink": "D7E3F2",
        "muted": "91A6BE",
        "accent": "4DB5FF",
        "accent_soft": "183C5A",
        "green": "40D6B1",
        "green_soft": "173D3A",
        "amber": "FFB35A",
        "amber_soft": "49351F",
        "red": "FF7385",
        "red_soft": "482631",
        "code": "070D16",
        "title_font": "Noto Sans SC",
        "title_east_asian_font": "Microsoft YaHei",
        "body_font": "Noto Sans SC",
        "body_east_asian_font": "Microsoft YaHei",
        "math_font": "Times New Roman",
    },
}

# The five v3 themes are authored once and consumed by both Vue and PPTX.
_SHARED_THEMES = load_slide_theme_pack()["themes"]
THEMES.update({name: dict(tokens) for name, tokens in _SHARED_THEMES.items()})

BODY_FONT = "Noto Sans SC"
BODY_EAST_ASIAN_FONT = "Microsoft YaHei"
CODE_FONT = "Aptos Mono"


class SlideDeckQualityError(ValueError):
    def __init__(self, report: dict[str, Any]) -> None:
        self.report = report
        codes = ", ".join(item["code"] for item in report["blockers"])
        super().__init__(f"Slide deck quality gate blocked export: {codes}")


def export_structured_slide_deck(
    content: dict[str, Any],
    output_path: str | Path,
    *,
    require_quality: bool = True,
    theme: str | dict[str, Any] = "qingfeng-classroom",
    asset_repository: SlideAssetRepository | None = None,
    course_data: dict[str, Any] | None = None,
) -> Path:
    """Render the same slide spec used by the browser preview into editable PPTX."""
    deck = SlideDeckContent.model_validate(content)
    payload = deck.model_dump(mode="json")
    if deck.schema_version == "slide_deck_v5":
        for unit in deck.slides:
            quality = unit.quality or {}
            if not quality.get("resolved_layout"):
                raise ValueError(
                    f"v5_final_layout_missing:{unit.unit_id}"
                )
            if str(quality["resolved_layout"]) not in V5_LAYOUT_RENDERER_NAMES:
                raise ValueError(
                    "v5_final_layout_unsupported:"
                    f"{unit.unit_id}:{quality['resolved_layout']}"
                )
    if require_quality:
        if deck.schema_version == "slide_deck_v5":
            from slide_deck_v5 import (
                v5_contract_issues,
                validate_slide_deck_v5,
            )

            if course_data is not None:
                report = validate_slide_deck_v5(
                    payload,
                    course_data=course_data,
                )
            elif payload.get("quality_report"):
                embedded = dict(payload["quality_report"])
                composition_issues = v5_contract_issues(
                    list(payload.get("slides") or [])
                )
                blockers = [
                    *(embedded.get("blockers") or []),
                    *composition_issues,
                ]
                report = {
                    **embedded,
                    "passed": bool(embedded.get("passed")) and not blockers,
                    "blockers": blockers,
                }
            else:
                report = validate_slide_deck_v5(payload)
        else:
            report = validate_slide_deck(payload)
        if not report["passed"]:
            raise SlideDeckQualityError(report)

    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    theme_config = validate_theme(theme)
    resolved_asset_repository = asset_repository or slide_asset_repository

    for index, unit in enumerate(deck.slides):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        _render_slide(
            slide,
            unit,
            index + 1,
            len(deck.slides),
            theme_config,
            resolved_asset_repository,
        )
        if unit.speaker_notes:
            slide.notes_slide.notes_text_frame.text = unit.speaker_notes

    path = Path(output_path)
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)
    return path


@lru_cache(maxsize=32)
def _audit_font(font_size_px: int) -> Any:
    from PIL import ImageFont

    candidates = [
        Path("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc"),
        Path("/usr/share/fonts/opentype/noto/NotoSansCJKsc-Regular.otf"),
        Path("C:/Windows/Fonts/NotoSansCJK-Regular.ttc"),
        Path("C:/Windows/Fonts/msyh.ttc"),
        Path("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf"),
    ]
    for path in candidates:
        if not path.is_file():
            continue
        try:
            return ImageFont.truetype(str(path), font_size_px)
        except OSError:
            continue
    return ImageFont.load_default()


def _paragraph_font_size_pt(paragraph: Any) -> float:
    sizes = [
        float(run.font.size.pt)
        for run in paragraph.runs
        if run.font.size is not None
    ]
    if paragraph.font.size is not None:
        sizes.append(float(paragraph.font.size.pt))
    return max(sizes, default=18.0)


def _paragraph_minimum_font_size_pt(paragraph: Any) -> float:
    sizes = [
        float(run.font.size.pt)
        for run in paragraph.runs
        if run.font.size is not None
    ]
    if paragraph.font.size is not None:
        sizes.append(float(paragraph.font.size.pt))
    return min(sizes, default=18.0)


def _wrapped_line_count(
    text: str,
    *,
    width_pt: float,
    font_size_pt: float,
    portable_width_safety_factor: float = 0.99,
) -> int:
    return wrapped_line_count(
        text,
        width_pt=width_pt,
        font_size_pt=font_size_pt,
        font_loader=_audit_font,
        portable_width_safety_factor=portable_width_safety_factor,
    )


def _wrapped_text_has_orphan_last_line(
    text: str,
    *,
    width_pt: float,
    font_size_pt: float,
    max_orphan_chars: int = 2,
) -> bool:
    """Return whether wrapping leaves only a tiny final title fragment."""

    value = str(text or "").strip()
    if not value:
        return False
    explicit_last_line = value.rsplit("\n", 1)[-1].strip()
    if "\n" in value and 0 < len(explicit_last_line) <= max_orphan_chars:
        return True
    wrapped_lines = _wrapped_line_count(
        value,
        width_pt=width_pt,
        font_size_pt=font_size_pt,
    )
    if wrapped_lines <= 1:
        return False
    for suffix_length in range(1, min(max_orphan_chars, len(value)) + 1):
        prefix = value[:-suffix_length].rstrip()
        if not prefix:
            break
        if _wrapped_line_count(
            prefix,
            width_pt=width_pt,
            font_size_pt=font_size_pt,
        ) < wrapped_lines:
            return True
    return False


def _text_frame_audit(shape: Any) -> dict[str, Any]:
    frame = shape.text_frame
    width_pt = max(
        1.0,
        (
            int(shape.width)
            - int(frame.margin_left or 0)
            - int(frame.margin_right or 0)
        ) / 12700,
    )
    height_pt = max(
        1.0,
        (
            int(shape.height)
            - int(frame.margin_top or 0)
            - int(frame.margin_bottom or 0)
        ) / 12700,
    )
    required_height = 0.0
    minimum_size = 10**9
    maximum_lines = 1
    horizontal_overflow = False
    portable_width_safety_factor = (
        BALANCED_TWO_COLUMN_BODY_FONT_WIDTH_SAFETY_FACTOR
        if f"[v6-body-capacity={BALANCED_TWO_COLUMN_BODY_V1}]"
        in str(shape.name or "")
        else 0.99
    )
    for paragraph in frame.paragraphs:
        text = paragraph.text or ""
        font_size = _paragraph_font_size_pt(paragraph)
        minimum_size = min(minimum_size, _paragraph_minimum_font_size_pt(paragraph))
        line_count = _wrapped_line_count(
            text,
            width_pt=width_pt,
            font_size_pt=font_size,
            portable_width_safety_factor=portable_width_safety_factor,
        )
        # A measured teaching scene writes explicit line breaks and disables
        # automatic wrapping. Rewrapping with the Unicode width floor invents
        # lines (even four Chinese characters in an exactly 80pt slot).
        # Still measure every real line against the actual declared font so a
        # no-wrap flag cannot hide horizontal overflow.
        fixed_font = Path(__file__).resolve().parents[1] / "frontend/public/presentation-assets/fonts/NotoSansCJKsc-Regular.otf"
        fixed_lines = (
            frame.word_wrap is False
            and paragraph.font.name == "Noto Sans CJK SC"
            and all(run.font.name in (None, paragraph.font.name) for run in paragraph.runs)
            and fixed_font.is_file()
        )
        if fixed_lines:
            from PIL import ImageFont
            measured_font = ImageFont.truetype(str(fixed_font), round(font_size * 4))
            actual_lines = text.replace("\v", "\n").split("\n")
            line_count = len(actual_lines)
            horizontal_overflow |= any(measured_font.getlength(line) / 4 > width_pt + 0.01 for line in actual_lines)
        maximum_lines = max(maximum_lines, line_count)
        before = float(paragraph.space_before.pt) if paragraph.space_before else 0.0
        after = float(paragraph.space_after.pt) if paragraph.space_after else 0.0
        spacing = paragraph.line_spacing
        line_height = (spacing.pt if hasattr(spacing, "pt") else font_size * spacing if spacing else font_size * 1.22) if fixed_lines else font_size * 1.22
        required_height += line_count * line_height + before + after
    return {
        # Keep only a small allowance for font-metric variance.  The previous
        # 18% tolerance let text visibly escape its semantic card while still
        # remaining inside the slide canvas (for example, dense two-column
        # bodies rendered by LibreOffice).
        "overflow": horizontal_overflow or required_height > max(height_pt * 1.02, height_pt + 2.0),
        "horizontal_overflow": horizontal_overflow,
        "required_height_pt": round(required_height, 2),
        "available_height_pt": round(height_pt, 2),
        "minimum_font_size_pt": 18.0 if minimum_size == 10**9 else minimum_size,
        "maximum_wrapped_lines": maximum_lines,
    }


def _table_cell_audits(shape: Any) -> list[dict[str, Any]]:
    if not getattr(shape, "has_table", False):
        return []
    table = shape.table
    audits: list[dict[str, Any]] = []
    for row_index, row in enumerate(table.rows):
        for column_index, cell in enumerate(row.cells):
            if not str(cell.text or "").strip():
                continue
            width_pt = max(
                1.0,
                (
                    int(table.columns[column_index].width)
                    - int(cell.margin_left or 0)
                    - int(cell.margin_right or 0)
                ) / 12700,
            )
            height_pt = max(
                1.0,
                (
                    int(row.height)
                    - int(cell.margin_top or 0)
                    - int(cell.margin_bottom or 0)
                ) / 12700,
            )
            required_height = 0.0
            minimum_size = 10**9
            maximum_lines = 1
            for paragraph in cell.text_frame.paragraphs:
                font_size = _paragraph_font_size_pt(paragraph)
                minimum_size = min(
                    minimum_size,
                    _paragraph_minimum_font_size_pt(paragraph),
                )
                line_count = _wrapped_line_count(
                    paragraph.text or "",
                    width_pt=width_pt,
                    font_size_pt=font_size,
                )
                maximum_lines = max(maximum_lines, line_count)
                before = (
                    float(paragraph.space_before.pt)
                    if paragraph.space_before
                    else 0.0
                )
                after = (
                    float(paragraph.space_after.pt)
                    if paragraph.space_after
                    else 0.0
                )
                required_height += line_count * font_size * 1.22 + before + after
            if required_height > max(height_pt * 1.02, height_pt + 2.0):
                audits.append({
                    "severity": "critical",
                    "code": "exported_table_cell_overflow",
                    "shape_name": str(shape.name or ""),
                    "row": row_index + 1,
                    "column": column_index + 1,
                    "required_height_pt": round(required_height, 2),
                    "available_height_pt": round(height_pt, 2),
                    "minimum_font_size_pt": (
                        18.0 if minimum_size == 10**9 else minimum_size
                    ),
                    "maximum_wrapped_lines": maximum_lines,
                })
    return audits


def _normalized_ocr_text(value: object) -> str:
    return re.sub(r"[^0-9a-zA-Z\u3400-\u9fff]+", "", str(value or "")).lower()


@lru_cache(maxsize=1)
def _rapidocr_engine() -> Any:
    from rapidocr_onnxruntime import RapidOCR

    return RapidOCR()


def _rapidocr_page_text(path: Path) -> str:
    result, _ = _rapidocr_engine()(str(path))
    return " ".join(
        str(item[1])
        for item in result or []
        if isinstance(item, (list, tuple)) and len(item) >= 2
    )


def _ocr_character_recall(expected: str, recognized: str) -> float:
    """Measure visible character coverage without assuming one-column OCR order."""

    if not expected:
        return 1.0
    expected_counts = Counter(expected)
    recognized_counts = Counter(recognized)
    matched = sum(
        min(count, recognized_counts[character])
        for character, count in expected_counts.items()
    )
    return matched / len(expected)


def audit_rendered_slide_images(
    presentation: Any,
    image_paths: list[Path],
    *,
    ocr_runner: Any | None = None,
) -> dict[str, Any]:
    """Compare OCR from every rendered page with the PPTX's expected visible text."""
    issues: list[dict[str, Any]] = []
    expected_count = len(presentation.slides)
    if len(image_paths) != expected_count:
        issues.append({
            "severity": "critical",
            "code": "rendered_page_count_mismatch",
            "expected": expected_count,
            "actual": len(image_paths),
        })
    runner = ocr_runner or _rapidocr_page_text
    checked_pages = 0
    for page_number, (slide, image_path) in enumerate(
        zip(presentation.slides, image_paths),
        start=1,
    ):
        expected = _normalized_ocr_text(" ".join(
            str(shape.text or "")
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and str(shape.text or "").strip()
        ))
        if len(expected) < 8:
            continue
        checked_pages += 1
        try:
            recognized = _normalized_ocr_text(runner(image_path))
        except Exception as exc:
            issues.append({
                "severity": "critical",
                "code": "rendered_page_ocr_failed",
                "page": page_number,
                "message": str(exc),
            })
            continue
        ordered_coverage = SequenceMatcher(
            None,
            expected,
            recognized,
            autojunk=False,
        ).ratio()
        character_recall = _ocr_character_recall(expected, recognized)
        coverage = max(ordered_coverage, character_recall)
        if coverage < 0.68:
            issues.append({
                "severity": "critical",
                "code": "exported_ocr_text_missing_or_clipped",
                "page": page_number,
                "coverage": round(coverage, 4),
                "ordered_coverage": round(ordered_coverage, 4),
                "character_recall": round(character_recall, 4),
                "expected_character_count": len(expected),
                "recognized_character_count": len(recognized),
            })
    blockers = [item for item in issues if item["severity"] == "critical"]
    return {
        "passed": not blockers,
        "page_count": expected_count,
        "checked_pages": checked_pages,
        "issues": issues,
        "blockers": blockers,
    }


def _libreoffice_render_audit(
    path: Path,
    presentation: Any,
) -> dict[str, Any]:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    pdftoppm = shutil.which("pdftoppm")
    if not soffice or not pdftoppm:
        issue = {
            "severity": "critical",
            "code": "libreoffice_render_audit_unavailable",
            "missing": [
                name
                for name, executable in (("libreoffice", soffice), ("pdftoppm", pdftoppm))
                if not executable
            ],
        }
        return {
            "passed": False,
            "page_count": len(presentation.slides),
            "checked_pages": 0,
            "issues": [issue],
            "blockers": [issue],
        }
    with tempfile.TemporaryDirectory(prefix="lingzhi-slide-pixel-audit-") as temp_dir:
        output_dir = Path(temp_dir)
        fontconfig_path = _write_libreoffice_fontconfig(output_dir)
        soffice_env = os.environ.copy()
        soffice_env["FONTCONFIG_FILE"] = str(fontconfig_path)
        soffice_env["FONTCONFIG_PATH"] = str(fontconfig_path.parent)
        soffice_env["XDG_CACHE_HOME"] = str(output_dir / "font-cache")
        subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(output_dir),
                str(path.resolve()),
            ],
            check=True,
            capture_output=True,
            timeout=90,
            env=soffice_env,
        )
        pdf_path = output_dir / f"{path.stem}.pdf"
        if not pdf_path.is_file():
            raise RuntimeError("LibreOffice did not produce the expected PDF")
        prefix = output_dir / "page"
        subprocess.run(
            [pdftoppm, "-png", "-r", "160", str(pdf_path), str(prefix)],
            check=True,
            capture_output=True,
            timeout=90,
        )
        image_paths = sorted(
            output_dir.glob("page-*.png"),
            key=lambda item: int(item.stem.rsplit("-", 1)[-1]),
        )
        return audit_rendered_slide_images(presentation, image_paths)


def _write_libreoffice_fontconfig(output_dir: Path) -> Path:
    """Give the isolated renderer a real CJK fallback without changing PPTX fonts."""

    from xml.sax.saxutils import escape

    project_font_dir = Path(__file__).resolve().parent / "assets" / "fonts"
    portable_cjk_font = project_font_dir / "WenQuanYiMicroHei.ttc"
    if not portable_cjk_font.is_file():
        raise RuntimeError(
            "LibreOffice pixel audit requires the bundled CJK fallback font"
        )

    candidates = [
        Path(__file__).resolve().parents[1] / "frontend/public/presentation-assets/fonts",
        project_font_dir,
        Path("/System/Library/Fonts"),
        Path("/System/Library/Fonts/Supplemental"),
        Path("/Library/Fonts"),
        Path.home() / "Library" / "Fonts",
        Path("/usr/share/fonts"),
        Path("/usr/local/share/fonts"),
    ]
    configured = os.getenv("SLIDE_AUDIT_FONT_DIRS", "")
    candidates.extend(
        Path(item).expanduser()
        for item in configured.split(os.pathsep)
        if item.strip()
    )
    font_dirs: list[Path] = []
    seen: set[str] = set()
    for candidate in candidates:
        if not candidate.is_dir():
            continue
        resolved = str(candidate.resolve())
        if resolved in seen:
            continue
        seen.add(resolved)
        font_dirs.append(Path(resolved))

    cache_dir = output_dir / "font-cache"
    cache_dir.mkdir(parents=True, exist_ok=True)
    fontconfig_path = output_dir / "fonts.conf"
    directories = "\n".join(
        f"  <dir>{escape(str(directory))}</dir>" for directory in font_dirs
    )
    fontconfig_path.write_text(
        "\n".join([
            '<?xml version="1.0"?>',
            '<!DOCTYPE fontconfig SYSTEM "urn:fontconfig:fonts.dtd">',
            "<fontconfig>",
            directories,
            f"  <cachedir>{escape(str(cache_dir))}</cachedir>",
            "</fontconfig>",
            "",
        ]),
        encoding="utf-8",
    )
    return fontconfig_path


def audit_exported_pptx(
    path: str | Path,
    *,
    expected_slide_count: int | None = None,
    require_pixel_audit: bool | None = None,
) -> dict[str, Any]:
    """Audit exported objects, then optionally render and OCR every page."""
    from pptx import Presentation

    presentation = Presentation(path)
    issues: list[dict[str, Any]] = []
    if expected_slide_count is not None and len(presentation.slides) != expected_slide_count:
        issues.append({
            "severity": "critical",
            "code": "exported_slide_count_mismatch",
            "expected": expected_slide_count,
            "actual": len(presentation.slides),
        })
    if round(presentation.slide_width / presentation.slide_height, 3) != round(16 / 9, 3):
        issues.append({
            "severity": "critical",
            "code": "exported_aspect_ratio_invalid",
        })
    for slide_index, slide in enumerate(presentation.slides, start=1):
        visible_object_count = 0
        text_shapes: list[Any] = []
        for shape in slide.shapes:
            left = int(shape.left)
            top = int(shape.top)
            right = left + int(shape.width)
            bottom = top + int(shape.height)
            if left < 0 or top < 0 or right > presentation.slide_width or bottom > presentation.slide_height:
                issues.append({
                    "severity": "critical",
                    "code": "exported_object_out_of_bounds",
                    "page": slide_index,
                    "shape_name": str(shape.name or ""),
                })
            if int(shape.width) > 0 and int(shape.height) > 0:
                visible_object_count += 1
            for table_issue in _table_cell_audits(shape):
                issues.append({"page": slide_index, **table_issue})
            if getattr(shape, "has_text_frame", False) and str(shape.text or "").strip():
                text_shapes.append(shape)
                raw_text = str(shape.text or "")
                if re.search(
                    r"(?:\\(?:begin|end)\{?(?:array|[bp]?matrix)|"
                    r"\\(?:frac|mathbf|mathbb|leftarrow|rightarrow|xrightarrow)\b|"
                    r"\b(?:beginarray|endarray|mathbf[a-z]|frac\d+|rightarrow)\b)",
                    raw_text,
                    re.IGNORECASE,
                ):
                    issues.append({
                        "severity": "critical",
                        "code": "exported_raw_latex_visible",
                        "page": slide_index,
                        "shape_name": str(shape.name or ""),
                    })
                if "⁄" in raw_text or (
                    any(symbol in raw_text for symbol in "⎡⎢⎣⎤⎥⎦")
                    and any(symbol in raw_text for symbol in "()⎛⎜⎝⎞⎟⎠")
                ):
                    issues.append({
                        "severity": "critical",
                        "code": "exported_formula_glyph_not_portable",
                        "page": slide_index,
                        "shape_name": str(shape.name or ""),
                    })
                text_audit = _text_frame_audit(shape)
                top_inches = int(shape.top) / 914400
                bottom_inches = (int(shape.top) + int(shape.height)) / 914400
                is_footer = top_inches >= 6.9
                is_eyebrow = top_inches < 0.62 and int(shape.height) / 914400 < 0.5
                is_title = (
                    0.6 <= top_inches < 1.95
                    and text_audit["minimum_font_size_pt"] >= 28
                )
                title_metric_variance_fits = bool(
                    is_title
                    and not text_audit.get("horizontal_overflow")
                    and text_audit["required_height_pt"]
                    <= max(
                        text_audit["available_height_pt"] * 1.06,
                        text_audit["available_height_pt"] + 4.0,
                    )
                )
                if text_audit["overflow"] and not title_metric_variance_fits:
                    issues.append({
                        "severity": "critical",
                        "code": "exported_text_frame_overflow",
                        "page": slide_index,
                        "shape_name": str(shape.name or ""),
                        **text_audit,
                    })
                body_line_match = re.search(
                    r"\[v6-body-max-lines=(\d+)\]",
                    str(shape.name or ""),
                )
                if (
                    body_line_match
                    and text_audit["maximum_wrapped_lines"]
                    > max(1, int(body_line_match.group(1)))
                ):
                    issues.append({
                        "severity": "critical",
                        "code": "exported_body_capacity_exceeded",
                        "page": slide_index,
                        "shape_name": str(shape.name or ""),
                        "maximum_wrapped_lines": text_audit[
                            "maximum_wrapped_lines"
                        ],
                        "allowed_wrapped_lines": max(
                            1,
                            int(body_line_match.group(1)),
                        ),
                    })
                title_line_match = re.search(
                    r"\[v6-title-max-lines=(\d+)\]",
                    str(shape.name or ""),
                )
                title_line_limit = (
                    max(1, int(title_line_match.group(1)))
                    if title_line_match
                    else 1
                )
                if (
                    is_title
                    and text_audit["maximum_wrapped_lines"] > title_line_limit
                ):
                    issues.append({
                        "severity": "critical",
                        "code": "exported_title_unexpected_wrap",
                        "page": slide_index,
                        "shape_name": str(shape.name or ""),
                        "maximum_wrapped_lines": text_audit["maximum_wrapped_lines"],
                        "allowed_wrapped_lines": title_line_limit,
                    })
                if is_title:
                    frame = shape.text_frame
                    width_pt = max(
                        1.0,
                        (
                            int(shape.width)
                            - int(frame.margin_left or 0)
                            - int(frame.margin_right or 0)
                        ) / 12700,
                    )
                    if any(
                        _wrapped_text_has_orphan_last_line(
                            paragraph.text or "",
                            width_pt=width_pt,
                            font_size_pt=_paragraph_font_size_pt(paragraph),
                        )
                        for paragraph in frame.paragraphs
                    ):
                        issues.append({
                            "severity": "critical",
                            "code": "exported_title_orphan_line",
                            "page": slide_index,
                            "shape_name": str(shape.name or ""),
                        })
                if (
                    not is_footer
                    and not is_eyebrow
                    and not is_title
                    and top_inches >= 1.9
                    and bottom_inches <= 7.05
                    and text_audit["minimum_font_size_pt"] < 16
                    and (
                        len(re.sub(r"\s+", "", str(shape.text or ""))) >= 20
                        or int(shape.height) / 914400 > 0.45
                    )
                ):
                    issues.append({
                        "severity": "critical",
                        "code": "exported_body_font_below_16pt",
                        "page": slide_index,
                        "shape_name": str(shape.name or ""),
                        "minimum_font_size_pt": text_audit["minimum_font_size_pt"],
                    })
            if getattr(shape, "shape_type", None) == 13:
                crop_total = sum(
                    float(getattr(shape, name, 0) or 0)
                    for name in ("crop_left", "crop_right", "crop_top", "crop_bottom")
                )
                if crop_total > 0.35:
                    issues.append({
                        "severity": "critical",
                        "code": "exported_image_subject_overcropped",
                        "page": slide_index,
                        "shape_name": str(shape.name or ""),
                        "crop_total": round(crop_total, 4),
                    })
        for left_index, left_shape in enumerate(text_shapes):
            for right_shape in text_shapes[left_index + 1:]:
                intersection_width = max(
                    0,
                    min(
                        int(left_shape.left) + int(left_shape.width),
                        int(right_shape.left) + int(right_shape.width),
                    ) - max(int(left_shape.left), int(right_shape.left)),
                )
                intersection_height = max(
                    0,
                    min(
                        int(left_shape.top) + int(left_shape.height),
                        int(right_shape.top) + int(right_shape.height),
                    ) - max(int(left_shape.top), int(right_shape.top)),
                )
                intersection = intersection_width * intersection_height
                smaller = min(
                    int(left_shape.width) * int(left_shape.height),
                    int(right_shape.width) * int(right_shape.height),
                )
                if not smaller or intersection / smaller <= 0.12:
                    continue
                left_top = int(left_shape.top) / 914400
                right_top = int(right_shape.top) / 914400
                code = (
                    "exported_footer_overlap"
                    if max(left_top, right_top) >= 6.9
                    else "exported_text_overlap"
                )
                issues.append({
                    "severity": "critical",
                    "code": code,
                    "page": slide_index,
                    "shape_names": [
                        str(left_shape.name or ""),
                        str(right_shape.name or ""),
                    ],
                })
        if visible_object_count == 0:
            issues.append({
                "severity": "critical",
                "code": "exported_page_has_no_objects",
                "page": slide_index,
            })
    pixel_audit: dict[str, Any] | None = None
    pixel_audit_enabled = (
        require_pixel_audit
        if require_pixel_audit is not None
        else os.getenv("SLIDE_LIBREOFFICE_AUDIT_ENABLED", "").strip().lower()
        in {"1", "true", "yes", "on"}
    )
    if pixel_audit_enabled:
        try:
            pixel_audit = _libreoffice_render_audit(Path(path), presentation)
        except Exception as exc:
            pixel_audit = {
                "passed": False,
                "page_count": len(presentation.slides),
                "checked_pages": 0,
                "issues": [{
                    "severity": "critical",
                    "code": "libreoffice_render_audit_failed",
                    "message": str(exc),
                }],
                "blockers": [{
                    "severity": "critical",
                    "code": "libreoffice_render_audit_failed",
                    "message": str(exc),
                }],
            }
        issues.extend(pixel_audit.get("issues") or [])
    blockers = [item for item in issues if item["severity"] == "critical"]
    return {
        "schema_version": "slide_render_review_v1",
        "reviewer": (
            "post_export_pptx_object_and_ocr_audit"
            if pixel_audit is not None
            else "post_export_pptx_object_audit"
        ),
        "passed": not blockers,
        "page_count": len(presentation.slides),
        "issues": issues,
        "blockers": blockers,
        "repair_attempts": 0,
        "pixel_audit": pixel_audit,
    }


def validate_theme(theme: str | dict[str, Any]) -> dict[str, Any]:
    if isinstance(theme, dict):
        required = {"surface", "title", "body_font", "title_font"}
        missing = sorted(required.difference(theme))
        if missing:
            raise ValueError(
                "Compiled slide theme is missing required tokens: "
                + ", ".join(missing)
            )
        return deepcopy(theme)
    try:
        return THEMES[theme]
    except KeyError as exc:
        choices = ", ".join(sorted(THEMES))
        raise ValueError(f"Unknown slide theme '{theme}'. Expected one of: {choices}") from exc


V5_LAYOUT_RENDERER_NAMES = {
    "cover-minimal": "_render_cover_minimal",
    "cover-editorial": "_render_cover_editorial",
    "agenda-linear": "_render_agenda_linear",
    "chapter-entry": "_render_chapter",
    "hero-claim": "_render_claim_only",
    "editorial-body": "_render_editorial_body",
    "balanced-two-column": "_render_two_column",
    "classification-3": "_render_classification_three",
    "process-sequence": "_render_process",
    "formula-explanation": "_render_editorial_body",
    "code": "_render_code",
    "figure-text": "_render_visual_directed",
    "diagram-full": "_render_visual_directed",
    "worked-example": "_render_worked_example",
    "parallel-examples": "_render_parallel_examples",
    "question-prompt": "_render_question_prompt",
    "practice-sequence": "_render_process",
    "practice-artifact": "_render_practice_artifact",
    "practice-feedback": "_render_practice_feedback",
    "chapter-recap": "_render_chapter_recap",
    "course-synthesis": "_render_course_synthesis",
}

_VISUAL_DIRECTED_V5_LAYOUTS = {
    "figure-text",
    "diagram-full",
    "formula-explanation",
}


def _uses_visual_directed_renderer(
    unit: SlideSpec,
    resolved_layout: str,
) -> bool:
    if not unit.visuals:
        return False
    if resolved_layout in _VISUAL_DIRECTED_V5_LAYOUTS:
        return True
    # Legacy decks predate the final V5 layout contract. Preserve their
    # visual-first behavior while allowing every V5 semantic composition to
    # remain authoritative over an optional visual anchor.
    return resolved_layout not in V5_LAYOUT_RENDERER_NAMES


def _render_slide(
    slide: Any,
    unit: SlideSpec,
    page_number: int,
    page_count: int,
    theme: dict[str, str],
    asset_repository: SlideAssetRepository,
) -> None:
    _fill_background(slide, theme["surface"])
    resolved_layout = str(
        unit.quality.get("resolved_layout")
        or unit.quality.get("requested_layout")
        or unit.layout
    )
    if resolved_layout not in {
        "cover",
        "cover-minimal",
        "cover-editorial",
        "chapter",
        "chapter-entry",
        "recap",
        "chapter-recap",
        "course-synthesis",
    }:
        _add_theme_page_background(slide, unit, theme, resolved_layout)
    if (
        _uses_visual_directed_renderer(unit, resolved_layout)
        and resolved_layout not in {
        "cover",
        "cover-minimal",
        "cover-editorial",
        "roadmap",
        "agenda-linear",
        "chapter",
        "chapter-entry",
        "recap",
        "chapter-recap",
        "course-synthesis",
        "appendix",
        }
    ):
        _render_visual_directed(slide, unit, theme, asset_repository)
        _footer(slide, unit, page_number, page_count, theme)
        return
    renderer_name = V5_LAYOUT_RENDERER_NAMES.get(resolved_layout)
    renderer = globals().get(renderer_name) if renderer_name else None
    renderer = renderer or {
        "cover": _render_cover,
        "roadmap": _render_roadmap,
        "chapter": _render_chapter,
        "objective": _render_objective,
        "concept": _render_concept,
        "comparison": _render_comparison,
        "comparison-matrix": _render_comparison,
        "process": _render_process,
        "code": _render_code,
        "misconception": _render_misconception,
        "practice": _render_practice,
        "recap": _render_recap,
        "appendix": _render_appendix,
        "hero-statement": _render_hero_statement,
        "two-column": _render_two_column,
        "case-study": _render_case_study,
        "question": _render_practice,
        "summary": _render_recap,
    }.get(resolved_layout) or {
        "cover": _render_cover,
        "roadmap": _render_roadmap,
        "chapter": _render_chapter,
        "objective": _render_objective,
        "concept": _render_concept,
        "comparison": _render_comparison,
        "process": _render_process,
        "code": _render_code,
        "misconception": _render_misconception,
        "practice": _render_practice,
        "recap": _render_recap,
        "appendix": _render_appendix,
    }.get(unit.layout, _render_concept)
    renderer(slide, unit, theme)
    if unit.layout != "cover":
        _footer(slide, unit, page_number, page_count, theme)


def _render_visual_directed(
    slide: Any,
    unit: SlideSpec,
    theme: dict[str, str],
    asset_repository: SlideAssetRepository,
) -> None:
    """Render one dominant visual and the complete source body on a flat canvas."""
    visual = dict(unit.visuals[0])
    kind = str(visual.get("kind") or "none")
    _heading(slide, unit, theme)
    if kind in {"source_image", "retrieved_image", "generated_illustration"}:
        if _render_image_visual(
            slide,
            unit,
            visual,
            theme,
            asset_repository,
        ):
            return
    if kind in {"relational_diagram", "rule_diagram"}:
        _render_relational_visual(slide, unit, visual, theme)
        return
    if kind == "coordinate_plot":
        _render_coordinate_visual(slide, unit, visual, theme)
        return
    if kind == "chart":
        _render_chart_visual(slide, unit, visual, theme)
        return
    if kind == "table":
        _render_table_visual(slide, unit, visual, theme)
        return
    if kind == "formula":
        _render_formula_visual(slide, unit, visual, theme)
        return
    if kind == "code":
        _render_code_visual(slide, unit, visual, theme)
        return
    if not _visible_source_text(unit):
        _render_navigation_statement(slide, unit, theme, heading_already_rendered=True)
        return
    _render_editorial_body(
        slide,
        unit,
        theme,
        heading_already_rendered=True,
    )


def _render_relational_visual(
    slide: Any,
    unit: SlideSpec,
    visual: dict[str, Any],
    theme: dict[str, str],
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.util import Inches

    nodes = list(visual.get("nodes") or [])[:6]
    edges = list(visual.get("edges") or [])[:10]
    composition = unit.composition or "split-visual"
    visual_left = composition in {"figure-first", "diagram-full"}
    diagram_x = 0.78 if visual_left else 7.0
    text_x = 7.22 if visual_left else 0.78
    diagram_w = 5.95
    text_w = 5.32
    _source_panel(slide, unit, text_x, 1.92, text_w, 4.62, theme)
    _shape(
        slide,
        diagram_x,
        1.92,
        diagram_w,
        4.62,
        theme["canvas"],
        radius=True,
        line=theme["chart_bg"],
    )
    caption = _visual_caption(visual)
    if " ".join(caption.split()).casefold() != " ".join(unit.title.split()).casefold():
        _text(
            slide,
            caption,
            diagram_x + 0.32,
            2.12,
            diagram_w - 0.64,
            0.42,
            16,
            theme["accent"],
            bold=True,
        )
    if not nodes:
        return
    direction = str(
        (visual.get("parameters") or {}).get("direction") or "vertical"
    )
    node_metrics = diagram_node_layout_metrics(
        [str(node.get("label") or "").strip() for node in nodes],
        direction=direction,
    )
    if not node_metrics["fits"]:
        raise ValueError("diagram_node_render_capacity_exceeded")
    positions: dict[str, tuple[float, float, float, float]] = {}
    for node, box in zip(nodes, node_metrics["node_boxes"]):
        positions[str(node.get("node_id"))] = (
            diagram_x + box["x"],
            1.92 + box["y"],
            box["width"],
            box["height"],
        )

    # Connectors are deliberately created first so they remain behind nodes.
    edge_labels: list[tuple[str, float, float]] = []
    for edge in edges:
        source = positions.get(str(edge.get("source") or ""))
        target = positions.get(str(edge.get("target") or ""))
        if not source or not target:
            continue
        x1 = source[0] + source[2] / 2
        y1 = source[1] + source[3] / 2
        x2 = target[0] + target[2] / 2
        y2 = target[1] + target[3] / 2
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x1),
            Inches(y1),
            Inches(x2),
            Inches(y2),
        )
        connector.line.color.rgb = RGBColor.from_string(theme["muted"])
        connector.line.width = Inches(0.018)
        edge_label = str(edge.get("label") or "").strip()
        if edge_label:
            edge_labels.append((edge_label, (x1 + x2) / 2, (y1 + y2) / 2))

    for edge_label, x, y in edge_labels:
        _text(
            slide,
            _diagram_label(edge_label),
            x - 0.68,
            y - 0.22,
            1.36,
            0.4,
            9,
            theme["muted"],
            bold=True,
            align="center",
        )

    for node in nodes:
        node_id = str(node.get("node_id") or "")
        x, y, width, height = positions[node_id]
        primary = str(node.get("emphasis") or "") == "primary"
        full_label = str(node.get("label") or "")
        visible_label = " ".join(full_label.split())
        label_size = 16
        shape = _shape(
            slide,
            x,
            y,
            width,
            height,
            theme["accent_soft"] if primary else theme["surface"],
            radius=True,
            line=theme["accent"] if primary else theme["chart_bg"],
        )
        _set_alt_text(shape, full_label)
        _text(
            slide,
            visible_label,
            x + 0.15,
            y + 0.09,
            width - 0.3,
            height - 0.16,
            label_size,
            theme["title"] if primary else theme["ink"],
            bold=primary or len(visible_label) < 18,
            align="center",
        )


def _formula_font_size(value: str, *, width_inches: float) -> int:
    """Keep short editable formulae prominent without overflowing long chains."""

    lines = [line.strip() for line in str(value or "").splitlines() if line.strip()]
    if not lines:
        return 25
    line_count = len(lines)
    longest_line = max(len(line) for line in lines)
    if line_count <= 3 and longest_line <= 32:
        return 32 if width_inches >= 9 else 30
    if line_count <= 4 and longest_line <= 40:
        return 29 if width_inches >= 6 else 27
    if line_count <= 6 and longest_line <= 54:
        return 25
    return 22


def _render_formula_visual(
    slide: Any,
    unit: SlideSpec,
    visual: dict[str, Any],
    theme: dict[str, str],
) -> None:
    formula = str((visual.get("parameters") or {}).get("formula") or "").strip()
    if not formula:
        formula = next(
            (
                block.content
                for block in unit.blocks
                if block.metadata.get("formula")
            ),
            "",
        )
    if not formula:
        _render_editorial_body(
            slide,
            unit,
            theme,
            heading_already_rendered=True,
        )
        return
    supporting_blocks = [
        block
        for block in unit.blocks
        if not block.metadata.get("formula")
    ]
    has_supporting_copy = bool(supporting_blocks)
    formula_width = 7.0 if has_supporting_copy else 11.2
    _shape(slide, 0.88, 2.05, 0.055, 3.95, theme["accent"], radius=False)
    _text(slide, "公式与推导", 1.18, 2.03, 1.5, 0.3, 11, theme["accent"], bold=True)
    display_formula = _format_formula_text(formula)
    display_formula = re.sub(r"\n\s*\n", "\n", display_formula)
    formula_size = _formula_font_size(
        display_formula,
        width_inches=formula_width - 0.42,
    )
    _text(
        slide,
        display_formula,
        1.18,
        2.62 if "\n" in display_formula else 2.94,
        formula_width - 0.42,
        3.05 if "\n" in display_formula else 1.7,
        formula_size,
        theme["title"],
        bold=False,
        align="center",
        font=theme["math_font"],
        east_asian_font=theme["body_east_asian_font"],
    )
    _shape(
        slide,
        1.18,
        6.08,
        formula_width - 0.42,
        0.025,
        theme["chart_bg"],
        radius=False,
    )
    if has_supporting_copy:
        supporting = SlideSpec.model_validate({
            **unit.model_dump(mode="json"),
            "blocks": [block.model_dump(mode="json") for block in supporting_blocks],
        })
        supporting_body = _visible_source_text(supporting)
        _shape(slide, 8.28, 2.05, 0.025, 4.03, theme["chart_bg"], radius=False)
        _text(
            slide,
            "条件与结论",
            8.62,
            2.03,
            1.7,
            0.3,
            11,
            theme["muted"],
            bold=True,
        )
        _text(
            slide,
            supporting_body,
            8.62,
            2.58,
            3.72,
            3.28,
            18 if len(supporting_body) <= 150 else 16,
            theme["ink"],
            font=theme["body_font"],
            east_asian_font=theme["body_east_asian_font"],
        )


def _plain_formula(value: str) -> str:
    clean = str(value or "").strip()
    clean = re.sub(r"^\s*(?:\$\$|\\\[)", "", clean)
    clean = re.sub(r"(?:\$\$|\\\])\s*$", "", clean)
    # Inline math delimiters are transport markup, not classroom content.
    # Remove every unescaped delimiter so mixed prose/formula titles do not
    # leak raw ``$`` characters into the exported deck.
    clean = re.sub(r"(?<!\\)\$+", "", clean)
    return clean.strip()


_SUBSCRIPT = str.maketrans("0123456789+-=()aehijklmnoprstuvxy", "₀₁₂₃₄₅₆₇₈₉₊₋₌₍₎ₐₑₕᵢⱼₖₗₘₙₒₚᵣₛₜᵤᵥₓᵧ")
_SUPERSCRIPT = str.maketrans({
    **dict(zip("0123456789+-=()", "⁰¹²³⁴⁵⁶⁷⁸⁹⁺⁻⁼⁽⁾", strict=True)),
    **{
        "a": "ᵃ", "b": "ᵇ", "c": "ᶜ", "d": "ᵈ", "e": "ᵉ",
        "f": "ᶠ", "g": "ᵍ", "h": "ʰ", "i": "ⁱ", "j": "ʲ",
        "k": "ᵏ", "l": "ˡ", "m": "ᵐ", "n": "ⁿ", "o": "ᵒ",
        "p": "ᵖ", "r": "ʳ", "s": "ˢ", "t": "ᵗ", "u": "ᵘ",
        "v": "ᵛ", "w": "ʷ", "x": "ˣ", "y": "ʸ", "z": "ᶻ",
    },
})
_FORMULA_SYMBOLS = {
    r"\Longleftrightarrow": "⟺",
    r"\Leftrightarrow": "⇔",
    r"\leftrightarrow": "↔",
    r"\longrightarrow": "⟶",
    r"\rightarrow": "→",
    r"\longleftarrow": "⟵",
    r"\leftarrow": "←",
    r"\gets": "←",
    r"\to": "→",
    r"\circ": "∘",
    r"\mapsto": "↦",
    r"\subseteq": "⊆",
    r"\supseteq": "⊇",
    r"\mathbf": "",
    r"\boldsymbol": "",
    r"\operatorname": "",
    r"\mathrm": "",
    r"\mathbb": "",
    r"\text": "",
    r"\rank": "rank",
    r"\dim": "dim",
    r"\min": "min",
    r"\max": "max",
    r"\ker": "ker",
    r"\det": "det",
    r"\cdots": "⋯",
    r"\vdots": "⋮",
    r"\times": "×",
    r"\cdot": "·",
    r"\ll": "≪",
    r"\gg": "≫",
    r"\approx": "≈",
    r"\neq": "≠",
    r"\ne": "≠",
    r"\leq": "≤",
    r"\le": "≤",
    r"\geq": "≥",
    r"\ge": "≥",
    r"\equiv": "≡",
    r"\subset": "⊂",
    r"\supset": "⊃",
    r"\forall": "∀",
    r"\exists": "∃",
    r"\lambda": "λ",
    r"\alpha": "α",
    r"\beta": "β",
    r"\gamma": "γ",
    r"\theta": "θ",
    r"\varepsilon": "ε",
    r"\epsilon": "ε",
    r"\delta": "δ",
    r"\phi": "φ",
    r"\pi": "π",
    r"\sigma": "σ",
    r"\sum": "∑",
    r"\Sigma": "Σ",
    r"\Delta": "Δ",
    r"\Omega": "Ω",
    r"\prod": "∏",
    r"\int": "∫",
    r"\iint": "∬",
    r"\infty": "∞",
    r"\partial": "∂",
    r"\pm": "±",
    r"\mp": "∓",
    r"\lceil": "⌈",
    r"\rceil": "⌉",
    r"\lfloor": "⌊",
    r"\rfloor": "⌋",
    r"\limsup": "lim sup",
    r"\liminf": "lim inf",
    r"\lim": "lim",
    r"\ln": "ln",
    r"\log": "log",
    r"\exp": "exp",
    r"\sin": "sin",
    r"\cos": "cos",
    r"\tan": "tan",
    r"\arctan": "arctan",
    r"\cap": "∩",
    r"\cup": "∪",
    r"\land": "∧",
    r"\lor": "∨",
    r"\in": "∈",
    r"\mid": "∣",
    r"\qquad": "    ",
    r"\quad": "  ",
    r"\ ": " ",
    r"\ldots": "…",
    r"\,": "",
    r"\!": "",
    r"\left": "",
    r"\right": "",
}


def _latex_group(value: str, start: int) -> tuple[str, int] | None:
    """Read one balanced LaTeX group, including nested braces."""

    if start >= len(value) or value[start] != "{":
        return None
    depth = 0
    for index in range(start, len(value)):
        if value[index] == "{" and (index == 0 or value[index - 1] != "\\"):
            depth += 1
        elif value[index] == "}" and (index == 0 or value[index - 1] != "\\"):
            depth -= 1
            if depth == 0:
                return value[start + 1:index], index + 1
    return None


def _replace_group_command(
    value: str,
    command: str,
    arity: int,
    builder: Any,
) -> str:
    r"""Replace balanced commands such as nested ``\frac`` without truncation."""

    result = value
    offset = 0
    replacements = 0
    while replacements < 128:
        start = result.find(command, offset)
        if start < 0:
            break
        cursor = start + len(command)
        arguments: list[str] = []
        for _ in range(arity):
            while cursor < len(result) and result[cursor].isspace():
                cursor += 1
            group = _latex_group(result, cursor)
            if group is not None:
                argument, cursor = group
            elif cursor < len(result) and result[cursor].isalnum():
                argument = result[cursor]
                cursor += 1
            else:
                break
            arguments.append(argument)
        if len(arguments) != arity:
            offset = start + len(command)
            continue
        result = result[:start] + str(builder(*arguments)) + result[cursor:]
        offset = max(0, start)
        replacements += 1
    return result


def _format_formula_cases(body: str) -> str:
    """Render a LaTeX cases environment as an editable piecewise system."""

    rows = [
        _format_formula_text(row.replace("&", "  "))
        for row in re.split(r"\\\\", body)
        if row.strip()
    ]
    if not rows:
        return ""
    if len(rows) == 1:
        return f"⎧ {rows[0]}"
    return "\n".join(
        f"{('⎧' if index == 0 else '⎩' if index == len(rows) - 1 else '⎨')} {row}"
        for index, row in enumerate(rows)
    )


def _format_formula_text(value: str) -> str:
    """Compile common course LaTeX into portable, editable mathematical text."""
    expression = _plain_formula(value)
    cases_pattern = re.compile(
        r"\\begin\{cases\}(?P<body>.*?)\\end\{cases\}",
        re.DOTALL,
    )
    expression = cases_pattern.sub(
        lambda match: _format_formula_cases(match.group("body")),
        expression,
    )
    array_pattern = re.compile(
        r"(?:\\left\s*\[\s*)?"
        r"\\begin\{array\}\{(?P<columns>[^{}]+)\}"
        r"(?P<body>.*?)"
        r"\\end\{array\}"
        r"(?:\s*\\right\s*\])?",
        re.DOTALL,
    )
    expression = array_pattern.sub(
        lambda match: _format_formula_array(
            match.group("body"),
            match.group("columns"),
        ),
        expression,
    )
    expression = _replace_group_command(
        expression,
        r"\frac",
        2,
        lambda numerator, denominator: f"({numerator})/({denominator})",
    )
    expression = _replace_group_command(
        expression,
        r"\sqrt",
        1,
        lambda radicand: f"√({radicand})",
    )
    expression = _replace_group_command(
        expression,
        r"\xrightarrow",
        1,
        lambda label: f" ⟶[{_format_formula_text(label)}] ",
    )
    expression = re.sub(
        r"\\frac\s*([A-Za-z0-9])\s*([A-Za-z0-9])",
        lambda match: f"({match.group(1)})/({match.group(2)})",
        expression,
    )
    expression = re.sub(
        r"\\sqrt\s*([A-Za-z0-9])",
        lambda match: f"√{match.group(1)}",
        expression,
    )
    expression = re.sub(
        r"\\sum_\{([^{}]+)\}\^\{([^{}]+)\}",
        lambda match: (
            "∑"
            + _script_text(match.group(1), _SUBSCRIPT, "_")
            + _script_text(match.group(2), _SUPERSCRIPT, "^")
        ),
        expression,
    )
    expression = re.sub(
        r"\\sum_\{([^{}]+)\}\^([A-Za-z0-9])",
        lambda match: (
            "∑"
            + _script_text(match.group(1), _SUBSCRIPT, "_")
            + _script_text(match.group(2), _SUPERSCRIPT, "^")
        ),
        expression,
    )
    matrix_pattern = re.compile(
        r"\\begin\{(?P<kind>[bp]?matrix)\}(?P<body>.*?)\\end\{(?P=kind)\}",
        re.DOTALL,
    )
    expression = matrix_pattern.sub(
        lambda match: _format_formula_matrix(match.group("body"), match.group("kind")),
        expression,
    )
    # Some imported course documents JSON-escape LaTeX commands twice. Matrix
    # row separators have already been consumed above, so the remaining paired
    # slashes can safely be normalized back to one command introducer.
    expression = re.sub(r"\\\\(?=[A-Za-z{}])", r"\\", expression)
    expression = re.sub(
        r"\\vec\{([^{}]+)\}",
        lambda match: f"{match.group(1)}⃗",
        expression,
    )
    expression = re.sub(
        r"\\vec\s*([A-Za-z])",
        lambda match: f"{match.group(1)}⃗",
        expression,
    )
    expression = re.sub(
        r"\\sqrt\{([^{}]+)\}",
        lambda match: f"√({match.group(1)})",
        expression,
    )
    expression = re.sub(
        r"\\frac\{([^{}]+)\}\{([^{}]+)\}",
        lambda match: f"({match.group(1)})/({match.group(2)})",
        expression,
    )
    for command, symbol in sorted(_FORMULA_SYMBOLS.items(), key=lambda item: -len(item[0])):
        expression = expression.replace(command, symbol)
    # Recover unambiguous operator words when an upstream projection lost only
    # the LaTeX command introducer. This changes presentation, not source facts.
    expression = re.sub(r"(?<=\w)leftrightarrow(?=\w)", "↔", expression)
    expression = re.sub(r"(?<=\w)rightarrow(?=\w)", "→", expression)
    expression = re.sub(r"(?<=\w)leftarrow(?=\w)", "←", expression)
    expression = re.sub(
        r"(?<![A-Za-z])frac\s*([A-Za-z0-9])\s*([A-Za-z0-9])",
        lambda match: f"({match.group(1)})/({match.group(2)})",
        expression,
    )
    expression = expression.replace(r"\{", "⦃").replace(r"\}", "⦄")
    expression = re.sub(
        r"([∑Σ])_([A-Za-z0-9]+)=([A-Za-z0-9]+)\^([A-Za-z0-9]+)",
        lambda match: (
            "∑"
            + _script_text(f"{match.group(2)}={match.group(3)}", _SUBSCRIPT, "_")
            + _script_text(match.group(4), _SUPERSCRIPT, "^")
        ),
        expression,
    )
    expression = re.sub(
        r"_\s*\{([^{}]+)\}",
        lambda match: _script_text(match.group(1), _SUBSCRIPT, "_"),
        expression,
    )
    expression = re.sub(
        r"\^\s*\{([^{}]+)\}",
        lambda match: _script_text(match.group(1), _SUPERSCRIPT, "^"),
        expression,
    )
    expression = re.sub(
        r"_\s*([A-Za-z0-9+\-=()])",
        lambda match: _script_text(match.group(1), _SUBSCRIPT, "_"),
        expression,
    )
    expression = re.sub(
        r"\^\s*([A-Za-z0-9+\-=()])",
        lambda match: _script_text(match.group(1), _SUPERSCRIPT, "^"),
        expression,
    )
    expression = (
        expression.replace("{", "")
        .replace("}", "")
        .replace("⦃", "{")
        .replace("⦄", "}")
    )
    expression = re.sub(r"[ \t]+", " ", expression)
    expression = re.sub(r"\s*([≪≫])\s*", r" \1 ", expression)
    expression = re.sub(r"[ \t]+", " ", expression)
    expression = re.sub(r" *\n *", "\n", expression)
    return expression.strip(" ,")


def _display_text(value: str) -> str:
    """Format math markup for display while retaining the source text in the model."""
    text = str(value or "")
    if "\\" not in text and "$" not in text:
        return text
    return _format_formula_text(text)


def _script_text(value: str, translation: dict[int, str], marker: str) -> str:
    rendered = str(value).translate(translation)
    grouped = re.sub(r"\s*([∘→⇒⇔])\s*", r"\1", str(value).strip())
    if all(ord(character) in translation for character in str(value)):
        return rendered
    if len(grouped) == 1:
        return f"{marker}{grouped}"
    if marker == "_":
        return f"₍{grouped}₎"
    if marker == "^":
        return f"⁽{grouped}⁾"
    return f"{marker}({grouped})"


def _format_formula_matrix(body: str, kind: str) -> str:
    rows = [
        _format_formula_text(row.replace("&", "  "))
        for row in re.split(r"\\\\", body)
        if row.strip()
    ]
    if not rows:
        return ""
    brackets = ("(", ")") if kind == "pmatrix" else ("⎡", "⎤")
    if len(rows) == 1:
        return f"{brackets[0]} {rows[0]} {brackets[1]}"
    if kind == "pmatrix":
        left_brackets = ("⎛", "⎜", "⎝")
        right_brackets = ("⎞", "⎟", "⎠")
    else:
        left_brackets = ("⎡", "⎢", "⎣")
        right_brackets = ("⎤", "⎥", "⎦")
    lines: list[str] = []
    for index, row in enumerate(rows):
        position = 0 if index == 0 else 2 if index == len(rows) - 1 else 1
        left = left_brackets[position]
        right = right_brackets[position]
        lines.append(f"{left} {row} {right}")
    return "\n".join(lines)


def _format_formula_array(body: str, column_spec: str) -> str:
    """Render LaTeX arrays, including augmented-matrix column dividers."""

    column_tokens = [token for token in str(column_spec) if token in "lcr|pmb"]
    divider_after: set[int] = set()
    logical_column = 0
    for token in column_tokens:
        if token == "|":
            if logical_column:
                divider_after.add(logical_column)
        elif token in "lcrpmb":
            logical_column += 1

    rendered_rows: list[str] = []
    for raw_row in re.split(r"\\\\", body):
        if not raw_row.strip():
            continue
        cells = [_format_formula_text(cell) for cell in raw_row.split("&")]
        pieces: list[str] = []
        for index, cell in enumerate(cells, start=1):
            pieces.append(cell)
            if index in divider_after and index < len(cells):
                pieces.append("│")
        rendered_rows.append("  ".join(piece for piece in pieces if piece))
    if not rendered_rows:
        return ""
    if len(rendered_rows) == 1:
        return f"⎡ {rendered_rows[0]} ⎤"
    lines: list[str] = []
    for index, row in enumerate(rendered_rows):
        left = "⎡" if index == 0 else "⎣" if index == len(rendered_rows) - 1 else "⎢"
        right = "⎤" if index == 0 else "⎦" if index == len(rendered_rows) - 1 else "⎥"
        lines.append(f"{left} {row} {right}")
    return "\n".join(lines)


def _render_code_visual(
    slide: Any,
    unit: SlideSpec,
    visual: dict[str, Any],
    theme: dict[str, str],
) -> None:
    code = _find_block(unit, "code")
    _shape(slide, 0.78, 1.92, 7.65, 4.62, theme["code"], radius=True)
    _text(slide, "SOURCE", 1.12, 2.18, 1.2, 0.28, 11, "AEB6D0", bold=True, font=CODE_FONT)
    _text(
        slide,
        code.content if code else "",
        1.12,
        2.7,
        6.92,
        3.38,
        16,
        "F5F7FF",
        font=CODE_FONT,
        east_asian_font=theme["body_east_asian_font"],
        literal=True,
    )
    supporting = SlideSpec.model_validate({
        **unit.model_dump(mode="json"),
        "blocks": [
            block.model_dump(mode="json")
            for block in unit.blocks
            if block is not code
        ],
    })
    _source_panel(slide, supporting, 8.7, 1.92, 3.86, 4.62, theme)


def _render_table_row_detail(
    slide: Any,
    headers: list[str],
    row: list[str],
    theme: dict[str, str],
    support_text: str = "",
) -> None:
    """Render one oversized source row as readable labeled evidence fields."""

    pairs = [
        (
            headers[index] if index < len(headers) else f"字段 {index + 1}",
            value,
        )
        for index, value in enumerate(row)
    ]
    if not pairs:
        return
    available_height = 3.34 if support_text else 4.48
    gap = 0.08
    usable_height = available_height - gap * max(0, len(pairs) - 1)
    value_fonts = [17 if len(value) <= 72 else 16 for _label, value in pairs]
    required_heights = [
        max(
            0.58,
            _wrapped_line_count(
                _display_text(label),
                width_pt=2.15 * 72.0,
                font_size_pt=16,
            )
            * 16
            * 1.22
            / 72.0
            + 0.28,
            _wrapped_line_count(
                _display_text(value),
                width_pt=8.75 * 72.0,
                font_size_pt=value_font,
            )
            * value_font
            * 1.22
            / 72.0
            + 0.22,
        )
        for (label, value), value_font in zip(pairs, value_fonts)
    ]
    required_total = sum(required_heights)
    if required_total > usable_height:
        raise ValueError("table_row_detail_render_capacity_exceeded")
    extra_height = (usable_height - required_total) / max(1, len(pairs))
    heights = [height + extra_height for height in required_heights]
    y = 1.94
    for index, ((label, value), value_font, height) in enumerate(
        zip(pairs, value_fonts, heights)
    ):
        _shape(
            slide,
            0.8,
            y,
            11.76,
            height,
            theme["surface"],
            radius=True,
            line=theme["chart_bg"],
        )
        _shape(slide, 0.8, y, 0.08, height, theme["accent"], radius=False)
        _text(
            slide,
            label,
            1.1,
            y + 0.16,
            2.15,
            max(0.3, height - 0.28),
            16,
            theme["accent"],
            bold=True,
        )
        _text(
            slide,
            value,
            3.35,
            y + 0.12,
            8.75,
            max(0.36, height - 0.22),
            value_font,
            theme["ink"],
            bold=len(value) <= 48,
        )
        y += height + gap
    if support_text:
        support_y = 1.94 + available_height + 0.12
        _shape(
            slide,
            0.8,
            support_y,
            11.76,
            1.02,
            theme["accent_soft"],
            radius=True,
            line=theme["chart_bg"],
        )
        _shape(slide, 0.8, support_y, 0.08, 1.02, theme["accent"], radius=False)
        _text(
            slide,
            support_text,
            1.1,
            support_y + 0.13,
            11.08,
            0.74,
            16,
            theme["ink"],
        )


def _table_support_text(blocks: list[SlideBlockSpec]) -> str:
    return "\n".join(
        [
            *(
                item
                for block in blocks
                for item in block.items
                if str(item).strip()
            ),
            *(
                block.content
                for block in blocks
                if str(block.content or "").strip()
            ),
        ]
    ).strip()


def _required_table_height_inches(
    headers: list[str],
    rows: list[list[str]],
    *,
    width: float,
    font_size: int = 16,
    cell_horizontal_margin: float = 0.1,
    cell_vertical_margin: float = 0.07,
) -> float:
    column_count = max(1, len(headers), max((len(row) for row in rows), default=0))
    values = [headers or ["Comparison"], *rows]
    column_text_width_pt = max(
        12.0,
        (width / column_count - cell_horizontal_margin * 2) * 72.0,
    )
    required_height_pt = 0.0
    for row_values in values:
        maximum_lines = max(
            1,
            max(
                (
                    _wrapped_line_count(
                        _display_text(str(row_values[column_index])),
                        width_pt=column_text_width_pt,
                        font_size_pt=font_size,
                    )
                    for column_index in range(min(column_count, len(row_values)))
                ),
                default=1,
            ),
        )
        required_height_pt += (
            maximum_lines * font_size * 1.22 + cell_vertical_margin * 2 * 72.0
        )
    return required_height_pt / 72.0


def _render_table_visual(
    slide: Any,
    unit: SlideSpec,
    visual: dict[str, Any],
    theme: dict[str, str],
) -> None:
    parameters = visual.get("parameters") or {}
    parameter_rows = parameters.get("rows") or []
    if parameter_rows:
        supporting_blocks = [
            block for block in unit.blocks
            if not block.metadata.get("table_source")
        ]
        if (
            str(unit.quality.get("v6_layout_variant") or "")
            == "table-row-detail"
            and len(parameter_rows) == 1
        ):
            _render_table_row_detail(
                slide,
                [str(value) for value in parameters.get("headers") or []],
                [str(value) for value in parameter_rows[0]],
                theme,
                _table_support_text(supporting_blocks),
            )
            return
        support_mode = str(unit.quality.get("v6_artifact_support_mode") or "")
        split_support = bool(supporting_blocks and support_mode == "split")
        band_support = bool(supporting_blocks and support_mode == "band")
        table_width = 7.18 if split_support else 11.78
        table_height = 3.28 if band_support else 4.62
        support_band_y = 5.30
        support_band_height = 1.58
        if band_support:
            required_table_height = _required_table_height_inches(
                [str(value) for value in parameters.get("headers") or []],
                [[str(value) for value in row] for row in parameter_rows],
                width=table_width,
            )
            # The template's summary band has more vertical capacity than most
            # summaries need.  Lend only that verified slack to the table while
            # preserving a gap and the fixed footer-safe lower boundary.
            support_text = _table_support_text(supporting_blocks)
            required_support_height = max(
                0.72,
                _wrapped_line_count(
                    support_text,
                    width_pt=9.55 * 72.0,
                    font_size_pt=16,
                ) * 16 * 1.22 / 72.0 + 0.20,
            )
            support_band_height = min(1.58, max(0.78, required_support_height))
            support_band_y = 6.88 - support_band_height
            table_height = min(
                required_table_height + 0.02,
                support_band_y - 1.92 - 0.10,
            )
        _table(
            slide,
            [str(value) for value in parameters.get("headers") or ["顺序", "课程原文要点"]],
            [[str(value) for value in row] for row in parameter_rows],
            0.78,
            1.92,
            table_width,
            table_height,
            theme,
        )
        if split_support:
            supporting = SlideSpec.model_validate({
                **unit.model_dump(mode="json"),
                "blocks": [
                    block.model_dump(mode="json") for block in supporting_blocks
                ],
            })
            _source_panel(slide, supporting, 8.2, 1.92, 4.36, 4.62, theme)
        elif band_support:
            support_text = _table_support_text(supporting_blocks)
            support_label = next(
                (
                    str(block.title or "").strip()
                    for block in supporting_blocks
                    if str(block.title or "").strip()
                ),
                "" if _uses_source_only_audience_labels(unit) else "SUMMARY",
            )
            _shape(
                slide,
                0.78,
                support_band_y,
                11.78,
                support_band_height,
                theme["surface"],
                radius=True,
                line=theme["chart_bg"],
            )
            if support_label:
                _text(
                    slide,
                    support_label.upper(),
                    1.08,
                    support_band_y + support_band_height / 2 - 0.14,
                    1.45,
                    0.28,
                    10,
                    theme["accent"],
                    bold=True,
                    font=theme["body_font"],
                    east_asian_font=theme["body_east_asian_font"],
                )
            _text(
                slide,
                support_text,
                2.6 if support_label else 1.08,
                support_band_y + 0.10,
                9.55 if support_label else 11.08,
                support_band_height - 0.18,
                16,
                theme["ink"],
                font=theme["body_font"],
                east_asian_font=theme["body_east_asian_font"],
            )
        return
    block = next(
        (
            item for item in unit.blocks
            if item.metadata.get("rows")
        ),
        None,
    )
    if block is None:
        _source_panel(slide, unit, 0.78, 1.92, 11.78, 4.62, theme)
        return
    _table(
        slide,
        [str(value) for value in block.metadata.get("headers") or []],
        [[str(value) for value in row] for row in block.metadata.get("rows") or []],
        0.78,
        1.92,
        11.78,
        4.62,
        theme,
    )


def _render_coordinate_visual(
    slide: Any,
    unit: SlideSpec,
    visual: dict[str, Any],
    theme: dict[str, str],
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.util import Inches

    _source_panel(slide, unit, 0.78, 1.92, 4.65, 4.62, theme)
    _shape(slide, 5.72, 1.92, 6.84, 4.62, theme["canvas"], radius=True, line=theme["chart_bg"])
    origin_x, origin_y = 8.98, 4.38
    for x1, y1, x2, y2 in (
        (6.25, origin_y, 12.05, origin_y),
        (origin_x, 2.36, origin_x, 6.0),
    ):
        axis = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x1),
            Inches(y1),
            Inches(x2),
            Inches(y2),
        )
        axis.line.color.rgb = RGBColor.from_string(theme["muted"])
    parameters = visual.get("parameters") or {}
    points = list(parameters.get("points") or [])
    labels = [str(item) for item in parameters.get("point_labels") or []]
    maximum = max(
        1.0,
        max(
            (abs(float(value)) for point in points[:10] for value in point[:2]),
            default=1.0,
        ),
    )
    x_scale = 2.25 / maximum
    y_scale = 1.52 / maximum
    rendered_points = [
        (
            origin_x + float(raw_x) * x_scale,
            origin_y - float(raw_y) * y_scale,
        )
        for raw_x, raw_y in points[:10]
    ]
    if parameters.get("connect_points") and len(rendered_points) >= 2:
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(rendered_points[0][0]),
            Inches(rendered_points[0][1]),
            Inches(rendered_points[1][0]),
            Inches(rendered_points[1][1]),
        )
        connector.line.color.rgb = RGBColor.from_string(theme["accent"])
        connector.line.width = Inches(0.025)
    for index, (x, y) in enumerate(rendered_points):
        shape = _shape(slide, x - 0.07, y - 0.07, 0.14, 0.14, theme["accent"], radius=True)
        label = labels[index] if index < len(labels) else ""
        _set_alt_text(shape, label or str(visual.get("alt_text") or "坐标数据点"))
        _text(
            slide,
            label,
            x + 0.12,
            y - 0.34,
            1.1,
            0.3,
            16,
            theme["ink"],
            bold=True,
        )
    axis_labels = [str(item) for item in parameters.get("axis_labels") or ["x", "y"]]
    _text(slide, axis_labels[0], 12.04, origin_y - 0.28, 0.22, 0.24, 12, theme["muted"], bold=True)
    _text(slide, axis_labels[1], origin_x + 0.12, 2.2, 0.22, 0.24, 12, theme["muted"], bold=True)


def _render_chart_visual(
    slide: Any,
    unit: SlideSpec,
    visual: dict[str, Any],
    theme: dict[str, str],
) -> None:
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    parameters = visual.get("parameters") or {}
    categories = [str(item) for item in parameters.get("categories") or []]
    series = list(parameters.get("series") or [])
    if not categories or not series:
        _source_panel(slide, unit, 0.78, 1.92, 11.78, 4.62, theme)
        return
    data = ChartData()
    data.categories = categories
    for item in series[:4]:
        values = item.get("values") or []
        if len(values) != len(categories) or not all(isinstance(value, (int, float)) for value in values):
            continue
        data.add_series(str(item.get("name") or "Series"), values)
    if not data._series:
        _source_panel(slide, unit, 0.78, 1.92, 11.78, 4.62, theme)
        return
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.78),
        Inches(1.92),
        Inches(7.4),
        Inches(4.62),
        data,
    ).chart
    chart.has_legend = len(series) > 1
    _source_panel(slide, unit, 8.48, 1.92, 4.08, 4.62, theme)


def _render_image_visual(
    slide: Any,
    unit: SlideSpec,
    visual: dict[str, Any],
    theme: dict[str, str],
    asset_repository: SlideAssetRepository,
) -> bool:
    from pptx.util import Inches

    asset_id = str(visual.get("asset_id") or "")
    if not asset_id:
        return False
    try:
        image_path = asset_repository.resolve(asset_id)
    except (FileNotFoundError, ValueError):
        return False
    from PIL import Image

    with Image.open(image_path) as image:
        image_width, image_height = image.size
    frame_x, frame_y, frame_width, frame_height = 0.78, 1.92, 7.2, 4.62
    image_ratio = image_width / max(1, image_height)
    frame_ratio = frame_width / frame_height
    if image_ratio >= frame_ratio:
        picture_width = frame_width
        picture_height = frame_width / image_ratio
        picture_x = frame_x
        picture_y = frame_y + (frame_height - picture_height) / 2
    else:
        picture_height = frame_height
        picture_width = frame_height * image_ratio
        picture_x = frame_x + (frame_width - picture_width) / 2
        picture_y = frame_y
    picture = slide.shapes.add_picture(
        str(image_path),
        Inches(picture_x),
        Inches(picture_y),
        width=Inches(picture_width),
        height=Inches(picture_height),
    )
    _set_alt_text(picture, str(visual.get("alt_text") or "课程视觉素材"))
    _source_panel(slide, unit, 8.28, 1.92, 4.28, 4.62, theme)
    return True


def _source_panel(
    slide: Any,
    unit: SlideSpec,
    x: float,
    y: float,
    width: float,
    height: float,
    theme: dict[str, str],
) -> None:
    body = _visible_source_text(unit)
    _shape(slide, x, y, width, height, theme["surface"], radius=True, line=theme["chart_bg"])
    _text(slide, "课程正文", x + 0.28, y + 0.25, 1.1, 0.28, 10, theme["muted"], bold=True)
    _text(
        slide,
        body,
        x + 0.28,
        y + 0.72,
        width - 0.56,
        height - 1.05,
        18 if len(body) <= 170 else 16,
        theme["ink"],
        font=theme["body_font"],
        east_asian_font=theme["body_east_asian_font"],
    )


def _visible_source_text(unit: SlideSpec) -> str:
    values: list[str] = []
    for block in unit.blocks:
        if block.title:
            values.append(block.title)
        if block.items:
            values.extend(f"• {item}" for item in block.items if item)
        elif block.content:
            values.append(block.content)
    return "\n\n".join(value for value in values if value)


def _diagram_label(value: str, maximum: int = 32) -> str:
    """Keep diagram nodes glanceable while preserving the full source in alt text."""
    clean = " ".join(str(value or "").split())
    if len(clean) <= maximum:
        return clean
    boundary = max(
        clean.rfind("。", 0, maximum),
        clean.rfind("；", 0, maximum),
        clean.rfind("，", 0, maximum),
        clean.rfind("）", 0, maximum),
        clean.rfind(")", 0, maximum),
        clean.rfind(" ", 0, maximum),
    )
    if boundary < maximum // 2:
        boundary = maximum
    elif clean[boundary] in "）)":
        boundary += 1
    return clean[:boundary].rstrip("。；， ")


def _visual_caption(visual: dict[str, Any]) -> str:
    diagram_type = str((visual.get("parameters") or {}).get("diagram_type") or "")
    if diagram_type:
        return {
            "process": "步骤关系",
            "cause-effect": "因果关系",
            "comparison": "对比关系",
            "mapping": "映射关系",
            "hierarchy": "概念层级",
            "reasoning": "推理关系",
        }.get(diagram_type, "概念关系")
    alt_text = str(visual.get("alt_text") or "").strip()
    if alt_text.startswith(("由课程原文", "按课程原文", "课程原文")):
        return "概念关系"
    return alt_text or "概念关系"


def _set_alt_text(shape: Any, value: str) -> None:
    try:
        non_visual = shape._element.xpath(".//p:cNvPr")[0]
        non_visual.set("descr", str(value or ""))
    except (AttributeError, IndexError):
        return


def _add_theme_visual_asset(
    slide: Any,
    theme: dict[str, Any],
    asset_name: str,
) -> bool:
    """Place a bundled 16:9 theme visual behind editable slide objects."""
    image_path = slide_theme_asset_path(theme, asset_name)
    if image_path is None:
        return False
    from pptx.util import Inches

    picture = slide.shapes.add_picture(
        str(image_path),
        Inches(0.01),
        Inches(0.01),
        width=Inches(13.30),
        height=Inches(7.47),
    )
    asset = (theme.get("visual_assets") or {}).get(asset_name) or {}
    _set_alt_text(picture, str(asset.get("alt") or "课程主题装饰背景"))
    return True


def _add_theme_page_background(
    slide: Any,
    unit: SlideSpec,
    theme: dict[str, Any],
    resolved_layout: str,
) -> bool:
    """Select a low-distraction authored background for an interior slide."""
    candidates = {
        str(resolved_layout or "").strip(),
        str(getattr(unit, "layout", "") or "").strip(),
        str(getattr(unit, "scene_kind", "") or "").strip(),
    }
    for profile in (theme.get("background_profiles") or {}).values():
        layouts = {str(value).strip() for value in profile.get("layouts") or []}
        if candidates & layouts:
            return _add_theme_visual_asset(slide, theme, str(profile.get("asset") or ""))
    return False


def _theme_text_box_style(
    theme: dict[str, Any],
    style_name: str,
    *,
    fill: str,
    border: str,
    accent: str,
    text: str | None = None,
) -> dict[str, str]:
    """Resolve one semantic text-box token with safe renderer fallbacks."""
    token = dict((theme.get("text_box_styles") or {}).get(style_name) or {})
    return {
        "fill": str(token.get("fill") or fill),
        "border": str(token.get("border") or border),
        "depth": str(token.get("depth") or token.get("border") or border),
        "accent": str(token.get("accent") or accent),
        "text": str(token.get("text") or text or theme["ink"]),
    }


def _semantic_panel(
    slide: Any,
    x: float,
    y: float,
    width: float,
    height: float,
    style: dict[str, str],
    *,
    rail: bool = True,
) -> Any:
    """Build an editable paper-card panel with restrained depth and highlight."""
    _shape(
        slide,
        x + 0.045,
        y + 0.055,
        width,
        height,
        style["depth"],
        radius=True,
    )
    panel = _shape(
        slide,
        x,
        y,
        width,
        height,
        style["fill"],
        radius=True,
        line=style["border"],
    )
    _shape(
        slide,
        x + 0.10,
        y + 0.025,
        max(0.12, width - 0.20),
        0.018,
        "FFFFFF",
        radius=False,
    )
    if rail:
        _shape(slide, x, y, 0.07, height, style["accent"], radius=False)
    return panel


def _uses_source_only_audience_labels(unit: SlideSpec) -> bool:
    return str(unit.quality.get("audience_label_policy") or "") == "source_only"


def _render_authored_cover(
    slide: Any,
    unit: SlideSpec,
    theme: dict[str, Any],
    *,
    minimal: bool = False,
) -> bool:
    """Render the authored Qizhi cover when the theme ships a cover visual."""
    if not _add_theme_visual_asset(slide, theme, "cover"):
        return False
    eyebrow = unit.eyebrow or (
        "" if _uses_source_only_audience_labels(unit) else "课堂演示"
    )
    if eyebrow:
        _text(
            slide,
            eyebrow,
            0.92,
            0.72,
            4.0,
            0.38,
            14,
            theme["accent"],
            bold=True,
        )
    title_size = 35 if len(unit.title) > 10 else 44 if len(unit.title) > 6 else 50
    _text(
        slide,
        unit.title,
        0.92,
        1.32 if minimal else 1.22,
        8.15,
        2.45,
        title_size,
        theme["title"],
        bold=True,
        font=theme["title_font"],
        east_asian_font=theme["title_east_asian_font"],
    )
    if unit.subtitle:
        subtitle_y = 3.80 if minimal else 3.86
        subtitle_width = 7.55 if minimal else 6.85
        subtitle_height = 1.26 if minimal else 0.58
        _text(
            slide,
            unit.subtitle,
            0.94,
            subtitle_y,
            subtitle_width,
            subtitle_height,
            17,
            theme["muted"],
        )
    if not minimal:
        message = unit.key_message or _block_content(unit.blocks, 0)
        if message:
            _shape(
                slide,
                0.92,
                4.62,
                6.88,
                1.02,
                theme["surface"],
                radius=True,
                line=theme["chart_bg"],
            )
            _shape(slide, 1.14, 4.84, 0.07, 0.58, theme["green"], radius=False)
            _text(slide, message, 1.43, 4.76, 5.98, 0.68, 17, theme["ink"], bold=True)
    _text(
        slide,
        str(theme.get("label") or "启智课堂"),
        11.10,
        0.72,
        1.32,
        0.34,
        13,
        "FFFFFF",
        bold=True,
        align="center",
    )
    _text(slide, "同源课程课件 · 可继续编辑", 0.94, 6.62, 4.8, 0.28, 10, theme["muted"])
    return True


def _render_cover(slide: Any, unit: SlideSpec, theme: dict[str, str]) -> None:
    if _render_authored_cover(slide, unit, theme):
        return
    _shape(slide, 0.58, 0.55, 0.12, 5.92, theme["accent"], radius=False)
    _shape(slide, 10.82, 0.0, 2.513, 7.5, theme["accent_soft"], radius=False)
    _shape(slide, 11.35, 0.72, 1.04, 1.04, theme["green"], radius=True)
    _text(slide, "灵知", 11.35, 0.99, 1.04, 0.34, 15, "FFFFFF", bold=True, align="center")
    eyebrow = unit.eyebrow or (
        "" if _uses_source_only_audience_labels(unit) else "课程演示"
    )
    if eyebrow:
        _text(slide, eyebrow, 0.92, 0.72, 4.0, 0.38, 14, theme["accent"], bold=True)
    cover_title_size = 42 if len(unit.title) > 32 else 46 if len(unit.title) > 22 else 50
    _text(
        slide, unit.title, 0.92, 1.22, 9.15, 2.6, cover_title_size, theme["title"], bold=True,
        font=theme["title_font"], east_asian_font=theme["title_east_asian_font"],
    )
    if unit.subtitle:
        _text(slide, unit.subtitle, 0.94, 3.77, 7.8, 0.42, 16, theme["muted"])
    _shape(slide, 0.92, 4.23, 8.85, 1.14, theme["canvas"], radius=True, line=theme["accent_soft"])
    _text(slide, "学习主线", 1.18, 4.48, 1.2, 0.28, 11, theme["green"], bold=True)
    _text(slide, unit.key_message or _block_content(unit.blocks, 0), 2.42, 4.35, 7.0, 0.68, 18, theme["ink"], bold=True)
    _text(slide, "同一课程结构 · 知识与能力绑定 · 可继续编辑", 0.94, 6.45, 7.0, 0.32, 11, theme["muted"])


def _render_cover_minimal(slide: Any, unit: SlideSpec, theme: dict[str, str]) -> None:
    """Render a restrained title page with one clear focal hierarchy."""
    if _render_authored_cover(slide, unit, theme, minimal=True):
        return
    _shape(slide, 0.82, 0.76, 0.12, 0.72, theme["accent"], radius=False)
    eyebrow = unit.eyebrow or (
        "" if _uses_source_only_audience_labels(unit) else "课程课件"
    )
    if eyebrow:
        _text(
            slide,
            eyebrow,
            1.12,
            0.86,
            3.6,
            0.34,
            13,
            theme["accent"],
            bold=True,
        )
    title_size = 38 if len(unit.title) > 34 else 44 if len(unit.title) > 22 else 50
    _text(
        slide,
        unit.title,
        0.86,
        2.02,
        10.9,
        2.05,
        title_size,
        theme["title"],
        bold=True,
        font=theme["title_font"],
        east_asian_font=theme["title_east_asian_font"],
    )
    _shape(slide, 0.88, 5.5, 3.25, 0.06, theme["accent"], radius=False)
    if unit.subtitle:
        _text(slide, unit.subtitle, 0.9, 5.66, 10.9, 1.18, 16, theme["muted"])


def _render_cover_editorial(slide: Any, unit: SlideSpec, theme: dict[str, str]) -> None:
    """Render a concise editorial cover with a balanced visual field."""
    if _render_authored_cover(slide, unit, theme):
        return
    _shape(slide, 9.25, 0.0, 4.083, 7.5, theme["accent_soft"], radius=False)
    _shape(slide, 0.88, 0.78, 0.12, 0.74, theme["accent"], radius=False)
    _text(
        slide,
        unit.eyebrow or "课程课件",
        1.18,
        0.88,
        3.8,
        0.36,
        13,
        theme["accent"],
        bold=True,
    )
    _text(
        slide,
        unit.title,
        0.9,
        2.0,
        7.65,
        2.0,
        50 if len(unit.title) <= 14 else 44,
        theme["title"],
        bold=True,
        font=theme["title_font"],
        east_asian_font=theme["title_east_asian_font"],
    )
    if unit.subtitle:
        _shape(slide, 0.92, 4.72, 5.65, 0.05, theme["accent"], radius=False)
        _text(slide, unit.subtitle, 0.94, 5.04, 7.3, 0.7, 20, theme["ink"], bold=True)
    if not _uses_source_only_audience_labels(unit):
        _text(slide, "COURSE", 9.72, 1.08, 2.4, 0.42, 14, theme["accent"], bold=True)
    _text(
        slide,
        "概念\n方法\n应用",
        9.72,
        2.05,
        2.25,
        2.45,
        30,
        theme["title"],
        bold=True,
        font=theme["title_font"],
        east_asian_font=theme["title_east_asian_font"],
    )


def _render_roadmap(slide: Any, unit: SlideSpec, theme: dict[str, str]) -> None:
    _heading(slide, unit, theme)
    items = _all_items(unit)[:8]
    columns = 2
    rows = max(1, (len(items) + 1) // columns)
    card_h = min(1.04, 4.35 / rows)
    for index, item in enumerate(items):
        row, column = divmod(index, columns)
        x = 0.78 + column * 6.0
        y = 2.05 + row * (card_h + 0.16)
        _shape(slide, x, y, 5.55, card_h, theme["canvas"], radius=True, line=theme["chart_bg"])
        _shape(slide, x + 0.22, y + 0.22, 0.5, 0.5, theme["accent_soft"], radius=True)
        _text(slide, f"{index + 1:02d}", x + 0.22, y + 0.33, 0.5, 0.2, 10, theme["accent"], bold=True, align="center")
        _text(slide, item, x + 0.9, y + 0.24, 4.35, card_h - 0.25, 17, theme["ink"], bold=True)


def _render_agenda_linear(slide: Any, unit: SlideSpec, theme: dict[str, str]) -> None:
    """Render the course route as an editorial sequence instead of a card grid."""
    _heading(slide, unit, theme)
    agenda_block = next(
        (
            block
            for block in unit.blocks
            if isinstance(block.metadata.get("agenda_entries"), list)
        ),
        None,
    )
    entries = list(agenda_block.metadata.get("agenda_entries") or []) if agenda_block else []
    if not entries:
        entries = [
            {"index": index, "title": item, "description": ""}
            for index, item in enumerate(_all_items(unit)[:4], start=1)
        ]
    if not entries:
        _render_navigation_statement(slide, unit, theme, heading_already_rendered=True)
        return
    row_h = min(1.03, 4.35 / max(1, len(entries)))
    for fallback_index, entry in enumerate(entries[:4], start=1):
        index = int(entry.get("index") or fallback_index)
        title = str(entry.get("title") or "").strip()
        description = str(entry.get("description") or "").strip()
        y = 1.9 + (fallback_index - 1) * row_h
        _shape(
            slide,
            0.82,
            y + row_h - 0.04,
            11.55,
            0.018,
            theme["chart_bg"],
            radius=False,
        )
        _text(
            slide,
            f"{index:02d}",
            0.88,
            y + 0.13,
            0.72,
            0.28,
            11,
            theme["accent"],
            bold=True,
            font="Aptos Mono",
        )
        _text(
            slide,
            title,
            1.78,
            y + 0.08,
            9.95 if not description else 4.15,
            row_h - 0.12,
            18 if len(title) <= 24 else 16,
            theme["ink"],
            bold=True,
        )
        if description:
            _text(
                slide,
                description,
                6.18,
                y + 0.08,
                5.55,
                row_h - 0.12,
                16,
                theme["muted"],
            )


def _render_chapter(slide: Any, unit: SlideSpec, theme: dict[str, str]) -> None:
    if _add_theme_visual_asset(slide, theme, "chapter"):
        chapter_number = _chapter_number(unit.title)
        if not _uses_source_only_audience_labels(unit):
            _text(slide, "CHAPTER", 0.64, 0.88, 2.5, 0.34, 12, "FFFFFF", bold=True)
        _text(slide, chapter_number, 0.64, 1.48, 2.62, 1.35, 54, "FFFFFF", bold=True)
        eyebrow = unit.eyebrow or (
            "" if _uses_source_only_audience_labels(unit) else "章节转场"
        )
        if eyebrow:
            _text(slide, eyebrow, 4.48, 1.02, 2.4, 0.32, 12, theme["green"], bold=True)
        _text(
            slide,
            unit.title,
            4.48,
            1.58,
            7.55,
            1.42,
            35,
            theme["title"],
            bold=True,
            font=theme["title_font"],
            east_asian_font=theme["title_east_asian_font"],
        )
        chapter_message = (
            unit.key_message
            or _block_content(unit.blocks, 0)
            or unit.teaching_job
            or unit.takeaway
        )
        _shape(slide, 4.48, 3.42, 0.08, 2.43, theme["accent"], radius=False)
        _text(slide, "本章主线", 4.82, 3.47, 1.5, 0.30, 11, theme["accent"], bold=True)
        _text(slide, chapter_message, 4.82, 3.92, 6.72, 1.75, 17, theme["ink"], bold=True)
        return
    _shape(slide, 0.0, 0.0, 4.05, 7.5, theme["accent_soft"], radius=False)
    chapter_number = _chapter_number(unit.title)
    _text(slide, chapter_number, 0.72, 1.15, 2.5, 1.35, 54, theme["accent"], bold=True)
    eyebrow = unit.eyebrow or (
        "" if _uses_source_only_audience_labels(unit) else "章节转场"
    )
    if eyebrow:
        _text(slide, eyebrow, 4.65, 1.08, 2.3, 0.32, 12, theme["green"], bold=True)
    _text(
        slide, unit.title, 4.65, 1.62, 7.55, 1.4, 35, theme["title"], bold=True,
        font=theme["title_font"], east_asian_font=theme["title_east_asian_font"],
    )
    _shape(slide, 4.65, 3.45, 6.95, 2.5, theme["canvas"], radius=True, line=theme["chart_bg"])
    _text(slide, "本章主线", 4.98, 3.78, 1.4, 0.3, 11, theme["accent"], bold=True)
    chapter_message = (
        unit.key_message
        or _block_content(unit.blocks, 0)
        or unit.teaching_job
        or unit.takeaway
    )
    _text(slide, chapter_message, 4.98, 4.12, 6.18, 1.75, 16, theme["ink"], bold=True)


def _render_objective(slide: Any, unit: SlideSpec, theme: dict[str, str]) -> None:
    if not _visible_source_text(unit) and not unit.key_message:
        _render_navigation_statement(slide, unit, theme)
        return
    _heading(slide, unit, theme)
    question = _find_block(unit, "callout") or (unit.blocks[0] if unit.blocks else None)
    question_style = _theme_text_box_style(
        theme,
        "message",
        fill=theme["accent_soft"],
        border=theme["chart_bg"],
        accent=theme["accent"],
    )
    _semantic_panel(slide, 0.76, 1.83, 5.0, 4.67, question_style)
    _text(slide, "要解决的问题", 1.08, 2.14, 2.0, 0.32, 12, question_style["accent"], bold=True)
    _text(slide, question.content if question else unit.key_message, 1.08, 2.72, 4.28, 2.1, 22, question_style["text"], bold=True)
    right_blocks = [block for block in unit.blocks if block is not question]
    if not right_blocks:
        right_blocks = [SlideBlockSpec(block_id="objective", type="bullets", items=[unit.key_message])]
    for index, block in enumerate(right_blocks[:2]):
        y = 1.83 + index * 2.35
        style = _theme_text_box_style(
            theme,
            "feedback" if index == 0 else "practice",
            fill=theme["green_soft"] if index == 0 else theme["amber_soft"],
            border=theme["chart_bg"],
            accent=theme["green"] if index == 0 else theme["amber"],
        )
        _semantic_panel(slide, 6.05, y, 6.48, 2.12, style)
        _text(slide, block.title or ("知识坐标" if index == 0 else "完成后能够"), 6.38, y + 0.27, 2.3, 0.3, 11, style["accent"], bold=True)
        _bullets(slide, block.items or [block.content], 6.38, y + 0.72, 5.72, 1.12, 14, style["text"], style["accent"])


def _render_concept(slide: Any, unit: SlideSpec, theme: dict[str, str]) -> None:
    _heading(slide, unit, theme)
    if unit.key_message:
        message_style = _theme_text_box_style(
            theme,
            "message",
            fill=theme["accent_soft"],
            border=theme["chart_bg"],
            accent=theme["accent"],
        )
        _semantic_panel(slide, 0.76, 1.72, 11.82, 0.86, message_style)
        _text(slide, unit.key_message, 1.05, 1.94, 11.22, 0.4, 16, message_style["text"], bold=True)
    blocks = unit.blocks
    width = 11.82 / max(1, len(blocks)) - 0.18
    for index, block in enumerate(blocks):
        x = 0.76 + index * (width + 0.27)
        style_name = {
            "misconception": "misconception",
            "exercise": "practice",
            "callout": "message",
            "code": "evidence",
            "process": "reasoning",
            "comparison": "boundary",
        }.get(block.type, ("definition", "boundary", "practice")[index % 3])
        style = _theme_text_box_style(
            theme,
            style_name,
            fill=theme["canvas"],
            border=theme["chart_bg"],
            accent=[theme["accent"], theme["green"], theme["amber"]][index % 3],
        )
        _semantic_panel(slide, x, 2.87, width, 3.52, style)
        _text(
            slide,
            block.title or unit.eyebrow or f"要点 {index + 1}",
            x + 0.3,
            3.14,
            width - 0.55,
            0.46,
            15,
            style["accent"],
            bold=True,
        )
        if block.items:
            item_size = 19 if len(blocks) == 1 and len(block.items) <= 3 else 16
            _bullets(slide, block.items, x + 0.3, 3.72, width - 0.58, 2.25, item_size, style["text"], style["accent"])
        else:
            is_formula = bool(block.metadata.get("formula"))
            body_size = (
                26
                if len(blocks) == 1 and len(block.content) <= 70
                else 21
                if len(blocks) == 1 and len(block.content) <= 150
                else 16
            )
            _text(
                slide,
                block.content,
                x + 0.3,
                3.76,
                width - 0.58,
                2.18,
                body_size,
                style["text"],
                font=theme["math_font"] if is_formula else theme["body_font"],
                east_asian_font=theme["body_east_asian_font"],
            )


def _render_hero_statement(slide: Any, unit: SlideSpec, theme: dict[str, str]) -> None:
    _heading(slide, unit, theme)
    body = " ".join(
        value
        for block in unit.blocks
        for value in (block.items or [block.content])
        if value
    )
    body = body or unit.takeaway or unit.teaching_job or unit.title
    _shape(slide, 0.82, 1.82, 11.68, 4.52, theme["accent_soft"], radius=True)
    _shape(slide, 1.16, 2.2, 0.12, 3.72, theme["accent"], radius=False)
    _text(slide, "核心判断", 1.58, 2.24, 2.0, 0.34, 12, theme["accent"], bold=True)
    _text(
        slide,
        body,
        1.58,
        2.92,
        9.9,
        2.42,
        27 if len(body) <= 90 else 21,
        theme["ink"],
        bold=len(body) <= 90,
        font=theme["title_font"],
        east_asian_font=theme["title_east_asian_font"],
    )


def _render_claim_only(slide: Any, unit: SlideSpec, theme: dict[str, str]) -> None:
    """Render one intentional claim as the dominant object, exactly once."""
    claim = " ".join(
        value
        for block in unit.blocks
        for value in (block.items or [block.content])
        if value
    )
    claim = claim or unit.key_message or unit.takeaway or unit.title
    _text(
        slide,
        unit.eyebrow or "核心判断",
        0.82,
        0.72,
        3.0,
        0.32,
        12,
        theme["accent"],
        bold=True,
    )
    _shape(slide, 0.82, 1.28, 11.68, 0.018, theme["chart_bg"], radius=False)
    _shape(slide, 0.9, 1.88, 0.12, 4.28, theme["accent"], radius=False)
    _text(
        slide,
        claim,
        1.4,
        2.18,
        10.25,
        2.78,
        29 if len(claim) <= 58 else 24,
        theme["ink"],
        bold=True,
        font=theme["title_font"],
        east_asian_font=theme["title_east_asian_font"],
    )
    _shape(slide, 1.4, 5.48, 4.25, 0.05, theme["chart_bg"], radius=False)


def _render_navigation_statement(
    slide: Any,
    unit: SlideSpec,
    theme: dict[str, str],
    *,
    heading_already_rendered: bool = False,
) -> None:
    """Render source-free navigation beats without empty UI-like containers."""
    if not heading_already_rendered:
        _heading(slide, unit, theme)
    topic = _heading_excerpt(unit.title, 42)
    if unit.layout in {"recap", "summary"}:
        prefix = "本章回顾"
        detail = unit.takeaway or unit.key_message or f"回到“{topic}”，串联本章的概念、方法与应用。"
    else:
        prefix = "本节学习问题"
        detail = unit.teaching_job or unit.takeaway or f"围绕“{topic}”，建立定义、判断方法与应用联系。"
    _shape(slide, 0.92, 2.18, 0.12, 3.5, theme["accent"], radius=False)
    _text(slide, prefix, 1.38, 2.22, 2.5, 0.36, 12, theme["accent"], bold=True)
    _text(
        slide,
        detail,
        1.38,
        2.92,
        10.3,
        1.72,
        27 if len(detail) <= 42 else 22,
        theme["ink"],
        bold=True,
        font=theme["title_font"],
        east_asian_font=theme["title_east_asian_font"],
    )
    _text(
        slide,
        "先明确问题，再连接概念、方法与检验。",
        1.38,
        5.18,
        8.8,
        0.42,
        14,
        theme["muted"],
    )


def _render_classification_three(
    slide: Any,
    unit: SlideSpec,
    theme: dict[str, str],
) -> None:
    """Render exactly three peer concepts as equal semantic columns."""
    _heading(slide, unit, theme)
    items = [
        value
        for block in unit.blocks
        for value in (block.items or [block.content])
        if value
    ]
    capacity_profile = str(unit.quality.get("v6_capacity_profile") or "")
    if (
        capacity_profile == CLASSIFICATION_THREE_CARDS_V1
        and not classification_three_card_metrics(items)["fits"]
    ):
        raise ValueError("template_slot_capacity_exceeded")
    if len(items) != 3:
        _render_editorial_body(
            slide,
            unit,
            theme,
            heading_already_rendered=True,
        )
        return
    accents = (theme["accent"], theme["green"], theme["amber"])
    for index, item in enumerate(items):
        x = 0.82 + index * 3.92
        _shape(slide, x, 2.02, 3.58, 0.08, accents[index], radius=False)
        _text(
            slide,
            f"{index + 1:02d}",
            x,
            2.35,
            0.72,
            0.3,
            11,
            accents[index],
            bold=True,
            font="Aptos Mono",
        )
        _text(
            slide,
            item,
            x,
            3.0,
            3.38,
            2.35,
            18 if len(item) <= 48 else 16,
            theme["ink"],
            bold=True,
        )
        if index < 2:
            _shape(
                slide,
                x + 3.72,
                2.02,
                0.018,
                3.95,
                theme["chart_bg"],
                radius=False,
            )


def _render_editorial_body(
    slide: Any,
    unit: SlideSpec,
    theme: dict[str, str],
    *,
    heading_already_rendered: bool = False,
    body_capacity_profile: str = "",
) -> None:
    if not _visible_source_text(unit):
        _render_navigation_statement(slide, unit, theme)
        return
    if not body_capacity_profile:
        body_capacity_profile = str(
            unit.quality.get("v6_capacity_profile") or ""
        )
    if not heading_already_rendered:
        _heading(slide, unit, theme)
    values = [
        value
        for block in unit.blocks
        for value in (block.items or [block.content])
        if value
    ]
    body = "\n\n".join(values)
    body_metrics = (
        balanced_two_column_body_metrics(body)
        if body_capacity_profile == BALANCED_TWO_COLUMN_BODY_V1
        else None
    )
    if body_metrics is not None and not body_metrics["fits"]:
        raise ValueError("template_slot_capacity_exceeded")
    if body_metrics is not None and body_metrics["mode"] == "two-column":
        _shape(slide, 0.86, 1.92, 0.1, 4.32, theme["accent"], radius=False)
        segments = list(body_metrics["segments"])
        for index, segment in enumerate(segments[:2]):
            x = 1.18 + index * 5.64
            body_shape = _text(
                slide,
                segment,
                x,
                2.12,
                5.3,
                4.4,
                16,
                theme["ink"],
            )
            body_shape.name = (
                f"{body_shape.name} [v6-body-capacity={body_capacity_profile}] "
                f"[v6-body-max-lines={body_metrics['maximum_safe_lines']}] "
                f"[v6-body-column={index + 1}]"
            )
        _shape(slide, 6.65, 2.1, 0.018, 4.35, theme["chart_bg"], radius=False)
        return
    explicit_lines = len([line for line in body.splitlines() if line.strip()])
    body_font_size = (
        16
        if explicit_lines >= 8 or len(body) > 180
        else 22
        if len(body) > 90 or explicit_lines >= 5
        else 26
    )
    _shape(slide, 0.86, 1.92, 0.1, 4.32, theme["accent"], radius=False)
    body_shape = _text(
        slide,
        body,
        1.34,
        2.3,
        10.75,
        3.55,
        body_font_size,
        theme["ink"],
    )
    if body_metrics is not None:
        body_shape.name = (
            f"{body_shape.name} [v6-body-capacity={body_capacity_profile}] "
            f"[v6-body-max-lines={body_metrics['maximum_safe_lines']}]"
        )
    _shape(slide, 1.34, 6.13, 4.35, 0.025, theme["chart_bg"], radius=False)


def _balanced_two_column_body(value: str) -> list[str]:
    """Use the same deterministic split already proved by template capacity."""

    return list(balanced_two_column_body_metrics(value)["segments"])


def _render_two_column(slide: Any, unit: SlideSpec, theme: dict[str, str]) -> None:
    _heading(slide, unit, theme)
    values = [
        value
        for block in unit.blocks
        for value in (block.items or [block.content])
        if value
    ]
    body_capacity_profile = str(
        unit.quality.get("v6_capacity_profile") or ""
    )
    body_metrics = None
    if len(values) == 1:
        body_metrics = balanced_two_column_body_metrics(values[0])
        if (
            body_capacity_profile == BALANCED_TWO_COLUMN_BODY_V1
            and not body_metrics["fits"]
        ):
            raise ValueError("template_slot_capacity_exceeded")
        values = list(body_metrics["segments"])
    if len(values) < 2:
        _render_editorial_body(
            slide,
            unit,
            theme,
            heading_already_rendered=True,
            body_capacity_profile=body_capacity_profile,
        )
        return
    labels = ("依据", "推论")
    styles = (
        _theme_text_box_style(
            theme,
            "definition",
            fill=theme["accent_soft"],
            border=theme["chart_bg"],
            accent=theme["accent"],
        ),
        _theme_text_box_style(
            theme,
            "boundary",
            fill=theme["green_soft"],
            border=theme["chart_bg"],
            accent=theme["green"],
        ),
    )
    for index, value in enumerate(values[:2]):
        x = 0.68 + index * 6.05
        style = styles[index]
        _semantic_panel(slide, x, 1.88, 5.92, 5.0, style)
        _text(slide, labels[index], x + 0.3, 2.04, 1.2, 0.25, 12, style["accent"], bold=True)
        body_shape = _text(
            slide,
            value,
            x + 0.3,
            2.34,
            5.32,
            4.42,
            16,
            style["text"],
        )
        if (
            body_capacity_profile == BALANCED_TWO_COLUMN_BODY_V1
            and body_metrics is not None
        ):
            body_shape.name = (
                f"{body_shape.name} "
                f"[v6-body-capacity={body_capacity_profile}] "
                f"[v6-body-max-lines={body_metrics['maximum_safe_lines']}] "
                f"[v6-body-column={index + 1}]"
            )


def _render_case_study(slide: Any, unit: SlideSpec, theme: dict[str, str]) -> None:
    _heading(slide, unit, theme)
    values = [
        value
        for block in unit.blocks
        for value in (block.items or [block.content])
        if value
    ]
    body = "\n\n".join(values)
    _shape(slide, 0.82, 1.82, 3.0, 4.55, theme["green_soft"], radius=True)
    _text(slide, "CASE", 1.18, 2.2, 1.2, 0.36, 13, theme["green"], bold=True)
    _text(slide, "从具体情境\n检验抽象结构", 1.18, 3.0, 2.24, 1.35, 24, theme["ink"], bold=True)
    _shape(slide, 4.12, 1.82, 8.38, 4.55, theme["canvas"], radius=True, line=theme["chart_bg"])
    _text(slide, "案例观察", 4.48, 2.18, 1.6, 0.32, 12, theme["accent"], bold=True)
    _text(slide, body, 4.48, 2.82, 7.3, 2.95, 18 if len(body) <= 180 else 16, theme["ink"])


def _render_parallel_examples(
    slide: Any,
    unit: SlideSpec,
    theme: dict[str, str],
) -> None:
    """Render peer examples without implying causality or reasoning order."""
    _heading(slide, unit, theme)
    values = _all_items(unit)[:4]
    if len(values) < 2:
        _render_editorial_body(
            slide,
            unit,
            theme,
            heading_already_rendered=True,
        )
        return
    gap = 0.28
    width = (11.55 - gap * (len(values) - 1)) / len(values)
    for index, value in enumerate(values):
        x = 0.88 + index * (width + gap)
        _shape(slide, x, 2.08, width, 0.07, theme["accent"], radius=False)
        _text(slide, f"{index + 1:02d}", x, 2.42, 0.68, 0.28, 11, theme["accent"], bold=True)
        _text(
            slide,
            value,
            x,
            3.0,
            width - 0.12,
            2.35,
            19 if len(value) <= 48 else 16,
            theme["ink"],
            bold=True,
        )


def _render_question_prompt(
    slide: Any,
    unit: SlideSpec,
    theme: dict[str, str],
) -> None:
    """Render one practice prompt as a flat, presentation-first composition."""
    _heading(slide, unit, theme)
    values = [
        value
        for block in unit.blocks
        for value in ([block.content] if block.content else []) + list(block.items)
        if value
    ]
    prompt = "\n".join(
        f"{index}. {value}" if len(values) > 1 else value
        for index, value in enumerate(values, start=1)
    ) or unit.key_message or unit.takeaway
    _shape(slide, 0.92, 2.18, 0.11, 3.5, theme["accent"], radius=False)
    _text(
        slide,
        str(unit.quality.get("prompt_label") or "先独立判断"),
        1.42,
        2.24,
        2.4,
        0.34,
        13,
        theme["accent"],
        bold=True,
    )
    _text(
        slide,
        prompt,
        1.42,
        3.0,
        10.45,
        2.8,
        24 if len(prompt) <= 88 else 18,
        theme["ink"],
        bold=True,
        font=theme["title_font"],
        east_asian_font=theme["title_east_asian_font"],
    )


def _worked_example_labels(quality: dict[str, Any], count: int) -> tuple[str, ...]:
    explicit = tuple(
        str(item).strip()
        for item in quality.get("worked_step_labels") or []
        if str(item).strip()
    )
    if len(explicit) >= count:
        return explicit[:count]
    return tuple(f"步骤 {index + 1}" for index in range(count))


def _render_worked_example(
    slide: Any,
    unit: SlideSpec,
    theme: dict[str, str],
) -> None:
    """Render a worked example as a visible reasoning path."""
    _heading(slide, unit, theme)
    values = [
        value
        for block in unit.blocks
        for value in (block.items or [block.content])
        if value
    ]
    if not values:
        _render_navigation_statement(slide, unit, theme)
        return
    if len(values) == 1:
        values = [
            item.strip()
            for item in values[0].split("\n\n")
            if item.strip()
        ]
    if len(values) < 2:
        _render_case_study(slide, unit, theme)
        return
    labels = _worked_example_labels(unit.quality, min(3, len(values)))
    accents = (theme["accent"], theme["green"], theme["amber"])
    _shape(slide, 1.14, 2.17, 0.035, 3.72, theme["chart_bg"], radius=False)
    for index, (label, value) in enumerate(zip(labels, values[:3])):
        y = 1.92 + index * 1.38
        accent = accents[index]
        _shape(slide, 0.88, y + 0.17, 0.55, 0.55, accent, radius=True)
        _text(
            slide,
            str(index + 1),
            0.88,
            y + 0.34,
            0.55,
            0.18,
            11,
            "FFFFFF",
            bold=True,
            align="center",
        )
        _text(slide, label, 1.72, y + 0.1, 1.2, 0.3, 12, accent, bold=True)
        _text(
            slide,
            value,
            3.0,
            y,
            9.02,
            0.96,
            18 if len(value) <= 90 else 16,
            theme["ink"],
            bold=index == 2,
        )
        if index < len(labels) - 1:
            _shape(
                slide,
                1.72,
                y + 1.1,
                10.28,
                0.02,
                theme["chart_bg"],
                radius=False,
            )


def _render_comparison(slide: Any, unit: SlideSpec, theme: dict[str, str]) -> None:
    _heading(slide, unit, theme)
    block = _find_block(unit, "comparison")
    if block and block.metadata.get("rows"):
        headers = [str(value) for value in block.metadata.get("headers") or []]
        rows = [[str(value) for value in row] for row in block.metadata.get("rows") or []]
        _table(slide, headers, rows, 0.78, 1.85, 11.78, 4.45, theme)
    else:
        _render_concept(slide, unit, theme)
        return
    if unit.key_message:
        _shape(slide, 0.78, 6.37, 11.78, 0.48, theme["amber_soft"], radius=True)
        _text(slide, unit.key_message, 1.02, 6.48, 11.25, 0.23, 11, theme["amber"], bold=True)


def _render_process(slide: Any, unit: SlideSpec, theme: dict[str, str]) -> None:
    _heading(slide, unit, theme)
    all_items = _all_items(unit) or [
        block.content for block in unit.blocks if block.content
    ]
    capacity_profile = str(unit.quality.get("v6_capacity_profile") or "")
    if capacity_profile == HORIZONTAL_PROCESS_CARDS_V1:
        items = all_items
    else:
        items = all_items[:5]
    if str(unit.quality.get("task_prompt_mode") or "") == "action":
        _text(
            slide,
            str(unit.quality.get("prompt_label") or "执行步骤"),
            0.86,
            1.9,
            2.0,
            0.22,
            12,
            theme["accent"],
            bold=True,
        )
        step_parts = [_split_ordered_step(item) for item in items]
        title_fonts = [18 if len(title) <= 28 else 16 for title, _detail in step_parts]
        detail_font = 16
        required_heights: list[float] = []
        for (title, detail), title_font in zip(step_parts, title_fonts):
            title_width = 3.08 if detail else 10.2
            title_lines = _wrapped_line_count(
                _display_text(title),
                width_pt=title_width * 72.0,
                font_size_pt=title_font,
            )
            detail_lines = (
                _wrapped_line_count(
                    _display_text(detail),
                    width_pt=7.05 * 72.0,
                    font_size_pt=detail_font,
                )
                if detail
                else 0
            )
            required_heights.append(max(
                0.58,
                title_lines * title_font * 1.22 / 72.0 + 0.16,
                detail_lines * detail_font * 1.22 / 72.0 + 0.16,
            ))
        # Use the full template-safe content area.  The former 4.08-inch limit
        # left unused space above the footer and then compressed valid multi-line
        # steps into frames shorter than their text on wider CJK fonts.
        available_height = 4.60
        gap = 0.04
        usable_height = available_height - gap * max(0, len(items) - 1)
        y = 2.20
        required_total = max(0.01, sum(required_heights))
        if required_total <= usable_height:
            extra_height = (usable_height - required_total) / max(1, len(items))
            heights = [height + extra_height for height in required_heights]
        else:
            raise ValueError("ordered_process_render_capacity_exceeded")
        centers: list[float] = []
        cursor = y
        for height in heights:
            centers.append(cursor + height / 2)
            cursor += height + gap
        if len(centers) > 1:
            _shape(
                slide,
                1.16,
                centers[0],
                0.035,
                centers[-1] - centers[0],
                theme["chart_bg"],
                radius=False,
            )
        for index, ((title, detail), height, title_font) in enumerate(
            zip(step_parts, heights, title_fonts),
            start=1,
        ):
            center_y = y + height / 2
            if index < len(items):
                _shape(
                    slide,
                    1.72,
                    y + height - 0.015,
                    10.38,
                    0.015,
                    theme["chart_bg"],
                    radius=False,
                )
            _circle(slide, 0.88, center_y - 0.28, 0.56, theme["accent"])
            _text(
                slide,
                f"{index:02d}",
                0.88,
                center_y - 0.10,
                0.56,
                0.22,
                11,
                "FFFFFF",
                bold=True,
                align="center",
                font="Aptos Mono",
            )
            _text(
                slide,
                title,
                1.78,
                y + 0.08,
                3.08 if detail else 10.2,
                max(0.46, height - 0.14),
                title_font,
                theme["ink"],
                bold=True,
            )
            if detail:
                _text(
                    slide,
                    detail,
                    5.02,
                    y + 0.08,
                    7.05,
                    max(0.48, height - 0.14),
                    detail_font,
                    theme["muted"],
                )
            y += height + gap
        return
    if capacity_profile == HORIZONTAL_PROCESS_CARDS_V1:
        metrics = horizontal_process_card_metrics(items)
        if not metrics["fits"]:
            raise ValueError("horizontal_process_render_capacity_exceeded")
    width = (11.7 - max(0, len(items) - 1) * 0.24) / max(1, len(items))
    style = _theme_text_box_style(
        theme,
        "reasoning",
        fill=theme["canvas"],
        border=theme["chart_bg"],
        accent=theme["accent"],
    )
    for index, item in enumerate(items):
        x = 0.82 + index * (width + 0.24)
        step_accent = [style["accent"], theme["green"], theme["amber"]][index % 3]
        step_style = {**style, "accent": step_accent}
        _semantic_panel(slide, x, 2.08, width, 3.72, step_style, rail=False)
        _shape(slide, x, 2.08, width, 0.07, step_accent, radius=False)
        _shape(slide, x + 0.22, 2.34, 0.58, 0.58, step_accent, radius=True)
        _text(slide, str(index + 1), x + 0.22, 2.52, 0.58, 0.2, 12, "FFFFFF", bold=True, align="center")
        _text(slide, item, x + 0.23, 3.25, width - 0.46, 1.95, 16, style["text"], bold=True)
        if index < len(items) - 1:
            _text(slide, "→", x + width + 0.01, 3.68, 0.22, 0.35, 17, theme["muted"], bold=True, align="center")


def _render_practice_artifact(
    slide: Any,
    unit: SlideSpec,
    theme: dict[str, str],
) -> None:
    """Render ordered actions beside their source-backed characteristic artifact."""

    _heading(slide, unit, theme)
    process = _find_block(unit, "process")
    steps = list(process.items if process else [])[:7]
    artifact_kind = str(unit.quality.get("v6_practice_artifact_kind") or "")

    _text(
        slide,
        str(unit.quality.get("prompt_label") or "执行并核验"),
        0.82,
        1.90,
        2.0,
        0.25,
        11,
        theme["accent"],
        bold=True,
    )
    step_line_counts = [
        _wrapped_line_count(value, width_pt=3.70 * 72, font_size_pt=16)
        for value in steps
    ]
    minimum_heights = [
        # Allocate the real text-frame line height. Dividing by 1.18 here used
        # to make every row smaller than its text and allowed adjacent steps
        # to overlap even while the slide itself stayed inside the canvas.
        max(0.44, line_count * 16 * 1.22 / 72 + 0.02)
        for line_count in step_line_counts
    ]
    remaining_height = max(0.0, 3.98 - sum(minimum_heights))
    row_heights = [
        height + remaining_height / max(1, len(minimum_heights))
        for height in minimum_heights
    ]
    first_center = 2.34 + (row_heights[0] / 2 if row_heights else 0)
    last_center = 2.34 + sum(row_heights[:-1]) + (row_heights[-1] / 2 if row_heights else 0)
    if len(steps) > 1:
        _shape(
            slide,
            1.06,
            first_center,
            0.025,
            last_center - first_center,
            theme["chart_bg"],
            radius=False,
        )
    y = 2.34
    for index, (value, row_height) in enumerate(zip(steps, row_heights), start=1):
        center_y = y + row_height / 2
        if index < len(steps):
            _shape(
                slide,
                1.48,
                y + row_height - 0.012,
                3.72,
                0.012,
                theme["chart_bg"],
                radius=False,
            )
        _circle(slide, 0.82, center_y - 0.22, 0.44, theme["accent"])
        _text(
            slide,
            f"{index:02d}",
            0.82,
            center_y - 0.07,
            0.44,
            0.16,
            9,
            "FFFFFF",
            bold=True,
            align="center",
            font="Aptos Mono",
        )
        _text(
            slide,
            value,
            1.48,
            y + 0.01,
            3.72,
            max(0.38, row_height - 0.02),
            16,
            theme["ink"],
            bold=len(value) <= 28,
        )
        y += row_height

    _shape(slide, 5.48, 1.94, 0.025, 4.48, theme["chart_bg"], radius=False)
    artifact_x, artifact_y, artifact_w, artifact_h = 5.82, 1.94, 6.70, 4.48
    if artifact_kind == "code":
        code = _find_block(unit, "code")
        _shape(slide, artifact_x, artifact_y, artifact_w, artifact_h, theme["code"], radius=True)
        _render_code_reading_frame(
            slide,
            unit,
            code,
            x=6.16,
            header_y=2.20,
            code_y=2.66,
            width=6.02,
            height=3.28,
            text_color="F5F7FF",
        )
    elif artifact_kind == "formula":
        visual = next((item for item in unit.visuals if item.get("kind") == "formula"), {})
        formula = str((visual.get("parameters") or {}).get("formula") or "")
        _shape(slide, artifact_x, artifact_y, artifact_w, artifact_h, theme["canvas"], radius=True, line=theme["chart_bg"])
        _text(slide, "关键公式", 6.16, 2.20, 1.4, 0.25, 11, theme["accent"], bold=True)
        _text(
            slide,
            _format_formula_text(formula),
            6.16,
            3.05,
            6.02,
            1.75,
            26 if len(formula) <= 72 else 21,
            theme["title"],
            align="center",
            font=theme["math_font"],
            east_asian_font=theme["body_east_asian_font"],
        )
    elif artifact_kind == "table":
        visual = next((item for item in unit.visuals if item.get("kind") == "table"), {})
        parameters = visual.get("parameters") or {}
        _text(slide, "核验对照", 5.86, 1.94, 1.4, 0.25, 11, theme["accent"], bold=True)
        _table(
            slide,
            [str(value) for value in parameters.get("headers") or []],
            [[str(value) for value in row] for row in parameters.get("rows") or []],
            artifact_x,
            2.34,
            artifact_w,
            4.08,
            theme,
        )


def _split_ordered_step(value: str) -> tuple[str, str]:
    clean = str(value or "").strip()
    parts = re.split(r"\s*[:：]\s*", clean, maxsplit=1)
    if len(parts) == 2 and parts[0] and parts[1] and len(parts[0]) <= 42:
        return parts[0].strip(), parts[1].strip()
    return clean, ""


def _display_code_language(value: str) -> str:
    normalized = str(value or "").strip().casefold()
    return {
        "csharp": "C#",
        "python": "Python",
        "javascript": "JavaScript",
        "typescript": "TypeScript",
        "java": "Java",
        "cpp": "C++",
        "c++": "C++",
        "sql": "SQL",
        "bash": "Bash",
    }.get(normalized, normalized.upper() if normalized else "")


def _v6_code_header(unit: SlideSpec, code: SlideBlockSpec | None) -> str:
    if code is None or not unit.quality.get("v6_template_layout_id"):
        return ""
    language = _display_code_language(str(code.metadata.get("code_language") or ""))
    continuation_index = int(
        code.metadata.get("code_chunk_index")
        or unit.quality.get("v6_continuation_index")
        or 1
    )
    continuation_count = int(
        code.metadata.get("code_chunk_count")
        or unit.quality.get("v6_continuation_count")
        or 1
    )
    progress = (
        f"{continuation_index}/{continuation_count}"
        if continuation_count > 1
        else ""
    )
    return " · ".join(value for value in (language, progress) if value)


def _render_code_reading_frame(
    slide: Any,
    unit: SlideSpec,
    code: SlideBlockSpec | None,
    *,
    x: float,
    header_y: float,
    code_y: float,
    width: float,
    height: float,
    text_color: str,
    label_color: str = "AEB6D0",
) -> None:
    code_text = str(code.content if code else "")
    is_v6 = bool(unit.quality.get("v6_template_layout_id"))
    header = _v6_code_header(unit, code)
    if not is_v6:
        header = str(code.metadata.get("language") or "code").upper() if code else "CODE"
    if header:
        _text(
            slide,
            header,
            x,
            header_y,
            max(1.4, width),
            0.28,
            10,
            label_color,
            bold=True,
            font=CODE_FONT,
        )
    if not is_v6:
        _text(
            slide,
            code_text,
            x,
            code_y,
            width,
            height,
            16,
            text_color,
            font=CODE_FONT,
            literal=True,
        )
        return
    start_line = max(1, int(code.metadata.get("code_start_line") or 1)) if code else 1
    line_count = len(code_text.split("\n"))
    line_numbers = "\n".join(
        str(start_line + index)
        for index in range(line_count)
    )
    gutter = _text(
        slide,
        line_numbers,
        x,
        code_y,
        0.46,
        height,
        16,
        "73839B",
        align="right",
        font=CODE_FONT,
        literal=True,
    )
    gutter.name = f"{gutter.name} [v6-code-line-numbers]"
    _shape(
        slide,
        x + 0.58,
        code_y,
        0.012,
        height,
        "34465C",
        radius=False,
    )
    _text(
        slide,
        code_text,
        x + 0.76,
        code_y,
        max(0.5, width - 0.76),
        height,
        16,
        text_color,
        font=CODE_FONT,
        literal=True,
    )


def _render_code(slide: Any, unit: SlideSpec, theme: dict[str, str]) -> None:
    _heading(slide, unit, theme)
    code = _find_block(unit, "code")
    code_text = code.content if code else ""
    insight_blocks = [block for block in unit.blocks if block is not code]
    items = [
        item
        for block in insight_blocks
        for item in (block.items or [block.content])
        if str(item or "").strip()
    ][:5]
    code_panel_width = 7.48 if items else 11.80
    code_text_width = 6.9 if items else 11.18
    evidence_style = _theme_text_box_style(
        theme,
        "evidence",
        fill=theme["code"],
        border=theme["code"],
        accent="74B4FF",
        text="F5F7FF",
    )
    _semantic_panel(slide, 0.76, 1.75, code_panel_width, 4.72, evidence_style, rail=False)
    _render_code_reading_frame(
        slide,
        unit,
        code,
        x=1.05,
        header_y=2.02,
        code_y=2.48,
        width=code_text_width,
        height=3.6,
        text_color=evidence_style["text"],
    )
    if not items:
        return
    note_style = _theme_text_box_style(
        theme,
        "note",
        fill=theme["canvas"],
        border=theme["chart_bg"],
        accent=theme["green"],
    )
    _semantic_panel(slide, 8.52, 1.75, 4.04, 4.72, note_style)
    _text(slide, "阅读线索", 8.86, 2.08, 1.7, 0.32, 12, note_style["accent"], bold=True)
    _bullets(slide, items, 8.86, 2.48, 3.32, 3.65, 16, note_style["text"], note_style["accent"])


def _render_misconception(slide: Any, unit: SlideSpec, theme: dict[str, str]) -> None:
    _heading(slide, unit, theme)
    mistakes = [
        item for block in unit.blocks if block.type == "misconception"
        for item in (block.items or [block.content]) if item
    ]
    correction = unit.key_message or "回到定义、条件和可验证证据进行判断。"
    wrong_style = _theme_text_box_style(
        theme,
        "misconception",
        fill=theme["red_soft"],
        border=theme["red_soft"],
        accent=theme["red"],
    )
    repair_style = _theme_text_box_style(
        theme,
        "feedback",
        fill=theme["green_soft"],
        border=theme["green_soft"],
        accent=theme["green"],
    )
    _semantic_panel(slide, 0.78, 1.93, 5.65, 4.28, wrong_style)
    _text(slide, "容易这样想", 1.13, 2.25, 2.0, 0.34, 12, wrong_style["accent"], bold=True)
    _bullets(slide, mistakes[:4] or ["忽略条件，只记结论"], 1.13, 2.88, 4.92, 2.62, 16, wrong_style["text"], wrong_style["accent"])
    _semantic_panel(slide, 6.71, 1.93, 5.84, 4.28, repair_style)
    _text(slide, "应当这样判断", 7.06, 2.25, 2.2, 0.34, 12, repair_style["accent"], bold=True)
    _text(slide, correction, 7.06, 2.96, 5.08, 1.62, 21, repair_style["text"], bold=True)
    _text(slide, "用反例、边界或独立作答确认理解。", 7.06, 5.08, 4.9, 0.38, 12, theme["muted"])


def _render_practice(slide: Any, unit: SlideSpec, theme: dict[str, str]) -> None:
    _heading(slide, unit, theme)
    exercise = _find_block(unit, "exercise") or (unit.blocks[0] if unit.blocks else None)
    question_style = _theme_text_box_style(
        theme,
        "standard",
        fill=theme["canvas"],
        border=theme["chart_bg"],
        accent=theme["accent"],
    )
    check_style = _theme_text_box_style(
        theme,
        "practice",
        fill=theme["amber_soft"],
        border=theme["amber_soft"],
        accent=theme["amber"],
    )
    _semantic_panel(slide, 0.78, 1.82, 7.52, 4.64, question_style)
    _text(slide, exercise.title if exercise and exercise.title else "先独立作答", 1.12, 2.13, 2.4, 0.35, 12, question_style["accent"], bold=True)
    if exercise and exercise.items:
        _bullets(slide, exercise.items, 1.12, 2.75, 6.78, 2.95, 15, question_style["text"], question_style["accent"])
    else:
        _text(slide, exercise.content if exercise else unit.key_message, 1.12, 2.75, 6.78, 2.95, 17, question_style["text"], bold=True)
    _semantic_panel(slide, 8.58, 1.82, 3.97, 4.64, check_style)
    _text(slide, "检查标准", 8.91, 2.13, 1.7, 0.34, 12, check_style["accent"], bold=True)
    checks = [item for block in unit.blocks if block is not exercise for item in (block.items or [block.content]) if item]
    _bullets(slide, checks[:4] or ["能说明理由", "能处理边界", "能独立完成"], 8.91, 2.82, 3.30, 2.82, 13, check_style["text"], check_style["accent"])


def _render_practice_feedback(
    slide: Any,
    unit: SlideSpec,
    theme: dict[str, str],
) -> None:
    """Align every task with its answer instead of creating imbalanced columns."""
    _heading(slide, unit, theme)
    exercise = _find_block(unit, "exercise") or (
        unit.blocks[0] if unit.blocks else None
    )
    feedback_blocks = [
        block
        for block in unit.blocks
        if block is not exercise
    ]
    prompts = []
    if exercise:
        prompts = list(exercise.items) if exercise.items else [exercise.content]
    prompts = [value for value in prompts if value][:3]
    checks = [
        value
        for block in feedback_blocks
        for value in (block.items or [block.content])
        if value
    ]
    checks = checks[:3]
    prompt_ids = [
        str(value)
        for value in (
            (exercise.metadata if exercise else {}).get("question_ids") or []
        )
    ]
    answers_by_question: dict[str, str] = {}
    for block in feedback_blocks:
        answer_ids = [
            str(value)
            for value in block.metadata.get("answer_for_question_ids") or []
        ]
        answer_values = list(block.items) if block.items else [block.content]
        for answer_id, answer_value in zip(answer_ids, answer_values):
            if answer_id and answer_value:
                answers_by_question[answer_id] = answer_value
    if not prompts:
        prompts = [unit.key_message or unit.takeaway]
    prompts = [value for value in prompts if value]
    feedback_mode = str((unit.quality or {}).get("feedback_mode") or "paired")
    question_style = _theme_text_box_style(
        theme,
        "message",
        fill=theme["accent_soft"],
        border=theme["chart_bg"],
        accent=theme["accent"],
    )
    answer_style = _theme_text_box_style(
        theme,
        "feedback",
        fill=theme["green_soft"],
        border=theme["chart_bg"],
        accent=theme["green"],
    )
    if feedback_mode == "shared_evidence":
        question_count = max(len(prompts), 1)
        question_gap = 0.34
        question_width = (11.0 - question_gap * (question_count - 1)) / question_count
        _shape(slide, 0.82, 1.9, 11.72, 0.018, theme["chart_bg"], radius=False)
        for index, prompt in enumerate(prompts):
            x = 1.05 + index * (question_width + question_gap)
            _semantic_panel(slide, x - 0.18, 2.02, question_width + 0.1, 1.64, question_style)
            _text(
                slide,
                f"问题 {index + 1:02d}",
                x,
                2.13,
                1.45,
                0.28,
                11,
                question_style["accent"],
                bold=True,
            )
            _text(
                slide,
                prompt,
                x,
                2.62,
                question_width - 0.12,
                1.02,
                17 if len(prompt) <= 80 else 16,
                question_style["text"],
                bold=True,
            )
        _semantic_panel(slide, 0.82, 4.04, 11.72, 2.16, answer_style)
        _text(
            slide,
            "判断依据",
            1.05,
            4.28,
            1.45,
            0.28,
            11,
            answer_style["accent"],
            bold=True,
        )
        evidence_count = max(len(checks), 1)
        evidence_gap = 0.42
        evidence_width = (11.0 - evidence_gap * (evidence_count - 1)) / evidence_count
        for index, evidence in enumerate(checks):
            x = 1.05 + index * (evidence_width + evidence_gap)
            _text(
                slide,
                evidence,
                x,
                4.78,
                evidence_width - 0.12,
                1.42,
                16,
                answer_style["text"],
            )
        return
    row_count = max(len(prompts), 1)
    row_height = 4.48 / row_count
    for index, prompt in enumerate(prompts):
        y = 1.9 + index * row_height
        question_id = prompt_ids[index] if index < len(prompt_ids) else ""
        answer = (
            answers_by_question.get(question_id, "")
            or (checks[index] if index < len(checks) else "")
        )
        _semantic_panel(slide, 0.82, y + 0.06, 5.9, row_height - 0.12, question_style)
        _semantic_panel(slide, 6.92, y + 0.06, 5.62, row_height - 0.12, answer_style)
        _text(
            slide,
            f"问题 {index + 1:02d}",
            1.05,
            y + 0.22,
            1.45,
            0.28,
            11,
            question_style["accent"],
            bold=True,
        )
        _text(
            slide,
            "回答与判断依据",
            7.25,
            y + 0.22,
            2.0,
            0.28,
            11,
            answer_style["accent"],
            bold=True,
        )
        text_top = y + 0.72
        _text(
            slide,
            prompt,
            1.05,
            text_top,
            5.55,
            max(0.7, row_height - 1.0),
            17 if len(prompt) <= 80 else 16,
            question_style["text"],
            bold=True,
        )
        if answer:
            _text(
                slide,
                answer,
                7.25,
                text_top,
                4.85,
                max(0.7, row_height - 1.0),
                16,
                answer_style["text"],
            )


def _render_recap(slide: Any, unit: SlideSpec, theme: dict[str, str]) -> None:
    _add_theme_visual_asset(slide, theme, "recap")
    visible_text = _visible_source_text(unit)
    if (
        (not visible_text and not unit.key_message)
        or (
            len(unit.blocks) <= 1
            and len(visible_text.strip()) < 40
            and len(unit.key_message.strip()) < 40
        )
    ):
        _render_navigation_statement(slide, unit, theme)
        return
    _heading(slide, unit, theme)
    blocks = unit.blocks[:3]
    if len(blocks) == 1 and len(blocks[0].items) > 1:
        for index, value in enumerate(blocks[0].items[:6]):
            column = index % 2
            row = index // 2
            x = 0.82 + column * 5.92
            y = 1.82 + row * 1.43
            accent = [theme["accent"], theme["green"], theme["amber"]][row % 3]
            soft = [theme["accent_soft"], theme["green_soft"], theme["amber_soft"]][row % 3]
            _shape(slide, x, y, 5.66, 1.18, soft, radius=True)
            _text(slide, f"{index + 1:02d}", x + 0.26, y + 0.44, 0.52, 0.22, 12, accent, bold=True, align="center")
            _text(slide, value, x + 0.92, y + 0.28, 4.38, 0.62, 17, theme["ink"], bold=True)
        return
    for index, block in enumerate(blocks):
        y = 1.8 + index * 1.58
        accent = [theme["accent"], theme["green"], theme["amber"]][index % 3]
        soft = [theme["accent_soft"], theme["green_soft"], theme["amber_soft"]][index % 3]
        _shape(slide, 0.82, y, 11.72, 1.36, soft, radius=True)
        _text(slide, block.title or f"带走 {index + 1}", 1.14, y + 0.24, 2.0, 0.31, 12, accent, bold=True)
        values = block.items or [block.content]
        _text(slide, " · ".join(value for value in values if value), 3.22, y + 0.18, 8.83, 0.94, 14, theme["ink"], bold=True)
    if unit.key_message:
        _text(slide, unit.key_message, 0.86, 6.27, 11.4, 0.38, 13, theme["muted"], bold=True, align="center")


def _render_chapter_recap(
    slide: Any,
    unit: SlideSpec,
    theme: dict[str, str],
) -> None:
    """Render a compact memory path instead of generic recap cards."""
    _add_theme_visual_asset(slide, theme, "recap")
    _heading(slide, unit, theme)
    items = _all_items(unit)
    if not items:
        items = [
            block.content
            for block in unit.blocks
            if block.content
        ]
    items = items[:4]
    if not items:
        _render_navigation_statement(slide, unit, theme)
        return
    # A 2x2 memory grid keeps complete recap claims readable. The previous
    # single horizontal rail squeezed four or five paragraphs into narrow
    # columns and visually clipped otherwise complete source text.
    accents = (theme["accent"], theme["green"], theme["amber"])
    for index, item in enumerate(items):
        column = index % 2
        row = index // 2
        x = 0.98 + column * 5.78
        y = 1.82 + row * 1.82
        accent = accents[index % len(accents)]
        _shape(slide, x, y, 0.58, 0.58, accent, radius=True)
        _text(
            slide,
            str(index + 1),
            x,
            y + 0.17,
            0.58,
            0.18,
            11,
            "FFFFFF",
            bold=True,
            align="center",
        )
        _text(
            slide,
            item,
            x + 0.82,
            y + 0.03,
            4.62,
            1.35,
            17 if len(item) <= 72 else 16,
            theme["ink"],
            bold=True,
        )
    takeaway = unit.key_message or unit.takeaway
    if takeaway:
        _shape(slide, 0.98, 5.62, 0.08, 0.62, theme["accent"], radius=False)
        _text(
            slide,
            takeaway,
            1.32,
            5.66,
            10.65,
            0.52,
            16,
            theme["muted"],
            bold=True,
        )


def _render_course_synthesis(
    slide: Any,
    unit: SlideSpec,
    theme: dict[str, str],
) -> None:
    """Render the whole-course route as one connected synthesis."""
    _add_theme_visual_asset(slide, theme, "recap")
    _heading(slide, unit, theme)
    items = _all_items(unit)
    if not items:
        items = [
            block.content
            for block in unit.blocks
            if block.content
        ]
    items = items[:6]
    if not items:
        _render_navigation_statement(slide, unit, theme)
        return
    core = unit.key_message or unit.takeaway or unit.title
    _text(slide, "课程主线", 0.92, 1.92, 1.4, 0.3, 12, theme["accent"], bold=True)
    _text(
        slide,
        core,
        0.92,
        2.46,
        3.08,
        2.2,
        23 if len(core) <= 36 else 18,
        theme["ink"],
        bold=True,
    )
    _shape(slide, 4.3, 1.86, 0.035, 4.5, theme["chart_bg"], radius=False)
    row_height = 3.95 / max(len(items), 1)
    for index, item in enumerate(items):
        y = 1.94 + index * row_height
        accent = (theme["accent"], theme["green"], theme["amber"])[index % 3]
        _text(
            slide,
            f"{index + 1:02d}",
            4.72,
            y + 0.08,
            0.52,
            0.24,
            11,
            accent,
            bold=True,
            align="center",
        )
        _text(
            slide,
            item,
            5.52,
            y,
            6.34,
            max(0.45, row_height - 0.08),
            16,
            theme["ink"],
            bold=True,
        )


def _render_appendix(slide: Any, unit: SlideSpec, theme: dict[str, str]) -> None:
    _heading(slide, unit, theme)
    _shape(
        slide,
        0.78,
        1.82,
        11.78,
        4.93,
        theme["canvas"],
        radius=True,
        line=theme["chart_bg"],
    )
    _shape(slide, 0.78, 1.82, 0.09, 4.93, theme["accent"], radius=False)
    values: list[str] = []
    for block in unit.blocks:
        if block.title:
            values.append(block.title)
        if block.items:
            values.extend(f"• {item}" for item in block.items if item)
        elif block.content:
            values.append(block.content)
    body = "\n\n".join(values)
    if len(body) > 380:
        left, right = _balanced_text_columns(body)
        _text(
            slide,
            left,
            1.13,
            2.12,
            5.22,
            4.3,
            16,
            theme["ink"],
            font=theme["body_font"],
            east_asian_font=theme["body_east_asian_font"],
        )
        _shape(slide, 6.55, 2.13, 0.02, 4.25, theme["chart_bg"], radius=False)
        _text(
            slide,
            right,
            6.83,
            2.12,
            5.18,
            4.3,
            16,
            theme["ink"],
            font=theme["body_font"],
            east_asian_font=theme["body_east_asian_font"],
        )
    else:
        _text(
            slide,
            body,
            1.13,
            2.12,
            11.02,
            4.3,
            16,
            theme["ink"],
            font=theme["body_font"],
            east_asian_font=theme["body_east_asian_font"],
        )


def _balanced_text_columns(value: str) -> tuple[str, str]:
    paragraphs = [item.strip() for item in value.split("\n\n") if item.strip()]
    if len(paragraphs) < 2:
        midpoint = max(1, len(value) // 2)
        split_at = value.rfind("。", 0, midpoint)
        if split_at < 1:
            split_at = value.rfind("；", 0, midpoint)
        if split_at < 1:
            split_at = midpoint
        return value[: split_at + 1].strip(), value[split_at + 1 :].strip()
    target = len(value) / 2
    left: list[str] = []
    right: list[str] = []
    left_size = 0
    for paragraph in paragraphs:
        if left and left_size + len(paragraph) > target:
            right.append(paragraph)
        elif right:
            right.append(paragraph)
        else:
            left.append(paragraph)
            left_size += len(paragraph) + 2
    return "\n\n".join(left), "\n\n".join(right)


def _heading_font_size(
    value: str,
    *,
    width_inches: float,
    max_lines: int,
) -> int:
    """Choose the largest title size that fits without a tiny final line."""

    width_pt = max(1.0, (width_inches - 0.02) * 72)
    allowed_lines = max(1, int(max_lines or 1))
    for size in (35, 33, 31, 29):
        if _wrapped_line_count(
            value,
            width_pt=width_pt,
            font_size_pt=size,
        ) > allowed_lines:
            continue
        if _wrapped_text_has_orphan_last_line(
            value,
            width_pt=width_pt,
            font_size_pt=size,
        ):
            continue
        return size
    return 29


def _heading(slide: Any, unit: SlideSpec, theme: dict[str, str]) -> None:
    heading_mode = str(unit.quality.get("heading_mode") or "full")
    section_label = str(unit.quality.get("section_label") or "").strip()
    if _uses_source_only_audience_labels(unit):
        eyebrow = section_label or str(unit.eyebrow or "").strip()
    else:
        eyebrow = str(
            section_label
            if heading_mode == "hidden"
            else unit.eyebrow or unit.slide_purpose
        )
    heading = _display_heading(unit)
    heading_size = _heading_font_size(
        heading,
        width_inches=11.72,
        max_lines=max(1, int(unit.quality.get("v6_title_max_lines") or 2)),
    )
    if eyebrow:
        _text(slide, eyebrow, 0.78, 0.42, 8.8, 0.22, 11, theme["accent"], bold=True)
    _text(
        slide, heading, 0.78, 0.70, 11.72, 1.16, heading_size, theme["title"], bold=True,
        font=theme["title_font"], east_asian_font=theme["title_east_asian_font"],
    )
    _shape(slide, 0.78, 1.86, 0.72, 0.04, theme["accent"], radius=False)
    _shape(slide, 1.58, 1.86, 0.08, 0.04, theme["green"], radius=False)


def _footer(slide: Any, unit: SlideSpec, page: int, total: int, theme: dict[str, str]) -> None:
    section = unit.section_id or (
        "" if _uses_source_only_audience_labels(unit) else "COURSE"
    )
    if section:
        _text(slide, section, 0.78, 7.1, 2.4, 0.2, 8, theme["muted"], font="Aptos Mono")
    image_source = str(unit.quality.get("image_source_short") or "").strip()
    if image_source:
        _text(
            slide,
            f"图源：{image_source}",
            3.12,
            7.1,
            7.9,
            0.2,
            8,
            theme["muted"],
            align="center",
        )
    _text(slide, f"{page:02d} / {total:02d}", 11.48, 7.1, 1.02, 0.2, 8, theme["muted"], align="right", font="Aptos Mono")


def _fill_background(slide: Any, color: str) -> None:
    from pptx.dml.color import RGBColor

    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor.from_string(color)
    # Keep the legacy editable background shape as the first object while
    # staying a fraction inside the canvas to avoid false overflow at full bleed.
    _shape(slide, 0.01, 0.01, 13.30, 7.47, color, radius=False)


def _shape(
    slide: Any,
    x: float,
    y: float,
    width: float,
    height: float,
    fill: str,
    *,
    radius: bool,
    line: str | None = None,
) -> Any:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    if line:
        shape.line.color.rgb = RGBColor.from_string(line)
        shape.line.width = Inches(0.008)
    else:
        shape.line.fill.background()
    if radius and hasattr(shape, "adjustments") and len(shape.adjustments):
        shape.adjustments[0] = 0.08
    return shape


def _circle(slide: Any, x: float, y: float, size: float, fill: str) -> Any:
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches

    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x),
        Inches(y),
        Inches(size),
        Inches(size),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(fill)
    shape.line.fill.background()
    return shape


def _text(
    slide: Any,
    value: str,
    x: float,
    y: float,
    width: float,
    height: float,
    size: int,
    color: str,
    *,
    bold: bool = False,
    align: str = "left",
    font: str = BODY_FONT,
    east_asian_font: str = BODY_EAST_ASIAN_FONT,
    literal: bool = False,
) -> Any:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.01)
    frame.margin_top = frame.margin_bottom = Inches(0.01)
    frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]
    paragraph.text = (
        str(value or "")
        if literal
        else _display_text(str(value or ""))
    )
    _configure_font(paragraph.font, font, east_asian_font)
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = RGBColor.from_string(color)
    paragraph.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    paragraph.space_after = Pt(0)
    return box


def _bullets(
    slide: Any,
    items: list[str],
    x: float,
    y: float,
    width: float,
    height: float,
    size: int,
    color: str,
    bullet_color: str,
) -> Any:
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.01)
    for index, value in enumerate([_display_text(str(item)) for item in items if str(item).strip()]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"•  {value}"
        _configure_font(paragraph.font, BODY_FONT)
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = RGBColor.from_string(color)
        paragraph.space_after = Pt(max(5, size * 0.55))
        if paragraph.runs:
            paragraph.runs[0].font.color.rgb = RGBColor.from_string(bullet_color if len(value) < 1 else color)
    return box


def _table(
    slide: Any,
    headers: list[str],
    rows: list[list[str]],
    x: float,
    y: float,
    width: float,
    height: float,
    theme: dict[str, str],
) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    column_count = max(1, len(headers), max((len(row) for row in rows), default=0))
    row_count = 1 + max(1, len(rows))
    table = slide.shapes.add_table(
        row_count,
        column_count,
        Inches(x),
        Inches(y),
        Inches(width),
        Inches(height),
    ).table
    for column in table.columns:
        column.width = Inches(width / column_count)
    values = [headers or ["比较项"], *rows]
    font_size = 16
    cell_horizontal_margin = 0.1
    cell_vertical_margin = 0.07
    column_text_width_pt = max(
        12.0,
        (width / column_count - cell_horizontal_margin * 2) * 72.0,
    )
    required_row_heights_pt: list[float] = []
    for row_index in range(row_count):
        row_values = values[row_index] if row_index < len(values) else []
        maximum_lines = max(
            1,
            max(
                (
                    _wrapped_line_count(
                        _display_text(str(row_values[column_index])),
                        width_pt=column_text_width_pt,
                        font_size_pt=font_size,
                    )
                    for column_index in range(min(column_count, len(row_values)))
                ),
                default=1,
            ),
        )
        required_row_heights_pt.append(
            maximum_lines * font_size * 1.22 + cell_vertical_margin * 2 * 72.0
        )
    available_height_pt = height * 72.0
    required_height_pt = sum(required_row_heights_pt)
    if required_height_pt <= available_height_pt:
        extra_per_row = (available_height_pt - required_height_pt) / row_count
        row_heights_pt = [value + extra_per_row for value in required_row_heights_pt]
    else:
        # Shrinking rows below their measured requirement creates a valid PPTX
        # file with clipped cells.  Capacity must instead be handled by the V6
        # compiler's safe pagination contract and surfaced if it is violated.
        raise ValueError("table_render_capacity_exceeded")
    for row, row_height_pt in zip(table.rows, row_heights_pt):
        row.height = Pt(row_height_pt)
    for row_index in range(row_count):
        for column_index in range(column_count):
            cell = table.cell(row_index, column_index)
            cell.text = (
                _display_text(values[row_index][column_index])
                if row_index < len(values) and column_index < len(values[row_index])
                else ""
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor.from_string(
                theme["accent_soft"]
                if row_index == 0
                else (theme["surface"] if row_index % 2 else theme["canvas"])
            )
            cell.margin_left = cell.margin_right = Inches(cell_horizontal_margin)
            cell.margin_top = cell.margin_bottom = Inches(cell_vertical_margin)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            paragraph = cell.text_frame.paragraphs[0]
            _configure_font(paragraph.font, BODY_FONT)
            paragraph.font.size = Pt(font_size)
            paragraph.font.bold = row_index == 0
            paragraph.font.color.rgb = RGBColor.from_string(theme["accent"] if row_index == 0 else theme["ink"])
            paragraph.alignment = PP_ALIGN.LEFT


def _display_heading(unit: SlideSpec) -> str:
    title = _format_formula_text(str(unit.title or "").strip())
    takeaway = str(unit.takeaway or "").strip()
    visual_kind = str(unit.visuals[0].get("kind") or "") if unit.visuals else ""
    if title and not _is_generic_heading(title):
        return _heading_excerpt(title)
    if not takeaway:
        return _heading_excerpt(title)
    if (
        visual_kind == "formula"
        or takeaway.startswith(("$", "\\[", "\\("))
        or re.search(r"\\[A-Za-z]+", takeaway)
        or re.fullmatch(r"[\d\s.、:：()（）-]+", takeaway)
    ):
        return _heading_excerpt(title, limit=46)
    return _heading_excerpt(_format_formula_text(takeaway))


def _is_generic_heading(value: str) -> bool:
    normalized = re.sub(r"[\s:：|/\\_-]+", "", str(value or "")).lower()
    return normalized in {
        "课程正文",
        "课程内容",
        "正文",
        "内容",
        "未命名",
        "body",
        "content",
    }


def _heading_excerpt(value: str, limit: int | None = None) -> str:
    """Choose a complete audience-facing title phrase without an ellipsis."""
    clean = " ".join(str(value or "").split()).strip("，,；;：:。… ")
    # V5 compiles and validates a complete audience-facing title before the
    # renderer runs. A second character cut here can create dangling fragments.
    return clean


def _configure_font(font: Any, latin_font: str, east_asian_font: str = BODY_EAST_ASIAN_FONT) -> None:
    """Write both Latin and East Asian typefaces into DrawingML.

    python-pptx only writes ``a:latin`` through ``Font.name``. Explicit
    ``a:ea`` prevents Chinese text from disappearing in renderers that do not
    inherit a usable theme font.
    """
    from pptx.oxml.ns import qn

    font.name = latin_font
    properties = font._element
    east_asian = properties.find(qn("a:ea"))
    if east_asian is None:
        east_asian = properties.makeelement(qn("a:ea"))
        properties.append(east_asian)
    east_asian.set("typeface", east_asian_font)


def _find_block(unit: SlideSpec, block_type: str) -> SlideBlockSpec | None:
    return next((block for block in unit.blocks if block.type == block_type), None)


def _all_items(unit: SlideSpec) -> list[str]:
    return [item for block in unit.blocks for item in block.items if item]


def _block_content(blocks: list[SlideBlockSpec], index: int) -> str:
    return blocks[index].content if len(blocks) > index else ""


def _chapter_number(title: str) -> str:
    import re

    match = re.search(r"(\d+)", title)
    return f"{int(match.group(1)):02d}" if match else "•"


__all__ = ["SlideDeckQualityError", "THEMES", "export_structured_slide_deck", "validate_theme"]
