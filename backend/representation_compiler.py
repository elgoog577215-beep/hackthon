"""Deterministic course projections into same-source teaching representations."""

from __future__ import annotations

import re
import tempfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from course_document import CourseBlock, CourseDocument, stable_hash
from course_revisions import revision_vector_for_document
from teaching_representations import (
    RepresentationPlan,
    SourceBinding,
    TeachingRepresentation,
    TeachingRepresentationRepository,
    TeachingRepresentationSpec,
    source_binding_for_document,
)

REPRESENTATION_COMPILER_VERSION = "same_source_compiler_v1"
CORE_TYPES = ("outline", "lesson_plan", "handout", "practice_sheet", "slide_deck")


def compile_core_representations(
    document: CourseDocument,
    course_data: dict[str, Any],
    repository: TeachingRepresentationRepository,
) -> dict[str, Any]:
    now = datetime.now(timezone.utc).isoformat()
    vector = revision_vector_for_document(document).revisions
    existing_by_type = {
        item.representation_type: item
        for item in repository.load(document.course_id).representations
    }
    plan = RepresentationPlan(
        plan_id=stable_hash({
            "course_id": document.course_id,
            "revision": document.document_revision,
            "types": CORE_TYPES,
        }, prefix="rpl_"),
        course_id=document.course_id,
        source_revision_vector=vector,
        target_scope={"kind": "course"},
        requested_representations=list(CORE_TYPES),
        knowledge_refs=_course_knowledge_refs(document, course_data),
        pedagogical_reasons=[
            "为学习者、课程维护者和课堂展示提供同一课程语义的不同用途表达",
            "保持大纲、讲义、教案、正式题目和幻灯之间的来源与修订一致",
        ],
        cost_class="medium",
        accessibility_requirements=["阅读顺序", "替代文本", "可导出文本"],
        quality_requirements=["来源完整", "跨产物术语一致", "正式题目只引用不复制所有权"],
        fallback_chain=["handout", "outline"],
        status="ready",
    )
    repository.register_plan(plan)

    payloads = {
        "outline": _outline_spec(document),
        "lesson_plan": _lesson_plan_spec(document, course_data),
        "handout": _handout_spec(document),
        "practice_sheet": _practice_sheet_spec(document, course_data),
        "slide_deck": _slide_deck_spec(document, course_data),
    }
    built: list[dict[str, Any]] = []
    for representation_type, payload in payloads.items():
        unit_bindings = _unit_bindings_for_payload(document, payload)
        bindings = _dedupe_bindings([
            binding
            for values in unit_bindings.values()
            for binding in values
        ])
        spec_payload = {
            "compiler_version": REPRESENTATION_COMPILER_VERSION,
            "representation_type": representation_type,
            "content": payload,
        }
        spec_id = stable_hash({
            "course_id": document.course_id,
            "type": representation_type,
            "source_revision_vector": _combined_revisions(bindings),
            "payload": spec_payload,
        }, prefix="trs_")
        spec_revision = stable_hash(spec_payload, prefix="tsr_")
        unit_count = len(payload.get("units") or payload.get("slides") or payload.get("sections") or [])
        representation_status = "ready" if unit_count else "failed"
        spec = TeachingRepresentationSpec(
            spec_id=spec_id,
            course_id=document.course_id,
            representation_type=representation_type,
            source_bindings=bindings,
            unit_bindings=unit_bindings,
            payload=spec_payload,
            revision=spec_revision,
            created_at=now,
            updated_at=now,
        )
        repository.register_spec(spec)
        representation_id = stable_hash({
            "course_id": document.course_id,
            "type": representation_type,
        }, prefix="trp_")
        representation = TeachingRepresentation(
            representation_id=representation_id,
            course_id=document.course_id,
            representation_type=representation_type,
            source_bindings=bindings,
            source_revision_vector=_combined_revisions(bindings),
            spec_id=spec_id,
            semantic_fingerprint=stable_hash(payload, prefix="sem_"),
            render_fingerprint=stable_hash({
                "spec_revision": spec_revision,
                "renderer": "structured_json_v1",
            }, prefix="rnd_"),
            quality_report_id=stable_hash({
                "spec_revision": spec_revision,
                "quality": "passed" if unit_count else "failed_empty_representation",
            }, prefix="rqr_"),
            revision=stable_hash({
                "spec_revision": spec_revision,
                "source_revision_vector": _combined_revisions(bindings),
            }, prefix="rpr_"),
            status=representation_status,
            stale_unit_ids=[],
            stale_reasons=[] if unit_count else ["empty_representation"],
            created_at=now,
            updated_at=now,
        )
        repository.register_representation(representation)
        built.append({
            "representation_id": representation_id,
            "representation_type": representation_type,
            "spec_id": spec_id,
            "status": representation_status,
            "unit_count": unit_count,
            "rebuilt_unit_ids": list(
                (existing_by_type.get(representation_type).stale_unit_ids)
                if existing_by_type.get(representation_type)
                else []
            ),
        })
    return {"plan_id": plan.plan_id, "representations": built}


def rebuild_core_representations_safely(
    document: CourseDocument,
    course_data: dict[str, Any],
    repository: TeachingRepresentationRepository,
) -> dict[str, Any]:
    """Compile in isolation and publish only a complete, quality-passing set.

    The active registry remains available in its stale state when compilation
    or validation fails. This keeps the last usable lesson plan, handout,
    practice sheet and slide deck accessible without pretending they match the
    latest course revision.
    """
    previous = repository.load(document.course_id)
    stale_before = [
        {
            "representation_id": item.representation_id,
            "representation_type": item.representation_type,
            "spec_id": item.spec_id,
            "stale_unit_ids": list(item.stale_unit_ids),
            "stale_reasons": list(item.stale_reasons),
        }
        for item in previous.representations
        if item.status == "stale"
    ]
    try:
        with tempfile.TemporaryDirectory(prefix="lingzhi-representation-build-") as temp_dir:
            shadow = TeachingRepresentationRepository(temp_dir)
            build = compile_core_representations(document, course_data, shadow)
            candidate = shadow.load(document.course_id)
            current_spec_ids = {item.spec_id for item in candidate.representations}
            current_specs = [
                item for item in candidate.specs if item.spec_id in current_spec_ids
            ]
            quality = validate_compiled_representations(current_specs)
            if not quality["passed"]:
                return {
                    "status": "failed_using_last_available",
                    "quality": quality,
                    "stale_before": stale_before,
                    "last_available": [
                        {
                            "representation_id": item.representation_id,
                            "representation_type": item.representation_type,
                            "spec_id": item.spec_id,
                            "status": item.status,
                        }
                        for item in previous.representations
                    ],
                }
            candidate.applied_revision_event_ids = list(previous.applied_revision_event_ids)
            committed = repository.save(candidate)
            stale_by_type = {
                item["representation_type"]: item["stale_unit_ids"]
                for item in stale_before
            }
            return {
                "status": "synchronized",
                "quality": quality,
                "stale_before": stale_before,
                "rebuilt": [
                    {
                        **item,
                        "rebuilt_unit_ids": stale_by_type.get(item["representation_type"], []),
                    }
                    for item in build["representations"]
                ],
                "registry_revision": committed.registry_revision,
            }
    except Exception as exc:
        return {
            "status": "failed_using_last_available",
            "quality": {
                "passed": False,
                "issues": [{
                    "severity": "critical",
                    "code": "representation_rebuild_failed",
                    "message": str(exc),
                }],
            },
            "stale_before": stale_before,
            "last_available": [
                {
                    "representation_id": item.representation_id,
                    "representation_type": item.representation_type,
                    "spec_id": item.spec_id,
                    "status": item.status,
                }
                for item in previous.representations
            ],
        }


def export_slide_deck_pptx(spec: TeachingRepresentationSpec, output_path: str | Path) -> Path:
    if spec.representation_type != "slide_deck":
        raise ValueError("Only slide deck specs can be exported to pptx")
    from pptx import Presentation
    from pptx.util import Inches, Pt

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slides = ((spec.payload.get("content") or {}).get("slides") or [])
    for unit in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.8), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = str(unit.get("title") or "")
        title_frame.paragraphs[0].font.size = Pt(28)
        title_frame.paragraphs[0].font.bold = True
        body_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.55), Inches(11.4), Inches(4.9))
        frame = body_box.text_frame
        frame.clear()
        for index, bullet in enumerate(unit.get("bullets") or []):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = str(bullet)
            paragraph.level = 0
            paragraph.font.size = Pt(20)
            paragraph.space_after = Pt(10)
        notes = str(unit.get("speaker_notes") or "")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
    path = Path(output_path)
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)
    return path


def validate_compiled_representations(specs: list[TeachingRepresentationSpec]) -> dict[str, Any]:
    issues: list[dict[str, Any]] = []
    required = set(CORE_TYPES)
    present = {spec.representation_type for spec in specs}
    for missing in sorted(required - present):
        issues.append({"severity": "critical", "code": "missing_representation", "target": missing})
    for spec in specs:
        content = spec.payload.get("content") or {}
        units = content.get("units") or content.get("slides") or content.get("sections") or []
        if not units:
            issues.append({
                "severity": "critical",
                "code": "empty_representation",
                "target": spec.representation_type,
            })
        for unit in units:
            unit_id = str(unit.get("unit_id") or "")
            if not unit_id or not spec.unit_bindings.get(unit_id):
                issues.append({
                    "severity": "critical",
                    "code": "missing_source_binding",
                    "target": unit_id or spec.representation_type,
                })
    return {
        "passed": not any(issue["severity"] == "critical" for issue in issues),
        "issues": issues,
        "representation_count": len(specs),
    }


def _outline_spec(document: CourseDocument) -> dict[str, Any]:
    blocks_by_section = _blocks_by_section(document)
    sections = []
    for section in sorted(document.sections, key=lambda item: item.position):
        blocks = blocks_by_section.get(section.section_id, [])
        sections.append({
            "unit_id": f"outline:{section.section_id}",
            "section_id": section.section_id,
            "parent_section_id": section.parent_section_id,
            "title": section.title,
            "level": section.level,
            "position": section.position,
            "learning_objective": section.learning_objective,
            "objective_id": section.objective_id,
            "source_section_ids": [section.section_id],
            "source_block_ids": [block.block_id for block in blocks],
            "block_roles": [block.role for block in blocks],
            "knowledge_refs": _knowledge_refs_for_blocks(blocks),
        })
    return {"title": document.title, "sections": sections}


def _lesson_plan_spec(document: CourseDocument, course_data: dict[str, Any]) -> dict[str, Any]:
    blocks_by_section = _blocks_by_section(document)
    units = []
    for section in _learning_sections(document):
        blocks = blocks_by_section.get(section.section_id, [])
        minutes = max(12, min(45, 6 + len(blocks) * 4))
        units.append({
            "unit_id": f"lesson:{section.section_id}",
            "section_id": section.section_id,
            "title": section.title,
            "learning_objective": section.learning_objective,
            "duration_minutes": minutes,
            "source_block_ids": [block.block_id for block in blocks],
            "activities": [
                {"phase": "导入", "minutes": max(2, minutes // 8), "prompt": f"从已有经验进入“{section.title}”"},
                {"phase": "建构", "minutes": max(6, minutes // 2), "prompt": _section_summary(blocks)},
                {"phase": "检查", "minutes": max(4, minutes // 4), "prompt": f"检查是否达到：{section.learning_objective or section.title}"},
            ],
            "misconceptions": _section_misconceptions(course_data, section.section_id),
            "practice_task_ids": _section_question_ids(course_data, section.section_id),
            "knowledge_refs": _knowledge_refs_for_blocks(blocks),
        })
    return {"title": f"{document.title} 教案", "units": units}


def _handout_spec(document: CourseDocument) -> dict[str, Any]:
    blocks_by_section = _blocks_by_section(document)
    units = []
    for section in _learning_sections(document):
        blocks = blocks_by_section.get(section.section_id, [])
        units.append({
            "unit_id": f"handout:{section.section_id}",
            "section_id": section.section_id,
            "title": section.title,
            "learning_objective": section.learning_objective,
            "source_block_ids": [block.block_id for block in blocks],
            "knowledge_refs": _knowledge_refs_for_blocks(blocks),
            "blocks": [
                {
                    "block_id": block.block_id,
                    "role": block.role,
                    "title": str(block.payload.get("title") or ""),
                    "markdown": str(block.payload.get("markdown") or block.payload.get("text") or ""),
                    "knowledge_refs": list(block.concept_refs),
                }
                for block in blocks if block.status != "retired"
            ],
        })
    return {"title": f"{document.title} 讲义", "units": units}


def _practice_sheet_spec(document: CourseDocument, course_data: dict[str, Any]) -> dict[str, Any]:
    questions = list((course_data.get("learning_assets") or {}).get("questions") or [])
    units = []
    by_section = {section.section_id: section for section in document.sections}
    for question in questions:
        section_id = str(question.get("node_id") or "")
        section = by_section.get(section_id)
        source_blocks = [block.block_id for block in document.blocks if block.section_id == section_id]
        units.append({
            "unit_id": f"practice:{question.get('revision_id') or question.get('question_id')}",
            "section_id": section_id,
            "section_title": section.title if section else section_id,
            "source_block_ids": source_blocks,
            "practice_task_id": question.get("question_id") or question.get("asset_id"),
            "practice_revision_id": question.get("revision_id"),
            "prompt": question.get("prompt"),
            "practice_level": question.get("practice_level"),
            "knowledge_refs": _unique([
                *(question.get("course_knowledge_refs") or []),
                *(question.get("concept_ids") or []),
                *[
                    knowledge_id
                    for block in document.blocks
                    if block.section_id == section_id
                    for knowledge_id in block.concept_refs
                ],
            ]),
            "answer_policy": "separate_answer_key",
        })
    return {"title": f"{document.title} 练习", "units": units}


def _slide_deck_spec(document: CourseDocument, course_data: dict[str, Any]) -> dict[str, Any]:
    blocks_by_section = _blocks_by_section(document)
    slides = [{
        "unit_id": "slide:title",
        "slide_purpose": "orientation",
        "title": document.title,
        "bullets": ["课程目标与学习路径"],
        "speaker_notes": "从课程目标和学习者已有经验开始。",
        "source_keys": ["course_title"],
        "source_block_ids": [],
        "knowledge_refs": [],
    }]
    for section in _learning_sections(document):
        blocks = blocks_by_section.get(section.section_id, [])
        if not blocks:
            continue
        bullets = []
        for block in blocks:
            text = _plain_text(str(block.payload.get("markdown") or block.payload.get("text") or ""))
            if text:
                bullets.append(text[:110])
            if len(bullets) == 4:
                break
        slides.append({
            "unit_id": f"slide:{section.section_id}",
            "slide_purpose": "concept_and_reasoning",
            "section_id": section.section_id,
            "title": section.title,
            "bullets": bullets or [section.learning_objective or section.title],
            "speaker_notes": f"本页目标：{section.learning_objective or section.title}",
            "source_block_ids": [block.block_id for block in blocks],
            "practice_task_ids": _section_question_ids(course_data, section.section_id)[:1],
            "knowledge_refs": _knowledge_refs_for_blocks(blocks),
        })
    return {"title": document.title, "slides": slides, "theme": "lingzhi-light-v1"}


def _unit_bindings_for_payload(
    document: CourseDocument,
    payload: dict[str, Any],
) -> dict[str, list[SourceBinding]]:
    vector = revision_vector_for_document(document).revisions
    blocks_by_id = {block.block_id: block for block in document.blocks}
    result: dict[str, list[SourceBinding]] = {}
    units = payload.get("units") or payload.get("slides") or payload.get("sections") or []
    for unit in units:
        unit_id = str(unit.get("unit_id") or "")
        if not unit_id:
            continue
        bindings: list[SourceBinding] = []
        for block_id in unit.get("source_block_ids") or []:
            block = blocks_by_id.get(str(block_id))
            bindings.append(source_binding_for_document(
                document,
                block_id=str(block_id),
                knowledge_node_ids=list(block.concept_refs) if block else [],
            ))
        for section_id in unit.get("source_section_ids") or []:
            bindings.append(source_binding_for_document(
                document,
                section_id=str(section_id),
                knowledge_node_ids=_unique(unit.get("knowledge_refs") or []),
            ))
        source_keys = [
            str(source_key) for source_key in unit.get("source_keys") or []
            if str(source_key) in vector
        ]
        if source_keys:
            bindings.append(SourceBinding(
                course_id=document.course_id,
                source_revisions={key: vector[key] for key in source_keys},
            ))
        practice_task_id = str(unit.get("practice_task_id") or "")
        practice_revision_id = str(unit.get("practice_revision_id") or "")
        if practice_task_id and practice_revision_id:
            bindings.append(SourceBinding(
                course_id=document.course_id,
                section_id=str(unit.get("section_id") or "") or None,
                knowledge_node_ids=_unique(unit.get("knowledge_refs") or []),
                practice_task_ids=[practice_task_id],
                source_revisions={f"practice:{practice_task_id}": practice_revision_id},
            ))
        result[unit_id] = _dedupe_bindings(bindings or [source_binding_for_document(document)])
    if not result:
        result["__whole__"] = [source_binding_for_document(document)]
    return result


def _dedupe_bindings(bindings: list[SourceBinding]) -> list[SourceBinding]:
    result: list[SourceBinding] = []
    seen: set[str] = set()
    for binding in bindings:
        key = stable_hash(binding.model_dump(mode="json"), prefix="sbd_")
        if key in seen:
            continue
        seen.add(key)
        result.append(binding)
    return result


def _combined_revisions(bindings: list[SourceBinding]) -> dict[str, str]:
    return {
        key: revision
        for binding in bindings
        for key, revision in binding.source_revisions.items()
    }


def _blocks_by_section(document: CourseDocument) -> dict[str, list[CourseBlock]]:
    result: dict[str, list[CourseBlock]] = {}
    for block in sorted(document.blocks, key=lambda item: (item.section_id, item.position)):
        result.setdefault(block.section_id, []).append(block)
    return result


def _learning_sections(document: CourseDocument):
    active_section_ids = {
        block.section_id for block in document.blocks if block.status != "retired"
    }
    return [
        section
        for section in sorted(document.sections, key=lambda item: item.position)
        if section.section_id in active_section_ids
    ]


def _plain_text(markdown: str) -> str:
    text = re.sub(r"```.*?```", "", markdown, flags=re.S)
    text = re.sub(r"[`*_#>\[\]()]", " ", text)
    return " ".join(text.split())


def _section_summary(blocks: list[CourseBlock]) -> str:
    for block in blocks:
        text = _plain_text(str(block.payload.get("summary") or block.payload.get("markdown") or ""))
        if text:
            return text[:180]
    return "围绕本节目标进行概念建构、推导和应用。"


def _section_question_ids(course_data: dict[str, Any], section_id: str) -> list[str]:
    return [
        str(item.get("question_id") or item.get("asset_id") or "")
        for item in ((course_data.get("learning_assets") or {}).get("questions") or [])
        if str(item.get("node_id") or "") == section_id
    ]


def _section_misconceptions(course_data: dict[str, Any], section_id: str) -> list[str]:
    return [
        str(item.get("error_pattern") or "")
        for item in ((course_data.get("learning_assets") or {}).get("misconceptions") or [])
        if str(item.get("node_id") or "") == section_id and item.get("error_pattern")
    ]


def _knowledge_refs_for_blocks(blocks: list[CourseBlock]) -> list[str]:
    return _unique([
        knowledge_id
        for block in blocks
        for knowledge_id in block.concept_refs
    ])


def _course_knowledge_refs(
    document: CourseDocument,
    course_data: dict[str, Any],
) -> list[str]:
    asset_refs = [
        knowledge_id
        for values in (course_data.get("learning_assets") or {}).values()
        if isinstance(values, list)
        for item in values
        if isinstance(item, dict)
        for knowledge_id in item.get("course_knowledge_refs") or []
    ]
    return _unique([
        *[knowledge_id for block in document.blocks for knowledge_id in block.concept_refs],
        *asset_refs,
    ])


def _unique(values: list[Any]) -> list[str]:
    return list(dict.fromkeys(
        str(value).strip() for value in values if str(value).strip()
    ))
