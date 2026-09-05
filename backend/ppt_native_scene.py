"""Execute and audit confirmed scenes. This module never calls a model."""
from __future__ import annotations

from copy import deepcopy
from pathlib import Path
import tempfile

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt

from ppt_layout_execution import file_digest, wrap_text
from ppt_teaching_content import relation_is_directed
from ppt_page_scene import ResolvedPageScene, verify_scene


def _color(value):
    return RGBColor.from_string(value)


def _shape_text(shape, obj, scene):
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = frame.margin_right = Pt(8)
    frame.margin_top = frame.margin_bottom = Pt(6)
    for index, source_paragraph in enumerate(obj.text.split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_before = paragraph.space_after = Pt(0)
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.line_spacing = Pt(obj.font_size * 1.3)
        font = paragraph.font
        font.name, font.size, font.bold, font.color.rgb = scene.execution.font_family, Pt(obj.font_size), obj.bold, _color(obj.color)
        for tag in ("a:ea", "a:cs"):
            el = OxmlElement(tag)
            el.set("typeface", scene.execution.font_family)
            font._rPr.append(el)
        lines = wrap_text(source_paragraph, obj.width - 16, obj.font_size, font_digest=scene.execution.font_sha256)
        for line_index, line in enumerate(lines):
            if line_index:
                paragraph.add_line_break()
            paragraph.add_run().text = line


def _native_target(slide, target):
    shapes = slide.shapes
    for group_id in target.group_path:
        matches = [s for s in shapes if s.shape_id == group_id]
        if len(matches) != 1 or not hasattr(matches[0], "shapes"):
            raise ValueError("native_group_target_missing")
        shapes = matches[0].shapes
    matches = [s for s in shapes if s.shape_id == target.shape_id]
    if len(matches) != 1:
        raise ValueError("native_shape_target_missing_or_ambiguous")
    shape = matches[0]
    if target.row is not None or target.column is not None:
        if target.row is None or target.column is None or not shape.has_table:
            raise ValueError("native_table_target_invalid")
        if target.row >= len(shape.table.rows) or target.column >= len(shape.table.columns):
            raise ValueError("native_table_target_out_of_range")
        cell = shape.table.cell(target.row, target.column)
        if cell.is_spanned:
            raise ValueError("native_table_target_merge_slave")
        return cell
    return shape


def native_target_geometry(slide, target, *, for_execution=False):
    """Measure an identified object in slide coordinates, including groups."""
    shapes = slide.shapes
    ox = oy = 0.0
    sx = sy = 1.0
    for group_id in target.group_path:
        matches = [s for s in shapes if s.shape_id == group_id]
        if len(matches) != 1 or not hasattr(matches[0], "shapes"):
            raise ValueError("native_group_target_missing")
        group = matches[0]
        xf = group._element.grpSpPr.xfrm
        if group.rotation or xf.get("flipH") == "1" or xf.get("flipV") == "1":
            raise ValueError("native_transformed_group_unsupported")
        if not xf.chExt.cx or not xf.chExt.cy:
            raise ValueError("native_group_geometry_invalid")
        gx, gy = xf.ext.cx / xf.chExt.cx, xf.ext.cy / xf.chExt.cy
        if for_execution and (abs(gx - 1) > 1e-6 or abs(gy - 1) > 1e-6):
            raise ValueError("native_scaled_group_unsupported")
        ox += sx * (xf.off.x - gx * xf.chOff.x)
        oy += sy * (xf.off.y - gy * xf.chOff.y)
        sx *= gx
        sy *= gy
        shapes = group.shapes
    matches = [s for s in shapes if s.shape_id == target.shape_id]
    if len(matches) != 1:
        raise ValueError("native_shape_target_missing_or_ambiguous")
    shape = matches[0]
    if for_execution and shape.rotation:
        raise ValueError("native_rotated_text_unsupported")
    x, y, w, h = shape.left, shape.top, shape.width, shape.height
    if target.row is not None or target.column is not None:
        cell = _native_target(slide, target)
        columns, rows = shape.table.columns, shape.table.rows
        x += sum(columns[i].width for i in range(target.column))
        y += sum(rows[i].height for i in range(target.row))
        w = sum(columns[i].width for i in range(target.column, target.column + cell.span_width))
        h = sum(rows[i].height for i in range(target.row, target.row + cell.span_height))
    return tuple(round(v / 12700, 5) for v in (ox + sx * x, oy + sy * y, sx * w, sy * h))


def inspect_native_bindings(slide, bindings):
    """Resolve explicit IDs only. Duplicate or guessed bindings cannot certify."""
    seen, resolved = set(), {}
    for slot, target in bindings.items():
        key = (tuple(target.group_path), target.shape_id, target.row, target.column)
        if key in seen:
            raise ValueError("native_target_bound_twice")
        seen.add(key)
        obj = _native_target(slide, target)
        if target.kind == "image":
            if target.group_path or target.row is not None or not hasattr(obj, "image"):
                raise ValueError("native_image_target_unsupported")
        elif not hasattr(obj, "text_frame"):
            raise ValueError("native_text_target_unsupported")
        geometry = native_target_geometry(slide, target, for_execution=True)
        if target.geometry_pt and any(abs(a - b) > 1 for a, b in zip(target.geometry_pt, geometry)):
            raise ValueError("native_target_geometry_mismatch")
        resolved[slot] = target.model_copy(update={"geometry_pt": geometry})
    return resolved


def inspect_native_connections(slide, execution):
    for key, binding in execution.connections.items():
        if key != f"{binding.source_slot}->{binding.target_slot}":
            raise ValueError("native_connection_key_mismatch")
        targets = [execution.targets.get(slot) for slot in (binding.source_slot, binding.target_slot)]
        if any(t is None or t.kind != "text" or t.group_path or t.row is not None for t in targets):
            raise ValueError("native_connection_endpoint_unsupported")
        matches = [s for s in slide.shapes if s.shape_id == binding.shape_id]
        if len(matches) != 1 or not hasattr(matches[0], 'begin_connect'):
            raise ValueError("native_connection_target_missing")
        shape = matches[0]
        for name, target, site in zip(('stCxn', 'endCxn'), targets, (binding.start_site, binding.end_site), strict=True):
            nodes = shape._element.xpath(f'.//a:{name}')
            if len(nodes) != 1 or nodes[0].get('id') != str(target.shape_id) or nodes[0].get('idx') != str(site):
                raise ValueError("native_connection_endpoint_mismatch")
        geometry = tuple(v / 12700 for v in (shape.begin_x, shape.begin_y, shape.end_x, shape.end_y))
        if any(abs(a - b) > 1 for a, b in zip(geometry, binding.geometry_pt)):
            raise ValueError("native_connection_geometry_mismatch")
        arrows = shape._element.xpath('.//a:tailEnd')
        if any(a.get('type') not in {None, 'none'} for a in shape._element.xpath('.//a:headEnd')) or any(a.get('type') not in {None, 'none', 'triangle'} for a in arrows):
            raise ValueError("native_connection_arrow_unsupported")
        if bool(arrows and arrows[0].get('type') == 'triangle') != binding.directed:
            raise ValueError("native_connection_direction_mismatch")


def _replace_picture(shape, data, geometry=None):
    from io import BytesIO
    _, relationship = shape.part.get_or_add_image_part(BytesIO(data))
    shape._element.blipFill.blip.rEmbed = relationship
    shape.crop_left = shape.crop_right = shape.crop_top = shape.crop_bottom = 0
    if geometry:
        shape.left, shape.top, shape.width, shape.height = (Pt(v) for v in geometry)


def _transparent_picture():
    from io import BytesIO
    from PIL import Image
    buffer = BytesIO()
    Image.new('RGBA', (1, 1), (0, 0, 0, 0)).save(buffer, format='PNG')
    return buffer.getvalue()


def render_scene(slide, scene: ResolvedPageScene, *, assets=None):
    verify_scene(scene)
    shapes_by_id = {}
    for obj in scene.objects:
        if scene.execution.mode == "native_fill":
            target = scene.execution.targets.get(obj.slot_id)
            if target is None:
                raise ValueError("native_template_target_missing")
            shape = _native_target(slide, target)
            if target.kind != obj.kind:
                raise ValueError("native_target_kind_mismatch")
            if obj.kind == "text" and not hasattr(shape, "text_frame"):
                raise ValueError("native_text_target_unsupported")
            # Native slot geometry is a contract, not a suggestion.
            if any(abs(a - b) > 1 for a, b in zip(
                    native_target_geometry(slide, target, for_execution=True),
                    (obj.x, obj.y, obj.width, obj.height))):
                raise ValueError("native_target_geometry_mismatch")
        if obj.kind == "image":
            if assets is None:
                raise ValueError("teaching_asset_repository_missing")
            path = assets.resolve(obj.asset_id)
            if file_digest(path) != obj.asset_digest:
                raise ValueError("teaching_asset_digest_mismatch")
            from PIL import Image
            with Image.open(path) as img:
                iw, ih = img.size
            scale = min(obj.width / iw, obj.height / ih)
            geometry = (obj.x + (obj.width - iw * scale) / 2, obj.y + (obj.height - ih * scale) / 2, iw * scale, ih * scale)
            if scene.execution.mode == "native_fill":
                _replace_picture(shape, path.read_bytes(), geometry)
            else:
                shape = slide.shapes.add_picture(str(path), Pt(geometry[0]), Pt(geometry[1]), width=Pt(geometry[2]), height=Pt(geometry[3]))
        elif scene.execution.mode != "native_fill":
            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(obj.x), Pt(obj.y), Pt(obj.width), Pt(obj.height))
            style = shape._element.find(qn("p:style"))
            if style is not None:
                shape._element.remove(style)
            shape._element.spPr.append(OxmlElement("a:effectLst"))
            shape.fill.solid()
            shape.fill.fore_color.rgb = _color(obj.fill)
            shape.line.color.rgb = _color(scene.accent_color if obj.element_id in scene.emphasized_element_ids else obj.stroke)
            shape.line.width = Pt(1.5 if obj.element_id in scene.emphasized_element_ids else 0)
        if obj.kind == "text":
            _shape_text(shape, obj, scene)
        if hasattr(shape, "name"):
            shape.name = f"teaching:{obj.object_id}"
            shapes_by_id[obj.object_id] = shape
    objects = {o.object_id: o for o in scene.objects}
    if scene.execution.mode == "native_fill":
        visible = {f"{objects[e.source_id].slot_id}->{objects[e.target_id].slot_id}" for e in scene.edges}
        for key, binding in scene.execution.connections.items():
            if key not in visible:
                shape = next(s for s in slide.shapes if s.shape_id == binding.shape_id)
                shape._element.getparent().remove(shape._element)
    for edge in scene.edges:
        key = f"{objects[edge.source_id].slot_id}->{objects[edge.target_id].slot_id}"
        if edge.source_id not in shapes_by_id or edge.target_id not in shapes_by_id:
            raise ValueError("native_cell_connector_unsupported")
        source, target = shapes_by_id[edge.source_id], shapes_by_id[edge.target_id]
        binding = scene.execution.connections.get(key) if scene.execution.mode == "native_fill" else None
        connector = next(s for s in slide.shapes if s.shape_id == binding.shape_id) if binding else slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Pt(edge.x1), Pt(edge.y1), Pt(edge.x2), Pt(edge.y2))
        connector.name = f"relation:{edge.relation_id}"
        style = connector._element.find(qn('p:style'))
        if style is not None:
            connector._element.remove(style)
        for effect in connector._element.spPr.findall(qn('a:effectLst')):
            connector._element.spPr.remove(effect)
        connector._element.spPr.append(OxmlElement('a:effectLst'))
        connector.begin_connect(source, edge.start_site)
        connector.end_connect(target, edge.end_site)
        connector.line.color.rgb = _color(scene.accent_color)
        connector.line.width = Pt(1.7)
        if relation_is_directed(edge.kind):
            for previous in connector._element.xpath('.//a:tailEnd'):
                previous.getparent().remove(previous)
            end = OxmlElement("a:tailEnd")
            end.set("type", "triangle")
            connector.line._get_or_add_ln().append(end)
        if edge.label:
            label = edge.label_object
            shape = _native_target(slide, scene.execution.targets[binding.label_slot]) if binding else slide.shapes.add_textbox(Pt(label.x), Pt(label.y), Pt(label.width), Pt(label.height))
            shape.name = f"relation-label:{edge.relation_id}"
            _shape_text(shape, label, scene)


def _iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if hasattr(shape, "shapes"):
            yield from _iter_shapes(shape.shapes)


def _source_text(shape):
    # Soft wrapping is a:br; original paragraph breaks are separate a:p nodes.
    return "\n".join("".join(p._p.xpath(".//a:t/text()")) for p in shape.text_frame.paragraphs)


def audit_scene(slide, scene: ResolvedPageScene) -> dict:
    shapes = list(_iter_shapes(slide.shapes))
    named = {s.name: s for s in shapes}
    dynamic_names = [s.name for s in shapes if s.name.startswith(("teaching:", "relation:", "relation-label:"))]
    if len(dynamic_names) != len(set(dynamic_names)):
        raise ValueError("export_object_name_ambiguous")
    native = scene.execution.mode == "native_fill"
    expected_names = {f"teaching:{o.object_id}" for o in scene.objects
                      if not native or scene.execution.targets[o.slot_id].row is None}
    actual_names = {n for n in named if n.startswith("teaching:")}
    if expected_names != actual_names:
        raise ValueError("export_visible_element_set_mismatch")
    for obj in scene.objects:
        binding = scene.execution.targets.get(obj.slot_id) if native else None
        shape = _native_target(slide, binding) if binding else named[f"teaching:{obj.object_id}"]
        if obj.kind == "text":
            if not hasattr(shape, "text_frame") or _source_text(shape) != obj.text:
                raise ValueError(f"export_element_text_mismatch:{obj.object_id}")
            if any(abs(a - b) > 1 for a, b in zip(
                native_target_geometry(slide, binding) if binding else (shape.left / 12700, shape.top / 12700, shape.width / 12700, shape.height / 12700),
                (obj.x, obj.y, obj.width, obj.height))):
                raise ValueError(f"export_element_position_mismatch:{obj.object_id}")
            if any(p.font.name != scene.execution.font_family or p.font.size != Pt(obj.font_size) for p in shape.text_frame.paragraphs):
                raise ValueError("export_font_mismatch")
        elif obj.kind == "shape":
            actual = (shape.left / 12700, shape.top / 12700, shape.width / 12700, shape.height / 12700)
            if any(abs(a - b) > 0.01 for a, b in zip(actual, (obj.x, obj.y, obj.width, obj.height))):
                raise ValueError(f"export_chart_scale_mismatch:{obj.object_id}")
            if str(shape.fill.fore_color.rgb) != obj.fill:
                raise ValueError(f"export_shape_fill_mismatch:{obj.object_id}")
        else:
            import hashlib
            if hashlib.sha256(shape.image.blob).hexdigest() != obj.asset_digest:
                raise ValueError("export_asset_identity_mismatch")
            iw, ih = shape.image.size
            scale = min(obj.width / iw, obj.height / ih)
            expected = (obj.x + (obj.width - iw * scale) / 2, obj.y + (obj.height - ih * scale) / 2, iw * scale, ih * scale)
            actual = (shape.left / 12700, shape.top / 12700, shape.width / 12700, shape.height / 12700)
            if any(abs(a - b) > 1 for a, b in zip(actual, expected)) or any(abs(getattr(shape, f"crop_{s}")) > 1e-6 for s in ("top", "bottom", "left", "right")):
                raise ValueError("export_asset_geometry_mismatch")
    if native:
        visible_slots = {o.slot_id for o in scene.objects}
        objects = {o.object_id: o for o in scene.objects}
        visible_connections = {f"{objects[e.source_id].slot_id}->{objects[e.target_id].slot_id}" for e in scene.edges}
        visible_slots.update(b.label_slot for k, b in scene.execution.connections.items() if k in visible_connections and b.label_slot)
        for slot, binding in scene.execution.targets.items():
            if slot not in visible_slots:
                target = _native_target(slide, binding)
                if binding.kind == "image":
                    if target.image.blob != _transparent_picture():
                        raise ValueError("export_hidden_native_content_visible")
                elif _source_text(target).strip():
                    raise ValueError("export_hidden_native_content_visible")
    relation_names = {f"relation:{e.relation_id}" for e in scene.edges}
    if {n for n in named if n.startswith("relation:")} != relation_names:
        raise ValueError("export_relation_set_mismatch")
    for edge in scene.edges:
        connector = named[f"relation:{edge.relation_id}"]
        starts, ends = connector._element.xpath(".//a:stCxn"), connector._element.xpath(".//a:endCxn")
        if len(starts) != 1 or len(ends) != 1 or starts[0].get("id") != str(named[f"teaching:{edge.source_id}"].shape_id) or ends[0].get("id") != str(named[f"teaching:{edge.target_id}"].shape_id):
            raise ValueError("export_relation_endpoint_mismatch")
        if starts[0].get("idx") != str(edge.start_site) or ends[0].get("idx") != str(edge.end_site):
            raise ValueError("export_relation_site_mismatch")
        if any(abs(a - b) > 1 for a, b in zip((connector.begin_x / 12700, connector.begin_y / 12700, connector.end_x / 12700, connector.end_y / 12700), (edge.x1, edge.y1, edge.x2, edge.y2))):
            raise ValueError("export_relation_geometry_mismatch")
        arrows = connector._element.xpath(".//a:tailEnd")
        if relation_is_directed(edge.kind) and (not arrows or arrows[0].get("type") != "triangle"):
            raise ValueError("export_relation_direction_missing")
        if not relation_is_directed(edge.kind) and arrows and arrows[0].get("type") not in {None, "none"}:
            raise ValueError("export_relation_direction_invented")
        if edge.label:
            label = named.get(f"relation-label:{edge.relation_id}")
            obj = edge.label_object
            if label is None or _source_text(label) != edge.label:
                raise ValueError("export_relation_label_mismatch")
            if any(abs(a - b) > 1 for a, b in zip((label.left / 12700, label.top / 12700, label.width / 12700, label.height / 12700), (obj.x, obj.y, obj.width, obj.height))):
                raise ValueError("export_relation_label_geometry_mismatch")
            if any(p.font.name != scene.execution.font_family or p.font.size != Pt(obj.font_size) for p in label.text_frame.paragraphs):
                raise ValueError("export_relation_label_font_mismatch")
    return {"passed": True, "elements": len(scene.objects), "relations": len(scene.edges), "scene_digest": scene.scene_digest,
            "editability": sorted({o.editability for o in scene.objects})}


def clone_native_slide(presentation, source, execution):
    slide = presentation.slides.add_slide(source.slide_layout)
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)
    from slide_deck_v6_personal_renderer import _copy_relationships, _replace_relationship_ids
    remap = _copy_relationships(source, slide)
    source_bg = source._element.find(qn("p:cSld")).find(qn("p:bg"))
    if source_bg is not None:
        bg = deepcopy(source_bg)
        _replace_relationship_ids(bg, remap)
        slide._element.find(qn("p:cSld")).insert(0, bg)
    for shape in source.shapes:
        node = deepcopy(shape._element)
        _replace_relationship_ids(node, remap)
        slide.shapes._spTree.insert_element_before(node, "p:extLst")
    inspect_native_bindings(slide, execution.targets)
    inspect_native_connections(slide, execution)
    for target in execution.targets.values():
        shape = _native_target(slide, target)
        if target.kind == "image":
            _replace_picture(shape, _transparent_picture())
        else:
            shape.text_frame.clear()
        if hasattr(shape, "name") and shape.name.startswith("teaching:"):
            shape.name = f"template-hidden:{shape.shape_id}"
    return slide


def render_teaching_deck(deck, output_path: Path, *, assets=None, source_path=None) -> Path:
    native = any(p.resolved_scene.execution.mode == "native_fill" for p in deck.pages)
    if native and source_path is None:
        raise ValueError("native_template_source_missing")
    presentation = Presentation(str(source_path)) if native else Presentation()
    originals = list(presentation.slides)
    if native and (presentation.slide_width != Pt(960) or presentation.slide_height != Pt(540)):
        raise ValueError("native_source_page_geometry_unsupported")
    presentation.slide_width, presentation.slide_height = Pt(960), Pt(540)
    for page in deck.pages:
        scene = page.resolved_scene
        if scene is None:
            raise ValueError("mixed_scene_contract_forbidden")
        if scene.execution.mode == "native_fill":
            if file_digest(Path(source_path)) != scene.execution.source_sha256:
                raise ValueError("native_template_source_changed")
            source = originals[scene.execution.source_slide_number - 1]
            slide = clone_native_slide(presentation, source, scene.execution)
        else:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = _color(scene.background)
        render_scene(slide, scene, assets=assets)
        from slide_speaker_notes import _speaker_notes
        slide.notes_slide.notes_text_frame.text = _speaker_notes(page)
    if native:
        from slide_deck_v6_personal_renderer import _remove_original_slides
        _remove_original_slides(presentation, len(originals))
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.NamedTemporaryFile(prefix=f".{output_path.stem}-", suffix=".pptx", dir=output_path.parent, delete=False) as pending:
        temporary = Path(pending.name)
    try:
        presentation.save(temporary)
        readback = Presentation(str(temporary))
        for slide, page in zip(readback.slides, deck.pages, strict=True):
            audit_scene(slide, page.resolved_scene)
            from slide_speaker_notes import _speaker_notes
            if slide.notes_slide.notes_text_frame.text != _speaker_notes(page):
                raise ValueError("export_notes_mismatch")
        from ppt_render_audit import render_evidence
        with tempfile.TemporaryDirectory(prefix="ppt-export-audit-") as audit_dir:
            render_evidence(temporary, [p.resolved_scene for p in deck.pages], output=audit_dir)
        temporary.replace(output_path)
    finally:
        temporary.unlink(missing_ok=True)
    return output_path
