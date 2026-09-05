from copy import deepcopy

import pytest
from pptx import Presentation
from pptx.util import Pt

from ppt_layout_execution import compile_teaching_template
from ppt_layout_samples import chart_sample, layout_sample
from ppt_native_scene import audit_scene, render_scene
from ppt_page_scene import resolve_page_scenes, verify_scene
from ppt_teaching_content import PageTeachingV2, chart_number


def scenes(content=None, theme="academic-editorial", slug="data-bars"):
    template = compile_teaching_template(theme, certification_required=False)
    return resolve_page_scenes(page_id="chart", title="比较记录时间", content=content or chart_sample(),
        layout=template.get_layout(template.layout_id(slug)), template=template, source_document_revision="sample-v1")


def test_shared_zero_scale_and_state_positions_roundtrip(tmp_path):
    before, after = scenes()
    first = next(o for o in before.objects if o.object_id == "chart-bar:av")
    second = next(o for o in after.objects if o.object_id == "chart-bar:bv")
    assert first.width * 2 == second.width
    assert first.x == second.x == 262
    assert first == next(o for o in after.objects if o.object_id == first.object_id)
    assert not any(o.element_id == "bv" for o in before.objects)
    p = Presentation()
    p.slide_width, p.slide_height = Pt(960), Pt(540)
    for scene in (before, after):
        verify_scene(scene)
        render_scene(p.slides.add_slide(p.slide_layouts[6]), scene)
    path = tmp_path / "chart.pptx"
    p.save(path)
    deck = Presentation(path)
    assert all(audit_scene(slide, scene)["passed"] for slide, scene in zip(deck.slides, (before, after)))
    bar = next(s for s in deck.slides[1].shapes if s.name == "teaching:chart-bar:bv")
    bar.width = Pt(250)
    with pytest.raises(ValueError, match="export_chart_scale_mismatch"):
        audit_scene(deck.slides[1], after)


@pytest.mark.parametrize("text", ["-1", "NaN", "Infinity", "12%", "1e3", "1,000", "1000000000001"])
def test_unsupported_chart_data_rejected(text):
    with pytest.raises(ValueError, match="chart_value_not_supported"):
        chart_number(text)


def test_chart_cannot_hide_unit_or_invent_units():
    value = chart_sample().model_dump()
    value["states"][0]["visible_element_ids"].remove("unit")
    with pytest.raises(ValueError, match="chart_context_hidden"):
        PageTeachingV2.model_validate(value)
    value = chart_sample().model_dump()
    value["elements"][0]["kind"] = "text"
    with pytest.raises(ValueError, match="chart_unit_must_be_source_exact"):
        PageTeachingV2.model_validate(value)


def test_second_theme_changes_style_only():
    content = layout_sample("concept-map")
    first = scenes(content, slug="concept-map")[0]
    second = scenes(content, "qizhi-classroom", "concept-map")[0]
    semantic = lambda scene: [(o.object_id, o.text, o.x, o.y, o.width, o.height) for o in scene.objects]
    assert semantic(first) == semantic(second)
    assert [e.model_dump(exclude={"label_object"}) for e in first.edges] == [e.model_dump(exclude={"label_object"}) for e in second.edges]
    assert first.accent_color != second.accent_color
    assert first.objects[1].fill != second.objects[1].fill
