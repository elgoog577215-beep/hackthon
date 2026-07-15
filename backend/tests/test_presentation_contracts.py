from pathlib import Path

import pytest
from pptx import Presentation
from pydantic import ValidationError

from presentation_models import (
    DeckRevision,
    PresentationDeck,
    PresentationEvent,
    PresentationScope,
    PreviewMessage,
    Slide,
    SlideBlock,
    SlideSourceRefs,
    SourceRef,
)
from presentation_quality import evaluate_presentation_quality, revision_checksum
from presentation_render import render_artifacts, speaker_notes_digest, title_digest
from presentation_templates import LAYOUTS, TEMPLATES, get_layout


def _slide(index: int, layout_id: str) -> Slide:
    block_type = {
        "L01": "text", "L04": "bullets", "L07": "code",
        "L08": "comparison", "L09": "exercise", "L10": "bullets",
    }[layout_id]
    return Slide(
        slide_id=f"slide-{index}", position=index, layout_id=layout_id, status="ready",
        title=f"第 {index + 1} 页：指针与内存", subtitle="用一个可追溯示例解释核心关系",
        key_message="地址是变量在内存中的位置",
        blocks=[SlideBlock(
            block_id=f"block-{index}", type=block_type, title="关键内容",
            content="int x = 10;\nint *p = &x;" if block_type == "code" else "地址与变量的关系",
            items=["取地址运算符 &", "解引用运算符 *"] if block_type in {"bullets", "exercise"} else [],
            metadata={"left": "错误：空指针可以直接使用", "right": "正确：先判空再解引用"} if block_type == "comparison" else {},
        )],
        speaker_notes=f"第 {index + 1} 页讲稿：联系课程来源说明。",
        source_refs=SlideSourceRefs(section_ids=["node-2"], block_ids=["cb-2"], block_revision_ids=["cbr-2"], objective_ids=["obj-2"]),
    )


def test_registry_is_frozen_to_three_styles_and_ten_layouts():
    assert set(TEMPLATES) == {"lingzhi-classroom", "lingzhi-engineering", "lingzhi-academic"}
    assert set(LAYOUTS) == {f"L{i:02d}" for i in range(1, 11)}
    with pytest.raises(ValueError, match="unknown_layout"):
        get_layout("L99")
    with pytest.raises(ValidationError):
        Slide(slide_id="bad", position=0, layout_id="L99")


def test_event_id_and_schema_are_ordered():
    event = PresentationEvent(event_type="deck_outline", deck_id="deck-1", generation_id="gen-1", event_seq=1, outline_revision=1)
    assert event.sse_id == "gen-1:1"
    assert event.schema_version == "presentation-event/v1"


def test_preview_message_is_bound_to_revision_checksum():
    message = PreviewMessage(
        type="render:measured",
        deck_id="deck-1",
        revision_id="rev-1",
        revision_checksum="sha256:preview-revision",
        payload={"slide_count": 1, "overflow": False, "collision": False},
    )
    assert message.revision_checksum == "sha256:preview-revision"
    with pytest.raises(ValidationError):
        PreviewMessage(type="slide:selected", deck_id="deck-1", revision_id="rev-1")


def test_render_issues_identify_the_exact_slide_and_fix():
    slide = _slide(0, "L04")
    revision = DeckRevision(
        revision_id="rev-render-issue",
        deck_id="deck-render-issue",
        reason="initial_generation",
        source_snapshot_id="source-render-issue",
        slide_order=[slide.slide_id],
        slides=[slide],
    )
    deck = PresentationDeck(
        deck_id="deck-render-issue",
        course_id="course-render-issue",
        title="排版检查课件",
        source_ref=SourceRef(
            course_id="course-render-issue",
            source_format="canonical",
            version_id="cv1",
            document_revision="cdr1",
            source_snapshot_id="source-render-issue",
            source_snapshot_sha256=f"sha256:{'a' * 64}",
        ),
        scope=PresentationScope(type="chapter", section_ids=["node-2"]),
    )
    report = evaluate_presentation_quality(
        deck,
        revision,
        {
            "publication_allowed": True,
            "sections": [{"section_id": "node-2"}],
            "blocks": [{"block_id": "cb-2", "block_revision_id": "cbr-2"}],
            "objectives": [{"objective_id": "obj-2"}],
        },
        render_measurement={
            "revision_checksum": revision_checksum(revision),
            "slide_count": 1,
            "overflow": True,
            "collision": True,
            "overflow_slide_ids": [slide.slide_id],
            "collision_slide_ids": [slide.slide_id],
        },
        require_publication=True,
        require_render_measurement=True,
    )

    overflow = next(issue for issue in report.issues if issue.code == "render_overflow")
    collision = next(issue for issue in report.issues if issue.code == "render_collision")
    assert (overflow.target_type, overflow.target_id) == ("slide", slide.slide_id)
    assert "第 1 页" in overflow.message and "精简本页" in overflow.fix_action
    assert (collision.target_type, collision.target_id) == ("slide", slide.slide_id)
    assert "第 1 页" in collision.message and "调整版式" in collision.fix_action


@pytest.mark.parametrize("template_id", list(TEMPLATES))
def test_representative_deck_renders_matching_html_and_pptx(tmp_path: Path, template_id: str):
    slides = [_slide(index, layout) for index, layout in enumerate(["L01", "L04", "L07", "L08", "L09", "L10"])]
    revision = DeckRevision(
        revision_id="rev-1", deck_id="deck-1", reason="initial_generation",
        source_snapshot_id="source-1", slide_order=[slide.slide_id for slide in slides], slides=slides,
    )
    result = render_artifacts(revision, template_id, tmp_path / template_id)
    assert result["page_count"] == 6
    assert result["title_digest"] == title_digest(slides)
    assert result["speaker_notes_digest"] == speaker_notes_digest(slides)
    assert Path(str(result["html_path"])).read_text(encoding="utf-8").count('class="slide ') == 6
    pptx = Presentation(str(result["pptx_path"]))
    assert len(pptx.slides) == 6
    assert all(slide.notes_slide.notes_text_frame.text for slide in pptx.slides)
    for expected, rendered in zip(slides, pptx.slides, strict=True):
        texts = [shape.text for shape in rendered.shapes if hasattr(shape, "text_frame")]
        assert expected.title in texts
