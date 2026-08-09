from __future__ import annotations

from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from slide_deck_renderer import audit_exported_pptx, audit_rendered_slide_images
from slide_deck_v5 import _page_density_metrics
from slide_quality_v5 import (
    build_slide_deck_quality_v5,
    repair_render_slides_v5,
    repair_semantic_slides_v5,
)


def _slide(
    page_id: str,
    *,
    title: str,
    content: str,
    episode_id: str = "episode-1",
    scene_kind: str = "concept",
    block_type: str = "statement",
) -> dict:
    return {
        "unit_id": page_id,
        "position": 0,
        "layout": "concept",
        "slide_purpose": scene_kind,
        "title": title,
        "key_message": "",
        "episode_id": episode_id,
        "scene_kind": scene_kind,
        "blocks": [{
            "block_id": f"{page_id}:body",
            "type": block_type,
            "title": "",
            "content": content,
            "items": [],
            "metadata": {"semantic_atom_id": f"atom:{page_id}"},
        }],
        "visuals": [],
        "quality": {
            "requested_layout": "editorial-body",
            "resolved_layout": "editorial-body",
        },
    }


def test_unique_quality_report_rejects_stale_scores_and_uses_final_slide_count() -> None:
    slides = [
        _slide(
            "slide:v5:episode-1:001",
            title="概念成立需要同时满足两个条件",
            content="条件一给出对象边界，条件二说明对象之间的关系。",
        )
    ]

    report = build_slide_deck_quality_v5(
        slides,
        planner="deterministic_v5",
        fallback_reason="ai_unavailable",
        legacy_quality={
            "score": 12,
            "pedagogical_score": 44,
            "presentation_score": 31,
        },
    )

    assert report["schema_version"] == "slide_deck_quality_v5"
    assert report["status"] == "passed"
    assert report["metrics"]["total_slide_count"] == len(slides)
    assert set(report["dimensions"]) == {
        "source_integrity",
        "teaching_closure",
        "pagination_narrative",
        "layout_export",
        "visual_effectiveness",
        "attribution_accessibility",
    }
    assert "pedagogical_score" not in report
    assert "presentation_score" not in report
    assert report["score"] >= 80


def test_quality_contract_blocks_split_atoms_duplicate_copy_and_empty_answers() -> None:
    first = _slide(
        "slide:v5:episode-1:001",
        title="用共同维度比较两个对象",
        content="比较时先固定同一个观察维度，再说明两个对象的差异结论。",
        scene_kind="comparison",
    )
    first["blocks"][0]["metadata"]["semantic_atom_id"] = "atom-shared"
    second = _slide(
        "slide:v5:episode-1:002",
        title="用共同维度比较两个对象",
        content="比较时先固定同一个观察维度，再说明两个对象的差异结论。",
        scene_kind="practice_feedback",
        block_type="exercise",
    )
    second["blocks"][0]["metadata"].update({
        "semantic_atom_id": "atom-shared",
        "question_mode": "closed",
        "question_id": "question-1",
    })
    second["blocks"].append({
        "block_id": "answer-1",
        "type": "callout",
        "title": "答案",
        "content": "",
        "items": [],
        "metadata": {"answer_for": "question-1"},
    })

    report = build_slide_deck_quality_v5([first, second])
    codes = {issue["code"] for issue in report["blockers"]}

    assert report["status"] == "blocked"
    assert "semantic_atom_split" in codes
    assert "duplicate_visible_content" in codes
    assert "duplicate_title" in codes
    assert "empty_answer" in codes


def test_paired_answer_ids_satisfy_closed_question_mapping() -> None:
    slide = _slide(
        "slide:v5:practice:paired",
        title="判断映射是否保持线性结构",
        content="判断给定映射是否同时保持加法与数乘。",
        scene_kind="practice_feedback",
        block_type="exercise",
    )
    slide["blocks"][0]["metadata"].update({
        "question_id": "question-linear",
        "question_ids": ["question-linear"],
        "question_mode": "closed",
        "semantic_role": "prompt",
    })
    slide["blocks"].append({
        "block_id": "answer-linear",
        "type": "callout",
        "title": "参考答案与判断依据",
        "content": "",
        "items": ["两项都成立，因此该映射是线性映射。"],
        "metadata": {
            "semantic_role": "answer",
            "direct_answer": True,
            "answer_for_question_ids": ["question-linear"],
            "source_fragment_ids": ["fragment-answer"],
        },
    })

    report = build_slide_deck_quality_v5([slide])

    assert "answer_mapping_missing" not in {
        issue["code"] for issue in report["blockers"]
    }


def test_semantic_repair_does_not_merge_past_final_body_budget() -> None:
    first = _slide(
        "slide:v5:dense:001",
        title="第一组完整依据",
        content="甲" * 150,
    )
    second = _slide(
        "slide:v5:dense:002",
        title="第二组待补充依据",
        content=("乙" * 100) + "包括",
    )
    for slide, count in ((first, 150), (second, 102)):
        slide["quality"].update({
            "body_character_count": count,
            "body_character_budget": 230,
            "visible_item_count": 0,
            "visible_item_budget": 5,
        })

    repaired, history = repair_semantic_slides_v5(
        [first, second],
        max_rounds=2,
    )

    assert len(repaired) == 2
    assert not any(
        item["action"] == "merge_sparse_or_dangling_page"
        for item in history
    )


def test_continuation_requires_parent_link_and_visible_sequence() -> None:
    continuation = _slide(
        "slide:v5:episode-1:continuation",
        title="完整主题标题",
        content="这是超长语义原子的后续完整分段。",
    )
    continuation["quality"].update({
        "continuation_of": "slide:v5:episode-1:root",
        "continuation_index": 2,
        "continuation_total": 3,
    })

    report = build_slide_deck_quality_v5([continuation])

    assert "continuation_sequence_missing" in {
        issue["code"] for issue in report["blockers"]
    }


def test_semantic_repair_restores_continuation_sequence_after_title_repair() -> None:
    root = _slide(
        "slide:v5:episode-1:root",
        title="纵隔四分法的平面划分依据",
        content="前三个区域分别说明其位置与主要结构。" * 20,
    )
    continuation = _slide(
        "slide:v5:episode-1:continuation",
        title="纵隔四分法的平面划分依据",
        content="后纵隔位于心包后壁与脊柱之间。" * 20,
    )
    continuation["quality"].update({
        "continuation_of": root["unit_id"],
        "continuation_index": 2,
        "continuation_total": 2,
        "title_character_budget": 18,
    })

    repaired, _history = repair_semantic_slides_v5(
        [root, continuation],
        max_rounds=2,
    )

    assert repaired[1]["title"].endswith("（续2/2）")
    assert "continuation_sequence_missing" not in {
        issue["code"] for issue in build_slide_deck_quality_v5(repaired)["issues"]
    }


def test_semantic_repair_renumbers_duplicate_continuation_children() -> None:
    root = _slide(
        "slide:v5:episode-1:root",
        title="纵隔四分法的平面划分依据",
        content="纵隔划分的来源说明。" * 30,
    )
    children = []
    for number, content in enumerate((
        "前纵隔与中纵隔分别说明其边界。",
        "后纵隔位于心包后壁与脊柱之间。",
    ), start=1):
        child = _slide(
            f"slide:v5:episode-1:continuation-{number}",
            title="纵隔四分法的平面划分依据（续2/2）",
            content=content * 20,
        )
        child["quality"].update({
            "continuation_of": root["unit_id"],
            "continuation_index": 2,
            "continuation_total": 2,
            "title_character_budget": 18,
        })
        children.append(child)

    repaired, _history = repair_semantic_slides_v5(
        [root, *children],
        max_rounds=2,
    )

    assert [slide["quality"]["continuation_index"] for slide in repaired] == [
        1,
        2,
        3,
    ]
    assert repaired[1]["title"].endswith("（续2/3）")
    assert repaired[2]["title"].endswith("（续3/3）")


def test_grounded_practice_feedback_is_not_accidental_duplicate_copy() -> None:
    concept = _slide(
        "slide:v5:episode-1:concept",
        title="四层结构形成稳定定位顺序",
        content="皮层、浅筋膜、深筋膜与肌层构成由浅入深的定位顺序。",
    )
    practice = _slide(
        "slide:v5:episode-2:practice",
        title="核对四层结构的绘制顺序",
        content="请核对绘图中四层结构的顺序。",
        episode_id="episode-2",
        scene_kind="practice_feedback",
        block_type="exercise",
    )
    practice["blocks"][0]["metadata"].update({
        "question_mode": "open_discussion",
        "semantic_role": "prompt",
    })
    practice["blocks"].append({
        "block_id": "slide:v5:episode-2:practice:feedback",
        "type": "callout",
        "title": "判断依据",
        "content": "皮层、浅筋膜、深筋膜与肌层构成由浅入深的定位顺序。",
        "items": [],
        "metadata": {
            "semantic_role": "feedback",
            "grounded": True,
        },
    })

    report = build_slide_deck_quality_v5([concept, practice])

    assert "duplicate_visible_content" not in {
        issue["code"] for issue in report["issues"]
    }


def test_deterministic_semantic_repair_is_targeted_and_never_invents_an_answer() -> None:
    healthy = _slide(
        "slide:v5:episode-1:001",
        title="证据支持这一判断",
        content="记录中的两个观察结果共同支持本页结论。",
        scene_kind="evidence",
    )
    question = _slide(
        "slide:v5:episode-2:001",
        title="哪一个选项符合材料？",
        content="请选择符合材料描述的选项。",
        episode_id="episode-2",
        scene_kind="practice_feedback",
        block_type="exercise",
    )
    question["blocks"][0]["metadata"].update({
        "question_mode": "closed",
        "question_id": "question-2",
        "source_answer": "",
    })
    question["blocks"].append({
        "block_id": "answer-2",
        "type": "callout",
        "title": "答案",
        "content": "",
        "items": [],
        "metadata": {"answer_for": "question-2"},
    })
    original_healthy = deepcopy(healthy)

    repaired, history = repair_semantic_slides_v5([healthy, question], max_rounds=2)

    assert repaired[0] == original_healthy
    assert repaired[1]["quality"]["question_mode"] == "open_discussion"
    assert all(block.get("title") != "答案" for block in repaired[1]["blocks"])
    assert "开放讨论" in repaired[1]["key_message"]
    assert history
    assert all(item["page_id"] == "slide:v5:episode-2:001" for item in history)


def test_non_exempt_sparse_page_and_internal_labels_are_blockers() -> None:
    slide = _slide(
        "slide:v5:episode-1:001",
        title="知识规范名称",
        content="定义",
    )

    report = build_slide_deck_quality_v5([slide])
    codes = {issue["code"] for issue in report["blockers"]}

    assert "raw_internal_label_visible" in codes
    assert "sparse_non_exempt_page" in codes


def test_repetitive_text_only_editorial_pages_block_publication() -> None:
    slides = []
    for index in range(6):
        slide = _slide(
            f"slide:v5:episode-{index}:001",
            title=f"Concept {index} has a distinct instructional claim",
            content=(
                f"Concept {index} is explained as one uninterrupted prose paragraph "
                "without a second information region, worked example, comparison, "
                "process, question, or visual anchor for classroom presentation."
            ),
            episode_id=f"episode-{index}",
        )
        slides.append(slide)

    report = build_slide_deck_quality_v5(slides)
    codes = {issue["code"] for issue in report["blockers"]}

    assert "text_only_editorial_ratio_exceeded" in codes
    assert "repetitive_text_only_editorial_run" in codes


def test_deck_allows_at_most_three_intentional_hero_claim_pages() -> None:
    slides = []
    for index in range(4):
        slide = _slide(
            f"slide:v5:hero-{index}:001",
            title=f"Claim {index} is the one idea to remember",
            content="",
            episode_id=f"hero-{index}",
        )
        slide["quality"].update({
            "requested_layout": "hero-claim",
            "resolved_layout": "hero-claim",
            "suppress_redundant_body": True,
        })
        slides.append(slide)

    report = build_slide_deck_quality_v5(slides)

    assert "hero_claim_page_limit_exceeded" in {
        issue["code"] for issue in report["blockers"]
    }


def test_chapter_boundaries_reset_text_only_layout_runs() -> None:
    slides = [
        _slide(
            "slide:v5:chapter-a:001",
            title="First claim",
            content="A complete paragraph explains the first classroom claim with enough detail.",
        ),
        _slide(
            "slide:v5:chapter-a:002",
            title="Second claim",
            content="A different paragraph explains the second classroom claim with enough detail.",
        ),
        {
            **_slide(
                "slide:v5:chapter-a:recap",
                title="Chapter recap",
                content="The first and second claims are connected.",
                scene_kind="chapter_recap",
            ),
            "layout": "recap",
            "quality": {
                "requested_layout": "chapter-recap",
                "resolved_layout": "chapter-recap",
            },
        },
        {
            **_slide(
                "slide:v5:chapter-b:entry",
                title="Chapter B",
                content="",
                scene_kind="chapter_entry",
            ),
            "layout": "chapter",
            "quality": {
                "requested_layout": "chapter-entry",
                "resolved_layout": "chapter-entry",
            },
        },
        _slide(
            "slide:v5:chapter-b:001",
            title="Third claim",
            content="A third paragraph starts the next chapter with a complete independent claim.",
        ),
        _slide(
            "slide:v5:chapter-b:002",
            title="Fourth claim",
            content="A fourth paragraph closes the local sequence with a complete independent claim.",
        ),
    ]

    report = build_slide_deck_quality_v5(slides)

    assert "repetitive_text_only_editorial_run" not in {
        issue["code"] for issue in report["blockers"]
    }


def test_chapter_entry_without_a_mainline_is_a_publication_blocker() -> None:
    entry = _slide(
        "slide:v5:chapter-entry",
        title="交互逻辑与物理系统基础",
        content="",
        scene_kind="chapter_entry",
    )
    entry["blocks"] = []
    entry["key_message"] = ""
    entry["takeaway"] = ""
    entry["quality"].update({
        "requested_layout": "chapter-entry",
        "resolved_layout": "chapter-entry",
    })

    report = build_slide_deck_quality_v5([entry])

    assert "chapter_entry_mainline_missing" in {
        issue["code"] for issue in report["blockers"]
    }


def test_chapter_recap_may_repeat_the_preceding_claim_as_a_summary() -> None:
    concept = _slide(
        "slide:v5:chapter:concept",
        title="Initialization order controls safe access",
        content="Awake completes before Start, so references can be prepared before use.",
    )
    recap = {
        **_slide(
            "slide:v5:chapter:recap",
            title="Chapter recap",
            content="Awake completes before Start, so references can be prepared before use.",
            scene_kind="chapter_recap",
        ),
        "layout": "recap",
        "quality": {
            "requested_layout": "chapter-recap",
            "resolved_layout": "chapter-recap",
        },
    }

    report = build_slide_deck_quality_v5([concept, recap])

    assert "duplicate_visible_content" not in {
        issue["code"] for issue in report["blockers"]
    }


def test_internal_label_cleanup_handles_continuation_suffixes() -> None:
    slide = _slide(
        "slide:v5:internal-continuation",
        title="A complete source-backed title",
        content="The continuation keeps one complete source-backed teaching claim.",
    )
    slide["key_message"] = "知识规范名称（续2/2）"

    repaired, _history = repair_semantic_slides_v5([slide], max_rounds=2)

    assert repaired[0]["key_message"] == ""
    assert "raw_internal_label_visible" not in {
        issue["code"] for issue in build_slide_deck_quality_v5(repaired)["issues"]
    }


def test_complete_continuation_title_fits_the_presentation_heading_contract() -> None:
    slide = _slide(
        "slide:v5:episode-1:continuation",
        title="生命周期回调的触发时序逻辑（续2/2）",
        content="生命周期回调遵循初始化、帧更新与销毁阶段的明确触发顺序。",
    )

    report = build_slide_deck_quality_v5([slide])
    density = _page_density_metrics(slide)

    assert density["title_character_count"] <= density["title_character_budget"]
    assert "slide_title_overflow" not in {
        issue["code"] for issue in report["issues"]
    }


def test_semantic_repair_removes_internal_labels_from_source_body_without_losing_text() -> None:
    slide = _slide(
        "slide:v5:episode-1:001",
        title="本页解释完整对象关系",
        content="知识规范名称：对象之间的完整关系",
    )

    repaired, history = repair_semantic_slides_v5([slide], max_rounds=2)

    assert repaired[0]["blocks"][0]["content"] == "对象之间的完整关系"
    assert any(item["action"] == "replace_internal_label" for item in history)
    assert "raw_internal_label_visible" not in {
        issue["code"] for issue in build_slide_deck_quality_v5(repaired)["issues"]
    }


def test_semantic_repair_replaces_promoted_internal_titles_from_markdown_sources() -> None:
    slides = [
        _slide(
            "slide:v5:898723642bdc07b279355281",
            title="本节知识规范名称",
            content=(
                "Unity 开发环境初始化与工程目录结构规范。"
                "学习者需完成以下可观察目标。"
            ),
        ),
        _slide(
            "slide:v5:29900c205643b89ca2458cc8",
            title="知识规范名称",
            content=(
                "MonoBehaviour 脚本命名规范与生命周期回调执行顺序。"
                "本节通过创建符合规范的脚本验证初始化时序。"
            ),
        ),
        _slide(
            "slide:v5:cac08a9027846cdfba776eef",
            title="**知识规范名称",
            content=(
                "在脚本中实现 Awake、Start 和 Update 三个生命周期回调方法，"
                "并观察日志执行顺序。"
            ),
        ),
    ]
    slides[0]["takeaway"] = "**Unity 开发环境初始化与工程目录结构规范**。"
    for slide in slides[1:]:
        slide["takeaway"] = (
            "**知识规范名称：MonoBehaviour 脚本命名规范与生命周期回调执行顺序**"
        )
    slides[2]["quality"].update({
        "continuation_of": slides[1]["unit_id"],
        "continuation_index": 2,
        "continuation_total": 2,
    })

    repaired, history = repair_semantic_slides_v5(slides, max_rounds=2)
    report = build_slide_deck_quality_v5(repaired)

    assert all(str(slide["title"]).strip("* ") for slide in repaired)
    assert all("知识规范名称" not in str(slide["title"]) for slide in repaired)
    assert all("知识规范名称" not in str(slide.get("takeaway") or "") for slide in repaired)
    assert any(item["action"] == "replace_internal_label" for item in history)
    assert "raw_internal_label_visible" not in {
        issue["code"] for issue in report["issues"]
    }


def test_semantic_repair_drops_a_source_bound_dangling_scaffold() -> None:
    slide = _slide(
        "slide:v5:source-scaffold",
        title="切换平台会重建目标资源",
        content=(
            "切换目标平台会触发编译器、资源压缩算法与脚本后端的重构。"
        ),
        scene_kind="reasoning",
    )
    slide["blocks"][0]["metadata"]["fragment_ids"] = ["fragment-claim"]
    slide["blocks"].append({
        "block_id": "dangling-label",
        "type": "process",
        "title": "",
        "content": "",
        "items": ["平台切换操作："],
        "metadata": {"fragment_ids": ["fragment-label"]},
    })

    repaired, history = repair_semantic_slides_v5([slide], max_rounds=2)
    report = build_slide_deck_quality_v5(repaired)

    assert repaired[0]["blocks"][0]["content"].startswith("切换目标平台")
    assert all(
        "平台切换操作" not in str(item)
        for block in repaired[0]["blocks"]
        for item in block.get("items") or []
    )
    assert any(item["action"] == "remove_dangling_scaffold" for item in history)
    assert "dangling_fragment" not in {
        issue["code"] for issue in report["issues"]
    }


def test_quality_does_not_treat_code_semicolon_as_dangling_prose() -> None:
    slide = _slide(
        "slide:v5:code-semicolon",
        title="Lifecycle callback",
        content="void Start() { Debug.Log(\"ready\"); }",
        scene_kind="method",
        block_type="code",
    )
    slide["quality"].update({
        "requested_layout": "code",
        "resolved_layout": "code",
        "subject_artifact_kinds": ["code"],
    })

    report = build_slide_deck_quality_v5([slide])

    assert "dangling_fragment" not in {
        issue["code"] for issue in report["issues"]
    }


def test_quality_allows_structured_practice_rows_to_end_with_semicolons() -> None:
    slide = _slide(
        "slide:v5:practice-semicolon",
        title="Environment check",
        content="",
        scene_kind="practice_feedback",
        block_type="exercise",
    )
    slide["blocks"][0].update({
        "items": ["Confirm the editor version;", "Run the callback;"],
        "metadata": {"semantic_role": "prompt"},
    })
    slide["blocks"].append({
        "block_id": "practice-feedback",
        "type": "callout",
        "title": "Feedback",
        "content": "",
        "items": ["The callback ran in the expected order;"],
        "metadata": {"semantic_role": "feedback"},
    })
    slide["quality"].update({
        "requested_layout": "practice-feedback",
        "resolved_layout": "practice-feedback",
    })

    report = build_slide_deck_quality_v5([slide])

    assert "dangling_fragment" not in {
        issue["code"] for issue in report["issues"]
    }


def test_semantic_repair_records_source_bound_process_and_example_closure() -> None:
    process = _slide(
        "slide:v5:process-result",
        title="对象序列化为 JSON 字符串",
        content=(
            "使用 JsonUtility.ToJson 方法遍历对象的公共字段，"
            "并将其转换为 JSON 字符串。"
        ),
        scene_kind="method",
    )
    process["blocks"][0]["metadata"]["fragment_ids"] = ["fragment-process"]
    example = _slide(
        "slide:v5:worked-conclusion",
        title="初始化顺序会影响实例访问",
        content=(
            "案例中 PlayerController 先于 GameManager 初始化，"
            "因此访问实例时会报错。"
        ),
        scene_kind="worked_example",
        block_type="exercise",
    )
    example["blocks"][0]["metadata"].update({
        "fragment_ids": ["fragment-example"],
        "question_mode": "open_discussion",
    })

    repaired, history = repair_semantic_slides_v5(
        [process, example],
        max_rounds=2,
    )
    report = build_slide_deck_quality_v5(repaired)

    assert repaired[0]["quality"]["process_result"]
    assert repaired[1]["quality"]["worked_example_conclusion"]
    assert any(item["action"] == "bind_source_closure" for item in history)
    assert not {
        "process_result_missing",
        "worked_example_conclusion_missing",
    }.intersection(issue["code"] for issue in report["issues"])


def test_context_only_method_and_open_case_prompt_do_not_require_fake_results() -> None:
    method_context = _slide(
        "slide:v5:method-context",
        title="线程安全单例需要考虑执行环境",
        content=(
            "Unity 主线程逻辑可能与网络包处理或后台计算交织，"
            "标准懒汉式单例需考虑线程安全性。"
        ),
        scene_kind="method",
    )
    case_prompt = _slide(
        "slide:v5:case-prompt",
        title="粒子增多伴随帧率下降",
        content=(
            "随着粒子特效增多，帧率从 60 FPS 下降至 15 FPS，"
            "并伴随周期性卡顿。"
        ),
        scene_kind="worked_example",
        block_type="exercise",
    )
    case_prompt["blocks"][0]["metadata"]["question_mode"] = "open_discussion"

    report = build_slide_deck_quality_v5([method_context, case_prompt])

    assert not {
        "process_result_missing",
        "worked_example_conclusion_missing",
    }.intersection(issue["code"] for issue in report["issues"])


def test_duplicate_gate_allows_a_shared_technical_term_on_distinct_pages() -> None:
    overview = _slide(
        "slide:v5:profiler-overview",
        title="发布清单连接诊断与回归",
        content=(
            "发布清单要求使用 Unity Profiler 的 Deep Profile 模式定位瓶颈，"
            "并记录复现步骤、根因与修复方案。"
        ),
    )
    diagnosis = _slide(
        "slide:v5:profiler-diagnosis",
        title="动态状态需要运行时诊断",
        content=(
            "静态分析无法覆盖动态状态，必须使用 Unity Profiler 的 Deep Profile "
            "模式观察真实耗时。"
        ),
    )

    report = build_slide_deck_quality_v5([overview, diagnosis])

    assert "duplicate_visible_content" not in {
        issue["code"] for issue in report["issues"]
    }


def test_semantic_repair_completes_a_hard_truncated_title_from_source_copy() -> None:
    slide = _slide(
        "slide:v5:episode-1:001",
        title="本节课的核心目标是建立局部解剖学的空",
        content="本节课的核心目标是建立局部解剖学的空间定位基础。",
    )
    slide["quality"]["title_character_budget"] = 18

    repaired, history = repair_semantic_slides_v5([slide], max_rounds=2)

    assert repaired[0]["title"] == "局部解剖学的空间定位基础"
    assert any(
        item["action"] == "complete_title_from_existing_copy"
        for item in history
    )


def test_export_audit_detects_inner_text_clipping_and_unreadable_body_font(
    tmp_path: Path,
) -> None:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(3.0), Inches(0.35))
    box.text_frame.word_wrap = True
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = "这是一段必然超过文本框内部容量并造成裁切的正文。" * 5
    paragraph.font.size = Pt(10)
    path = tmp_path / "clipped.pptx"
    presentation.save(path)

    report = audit_exported_pptx(path, expected_slide_count=1)
    codes = {issue["code"] for issue in report["blockers"]}

    assert "exported_text_frame_overflow" in codes
    assert "exported_body_font_below_16pt" in codes


def test_rendered_page_ocr_checks_every_page_for_missing_visible_text(
    tmp_path: Path,
) -> None:
    presentation = Presentation()
    first = presentation.slides.add_slide(presentation.slide_layouts[6])
    first.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(8), Inches(1)
    ).text = "第一页标题与完整正文都必须可见"
    second = presentation.slides.add_slide(presentation.slide_layouts[6])
    second.shapes.add_textbox(
        Inches(0.5), Inches(0.5), Inches(8), Inches(1)
    ).text = "第二页标题与完整正文也必须可见"
    images = [tmp_path / "page-1.png", tmp_path / "page-2.png"]
    for image in images:
        image.touch()

    report = audit_rendered_slide_images(
        presentation,
        images,
        ocr_runner=lambda path: (
            "第一页标题与完整正文都必须可见"
            if path.name == "page-1.png"
            else "第二页标题"
        ),
    )

    assert report["page_count"] == 2
    assert report["checked_pages"] == 2
    assert {
        (issue["code"], issue.get("page")) for issue in report["blockers"]
    } == {("exported_ocr_text_missing_or_clipped", 2)}


def test_render_repair_changes_only_the_audited_page() -> None:
    first = _slide(
        "slide:v5:episode-1:001",
        title="这个标题在导出后发生了意外换行并需要压缩",
        content="同一段正文。同一段正文。",
    )
    first["blocks"].append(deepcopy(first["blocks"][0]))
    first["blocks"][1]["block_id"] = "duplicate"
    second = _slide(
        "slide:v5:episode-2:001",
        title="第二页保持不变",
        content="这一页的文字容量和布局都满足要求。",
        episode_id="episode-2",
    )
    original_second = deepcopy(second)
    review = {
        "issues": [
            {"severity": "critical", "code": "exported_text_frame_overflow", "page": 1},
            {"severity": "critical", "code": "exported_title_unexpected_wrap", "page": 1},
        ]
    }

    repaired, history = repair_render_slides_v5([first, second], review, round_index=1)

    assert repaired[1] == original_second
    assert repaired[0]["quality"]["requested_layout"] == "editorial-body"
    assert len(repaired[0]["blocks"]) == 1
    assert len(repaired[0]["title"]) <= 24
    assert {item["page_id"] for item in history} == {"slide:v5:episode-1:001"}


def test_render_repair_shortens_a_twenty_character_chinese_title() -> None:
    slide = _slide(
        "slide:v5:chapter:long-title",
        title="线性映射的定义，并能判断给定映射是否线性",
        content="什么样的映射不会破坏向量空间中的线性结构？",
    )
    assert len(slide["title"]) == 20
    review = {
        "issues": [{
            "severity": "critical",
            "code": "exported_title_unexpected_wrap",
            "page": 1,
        }],
    }

    repaired, history = repair_render_slides_v5(
        [slide],
        review,
        round_index=1,
    )

    assert len(repaired[0]["title"]) <= 18
    assert history[0]["actions"] == ["shorten_title_from_existing_copy"]


def test_process_result_in_source_grounded_takeaway_satisfies_contract() -> None:
    slide = _slide(
        "slide:v5:process:source-result",
        title="先检查条件，再形成判断",
        content="先检查两个条件；两项都成立才能形成最终判断。",
        scene_kind="process",
    )
    slide["takeaway"] = "两项都成立才能形成最终判断。"

    report = build_slide_deck_quality_v5([slide])

    assert "process_result_missing" not in {
        issue["code"] for issue in report["issues"]
    }


def test_worked_example_accepts_source_backed_linked_answer_as_conclusion() -> None:
    prompt = _slide(
        "slide:v5:example:prompt",
        title="根据条件完成判断",
        content="请依据给定条件完成推导并判断结论。",
        scene_kind="worked_example",
        block_type="exercise",
    )
    prompt["blocks"][0]["metadata"].update({
        "question_id": "question-linked",
        "question_mode": "closed",
    })
    answer = _slide(
        "slide:v5:example:answer",
        title="推导结果与理由",
        content="结论成立，因为两个来源条件均已满足。",
        scene_kind="practice_feedback",
        block_type="callout",
    )
    answer["blocks"][0]["title"] = "答案"
    answer["blocks"][0]["metadata"].update({
        "answer_for": "question-linked",
        "source_fragment_id": "fragment-answer",
    })

    report = build_slide_deck_quality_v5([prompt, answer])

    assert "worked_example_conclusion_missing" not in {
        issue["code"] for issue in report["issues"]
    }


def test_open_worked_example_accepts_source_backed_analysis_as_conclusion() -> None:
    slide = _slide(
        "slide:v5:example:analysis",
        title="根据条件解释判断依据",
        content="推导依据：来源材料中的两个条件共同支持这一判断。",
        scene_kind="worked_example",
        block_type="exercise",
    )
    slide["blocks"][0]["metadata"].update({
        "question_mode": "open_discussion",
        "fragment_ids": ["fragment-analysis"],
    })

    report = build_slide_deck_quality_v5([slide])

    assert "worked_example_conclusion_missing" not in {
        issue["code"] for issue in report["issues"]
    }


def test_four_domain_styles_use_the_same_generic_quality_contract() -> None:
    structure = _slide(
        "slide:v5:structure:001",
        title="对象之间的空间关系决定观察顺序",
        content="先辨认两个对象，再依据材料说明它们的相对位置与连接关系。",
        scene_kind="structure",
    )
    structure["visuals"] = [{
        "kind": "relational_diagram",
        "alt_text": "两个对象及其空间关系",
    }]
    structure["quality"]["need_visual"] = True

    formula = _slide(
        "slide:v5:formula:001",
        title="公式与解释共同支持数量结论",
        content="公式给出数量关系，正文解释变量含义以及结论成立的条件。",
        scene_kind="concept",
    )
    formula["quality"]["resolved_layout"] = "formula-explanation"
    formula["visuals"] = [{"kind": "formula", "alt_text": "来源公式"}]

    narrative = _slide(
        "slide:v5:narrative:001",
        title="事件顺序解释了结果如何形成",
        content="材料先交代背景，再呈现关键转折，最后说明该转折带来的影响。",
        scene_kind="concept",
    )

    process = _slide(
        "slide:v5:process:001",
        title="完整流程从输入走向可检查结果",
        content="流程从明确输入开始，依次完成关键步骤，并产出可以复核的结果。",
        scene_kind="process",
    )
    process["quality"]["process_result"] = "产出可以复核的结果"

    report = build_slide_deck_quality_v5([
        structure,
        formula,
        narrative,
        process,
    ])

    assert report["passed"] is True
    assert report["blockers"] == []
