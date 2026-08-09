from __future__ import annotations

import json
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import Mock, patch

import pytest
from pptx import Presentation

import slide_deck_renderer
from course_document import CourseBlock, CourseDocument, CourseSection
from slide_deck import SlideDeckContent, validate_slide_deck
from slide_deck_renderer import (
    V5_LAYOUT_RENDERER_NAMES,
    _render_editorial_body,
    _render_slide,
    export_structured_slide_deck,
)
from slide_deck_v3 import (
    ContentFragmentV1,
    PlannedPageV2,
    SlideAllocationPlanV2,
    fragment_course_document,
    slide_deck_variant_key,
)
from slide_deck_v4 import allocation_from_story_plan_v2
from slide_deck_v5 import (
    _bounded_title,
    _chapter_recap_slide,
    _enrich_practice_feedback_slides_v5,
    _split_practice_feedback_capacity_v5,
    _title_with_continuation_sequence,
    _v5_fragment_groups_for_profile,
    _v5_group_kind_for_profile,
    allocation_from_story_plan_v5,
    apply_page_contract_v5,
    build_signature_v5,
    compact_story_plan_v5,
    compile_deck_outline_v5,
    compile_page_title_v5,
    compile_slide_deck_v5,
    finalize_v5_candidate_contract,
    finalize_v5_quality_report,
    repair_final_page_contracts_v5,
    resolve_page_contract_v5,
    summarize_v5_slide_counts,
    v5_contract_issues,
)
from slide_quality_v5 import _concise_existing_title, build_slide_deck_quality_v5
from slide_story_plan import (
    ChapterStoryV2,
    ClaimSourceV2,
    CommunicationBriefV2,
    SlideStoryPlanV2,
    StoryBeatV2,
    StorySourceRevisionsV2,
    TeachingEpisodeV2,
)
from slide_visuals import deterministic_visual_plan


def _beat(chapter_index: int, scene: str) -> StoryBeatV2:
    return StoryBeatV2(
        beat_id=f"beat-{chapter_index}-{scene}",
        beat_role="driving_question" if scene == "chapter_entry" else "closure",
        teaching_job="建立学习问题" if scene == "chapter_entry" else "完成章节回顾",
        primary_claim_source=ClaimSourceV2(
            kind="learning_objective",
            text=f"理解第{chapter_index}章的核心关系",
            objective_id=f"objective-{chapter_index}",
        ),
        layout_intent="hero-statement",
        renderer_layout="hero-statement",
        layout_family="statement",
        layout_selection_reason="test fixture",
    )


def _chapter(index: int) -> ChapterStoryV2:
    return ChapterStoryV2(
        chapter_id=f"chapter-{index}",
        title=f"第{index}章 主题{index}",
        driving_question=f"主题{index}解决什么问题？",
        learning_objective=f"理解主题{index}的核心关系",
        episodes=[
            TeachingEpisodeV2(
                episode_id=f"episode-{index}-entry",
                scene_kind="chapter_entry",
                teaching_job="建立学习问题",
                beats=[_beat(index, "chapter_entry")],
            ),
            TeachingEpisodeV2(
                episode_id=f"episode-{index}-recap",
                scene_kind="chapter_recap",
                teaching_job="完成章节回顾",
                beats=[_beat(index, "chapter_recap")],
            ),
        ],
    )


def _story(chapter_count: int) -> SlideStoryPlanV2:
    return SlideStoryPlanV2(
        plan_id="story-v5-test",
        mode="teaching",
        theme="qizhi-classroom",
        communication_brief=CommunicationBriefV2(
            audience="本科生",
            course_goal="建立一套可迁移的分析框架",
            central_question="如何从基本概念走向完整分析？",
            expected_learning_results=["解释概念", "应用方法"],
        ),
        source_revisions=StorySourceRevisionsV2(
            course_document_revision="doc-rev-1",
            teaching_plan_revision="plan-rev-1",
            knowledge_base_revision="kb-rev-1",
            coherence_contract_revision="coherence-rev-1",
        ),
        chapters=[_chapter(index) for index in range(1, chapter_count + 1)],
    )


def _document(chapter_count: int) -> CourseDocument:
    return CourseDocument(
        course_id="course-v5-test",
        title="课程 V5 测试",
        document_revision="doc-rev-1",
        sections=[
            CourseSection(
                section_id=f"chapter-{index}",
                title=f"第{index}章 主题{index}",
                position=index - 1,
                level=1,
                learning_objective=f"理解主题{index}的核心关系",
            )
            for index in range(1, chapter_count + 1)
        ],
    )


def test_v5_allocation_stabilizes_reversed_story_beats_by_source_order() -> None:
    document = _document(1)
    early = ContentFragmentV1(
        fragment_id="fragment-early",
        section_id="chapter-1",
        block_id="block-early",
        kind="paragraph",
        text="Earlier source statement with enough detail for a complete teaching claim.",
        ordinal=10,
        source_hash="hash-early",
        role="concept",
        source_kind="course_block",
    )
    late = ContentFragmentV1(
        fragment_id="fragment-late",
        section_id="chapter-1",
        block_id="block-late",
        kind="paragraph",
        text="Later source statement with enough detail for a complete teaching claim.",
        ordinal=20,
        source_hash="hash-late",
        role="concept",
        source_kind="course_block",
    )
    story = _story(1)
    base_chapter = story.chapters[0]
    late_beat = _beat(1, "concept").model_copy(update={
        "beat_id": "beat-late",
        "fragment_ids": [late.fragment_id],
    })
    early_beat = _beat(1, "concept").model_copy(update={
        "beat_id": "beat-early",
        "fragment_ids": [early.fragment_id],
    })
    chapter = base_chapter.model_copy(update={
        "episodes": [
            base_chapter.episodes[0],
            TeachingEpisodeV2(
                episode_id="episode-late",
                scene_kind="concept",
                teaching_job="Explain the later source statement",
                beats=[late_beat],
            ),
            TeachingEpisodeV2(
                episode_id="episode-early",
                scene_kind="concept",
                teaching_job="Explain the earlier source statement",
                beats=[early_beat],
            ),
            base_chapter.episodes[-1],
        ],
    })
    story = story.model_copy(update={"chapters": [chapter]})

    allocation, _ = allocation_from_story_plan_v5(
        document,
        [early, late],
        story,
    )

    assert [
        fragment_id
        for page in allocation.pages
        for fragment_id in page.fragment_ids
    ] == [early.fragment_id, late.fragment_id]


def test_v5_allocation_closes_source_lists_and_uses_continuations() -> None:
    document = _document(1)
    fragments = [
        ContentFragmentV1(
            fragment_id="fragment-lead",
            section_id="chapter-1",
            block_id="block-list",
            kind="list_item",
            text="来源明确包括5个分支：",
            ordinal=10,
            source_hash="hash-lead",
            role="concept",
            source_kind="course_block",
        ),
        *[
            ContentFragmentV1(
                fragment_id=f"fragment-item-{index}",
                section_id="chapter-1",
                block_id="block-list",
                kind="list_item",
                text=f"Branch {index} is explicitly described by the source.",
                ordinal=10 + index,
                source_hash=f"hash-item-{index}",
                role="concept",
                source_kind="course_block",
            )
            for index in range(1, 6)
        ],
        ContentFragmentV1(
            fragment_id="fragment-next-heading",
            section_id="chapter-1",
            block_id="block-next",
            kind="heading",
            text="Next independent topic",
            ordinal=20,
            source_hash="hash-next",
            role="concept",
            source_kind="course_block",
        ),
    ]
    story = _story(1)
    chapter = story.chapters[0]
    source_beat = _beat(1, "concept").model_copy(update={
        "beat_id": "beat-list",
        "fragment_ids": ["fragment-lead"],
        "renderer_layout": "question",
    })
    story = story.model_copy(update={
        "chapters": [chapter.model_copy(update={
            "episodes": [
                chapter.episodes[0],
                TeachingEpisodeV2(
                    episode_id="episode-list",
                    scene_kind="concept",
                    teaching_job="Explain the source list",
                    beats=[source_beat],
                ),
                chapter.episodes[-1],
            ],
        })],
    })

    allocation, _ = allocation_from_story_plan_v5(document, fragments, story)
    content_pages = [page for page in allocation.pages if page.fragment_ids]
    allocated_ids = [
        fragment_id
        for page in content_pages
        for fragment_id in page.fragment_ids
    ]

    assert allocated_ids == [
        "fragment-lead",
        "fragment-item-1",
        "fragment-item-2",
        "fragment-item-3",
        "fragment-item-4",
        "fragment-item-5",
    ]
    assert len(content_pages) >= 2
    assert content_pages[1].continuation_of == content_pages[0].page_id


def test_quality_fallback_prefers_explicit_source_group_kind() -> None:
    fragment = ContentFragmentV1(
        fragment_id="fragment-clinical-context",
        section_id="chapter-1",
        block_id="block-clinical-context",
        kind="heading",
        text="Clinical context",
        ordinal=10,
        source_hash="hash-clinical-context",
        role="concept",
        source_kind="course_block",
    )
    inferred_practice = SimpleNamespace(
        semantic_unit_id="semantic-clinical-context",
        primary_role="checkpoint",
        presentation_intent="practice_feedback",
    )

    assert _v5_group_kind_for_profile(
        [fragment],
        {fragment.fragment_id: inferred_practice},
        profile="quality_fallback",
    ) == "concept"

    continuation = fragment.model_copy(update={
        "fragment_id": "fragment-clinical-context-detail",
        "block_id": "block-clinical-context-detail",
        "kind": "paragraph",
        "text": "Source-bound detail stored in a separate legacy block.",
        "ordinal": 11,
        "source_hash": "hash-clinical-context-detail",
    })
    assert _v5_fragment_groups_for_profile(
        [fragment, continuation],
        profile="quality_fallback",
    ) == [[fragment, continuation]]


def test_v5_compiles_directly_to_final_ids_and_rebuilds_a_stale_visual_plan() -> None:
    document = _document(1)
    variant_key = slide_deck_variant_key("teaching", "qizhi-classroom")
    supplied_allocation = SlideAllocationPlanV2(
        title=document.title,
        mode="teaching",
        theme="qizhi-classroom",
        variant_key=variant_key,
        source_document_revision=document.document_revision,
        pages=[
            PlannedPageV2(page_id="slide:title", layout="cover"),
            PlannedPageV2(page_id="slide:roadmap", layout="roadmap"),
            PlannedPageV2(
                page_id="slide:v4:0001",
                layout="editorial-body",
                chapter_id="chapter-1",
            ),
        ],
    )
    supplied_visual_plan = deterministic_visual_plan(
        document,
        supplied_allocation,
        [],
    )
    content = compile_slide_deck_v5(
        document,
        {},
        story_plan=_story(1),
        allocation_plan=supplied_allocation,
        visual_plan=supplied_visual_plan,
    )

    final_page_ids = [slide["unit_id"] for slide in content["slides"]]
    allocation_page_ids = [
        page["page_id"] for page in content["allocation_plan"]["pages"]
    ]
    visual_page_ids = [
        page["page_id"] for page in content["visual_plan"]["pages"]
    ]
    assert all(page_id.startswith("slide:v5:") for page_id in final_page_ids)
    assert allocation_page_ids == visual_page_ids
    assert content["deck_brief"]["fallback_reason"] == (
        "v5_final_page_ids_visual_plan_rebuilt"
    )


def test_v5_only_streams_final_contract_candidate_slides() -> None:
    events: list[dict] = []

    content = compile_slide_deck_v5(
        _document(1),
        {},
        story_plan=_story(1),
        progress_callback=events.append,
    )

    candidate_events = [event for event in events if event.get("event") == "slide_upsert"]
    reset_events = [event for event in events if event.get("event") == "slide_reset"]
    assert reset_events == [{
        "event": "slide_reset",
        "progress": 97,
        "stage": "v5_candidate",
        "engine_schema": "slide_deck_v5",
        "candidate_stage": "final_contract",
    }]
    assert len(candidate_events) == len(content["slides"])
    assert [event["slide"] for event in candidate_events] == content["slides"]
    assert all(
        event.get("engine_schema") == "slide_deck_v5"
        and event.get("candidate_stage") == "final_contract"
        for event in candidate_events
    )


def test_v5_candidate_exposes_source_contract_dispositions_and_terminal_state() -> None:
    document = _document(1).model_copy(update={
        "blocks": [CourseBlock(
            block_id="block-source-contract",
            section_id="chapter-1",
            position=0,
            role="concept",
            payload={
                "markdown": (
                    "系统边界决定系统与环境之间能够发生的交换。"
                    "识别边界后，需要分别检查物质交换与能量交换，"
                    "再根据两类交换是否存在判断系统类型。"
                ),
            },
        )],
    })
    fragment = fragment_course_document(document)[0]
    story = _story(1)
    chapter = story.chapters[0]
    source_beat = _beat(1, "concept").model_copy(update={
        "beat_id": "beat-source-contract",
        "fragment_ids": [fragment.fragment_id],
    })
    story = story.model_copy(update={
        "chapters": [chapter.model_copy(update={
            "episodes": [
                chapter.episodes[0],
                TeachingEpisodeV2(
                    episode_id="episode-source-contract",
                    scene_kind="concept",
                    teaching_job="解释系统边界的判断方法",
                    beats=[source_beat],
                ),
                chapter.episodes[-1],
            ],
        })],
    })

    content = compile_slide_deck_v5(
        document,
        {},
        story_plan=story,
    )

    assert content["ppt_source_contract_v1"]["source_document_revision"] == (
        document.document_revision
    )
    assert content["candidate_status"] in {"v5_ready", "v5_needs_manual_edit"}
    dispositions = content["source_dispositions"]
    fragment_ids = {
        str(item["fragment_id"])
        for item in content.get("fragment_manifest") or []
    }
    assert {str(item["fragment_id"]) for item in dispositions} == fragment_ids
    assert all(item["disposition"] in {
        "rendered",
        "rendered_in_safe_layout",
        "moved_to_appendix",
        "needs_manual_edit",
        "intentionally_excluded_with_reason",
    } for item in dispositions)


def test_v5_readable_layout_warning_publishes_manual_edit_candidate() -> None:
    content = compile_slide_deck_v5(
        _document(1),
        {},
        story_plan=_story(1),
    )
    page_id = content["slides"][0]["unit_id"]
    quality = {
        **content["quality_report"],
        "passed": True,
        "status": "ready",
        "blockers": [],
        "warnings": [
            *(content["quality_report"].get("warnings") or []),
            {
                "severity": "warning",
                "dimension": "layout_export",
                "code": "render_review_manual_adjustment",
                "page_id": page_id,
                "message": "本页内容完整，但建议人工微调视觉间距。",
            },
        ],
    }

    finalized = finalize_v5_candidate_contract(content, quality)

    assert finalized["passed"] is True
    assert finalized["candidate_status"] == "v5_needs_manual_edit"
    assert content["candidate_status"] == "v5_needs_manual_edit"
    assert content["manual_edit_required"] == [{
        "page_id": page_id,
        "reasons": [{
            "code": "render_review_manual_adjustment",
            "message": "本页内容完整，但建议人工微调视觉间距。",
        }],
    }]
    assert content["slides"][0]["quality"]["manual_edit_required"] is True


def test_v5_source_disposition_keeps_the_strongest_page_outcome() -> None:
    content = {
        "fragment_manifest": [{"fragment_id": "fragment-manual"}],
        "allocation_plan": {"pages": []},
        "exclusions": [],
        "slides": [
            {
                "unit_id": "manual-page",
                "blocks": [],
                "quality": {
                    "fragment_ids": ["fragment-manual"],
                    "manual_edit_required": True,
                    "manual_edit_reasons": [{
                        "code": "layout_spacing",
                        "message": "需要人工微调间距。",
                    }],
                },
            },
            {
                "unit_id": "derived-recap",
                "blocks": [],
                "quality": {"fragment_ids": ["fragment-manual"]},
            },
        ],
    }

    finalize_v5_candidate_contract(content, {
        "passed": True,
        "score": 100,
        "issues": [],
        "warnings": [],
        "blockers": [],
    })

    assert content["source_dispositions"] == [{
        "fragment_id": "fragment-manual",
        "disposition": "needs_manual_edit",
        "page_id": "manual-page",
    }]


def test_v5_candidate_normalizes_string_manual_edit_reasons() -> None:
    content = {
        "fragment_manifest": [],
        "allocation_plan": {"pages": []},
        "exclusions": [],
        "slides": [{
            "unit_id": "subject-review-page",
            "blocks": [],
            "quality": {
                "manual_edit_required": True,
                "manual_edit_reasons": [
                    "presentation_grammar_mismatch",
                    "required_subject_source_missing:code",
                ],
            },
        }],
    }

    report = finalize_v5_candidate_contract(content, {
        "passed": True,
        "score": 100,
        "issues": [],
        "warnings": [],
        "blockers": [],
    })

    assert report["candidate_status"] == "v5_needs_manual_edit"
    assert content["manual_edit_required"] == [{
        "page_id": "subject-review-page",
        "reasons": [
            {
                "code": "presentation_grammar_mismatch",
                "message": "页面版式未完全匹配教学意图，请手动检查视觉表达。",
            },
            {
                "code": "required_subject_source_missing",
                "message": "课程原文缺少建议的 code 学科工件，请手动补充或确认。",
                "representation_kind": "code",
            },
        ],
    }]


def test_outline_groups_eight_chapters_into_at_most_six_source_bound_sections() -> None:
    outline = compile_deck_outline_v5(_document(8), _story(8))

    assert 3 <= len(outline.agenda_sections) <= 6
    assert [
        chapter_id
        for section in outline.agenda_sections
        for chapter_id in section.source_chapter_ids
    ] == [f"chapter-{index}" for index in range(1, 9)]
    assert outline.cover.subtitle == "建立一套可迁移的分析框架"
    assert outline.closing.kind == "course_synthesis"


def test_outline_does_not_invent_sections_for_a_single_chapter_course() -> None:
    outline = compile_deck_outline_v5(_document(1), _story(1))

    assert len(outline.agenda_sections) == 1
    assert outline.agenda_sections[0].source_chapter_ids == ["chapter-1"]


def test_v5_story_compaction_selects_complete_semantic_groups_per_section() -> None:
    document = CourseDocument(
        course_id="course-v5-compaction",
        title="课程压缩测试",
        document_revision="doc-rev-1",
        sections=[
            CourseSection(
                section_id="chapter-1",
                title="第一章",
                position=0,
                level=1,
            ),
            CourseSection(
                section_id="section-1",
                parent_section_id="chapter-1",
                title="1.1 核心主题",
                position=1,
                level=2,
            ),
        ],
    )
    raw = [
        ("title", "heading", "1.1 核心主题"),
        ("core-heading", "heading", "核心概念与背景"),
        ("core-body", "paragraph", "核心概念通过已有正文建立。"),
        ("method-heading", "heading", "技术实现与方法"),
        ("method-body", "paragraph", "先识别条件，再选择步骤。"),
        ("case-heading", "heading", "实战案例"),
        ("case-body", "paragraph", "案例用于检验抽象判断。"),
        ("practice-heading", "heading", "思考与挑战"),
        ("practice-body", "list_item", "请说明结论成立的边界。"),
    ]
    fragments = [
        ContentFragmentV1(
            fragment_id=fragment_id,
            section_id="section-1",
            block_id="section-1-body",
            kind=kind,  # type: ignore[arg-type]
            text=text,
            ordinal=index,
            source_hash=f"hash-{index}",
            role="concept",
            source_kind="course_block",
        )
        for index, (fragment_id, kind, text) in enumerate(raw)
    ]

    compact = compact_story_plan_v5(document, _story(1), fragments)
    chapter = compact.chapters[0]

    assert [episode.scene_kind for episode in chapter.episodes] == [
        "chapter_entry",
        "concept",
        "method",
        "worked_example",
        "practice_feedback",
        "chapter_recap",
    ]
    selected_ids = {
        fragment_id
        for episode in chapter.episodes
        for beat in episode.beats
        for fragment_id in beat.fragment_ids
    }
    assert {
        "core-body",
        "method-body",
        "case-body",
        "practice-body",
    } <= selected_ids
    allocation, _ = allocation_from_story_plan_v2(
        document,
        fragments,
        compact,
    )
    teaching_pages = [
        page
        for page in allocation.pages
        if page.fragment_ids
    ]
    assert 4 <= len(teaching_pages) <= 6
    assert all(
        exclusion.reason == "v5_semantic_core"
        for exclusion in allocation.exclusions
    )
    assert "method-body" not in {
        exclusion.fragment_id for exclusion in allocation.exclusions
    }
    refined = compact.model_copy(update={
        "planner": "ai",
        "chapters": [
            compact.chapters[0].model_copy(update={
                "episodes": [
                    episode.model_copy(update={
                        "beats": [
                            beat.model_copy(update={
                                "layout_selection_reason": (
                                    "ai_source_bound_directive"
                                ),
                            })
                            for beat in episode.beats
                        ],
                    })
                    for episode in compact.chapters[0].episodes
                ],
            }),
        ],
    })

    recompacted = compact_story_plan_v5(document, refined, fragments)

    assert all(
        beat.layout_selection_reason == "ai_source_bound_directive"
        for chapter in recompacted.chapters
        for episode in chapter.episodes[1:-1]
        for beat in episode.beats
    )
    refined_allocation, _ = allocation_from_story_plan_v2(
        document,
        fragments,
        recompacted,
    )

    assert not any(page.appendix for page in refined_allocation.pages)
    assert "method-body" not in {
        exclusion.fragment_id for exclusion in refined_allocation.exclusions
    }
    assert all(
        exclusion.reason == "v5_semantic_core"
        for exclusion in refined_allocation.exclusions
    )


def test_quality_fallback_compacts_sibling_concepts_into_one_presentation_page() -> None:
    document = CourseDocument(
        course_id="course-v5-presentation-native-fallback",
        title="Presentation native fallback",
        document_revision="doc-rev-1",
        sections=[
            CourseSection(
                section_id="chapter-1",
                title="Chapter one",
                position=0,
                level=1,
            ),
            CourseSection(
                section_id="section-1",
                parent_section_id="chapter-1",
                title="1.1 Runtime state",
                position=1,
                level=2,
            ),
        ],
    )
    raw = [
        ("heading-a", "heading", "State ownership"),
        ("body-a", "paragraph", "One component owns the authoritative runtime state."),
        ("heading-b", "heading", "State updates"),
        ("body-b", "paragraph", "Commands update that state through one controlled path."),
        ("heading-c", "heading", "State verification"),
        ("body-c", "paragraph", "Observable output verifies the update after each command."),
    ]
    fragments = [
        ContentFragmentV1(
            fragment_id=fragment_id,
            section_id="section-1",
            block_id="section-1-body",
            kind=kind,  # type: ignore[arg-type]
            text=text,
            ordinal=index,
            source_hash=f"hash-{index}",
            role="concept",
            source_kind="course_block",
        )
        for index, (fragment_id, kind, text) in enumerate(raw)
    ]

    compact = compact_story_plan_v5(
        document,
        _story(1),
        fragments,
        profile="quality_fallback",
    )
    concept_episodes = [
        episode
        for episode in compact.chapters[0].episodes
        if episode.scene_kind == "concept"
    ]

    assert len(concept_episodes) == 1
    assert {
        "body-a",
        "body-b",
        "body-c",
    } <= set(concept_episodes[0].beats[0].fragment_ids)


def test_none_visual_does_not_occupy_a_visual_slot() -> None:
    slide = {
        "unit_id": "slide:v5:none-visual",
        "layout": "concept",
        "slide_purpose": "concept",
        "scene_kind": "concept",
        "title": "One source claim remains text only",
        "blocks": [{
            "block_id": "claim",
            "type": "statement",
            "content": "One complete source-backed claim.",
            "items": [],
        }],
        "visuals": [{"kind": "none", "alt_text": ""}],
        "quality": {"requested_layout": "classification-3"},
    }

    contract = resolve_page_contract_v5(slide)

    assert contract.visual_decision == "none"
    assert contract.resolved_layout == "editorial-body"


def test_v5_compaction_excludes_formula_without_source_explanation() -> None:
    document = CourseDocument(
        course_id="course-v5-formula-compaction",
        title="公式压缩测试",
        document_revision="doc-rev-1",
        sections=[
            CourseSection(
                section_id="chapter-1",
                title="第一章",
                position=0,
                level=1,
            ),
            CourseSection(
                section_id="section-1",
                parent_section_id="chapter-1",
                title="1.1 核心主题",
                position=1,
                level=2,
            ),
        ],
    )
    fragments = [
        ContentFragmentV1(
            fragment_id=fragment_id,
            section_id="section-1",
            block_id="section-1-body",
            kind=kind,  # type: ignore[arg-type]
            text=text,
            ordinal=index,
            source_hash=f"hash-{index}",
            role="concept",
            source_kind="course_block",
        )
        for index, (fragment_id, kind, text) in enumerate([
            ("core-heading", "heading", "核心概念"),
            ("core-body", "paragraph", "正文解释用于建立概念。"),
            ("formula-heading", "heading", "热力学第一定律"),
            ("formula-only", "formula", r"$$ \Delta U = Q - W $$"),
        ])
    ]

    compact = compact_story_plan_v5(document, _story(1), fragments)
    allocation, _ = allocation_from_story_plan_v2(
        document,
        fragments,
        compact,
    )

    allocated_ids = {
        fragment_id
        for page in allocation.pages
        for fragment_id in page.fragment_ids
    }
    excluded_ids = {
        exclusion.fragment_id for exclusion in allocation.exclusions
    }
    assert "core-body" in allocated_ids
    assert {"formula-heading", "formula-only"} <= excluded_ids
    assert not any(page.appendix for page in allocation.pages)


def test_ai_refinement_keeps_formula_and_source_explanation_on_one_page() -> None:
    document = CourseDocument(
        course_id="course-v5-formula-binding",
        title="公式绑定测试",
        document_revision="doc-rev-1",
        sections=[
            CourseSection(
                section_id="chapter-1",
                title="第一章",
                position=0,
                level=1,
            ),
            CourseSection(
                section_id="section-1",
                parent_section_id="chapter-1",
                title="1.1 能量守恒",
                position=1,
                level=2,
            ),
        ],
    )
    fragments = [
        ContentFragmentV1(
            fragment_id=fragment_id,
            section_id="section-1",
            block_id="section-1-body",
            kind=kind,  # type: ignore[arg-type]
            text=text,
            ordinal=index,
            source_hash=f"hash-{index}",
            role="concept",
            source_kind="course_block",
        )
        for index, (fragment_id, kind, text) in enumerate([
            ("formula-heading", "heading", "热力学第一定律"),
            ("formula", "formula", r"$$ \Delta U = Q - W $$"),
            ("formula-explanation", "paragraph", "其中各符号分别表示内能、热量和功。"),
        ])
    ]
    compact = compact_story_plan_v5(document, _story(1), fragments)
    refined = compact.model_copy(update={
        "planner": "ai",
        "chapters": [
            chapter.model_copy(update={
                "episodes": [
                    episode.model_copy(update={
                        "beats": [
                            beat.model_copy(update={
                                "layout_selection_reason": (
                                    "ai_source_bound_directive"
                                ),
                            })
                            for beat in episode.beats
                        ],
                    })
                    for episode in chapter.episodes
                ],
            })
            for chapter in compact.chapters
        ],
    })

    allocation, _ = allocation_from_story_plan_v2(
        document,
        fragments,
        refined,
    )
    formula_pages = [
        page for page in allocation.pages
        if "formula" in page.fragment_ids
    ]

    assert len(formula_pages) == 1
    assert {
        "formula",
        "formula-explanation",
    } <= set(formula_pages[0].fragment_ids)


def test_ai_refinement_missing_required_formula_is_recompacted_from_source() -> None:
    document = CourseDocument(
        course_id="course-v5-required-formula",
        title="Required formula recovery",
        document_revision="doc-rev-1",
        sections=[
            CourseSection(
                section_id="chapter-1",
                title="Linear relations",
                position=0,
                level=1,
            ),
            CourseSection(
                section_id="section-1",
                parent_section_id="chapter-1",
                title="Matrix action",
                position=1,
                level=2,
            ),
        ],
    )
    fragments = [
        ContentFragmentV1(
            fragment_id=fragment_id,
            section_id="section-1",
            block_id="section-1-body",
            kind=kind,  # type: ignore[arg-type]
            text=text,
            ordinal=index,
            source_hash=f"hash-{index}",
            role="concept",
            source_kind="course_block",
        )
        for index, (fragment_id, kind, text) in enumerate([
            ("prose", "paragraph", "A linear map preserves addition and scaling."),
            ("formula-heading", "heading", "Matrix representation"),
            ("formula", "formula", r"$$ T(x) = Ax $$"),
            ("formula-explanation", "paragraph", "The matrix A records the action of T."),
        ])
    ]
    story = _story(1)
    chapter = story.chapters[0]
    prose_beat = _beat(1, "concept").model_copy(update={
        "beat_id": "beat-ai-prose-only",
        "fragment_ids": ["prose"],
        "layout_selection_reason": "ai_source_bound_directive",
    })
    refined = story.model_copy(update={
        "planner": "ai",
        "planning_diagnostics": {
            "subject_presentation_contract": {
                "schema_version": "subject_presentation_contract_v1",
                "profile_id": "math_formal",
                "primary_mode": "math_formal",
                "required_representation_kinds": ["formula"],
                "optional_representation_kinds": ["diagram", "table"],
                "characteristic_fragment_ids": {"formula": ["formula"]},
                "chapter_requirements": [{
                    "chapter_id": "chapter-1",
                    "required_representation_kinds": ["formula"],
                    "minimum_artifact_count": 1,
                }],
                "classification_confidence": 1.0,
                "classification_source": "test",
            },
        },
        "chapters": [chapter.model_copy(update={
            "episodes": [
                chapter.episodes[0],
                TeachingEpisodeV2(
                    episode_id="episode-ai-prose-only",
                    scene_kind="concept",
                    teaching_job="Explain the source prose",
                    beats=[prose_beat],
                ),
                chapter.episodes[-1],
            ],
        })],
    })

    recompacted = compact_story_plan_v5(document, refined, fragments)
    formula_beats = [
        beat
        for episode in recompacted.chapters[0].episodes[1:-1]
        for beat in episode.beats
        if "formula" in beat.subject_artifact_kinds
    ]

    assert len(formula_beats) == 1
    assert {
        "formula",
        "formula-explanation",
    } <= set(formula_beats[0].fragment_ids)


def test_v5_compaction_includes_required_formula_from_chapter_root() -> None:
    document = CourseDocument(
        course_id="course-v5-root-formula",
        title="Root formula recovery",
        document_revision="doc-rev-1",
        sections=[
            CourseSection(
                section_id="chapter-1",
                title="Linear transformations",
                position=0,
                level=1,
            ),
            CourseSection(
                section_id="section-1",
                parent_section_id="chapter-1",
                title="Geometric interpretation",
                position=1,
                level=2,
            ),
        ],
    )
    fragments = [
        ContentFragmentV1(
            fragment_id="root-formula-heading",
            section_id="chapter-1",
            block_id="chapter-1-body",
            kind="heading",
            text="Matrix representation",
            ordinal=0,
            source_hash="hash-root-heading",
            role="concept",
            source_kind="course_block",
        ),
        ContentFragmentV1(
            fragment_id="root-formula",
            section_id="chapter-1",
            block_id="chapter-1-body",
            kind="formula",
            text=r"$$ T(x) = Ax $$",
            ordinal=1,
            source_hash="hash-root-formula",
            role="concept",
            source_kind="course_block",
        ),
        ContentFragmentV1(
            fragment_id="root-formula-explanation",
            section_id="chapter-1",
            block_id="chapter-1-body",
            kind="paragraph",
            text="The matrix A records the action of T in the selected basis.",
            ordinal=2,
            source_hash="hash-root-explanation",
            role="concept",
            source_kind="course_block",
        ),
        ContentFragmentV1(
            fragment_id="child-prose",
            section_id="section-1",
            block_id="section-1-body",
            kind="paragraph",
            text="The image of each basis vector determines the transformation.",
            ordinal=3,
            source_hash="hash-child-prose",
            role="concept",
            source_kind="course_block",
        ),
    ]
    story = _story(1).model_copy(update={
        "planning_diagnostics": {
            "subject_presentation_contract": {
                "schema_version": "subject_presentation_contract_v1",
                "profile_id": "math_formal",
                "primary_mode": "math_formal",
                "required_representation_kinds": ["formula"],
                "optional_representation_kinds": ["diagram", "table"],
                "characteristic_fragment_ids": {
                    "formula": ["root-formula"],
                },
                "chapter_requirements": [{
                    "chapter_id": "chapter-1",
                    "required_representation_kinds": ["formula"],
                    "minimum_artifact_count": 1,
                }],
                "classification_confidence": 1.0,
                "classification_source": "test",
            },
        },
    })

    compact = compact_story_plan_v5(document, story, fragments)
    formula_beats = [
        beat
        for episode in compact.chapters[0].episodes[1:-1]
        for beat in episode.beats
        if "formula" in beat.subject_artifact_kinds
    ]

    assert len(formula_beats) == 1
    assert {
        "root-formula",
        "root-formula-explanation",
    } <= set(formula_beats[0].fragment_ids)


def test_v5_compaction_includes_required_table_from_nested_section() -> None:
    document = CourseDocument(
        course_id="course-v5-nested-table",
        title="Customer operations",
        document_revision="doc-rev-1",
        sections=[
            CourseSection(
                section_id="chapter-1",
                title="Retention planning",
                position=0,
                level=1,
            ),
            CourseSection(
                section_id="section-1",
                parent_section_id="chapter-1",
                title="Cohort review",
                position=1,
                level=2,
            ),
            CourseSection(
                section_id="detail-1",
                parent_section_id="section-1",
                title="Channel comparison",
                position=2,
                level=3,
            ),
        ],
    )
    fragments = [
        ContentFragmentV1(
            fragment_id="section-prose",
            section_id="section-1",
            block_id="section-1-body",
            kind="paragraph",
            text="Compare customer groups before selecting a retention action.",
            ordinal=0,
            source_hash="hash-section-prose",
            role="concept",
            source_kind="course_block",
        ),
        ContentFragmentV1(
            fragment_id="table-heading",
            section_id="detail-1",
            block_id="detail-1-body",
            kind="heading",
            text="Retention by acquisition channel",
            ordinal=1,
            source_hash="hash-table-heading",
            role="case",
            source_kind="course_block",
        ),
        ContentFragmentV1(
            fragment_id="retention-table",
            section_id="detail-1",
            block_id="detail-1-body",
            kind="table",
            text="| Channel | Retained | Churned |\n| --- | ---: | ---: |\n| Referral | 72 | 18 |",
            ordinal=2,
            source_hash="hash-retention-table",
            role="case",
            source_kind="course_block",
        ),
        ContentFragmentV1(
            fragment_id="table-explanation",
            section_id="detail-1",
            block_id="detail-1-body",
            kind="paragraph",
            text="Referral customers show the strongest retention in this cohort.",
            ordinal=3,
            source_hash="hash-table-explanation",
            role="case",
            source_kind="course_block",
        ),
    ]
    story = _story(1).model_copy(update={
        "planning_diagnostics": {
            "subject_presentation_contract": {
                "schema_version": "subject_presentation_contract_v1",
                "profile_id": "business_career",
                "primary_mode": "business_career",
                "required_representation_kinds": ["table"],
                "optional_representation_kinds": ["case", "data"],
                "characteristic_fragment_ids": {
                    "table": ["retention-table"],
                },
                "chapter_requirements": [{
                    "chapter_id": "chapter-1",
                    "required_representation_kinds": ["table"],
                    "minimum_artifact_count": 1,
                }],
                "classification_confidence": 1.0,
                "classification_source": "test",
            },
        },
    })

    compact = compact_story_plan_v5(document, story, fragments)
    table_beats = [
        beat
        for episode in compact.chapters[0].episodes[1:-1]
        for beat in episode.beats
        if "table" in beat.subject_artifact_kinds
    ]

    assert len(table_beats) == 1
    assert {
        "retention-table",
        "table-explanation",
    } <= set(table_beats[0].fragment_ids)


def test_v5_compaction_keeps_a_complete_enumeration_over_optional_background() -> None:
    document = CourseDocument(
        course_id="course-enumeration",
        title="热力学",
        document_revision="doc-enumeration",
        sections=[
            CourseSection(
                section_id="chapter-1",
                title="第一章 热力学基础",
                position=0,
                level=1,
            ),
            CourseSection(
                section_id="section-1",
                parent_section_id="chapter-1",
                title="1.1 热力学系统的分类与描述",
                position=1,
                level=2,
            ),
        ],
    )
    raw = [
        ("title", "heading", "1.1 热力学系统的分类与描述"),
        ("core-heading", "heading", "核心概念与背景"),
        (
            "background",
            "paragraph",
            "在热力学中，系统（System）是指我们研究的物理对象或区域，"
            "而环境（Surroundings）则是系统以外的部分。"
            "系统和环境之间的边界可以是实际存在的（如容器壁），"
            "也可以是想象的（如一个气球内的气体）。"
            "系统与环境之间可能有物质、能量甚至信息的交换。",
        ),
        ("promise", "paragraph", "根据系统与环境之间的交互方式，热力学将系统分为三类："),
        ("isolated", "list_item", "孤立系统：既不交换物质，也不交换能量。"),
        ("closed", "list_item", "封闭系统：不交换物质，但可以交换能量。"),
        ("open", "list_item", "开放系统：既可以交换物质，也可以交换能量。"),
        ("summary", "paragraph", "三种类型为后续建模提供边界条件。"),
    ]
    fragments = [
        ContentFragmentV1(
            fragment_id=fragment_id,
            section_id="section-1",
            block_id="section-1-body",
            kind=kind,  # type: ignore[arg-type]
            text=text,
            ordinal=index,
            source_hash=f"hash-{index}",
            role="concept",
            source_kind="course_block",
        )
        for index, (fragment_id, kind, text) in enumerate(raw)
    ]

    compact = compact_story_plan_v5(document, _story(1), fragments)
    concept_beat = next(
        beat
        for episode in compact.chapters[0].episodes
        if episode.scene_kind == "concept"
        for beat in episode.beats
    )

    assert {"promise", "isolated", "closed", "open"} <= set(
        concept_beat.fragment_ids
    )
    assert "background" not in concept_beat.fragment_ids


def test_v5_compaction_finds_enumeration_members_after_intervening_context() -> None:
    document = CourseDocument(
        course_id="course-deferred-enumeration",
        title="局部解剖学",
        document_revision="doc-rev-1",
        sections=[
            CourseSection(
                section_id="chapter-1",
                title="第一章 胸部",
                position=0,
                level=1,
            ),
            CourseSection(
                section_id="section-1",
                parent_section_id="chapter-1",
                title="1.1 纵隔分区",
                position=1,
                level=2,
            ),
        ],
    )
    raw = [
        ("heading", "heading", "纵隔四分法"),
        (
            "promise",
            "paragraph",
            "纵隔四分法将纵隔划分为四个区域，划分依赖两个关键平面。",
        ),
        ("plane-one", "list_item", "第一平面经过胸骨角与第四胸椎下缘。"),
        ("plane-two", "list_item", "第二平面沿心包前缘界定前方区域。"),
        ("member-label", "paragraph", "**纵隔四分法的具体构成**"),
        ("superior", "list_item", "上纵隔：位于横断面以上。"),
        ("anterior", "list_item", "前纵隔：位于心包与胸骨之间。"),
        ("middle", "list_item", "中纵隔：主要容纳心包和心脏。"),
        ("posterior", "list_item", "后纵隔：位于心包后方。"),
    ]
    fragments = [
        ContentFragmentV1(
            fragment_id=fragment_id,
            section_id="section-1",
            block_id="section-1-body",
            kind=kind,  # type: ignore[arg-type]
            text=text,
            ordinal=index,
            source_hash=f"hash-{index}",
            role="concept",
            source_kind="course_block",
        )
        for index, (fragment_id, kind, text) in enumerate(raw)
    ]

    compact = compact_story_plan_v5(document, _story(1), fragments)
    concept_beat = next(
        beat
        for episode in compact.chapters[0].episodes
        if episode.scene_kind == "concept"
        for beat in episode.beats
    )

    assert {"promise", "superior", "anterior", "middle", "posterior"} <= set(
        concept_beat.fragment_ids
    )


def test_v5_compaction_closes_enumeration_across_adjacent_semantic_blocks() -> None:
    document = CourseDocument(
        course_id="course-cross-block-enumeration",
        title="Regional anatomy",
        document_revision="doc-rev-1",
        sections=[
            CourseSection(
                section_id="chapter-1",
                title="Thorax",
                position=0,
                level=1,
            ),
            CourseSection(
                section_id="section-1",
                parent_section_id="chapter-1",
                title="Mediastinal regions",
                position=1,
                level=2,
            ),
        ],
    )
    raw = [
        ("core-heading", "core-block", "heading", "Mediastinal partition"),
        (
            "promise",
            "core-block",
            "paragraph",
            "The mediastinum is divided into four regions by two planes.",
        ),
        ("plane-one", "core-block", "list_item", "The upper boundary plane."),
        ("plane-two", "core-block", "list_item", "The lower boundary plane."),
        ("regions-heading", "regions-block", "heading", "Four regions"),
        ("superior", "regions-block", "list_item", "Superior region"),
        ("anterior", "regions-block", "list_item", "Anterior region"),
        ("middle", "regions-block", "list_item", "Middle region"),
        ("posterior", "regions-block", "list_item", "Posterior region"),
    ]
    fragments = [
        ContentFragmentV1(
            fragment_id=fragment_id,
            section_id="section-1",
            block_id=block_id,
            kind=kind,  # type: ignore[arg-type]
            text=text,
            ordinal=index,
            source_hash=f"hash-{index}",
            role="concept",
            source_kind="course_block",
        )
        for index, (fragment_id, block_id, kind, text) in enumerate(raw)
    ]

    compact = compact_story_plan_v5(document, _story(1), fragments)
    concept_beat = next(
        beat
        for episode in compact.chapters[0].episodes
        if episode.scene_kind == "concept" and "promise" in episode.beats[0].fragment_ids
        for beat in episode.beats
    )

    assert {"promise", "superior", "anterior", "middle", "posterior"} <= set(
        concept_beat.fragment_ids
    )


def test_v5_compaction_does_not_publish_an_unresolved_enumeration_promise() -> None:
    document = CourseDocument(
        course_id="course-unresolved-enumeration",
        title="Regional anatomy",
        document_revision="doc-rev-1",
        sections=[
            CourseSection(
                section_id="chapter-1",
                title="Abdomen",
                position=0,
                level=1,
            ),
            CourseSection(
                section_id="section-1",
                parent_section_id="chapter-1",
                title="Abdominal regions",
                position=1,
                level=2,
            ),
        ],
    )
    fragments = [
        ContentFragmentV1(
            fragment_id="promise",
            section_id="section-1",
            block_id="regions",
            kind="paragraph",
            text="The abdomen is divided into nine regions.",
            ordinal=0,
            source_hash="promise",
            role="concept",
            source_kind="course_block",
        ),
        ContentFragmentV1(
            fragment_id="plane-definition",
            section_id="section-1",
            block_id="regions",
            kind="paragraph",
            text="The subcostal plane is one of the reference planes.",
            ordinal=1,
            source_hash="definition",
            role="concept",
            source_kind="course_block",
        ),
        ContentFragmentV1(
            fragment_id="plane-landmark",
            section_id="section-1",
            block_id="regions",
            kind="list_item",
            text="It passes through the inferior border of the tenth costal cartilage.",
            ordinal=2,
            source_hash="landmark",
            role="concept",
            source_kind="course_block",
        ),
    ]

    compact = compact_story_plan_v5(document, _story(1), fragments)
    selected_ids = {
        fragment_id
        for episode in compact.chapters[0].episodes
        for beat in episode.beats
        for fragment_id in beat.fragment_ids
    }

    assert "promise" not in selected_ids
    assert {"plane-definition", "plane-landmark"} <= selected_ids


@pytest.mark.parametrize(
    ("renderer_layout", "item_count", "page_capacity"),
    [
        ("question", 8, 5),
        # Concept pages may receive a visual after allocation and then resolve
        # to figure-text, whose hard visible-item capacity is five.
        ("editorial-body", 6, 5),
    ],
)
def test_v5_semantic_core_paginates_to_final_renderer_capacity(
    renderer_layout: str,
    item_count: int,
    page_capacity: int,
) -> None:
    document = CourseDocument(
        course_id="course-semantic-pagination",
        title="课堂核对",
        document_revision="doc-rev-1",
        sections=[
            CourseSection(
                section_id="chapter-1",
                title="第一章",
                position=0,
                level=1,
            ),
            CourseSection(
                section_id="section-1",
                parent_section_id="chapter-1",
                title="1.1 核对练习",
                position=1,
                level=2,
            ),
        ],
    )
    fragments = [
        ContentFragmentV1(
            fragment_id=f"question-{index}",
            section_id="section-1",
            block_id="practice-block",
            kind="list_item",
            text=f"核对问题 {index}",
            ordinal=index,
            source_hash=f"hash-{index}",
            role="checkpoint",
            source_kind="course_block",
        )
        for index in range(1, item_count + 1)
    ]
    story = _story(1)
    concept_beat = _beat(1, "concept").model_copy(update={
        "beat_id": "beat-semantic-question",
        "fragment_ids": [item.fragment_id for item in fragments],
        "renderer_layout": renderer_layout,
        "layout_selection_reason": "v5_semantic_grouping",
    })
    chapter = story.chapters[0]
    story = story.model_copy(update={
        "chapters": [chapter.model_copy(update={
            "episodes": [
                chapter.episodes[0],
                TeachingEpisodeV2(
                    episode_id="episode-semantic-question",
                    scene_kind="concept",
                    teaching_job="完成问题核对",
                    beats=[concept_beat],
                ),
                chapter.episodes[-1],
            ],
        })],
    })

    allocation, _ = allocation_from_story_plan_v2(document, fragments, story)
    teaching_pages = [page for page in allocation.pages if page.fragment_ids]

    assert len(teaching_pages) == 2
    assert all(len(page.fragment_ids) <= page_capacity for page in teaching_pages)
    assert [
        fragment_id
        for page in teaching_pages
        for fragment_id in page.fragment_ids
    ] == [item.fragment_id for item in fragments]


def test_one_text_group_cannot_keep_a_two_column_layout() -> None:
    contract = resolve_page_contract_v5({
        "layout": "concept",
        "composition": "split-visual",
        "visuals": [],
        "blocks": [
            {
                "block_id": "body",
                "type": "rich_text",
                "content": "只有一个完整的语义段落。",
                "items": [],
            },
        ],
        "quality": {"requested_layout": "two-column"},
    })

    assert contract.resolved_layout == "editorial-body"
    assert contract.resolved_composition == "statement"
    assert contract.occupied_major_region_count == 1
    assert contract.layout_fallback_reason == "single_group_two_column"


def test_diagram_full_with_source_text_resolves_to_figure_text() -> None:
    contract = resolve_page_contract_v5({
        "layout": "concept",
        "composition": "diagram-full",
        "visuals": [{"visual_id": "anatomy-diagram"}],
        "blocks": [{
            "block_id": "explanation",
            "type": "rich_text",
            "content": "图解旁必须保留这段来源解释。",
            "items": [],
        }],
        "quality": {"requested_layout": "diagram-full"},
    })

    assert contract.resolved_layout == "figure-text"
    assert contract.resolved_composition == "split-visual"
    assert contract.layout_fallback_reason == "diagram_full_with_source_text"


def test_editorial_body_with_an_effective_visual_resolves_to_figure_text() -> None:
    contract = resolve_page_contract_v5({
        "layout": "concept",
        "blocks": [{
            "block_id": "explanation",
            "type": "rich_text",
            "content": "来源解释与图示共同构成这一页的完整论证。",
            "items": [],
        }],
        "visuals": [{"kind": "relational_diagram", "visual_id": "diagram-1"}],
        "quality": {"requested_layout": "editorial-body"},
    })

    assert contract.resolved_layout == "figure-text"


def test_one_prompt_block_cannot_fabricate_a_practice_feedback_region() -> None:
    contract = resolve_page_contract_v5({
        "layout": "practice",
        "blocks": [{
            "block_id": "prompt-only",
            "type": "exercise",
            "items": ["判断系统类型", "说明判断依据"],
        }],
        "quality": {"requested_layout": "practice-feedback"},
    })

    assert contract.resolved_layout == "editorial-body"
    assert contract.layout_fallback_reason == "practice_without_feedback"
    assert contract.major_region_count == 1


def test_three_sibling_items_select_a_classification_layout() -> None:
    contract = resolve_page_contract_v5({
        "layout": "concept",
        "composition": "statement",
        "visuals": [],
        "blocks": [
            {
                "block_id": "classification",
                "type": "bullets",
                "content": "",
                "items": ["孤立系统", "封闭系统", "开放系统"],
            },
        ],
        "quality": {"requested_layout": "two-column"},
    })

    assert contract.resolved_layout == "classification-3"
    assert [binding.semantic_role for binding in contract.slot_bindings] == [
        "classification_item",
        "classification_item",
        "classification_item",
    ]
    assert contract.occupied_major_region_count == 3


def test_two_regions_requested_as_classification_use_two_column_layout() -> None:
    contract = resolve_page_contract_v5({
        "layout": "concept",
        "composition": "statement",
        "visuals": [],
        "blocks": [
            {"block_id": "left", "type": "statement", "content": "输入条件"},
            {"block_id": "right", "type": "statement", "content": "输出结果"},
        ],
        "quality": {"requested_layout": "classification-3"},
    })

    assert contract.resolved_layout == "balanced-two-column"
    assert contract.major_region_count == 2


def test_four_sibling_items_select_a_two_by_two_parallel_layout() -> None:
    contract = resolve_page_contract_v5({
        "layout": "concept",
        "composition": "statement",
        "visuals": [],
        "blocks": [{
            "block_id": "four-points",
            "type": "bullets",
            "content": "",
            "items": ["输入", "处理", "输出", "验证"],
        }],
        "quality": {"requested_layout": "classification-3"},
    })

    assert contract.resolved_layout == "parallel-examples"
    assert contract.occupied_major_region_count == 4


@pytest.mark.parametrize(
    ("requested_layout", "expected_regions"),
    [
        ("worked-example", 3),
        ("practice-feedback", 2),
    ],
)
def test_instructional_layouts_survive_final_contract_resolution(
    requested_layout: str,
    expected_regions: int,
) -> None:
    blocks = (
        [
            {
                "block_id": "practice-prompt",
                "type": "exercise",
                "content": "先判断系统类型。",
            },
            {
                "block_id": "practice-feedback",
                "type": "bullets",
                "items": ["边界条件一致", "排除其他类型"],
            },
        ]
        if requested_layout == "practice-feedback"
        else [
            {
                "block_id": "instructional-sequence",
                "type": "process",
                "items": ["识别条件", "选择方法", "检查结论"],
            },
        ]
    )
    contract = resolve_page_contract_v5({
        "layout": "concept",
        "blocks": blocks,
        "quality": {"requested_layout": requested_layout},
    })

    assert contract.resolved_layout == requested_layout
    assert contract.major_region_count == expected_regions


def test_rejected_visual_reflows_to_a_text_native_composition() -> None:
    contract = resolve_page_contract_v5({
        "layout": "concept",
        "composition": "split-visual",
        "visuals": [],
        "blocks": [
            {
                "block_id": "definition",
                "type": "rich_text",
                "content": "系统边界决定可发生的交换。",
                "items": [],
            },
        ],
        "quality": {"requested_layout": "figure-text"},
    })

    assert contract.visual_decision == "none"
    assert contract.resolved_layout == "editorial-body"
    assert contract.resolved_composition == "statement"
    assert contract.layout_fallback_reason == "visual_layout_without_visual"


def test_shared_slide_model_accepts_v5_outline_contract() -> None:
    deck = SlideDeckContent.model_validate({
        "schema_version": "slide_deck_v5",
        "title": "V5",
        "slides": [],
        "deck_outline": {
            "schema_version": "deck_outline_v5",
            "outline_id": "outline-1",
        },
    })

    assert deck.schema_version == "slide_deck_v5"
    assert deck.deck_outline["schema_version"] == "deck_outline_v5"


def test_shared_v5_layout_catalog_matches_pptx_renderer_contract() -> None:
    catalog = json.loads(
        (
            Path(__file__).resolve().parents[2]
            / "shared"
            / "slide-layout-contract-v5.json"
        ).read_text(encoding="utf-8")
    )
    expected = {
        item["layout"]: item["pptx_renderer"]
        for item in catalog["layouts"]
    }

    assert expected == V5_LAYOUT_RENDERER_NAMES
    assert catalog["minimum_title_font_pt"] >= 35
    assert catalog["minimum_body_font_pt"] >= 16
    assert all(
        callable(getattr(slide_deck_renderer, renderer_name))
        for renderer_name in expected.values()
    )


def test_v5_render_review_fixture_exports_all_semantic_compositions(
    tmp_path,
) -> None:
    fixture = json.loads(
        (
            Path(__file__).resolve().parent
            / "fixtures"
            / "slide_deck_v5_render_review.json"
        ).read_text(encoding="utf-8")
    )
    output = export_structured_slide_deck(
        fixture,
        tmp_path / "slide-deck-v5-review.pptx",
        require_quality=False,
    )
    presentation = Presentation(output)
    visible_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )

    assert len(presentation.slides) == 5
    visible_lines = {line.strip() for line in visible_text.splitlines()}
    assert {"步骤 1", "步骤 2", "步骤 3", "回答与判断依据", "课程主线"} <= visible_lines
    assert {"已知", "推理", "结论"}.isdisjoint(visible_lines)


def test_pptx_export_uses_v5_quality_contract_for_v5_content(tmp_path) -> None:
    content = {
        "schema_version": "slide_deck_v5",
        "title": "V5 quality routing",
        "slides": [],
    }

    with (
        patch(
            "slide_deck_v5.validate_slide_deck_v5",
            return_value={"passed": True, "blockers": []},
        ) as validate_v5,
        patch("slide_deck_renderer.validate_slide_deck") as validate_legacy,
    ):
        export_structured_slide_deck(
            content,
            tmp_path / "v5-quality-routing.pptx",
            require_quality=True,
            course_data={},
        )

    validate_v5.assert_called_once()
    validate_legacy.assert_not_called()


def test_v5_navigation_slide_can_bind_a_chapter_without_inventing_body_sources() -> None:
    report = validate_slide_deck({
        "schema_version": "slide_deck_v5",
        "title": "V5",
        "slides": [
            {
                "unit_id": "cover",
                "position": 0,
                "layout": "cover",
                "slide_purpose": "orientation",
                "title": "V5",
                "blocks": [],
            },
            {
                "unit_id": "chapter-entry",
                "position": 1,
                "layout": "chapter",
                "slide_purpose": "chapter_open",
                "title": "第一章",
                "section_id": "chapter-1",
                "source_section_ids": ["chapter-1"],
                "source_block_ids": [],
                "blocks": [],
                "quality": {"navigation_only": True},
            },
            {
                "unit_id": "recap",
                "position": 2,
                "layout": "recap",
                "slide_purpose": "course_recap",
                "title": "课程总结",
                "blocks": [{
                    "block_id": "summary",
                    "type": "bullets",
                    "items": ["关键结论"],
                }],
            },
        ],
    })

    assert "slide_source_missing" not in {
        issue["code"] for issue in report["blockers"]
    }


def test_pptx_renderer_uses_resolved_layout_instead_of_requested_layout() -> None:
    unit = SimpleNamespace(
        visuals=[],
        layout="concept",
        quality={
            "requested_layout": "two-column",
            "resolved_layout": "editorial-body",
        },
    )
    editorial_renderer = Mock()
    two_column_renderer = Mock()

    with (
        patch("slide_deck_renderer._fill_background"),
        patch("slide_deck_renderer._footer"),
        patch("slide_deck_renderer._render_editorial_body", editorial_renderer),
        patch("slide_deck_renderer._render_two_column", two_column_renderer),
    ):
        _render_slide(
            Mock(),
            unit,
            1,
            1,
            {"surface": "FFFFFF"},
            Mock(),
        )

    editorial_renderer.assert_called_once()
    two_column_renderer.assert_not_called()


def test_pptx_renderer_rejects_a_final_v5_page_without_resolved_layout(
    tmp_path: Path,
) -> None:
    content = compile_slide_deck_v5(
        _document(1),
        {},
        story_plan=_story(1),
    )
    content["slides"][0]["quality"].pop("resolved_layout", None)

    with pytest.raises(ValueError, match="v5_final_layout_missing"):
        export_structured_slide_deck(
            content,
            tmp_path / "invalid-v5-layout.pptx",
            require_quality=False,
        )


@pytest.mark.parametrize(
    ("resolved_layout", "renderer_name"),
    [
        ("worked-example", "_render_worked_example"),
        ("practice-feedback", "_render_practice_feedback"),
        ("chapter-recap", "_render_chapter_recap"),
        ("course-synthesis", "_render_course_synthesis"),
    ],
)
def test_v5_semantic_layouts_use_dedicated_pptx_renderers(
    resolved_layout: str,
    renderer_name: str,
) -> None:
    unit = SimpleNamespace(
        visuals=[],
        layout="concept",
        quality={"resolved_layout": resolved_layout},
    )
    renderer = Mock()

    with (
        patch("slide_deck_renderer._fill_background"),
        patch("slide_deck_renderer._footer"),
        patch(f"slide_deck_renderer.{renderer_name}", renderer),
    ):
        _render_slide(
            Mock(),
            unit,
            1,
            1,
            {"surface": "FFFFFF"},
            Mock(),
        )

    renderer.assert_called_once()


def test_v5_editorial_body_does_not_reserve_a_fake_right_sidebar() -> None:
    unit = SimpleNamespace(
        blocks=[
            SimpleNamespace(
                items=[],
                content="系统边界决定系统与环境之间可以发生的交换。",
                title="",
            ),
        ],
        key_message="",
        eyebrow="核心概念",
        title="系统边界决定交换方式",
    )
    shape = Mock()
    text = Mock()

    with (
        patch("slide_deck_renderer._heading"),
        patch("slide_deck_renderer._shape", shape),
        patch("slide_deck_renderer._text", text),
    ):
        _render_editorial_body(
            Mock(),
            unit,
            {
                "accent": "00F",
                "accent_soft": "DDF",
                "canvas": "FFF",
                "chart_bg": "DDD",
                "ink": "000",
            },
        )

    assert all(float(call.args[1]) < 9.5 for call in shape.call_args_list)
    body_call = next(
        call for call in text.call_args_list
        if call.args[1].startswith("系统边界")
    )
    assert float(body_call.args[4]) >= 10
    assert int(body_call.args[6]) >= 24


def test_title_compiler_rejects_raw_diagram_identifiers() -> None:
    title = compile_page_title_v5(
        explicit_title='> ID: "ThermodynamicSystemClassification"',
        primary_claim="热力学系统按交换方式分为三类",
        body_text="孤立、封闭和开放系统的边界条件不同。",
    )

    assert title == "热力学系统按交换方式分为三类"
    assert "ID:" not in title


def test_title_compiler_replaces_numbered_section_heading_with_visible_claim() -> None:
    title = compile_page_title_v5(
        explicit_title="1.1 热力学系统的分类与描述",
        primary_claim="1.1 热力学系统的分类与描述",
        body_text=(
            "核心概念与背景\n"
            "在热力学中，系统是我们研究的对象，环境是系统以外的部分。\n"
            "根据系统与环境之间的交互方式，热力学将系统分为三类：\n"
            "孤立系统\n封闭系统\n开放系统"
        ),
    )

    assert title == "热力学系统的三种类型"


def test_title_compiler_prefers_supported_claim_over_source_topic_heading() -> None:
    title = compile_page_title_v5(
        explicit_title="内能的本质",
        primary_claim="内能的本质",
        body_text="内能是系统内所有微观粒子能量的总和。",
        prefer_body_claim=True,
    )

    assert title == "内能是系统内所有微观粒子能量的总和"


def test_title_compiler_prefers_the_lead_definition_over_later_detail() -> None:
    title = compile_page_title_v5(
        explicit_title="内能的本质",
        primary_claim="内能的本质",
        body_text=(
            "内能是系统内所有微观粒子能量的总和。"
            "它由分子平动、转动、振动以及相互作用势能共同构成。"
        ),
        prefer_body_claim=True,
    )

    assert title == "内能是系统内所有微观粒子能量的总和"


def test_title_compiler_strips_template_label_before_selecting_body_claim() -> None:
    title = compile_page_title_v5(
        explicit_title="内能的本质",
        primary_claim="内能的本质",
        body_text=(
            "💡 核心概念与背景 "
            "内能是系统内所有微观粒子能量的总和。"
            "它由分子平动、转动和振动能共同构成。"
        ),
        prefer_body_claim=True,
    )

    assert title == "内能是系统内所有微观粒子能量的总和"


def test_title_compiler_keeps_an_existing_takeaway_title() -> None:
    title = compile_page_title_v5(
        explicit_title="系统边界决定可发生的交换",
        primary_claim="系统边界决定可发生的交换",
        body_text="系统和环境之间存在边界。",
        prefer_body_claim=True,
    )

    assert title == "系统边界决定可发生的交换"


@pytest.mark.parametrize(
    ("template_title", "body"),
    [
        (
            "🏭 实战案例/行业应用",
            "在航天器设计中，速度分布用于检验喷嘴方案。",
        ),
        (
            "✅ 思考与挑战",
            "为什么封闭系统仍然可以和环境交换能量？",
        ),
    ],
)
def test_title_compiler_demotes_template_labels_to_eyebrow(
    template_title: str,
    body: str,
) -> None:
    title = compile_page_title_v5(
        explicit_title=template_title,
        primary_claim=template_title,
        body_text=body,
    )

    assert title not in {template_title, "实战案例/行业应用", "思考与挑战"}
    assert title


def test_v5_quality_gate_rejects_orphan_formula_and_title_duplication() -> None:
    issues = v5_contract_issues([
        {
            "unit_id": "formula-only",
            "title": "内能变化",
            "visuals": [{"kind": "formula", "latex": r"\Delta U=U_2-U_1"}],
            "blocks": [],
            "quality": {
                "requested_layout": "formula-explanation",
                "resolved_layout": "figure-text",
                "resolved_composition": "split-visual",
                "major_region_count": 2,
                "occupied_major_region_count": 1,
            },
        },
        {
            "unit_id": "duplicate-title",
            "title": "系统边界决定可发生的交换",
            "visuals": [],
            "blocks": [{
                "block_id": "body",
                "type": "rich_text",
                "content": "系统边界决定可发生的交换。",
                "items": [],
            }],
            "quality": {
                "requested_layout": "editorial-body",
                "resolved_layout": "editorial-body",
                "resolved_composition": "statement",
                "major_region_count": 1,
                "occupied_major_region_count": 1,
            },
        },
    ])

    assert {issue["code"] for issue in issues} >= {
        "orphan_formula",
        "title_body_duplication",
    }


def test_v5_enumeration_gate_counts_visible_bullet_lines_and_ignores_singular_phrases() -> None:
    slides = [
        {
            "unit_id": "classification",
            "title": "热力学将系统分为三类",
            "blocks": [{
                "type": "statement",
                "content": (
                    "热力学将系统分为三类：\n"
                    "• 孤立系统\n"
                    "• 封闭系统\n"
                    "• 开放系统"
                ),
                "items": [],
            }],
            "visuals": [],
            "quality": {
                "resolved_layout": "editorial-body",
                "occupied_major_region_count": 1,
                "major_region_count": 1,
            },
        },
        {
            "unit_id": "singular-example",
            "title": "一个系统的宏观状态",
            "blocks": [{
                "type": "statement",
                "content": "这里解释一个系统如何由状态变量描述。",
                "items": [],
            }],
            "visuals": [],
            "quality": {
                "resolved_layout": "editorial-body",
                "occupied_major_region_count": 1,
                "major_region_count": 1,
            },
        },
    ]

    issues = v5_contract_issues(slides)

    assert not any(
        issue["code"] == "enumeration_cardinality_mismatch"
        for issue in issues
    )


def test_v5_quality_gate_rejects_incomplete_enumeration_and_section_title() -> None:
    issues = v5_contract_issues([
        {
            "unit_id": "incomplete-classification",
            "title": "1.1 热力学系统的分类与描述",
            "visuals": [],
            "blocks": [
                {
                    "block_id": "promise",
                    "type": "statement",
                    "content": "根据交换方式，热力学系统分为三类：",
                    "items": [],
                },
                {
                    "block_id": "classification",
                    "type": "bullets",
                    "content": "",
                    "items": ["孤立系统"],
                },
            ],
            "quality": {
                "requested_layout": "editorial-body",
                "resolved_layout": "editorial-body",
                "resolved_composition": "statement",
                "major_region_count": 1,
                "occupied_major_region_count": 1,
            },
        },
    ])

    assert {issue["code"] for issue in issues} >= {
        "enumeration_cardinality_mismatch",
        "source_section_heading_as_title",
    }


def test_v5_quality_gate_does_not_treat_quantities_as_visible_list_contracts() -> None:
    slides = [
        {
            "unit_id": "two-vector-spaces",
            "title": "我们从两个向量空间 V 和 W 开始",
            "blocks": [{
                "type": "process",
                "content": "",
                "items": ["可加性"],
            }],
        },
        {
            "unit_id": "three-input-vectors",
            "title": "假设我们有三个线性无关的向量",
            "blocks": [{
                "type": "process",
                "content": "",
                "items": ["初始化", "第二步"],
            }],
        },
        {
            "unit_id": "inline-two-classes",
            "title": "为什么需要 QR 分解",
            "blocks": [{
                "type": "statement",
                "content": (
                    "把食材分成两类：一类是整齐摆放的原料（Q），"
                    "另一类是精确用量的调味料（R）。"
                ),
                "items": [],
            }],
        },
        {
            "unit_id": "pixel-count",
            "title": "应用场景：人脸识别",
            "blocks": [{
                "type": "statement",
                "content": "一张 64x64 像素的灰度图像共有 4096 个像素点。",
                "items": [],
            }],
        },
    ]
    for slide in slides:
        slide["visuals"] = []
        slide["quality"] = {
            "resolved_layout": "editorial-body",
            "major_region_count": 1,
            "occupied_major_region_count": 1,
        }

    issues = v5_contract_issues(slides)

    assert not any(
        issue["code"] == "enumeration_cardinality_mismatch"
        for issue in issues
    )


def test_v5_quality_gate_keeps_explicit_title_enumerations_strict() -> None:
    issues = v5_contract_issues([{
        "unit_id": "incomplete-three-types",
        "title": "热力学系统的三类交换方式",
        "blocks": [{
            "type": "bullets",
            "content": "",
            "items": ["孤立系统", "封闭系统"],
        }],
        "visuals": [],
        "quality": {
            "resolved_layout": "editorial-body",
            "major_region_count": 1,
            "occupied_major_region_count": 1,
        },
    }])

    assert any(
        issue["code"] == "enumeration_cardinality_mismatch"
        and issue["expected_count"] == 3
        and issue["visible_item_count"] == 2
        for issue in issues
    )


def test_v5_compiler_removes_repeated_lead_claim_but_keeps_supporting_items() -> None:
    slide = apply_page_contract_v5({
        "unit_id": "deduplicated-claim",
        "layout": "concept",
        "title": "系统按交换方式分为三类",
        "blocks": [
            {
                "block_id": "lead",
                "type": "rich_text",
                "content": "系统按交换方式分为三类。",
                "items": [],
            },
            {
                "block_id": "classification",
                "type": "bullets",
                "items": ["孤立系统", "封闭系统", "开放系统"],
            },
        ],
        "quality": {"requested_layout": "editorial-body"},
    })

    assert [block["block_id"] for block in slide["blocks"]] == [
        "classification"
    ]
    assert slide["quality"]["resolved_layout"] == "classification-3"
    assert "title_body_duplication" not in {
        issue["code"] for issue in v5_contract_issues([slide])
    }


def test_v5_compiler_removes_only_the_repeated_question_clause() -> None:
    slide = apply_page_contract_v5({
        "unit_id": "question-clause",
        "layout": "practice",
        "title": "✅ 思考与挑战",
        "blocks": [{
            "block_id": "questions",
            "type": "bullets",
            "items": [
                "为什么沙漠地区昼夜温差大？这是否与沙子的比热容有关？",
                "水和油在相同条件下谁需要更多热量？",
            ],
        }],
        "quality": {"requested_layout": "editorial-body"},
    })

    assert slide["title"] == "为什么沙漠地区昼夜温差大"
    assert slide["blocks"][0]["items"] == [
        "这是否与沙子的比热容有关？",
        "水和油在相同条件下谁需要更多热量？",
    ]
    assert not v5_contract_issues([slide])


def test_v5_density_contract_rejects_overflow_without_reducing_font_floor() -> None:
    slide = apply_page_contract_v5({
        "unit_id": "dense-page",
        "layout": "concept",
        "title": "热力学系统的分类",
        "key_message": "",
        "blocks": [{
            "block_id": "dense-body",
            "type": "rich_text",
            "content": "正文" * 220,
            "items": [],
        }],
        "quality": {"requested_layout": "editorial-body"},
    })

    assert slide["quality"]["density_band"] == "overflow"
    assert slide["quality"]["minimum_body_font_pt"] >= 16
    assert slide["quality"]["minimum_title_font_pt"] >= 35
    assert {
        issue["code"] for issue in v5_contract_issues([slide])
    } >= {"body_density_overflow"}


def test_worked_example_with_more_than_three_regions_reflows_without_dropping_items() -> None:
    slide = apply_page_contract_v5({
        "unit_id": "application-list",
        "layout": "case-study",
        "title": "分子速度分布支持多类工程判断",
        "blocks": [
            {
                "block_id": "context",
                "type": "statement",
                "content": "同一统计模型可以服务不同工程场景。",
            },
            {
                "block_id": "applications",
                "type": "bullets",
                "items": ["航天推进", "真空系统", "材料扩散", "稀薄气体流动"],
            },
        ],
        "quality": {"requested_layout": "worked-example"},
    })

    assert slide["quality"]["resolved_layout"] == "editorial-body"
    assert slide["quality"]["layout_fallback_reason"] == (
        "worked_example_item_overflow"
    )
    assert slide["blocks"][1]["items"] == [
        "航天推进",
        "真空系统",
        "材料扩散",
        "稀薄气体流动",
    ]
    assert not v5_contract_issues([slide])


def test_v5_quality_cannot_publish_when_a_retained_nested_gate_is_critical() -> None:
    report = finalize_v5_quality_report(
        previous_quality={
            "passed": True,
            "score": 92,
            "semantic": {
                "passed": False,
                "issues": [{
                    "severity": "critical",
                    "code": "official_source_revision_mismatch",
                    "target": "deck",
                }],
            },
            "visual": {"passed": True, "issues": []},
            "blockers": [],
        },
        slides=[],
        planner="ai",
        fallback_reason="",
    )

    assert report["passed"] is False
    assert report["semantic"]["passed"] is False
    assert report["visual"]["passed"] is True
    assert {
        issue["code"] for issue in report["blockers"]
    } == {"official_source_revision_mismatch"}


def test_v5_quality_retains_non_superseded_presentation_blockers() -> None:
    report = finalize_v5_quality_report(
        previous_quality={
            "passed": True,
            "score": 95,
            "slide_count": 1,
            "presentation": {
                "passed": False,
                "score": 80,
                "issues": [{
                    "severity": "critical",
                    "code": "rendered_text_clipped",
                    "target": "slide:v4:0001",
                }],
                "blockers": [{
                    "severity": "critical",
                    "code": "rendered_text_clipped",
                    "target": "slide:v4:0001",
                }],
            },
            "blockers": [],
        },
        slides=[],
        planner="ai",
        fallback_reason="",
    )

    assert report["passed"] is False
    assert report["presentation"]["passed"] is False
    assert report["visual"]["passed"] is True
    assert {
        issue["code"] for issue in report["blockers"]
    } == {"rendered_text_clipped"}


def test_v5_quality_recomputes_stale_presentation_and_final_slide_count() -> None:
    slides = [
        {
            "unit_id": "slide:v5:0001",
            "title": "结构关系决定判断顺序",
            "blocks": [{
                "type": "rich_text",
                "content": "先定位目标结构，再辨认相邻关系，最后依据边界完成判断。",
            }],
            "quality": {"passed": True},
        },
        {
            "unit_id": "slide:v5:0002",
            "title": "安全边界约束操作路径",
            "blocks": [{
                "type": "rich_text",
                "content": "先识别安全边界，再选择操作路径，并用最终结果检查判断。",
            }],
            "quality": {"passed": True},
        },
    ]
    report = finalize_v5_quality_report(
        previous_quality={
            "passed": False,
            "score": 80,
            "slide_count": 1,
            "presentation": {
                "passed": False,
                "score": 80,
                "issues": [{
                    "severity": "critical",
                    "code": "layout_family_repeated_more_than_twice",
                    "target": "deck",
                }],
                "blockers": [{
                    "severity": "critical",
                    "code": "layout_family_repeated_more_than_twice",
                    "target": "deck",
                }],
            },
            "blockers": [],
        },
        slides=slides,
        planner="ai",
        fallback_reason="",
    )

    assert report["passed"] is True
    assert report["slide_count"] == len(slides)
    assert report["presentation"]["passed"] is True
    assert report["presentation"]["issues"] == []
    assert report["presentation"]["blockers"] == []


def test_v5_quality_cannot_publish_when_any_final_slide_is_critical() -> None:
    report = finalize_v5_quality_report(
        previous_quality={
            "passed": True,
            "score": 100,
            "semantic": {"passed": True, "issues": []},
            "visual": {"passed": True, "issues": []},
            "blockers": [],
        },
        slides=[{
            "unit_id": "slide:v4:0001",
            "quality": {
                "passed": False,
                "issues": [{
                    "severity": "critical",
                    "code": "slide_block_overflow",
                    "slide_id": "slide:v4:0001",
                    "message": "The final slide still exceeds its resolved layout.",
                }],
            },
        }],
        planner="ai",
        fallback_reason="",
    )

    assert report["passed"] is False
    assert "legacy_slide_id" in {
        issue["code"] for issue in report["blockers"]
    }


def test_v5_publishable_ai_fallback_is_a_warning_not_a_blocker() -> None:
    report = finalize_v5_quality_report(
        previous_quality={
            "passed": True,
            "score": 100,
            "semantic": {"passed": True, "issues": []},
            "visual": {"passed": True, "issues": []},
            "blockers": [],
        },
        slides=[],
        planner="deterministic_fallback",
        fallback_reason="invalid_or_failed_ai_story_plan",
    )

    assert report["passed"] is True
    assert report["blockers"] == []
    assert {
        issue["code"] for issue in report["warnings"]
    } == {"ai_story_planner_fallback"}


def test_v5_reports_partial_visual_ai_fallback_as_degraded() -> None:
    report = finalize_v5_quality_report(
        previous_quality={"passed": True, "score": 100, "blockers": []},
        slides=[],
        planner="ai",
        fallback_reason="",
        visual_planning={
            "planner": "ai",
            "fallback_reason": "partial_ai_visual_plan",
            "ai_visual_pages_accepted": 65,
            "ai_visual_pages_fallback": 24,
        },
    )

    assert report["passed"] is True
    assert report["visual_planning"]["degraded"] is True
    assert {
        issue["code"] for issue in report["warnings"]
    } == {"ai_visual_planner_partial_fallback"}


def test_v5_contract_discards_superseded_v4_capacity_blockers() -> None:
    slide = apply_page_contract_v5({
        "unit_id": "slide:v5:long-course",
        "layout": "concept",
        "composition": "statement",
        "title": "状态变量只取决于系统当前状态",
        "blocks": [{
            "block_id": "definition",
            "type": "rich_text",
            "content": "状态变量与过程路径无关，因此判断时只比较系统的初态与终态。",
            "items": [],
        }],
        "visuals": [],
        "quality": {
            "passed": False,
            "issues": [{
                "severity": "critical",
                "code": "concept_card_overflow",
                "slide_id": "slide:v5:long-course",
            }],
            "blockers": [{
                "severity": "critical",
                "code": "slide_block_overflow",
                "slide_id": "slide:v5:long-course",
            }],
        },
    })

    assert slide["quality"]["passed"] is True
    assert slide["quality"]["issues"] == []
    assert slide["quality"]["blockers"] == []
    report = finalize_v5_quality_report(
        previous_quality={
            "passed": False,
            "score": 0,
            "semantic": {"passed": True, "issues": []},
            "visual": {"passed": True, "issues": []},
            "blockers": [{
                "severity": "critical",
                "code": "slide_block_overflow",
                "slide_id": "slide:v5:long-course",
            }],
        },
        slides=[slide],
        planner="deterministic_fallback",
        fallback_reason="invalid_or_failed_ai_story_plan",
    )
    assert report["passed"] is True
    assert report["blockers"] == []


def test_v5_slide_counts_include_inserted_navigation_and_chapter_pages() -> None:
    summary = summarize_v5_slide_counts([
        {"unit_id": "slide:title", "quality": {}},
        {"unit_id": "slide:roadmap", "quality": {}},
        {"unit_id": "slide:v5:chapter:1", "quality": {}},
        {"unit_id": "slide:v4:0001", "quality": {}},
        {"unit_id": "slide:v4:leftover:0001", "quality": {"appendix": True}},
    ])

    assert summary == {
        "main_slide_count": 4,
        "appendix_slide_count": 1,
        "total_slide_count": 5,
    }


def test_v5_title_and_item_budgets_are_hard_quality_gates() -> None:
    issues = v5_contract_issues([{
        "unit_id": "over-budget",
        "layout": "recap",
        "title": "这是一个明显超出投影扫读容量且没有进行结构化压缩的页面标题" * 2,
        "blocks": [{
            "block_id": "too-many",
            "type": "bullets",
            "items": [f"结论 {index}" for index in range(8)],
        }],
        "quality": {
            "requested_layout": "chapter-recap",
            "resolved_layout": "chapter-recap",
            "resolved_composition": "statement",
            "major_region_count": 1,
            "occupied_major_region_count": 1,
        },
    }])

    assert {issue["code"] for issue in issues} >= {
        "slide_title_overflow",
        "visible_item_overflow",
    }


def test_v5_comparison_density_counts_rendered_rows_not_flattened_cells() -> None:
    issues = v5_contract_issues([{
        "unit_id": "comparison-table",
        "layout": "concept",
        "title": "浅筋膜与深筋膜",
        "blocks": [{
            "block_id": "comparison",
            "type": "comparison",
            "items": [
                "疏松结缔组织",
                "致密不规则结缔组织",
                "脂肪细胞",
                "胶原纤维",
                "易分离",
                "难分离",
                "出血少",
                "张力大",
            ],
            "metadata": {
                "headers": ["维度", "浅筋膜", "深筋膜"],
                "rows": [
                    ["组织学", "疏松", "致密"],
                    ["成分", "脂肪", "胶原"],
                    ["分离", "容易", "困难"],
                    ["临床", "出血少", "张力大"],
                ],
            },
        }],
        "visuals": [],
        "quality": {
            "resolved_layout": "figure-text",
            "major_region_count": 1,
            "occupied_major_region_count": 1,
        },
    }])

    assert not any(
        issue["code"] == "visible_item_overflow"
        for issue in issues
    )


def test_v5_paired_practice_feedback_counts_aligned_rows_once() -> None:
    slide = apply_page_contract_v5({
        "unit_id": "practice-three-pairs",
        "layout": "practice",
        "title": "Check the three clinical decisions",
        "blocks": [
            {
                "block_id": "prompts",
                "type": "question",
                "items": ["Question A", "Question B", "Question C"],
                "metadata": {
                    "semantic_role": "prompt",
                    "question_ids": ["q-a", "q-b", "q-c"],
                },
            },
            {
                "block_id": "answers",
                "type": "callout",
                "items": ["Answer A", "Answer B", "Answer C"],
                "metadata": {
                    "semantic_role": "answer",
                    "answer_for_question_ids": ["q-a", "q-b", "q-c"],
                },
            },
        ],
        "visuals": [],
        "quality": {
            "requested_layout": "practice-feedback",
            "feedback_mode": "paired",
            "feedback_pair_count": 3,
        },
    })

    assert slide["quality"]["resolved_layout"] == "practice-feedback"
    assert slide["quality"]["visible_item_count"] == 3
    assert not any(
        issue["code"] == "visible_item_overflow"
        for issue in v5_contract_issues([slide])
    )


def test_v5_paired_practice_discards_non_rendered_source_checklist() -> None:
    [enriched] = _enrich_practice_feedback_slides_v5([{
        "unit_id": "practice-production-shape",
        "layout": "question",
        "scene_kind": "practice_feedback",
        "beat_role": "prompt",
        "title": "Check the three anatomical decisions",
        "blocks": [
            {
                "block_id": "prompts",
                "type": "question",
                "items": ["Question A", "Question B", "Question C"],
                "metadata": {"semantic_role": "prompt"},
            },
            {
                "block_id": "source-checklist",
                "type": "bullets",
                "items": ["Check A", "Check B", "Check C"],
                "metadata": {"semantic_role": "support"},
            },
        ],
        "visuals": [],
        "quality": {
            "requested_layout": "question-prompt",
            "question_ids": ["q-a", "q-b", "q-c"],
            "generated_practice_answers": [
                {"question_id": "q-a", "answer_text": "Answer A"},
                {"question_id": "q-b", "answer_text": "Answer B"},
                {"question_id": "q-c", "answer_text": "Answer C"},
            ],
        },
    }])

    slide = apply_page_contract_v5(enriched)

    assert [
        block["metadata"]["semantic_role"]
        for block in slide["blocks"]
    ] == ["prompt", "answer"]
    assert slide["quality"]["resolved_layout"] == "practice-feedback"
    assert slide["quality"]["visible_item_count"] == 3
    assert not any(
        issue["code"] == "visible_item_overflow"
        for issue in v5_contract_issues([slide])
    )


def test_v5_shared_feedback_matches_the_six_item_renderer_contract() -> None:
    enriched = _enrich_practice_feedback_slides_v5([
        {
            "unit_id": "concept-source",
            "chapter_id": "chapter-1",
            "scene_kind": "concept",
            "knowledge_refs": ["knowledge-1"],
            "blocks": [{
                "block_id": "concept-evidence",
                "type": "bullets",
                "items": [
                    "Evidence A explains the first decision.",
                    "Evidence B explains the second decision.",
                    "Evidence C explains the third decision.",
                ],
            }],
        },
        {
            "unit_id": "practice-shared-evidence",
            "layout": "question",
            "chapter_id": "chapter-1",
            "scene_kind": "practice_feedback",
            "beat_role": "prompt",
            "knowledge_refs": ["knowledge-1"],
            "title": "Check the three decisions",
            "blocks": [
                {
                    "block_id": "prompts",
                    "type": "question",
                    "items": ["Question A", "Question B", "Question C"],
                    "metadata": {"semantic_role": "prompt"},
                },
                {
                    "block_id": "source-checklist",
                    "type": "bullets",
                    "items": ["Check A", "Check B", "Check C"],
                    "metadata": {"semantic_role": "support"},
                },
            ],
            "visuals": [],
            "quality": {"requested_layout": "question-prompt"},
        },
    ])

    slide = apply_page_contract_v5(enriched[-1])

    assert [
        block["metadata"]["semantic_role"]
        for block in slide["blocks"]
    ] == ["prompt", "feedback"]
    assert slide["quality"]["feedback_mode"] == "shared_evidence"
    assert slide["quality"]["visible_item_count"] == 6
    assert slide["quality"]["visible_item_budget"] == 6
    assert not any(
        issue["code"] == "visible_item_overflow"
        for issue in v5_contract_issues([slide])
    )


def test_v5_paginates_every_practice_question_instead_of_hiding_overflow() -> None:
    questions = [f"Question {index}" for index in range(1, 6)]
    pages = _split_practice_feedback_capacity_v5([{
        "unit_id": "practice-five-questions",
        "layout": "practice",
        "title": "Check all five decisions",
        "blocks": [
            {
                "block_id": "prompts",
                "type": "question",
                "items": questions,
                "metadata": {
                    "semantic_role": "prompt",
                    "question_ids": [f"q-{index}" for index in range(1, 6)],
                },
            },
            {
                "block_id": "feedback",
                "type": "callout",
                "items": ["Evidence A", "Evidence B", "Evidence C"],
                "metadata": {"semantic_role": "feedback"},
            },
        ],
        "visuals": [],
        "quality": {
            "requested_layout": "practice-feedback",
            "feedback_mode": "shared_evidence",
            "semantic_atom_ids": ["atom-practice-five-questions"],
        },
    }])
    resolved = [apply_page_contract_v5(page) for page in pages]

    assert len(resolved) == 2
    assert resolved[0]["unit_id"] == "practice-five-questions"
    assert resolved[1]["unit_id"] == "practice-five-questions:practice:2"
    assert [
        question
        for page in resolved
        for block in page["blocks"]
        if block["metadata"]["semantic_role"] == "prompt"
        for question in block["items"]
    ] == questions
    assert all(
        page["quality"]["visible_item_count"]
        <= page["quality"]["visible_item_budget"]
        for page in resolved
    )
    issue_codes = {
        issue["code"]
        for issue in build_slide_deck_quality_v5(resolved)["issues"]
    }
    assert "visible_item_overflow" not in issue_codes
    assert "semantic_atom_split" not in issue_codes


def test_v5_promotes_feedback_group_labels_instead_of_counting_them_as_items() -> None:
    slide = apply_page_contract_v5({
        "unit_id": "practice-feedback-groups",
        "layout": "practice",
        "title": "任务一：绘图练习核对",
        "blocks": [{
            "block_id": "feedback",
            "type": "bullets",
            "items": [
                "**核对标准**：",
                "标出结构边界",
                "核对层次关系",
                "说明临床意义",
                "**典型错误**：",
                "遗漏关键结构",
                "混淆相邻层次",
            ],
        }],
        "visuals": [],
        "quality": {"requested_layout": "practice-feedback"},
    })

    assert slide["quality"]["resolved_layout"] == "practice-feedback"
    assert [block["title"] for block in slide["blocks"]] == [
        "核对标准",
        "典型错误",
    ]
    assert sum(len(block["items"]) for block in slide["blocks"]) == 5
    assert not any(
        issue["code"] == "visible_item_overflow"
        for issue in v5_contract_issues([slide])
    )


@pytest.mark.parametrize(
    ("course_type", "requested_layout", "blocks", "expected_layout"),
    [
        (
            "quantitative",
            "formula-explanation",
            [{
                "block_id": "energy-balance",
                "type": "formula",
                "content": "ΔU = Q + W 表示封闭系统的能量收支。",
                "metadata": {"formula": r"\Delta U = Q + W"},
            }],
            "formula-explanation",
        ),
        (
            "programming",
            "balanced-two-column",
            [
                {
                    "block_id": "source-code",
                    "type": "code",
                    "content": "def total(values):\n    return sum(values)",
                },
                {
                    "block_id": "reading",
                    "type": "rich_text",
                    "content": "先确认输入，再检查返回值和空列表边界。",
                },
            ],
            "balanced-two-column",
        ),
        (
            "humanities",
            "editorial-body",
            [{
                "block_id": "argument",
                "type": "rich_text",
                "content": "史料的作者、语境和受众共同决定证据能够支持的解释范围。",
            }],
            "editorial-body",
        ),
        (
            "business",
            "editorial-body",
            [{
                "block_id": "segments",
                "type": "bullets",
                "items": ["成本领先", "差异化", "聚焦细分市场"],
            }],
            "classification-3",
        ),
        (
            "medical-structural",
            "process-sequence",
            [{
                "block_id": "clinical-path",
                "type": "process",
                "items": ["采集症状", "识别危险信号", "形成鉴别诊断", "安排检查"],
            }],
            "process-sequence",
        ),
    ],
)
def test_v5_structural_evaluation_across_course_types(
    course_type: str,
    requested_layout: str,
    blocks: list[dict[str, object]],
    expected_layout: str,
) -> None:
    slide = apply_page_contract_v5({
        "unit_id": f"evaluation-{course_type}",
        "layout": "concept",
        "title": {
            "quantitative": "能量守恒连接热、功与内能",
            "programming": "函数契约决定边界行为",
            "humanities": "证据必须放回历史语境",
            "business": "竞争战略有三种基本选择",
            "medical-structural": "临床判断沿证据路径推进",
        }[course_type],
        "blocks": blocks,
        "quality": {"requested_layout": requested_layout},
    })

    assert slide["quality"]["resolved_layout"] == expected_layout
    assert slide["quality"]["density_band"] != "overflow"
    assert not v5_contract_issues([slide])


def test_final_repair_discards_stale_intermediate_capacity_findings() -> None:
    slides = repair_final_page_contracts_v5([{
        "unit_id": "slide:v5:repaired-page",
        "layout": "concept",
        "title": "三种系统具有不同交换边界",
        "blocks": [{
            "block_id": "classification",
            "type": "bullets",
            "items": ["孤立系统", "封闭系统", "开放系统"],
        }],
        "quality": {"requested_layout": "classification-3"},
    }])
    report = finalize_v5_quality_report(
        previous_quality={
            "passed": False,
            "blockers": [
                {
                    "severity": "critical",
                    "code": "visible_item_overflow",
                    "target": "slide:v5:repaired-page",
                },
                {
                    "severity": "critical",
                    "code": "enumeration_cardinality_mismatch",
                    "target": "slide:v5:repaired-page",
                },
            ],
            "semantic": {"issues": [{
                "severity": "major",
                "code": "slide_title_too_long",
                "target": "slide:v5:repaired-page",
            }]},
        },
        slides=slides,
        planner="ai",
        fallback_reason="",
    )

    assert report["passed"] is True
    assert slides[0]["quality"]["final_page_contract_v2"]["passed"] is True
    assert (
        slides[0]["quality"]["final_page_contract_v2"]["resolved_layout"]
        == "classification-3"
    )
    assert not {
        "visible_item_overflow",
        "enumeration_cardinality_mismatch",
        "slide_title_too_long",
    } & {issue["code"] for issue in report["issues"]}


def test_final_repair_paginates_promoted_groups_without_losing_content() -> None:
    labels = [f"Math relation {index}" for index in range(1, 9)]
    grouped_items = [
        value
        for index, label in enumerate(labels, start=1)
        for value in (f"**{label}**", f"Source-backed explanation {index}")
    ]

    slides = repair_final_page_contracts_v5([{
        "unit_id": "slide:v5:math-groups",
        "position": 0,
        "layout": "concept",
        "slide_purpose": "teach",
        "title": "Eight related mathematical statements",
        "blocks": [{
            "block_id": "math-groups",
            "type": "bullets",
            "title": "",
            "content": "",
            "items": grouped_items,
            "metadata": {
                "semantic_atom_id": "atom-math-groups",
                "source_fragment_ids": ["fragment-math-groups"],
            },
        }],
        "quality": {
            "requested_layout": "parallel-examples",
            "semantic_atom_ids": ["atom-math-groups"],
        },
    }])

    assert [len(slide["blocks"]) for slide in slides] == [4, 4]
    assert [
        block["title"]
        for slide in slides
        for block in slide["blocks"]
    ] == labels
    assert len({slide["unit_id"] for slide in slides}) == 2
    assert slides[1]["quality"]["continuation_of"] == slides[0]["unit_id"]
    assert slides[1]["quality"]["continuation_index"] == 2
    assert slides[1]["quality"]["continuation_total"] == 2
    atom_sets = [
        set(slide["quality"]["semantic_atom_ids"])
        for slide in slides
    ]
    assert atom_sets[0].isdisjoint(atom_sets[1])
    assert all(
        slide["quality"]["parent_semantic_atom_ids"] == ["atom-math-groups"]
        for slide in slides
    )
    quality = build_slide_deck_quality_v5(slides)
    assert "semantic_atom_split" not in {
        issue["code"] for issue in quality["issues"]
    }
    SlideDeckContent.model_validate({
        "schema_version": "slide_deck_v5",
        "title": "Math capacity regression",
        "slides": slides,
    })


def test_v5_reports_course_input_semantic_gaps_without_blocking_safe_pages() -> None:
    slide = apply_page_contract_v5({
        "unit_id": "slide:v5:safe-concept",
        "layout": "concept",
        "title": "结构关系决定判断顺序",
        "blocks": [{
            "block_id": "concept",
            "type": "rich_text",
            "content": "先确认对象的位置关系，再检查相邻结构，最终形成来源支持的完整判断。",
        }],
        "quality": {"requested_layout": "editorial-body"},
    })

    report = finalize_v5_quality_report(
        previous_quality={"passed": True, "issues": []},
        slides=[slide],
        planner="ai",
        fallback_reason="",
        planning_diagnostics={
            "structured_semantic_unit_count": 7,
            "balanced_composition_unit_count": 7,
            "semantic_role_counts": {
                "concept": 5,
                "activity": 1,
                "feedback": 1,
            },
            "knowledge_binding_unmapped_count": 7,
            "question_answer_binding_coverage": 1.0,
        },
    )

    assert report["passed"] is True
    issues = {issue["code"]: issue for issue in report["issues"]}
    assert issues["course_input_example_missing"]["responsibility"] == (
        "course_generation"
    )
    assert issues["course_input_knowledge_unmapped"]["responsibility"] == (
        "course_generation"
    )


def test_title_compiler_keeps_explicit_title_and_never_promotes_takeaway() -> None:
    title = compile_page_title_v5(
        explicit_title="热力学系统的三种类型",
        primary_claim="根据系统与环境之间的交互方式，热力学将系统分为三类。",
        body_text="根据系统与环境之间的交互方式，热力学将系统分为三类。",
    )

    assert title == "热力学系统的三种类型"


def test_v5_cover_splits_a_long_course_name_into_title_and_subtitle() -> None:
    document = _document(1).model_copy(update={
        "title": "《热力学与统计物理：原理、方法与应用》",
    })

    outline = compile_deck_outline_v5(document, _story(1))
    cover_contract = resolve_page_contract_v5({
        "layout": "cover",
        "title": outline.cover.title,
        "subtitle": outline.cover.subtitle,
        "blocks": [],
        "quality": {"requested_layout": "cover"},
    })

    assert outline.cover.title == "热力学与统计物理"
    assert outline.cover.subtitle == "原理、方法与应用"
    assert cover_contract.resolved_layout == "cover-editorial"


def test_v5_build_signature_invalidates_all_teaching_semantic_policies() -> None:
    signature = build_signature_v5(
        document=_document(1),
        course_data={},
        mode="teaching",
        theme="qizhi-classroom",
    )

    assert signature["semantic_compiler_version"].startswith(
        "ppt_teaching_semantics_v2"
    )
    assert signature["domain_presentation_profile_version"].startswith(
        "domain_presentation_profiles_v1"
    )
    assert signature["visual_planning_batch_version"].startswith(
        "chapter_visual_batches_v2"
    )


def test_v5_promotes_long_title_detail_into_supporting_copy() -> None:
    slide = apply_page_contract_v5({
        "unit_id": "classification-title",
        "layout": "concept",
        "title": "热力学将系统分为三类：这些分类帮助我们理解不同",
        "key_message": "",
        "blocks": [{
            "block_id": "classification",
            "type": "bullets",
            "items": ["孤立系统", "封闭系统", "开放系统"],
        }],
        "quality": {"requested_layout": "classification-3"},
    })

    assert slide["title"] == "热力学将系统分为三类"
    assert slide["key_message"] == "这些分类帮助我们理解不同"
    assert len(slide["title"]) <= 18


def test_parallel_application_items_never_become_a_worked_reasoning_chain() -> None:
    slide = apply_page_contract_v5({
        "unit_id": "three-parallel-applications",
        "layout": "case-study",
        "scene_kind": "application",
        "beat_role": "mapping",
        "title": "第零定律的实际应用",
        "blocks": [{
            "block_id": "applications",
            "type": "bullets",
            "items": ["空调温控", "冷链运输", "体温测量"],
        }],
        "quality": {"requested_layout": "worked-example"},
    })

    assert slide["quality"]["resolved_layout"] == "parallel-examples"
    assert slide["quality"]["resolved_composition"] == "parallel"
    assert not v5_contract_issues([slide])


def test_chapter_recap_uses_claims_and_a_retrieval_prompt_not_slide_titles() -> None:
    outline = compile_deck_outline_v5(_document(1), _story(1))
    source_slides = [
        {
            "unit_id": "classification",
            "title": "1.1 热力学系统的分类与描述",
            "takeaway": "系统类型取决于它与环境交换物质和能量的方式。",
            "blocks": [],
        },
        {
            "unit_id": "application",
            "title": "实践案例/行业应用",
            "key_message": "第零定律支撑温度测量和温度控制。",
            "blocks": [],
        },
    ]

    recap = _chapter_recap_slide(outline.chapters[0], source_slides)

    assert recap["title"] == "回顾：主题1"
    assert recap["blocks"][0]["items"] == [
        "系统类型取决于它与环境交换物质和能量的方式。",
        "第零定律支撑温度测量和温度控制。",
    ]
    assert "？" in recap["key_message"]
    assert all(len(item) <= 24 for item in recap["blocks"][0]["items"])
    assert "不看前文" in recap["key_message"]
    assert recap["quality"]["navigation_only"] is False
    assert recap["quality"]["retrieval_recap"] is True


@pytest.mark.parametrize(
    ("source", "expected"),
    [
        ("本节课的核心目标是建立局部解剖学的空间定位基础", "局部解剖学的空间定位基础"),
        ("起点：面神经主干自茎乳孔（Stylomastoid foramen）", "起点：面神经主干自茎乳孔"),
        ("本节课旨在掌握纵隔的“四分法”分区逻辑", "纵隔的四分法分区逻辑"),
        (
            "本节的核心任务是建立骨盆作为“骨性容器”与盆底肌群作为“功能性底板”的空间对应关系",
            "骨盆与盆底肌群的空间对应关系",
        ),
        (
            "本节旨在建立上肢近端至中段的“骨 - 肌 - 神经”空间对应关系",
            "骨 - 肌 - 神经空间对应关系",
        ),
        (
            "本模块依据下肢肌群配布的功能分区与协同拮抗机制开展教学",
            "下肢肌群配布的功能分区",
        ),
    ],
)
def test_concise_title_uses_complete_existing_phrases(source: str, expected: str) -> None:
    assert _concise_existing_title(source, maximum=18) == expected


def test_bounded_title_keeps_a_complete_capability_phrase() -> None:
    assert _bounded_title(
        "能够编写脚本动态调整 UI 适配参数，以应对非标准屏幕比例。",
        limit=24,
    ) == "编写脚本动态调整 UI 适配参数"
    assert _bounded_title(
        "编写脚本动态调整 UI 适配参数",
        limit=18,
    ) == "编写脚本动态调整 UI 适配参数"
    assert _title_with_continuation_sequence(
        "编写脚本动态调整 UI 适配参数",
        {
            "continuation_of": "slide:v5:root",
            "continuation_index": 2,
            "continuation_total": 2,
            "title_character_budget": 18,
        },
    ) == "编写脚本动态调整 UI 适配参数（续2/2）"


def test_quality_gate_blocks_mixed_question_and_chapter_transition() -> None:
    slide = apply_page_contract_v5({
        "unit_id": "mixed-question-transition",
        "layout": "practice",
        "title": "判断水壶属于哪类系统",
        "blocks": [
            {
                "block_id": "question",
                "type": "exercise",
                "content": "水壶盖子没有打开，这个系统属于哪种类型？",
            },
            {
                "block_id": "transition",
                "type": "statement",
                "content": "本节介绍了系统分类。下一节将深入探讨热力学第一定律。",
            },
        ],
        "quality": {"requested_layout": "two-column"},
    })

    assert "mixed_narrative_jobs" in {
        issue["code"] for issue in v5_contract_issues([slide])
    }
