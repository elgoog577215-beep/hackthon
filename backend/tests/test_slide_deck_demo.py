from __future__ import annotations

import asyncio
import time
from copy import deepcopy
from pathlib import Path
from zipfile import ZipFile

import pytest
from fastapi import FastAPI, Request
from fastapi.testclient import TestClient
from pptx import Presentation

import slide_deck
from course_document import (
    COURSE_DOCUMENT_SCHEMA,
    document_from_legacy_course,
    refresh_document_revision,
)
from course_repository import CourseDocumentRepository
from course_revisions import revision_event_for_documents
from representation_compiler import compile_core_representations
from slide_deck_renderer import SlideDeckQualityError, export_structured_slide_deck
from teaching_representations import TeachingRepresentationRepository


class MemoryStorage:
    def __init__(self, course: dict) -> None:
        self.course = deepcopy(course)

    def load_course(self, _course_id: str) -> dict:
        return deepcopy(self.course)

    async def save_course(self, _course_id: str, data: dict) -> None:
        self.course = deepcopy(data)


def advanced_python_course(section_count: int = 25) -> dict:
    roles = ("concept", "reasoning", "example", "application")
    nodes = []
    for index in range(section_count):
        section_number = index + 1
        role = roles[index % len(roles)]
        nodes.append({
            "node_id": f"advanced-python-{section_number}",
            "parent_node_id": "root",
            "node_name": f"Advanced Python 核心主题 {section_number}",
            "node_level": 1,
            "learning_objective": f"能够解释并应用 Python 主题 {section_number}",
            "objective_id": f"objective-{section_number}",
            "concept_refs": [f"python-concept-{section_number}"],
            "content_blocks": [{
                "block_id": f"advanced-python-{section_number}-block",
                "title": f"主题 {section_number} 的关键判断",
                "content": (
                    f"主题 {section_number} 解决一个具体的 Python 工程问题。\n\n"
                    f"先识别主题 {section_number} 的适用条件。\n\n"
                    f"再比较常见做法与推荐做法。\n\n"
                    f"最后用一个小例子验证主题 {section_number}。"
                ),
                "metadata": {"role": role},
            }],
        })
    return {
        "course_id": "advanced-python",
        "course_name": "Advanced Python",
        "nodes": nodes,
        "learning_assets": {
            "questions": [{
                "question_id": "python-practice-1",
                "revision_id": "python-practice-1-r1",
                "node_id": "advanced-python-8",
                "prompt": "说明主题 8 的适用边界，并给出一个最小示例。",
            }],
            "misconceptions": [],
        },
    }


def _compiled_slide_content(tmp_path: Path, course_data: dict | None = None) -> tuple[dict, object, object]:
    course = course_data or advanced_python_course(3)
    document = document_from_legacy_course(course)
    repository = TeachingRepresentationRepository(tmp_path / "representations")
    compile_core_representations(document, course, repository)
    registry = repository.load(document.course_id)
    representation = next(
        item for item in registry.representations
        if item.representation_type == "slide_deck"
    )
    spec = next(item for item in registry.specs if item.spec_id == representation.spec_id)
    return deepcopy(spec.payload["content"]), representation, spec


def test_advanced_python_course_compiles_to_demo_sized_source_bound_deck(tmp_path):
    course = advanced_python_course()
    document = document_from_legacy_course(course)

    content = slide_deck.compile_slide_deck(document, course)

    assert 12 <= len(content["slides"]) <= 18
    assert content["quality_summary"]["passed"] is True
    teaching_slides = [item for item in content["slides"] if item.get("section_id")]
    assert teaching_slides
    assert all(item["source_section_ids"] for item in teaching_slides)
    assert all(item["source_block_ids"] for item in teaching_slides)
    assert all(item["speaker_notes"].strip() for item in teaching_slides)
    assert sum(item["layout"] == "practice" for item in teaching_slides) <= 2
    for item in teaching_slides:
        visible_points = sum(
            len(block["items"]) if block["items"] else bool(block["content"].strip())
            for block in item["blocks"]
        )
        assert 3 <= visible_points <= 5


def test_production_plan_progress_exposes_explicit_fallback_reason():
    course = advanced_python_course()
    document = document_from_legacy_course(course)
    events: list[dict] = []

    slide_deck.compile_slide_deck(document, course, progress_callback=events.append)

    planned = next(event for event in events if event["event"] == "deck_plan")
    assert planned["planner"] == "deterministic_fallback"
    assert planned["fallback_reason"] == "no_ai_planner"


@pytest.mark.asyncio
async def test_compile_core_representations_consumes_injected_validated_plan(tmp_path):
    course = advanced_python_course()
    document = document_from_legacy_course(course)
    repository = TeachingRepresentationRepository(tmp_path)
    plan = await slide_deck.plan_slide_deck(document, course)
    plan.slides[2].slide_id = "slide:injected-production-plan"

    compile_core_representations(document, course, repository, deck_plan=plan)

    registry = repository.load(document.course_id)
    representation = next(
        item for item in registry.representations
        if item.representation_type == "slide_deck"
    )
    spec = next(item for item in registry.specs if item.spec_id == representation.spec_id)
    assert any(
        item["unit_id"] == "slide:injected-production-plan"
        for item in spec.payload["content"]["slides"]
    )


@pytest.mark.asyncio
async def test_build_stream_awaits_injected_planner_before_sync_compile(tmp_path, monkeypatch):
    from routers import teaching_representations as representation_router

    course = advanced_python_course()
    document = document_from_legacy_course(course)
    canonical = {
        **course,
        "course_schema_version": COURSE_DOCUMENT_SCHEMA,
        "course_document": document.model_dump(mode="json"),
        "course_document_authoritative": True,
        "course_operation_log": [],
    }
    course_repository = CourseDocumentRepository(MemoryStorage(canonical))
    representation_repository = TeachingRepresentationRepository(tmp_path)
    injected = await slide_deck.plan_slide_deck(document, course)
    injected.slides[2].slide_id = "slide:stream-injected-plan"
    raw_plan = injected.model_dump(mode="json")
    planner_calls: list[dict] = []

    async def injected_planner(request: dict) -> dict:
        planner_calls.append(request)
        return raw_plan

    monkeypatch.setattr(
        representation_router,
        "get_teaching_representation_repository",
        lambda: representation_repository,
    )
    monkeypatch.setattr(
        representation_router,
        "get_course_document_repository",
        lambda: course_repository,
    )
    monkeypatch.setattr(
        representation_router,
        "get_slide_deck_ai_planner",
        lambda: injected_planner,
        raising=False,
    )

    async def existing_course(_course_id: str):
        return course_repository.load_course_view("advanced-python")

    monkeypatch.setattr(representation_router, "get_course_or_404", existing_course)
    app = FastAPI()
    app.include_router(representation_router.router, prefix="/api")
    client = TestClient(app)

    with client.stream(
        "POST",
        "/api/courses/advanced-python/teaching-representations/build/stream",
        headers={"X-User-Id": "teacher-1"},
    ) as response:
        body = "".join(response.iter_text())

    assert response.status_code == 200
    assert len(planner_calls) == 1
    assert "event: build_complete" in body
    registry = representation_repository.load("advanced-python")
    representation = next(
        item for item in registry.representations
        if item.representation_type == "slide_deck"
    )
    spec = next(item for item in registry.specs if item.spec_id == representation.spec_id)
    assert any(
        item["unit_id"] == "slide:stream-injected-plan"
        for item in spec.payload["content"]["slides"]
    )


@pytest.mark.asyncio
async def test_build_stream_emits_planner_started_before_slow_planner_finishes(monkeypatch):
    from routers import teaching_representations as representation_router

    async def existing_course(_course_id: str):
        return {"course_id": "course-1"}

    planner_finished = asyncio.Event()

    async def slow_plan(_course_id: str):
        await asyncio.sleep(0.2)
        planner_finished.set()
        return object()

    monkeypatch.setattr(representation_router, "get_course_or_404", existing_course)
    monkeypatch.setattr(representation_router, "_plan_registry_slide_deck", slow_plan)
    request = Request({
        "type": "http",
        "method": "POST",
        "path": "/api/courses/course-1/teaching-representations/build/stream",
        "headers": [(b"x-user-id", b"teacher-1")],
        "query_string": b"",
        "scheme": "http",
        "server": ("testserver", 80),
        "client": ("testclient", 50000),
    })

    response = await asyncio.wait_for(
        representation_router.stream_teaching_representation_build("course-1", request),
        timeout=0.05,
    )
    first_chunk = await asyncio.wait_for(anext(response.body_iterator), timeout=0.05)

    assert "event: planner_started" in first_chunk
    assert not planner_finished.is_set()
    await response.body_iterator.aclose()


def test_demo_practice_page_caps_generated_checks_at_five_points():
    course = advanced_python_course()
    course["nodes"][-1]["misconceptions"] = [
        f"主题 25 的常见误区 {index}"
        for index in range(3)
    ]
    course["learning_assets"] = {
        "questions": [],
        "misconceptions": [],
    }
    document = document_from_legacy_course(course)

    content = slide_deck.compile_slide_deck(document, course)
    practice = next(item for item in content["slides"] if item["layout"] == "practice")
    visible_points = sum(
        len(block["items"]) if block["items"] else bool(block["content"].strip())
        for block in practice["blocks"]
    )

    assert 3 <= visible_points <= 5


def test_practice_misconceptions_keep_a_complete_readable_phrase():
    course = advanced_python_course(1)
    course["course_knowledge_base"] = {
        "knowledge_points": [{"knowledge_id": "knowledge-1", "name": "scope rules"}],
        "skill_units": [],
        "misconceptions": [{
            "misconception_id": "misconception-1",
            "primary_knowledge_id": "knowledge-1",
            "observable_error_pattern": (
                "Assuming that global and nonlocal are interchangeable or equivalent in nested scopes"
            ),
        }],
        "mastery_criteria": [],
        "bindings": [{
            "target_type": "section",
            "target_id": "advanced-python-1",
            "knowledge_ids": ["knowledge-1"],
            "skill_ids": [],
        }],
    }
    document = document_from_legacy_course(course)
    section = next(item for item in document.sections if item.section_id == "advanced-python-1")
    context = slide_deck._KnowledgeIndex(course).for_section(section.section_id, section)
    source = slide_deck.SlideBlockSpec(
        block_id="misconception",
        type="misconception",
        items=context["misconception_labels"],
    )

    fitted = slide_deck._fit_blocks_for_layout("practice", [source])
    visible = fitted[0].items[0]

    assert visible == (
        "Assuming that global and nonlocal are interchangeable or equivalent in nested scopes"
    )
    assert len(visible) <= 110


def test_demo_visible_copy_uses_complete_short_points_without_ellipsis():
    course = advanced_python_course()
    for node in course["nodes"]:
        node["learning_objective"] = (
            "Construct and debug multithreaded Python applications for I/O-bound "
            "workloads under production pressure"
        )
        node["content_blocks"][0]["content"] = (
            "You are preparing a production service that must remain responsive under load, "
            "recover from partial failures, preserve resources, and explain every design tradeoff."
        )
    course["learning_assets"]["misconceptions"] = [
        {
            "node_id": node["node_id"],
            "error_pattern": (
                "Assuming that every production scheduling mechanism behaves identically "
                "under cancellation pressure and resource exhaustion"
            ),
        }
        for node in course["nodes"]
    ]
    document = document_from_legacy_course(course)

    content = slide_deck.compile_slide_deck(document, course)
    report = slide_deck.validate_slide_deck(content)
    concept_slides = [slide for slide in content["slides"] if slide["layout"] == "concept"]

    visible: list[str] = []
    for slide in content["slides"]:
        visible.extend([slide["title"], slide["key_message"]])
        for block in slide["blocks"]:
            visible.extend([block.get("title") or "", block.get("content") or ""])
            visible.extend(block.get("items") or [])
    truncated = [
        value for value in visible
        if value.strip().endswith(("…", "..."))
    ]
    assert truncated == []
    assert concept_slides
    assert all(
        slide["key_message"].startswith("本页聚焦：")
        for slide in concept_slides
    )
    assert report["passed"] is True


@pytest.mark.asyncio
async def test_invalid_ai_planner_uses_deterministic_fallback():
    course = advanced_python_course()
    document = document_from_legacy_course(course)

    async def invalid_planner(_request: dict) -> dict:
        return {
            "schema_version": "slide_deck_plan_v1",
            "title": "invalid",
            "target_slide_count": 15,
            "slides": [],
        }

    plan = await slide_deck.plan_slide_deck(
        document,
        course,
        ai_planner=invalid_planner,
        timeout_seconds=0.1,
    )

    assert isinstance(plan, slide_deck.SlideDeckPlanV1)
    assert plan.planner == "deterministic_fallback"
    assert plan.fallback_reason == "invalid_plan"
    assert len(plan.slides) == 15
    default_plan = await slide_deck.plan_slide_deck(document, course)
    assert default_plan.fallback_reason == "no_ai_planner"
    assert plan.model_dump(mode="json") == default_plan.model_dump(mode="json") | {
        "fallback_reason": "invalid_plan",
    }


@pytest.mark.asyncio
async def test_timed_out_ai_planner_uses_deterministic_fallback():
    course = advanced_python_course()
    document = document_from_legacy_course(course)

    async def slow_planner(_request: dict) -> dict:
        await asyncio.sleep(0.05)
        return {}

    plan = await slide_deck.plan_slide_deck(
        document,
        course,
        ai_planner=slow_planner,
        timeout_seconds=0.001,
    )

    assert plan.planner == "deterministic_fallback"
    assert plan.fallback_reason == "timeout"
    assert len(plan.slides) == 15


@pytest.mark.asyncio
async def test_sync_ai_planner_is_preempted_by_timeout():
    course = advanced_python_course()
    document = document_from_legacy_course(course)

    def blocking_planner(_request: dict) -> dict:
        time.sleep(0.05)
        return {}

    started = time.perf_counter()
    plan = await slide_deck.plan_slide_deck(
        document,
        course,
        ai_planner=blocking_planner,
        timeout_seconds=0.001,
    )
    elapsed = time.perf_counter() - started

    assert plan.fallback_reason == "timeout"
    assert elapsed < 0.03


def test_compressed_formal_practice_warnings_are_aggregated_without_zeroing_score():
    course = advanced_python_course()
    course["learning_assets"]["questions"] = [
        {
            "question_id": f"python-practice-{index}",
            "revision_id": f"python-practice-{index}-r1",
            "node_id": f"advanced-python-{index}",
            "prompt": f"说明主题 {index} 的适用边界。",
        }
        for index in range(1, 26)
    ]
    document = document_from_legacy_course(course)

    content = slide_deck.compile_slide_deck(document, course)
    report = slide_deck.validate_slide_deck(content, course_data=course)
    compression_warnings = [
        issue for issue in report["warnings"]
        if issue["code"] == "formal_practice_not_represented"
    ]

    assert report["passed"] is True
    assert report["score"] >= 0.8
    assert len(compression_warnings) == 1
    assert compression_warnings[0]["count"] >= 20
    assert len(compression_warnings[0]["section_ids"]) == compression_warnings[0]["count"]
    assert not any(
        issue["code"] == "formal_practice_not_represented"
        for issue in report["blockers"]
    )


def test_two_export_themes_render_without_regenerating_content(tmp_path):
    content, _, _ = _compiled_slide_content(tmp_path)
    original = deepcopy(content)

    qingfeng_path = export_structured_slide_deck(
        content,
        tmp_path / "qingfeng.pptx",
        theme="qingfeng-classroom",
    )
    academic_path = export_structured_slide_deck(
        content,
        tmp_path / "academic.pptx",
        theme="academic-bluegray",
    )

    assert content == original
    assert qingfeng_path.exists() and academic_path.exists()
    qingfeng = Presentation(qingfeng_path)
    academic = Presentation(academic_path)
    assert str(qingfeng.slides[0].shapes[0].fill.fore_color.rgb) == "F7FAFC"
    assert str(academic.slides[0].shapes[0].fill.fore_color.rgb) == "FCFCFD"
    with pytest.raises(ValueError, match="Unknown slide theme"):
        export_structured_slide_deck(content, tmp_path / "invalid.pptx", theme="neon")


def test_academic_theme_writes_title_body_and_formula_fonts(tmp_path):
    course = advanced_python_course(3)
    course["nodes"][0]["content_blocks"][0]["content"] = (
        r"向量满足 $\mathbf{x} \in \mathbb{R}^3$，并可用于表示三维坐标。"
    )
    content, _, _ = _compiled_slide_content(tmp_path, course)

    output = export_structured_slide_deck(
        content,
        tmp_path / "academic-formula.pptx",
        theme="academic-bluegray",
    )

    with ZipFile(output) as archive:
        slide_xml = b"\n".join(
            archive.read(name)
            for name in archive.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        )
    assert b"Noto Serif SC" in slide_xml
    assert b"SimSun" in slide_xml
    assert b"Noto Sans SC" in slide_xml
    assert b"Microsoft YaHei" in slide_xml
    assert b"Times New Roman" in slide_xml


def test_export_api_accepts_and_validates_theme_query(tmp_path, monkeypatch):
    from routers import teaching_representations as representation_router

    _, representation, spec = _compiled_slide_content(tmp_path)
    exported_themes: list[str] = []

    async def fake_spec(*_args, **_kwargs):
        return {
            "status": "success",
            "representation": representation.model_dump(mode="json"),
            "spec": spec.model_dump(mode="json"),
        }

    def fake_export(_spec, output_path, *, theme):
        exported_themes.append(theme)
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        Path(output_path).write_bytes(b"pptx")
        return Path(output_path)

    monkeypatch.setattr(representation_router, "get_teaching_representation_spec", fake_spec)
    monkeypatch.setattr(representation_router, "export_slide_deck_pptx", fake_export)
    monkeypatch.setattr(representation_router, "DATA_DIR", tmp_path)
    app = FastAPI()
    app.include_router(representation_router.router, prefix="/api")
    client = TestClient(app)
    url = (
        f"/api/courses/advanced-python/teaching-representations/"
        f"{representation.representation_id}/export.pptx"
    )

    response = client.get(url, params={"theme": "academic-bluegray"})
    invalid = client.get(url, params={"theme": "neon"})

    assert response.status_code == 200
    assert exported_themes == ["academic-bluegray"]
    assert invalid.status_code == 422
    assert invalid.json()["detail"]["code"] == "invalid_slide_theme"


def test_export_api_returns_actionable_quality_report(tmp_path, monkeypatch):
    from routers import teaching_representations as representation_router

    content, representation, spec = _compiled_slide_content(tmp_path)
    target = next(item for item in content["slides"] if item.get("section_id"))
    target["blocks"][0]["content"] = "整段正文复制。" * 80
    spec.payload["content"] = content

    async def fake_spec(*_args, **_kwargs):
        return {
            "status": "success",
            "representation": representation.model_dump(mode="json"),
            "spec": spec.model_dump(mode="json"),
        }

    monkeypatch.setattr(representation_router, "get_teaching_representation_spec", fake_spec)
    monkeypatch.setattr(representation_router, "DATA_DIR", tmp_path)
    app = FastAPI()
    app.include_router(representation_router.router, prefix="/api")
    client = TestClient(app)

    response = client.get(
        f"/api/courses/{spec.course_id}/teaching-representations/"
        f"{representation.representation_id}/export.pptx",
        params={"theme": "qingfeng-classroom"},
    )

    assert response.status_code == 422
    detail = response.json()["detail"]
    assert detail["code"] == "slide_export_quality_blocked"
    assert detail["blockers"]
    assert all(
        issue["slide_id"] and issue["layout"] and issue["suggestion"]
        for issue in detail["blockers"]
    )
    assert "warnings" in detail


def test_quality_report_separates_blockers_and_actionable_warnings(tmp_path):
    content, _, _ = _compiled_slide_content(tmp_path)
    target = next(item for item in content["slides"] if item.get("section_id"))
    target["title"] = "这是一个明显过长但不应该阻断导出的课堂演示标题" * 3
    target["blocks"][0]["content"] = "整段正文复制。" * 80

    report = slide_deck.validate_slide_deck(content)

    assert report["passed"] is False
    assert report["blockers"]
    assert report["warnings"]
    assert any(item["code"] == "slide_title_too_long" for item in report["warnings"])
    for issue in [*report["blockers"], *report["warnings"]]:
        assert issue["slide_id"]
        assert issue["layout"]
        assert issue["message"]
        assert issue["suggestion"]


def _rendered_slide_text(presentation: Presentation, position: int) -> str:
    return "\n".join(
        shape.text
        for shape in presentation.slides[position].shapes
        if hasattr(shape, "text_frame") and shape.has_text_frame
    )


def _raw_deck_with_instructional_slide(slide: slide_deck.SlideSpec) -> dict:
    slide.position = 1
    return slide_deck.SlideDeckContent(
        title="Raw quality contract",
        slides=[
            slide_deck.SlideSpec(
                unit_id="slide:raw:cover",
                position=0,
                layout="cover",
                slide_purpose="orientation",
                title="Raw quality contract",
                blocks=[slide_deck.SlideBlockSpec(
                    block_id="slide:raw:cover:promise",
                    type="callout",
                    content="验证未经过生产 fit 的共享 spec。",
                )],
            ),
            slide,
        ],
    ).model_dump(mode="json")


def test_raw_objective_region_overflow_blocks_direct_export(tmp_path):
    content = _raw_deck_with_instructional_slide(slide_deck.SlideSpec(
        unit_id="slide:raw:objective",
        position=0,
        layout="objective",
        slide_purpose="learning_objective",
        title="目标页区域容量",
        section_id="section-raw",
        source_section_ids=["section-raw"],
        source_block_ids=["block-raw"],
        blocks=[
            slide_deck.SlideBlockSpec(
                block_id="objective-question",
                type="callout",
                content="需要解决什么问题？",
            ),
            slide_deck.SlideBlockSpec(
                block_id="objective-knowledge",
                type="bullets",
                items=["知识一", "知识二", "知识三"],
            ),
            slide_deck.SlideBlockSpec(
                block_id="objective-ability",
                type="bullets",
                items=["能力一"],
            ),
        ],
    ))

    report = slide_deck.validate_slide_deck(content)

    assert "objective_content_overflow" in {item["code"] for item in report["blockers"]}
    with pytest.raises(SlideDeckQualityError) as exc_info:
        export_structured_slide_deck(content, tmp_path / "raw-objective-overflow.pptx")
    assert "objective_content_overflow" in {
        item["code"] for item in exc_info.value.report["blockers"]
    }


def test_raw_concept_card_overflow_blocks_direct_export(tmp_path):
    content = _raw_deck_with_instructional_slide(slide_deck.SlideSpec(
        unit_id="slide:raw:concept",
        position=0,
        layout="concept",
        slide_purpose="concept_and_reasoning",
        title="概念卡区域容量",
        section_id="section-raw",
        source_section_ids=["section-raw"],
        source_block_ids=["block-raw"],
        blocks=[
            slide_deck.SlideBlockSpec(
                block_id="concept-skewed-card",
                type="bullets",
                title="单卡偏斜",
                items=["卡片要点一", "卡片要点二", "卡片要点三", "卡片要点四"],
            ),
            slide_deck.SlideBlockSpec(
                block_id="concept-short-card",
                type="statement",
                title="短卡",
                content="另一张卡保持简短。",
            ),
        ],
    ))

    report = slide_deck.validate_slide_deck(content)

    assert "concept_card_overflow" in {item["code"] for item in report["blockers"]}
    with pytest.raises(SlideDeckQualityError) as exc_info:
        export_structured_slide_deck(content, tmp_path / "raw-concept-overflow.pptx")
    assert "concept_card_overflow" in {
        item["code"] for item in exc_info.value.report["blockers"]
    }


def test_compiled_objective_preserves_skewed_knowledge_and_ability_in_notes(tmp_path):
    course = advanced_python_course(1)
    knowledge_labels = [f"知识坐标 {index}" for index in range(1, 5)]
    ability_labels = [f"能力目标 {index}" for index in range(1, 4)]
    course["course_knowledge_base"] = {
        "knowledge_points": [
            {"knowledge_id": f"knowledge-{index}", "name": label}
            for index, label in enumerate(knowledge_labels, start=1)
        ],
        "skill_units": [
            {"skill_id": f"skill-{index}", "observable_behavior": label}
            for index, label in enumerate(ability_labels, start=1)
        ],
        "misconceptions": [],
        "mastery_criteria": [],
        "bindings": [{
            "target_type": "section",
            "target_id": "advanced-python-1",
            "knowledge_ids": [f"knowledge-{index}" for index in range(1, 5)],
            "skill_ids": [f"skill-{index}" for index in range(1, 4)],
        }],
    }
    document = document_from_legacy_course(course)

    content = slide_deck.compile_slide_deck(document, course)
    objective = next(item for item in content["slides"] if item["layout"] == "objective")
    output = export_structured_slide_deck(content, tmp_path / "objective-notes.pptx")
    presentation = Presentation(output)
    rendered = _rendered_slide_text(presentation, objective["position"])
    notes = presentation.slides[objective["position"]].notes_slide.notes_text_frame.text

    assert sum(len(block["items"]) for block in objective["blocks"]) == 4
    assert knowledge_labels[-1] not in rendered
    assert ability_labels[-1] not in rendered
    assert knowledge_labels[-1] in objective["speaker_notes"]
    assert ability_labels[-1] in objective["speaker_notes"]
    assert knowledge_labels[-1] in notes
    assert ability_labels[-1] in notes


def test_compiled_concept_preserves_fourth_source_block_in_pptx_and_notes(tmp_path):
    course = advanced_python_course(1)
    overflow_detail = "闭包捕获变量绑定而非循环当下值，循环创建函数时应显式绑定参数快照。"
    course["nodes"][0]["content_blocks"] = [
        {
            "block_id": "scope-foundation",
            "title": "作用域基础",
            "content": (
                "名称查找遵循词法作用域。\n\n"
                "- 局部名称优先\n"
                "- 外层名称按需查找"
            ),
            "metadata": {"role": "concept"},
        },
        {
            "block_id": "closure-capture",
            "title": "闭包捕获语义",
            "content": (
                f"{overflow_detail}\n\n"
                "- 检查自由变量来自哪个作用域\n"
                "- 循环创建函数时验证最终绑定值"
            ),
            "metadata": {"role": "concept"},
        },
    ]
    document = document_from_legacy_course(course)
    original_slide_blocks = [
        slide_block
        for course_block in document.blocks
        for slide_block in slide_deck._block_to_slide_blocks(course_block)
    ]

    assert len(document.blocks) == 2
    assert len(original_slide_blocks) == 4
    assert overflow_detail in original_slide_blocks[3].content

    content = slide_deck.compile_slide_deck(document, course)
    concept = next(
        item for item in content["slides"]
        if item["layout"] == "concept" and item["slide_purpose"] == "concept_and_reasoning"
    )
    overflow_summaries = [
        item
        for block in concept["blocks"]
        for item in block["items"]
        if "闭包捕获变量绑定" in item
    ]

    assert overflow_summaries
    assert len(concept["blocks"]) == slide_deck.LAYOUT_CAPACITY["concept"]["blocks"]
    assert overflow_detail in concept["speaker_notes"]
    output = export_structured_slide_deck(content, tmp_path / "concept-contract.pptx")
    presentation = Presentation(output)
    rendered = _rendered_slide_text(presentation, concept["position"])

    assert overflow_summaries[0] in rendered
    assert overflow_detail not in rendered
    assert overflow_detail in presentation.slides[concept["position"]].notes_slide.notes_text_frame.text


def test_compiled_practice_preserves_trimmed_prompt_and_hidden_checks_in_notes(tmp_path):
    course = advanced_python_course(1)
    full_prompt = (
        "请分析闭包在循环中捕获变量绑定时的行为，并比较默认参数快照与延迟查找两种机制；"
        "随后给出一个能够稳定复现错误结果的最小示例，解释每次调用为何得到相同值，"
        "最后提出修复方案、验证修复结果，并说明方案在异步回调和事件处理器中的适用边界。"
        "作答时还需要逐步记录自由变量的查找路径、函数对象创建时机、调用发生时的环境状态，"
        "并比较列表推导式、普通循环和异步任务调度下的差异；请给出至少两个反例，"
        "说明哪些修改只是偶然改变输出顺序，哪些修改真正固定了每个回调所需的数据，"
        "最后写出可重复执行的验证步骤，并解释为什么这些步骤足以排除缓存和共享状态干扰。"
    )
    checks = [f"隐藏检查语义 {index}" for index in range(1, 6)]
    course["nodes"][0]["misconceptions"] = checks
    course["learning_assets"] = {
        "questions": [{
            "question_id": "closure-practice",
            "revision_id": "closure-practice-r1",
            "node_id": "advanced-python-1",
            "prompt": full_prompt,
        }],
        "misconceptions": [],
    }
    document = document_from_legacy_course(course)

    content = slide_deck.compile_slide_deck(document, course)
    practice = next(item for item in content["slides"] if item["layout"] == "practice")
    output = export_structured_slide_deck(content, tmp_path / "practice-notes.pptx")
    presentation = Presentation(output)
    rendered = _rendered_slide_text(presentation, practice["position"])
    notes = presentation.slides[practice["position"]].notes_slide.notes_text_frame.text

    assert sum(
        len(block["items"]) if block["items"] else bool(block["content"])
        for block in practice["blocks"]
    ) <= 5
    assert full_prompt not in rendered
    assert checks[-1] not in rendered
    assert full_prompt in practice["speaker_notes"]
    assert checks[-1] in practice["speaker_notes"]
    assert full_prompt in notes
    assert checks[-1] in notes


def test_practice_quality_gate_blocks_five_checks_even_when_total_items_is_eight(tmp_path):
    course = advanced_python_course(1)
    document = document_from_legacy_course(course)
    content = slide_deck.compile_slide_deck(document, course)
    practice = next(item for item in content["slides"] if item["layout"] == "practice")
    practice["blocks"] = [
        slide_deck.SlideBlockSpec(
            block_id="practice-exercise",
            type="exercise",
            title="独立作答",
            items=["完成判断", "说明依据", "给出示例"],
        ).model_dump(mode="json"),
        slide_deck.SlideBlockSpec(
            block_id="practice-checks",
            type="bullets",
            title="检查标准",
            items=[f"检查项 {index}" for index in range(1, 6)],
        ).model_dump(mode="json"),
    ]

    report = slide_deck.validate_slide_deck(content)

    assert report["passed"] is False
    assert "practice_check_overflow" in {item["code"] for item in report["blockers"]}
    with pytest.raises(SlideDeckQualityError) as exc_info:
        export_structured_slide_deck(content, tmp_path / "invalid-practice.pptx")
    assert "practice_check_overflow" in {
        item["code"] for item in exc_info.value.report["blockers"]
    }


def test_practice_fit_exports_every_legal_exercise_and_check_item(tmp_path):
    course = advanced_python_course(1)
    document = document_from_legacy_course(course)
    content = slide_deck.compile_slide_deck(document, course)
    practice = next(item for item in content["slides"] if item["layout"] == "practice")
    original_blocks = [
        slide_deck.SlideBlockSpec(
            block_id="practice-exercise",
            type="exercise",
            title="独立作答",
            items=[f"作答步骤 {index}" for index in range(1, 5)],
        ),
        slide_deck.SlideBlockSpec(
            block_id="practice-checks",
            type="bullets",
            title="检查标准",
            items=[f"可见检查项 {index}" for index in range(1, 6)],
        ),
    ]

    fitted = slide_deck._fit_blocks_for_layout("practice", original_blocks)

    assert fitted[1].items == [f"可见检查项 {index}" for index in range(1, 5)]
    practice["blocks"] = [block.model_dump(mode="json") for block in fitted]
    output = export_structured_slide_deck(content, tmp_path / "legal-practice.pptx")
    rendered = _rendered_slide_text(Presentation(output), practice["position"])

    for block in fitted:
        for item in block.items:
            assert item in rendered


@pytest.mark.parametrize(
    ("layout", "max_blocks", "max_items"),
    [
        ("cover", 1, 0),
        ("roadmap", 1, 8),
        ("chapter", 1, 0),
        ("objective", 3, 4),
        ("concept", 3, 9),
        ("comparison", 1, 0),
        ("process", 1, 5),
        ("code", 2, 4),
        ("misconception", 1, 4),
        ("practice", 2, 8),
        ("recap", 3, 8),
    ],
)
def test_layout_capacity_matches_renderer_consumption(layout, max_blocks, max_items):
    capacity = slide_deck.LAYOUT_CAPACITY[layout]

    assert capacity["blocks"] == max_blocks
    assert capacity["items"] == max_items

    block_overflow = slide_deck.SlideSpec(
        unit_id=f"slide:{layout}:blocks",
        position=0,
        layout=layout,
        slide_purpose="capacity_check",
        title="容量检查",
        blocks=[
            slide_deck.SlideBlockSpec(block_id=f"block-{index}", type="statement", content="可见内容")
            for index in range(max_blocks + 1)
        ],
    )
    assert "slide_block_overflow" in {
        issue["code"] for issue in slide_deck.slide_quality(block_overflow)["issues"]
    }

    remaining = max_items + 1
    item_blocks = []
    for index in range(max_blocks):
        count = min(8, remaining)
        if count <= 0:
            break
        item_blocks.append(slide_deck.SlideBlockSpec(
            block_id=f"items-{index}",
            type="bullets",
            items=[f"要点 {index}-{item}" for item in range(count)],
        ))
        remaining -= count
    if remaining:
        with pytest.raises(ValueError):
            slide_deck.SlideBlockSpec(
                block_id="model-item-overflow",
                type="bullets",
                items=[f"要点 {item}" for item in range(9)],
            )
        return

    item_overflow = slide_deck.SlideSpec(
        unit_id=f"slide:{layout}:items",
        position=0,
        layout=layout,
        slide_purpose="capacity_check",
        title="容量检查",
        blocks=item_blocks,
    )
    assert "slide_item_overflow" in {
        issue["code"] for issue in slide_deck.slide_quality(item_overflow)["issues"]
    }


def test_concept_fit_enforces_readable_per_card_limits():
    blocks = [
        slide_deck.SlideBlockSpec(
            block_id=f"concept-{index}",
            type="bullets",
            title=f"概念卡 {index}",
            content="不会用总字符数掩盖单卡溢出" * 12,
            items=["单个要点也必须适合三列投影版式" * 4 for _ in range(5)],
        )
        for index in range(4)
    ]

    fitted = slide_deck._fit_blocks_for_layout("concept", blocks)

    assert len(fitted) == 3
    assert all(len(block.items) <= 3 for block in fitted)
    assert all(len(item) <= 32 for block in fitted for item in block.items)
    assert all(len(block.content) <= 96 for block in fitted)
    assert not any(
        issue["code"] == "concept_card_overflow"
        for issue in slide_deck.slide_quality(slide_deck.SlideSpec(
            unit_id="slide:concept:fit",
            position=0,
            layout="concept",
            slide_purpose="concept",
            title="概念卡容量",
            blocks=fitted,
        ))["issues"]
    )


def test_objective_slide_fits_callout_before_validation():
    course = advanced_python_course(1)
    course["nodes"][0]["learning_objective"] = (
        "能够在复杂工程约束下解释高级 Python 机制的适用条件、运行边界、失败模式、"
        "资源释放策略以及不同实现之间的取舍，并能独立验证最终方案"
    )
    document = document_from_legacy_course(course)
    section = document.sections[0]
    blocks = [block for block in document.blocks if block.section_id == section.section_id]
    context = slide_deck._merge_block_knowledge(
        slide_deck._KnowledgeIndex(course).for_section(section.section_id, section),
        blocks,
    )

    objective = slide_deck._objective_slide(section, blocks, context)
    callout = next(block for block in objective.blocks if block.type == "callout")

    assert len(callout.content) <= 84
    assert "objective_content_overflow" not in {
        issue["code"] for issue in slide_deck.slide_quality(objective)["issues"]
    }


def test_assessment_slide_fits_three_candidates_to_two_renderer_blocks():
    course = advanced_python_course(1)
    course["nodes"][0]["misconceptions"] = ["忽略主题的适用边界"]
    course["learning_assets"]["questions"] = [{
        "question_id": "formal-check",
        "revision_id": "formal-check-r1",
        "node_id": "advanced-python-1",
        "prompt": "说明主题 1 的适用边界。",
    }]
    document = document_from_legacy_course(course)
    section = document.sections[0]
    blocks = [block for block in document.blocks if block.section_id == section.section_id]
    context = slide_deck._merge_block_knowledge(
        slide_deck._KnowledgeIndex(course).for_section(section.section_id, section),
        blocks,
    )

    assessment = slide_deck._assessment_slide(section, blocks, context, course)

    assert len(assessment.blocks) <= slide_deck.LAYOUT_CAPACITY["practice"]["blocks"]
    assert "slide_block_overflow" not in {
        issue["code"] for issue in slide_deck.slide_quality(assessment)["issues"]
    }


def test_demo_practice_padding_stays_inside_two_renderer_blocks():
    fitted = slide_deck._fit_blocks_for_layout("practice", [
        slide_deck.SlideBlockSpec(
            block_id="exercise",
            type="exercise",
            content="完成一个判断",
        ),
        slide_deck.SlideBlockSpec(
            block_id="check",
            type="callout",
            content="说明理由",
        ),
    ])

    padded = slide_deck._cap_demo_practice_points(fitted)

    assert len(padded) <= slide_deck.LAYOUT_CAPACITY["practice"]["blocks"]
    assert sum(len(block.items) if block.items else bool(block.content) for block in padded) >= 3


@pytest.mark.parametrize(
    ("section_count", "expected_strategy"),
    [(1, "plan_then_fill"), (2, "plan_then_fill"), (3, "plan_then_fill"), (4, "demo_plan_v1")],
)
def test_one_to_four_active_sections_use_the_expected_plan_boundary(section_count, expected_strategy):
    course = advanced_python_course(section_count)
    document = document_from_legacy_course(course)
    events: list[dict] = []

    content = slide_deck.compile_slide_deck(document, course, progress_callback=events.append)
    plan_event = next(item for item in events if item["event"] == "deck_plan")

    assert plan_event["strategy"] == expected_strategy
    if section_count == 4:
        assert 12 <= len(content["slides"]) <= 18


def test_large_course_remains_demo_sized():
    course = advanced_python_course(31)
    content = slide_deck.compile_slide_deck(document_from_legacy_course(course), course)

    assert 12 <= len(content["slides"]) <= 18
    assert len(content["slides"]) != 125


def test_cover_binds_every_section_structure_and_learning_objective():
    course = advanced_python_course(3)
    content = slide_deck.compile_slide_deck(document_from_legacy_course(course), course)
    cover = content["slides"][0]
    expected = {
        "course_title",
        *{f"section_structure:advanced-python-{index}" for index in range(1, 4)},
        *{f"objective:objective-{index}" for index in range(1, 4)},
    }

    assert expected <= set(cover["source_keys"])


@pytest.mark.parametrize("change_kind", ["section", "objective"])
def test_section_or_objective_change_stales_and_rebuilds_cover(tmp_path, change_kind):
    course = advanced_python_course(3)
    before = document_from_legacy_course(course)
    repository = TeachingRepresentationRepository(tmp_path / change_kind)
    compile_core_representations(before, course, repository)
    after = before.model_copy(deep=True)
    target = after.sections[-1]
    if change_kind == "section":
        target.title = "更新后的课程收束主题"
        expected_source_key = f"section_structure:{target.section_id}"
    else:
        target.learning_objective = "能够独立解释更新后的课程学习承诺"
        target.objective_revision_id = ""
        expected_source_key = f"objective:{target.objective_id}"
    refresh_document_revision(after)
    event = revision_event_for_documents(before, after, command_id=f"cover-{change_kind}-change")

    assert expected_source_key in event.changed_source_keys
    stale = repository.apply_revision_event(before.course_id, event)
    slide_representation = next(
        item for item in stale.representations if item.representation_type == "slide_deck"
    )
    assert "slide:title" in slide_representation.stale_unit_ids

    build = compile_core_representations(after, course, repository)
    slide_build = next(
        item for item in build["representations"] if item["representation_type"] == "slide_deck"
    )
    assert "slide:title" in slide_build["rebuilt_unit_ids"]


def test_full_slide_quality_report_persists_in_content_and_spec(tmp_path):
    course = advanced_python_course(4)
    document = document_from_legacy_course(course)
    root = tmp_path / "quality-refresh"
    repository = TeachingRepresentationRepository(root)
    compile_core_representations(document, course, repository)

    refreshed = TeachingRepresentationRepository(root).load(document.course_id)
    representation = next(
        item for item in refreshed.representations if item.representation_type == "slide_deck"
    )
    spec = next(item for item in refreshed.specs if item.spec_id == representation.spec_id)
    content_report = spec.payload["content"]["quality_report"]
    spec_report = spec.payload["quality_report"]
    required = {"score", "passed", "issues", "blockers", "warnings", "slide_count"}

    assert required <= content_report.keys()
    assert required <= spec_report.keys()
    assert content_report == spec_report
    assert content_report["slide_count"] == len(spec.payload["content"]["slides"])
    assert "quality_report" not in content_report
