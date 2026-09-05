from copy import deepcopy

import pytest
from pptx import Presentation
from pptx.util import Pt

from ppt_layout_execution import compile_teaching_template
from ppt_native_scene import audit_scene, render_scene
from ppt_page_scene import resolve_page_scenes
from ppt_teaching_content import PageTeachingV2, validate_source_bindings


def comparison_fixture():
    text = "在相同任务条件下，串行按顺序完成工作，并行同时处理独立工作。比较执行方式：串行逐项执行，并行同时执行。"
    source = {"block_id": "b", "block_revision": "r1", "start": 0, "end": len(text), "quote": text}

    def element(key, text, **kwargs):
        return {"element_id": key, "text": text, "sources": [source], **kwargs}

    elements = [
        element("condition", "相同任务条件", role="condition"),
        element("serial", "串行", subject_id="a", role="label"),
        element("parallel", "并行", subject_id="b", role="label"),
        element("dimension", "执行方式", dimension_id="mode", role="label"),
        element("a-mode", "逐项执行", subject_id="a", dimension_id="mode"),
        element("b-mode", "同时执行", subject_id="b", dimension_id="mode"),
    ]
    ids = [e["element_id"] for e in elements]
    return {
        "elements": elements,
        "expression": {"kind": "comparison", "subjects": [{"subject_id": "a", "label_element_id": "serial"}, {"subject_id": "b", "label_element_id": "parallel"}],
            "dimensions": [{"dimension_id": "mode", "label_element_id": "dimension"}],
            "cells": [{"subject_id": "a", "dimension_id": "mode", "element_ids": ["a-mode"]}, {"subject_id": "b", "dimension_id": "mode", "element_ids": ["b-mode"]}],
            "condition_element_ids": ["condition"]},
        "must_show": ids,
        "source_dispositions": [{"block_id": "b", "purpose": "screen", "element_ids": ids, "reason": "保留比较对象、共同条件和执行差异"}],
        "states": [{"state_id": "all", "visible_element_ids": ids, "teaching_note": "比较两种执行方式"}],
    }, {"b": ("r1", text)}


def branch_fixture():
    value, notes = comparison_fixture()
    source = value["elements"][0]["sources"]
    value["elements"] = [{"element_id": key, "text": label, "sources": source} for key, label in [("a", "任务"), ("b", "串行"), ("c", "并行")]]
    value["expression"] = {"kind": "concept", "node_element_ids": ["a", "b", "c"], "relations": [
        {"relation_id": "ab", "source_id": "a", "target_id": "b", "kind": "association", "sources": source},
        {"relation_id": "ac", "source_id": "a", "target_id": "c", "kind": "association", "sources": source}]}
    value["must_show"] = ["a", "b", "c"]
    value["source_dispositions"][0]["element_ids"] = ["a", "b", "c"]
    value["states"][0]["visible_element_ids"] = ["a", "b", "c"]
    return value, notes


def scene_for(value, slug="compare-matrix"):
    template = compile_teaching_template("academic-editorial", certification_required=False)
    return resolve_page_scenes(page_id="p1", title="比较执行方式", content=PageTeachingV2.model_validate(value),
        layout=template.get_layout(template.layout_id(slug)), template=template, source_document_revision="doc1")[0]


def test_comparison_retains_notes_and_aligned_objects():
    value, notes = comparison_fixture()
    content = PageTeachingV2.model_validate(value)
    validate_source_bindings(content, notes)
    scene = scene_for(value)
    left = next(o for o in scene.objects if o.element_id == "a-mode")
    right = next(o for o in scene.objects if o.element_id == "b-mode")
    assert left.y == right.y and left.x < right.x
    assert left.subject_id == "a" and right.subject_id == "b"


def test_comparison_allocates_matrix_rows_without_shrinking_or_moving_reveal_states():
    from ppt_layout_samples import matrix_boundary_sample
    value = matrix_boundary_sample(2, 3).model_dump()
    for element in value["elements"]:
        if element["element_id"] in {"cell-0-0", "cell-1-0"}:
            element.update(kind="formula", text=r"\begin{bmatrix}1&2&3\\0&4&5\\0&0&6\end{bmatrix}")
    scene = scene_for(value)
    objects = {obj.element_id: obj for obj in scene.objects}
    assert objects["cell-0-0"].height > objects["cell-0-1"].height
    assert len(objects["cell-0-0"].lines) == 3
    assert objects["cell-0-0"].y == objects["cell-1-0"].y
    assert objects["cell-0-0"].font_size == 20
    value["states"].insert(0, {**value["states"][0], "state_id": "before", "visible_element_ids": [
        key for key in value["states"][0]["visible_element_ids"] if key != "cell-1-0"]})
    earlier = scene_for(value)
    assert next(obj for obj in earlier.objects if obj.element_id == "cell-0-0") == objects["cell-0-0"]


def test_intentionally_ragged_matrix_counterexample_preserves_missing_entry():
    from ppt_formula_projection import project_matrix
    source = r"\begin{bmatrix}2&3&-1&5\\1&-1&0&2\\4&5&-1\end{bmatrix}"
    projected = project_matrix(source)
    assert len(projected.splitlines()) == 3
    assert projected.count("0") == 1
    assert "".join(c for c in projected.splitlines()[-1] if not c.isspace()) == "⎣45-1⎦"


def test_inline_math_in_exact_quote_is_projected_without_changing_source():
    from ppt_page_scene import display_element_text
    from ppt_teaching_content import ScreenElement
    original = r"将 $[1, 0, 4 \mid -1]$ 与 $[b_1, \dots, b_m]^T$ 对照。"
    element = ScreenElement(element_id="quote", kind="quote", text=original,
        sources=[{"block_id": "b", "block_revision": "r", "start": 0, "end": len(original), "quote": original}])
    displayed = display_element_text(element)
    assert "$" not in displayed and "\\" not in displayed
    assert "∣" in displayed and "…" in displayed
    assert element.text == element.sources[0].quote == original
    element.kind = "code"
    assert display_element_text(element) == original


def test_math_projection_preserves_prose_and_refuses_unknown_commands():
    from ppt_formula_projection import project_prose_math
    assert project_prose_math("集合 {甲, 乙}，价格 $5 和 $10") == "集合 {甲, 乙}，价格 $5 和 $10"
    assert project_prose_math(r"文件位于 C:\notes\lesson.txt") == r"文件位于 C:\notes\lesson.txt"
    with pytest.raises(ValueError, match="teaching_formula_not_supported"):
        project_prose_math(r"观察 $\unknown{x}$")


@pytest.mark.parametrize("mutation,code", [
    (lambda v: v["expression"]["cells"].pop(), "comparison_matrix_incomplete"),
    (lambda v: v["expression"]["cells"][0].update(element_ids=["b-mode"]), "comparison_cell_identity_mismatch"),
    (lambda v: v["expression"]["subjects"][0].update(label_element_id="parallel"), "comparison_subject_binding_mismatch"),
    (lambda v: v["states"][0]["visible_element_ids"].remove("condition"), "screen_element_never_visible"),
])
def test_invalid_comparisons_are_rejected(mutation, code):
    value, _ = comparison_fixture()
    mutation(value)
    with pytest.raises(ValueError, match=code):
        PageTeachingV2.model_validate(value)


def test_stale_source_and_selected_quote_must_match_exactly():
    value, notes = comparison_fixture()
    notes["b"] = ("r2", notes["b"][1])
    with pytest.raises(ValueError, match="source_revision_stale"):
        validate_source_bindings(PageTeachingV2.model_validate(value), notes)
    value["elements"][4]["kind"] = "quote"
    notes["b"] = ("r1", notes["b"][1])
    with pytest.raises(ValueError, match="selected_artifact_not_exact"):
        validate_source_bindings(PageTeachingV2.model_validate(value), notes)


def test_measured_long_chinese_fails_without_shrinking():
    value, _ = comparison_fixture()
    value["elements"][4]["text"] = "这是一段不能通过缩小字体塞入单元格的很长中文。" * 30
    with pytest.raises(ValueError, match="teaching_text_capacity_exceeded"):
        scene_for(value)


def test_real_native_readback_rejects_swapped_cells(tmp_path):
    value, _ = comparison_fixture()
    scene = scene_for(value)
    p = Presentation()
    p.slide_width, p.slide_height = Pt(960), Pt(540)
    slide = p.slides.add_slide(p.slide_layouts[6])
    render_scene(slide, scene)
    path = tmp_path / "comparison.pptx"
    p.save(path)
    slide = Presentation(path).slides[0]
    assert audit_scene(slide, scene)["passed"]
    shapes = {s.name: s for s in slide.shapes}
    shapes["teaching:a-mode"].left, shapes["teaching:b-mode"].left = shapes["teaching:b-mode"].left, shapes["teaching:a-mode"].left
    with pytest.raises(ValueError, match="export_element_position_mismatch"):
        audit_scene(slide, scene)


def test_branch_topology_survives_native_export_and_rejects_chain(tmp_path):
    value, _ = branch_fixture()
    scene = scene_for(value, "concept-map")
    p = Presentation()
    p.slide_width, p.slide_height = Pt(960), Pt(540)
    slide = p.slides.add_slide(p.slide_layouts[6])
    render_scene(slide, scene)
    path = tmp_path / "branch.pptx"
    p.save(path)
    slide = Presentation(path).slides[0]
    assert audit_scene(slide, scene)["relations"] == 2
    shapes = {s.name: s for s in slide.shapes}
    shapes["relation:ac"]._element.xpath(".//a:stCxn")[0].set("id", str(shapes["teaching:b"].shape_id))
    with pytest.raises(ValueError, match="export_relation_endpoint_mismatch"):
        audit_scene(slide, scene)


def test_old_template_serialization_does_not_gain_new_hash_fields():
    from template_layout_contract import compile_builtin_template_layout_contract_v1
    pack = compile_builtin_template_layout_contract_v1("academic-editorial")
    assert all("execution" not in l.model_dump() for l in pack.layouts)


def compiled_manuscript(value=None):
    from course_document import CourseBlock, CourseDocument, CourseSection
    from course_presentation_graph import compile_course_presentation_graph
    from ppt_teaching_manuscript import compile_teaching_manuscript
    content, notes = comparison_fixture()
    if value is not None:
        content = value
    doc = CourseDocument(course_id="fixture", title="执行方式", document_revision="doc1",
        sections=[CourseSection(section_id="s", title="执行方式", position=0)],
        blocks=[CourseBlock(block_id="b", section_id="s", position=0, payload={"markdown": notes["b"][1]}, internal_revision="r1")])
    graph = compile_course_presentation_graph(doc, teaching_plan={})
    template = compile_teaching_template("academic-editorial", certification_required=False)
    manuscript = compile_teaching_manuscript(doc, graph, template, {"central_question": "比较执行方式"}, [{
        "page_id": "p1", "teaching_unit_id": graph.units[0].teaching_unit_id,
        "source_block_ids": ["b"], "title": "比较执行方式", "layout_id": template.layout_id("compare-matrix"),
        "page_goal": "理解串行与并行的区别", "primary_claim": "串行逐项执行，并行同时执行", "teaching": content}])
    return doc, graph, template, manuscript


def test_confirmed_manuscript_exports_repeatedly_without_models(tmp_path):
    from slide_deck_v6 import compile_slide_deck_v6_from_manuscript, validate_deck_matches_ppt_manuscript_v1
    from slide_deck_v6_renderer import export_slide_deck_v6_pptx
    doc, graph, template, manuscript = compiled_manuscript()
    first = compile_slide_deck_v6_from_manuscript(doc, graph, manuscript, template)
    second = compile_slide_deck_v6_from_manuscript(doc, graph, manuscript, template)
    assert first.model_dump() == second.model_dump()
    assert validate_deck_matches_ppt_manuscript_v1(first, manuscript)
    export_slide_deck_v6_pptx(first, tmp_path / "first.pptx")
    export_slide_deck_v6_pptx(second, tmp_path / "second.pptx")
    a, b = Presentation(tmp_path / "first.pptx"), Presentation(tmp_path / "second.pptx")
    assert [s.text for s in a.slides[0].shapes] == [s.text for s in b.slides[0].shapes]
    assert a.slides[0].notes_slide.notes_text_frame.text == b.slides[0].notes_slide.notes_text_frame.text


def test_orchestrator_final_generation_never_invokes_planners(tmp_path, monkeypatch):
    import asyncio
    from slide_deck_v6_orchestrator import SlideDeckV6Orchestrator, SlideDeckV6CandidateRepository
    from teaching_representations import TeachingRepresentationRepository
    doc, _, template, manuscript = compiled_manuscript()
    async def forbidden(_request):
        pytest.fail("final generation called a planning model")
    # Rendering pixels is covered separately; keep this test focused on the
    # orchestration/confirmation boundary and real native-object export.
    monkeypatch.setattr("slide_deck_v6_orchestrator.audit_exported_pptx", lambda *a, **kw: {"passed": True, "blockers": []})
    orchestrator = SlideDeckV6Orchestrator(representation_repository=TeachingRepresentationRepository(tmp_path / "representations"),
        candidate_repository=SlideDeckV6CandidateRepository(tmp_path / "candidates"), progress_root=tmp_path / "progress")
    result = asyncio.run(orchestrator.build(task_id="final-v2", document=doc, course_data={}, mode="teaching", theme="academic-editorial",
        story_planner=forbidden, visual_planner=forbidden, source_revision_provider=lambda: doc.document_revision,
        template_contract=template, confirmed_manuscript=manuscript, publish_result=False))
    assert result["status"] == "v6_ready"


def test_render_failure_preserves_last_ppt_and_cleans_temporary_file(tmp_path, monkeypatch):
    from types import SimpleNamespace
    from ppt_teaching_manuscript import physical_pages
    from ppt_native_scene import render_teaching_deck
    _, _, _, manuscript = compiled_manuscript()
    target = tmp_path / 'last-good.pptx'
    target.write_bytes(b'last-known-good')
    def fail(*args, **kwargs):
        raise ValueError('simulated_render_failure')
    monkeypatch.setattr('ppt_render_audit.render_evidence', fail)
    with pytest.raises(ValueError, match='simulated_render_failure'):
        render_teaching_deck(SimpleNamespace(pages=physical_pages(manuscript)), target)
    assert target.read_bytes() == b'last-known-good'
    assert list(tmp_path.glob('*.pptx')) == [target]
    assert not list(tmp_path.glob('.last-good-*'))


def test_teacher_edit_and_targeted_repair_preserve_other_pages(monkeypatch):
    import asyncio
    from ppt_teaching_manuscript import refresh_manuscript
    from slide_ai_planning_v6 import regenerate_ppt_manuscript_pages_v1
    from slide_deck_v6 import revise_ppt_manuscript_v1, V6BuildError
    _, _, template, manuscript = compiled_manuscript()
    monkeypatch.setattr("ppt_teaching_manuscript.template_for_manuscript", lambda m: template)
    other = manuscript.pages[0].model_copy(deep=True, update={"page_id": "p2", "page_number": 2, "teacher_locked": True})
    manuscript.pages.append(other)
    manuscript = refresh_manuscript(manuscript)
    before = manuscript.model_dump()
    calls = []
    async def planner(request):
        calls.append(request)
        page = request["current_page"]
        content = deepcopy(page["teaching"])
        # First response misses a dimension cell; second repairs that page.
        if len(calls) == 1:
            content["expression"]["cells"].pop()
        return {"title": "串行与并行", "page_goal": page["page_goal"], "primary_claim": page["primary_claim"], "teaching": content}
    revised = asyncio.run(regenerate_ppt_manuscript_pages_v1(manuscript, target_page_ids=["p1"], ai_planner=planner))
    assert len(calls) == 2 and "comparison_matrix_incomplete" in calls[1]["validation_error"]
    assert len(calls[1]["previous_candidate"]["teaching"]["expression"]["cells"]) == 1
    assert revised.pages[0].title == "串行与并行"
    assert revised.pages[1].model_dump() == before["pages"][1]
    assert manuscript.model_dump() == before
    assert revised.manuscript_revision != manuscript.manuscript_revision
    with pytest.raises(V6BuildError, match="target_locked"):
        asyncio.run(regenerate_ppt_manuscript_pages_v1(manuscript, target_page_ids=["p2"], ai_planner=planner))
    assert len(calls) == 2
    locked = revise_ppt_manuscript_v1(revised, [{"page_id": "p1", "teacher_locked": True}])
    assert locked.pages[0].teacher_locked and locked.manuscript_revision != revised.manuscript_revision


def test_page_capacity_failure_is_locally_repaired(monkeypatch):
    import asyncio
    from ppt_teaching_planner import plan_teaching_manuscript
    doc, graph, template, manuscript = compiled_manuscript()
    responses = []
    value = manuscript.pages[0].teaching.model_dump(mode="json")
    plan = {"page_id": "p1", "teaching_unit_id": graph.units[0].teaching_unit_id, "source_block_ids": ["b"],
            "title": "执行方式", "layout_id": template.layout_id("compare-matrix"),
            "page_goal": "理解执行方式", "primary_claim": "串行逐项执行，并行同时执行"}
    async def planner(request):
        responses.append(request)
        if request["teaching_request"] == "narrative":
            return {"pacing": {"max_physical_pages": 12, "rationale": "保留比较与必要推理停顿"}, "narrative_brief": {"central_question": "比较执行方式"}, "pages": [plan]}
        content = deepcopy(value)
        if len(responses) == 2:
            content["elements"][4]["text"] = "逐项执行" * 150
        return {"title": plan["title"], "page_goal": plan["page_goal"], "primary_claim": plan["primary_claim"], "teaching": content}
    result, checkpoint = asyncio.run(plan_teaching_manuscript(doc, graph, template, planner))
    assert len(responses) == 3
    assert "capacity" in responses[-1]["validation_error"]
    assert result.pages[0].teaching.model_dump(exclude={"presentation"}) == manuscript.pages[0].teaching.model_dump(exclude={"presentation"})
    assert checkpoint["pages"]["p1"]


def test_compact_comparison_preserves_identity_sources_and_reveal_order():
    from ppt_comparison_draft import lower_comparison_draft
    value, notes = comparison_fixture()
    quote = {"block_id": "b", "quote": notes["b"][1]}
    def item(text, **fields):
        return {"text": text, "sources": [quote], **fields}
    draft = {"title": "执行方式", "page_goal": "辨别串行与并行", "audience_question": "怎样执行？", "expected_response": "逐项执行或同时执行",
             "conditions": [item("相同任务条件")], "subjects": [item("串行", key="serial"), item("并行", key="parallel")],
             "dimensions": [item("执行方式", key="mode")],
             "cells": [{"subject_key": "serial", "dimension_key": "mode", "content": [item("逐项执行", show_from=2, role="answer")]},
                       {"subject_key": "parallel", "dimension_key": "mode", "content": [item("同时执行", show_from=2, role="answer")]}],
             "screen_question": item("怎样执行？"), "reveal_notes": ["先判断", "再核对答案"]}
    sources = {"b": {"block_id": "b", "block_revision": "r1", "full_text": notes["b"][1]}}
    lowered = lower_comparison_draft(draft, sources)
    teaching = PageTeachingV2.model_validate(lowered["teaching"])
    validate_source_bindings(teaching, notes)
    assert teaching.expression.cells[0].subject_id == "subject-0"
    assert not any(e.element_id in teaching.states[0].visible_element_ids for e in teaching.elements if e.role == "answer")
    broken = deepcopy(draft)
    broken["cells"].pop()
    with pytest.raises(ValueError):
        lower_comparison_draft(broken, sources)
    broken = deepcopy(draft)
    broken["cells"][0]["content"][0]["show_from"] = 1
    with pytest.raises(ValueError, match="answer_revealed_before_question"):
        lower_comparison_draft(broken, sources)


def test_manuscript_rejects_unadopted_image_before_asset_lookup():
    from slide_deck_v6 import V6BuildError
    value, _ = comparison_fixture()
    value["elements"][4].update(kind="image", asset_id="sva_other_course_asset", asset_digest="untrusted")
    with pytest.raises(V6BuildError, match="teaching_asset_not_adopted"):
        compiled_manuscript(value)
