from __future__ import annotations

import json
from pathlib import Path
from types import SimpleNamespace

import pytest
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from course_document import document_from_legacy_course
from slide_asset_repository import SlideAssetRepository, finalize_visual_assets
from slide_deck import SlideSpec, _plain_math_text
from slide_deck_renderer import (
    _display_heading,
    _format_formula_text,
    export_structured_slide_deck,
)
from slide_deck_v3 import (
    compile_slide_deck_v3,
    deterministic_slide_allocation,
    fragment_course_document,
)
from slide_visuals import (
    SlideVisualPlanV1,
    VisualAnchorV1,
    _semantic_relation_spec,
    _source_clauses,
    _visual_anchor,
    _visual_plan_batches,
    _visual_plan_request,
    deterministic_visual_plan,
    plan_slide_visuals,
    rebalance_visual_plan_pages,
    validate_visual_plan,
    visual_integrity_issues,
)
from teaching_storyboard import build_teaching_storyboard


def visual_course() -> dict:
    return {
        "course_id": "visual-course",
        "course_name": "线性映射：结构与应用",
        "nodes": [
            {
                "node_id": "chapter-1",
                "parent_node_id": "root",
                "node_name": "第一章 线性映射",
                "node_level": 1,
                "content_blocks": [
                    {
                        "block_id": "concept",
                        "title": "线性映射保持两类运算结构",
                        "content": (
                            "线性映射同时保持向量加法与数乘。"
                            "\n\n- 先验证加法保持性"
                            "\n- 再验证数乘保持性"
                            "\n- 最后检查零向量"
                        ),
                        "metadata": {"role": "concept"},
                    },
                    {
                        "block_id": "formula",
                        "title": "定义式",
                        "content": "$$T(au+bv)=aT(u)+bT(v)$$",
                        "metadata": {"role": "reasoning", "kind": "formula"},
                    },
                    {
                        "block_id": "example",
                        "title": "平面旋转是线性映射",
                        "content": "旋转把每个向量映射到同角度的新方向，并保持向量组合关系。",
                        "metadata": {"role": "example"},
                    },
                    {
                        "block_id": "coordinate-example",
                        "title": "二维向量的旋转结果",
                        "content": "向量坐标 (1, 2) 经过旋转后映射为 (-2, 1)。",
                        "metadata": {"role": "example"},
                    },
                    {
                        "block_id": "check",
                        "title": "判断练习",
                        "content": "平移为什么通常不是线性映射？",
                        "metadata": {"role": "checkpoint"},
                    },
                ],
            }
        ],
    }


def test_compiler_adds_grounded_visual_director_plan() -> None:
    course = visual_course()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    allocation = deterministic_slide_allocation(
        document,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )

    content = compile_slide_deck_v3(
        document,
        course,
        mode="teaching",
        theme="qizhi-classroom",
        allocation_plan=allocation,
    )

    assert content["visual_plan"]["schema_version"] == "slide_visual_plan_v1"
    assert content["build_signature"]["visual_policy_version"]
    assert content["visual_quality_report"]["passed"] is True
    assert content["visual_quality_report"]["effective_visual_coverage_ratio"] >= 0.60
    assert all(slide["teaching_job"] for slide in content["slides"])
    assert all(slide["takeaway"] for slide in content["slides"])


def test_classification_page_does_not_use_next_heading_as_diagram_root() -> None:
    course = visual_course()
    concept = course["nodes"][0]["content_blocks"][0]
    concept["content"] = (
        "#### 核心概念与背景\n\n"
        "系统是研究对象，环境是系统以外的部分。\n\n"
        "根据系统与环境之间的交互方式，热力学将系统分为三类：\n\n"
        "- 孤立系统：既不交换物质，也不交换能量。\n"
        "- 封闭系统：不交换物质，但可以交换能量。\n"
        "- 开放系统：既可以交换物质，也可以交换能量。\n\n"
        "这些分类帮助我们理解不同条件下热力学行为的变化，是后续建立模型和分析的基础。\n\n"
        "#### 深度原理/底层机制\n\n"
        "系统边界决定了物质和能量能否通过。"
    )
    course["nodes"][0]["content_blocks"] = [concept]
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    allocation = deterministic_slide_allocation(
        document,
        fragments,
        mode="full",
        theme="qizhi-classroom",
    )
    fragment_by_id = {
        fragment.fragment_id: fragment
        for fragment in fragments
    }
    classification_page = next(
        page
        for page in allocation.pages
        if {
            "孤立系统：既不交换物质，也不交换能量。",
            "封闭系统：不交换物质，但可以交换能量。",
            "开放系统：既可以交换物质，也可以交换能量。",
        } <= {
            fragment_by_id[fragment_id].text
            for fragment_id in page.fragment_ids
        }
    )
    classification_text = {
        fragment_by_id[fragment_id].text
        for fragment_id in classification_page.fragment_ids
    }

    plan = deterministic_visual_plan(document, allocation, fragments)
    visual_page = next(
        page for page in plan.pages
        if page.page_id == classification_page.page_id
    )

    assert "深度原理/底层机制" not in classification_text
    assert visual_page.visual_anchor.kind == "none"
    assert visual_page.composition == "statement"


def test_subject_profile_drives_a_source_bound_hierarchy_without_heading_guessing() -> None:
    course = visual_course()
    course["nodes"][0]["content_blocks"] = [{
        "block_id": "anatomy-layers",
        "title": "课程内容",
        "content": (
            "人体局部结构按照由浅入深的空间关系组织，观察时需要逐层确认边界。\n\n"
            "- 皮肤\n"
            "- 浅筋膜\n"
            "- 深筋膜"
        ),
        "metadata": {
            "role": "concept",
            "module_id": "life_location_structure",
            "module_instance_id": "life-location-1",
            "composition_style": "balanced",
        },
    }]
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    allocation = deterministic_slide_allocation(
        document,
        fragments,
        mode="full",
        theme="qizhi-classroom",
    )

    plan = deterministic_visual_plan(document, allocation, fragments)
    page = next(
        item
        for item in plan.pages
        if any(
            fragment.text == "浅筋膜"
            for fragment in fragments
            if fragment.fragment_id
            in next(
                allocation_page.fragment_ids
                for allocation_page in allocation.pages
                if allocation_page.page_id == item.page_id
            )
        )
    )

    assert page.visual_anchor.kind == "relational_diagram"
    assert page.visual_anchor.parameters["diagram_type"] == "hierarchy"
    assert [node.label for node in page.visual_anchor.nodes][1:] == [
        "皮肤",
        "浅筋膜",
        "深筋膜",
    ]


def test_display_heading_prefers_local_title_and_complete_short_phrase() -> None:
    explicit_heading = SlideSpec(
        unit_id="explicit-heading",
        position=0,
        layout="concept",
        slide_purpose="reasoning",
        title="🔍 深度原理/底层机制",
        takeaway=(
            "热力学系统的核心在于其“边界”及其对物质与能量的控制能力。"
        ),
    )
    classification_heading = SlideSpec(
        unit_id="classification-heading",
        position=0,
        layout="concept",
        slide_purpose="concept",
        title="根据系统与环境之间的交互方式，热力学将系统分为三类",
        takeaway="根据系统与环境之间的交互方式，热力学将系统分为三类",
    )

    assert _display_heading(explicit_heading) == "🔍 深度原理/底层机制"
    assert _display_heading(classification_heading) == (
        "根据系统与环境之间的交互方式，热力学将系统分为三类"
    )


def test_dense_prose_relation_is_suppressed_when_it_repeats_the_source() -> None:
    course = visual_course()
    concept = course["nodes"][0]["content_blocks"][0]
    concept["content"] = (
        "#### 🔍 深度原理/底层机制\n\n"
        "热力学系统的核心在于其“边界”及其对物质与能量的控制能力。"
        "边界决定了系统是否能与环境发生互动。"
        "例如，在一个密封保温杯中的水就是一个封闭系统——"
        "水不会流出，但可以通过杯子壁传递热量。\n\n"
        "从微观角度看，系统的状态由大量粒子的运动构成。"
        "热力学通过宏观变量来描述这些微观行为的统计结果。"
        "因此，系统分类不仅影响理论建模，也直接影响实验设计。"
    )
    course["nodes"][0]["content_blocks"] = [concept]
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    allocation = deterministic_slide_allocation(
        document,
        fragments,
        mode="full",
        theme="qizhi-classroom",
    )
    page = next(
        page
        for page in allocation.pages
        if any(
            fragment.fragment_id in page.fragment_ids
            and fragment.kind == "heading"
            and "深度原理" in fragment.text
            for fragment in fragments
        )
    )

    plan = deterministic_visual_plan(document, allocation, fragments)
    visual_page = next(
        item for item in plan.pages
        if item.page_id == page.page_id
    )

    assert visual_page.visual_anchor.kind == "none"
    assert visual_page.composition == "statement"


def test_relation_diagram_requires_a_local_bounded_relationship() -> None:
    fragments = [
        SimpleNamespace(
            fragment_id="f1",
            kind="paragraph",
            text="系统边界用于区分研究对象和环境。",
        ),
        SimpleNamespace(
            fragment_id="f2",
            kind="paragraph",
            text="实验记录用于检查测量结果。",
        ),
        SimpleNamespace(
            fragment_id="f3",
            kind="paragraph",
            text="但是本节还需要讨论单位换算。",
        ),
    ]
    clauses = [
        (fragment.text, fragment.fragment_id)
        for fragment in fragments
    ]

    relation = _semantic_relation_spec(
        SimpleNamespace(narrative_role="concept"),
        fragments,
        clauses,
        [],
    )

    assert relation is None


def test_relation_connector_inside_parentheses_cannot_split_diagram_nodes() -> None:
    source = (
        "但是，空调在这个过程中消耗的电能（转化为机械功）和排放的热量，"
        "会因控制策略不同而不同。"
    )
    fragments = [SimpleNamespace(
        fragment_id="air-conditioner",
        kind="paragraph",
        text=source,
    )]

    clauses = _source_clauses(fragments)
    relation = _semantic_relation_spec(
        SimpleNamespace(narrative_role="example"),
        fragments,
        clauses,
        [],
    )

    assert clauses == [(source.rstrip("。"), "air-conditioner")]
    assert relation is None

    compact_source = "电能（转化为机械功）"
    compact_fragments = [SimpleNamespace(
        fragment_id="compact-parenthetical",
        kind="paragraph",
        text=compact_source,
    )]
    assert _semantic_relation_spec(
        SimpleNamespace(narrative_role="example"),
        compact_fragments,
        _source_clauses(compact_fragments),
        [],
    ) is None


def test_visual_anchor_rejects_unbalanced_relation_labels() -> None:
    with pytest.raises(ValueError, match="balanced punctuation"):
        VisualAnchorV1(
            visual_id="malformed-relation",
            kind="relational_diagram",
            purpose="process",
            source_fragment_ids=["coordinate-example"],
            alt_text="错误截断的映射关系",
            nodes=[
                {
                    "node_id": "input",
                    "label": "向量坐标 (1, 2",
                    "source_fragment_ids": ["coordinate-example"],
                },
                {
                    "node_id": "output",
                    "label": "映射为 (-2, 1",
                    "source_fragment_ids": ["coordinate-example"],
                },
            ],
            edges=[{
                "source": "input",
                "target": "output",
                "relation": "maps_to",
            }],
            parameters={"relation_evidence": "explicit_mapping_connector"},
        )


def test_plain_math_text_keeps_degree_symbol() -> None:
    assert _plain_math_text(r"$30^\circ \text{C}$") == "30° C"


def test_template_heading_cannot_become_a_diagram_node() -> None:
    fragments = [
        SimpleNamespace(
            fragment_id="heading",
            kind="heading",
            text="核心概念与背景",
        ),
        SimpleNamespace(
            fragment_id="item-a",
            kind="list_item",
            text="孤立系统不交换物质和能量。",
        ),
        SimpleNamespace(
            fragment_id="item-b",
            kind="list_item",
            text="封闭系统可以交换能量。",
        ),
    ]
    list_clauses = [
        (fragment.text, fragment.fragment_id)
        for fragment in fragments[1:]
    ]

    relation = _semantic_relation_spec(
        SimpleNamespace(narrative_role="concept"),
        fragments,
        [(fragment.text, fragment.fragment_id) for fragment in fragments],
        list_clauses,
    )

    assert relation is None


def test_deterministic_director_uses_source_bound_visual_variety() -> None:
    course = visual_course()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    allocation = deterministic_slide_allocation(
        document,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )

    plan = deterministic_visual_plan(document, allocation, fragments)
    visual_kinds = {
        page.visual_anchor.kind
        for page in plan.pages
        if page.visual_anchor.kind != "none"
    }

    assert {
        "coordinate_plot",
        "formula",
        "generated_illustration",
    } <= visual_kinds
    assert all("**" not in page.takeaway and "$$" not in page.takeaway for page in plan.pages)
    coordinate_pages = [
        page for page in plan.pages
        if page.visual_anchor.kind == "coordinate_plot"
    ]
    assert coordinate_pages
    assert all(page.visual_anchor.parameters.get("points") for page in coordinate_pages)
    assert all(
        not page.visual_anchor.parameters.get("not_to_scale")
        for page in coordinate_pages
    )
    assert all(
        page.visual_anchor.kind != "coordinate_plot"
        for page in plan.pages
        if "旋转把每个向量映射到同角度" in page.takeaway
    )
    formula_page = next(
        page for page in plan.pages
        if page.visual_anchor.kind == "formula"
    )
    assert formula_page.visual_anchor.parameters["formula"].startswith("$$")
    relation_pages = [
        page for page in plan.pages
        if page.visual_anchor.kind == "relational_diagram"
    ]
    assert all(len(page.visual_anchor.nodes) >= 2 for page in relation_pages)
    assert all(page.visual_anchor.edges for page in relation_pages)
    assert all(
        page.visual_anchor.parameters.get("relation_evidence")
        for page in relation_pages
    )


def test_mermaid_fragment_compiles_to_source_bound_rule_diagram() -> None:
    course = {
        "course_id": "rule-diagram-course",
        "course_name": "System classification",
        "nodes": [{
            "node_id": "chapter-system",
            "parent_node_id": "root",
            "node_name": "System classification",
            "node_level": 1,
            "content_blocks": [{
                "block_id": "system-flow",
                "title": "Closed system",
                "content": (
                    "```mermaid\n"
                    "graph TD\n"
                    "A[Closed system] -->|cannot exchange matter| B[Environment]\n"
                    "```"
                ),
                "metadata": {"role": "concept"},
            }],
        }],
    }
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    allocation = deterministic_slide_allocation(
        document,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )

    plan = deterministic_visual_plan(document, allocation, fragments)
    anchor = next(
        page.visual_anchor
        for page in plan.pages
        if page.visual_anchor.kind == "rule_diagram"
    )

    assert anchor.parameters["template"] == "process_flow"
    assert anchor.parameters["relation_evidence"]
    assert [node.label for node in anchor.nodes] == [
        "Closed system",
        "Environment",
    ]
    assert anchor.edges[0].label == "cannot exchange matter"
    assert all(node.source_fragment_ids for node in anchor.nodes)


def test_rule_diagram_exports_as_editable_ppt_shapes(tmp_path: Path) -> None:
    course = {
        "course_id": "editable-rule-diagram",
        "course_name": "System classification",
        "nodes": [{
            "node_id": "chapter-system",
            "parent_node_id": "root",
            "node_name": "System classification",
            "node_level": 1,
            "content_blocks": [{
                "block_id": "system-flow",
                "title": "Closed system",
                "content": (
                    "A closed system cannot exchange matter with its environment."
                    "\n\n```mermaid\n"
                    "flowchart LR\n"
                    "A[Closed system] -->|cannot exchange matter| B[Environment]\n"
                    "```"
                ),
                "metadata": {"role": "concept"},
            }],
        }],
    }
    document = document_from_legacy_course(course)
    content = compile_slide_deck_v3(
        document,
        course,
        mode="teaching",
        theme="qizhi-classroom",
    )
    output = export_structured_slide_deck(
        content,
        tmp_path / "rule-diagram.pptx",
    )
    presentation = Presentation(output)
    slide_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )

    assert "Closed system" in slide_text
    assert "Environment" in slide_text
    assert "cannot exchange matter" in slide_text
    assert "flowchart LR" not in slide_text
    assert any(
        shape.shape_type == MSO_SHAPE_TYPE.LINE
        for slide in presentation.slides
        for shape in slide.shapes
    )


def test_unsupported_mermaid_degrades_without_raw_source_or_placeholder() -> None:
    course = {
        "course_id": "unsupported-rule-diagram",
        "course_name": "Interaction sequence",
        "nodes": [{
            "node_id": "chapter-sequence",
            "parent_node_id": "root",
            "node_name": "Interaction sequence",
            "node_level": 1,
            "content_blocks": [{
                "block_id": "sequence",
                "title": "Unsupported sequence",
                "content": (
                    "The learner should focus on the request and response."
                    "\n\n```mermaid\n"
                    "sequenceDiagram\n"
                    "Client->>Server: Request\n"
                    "Server-->>Client: Response\n"
                    "```"
                ),
                "metadata": {"role": "remediation"},
            }],
        }],
    }
    document = document_from_legacy_course(course)
    content = compile_slide_deck_v3(
        document,
        course,
        mode="teaching",
        theme="qizhi-classroom",
    )
    visible_text = "\n".join(
        str(block.get("content") or "")
        for slide in content["slides"]
        for block in slide["blocks"]
    )

    assert "sequenceDiagram" not in visible_text
    assert "Client->>Server" not in visible_text
    assert not any(
        visual["kind"] in {"code", "generated_illustration"}
        for slide in content["slides"]
        for visual in slide["visuals"]
    )
    assert not any(
        issue["code"] == "raw_mermaid_visible"
        for issue in content["quality_report"]["issues"]
    )


def test_visual_integrity_rejects_visible_raw_mermaid() -> None:
    issues = visual_integrity_issues({
        "fragment_manifest": [],
        "visual_asset_manifest": [],
        "slides": [{
            "unit_id": "slide-raw-mermaid",
            "quality": {"fragment_ids": []},
            "blocks": [{
                "kind": "code",
                "content": "graph TD\nA[Closed system] --> B[Environment]",
            }],
            "visuals": [],
        }],
    })

    assert any(issue["code"] == "raw_mermaid_visible" for issue in issues)


def test_hierarchy_diagram_keeps_every_required_source_sibling() -> None:
    page = SimpleNamespace(
        page_id="system-classification",
        narrative_role="concept",
        layout="concept",
    )
    fragments = [
        SimpleNamespace(
            fragment_id="heading",
            kind="heading",
            text="热力学系统的三种分类",
            source_kind="text",
            asset_refs=[],
        ),
        SimpleNamespace(
            fragment_id="isolated",
            kind="list_item",
            text="孤立系统（Isolated System）：既不与外界交换物质，也不交换能量。",
            source_kind="text",
            asset_refs=[],
        ),
        SimpleNamespace(
            fragment_id="closed",
            kind="list_item",
            text="封闭系统（Closed System）：不与外界交换物质，但可以交换能量（如热量）。",
            source_kind="text",
            asset_refs=[],
        ),
        SimpleNamespace(
            fragment_id="open",
            kind="list_item",
            text="开放系统（Open System）：既可以交换物质，也可以交换能量。",
            source_kind="text",
            asset_refs=[],
        ),
    ]

    visual = _visual_anchor(page, fragments, 0).model_dump(mode="json")
    node_source_ids = {
        source_id
        for node in visual["nodes"]
        for source_id in node["source_fragment_ids"]
    }

    assert visual["kind"] == "relational_diagram"
    assert {"isolated", "closed", "open"} <= node_source_ids
    assert set(visual["parameters"]["required_node_fragment_ids"]) == {
        "heading",
        "isolated",
        "closed",
        "open",
    }
    assert any("封闭系统" in node["label"] for node in visual["nodes"])


def test_visual_integrity_blocks_a_diagram_that_loses_required_sibling() -> None:
    issues = visual_integrity_issues({
        "fragment_manifest": [
            {"fragment_id": "heading", "kind": "heading", "text": "三类系统"},
            {"fragment_id": "isolated", "kind": "list_item", "text": "孤立系统"},
            {"fragment_id": "closed", "kind": "list_item", "text": "封闭系统"},
            {"fragment_id": "open", "kind": "list_item", "text": "开放系统"},
        ],
        "visual_asset_manifest": [],
        "slides": [{
            "unit_id": "system-classification",
            "quality": {
                "fragment_ids": ["heading", "isolated", "closed", "open"],
            },
            "blocks": [],
            "visuals": [{
                "kind": "relational_diagram",
                "source_fragment_ids": ["heading", "isolated", "closed", "open"],
                "nodes": [
                    {"label": "三类系统", "source_fragment_ids": ["heading"]},
                    {"label": "孤立系统", "source_fragment_ids": ["isolated"]},
                    {"label": "开放系统", "source_fragment_ids": ["open"]},
                ],
                "parameters": {
                    "required_node_fragment_ids": [
                        "heading", "isolated", "closed", "open",
                    ],
                },
            }],
        }],
    })

    missing = next(
        issue for issue in issues
        if issue["code"] == "diagram_required_item_missing"
    )
    assert missing["missing_fragment_ids"] == ["closed"]


@pytest.mark.asyncio
async def test_ai_visual_request_exposes_only_safe_rule_diagram_controls() -> None:
    course = visual_course()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    allocation = deterministic_slide_allocation(
        document,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    captured: dict = {}

    async def capture_request(request):
        captured.update(request)
        return deterministic_visual_plan(
            document,
            allocation,
            fragments,
        ).model_dump(mode="json")

    await plan_slide_visuals(
        document,
        allocation,
        fragments,
        ai_planner=capture_request,
    )

    assert captured["allowed_rule_diagram_templates"] == [
        "apparatus",
        "cycle",
        "energy_balance",
        "process_flow",
        "qualitative_plot",
        "relation_graph",
        "system_boundary",
    ]
    assert captured["rules"]["arbitrary_drawing_code_forbidden"] is True
    assert captured["rules"]["uncertain_visual_must_be_none"] is True
    assert captured["rules"]["raster_generation_default"] == "disabled"
    retrieval_schema = captured["optional_visual_search_output"]
    assert retrieval_schema["location"] == "deck_brief.visual_search_requests"
    assert retrieval_schema["maximum_queries_per_page"] == 2
    assert "spatial_relation" in retrieval_schema["visual_intents"]


def test_rebalance_breaks_visual_kind_runs_longer_than_quality_gate_limit() -> None:
    course = visual_course()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    allocation = deterministic_slide_allocation(
        document,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    plan = deterministic_visual_plan(document, allocation, fragments)
    relation = VisualAnchorV1.model_validate({
        "visual_id": "visual-repeated-relation",
        "kind": "relational_diagram",
        "purpose": "structure",
        "source_fragment_ids": ["fragment-a", "fragment-b"],
        "alt_text": "source-bound relationship",
        "nodes": [
            {
                "node_id": "a",
                "label": "A",
                "source_fragment_ids": ["fragment-a"],
            },
            {
                "node_id": "b",
                "label": "B",
                "source_fragment_ids": ["fragment-b"],
            },
        ],
        "edges": [{
            "source": "a",
            "target": "b",
            "relation": "sequence",
        }],
        "parameters": {"relation_evidence": "source_order"},
    })
    allocation_by_id = {page.page_id: page for page in allocation.pages}
    candidates = [
        page
        for page in plan.pages
        if (
            not page.appendix
            and allocation_by_id[page.page_id].fragment_ids
            and allocation_by_id[page.page_id].layout != "section-divider"
        )
    ][:4]
    assert len(candidates) == 4
    for page in candidates:
        page.visual_anchor = relation.model_copy(deep=True)

    rebalance_visual_plan_pages(plan.pages, allocation, fragments)

    run = 0
    previous = ""
    maximum_run = 0
    for page in candidates:
        kind = page.visual_anchor.kind
        run = run + 1 if kind == previous else 1
        previous = kind
        maximum_run = max(maximum_run, run)
    assert maximum_run <= 3


def test_storyboard_groups_pages_into_source_neutral_teaching_episodes() -> None:
    course = visual_course()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    allocation = deterministic_slide_allocation(
        document,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )

    storyboard = build_teaching_storyboard(document, allocation)

    assert storyboard.policy_version.startswith("knowledge_episode_director_")
    assert storyboard.episodes
    episode = storyboard.episodes[0]
    assert episode.learning_question
    assert episode.beats
    assert all(beat.source_fragment_ids for beat in episode.beats)


def test_single_node_diagram_is_rejected_as_fake_visual() -> None:
    with pytest.raises(ValueError, match="at least two nodes"):
        VisualAnchorV1.model_validate({
            "visual_id": "visual-one-node",
            "kind": "relational_diagram",
            "purpose": "structure",
            "source_fragment_ids": ["fragment-1"],
            "alt_text": "single node",
            "nodes": [{
                "node_id": "n1",
                "label": "Only node",
                "source_fragment_ids": ["fragment-1"],
            }],
            "edges": [],
            "parameters": {"relation_evidence": "source_order"},
        })


def test_visual_plan_rejects_unknown_fragment_bindings() -> None:
    course = visual_course()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    allocation = deterministic_slide_allocation(
        document,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    plan = deterministic_visual_plan(document, allocation, fragments)
    raw = plan.model_dump(mode="json")
    page = next(item for item in raw["pages"] if item["visual_anchor"]["kind"] != "none")
    page["visual_anchor"]["source_fragment_ids"] = ["unknown-fragment"]

    with pytest.raises(ValueError, match="unknown fragment"):
        validate_visual_plan(
            SlideVisualPlanV1.model_validate(raw),
            allocation,
            fragments,
        )


def test_visual_plan_takeaway_cannot_add_an_unbound_number() -> None:
    course = visual_course()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    allocation = deterministic_slide_allocation(
        document,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    plan = deterministic_visual_plan(document, allocation, fragments)
    raw = plan.model_dump(mode="json")
    page = next(item for item in raw["pages"] if item["takeaway_source_fragment_ids"])
    page["takeaway"] = "该方法可以把误差降低 99%"

    with pytest.raises(ValueError, match="ungrounded number"):
        validate_visual_plan(
            SlideVisualPlanV1.model_validate(raw),
            allocation,
            fragments,
        )


@pytest.mark.asyncio
async def test_ai_visual_plan_discards_rewritten_body_without_losing_visual() -> None:
    course = visual_course()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    allocation = deterministic_slide_allocation(
        document,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    valid = deterministic_visual_plan(document, allocation, fragments)
    raw = valid.model_dump(mode="json")
    page = next(item for item in raw["pages"] if item["takeaway_source_fragment_ids"])
    page_id = page["page_id"]
    expected_takeaway = next(
        item.takeaway for item in valid.pages if item.page_id == page_id
    )
    page["takeaway"] = "模型擅自改写出的新结论"
    page["body"] = "不允许出现的正文"

    async def invalid_planner(_request):
        return raw

    resolved = await plan_slide_visuals(
        document,
        allocation,
        fragments,
        ai_planner=invalid_planner,
    )

    assert resolved.deck_brief["planner"] == "ai"
    resolved_page = next(item for item in resolved.pages if item.page_id == page_id)
    assert resolved_page.takeaway == expected_takeaway
    assert resolved_page.planner == "ai"


@pytest.mark.asyncio
async def test_visual_planner_cannot_replace_compiler_owned_page_copy() -> None:
    course = visual_course()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    allocation = deterministic_slide_allocation(
        document,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    valid = deterministic_visual_plan(document, allocation, fragments)
    raw = valid.model_dump(mode="json")
    raw["pages"][0]["takeaway"] = "Provider-authored ungrounded body copy"

    async def planner(_request):
        return raw

    resolved = await plan_slide_visuals(
        document,
        allocation,
        fragments,
        ai_planner=planner,
    )

    assert resolved.deck_brief["ai_visual_batches_failed"] == 0
    assert resolved.pages[0].takeaway == valid.pages[0].takeaway
    assert resolved.pages[0].planner == "ai"


@pytest.mark.asyncio
async def test_visual_planner_sanitizes_invalid_choices_to_compiler_defaults() -> None:
    course = visual_course()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    allocation = deterministic_slide_allocation(
        document,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    baseline = deterministic_visual_plan(document, allocation, fragments)
    raw = baseline.model_dump(mode="json")
    for page in raw["pages"]:
        page["transition_from"] = None
        page["composition"] = "medical-atlas"
        page["role_layout_variant"] = "clinical"
        page["visual_anchor"]["purpose"] = "teaching"
        page["visual_anchor"]["asset_id"] = None

    async def planner(_request):
        return raw

    resolved = await plan_slide_visuals(
        document,
        allocation,
        fragments,
        ai_planner=planner,
    )

    assert resolved.deck_brief["ai_visual_batches_failed"] == 0
    assert resolved.deck_brief["ai_visual_pages_accepted"] == len(allocation.pages)
    assert [page.composition for page in resolved.pages] == [
        page.composition for page in baseline.pages
    ]
    assert [page.visual_anchor for page in resolved.pages] == [
        page.visual_anchor for page in baseline.pages
    ]


def test_visual_request_declares_the_exact_response_contract() -> None:
    course = visual_course()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    allocation = deterministic_slide_allocation(
        document,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )

    request = _visual_plan_request(
        document,
        allocation.model_copy(update={"pages": allocation.pages[:2]}),
        fragments,
        raster_generation_enabled=False,
        allowed_visual_kinds={"none", "relational_diagram"},
        batch_index=0,
        batch_count=1,
    )

    assert request["rules"]["return_every_requested_page"] is True
    assert request["response_contract"]["root"] == "slide_visual_plan_v1"
    assert request["response_contract"]["pages_item_required"] == [
        "page_id",
        "teaching_job",
        "takeaway",
        "takeaway_source_fragment_ids",
        "transition_from",
        "composition",
        "visual_anchor",
        "role_layout_variant",
    ]


@pytest.mark.asyncio
async def test_long_deck_uses_bounded_visual_planning_batches() -> None:
    course = visual_course()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    allocation = deterministic_slide_allocation(
        document,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    pages = list(allocation.pages)
    source_page = next(page for page in reversed(pages) if page.fragment_ids)
    while len(pages) < 30:
        pages.append(source_page.model_copy(update={
            "page_id": f"slide:long:{len(pages):04d}",
        }))
    long_allocation = allocation.model_copy(update={"pages": pages})
    calls = 0

    batch_sizes: list[int] = []

    async def planner(request: dict) -> dict:
        nonlocal calls
        calls += 1
        page_ids = [page["page_id"] for page in request["pages"]]
        batch_sizes.append(len(page_ids))
        batch_allocation = long_allocation.model_copy(update={
            "pages": [
                page for page in long_allocation.pages
                if page.page_id in set(page_ids)
            ],
        })
        return deterministic_visual_plan(
            document,
            batch_allocation,
            fragments,
        ).model_dump(mode="json")

    resolved = await plan_slide_visuals(
        document,
        long_allocation,
        fragments,
        ai_planner=planner,
    )

    assert calls == resolved.deck_brief["ai_visual_batches_total"]
    assert max(batch_sizes) <= 12
    assert resolved.deck_brief["planner"] == "ai"
    assert resolved.deck_brief["ai_visual_batches_total"] >= 2
    assert resolved.deck_brief["ai_visual_batches_successful"] == calls
    assert resolved.deck_brief["ai_visual_batches_failed"] == 0
    assert all(page.planner == "ai" for page in resolved.pages)


@pytest.mark.asyncio
async def test_visual_batch_accepts_pages_only_provider_envelope() -> None:
    course = visual_course()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    allocation = deterministic_slide_allocation(
        document,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )

    async def planner(request: dict) -> dict:
        page_ids = {page["page_id"] for page in request["pages"]}
        batch_allocation = allocation.model_copy(update={
            "pages": [
                page for page in allocation.pages
                if page.page_id in page_ids
            ],
        })
        full = deterministic_visual_plan(
            document,
            batch_allocation,
            fragments,
        ).model_dump(mode="json")
        return {"pages": full["pages"]}

    resolved = await plan_slide_visuals(
        document,
        allocation,
        fragments,
        ai_planner=planner,
    )

    assert resolved.deck_brief["planner"] == "ai"
    assert resolved.deck_brief["ai_visual_batches_failed"] == 0
    assert all(page.planner == "ai" for page in resolved.pages)


@pytest.mark.asyncio
async def test_visual_batches_salvage_partial_pages_from_string_envelope() -> None:
    course = visual_course()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    allocation = deterministic_slide_allocation(
        document,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )

    async def planner(request: dict) -> dict:
        requested_ids = [page["page_id"] for page in request["pages"]]
        batch_allocation = allocation.model_copy(update={
            "pages": [
                page for page in allocation.pages
                if page.page_id == requested_ids[0]
            ],
        })
        first_page = deterministic_visual_plan(
            document,
            batch_allocation,
            fragments,
        ).model_dump(mode="json")["pages"]
        return {
            "slide_visual_plan_v1": json.dumps(
                {"pages": first_page},
                ensure_ascii=False,
            ),
        }

    resolved = await plan_slide_visuals(
        document,
        allocation,
        fragments,
        ai_planner=planner,
    )

    expected_batches = len(_visual_plan_batches(allocation, 12))
    assert resolved.deck_brief["ai_visual_batches_successful"] == expected_batches
    assert resolved.deck_brief["ai_visual_batches_failed"] == 0
    assert resolved.deck_brief["ai_visual_pages_accepted"] == expected_batches
    assert resolved.deck_brief["ai_visual_pages_fallback"] == (
        len(allocation.pages) - expected_batches
    )
    assert sum(page.planner == "ai" for page in resolved.pages) == expected_batches


def test_visual_planning_batches_never_mix_chapters() -> None:
    pages = [
        SimpleNamespace(page_id="cover", chapter_id=""),
        *[
            SimpleNamespace(page_id=f"chapter-1-{index}", chapter_id="chapter-1")
            for index in range(8)
        ],
        *[
            SimpleNamespace(page_id=f"chapter-2-{index}", chapter_id="chapter-2")
            for index in range(7)
        ],
        *[
            SimpleNamespace(page_id=f"chapter-3-{index}", chapter_id="chapter-3")
            for index in range(4)
        ],
        SimpleNamespace(page_id="summary", chapter_id=""),
    ]

    batches = _visual_plan_batches(SimpleNamespace(pages=pages), 12)

    assert max(len(batch) for batch in batches) <= 12
    assert all(
        len({page.chapter_id for page in batch if page.chapter_id}) <= 1
        for batch in batches
    )


@pytest.mark.asyncio
async def test_long_deck_visual_batch_failure_preserves_successful_batches() -> None:
    course = visual_course()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    allocation = deterministic_slide_allocation(
        document,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    pages = list(allocation.pages)
    source_page = next(page for page in reversed(pages) if page.fragment_ids)
    while len(pages) < 30:
        pages.append(source_page.model_copy(update={
            "page_id": f"slide:partial:{len(pages):04d}",
        }))
    long_allocation = allocation.model_copy(update={"pages": pages})
    calls = 0

    async def planner(request: dict) -> dict:
        nonlocal calls
        calls += 1
        if calls == 2:
            raise RuntimeError("provider unavailable for this batch")
        page_ids = {page["page_id"] for page in request["pages"]}
        batch_allocation = long_allocation.model_copy(update={
            "pages": [page for page in pages if page.page_id in page_ids],
        })
        return deterministic_visual_plan(
            document,
            batch_allocation,
            fragments,
        ).model_dump(mode="json")

    resolved = await plan_slide_visuals(
        document,
        long_allocation,
        fragments,
        ai_planner=planner,
    )

    assert resolved.deck_brief["planner"] == "ai"
    assert resolved.deck_brief["fallback_reason"] == "partial_ai_visual_plan"
    assert resolved.deck_brief["ai_visual_batches_successful"] == calls - 1
    assert resolved.deck_brief["ai_visual_batches_failed"] == 1
    assert {page.planner for page in resolved.pages} == {
        "ai",
        "deterministic_fallback",
    }


def test_asset_repository_validates_and_promotes_content_addressed_images(
    tmp_path: Path,
) -> None:
    repository = SlideAssetRepository(tmp_path / "assets")
    source = tmp_path / "source.png"
    Image.new("RGB", (640, 360), "#2F6FE4").save(source)

    staged = repository.stage_image(
        source,
        course_id="visual-course",
        source_fragment_ids=["fragment-1"],
        alt_text="蓝色抽象教学背景",
        purpose="application",
    )
    assert repository.get(staged.asset_id) is None
    assert repository.get_staged(staged.asset_id) == staged
    published = repository.promote(staged)

    assert published.asset_id.startswith("sva_")
    assert published.sha256
    assert repository.resolve(published.asset_id).read_bytes() == source.read_bytes()


def test_asset_repository_rejects_bad_images(tmp_path: Path) -> None:
    repository = SlideAssetRepository(tmp_path / "assets")
    bad = tmp_path / "bad.png"
    bad.write_bytes(b"not-an-image")

    with pytest.raises(ValueError, match="valid raster image"):
        repository.stage_image(
            bad,
            course_id="visual-course",
            source_fragment_ids=["fragment-1"],
            alt_text="损坏图片",
            purpose="application",
        )


def test_failed_quality_discards_staged_assets_without_publishing(tmp_path: Path) -> None:
    repository = SlideAssetRepository(tmp_path / "assets")
    source = tmp_path / "source.png"
    Image.new("RGB", (640, 360), "#2F6FE4").save(source)
    staged = repository.stage_image(
        source,
        course_id="visual-course",
        source_fragment_ids=["fragment-1"],
        alt_text="课程应用场景",
        purpose="application",
        kind="generated_illustration",
    )

    finalize_visual_assets(
        [staged.model_dump(mode="json")],
        repository=repository,
        publish=False,
    )

    assert repository.get(staged.asset_id) is None
    assert not list(repository.staging.glob(f"{staged.asset_id}-*"))


def test_pptx_uses_editable_native_connectors_for_visual_diagrams(
    tmp_path: Path,
) -> None:
    course = visual_course()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    allocation = deterministic_slide_allocation(
        document,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    content = compile_slide_deck_v3(
        document,
        course,
        mode="teaching",
        theme="qizhi-classroom",
        allocation_plan=allocation,
    )

    output = export_structured_slide_deck(content, tmp_path / "visual-deck.pptx")
    presentation = Presentation(output)

    assert any(
        shape.shape_type == MSO_SHAPE_TYPE.LINE
        for slide in presentation.slides
        for shape in slide.shapes
    )
    assert any(
        shape.has_text_frame and "课程正文" in shape.text
        for slide in presentation.slides
        for shape in slide.shapes
    )
    assert all(
        not (shape.has_text_frame and "SOURCE" in shape.text)
        for slide in presentation.slides
        for shape in slide.shapes
    )
    slide_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "$$" not in slide_text
    assert "T(au+bv)=aT(u)+bT(v)" in slide_text


def test_formula_text_compiler_outputs_portable_mathematical_notation() -> None:
    assert _format_formula_text(
        r"$$ \dim(U_1 \cap U_2) \leq \min(\dim U_1, \dim U_2) $$"
    ) == "dim(U₁ ∩ U₂) ≤ min(dim U₁, dim U₂)"
    assert _format_formula_text(
        r"$$ \begin{bmatrix} a_1 \\ a_2 \\ \vdots \\ a_n \end{bmatrix} $$"
    ) == "⎡ a₁ ⎤\n⎢ a₂ ⎥\n⎢ ⋮ ⎥\n⎣ aₙ ⎦"
    assert _format_formula_text(
        r"$$ U_1 \cap U_2 = \{v \in V \mid v \in U_1 \land v \in U_2\} $$"
    ) == "U₁ ∩ U₂ = {v ∈ V ∣ v ∈ U₁ ∧ v ∈ U₂}"
    assert _format_formula_text(
        r"$$ I_j \approx \sum_{i=1}^{k} c_{ij}u_i $$"
    ) == "Iⱼ ≈ ∑ᵢ₌₁ᵏ cᵢⱼuᵢ"
    assert _format_formula_text(
        r"U_1 \\cap U_2 = \\{v \\in V \\mid v \\in U_1 \\land v \\in U_2\\}"
    ) == "U₁ ∩ U₂ = {v ∈ V ∣ v ∈ U₁ ∧ v ∈ U₂}"
    assert _format_formula_text(r"\Sigma_i=1^k c_i") == "∑ᵢ₌₁ᵏ cᵢ"


def test_image_provider_failure_degrades_to_deterministic_diagram(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setenv("SLIDE_IMAGE_API_BASE", "https://images.invalid/v1")
    monkeypatch.setenv("SLIDE_IMAGE_API_KEY", "test-key")
    monkeypatch.setenv("SLIDE_IMAGE_MODEL", "test-image-model")
    monkeypatch.setenv("SLIDE_GENERATED_ILLUSTRATIONS_ENABLED", "true")

    def fail_generation(*_args, **_kwargs):
        raise TimeoutError("provider timeout")

    monkeypatch.setattr(
        "slide_asset_repository.SlideImageProvider.generate",
        fail_generation,
    )
    course = visual_course()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    allocation = deterministic_slide_allocation(
        document,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )

    content = compile_slide_deck_v3(
        document,
        course,
        mode="teaching",
        theme="qizhi-classroom",
        allocation_plan=allocation,
    )

    assert content["quality_report"]["passed"] is True
    assert not content["visual_asset_manifest"]
    assert all(
        visual["kind"] != "generated_illustration"
        for slide in content["slides"]
        for visual in slide["visuals"]
    )


def test_raster_generation_requires_explicit_opt_in(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    monkeypatch.setenv("SLIDE_IMAGE_API_BASE", "https://images.example/v1")
    monkeypatch.setenv("SLIDE_IMAGE_API_KEY", "test-key")
    monkeypatch.setenv("SLIDE_IMAGE_MODEL", "test-image-model")
    monkeypatch.delenv("SLIDE_GENERATED_ILLUSTRATIONS_ENABLED", raising=False)
    generation_calls = 0

    def unexpected_generation(*_args, **_kwargs):
        nonlocal generation_calls
        generation_calls += 1
        raise AssertionError("raster generation must be opt-in")

    monkeypatch.setattr(
        "slide_asset_repository.SlideImageProvider.generate",
        unexpected_generation,
    )
    course = visual_course()
    document = document_from_legacy_course(course)
    content = compile_slide_deck_v3(
        document,
        course,
        mode="teaching",
        theme="qizhi-classroom",
    )

    assert generation_calls == 0
    assert not content["visual_asset_manifest"]
    assert all(
        visual["kind"] != "generated_illustration"
        for slide in content["slides"]
        for visual in slide["visuals"]
    )


def test_configured_image_provider_exports_a_real_picture(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    monkeypatch.setenv("SLIDE_IMAGE_API_BASE", "https://images.example/v1")
    monkeypatch.setenv("SLIDE_IMAGE_API_KEY", "test-key")
    monkeypatch.setenv("SLIDE_IMAGE_MODEL", "test-image-model")
    monkeypatch.setenv("SLIDE_GENERATED_ILLUSTRATIONS_ENABLED", "true")

    def generate_image(_provider, *, output_path, **_kwargs):
        Image.new("RGB", (960, 640), "#B9DCF4").save(output_path)
        return Path(output_path)

    monkeypatch.setattr(
        "slide_asset_repository.SlideImageProvider.generate",
        generate_image,
    )
    monkeypatch.setattr(
        "slide_asset_repository.SlideImageProvider.plan_prompt",
        lambda _provider, **_kwargs: "layered geometric objects in a clear process",
    )
    course = visual_course()
    course["nodes"][0]["content_blocks"].append({
        "block_id": "application-context",
        "title": "设备状态流转案例",
        "content": "设备先接收输入，然后执行校验，最后输出处理结果。",
        "metadata": {"role": "example"},
    })
    document = document_from_legacy_course(course)
    repository = SlideAssetRepository(tmp_path / "assets")
    content = compile_slide_deck_v3(
        document,
        course,
        mode="teaching",
        theme="qizhi-classroom",
        asset_repository=repository,
    )
    output = export_structured_slide_deck(
        content,
        tmp_path / "visual-image-deck.pptx",
        asset_repository=repository,
    )
    presentation = Presentation(output)

    assert content["visual_asset_manifest"]
    assert any(
        visual["kind"] == "generated_illustration" and visual["asset_id"]
        for slide in content["slides"]
        for visual in slide["visuals"]
    )
    assert any(
        shape.shape_type == MSO_SHAPE_TYPE.PICTURE
        for slide in presentation.slides
        for shape in slide.shapes
    )
