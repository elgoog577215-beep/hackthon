from __future__ import annotations

from copy import deepcopy
from datetime import datetime, timezone
from zipfile import ZipFile

import pytest
from fastapi import FastAPI
from fastapi.testclient import TestClient
from pptx import Presentation

from course_commands import CourseCommandService
from course_document import COURSE_DOCUMENT_SCHEMA, document_from_legacy_course
from course_repository import CourseDocumentRepository
from course_revisions import revision_event_for_documents, revision_vector_for_document
from representation_compiler import (
    CORE_TYPES,
    compile_core_representations,
    export_slide_deck_pptx,
    rebuild_core_representations_safely,
    validate_compiled_representations,
)
from representation_edits import (
    apply_representation_only_edit,
    classify_representation_edit,
    representation_edit_impact,
)
from slide_deck import (
    LAYOUT_CAPACITY,
    SlideBlockSpec,
    _fit_blocks_for_layout,
    _plain_text,
    validate_slide_deck,
)
from teaching_representations import (
    RepresentationConflict,
    SourceBinding,
    TeachingRepresentation,
    TeachingRepresentationRepository,
    source_binding_for_document,
)


class MemoryStorage:
    def __init__(self, course: dict | None) -> None:
        self.course = deepcopy(course)
        self.save_count = 0

    def load_course(self, _course_id: str) -> dict | None:
        return deepcopy(self.course)

    async def save_course(self, _course_id: str, data: dict) -> None:
        self.course = deepcopy(data)
        self.save_count += 1


def legacy_course() -> dict:
    return {
        "course_id": "course-1",
        "course_name": "线性代数",
        "nodes": [
            {
                "node_id": "section-a",
                "parent_node_id": "root",
                "node_name": "向量",
                "node_level": 1,
                "learning_objective": "理解向量",
                "objective_id": "objective-a",
                "concept_refs": ["ckp-vector-definition"],
                "node_content": "## 定义\n\n向量有大小和方向。",
            },
            {
                "node_id": "section-b",
                "parent_node_id": "root",
                "node_name": "矩阵",
                "node_level": 1,
                "learning_objective": "理解矩阵",
                "objective_id": "objective-b",
                "concept_refs": ["ckp-matrix-definition"],
                "node_content": "## 定义\n\n矩阵是数字的矩形阵列。",
            },
        ],
    }


def course_data_with_practice() -> dict:
    return {
        **legacy_course(),
        "learning_assets": {
            "questions": [{
                "question_id": "question-vector-1",
                "revision_id": "question-vector-1-r1",
                "node_id": "section-a",
                "prompt": "向量由哪些要素确定？",
                "practice_level": "understand",
                "course_knowledge_refs": ["ckp-vector-definition"],
            }],
            "misconceptions": [],
        },
    }


def representation(
    document,
    *,
    representation_id: str,
    block_id: str | None = None,
    representation_type: str = "slide_deck",
) -> TeachingRepresentation:
    now = datetime.now(timezone.utc).isoformat()
    binding = source_binding_for_document(document, block_id=block_id)
    return TeachingRepresentation(
        representation_id=representation_id,
        course_id=document.course_id,
        representation_type=representation_type,
        source_bindings=[binding],
        source_revision_vector=binding.source_revisions,
        spec_id=f"spec-{representation_id}",
        semantic_fingerprint=f"semantic-{representation_id}",
        revision=f"revision-{representation_id}",
        status="ready",
        created_at=now,
        updated_at=now,
    )


def test_revision_vector_changes_only_target_block_and_parent_section():
    before = document_from_legacy_course(legacy_course())
    after = before.model_copy(deep=True)
    after.blocks[0].payload["markdown"] = "向量同时具有大小和方向。"
    from course_document import refresh_document_revision

    refresh_document_revision(after)
    event = revision_event_for_documents(before, after, command_id="command-1")

    assert "course_document" in event.changed_source_keys
    assert f"block:{before.blocks[0].block_id}" in event.changed_source_keys
    assert f"section:{before.blocks[0].section_id}" in event.changed_source_keys
    assert f"block:{before.blocks[1].block_id}" not in event.changed_source_keys
    assert f"section:{before.blocks[1].section_id}" not in event.changed_source_keys


def test_registry_marks_only_dependent_representation_stale(tmp_path):
    before = document_from_legacy_course(legacy_course())
    repository = TeachingRepresentationRepository(tmp_path)
    first = representation(
        before,
        representation_id="slides-a",
        block_id=before.blocks[0].block_id,
    )
    second = representation(
        before,
        representation_id="slides-b",
        block_id=before.blocks[1].block_id,
    )
    repository.register_representation(first)
    repository.register_representation(second)

    after = before.model_copy(deep=True)
    after.blocks[0].payload["markdown"] = "向量同时具有大小和方向。"
    from course_document import refresh_document_revision

    refresh_document_revision(after)
    event = revision_event_for_documents(before, after, command_id="command-1")
    updated = repository.apply_revision_event(before.course_id, event)
    states = {item.representation_id: item for item in updated.representations}

    assert states["slides-a"].status == "stale"
    assert states["slides-b"].status == "ready"
    assert any(reason.startswith("source_revision_changed:block:") for reason in states["slides-a"].stale_reasons)


def test_course_bound_representation_stales_on_any_semantic_change(tmp_path):
    before = document_from_legacy_course(legacy_course())
    repository = TeachingRepresentationRepository(tmp_path)
    outline = representation(before, representation_id="outline", representation_type="outline")
    repository.register_representation(outline)

    after = before.model_copy(deep=True)
    after.blocks[1].payload["markdown"] = "矩阵表示线性映射。"
    from course_document import refresh_document_revision

    refresh_document_revision(after)
    event = revision_event_for_documents(before, after, command_id="command-2")
    updated = repository.apply_revision_event(before.course_id, event)

    assert updated.representations[0].status == "stale"
    assert "source_revision_changed:course_document" in updated.representations[0].stale_reasons


def test_revision_event_replay_is_idempotent_and_course_isolated(tmp_path):
    before = document_from_legacy_course(legacy_course())
    repository = TeachingRepresentationRepository(tmp_path)
    repository.register_representation(representation(
        before,
        representation_id="slides-a",
        block_id=before.blocks[0].block_id,
    ))
    after = before.model_copy(deep=True)
    after.blocks[0].payload["markdown"] = "变化后的内容"
    from course_document import refresh_document_revision

    refresh_document_revision(after)
    event = revision_event_for_documents(before, after, command_id="command-3")

    first = repository.apply_revision_event(before.course_id, event)
    second = repository.apply_revision_event(before.course_id, event)
    assert second.registry_revision == first.registry_revision
    assert second.representations[0].stale_reasons == first.representations[0].stale_reasons

    with pytest.raises(RepresentationConflict):
        repository.apply_revision_event("course-2", event)


@pytest.mark.asyncio
async def test_course_command_persists_replayable_revision_event(tmp_path):
    storage = MemoryStorage(legacy_course())
    course_repository = CourseDocumentRepository(storage)
    preview = course_repository.document_envelope("course-1")
    await course_repository.migrate_legacy_course(
        "course-1",
        expected_source_checksum=preview["migration"]["source_checksum"],
    )
    document, _ = course_repository.load_document("course-1")
    target = document.blocks[0]
    representation_repository = TeachingRepresentationRepository(tmp_path)
    representation_repository.register_representation(representation(
        document,
        representation_id="slides-a",
        block_id=target.block_id,
    ))

    receipt = await CourseCommandService(course_repository).replace_block(
        "course-1",
        command_id="replace-a",
        expected_document_revision=document.document_revision,
        expected_block_revision=target.internal_revision,
        block_id=target.block_id,
        payload={"title": "定义", "markdown": "向量同时具有大小和方向。"},
    )

    assert receipt["revision_change"]["event_id"].startswith("cre_")
    assert storage.course["course_revision_vector"] == receipt["revision_change"]["current"]
    reconciled = representation_repository.reconcile_course_operation_log(
        "course-1",
        storage.course["course_operation_log"],
    )
    assert reconciled.representations[0].status == "stale"


def test_source_binding_rejects_missing_block():
    document = document_from_legacy_course(legacy_course())
    with pytest.raises(RepresentationConflict):
        source_binding_for_document(document, block_id="missing-block")


def test_representation_rejects_revision_vector_that_disagrees_with_bindings():
    now = datetime.now(timezone.utc).isoformat()
    binding = SourceBinding(course_id="course-1", source_revisions={"block:a": "revision-a"})
    with pytest.raises(ValueError):
        TeachingRepresentation(
            representation_id="slides-a",
            course_id="course-1",
            representation_type="slide_deck",
            source_bindings=[binding],
            source_revision_vector={"block:a": "different-revision"},
            spec_id="spec-a",
            revision="representation-revision-a",
            created_at=now,
            updated_at=now,
        )


def test_revision_vector_contains_document_sections_blocks_and_objectives():
    document = document_from_legacy_course(legacy_course())
    vector = revision_vector_for_document(document)

    assert vector.revisions["course_document"] == document.document_revision
    assert f"section:{document.sections[0].section_id}" in vector.revisions
    assert f"block:{document.blocks[0].block_id}" in vector.revisions
    assert "objective:objective-a" in vector.revisions


def test_representation_router_reconciles_and_returns_graph(tmp_path, monkeypatch):
    from routers import teaching_representations as representation_router

    course = legacy_course()
    storage = MemoryStorage(course)
    course_repository = CourseDocumentRepository(storage)
    document = document_from_legacy_course(course)
    representation_repository = TeachingRepresentationRepository(tmp_path)
    representation_repository.register_representation(representation(
        document,
        representation_id="slides-a",
        block_id=document.blocks[0].block_id,
    ))

    monkeypatch.setattr(
        representation_router,
        "get_teaching_representation_repository",
        lambda: representation_repository,
    )
    monkeypatch.setattr(
        representation_router,
        "get_course_document_repository",
        lambda: course_repository,
    )

    async def existing_course(_course_id: str):
        return course

    monkeypatch.setattr(representation_router, "get_course_or_404", existing_course)
    app = FastAPI()
    app.include_router(representation_router.router, prefix="/api")
    client = TestClient(app)

    missing_identity = client.get("/api/courses/course-1/teaching-representations")
    assert missing_identity.status_code == 400

    response = client.get(
        "/api/courses/course-1/teaching-representations/derivation-graph",
        headers={"X-User-Id": "user-1"},
    )
    assert response.status_code == 200
    payload = response.json()
    assert payload["course_id"] == "course-1"
    assert payload["derivation_graph"]["nodes"]


def test_compiler_builds_five_bound_representations_and_exports_pptx(tmp_path):
    document = document_from_legacy_course(legacy_course())
    repository = TeachingRepresentationRepository(tmp_path / "registry")
    course_data = legacy_course()
    course_data["learning_assets"] = {
        "questions": [{
            "question_id": "question-a",
            "revision_id": "question-revision-a",
            "node_id": "section-a",
            "prompt": "向量有哪些基本属性？",
            "practice_level": "understanding",
            "course_knowledge_refs": ["ckp-vector-definition"],
        }],
        "misconceptions": [{
            "node_id": "section-a",
            "error_pattern": "把方向相反的向量当成相同向量",
        }],
    }

    result = compile_core_representations(document, course_data, repository)
    registry = repository.load(document.course_id)
    current_spec_ids = {item.spec_id for item in registry.representations}
    specs = [item for item in registry.specs if item.spec_id in current_spec_ids]

    assert {item.representation_type for item in registry.representations} == set(CORE_TYPES)
    assert len(result["representations"]) == len(CORE_TYPES)
    assert all(spec.unit_bindings for spec in specs)
    assert validate_compiled_representations(specs)["passed"] is True
    assert set(registry.plans[0].knowledge_refs) == {
        "ckp-vector-definition",
        "ckp-matrix-definition",
    }
    for spec in specs:
        units = (
            spec.payload["content"].get("units")
            or spec.payload["content"].get("slides")
            or spec.payload["content"].get("sections")
            or []
        )
        content_units = [item for item in units if item.get("section_id")]
        assert content_units
        assert all(item.get("knowledge_refs") for item in content_units)
        assert any(
            binding.knowledge_node_ids
            for bindings in spec.unit_bindings.values()
            for binding in bindings
        )
    practice = next(spec for spec in specs if spec.representation_type == "practice_sheet")
    assert practice.payload["content"]["units"][0]["practice_task_id"] == "question-a"

    diagram = next(spec for spec in specs if spec.representation_type == "diagram")
    diagram_content = diagram.payload["content"]
    assert diagram_content["schema_version"] == "diagram_spec_v1"
    assert diagram_content["quality_report"]["passed"] is True
    assert all(unit["mermaid"].startswith("flowchart LR") for unit in diagram_content["units"])
    assert all(unit["nodes"] and unit["source_block_ids"] for unit in diagram_content["units"])

    slides = next(spec for spec in specs if spec.representation_type == "slide_deck")
    slide_content = slides.payload["content"]
    assert slide_content["schema_version"] == "slide_deck_v2"
    assert {item["layout"] for item in slide_content["slides"]} >= {"cover", "roadmap", "concept", "practice", "recap"}
    assert all("blocks" in item and "bullets" not in item for item in slide_content["slides"])
    assert len({item["unit_id"] for item in slide_content["slides"]}) == len(slide_content["slides"])
    recap = next(item for item in slide_content["slides"] if item["layout"] == "recap")
    assert all(block["items"] for block in recap["blocks"])
    for section in document.sections:
        section_slides = [
            item for item in slide_content["slides"]
            if item.get("section_id") == section.section_id
        ]
        assert {item["slide_purpose"] for item in section_slides} >= {
            "learning_objective",
            "concept_and_reasoning",
            "mastery_check",
        }
        assert all(item["speaker_notes"].strip() for item in section_slides)
    practice_slides = [item for item in slide_content["slides"] if item["layout"] == "practice"]
    assert any("question-a" in item["practice_task_ids"] for item in practice_slides)
    output = export_slide_deck_pptx(slides, tmp_path / "course.pptx")
    assert output.exists()
    assert output.stat().st_size > 0
    presentation = Presentation(output)
    assert len(presentation.slides) == len(slide_content["slides"])
    rendered_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    assert "线性代数" in rendered_text
    assert "向量" in rendered_text
    with ZipFile(output) as archive:
        first_slide_xml = archive.read("ppt/slides/slide1.xml")
    assert b"<a:ea" in first_slide_xml
    assert "Noto Sans SC".encode() in first_slide_xml
    assert "Microsoft YaHei".encode() in first_slide_xml


def test_slide_quality_gate_rejects_raw_course_copy_and_markdown(tmp_path):
    document = document_from_legacy_course(legacy_course())
    repository = TeachingRepresentationRepository(tmp_path)
    course_data = course_data_with_practice()
    result = compile_core_representations(document, course_data, repository)
    registry = repository.load(document.course_id)
    slides = next(
        spec for spec in registry.specs
        if spec.spec_id == next(
            item.spec_id for item in registry.representations
            if item.representation_type == "slide_deck"
        )
    )
    content = deepcopy(slides.payload["content"])
    target = next(item for item in content["slides"] if item.get("section_id"))
    target["blocks"][0]["content"] = "# 未转译标题\n$A \\times B$\n" + ("整段正文复制。" * 80)

    report = validate_slide_deck(content, course_data=course_data)

    assert report["passed"] is False
    codes = {item["code"] for item in report["issues"] if item["severity"] == "critical"}
    assert {"raw_markdown_leaked", "raw_latex_leaked", "paragraph_copy_detected"} <= codes


def test_cross_product_quality_rejects_objective_drift(tmp_path):
    document = document_from_legacy_course(legacy_course())
    repository = TeachingRepresentationRepository(tmp_path)
    compile_core_representations(document, course_data_with_practice(), repository)
    registry = repository.load(document.course_id)
    current_ids = {item.spec_id for item in registry.representations}
    specs = [deepcopy(item) for item in registry.specs if item.spec_id in current_ids]
    handout = next(item for item in specs if item.representation_type == "handout")
    binding = next(
        item for item in handout.unit_bindings["handout:section-a"]
        if "objective:objective-a" in item.source_revisions
    )
    objective_revision = binding.source_revisions.pop("objective:objective-a")
    binding.source_revisions["objective:objective-drifted"] = objective_revision

    report = validate_compiled_representations(specs)

    assert report["passed"] is False
    assert any(
        issue["code"] == "cross_product_objective_mismatch"
        and issue["target"] == "section-a"
        for issue in report["issues"]
    )
def test_plain_text_converts_common_latex_into_classroom_readable_math():
    converted = _plain_text(
        r"$3 \times 3$；$\mathbf{x} \in \mathbb{R}^3$；"
        r"\begin{bmatrix}1 & 2 \\ 3 & 4\end{bmatrix}；"
        r"begin{cases}x_1=1 \\ x_2=2end{cases}；"
        r"\alpha+\beta\approx\sqrt{x^{-1}}"
    )

    assert "3 × 3" in converted
    assert "x ∈ ℝ³" in converted
    assert "[1 , 2 ； 3 , 4]" in converted
    assert "[x₁=1 ； x₂=2]" in converted
    assert "α+β≈√(x⁻¹)" in converted
    assert "$" not in converted
    assert "\\" not in converted


def test_slide_quality_gate_rejects_text_encoding_corruption(tmp_path):
    document = document_from_legacy_course(legacy_course())
    repository = TeachingRepresentationRepository(tmp_path)
    course_data = course_data_with_practice()
    compile_core_representations(document, course_data, repository)
    registry = repository.load(document.course_id)
    slides = next(
        spec for spec in registry.specs
        if spec.spec_id == next(
            item.spec_id for item in registry.representations
            if item.representation_type == "slide_deck"
        )
    )
    content = deepcopy(slides.payload["content"])
    target = next(item for item in content["slides"] if item.get("section_id"))
    target["key_message"] = "矩阵乘法�表示线性变换复合"

    report = validate_slide_deck(content, course_data=course_data)

    assert report["passed"] is False
    assert "text_encoding_corrupted" in {
        item["code"] for item in report["issues"] if item["severity"] == "critical"
    }


def test_slide_quality_gate_rejects_empty_recap_cards(tmp_path):
    document = document_from_legacy_course(legacy_course())
    repository = TeachingRepresentationRepository(tmp_path)
    course_data = course_data_with_practice()
    compile_core_representations(document, course_data, repository)
    registry = repository.load(document.course_id)
    slides = next(
        spec for spec in registry.specs
        if spec.spec_id == next(
            item.spec_id for item in registry.representations
            if item.representation_type == "slide_deck"
        )
    )
    content = deepcopy(slides.payload["content"])
    recap = next(item for item in content["slides"] if item["layout"] == "recap")
    recap["blocks"][0]["items"] = []

    report = validate_slide_deck(content, course_data=course_data)

    assert report["passed"] is False
    assert "recap_evidence_missing" in {
        item["code"] for item in report["issues"] if item["severity"] == "critical"
    }


def test_layout_fitter_rechecks_item_budget_after_teaching_support_enrichment():
    blocks = [
        SlideBlockSpec(
            block_id="source",
            type="bullets",
            title="错误分析",
            items=[f"错误线索 {index}" for index in range(6)],
        ),
        SlideBlockSpec(
            block_id="support",
            type="bullets",
            title="辨析边界",
            items=[f"边界 {index}" for index in range(3)],
        ),
    ]

    fitted = _fit_blocks_for_layout("misconception", blocks)

    assert sum(len(block.items) for block in fitted) == LAYOUT_CAPACITY["misconception"]["items"]
    assert len(fitted) == 1
    assert fitted[0].items == [f"错误线索 {index}" for index in range(4)]


def test_compiled_representations_track_stale_units_instead_of_whole_course(tmp_path):
    before = document_from_legacy_course(legacy_course())
    repository = TeachingRepresentationRepository(tmp_path)
    course_data = legacy_course()
    course_data["learning_assets"] = {"questions": [], "misconceptions": []}
    compile_core_representations(before, course_data, repository)

    after = before.model_copy(deep=True)
    changed_block = after.blocks[0]
    changed_block.payload["markdown"] = "向量同时具有大小、方向和线性组合语义。"
    from course_document import refresh_document_revision

    refresh_document_revision(after)
    event = revision_event_for_documents(before, after, command_id="unit-change")
    updated = repository.apply_revision_event(before.course_id, event)
    by_type = {item.representation_type: item for item in updated.representations}

    assert by_type["slide_deck"].status == "stale"
    assert by_type["slide_deck"].stale_unit_ids == [
        "slide:recap",
        "slide:section-a",
        "slide:section-a:check",
        "slide:section-a:content:1",
    ]
    assert "slide:title" not in by_type["slide_deck"].stale_unit_ids
    assert by_type["handout"].stale_unit_ids == ["handout:section-a"]


def test_representation_edits_classify_semantic_boundary_and_preserve_course_source(tmp_path):
    document = document_from_legacy_course(legacy_course())
    repository = TeachingRepresentationRepository(tmp_path)
    compile_core_representations(
        document,
        {**legacy_course(), "learning_assets": {"questions": [], "misconceptions": []}},
        repository,
    )
    registry = repository.load(document.course_id)
    original_block_payload = deepcopy(document.blocks[0].payload)
    slides = next(item for item in registry.representations if item.representation_type == "slide_deck")
    spec = next(item for item in registry.specs if item.spec_id == slides.spec_id)

    assert classify_representation_edit(
        field="layout", before="left", after="right",
    )["classification"] == "presentation"
    assert classify_representation_edit(
        field="example", before="旧例子", after="具有不同数学含义的新例子",
    )["classification"] == "ambiguous"
    assert classify_representation_edit(
        field="example", before="旧例子", after="新例子", semantic_intent=True,
    )["classification"] == "semantic"
    detected_goal_shift = classify_representation_edit(
        field="key_message",
        before="掌握向量加法的计算规则",
        after="理解向量加法为什么表示位移的复合",
    )
    assert detected_goal_shift["classification"] == "semantic"
    assert detected_goal_shift["semantic_change"]["from_label"] == "计算技能"
    assert detected_goal_shift["semantic_change"]["to_label"] == "概念理解"
    assert "不只是措辞调整" in detected_goal_shift["semantic_change"]["interpretation"]
    assert len(detected_goal_shift["semantic_change"]["instructional_implications"]) == 3

    impact = representation_edit_impact(registry, spec, unit_id="slide:section-a")
    assert document.blocks[0].block_id in impact["block_ids"]
    assert {item["representation_type"] for item in impact["affected_representations"]} >= {
        "outline", "lesson_plan", "handout", "slide_deck",
    }
    assert any(item["origin"] for item in impact["change_items"])
    assert any(item["role"] == "教案重点" for item in impact["change_items"])
    assert impact["protected_items"]

    updated = apply_representation_only_edit(
        repository,
        registry,
        slides,
        spec,
        unit_id="slide:section-a",
        field="title",
        after="向量意味着什么",
    )
    updated_slides = next(item for item in updated.representations if item.representation_type == "slide_deck")
    updated_spec = next(item for item in updated.specs if item.spec_id == updated_slides.spec_id)
    unit = next(item for item in updated_spec.payload["content"]["slides"] if item["unit_id"] == "slide:section-a")
    assert unit["title"] == "向量意味着什么"
    assert document.blocks[0].payload == original_block_payload


def test_semantic_representation_edit_creates_authoring_change_without_writing_course(
    tmp_path,
    monkeypatch,
):
    from change_proposals import ChangeProposalRepository
    from routers import teaching_representations as representation_router

    course = legacy_course()
    document = document_from_legacy_course(course)
    canonical = {
        **course,
        "course_schema_version": COURSE_DOCUMENT_SCHEMA,
        "course_document": document.model_dump(mode="json"),
        "course_document_authoritative": True,
        "course_operation_log": [],
    }
    storage = MemoryStorage(canonical)
    course_repository = CourseDocumentRepository(storage)
    representation_repository = TeachingRepresentationRepository(tmp_path / "representations")
    compile_core_representations(
        document,
        {**course, "learning_assets": {"questions": [], "misconceptions": []}},
        representation_repository,
    )
    proposal_repository = ChangeProposalRepository(tmp_path / "authoring_changes")
    slides = next(
        item for item in representation_repository.load("course-1").representations
        if item.representation_type == "slide_deck"
    )

    monkeypatch.setattr(
        representation_router,
        "get_teaching_representation_repository",
        lambda: representation_repository,
    )
    monkeypatch.setattr(
        representation_router,
        "get_course_document_repository",
        lambda: course_repository,
    )
    monkeypatch.setattr(
        representation_router,
        "change_proposal_repository",
        proposal_repository,
    )

    async def existing_course(_course_id: str):
        return course_repository.load_course_view("course-1")

    monkeypatch.setattr(representation_router, "get_course_or_404", existing_course)
    app = FastAPI()
    app.include_router(representation_router.router, prefix="/api")
    client = TestClient(app)
    before_course = deepcopy(storage.course)

    response = client.post(
        f"/api/courses/course-1/teaching-representations/{slides.representation_id}/edits/apply",
        headers={"X-User-Id": "teacher-1"},
        json={
            "unit_id": "slide:section-a",
            "field": "title",
            "before": "向量",
            "after": "向量的几何含义",
            "semantic_intent": True,
            "decision": "course_semantic",
        },
    )

    assert response.status_code == 200
    change = response.json()["authoring_change"]
    assert change["change_kind"] == "course_authoring_change"
    assert change["write_target"] == "base_course"
    assert change["source"] == "representation_semantic"
    assert storage.course == before_course


def test_objective_edit_updates_course_truth_and_reuses_unaffected_representation_units(
    tmp_path,
    monkeypatch,
):
    from change_proposals import ChangeProposalRepository
    from routers import change_proposals as changes_router
    from routers import teaching_representations as representation_router

    course = legacy_course()
    course["nodes"][0]["learning_objective"] = "掌握向量加法的计算规则"
    document = document_from_legacy_course(course)
    canonical = {
        **course,
        "course_schema_version": COURSE_DOCUMENT_SCHEMA,
        "course_document": document.model_dump(mode="json"),
        "course_document_revision": document.document_revision,
        "course_document_authoritative": True,
        "course_operation_log": [],
    }
    storage = MemoryStorage(canonical)
    course_repository = CourseDocumentRepository(storage)
    representation_repository = TeachingRepresentationRepository(tmp_path / "representations")
    compile_core_representations(
        document,
        course_data_with_practice(),
        representation_repository,
    )
    proposal_repository = ChangeProposalRepository(tmp_path / "authoring_changes")
    slides = next(
        item for item in representation_repository.load("course-1").representations
        if item.representation_type == "slide_deck"
    )

    monkeypatch.setattr(
        representation_router,
        "get_teaching_representation_repository",
        lambda: representation_repository,
    )
    monkeypatch.setattr(
        representation_router,
        "get_course_document_repository",
        lambda: course_repository,
    )
    monkeypatch.setattr(
        representation_router,
        "change_proposal_repository",
        proposal_repository,
    )
    monkeypatch.setattr(
        changes_router,
        "get_change_proposal_repository",
        lambda: proposal_repository,
    )
    monkeypatch.setattr(
        changes_router,
        "get_course_document_repository",
        lambda: course_repository,
    )
    monkeypatch.setattr(
        changes_router,
        "teaching_representation_repository",
        representation_repository,
    )

    async def existing_course(_course_id: str):
        return course_repository.load_course_view("course-1")

    monkeypatch.setattr(representation_router, "get_course_or_404", existing_course)
    app = FastAPI()
    app.include_router(representation_router.router, prefix="/api")
    app.include_router(changes_router.authoring_router, prefix="/api")
    client = TestClient(app)
    after_objective = "理解向量加法为什么表示位移的复合，并能够解释运算顺序"

    preview = client.post(
        f"/api/courses/course-1/teaching-representations/{slides.representation_id}/edits/preview",
        headers={"X-User-Id": "teacher-1"},
        json={
            "unit_id": "slide:section-a",
            "field": "key_message",
            "before": "掌握向量加法的计算规则",
            "after": after_objective,
        },
    )
    assert preview.status_code == 200
    assert preview.json()["classification"] == "semantic"
    assert preview.json()["semantic_change"]["from_label"] == "计算技能"
    assert preview.json()["semantic_change"]["to_label"] == "概念理解"
    assert preview.json()["impact"]["affected_unit_count"] > 0
    assert preview.json()["impact"]["unaffected_unit_count"] > 0
    assert any(
        item["role"] == "课堂理解检查"
        for item in preview.json()["impact"]["change_items"]
    )

    proposed = client.post(
        f"/api/courses/course-1/teaching-representations/{slides.representation_id}/edits/apply",
        headers={"X-User-Id": "teacher-1"},
        json={
            "unit_id": "slide:section-a",
            "field": "key_message",
            "before": "掌握向量加法的计算规则",
            "after": after_objective,
            "semantic_intent": True,
            "decision": "course_semantic",
        },
    )
    assert proposed.status_code == 200
    change = proposed.json()["authoring_change"]
    item = change["items"][0]
    assert change["scope"] == "section"
    assert item["target_kind"] == "course_objective"
    assert item["block_id"] == "section-a"
    before_apply, _canonical = course_repository.load_document("course-1")
    assert before_apply.sections[0].learning_objective == "掌握向量加法的计算规则"

    applied = client.post(
        f"/api/courses/course-1/authoring-changes/{change['proposal_id']}/items/{item['item_id']}/apply",
        headers={"X-User-Id": "teacher-1"},
    )
    assert applied.status_code == 200
    receipt = applied.json()["representation_sync"]
    assert receipt["status"] == "synchronized"
    assert receipt["rebuilt_unit_count"] > 0
    assert receipt["reused_unit_count"] > 0
    assert receipt["changed_unit_count"] >= 4
    changed_types = {
        item["representation_type"]
        for item in receipt["changes"]
        if any(unit["change_kind"] == "content_changed" for unit in item["units"])
    }
    assert {"outline", "lesson_plan", "handout", "practice_sheet", "slide_deck", "diagram"} <= changed_types
    lesson_change = next(
        unit
        for item in receipt["changes"]
        if item["representation_type"] == "lesson_plan"
        for unit in item["units"]
        if unit["change_kind"] == "content_changed"
    )
    assert "规则、步骤" in lesson_change["before"]
    assert "概念含义、关系" in lesson_change["after"]
    updated, _canonical = course_repository.load_document("course-1")
    assert updated.sections[0].learning_objective == after_objective
    assert updated.sections[1].learning_objective == "理解矩阵"
    assert updated.blocks == document.blocks
    revision_change = applied.json()["items"][0]["receipt"]["revision_change"]
    assert "objective:objective-a" in revision_change["changed_source_keys"]

    reverse = client.post(
        f"/api/courses/course-1/teaching-representations/{slides.representation_id}/edits/apply",
        headers={"X-User-Id": "teacher-1"},
        json={
            "unit_id": "slide:section-a",
            "field": "key_message",
            "before": after_objective,
            "after": "掌握向量加法的计算规则",
            "semantic_intent": True,
            "decision": "course_semantic",
        },
    )
    assert reverse.status_code == 200
    reverse_change = reverse.json()["authoring_change"]
    reverse_item = reverse_change["items"][0]
    reverse_applied = client.post(
        f"/api/courses/course-1/authoring-changes/{reverse_change['proposal_id']}/items/{reverse_item['item_id']}/apply",
        headers={"X-User-Id": "teacher-1"},
    )
    assert reverse_applied.status_code == 200

    repeated_value_on_new_revision = client.post(
        f"/api/courses/course-1/teaching-representations/{slides.representation_id}/edits/apply",
        headers={"X-User-Id": "teacher-1"},
        json={
            "unit_id": "slide:section-a",
            "field": "key_message",
            "before": "掌握向量加法的计算规则",
            "after": after_objective,
            "semantic_intent": True,
            "decision": "course_semantic",
        },
    )
    assert repeated_value_on_new_revision.status_code == 200
    repeated_change = repeated_value_on_new_revision.json()["authoring_change"]
    assert repeated_change["proposal_id"] != change["proposal_id"]
    assert repeated_change["status"] == "pending"


def test_safe_rebuild_publishes_only_after_quality_passes(tmp_path):
    document = document_from_legacy_course(legacy_course())
    repository = TeachingRepresentationRepository(tmp_path)
    course_data = course_data_with_practice()
    compile_core_representations(document, course_data, repository)
    before = repository.load(document.course_id)

    changed = document.model_copy(deep=True)
    changed.blocks[0].payload["markdown"] = "向量表达大小和方向。"
    from course_document import refresh_document_revision

    refresh_document_revision(changed)
    repository.apply_revision_event(
        document.course_id,
        revision_event_for_documents(document, changed, command_id="edit-1"),
    )

    result = rebuild_core_representations_safely(changed, course_data, repository)
    current = repository.load(document.course_id)

    assert result["status"] == "synchronized"
    assert result["quality"]["passed"] is True
    assert any(item["stale_unit_ids"] for item in result["stale_before"])
    assert all(item.status == "ready" for item in current.representations)
    assert current.registry_revision != before.registry_revision


def test_safe_rebuild_failure_keeps_last_available_registry(tmp_path, monkeypatch):
    import representation_compiler

    document = document_from_legacy_course(legacy_course())
    repository = TeachingRepresentationRepository(tmp_path)
    course_data = course_data_with_practice()
    compile_core_representations(document, course_data, repository)
    before = repository.load(document.course_id).model_dump(mode="json")

    def fail_compile(*_args, **_kwargs):
        raise RuntimeError("renderer unavailable")

    monkeypatch.setattr(representation_compiler, "compile_core_representations", fail_compile)
    result = rebuild_core_representations_safely(document, course_data, repository)

    assert result["status"] == "failed_using_last_available"
    assert result["quality"]["passed"] is False
    assert repository.load(document.course_id).model_dump(mode="json") == before


def test_presentation_only_override_survives_compatible_course_rebuild(tmp_path):
    document = document_from_legacy_course(legacy_course())
    repository = TeachingRepresentationRepository(tmp_path)
    course_data = course_data_with_practice()
    compile_core_representations(document, course_data, repository)
    registry = repository.load(document.course_id)
    slides = next(item for item in registry.representations if item.representation_type == "slide_deck")
    spec = next(item for item in registry.specs if item.spec_id == slides.spec_id)

    apply_representation_only_edit(
        repository,
        registry,
        slides,
        spec,
        unit_id="slide:section-a",
        field="title",
        after="向量：大小、方向与表示",
    )
    changed = document.model_copy(deep=True)
    changed.blocks[1].payload["markdown"] = "矩阵既是数字阵列，也可以表示线性映射。"
    from course_document import refresh_document_revision

    refresh_document_revision(changed)
    repository.apply_revision_event(
        document.course_id,
        revision_event_for_documents(document, changed, command_id="edit-unrelated-section"),
    )
    result = rebuild_core_representations_safely(changed, course_data, repository)
    current = repository.load(document.course_id)
    current_slides = next(item for item in current.representations if item.representation_type == "slide_deck")
    current_spec = next(item for item in current.specs if item.spec_id == current_slides.spec_id)
    edited_slide = next(
        item for item in current_spec.payload["content"]["slides"]
        if item["unit_id"] == "slide:section-a"
    )

    assert result["status"] == "synchronized"
    assert edited_slide["title"] == "向量：大小、方向与表示"
    assert current_spec.payload["content"]["override_conflicts"] == []


def test_compiler_upgrade_rebuilds_ready_units_instead_of_reusing_old_payload(tmp_path):
    document = document_from_legacy_course(legacy_course())
    repository = TeachingRepresentationRepository(tmp_path)
    course_data = course_data_with_practice()
    compile_core_representations(document, course_data, repository)
    registry = repository.load(document.course_id)
    slides = next(item for item in registry.representations if item.representation_type == "slide_deck")
    spec = next(item for item in registry.specs if item.spec_id == slides.spec_id)
    original_title = spec.payload["content"]["slides"][0]["title"]
    spec.payload["compiler_version"] = "same_source_compiler_v2:structured_slide_compiler_v2"
    spec.payload["content"]["slides"][0]["title"] = "旧生成器残留标题"
    repository.save(registry)

    result = compile_core_representations(document, course_data, repository)
    current = repository.load(document.course_id)
    current_slides = next(item for item in current.representations if item.representation_type == "slide_deck")
    current_spec = next(item for item in current.specs if item.spec_id == current_slides.spec_id)
    slide_build = next(
        item for item in result["representations"]
        if item["representation_type"] == "slide_deck"
    )

    assert current_spec.payload["content"]["slides"][0]["title"] == original_title
    assert current_spec.payload["compiler_version"].endswith("structured_slide_compiler_v9")
    assert slide_build["reused_unit_ids"] == []


def test_progressive_build_streams_each_slide_before_atomic_completion(tmp_path, monkeypatch):
    from routers import teaching_representations as representation_router

    course = legacy_course()
    document = document_from_legacy_course(course)
    canonical = {
        **course,
        "course_schema_version": COURSE_DOCUMENT_SCHEMA,
        "course_document": document.model_dump(mode="json"),
        "course_document_authoritative": True,
        "course_operation_log": [],
        "learning_assets": {"questions": [], "misconceptions": []},
    }
    course_repository = CourseDocumentRepository(MemoryStorage(canonical))
    representation_repository = TeachingRepresentationRepository(tmp_path)
    monkeypatch.setattr(
        representation_router,
        "get_teaching_representation_repository",
        lambda: representation_repository,
    )
    monkeypatch.setattr(
        representation_router,
        "get_course_document_repository",
        lambda: course_repository,
    )

    async def existing_course(_course_id: str):
        return course_repository.load_course_view("course-1")

    monkeypatch.setattr(representation_router, "get_course_or_404", existing_course)
    app = FastAPI()
    app.include_router(representation_router.router, prefix="/api")
    client = TestClient(app)

    with client.stream(
        "POST",
        "/api/courses/course-1/teaching-representations/build/stream",
        headers={"X-User-Id": "teacher-1"},
    ) as response:
        body = "".join(response.iter_text())

    assert response.status_code == 200
    assert "event: deck_plan" in body
    assert '"strategy": "plan_then_fill"' in body
    assert '"estimated_slide_count":' in body
    assert "event: slide_upsert" in body
    assert "event: slide_quality" in body
    assert "event: build_complete" in body
    assert body.index("event: slide_upsert") < body.index("event: build_complete")
    registry = representation_repository.load("course-1")
    assert all(item.status == "ready" for item in registry.representations)


@pytest.mark.asyncio
async def test_compile_registry_reconciles_before_reusing_units(tmp_path, monkeypatch):
    from routers import teaching_representations as representation_router

    storage = MemoryStorage(legacy_course())
    course_repository = CourseDocumentRepository(storage)
    preview = course_repository.document_envelope("course-1")
    await course_repository.migrate_legacy_course(
        "course-1",
        expected_source_checksum=preview["migration"]["source_checksum"],
    )
    document, _ = course_repository.load_document("course-1")
    representation_repository = TeachingRepresentationRepository(tmp_path / "representations")
    compile_core_representations(
        document,
        course_repository.load_course_view("course-1"),
        representation_repository,
    )
    target = document.blocks[0]
    marker = "fresh canonical unit text"
    await CourseCommandService(course_repository).replace_block(
        "course-1",
        command_id="canonical-change-before-first-build",
        expected_document_revision=document.document_revision,
        expected_block_revision=target.internal_revision,
        block_id=target.block_id,
        payload={**target.payload, "markdown": marker},
    )
    monkeypatch.setattr(
        representation_router,
        "get_course_document_repository",
        lambda: course_repository,
    )
    monkeypatch.setattr(
        representation_router,
        "get_teaching_representation_repository",
        lambda: representation_repository,
    )

    representation_router._compile_registry("course-1")

    registry = representation_repository.load("course-1")
    handout = next(
        item for item in registry.representations
        if item.representation_type == "handout"
    )
    handout_spec = next(item for item in registry.specs if item.spec_id == handout.spec_id)
    assert marker in str(handout_spec.payload)
    assert handout.status == "ready"


@pytest.mark.asyncio
async def test_durable_representation_task_builds_and_recovers_after_restart(tmp_path, monkeypatch):
    import task_manager as task_manager_module
    from task_manager import TaskManager

    course = legacy_course()
    document = document_from_legacy_course(course)
    storage = MemoryStorage({
        **course,
        "course_schema_version": COURSE_DOCUMENT_SCHEMA,
        "course_document": document.model_dump(mode="json"),
        "course_document_revision": document.document_revision,
        "course_document_authoritative": True,
        "course_operation_log": [],
        "learning_assets": course_data_with_practice()["learning_assets"],
    })
    document_repository = CourseDocumentRepository(storage)
    representation_repository = TeachingRepresentationRepository(tmp_path / "representations")
    tasks_file = tmp_path / "generation_jobs.json"
    monkeypatch.setattr(task_manager_module, "TASKS_FILE", tasks_file)
    monkeypatch.setattr(
        task_manager_module,
        "teaching_representation_repository",
        representation_repository,
    )
    manager = TaskManager(
        storage,
        course_service=None,
        ws_service=None,
        document_repository=document_repository,
    )
    task_id = await manager.create_task(
        "course-1", "teaching_representation_build", enqueue=False,
    )

    await manager._process_task(task_id)

    completed = manager.get_task(task_id)
    assert completed["status"] == "completed"
    assert set(completed["completed_representation_types"]) == set(CORE_TYPES)
    assert completed["recovery"]["state"] == "completed"
    assert any(item["event"] == "slide_upsert" for item in completed["event_history"])

    interrupted_id = await manager.create_task(
        "course-1", "teaching_representation_build", enqueue=False,
    )
    manager.tasks[interrupted_id]["status"] = "running"
    manager.tasks[interrupted_id]["progress"] = 42
    manager.save_tasks(strict=True)
    restarted = TaskManager(
        storage,
        course_service=None,
        ws_service=None,
        document_repository=document_repository,
    )
    assert await restarted._reconcile_task_after_restart(interrupted_id) is True
    assert restarted.tasks[interrupted_id]["status"] == "pending"
    assert restarted.tasks[interrupted_id]["restart_recovery_count"] == 1
    await restarted.pause_task(interrupted_id)
    resumed = await restarted.resume_task(interrupted_id)
    assert resumed["status"] == "resumed"
    assert restarted.tasks[interrupted_id]["status"] == "pending"


@pytest.mark.asyncio
async def test_cancelled_build_is_not_overwritten_as_failed_by_the_late_worker(
    tmp_path, monkeypatch,
):
    """Cancelling a running build must survive the worker's own late exception.

    Cancellation makes the in-flight progress callback raise an ordinary
    ``RuntimeError``. ``_run_job``'s generic exception handler used to write
    ``failed`` unconditionally, so a user's deliberate cancel was reported back
    to them as a build error.
    """
    import task_manager as task_manager_module
    from task_manager import TaskManager

    course = legacy_course()
    document = document_from_legacy_course(course)
    storage = MemoryStorage({
        **course,
        "course_schema_version": COURSE_DOCUMENT_SCHEMA,
        "course_document": document.model_dump(mode="json"),
        "course_document_revision": document.document_revision,
        "course_document_authoritative": True,
        "course_operation_log": [],
        "learning_assets": course_data_with_practice()["learning_assets"],
    })
    monkeypatch.setattr(task_manager_module, "TASKS_FILE", tmp_path / "generation_jobs.json")
    monkeypatch.setattr(
        task_manager_module,
        "teaching_representation_repository",
        TeachingRepresentationRepository(tmp_path / "representations"),
    )
    manager = TaskManager(
        storage,
        course_service=None,
        ws_service=None,
        document_repository=CourseDocumentRepository(storage),
    )
    task_id = await manager.create_task(
        "course-1", "teaching_representation_build", enqueue=False,
    )

    # Cancel the task the moment the first page is emitted, exactly as a user
    # pressing cancel mid-build would.
    original_record = manager._record_representation_event
    cancelled_at: list[str] = []

    async def cancel_on_first_slide(inner_task_id: str, payload: dict) -> None:
        await original_record(inner_task_id, payload)
        if payload.get("event") == "slide_upsert" and not cancelled_at:
            cancelled_at.append(str(payload.get("event")))
            manager.tasks[inner_task_id]["status"] = "cancelled"
            manager.tasks[inner_task_id]["message"] = "任务已取消，正在清理生成状态"

    monkeypatch.setattr(manager, "_record_representation_event", cancel_on_first_slide)

    await manager._run_job(task_id)

    task = manager.get_task(task_id)
    assert cancelled_at, "the build never reached a slide_upsert event"
    assert task["status"] == "cancelled"
    assert task["status"] != "failed"
    assert not task.get("error")
    # The cancelled build must not have published a completion event either.
    assert not any(
        item["event"] == "build_complete"
        for item in task.get("event_history") or []
    )


@pytest.mark.asyncio
async def test_revision_listener_asynchronously_reconciles_stale_units(tmp_path):
    from course_repository import (
        register_course_revision_listener,
        unregister_course_revision_listener,
    )
    from representation_reconciliation import RepresentationReconciliationService

    course = legacy_course()
    document = document_from_legacy_course(course)
    storage = MemoryStorage({
        **course,
        "course_schema_version": COURSE_DOCUMENT_SCHEMA,
        "course_document": document.model_dump(mode="json"),
        "course_document_revision": document.document_revision,
        "course_document_authoritative": True,
        "course_operation_log": [],
    })
    course_repository = CourseDocumentRepository(storage)
    representation_repository = TeachingRepresentationRepository(tmp_path / "representations")
    compile_core_representations(document, course_data_with_practice(), representation_repository)
    service = RepresentationReconciliationService(course_repository, representation_repository)
    register_course_revision_listener(service.enqueue)
    await service.start()
    try:
        target = document.blocks[0]
        await CourseCommandService(course_repository).replace_block(
            "course-1",
            command_id="async-reconcile-1",
            expected_document_revision=document.document_revision,
            expected_block_revision=target.internal_revision,
            block_id=target.block_id,
            payload={**target.payload, "markdown": "异步消费者应标记这一段的派生产物。"},
        )
        await service._queue.join()
        registry = representation_repository.load("course-1")
        handout = next(
            item for item in registry.representations
            if item.representation_type == "handout"
        )
        assert handout.status == "stale"
        assert handout.stale_unit_ids == ["handout:section-a"]
    finally:
        unregister_course_revision_listener(service.enqueue)
        await service.shutdown()
