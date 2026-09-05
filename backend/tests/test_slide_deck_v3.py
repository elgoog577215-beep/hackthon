from __future__ import annotations

from copy import deepcopy
from pathlib import Path
from tempfile import TemporaryDirectory

import pytest
from fastapi import FastAPI
from fastapi.testclient import TestClient
from pptx import Presentation

from course_document import COURSE_DOCUMENT_SCHEMA, document_from_legacy_course
from course_repository import CourseDocumentRepository
from course_revisions import revision_vector_for_document
from representation_compiler import (
    compile_core_representations,
    export_slide_deck_pptx,
    rebuild_slide_deck_variant_bundle_safely,
    rebuild_slide_deck_variant_safely,
)
from slide_deck_v3 import (
    MAX_GENERATED_SLIDE_COUNT,
    SLIDE_DECK_THEMES,
    SlideAllocationPlanV2,
    _paginate_fragments,
    compile_slide_deck_v3,
    deterministic_slide_allocation,
    fragment_course_document,
    plan_slide_deck_v3,
    slide_deck_preflight_quality,
    split_slide_deck_plan_by_chapter,
)
from teaching_representations import TeachingRepresentationRepository


class MemoryStorage:
    def __init__(self, course: dict) -> None:
        self.course = deepcopy(course)

    def load_course(self, _course_id: str) -> dict:
        return deepcopy(self.course)

    async def save_course(self, _course_id: str, data: dict) -> None:
        self.course = deepcopy(data)


def source_course() -> dict:
    return {
        "course_id": "source-first-course",
        "course_name": "正文驱动课件",
        "nodes": [
            {
                "node_id": "section-core",
                "parent_node_id": "root",
                "node_name": "第一章 核心内容",
                "node_level": 1,
                "learning_objective": "理解正文驱动的课件生成",
                "objective_id": "objective-core",
                "content_blocks": [
                    {
                        "block_id": "block-core",
                        "title": "核心概念",
                        "content": (
                            "课程正文是唯一内容源。\n\n"
                            "- 页面只能引用原文片段\n"
                            "- 模型不能重写教学正文\n\n"
                            "页面过长时应当在安全边界继续分页。"
                        ),
                        "metadata": {"role": "concept"},
                    },
                    {
                        "block_id": "block-check",
                        "title": "理解检查",
                        "content": "为什么页面计划只能保存 fragment_id？",
                        "metadata": {"role": "checkpoint"},
                    },
                    {
                        "block_id": "block-supplement",
                        "title": "补充说明",
                        "content": "失败版本不能覆盖上一份可用课件。",
                        "metadata": {"role": "remediation"},
                    },
                ],
            },
        ],
    }


def multi_chapter_course(chapter_count: int = 4) -> dict:
    return {
        "course_id": "chapter-bundle-course",
        "course_name": "分册课程",
        "nodes": [
            {
                "node_id": f"chapter-{index}",
                "parent_node_id": "root",
                "node_name": f"第 {index} 章",
                "node_level": 1,
                "learning_objective": f"掌握第 {index} 章",
                "objective_id": f"objective-{index}",
                "content_blocks": [{
                    "block_id": f"chapter-{index}-body",
                    "title": f"第 {index} 章核心内容",
                    "content": (
                        f"## 第 {index} 章概念\n\n"
                        f"这是第 {index} 章的课程正文与方法说明。\n\n"
                        f"## 第 {index} 章深入理解\n\n"
                        f"继续完成第 {index} 章的知识解释与理解检查。"
                    ),
                    "metadata": {"role": "concept"},
                }],
            }
            for index in range(1, chapter_count + 1)
        ],
    }


def course_with_ready_slide_story_inputs(course: dict) -> dict:
    """Promote a legacy fixture without weakening the production preflight."""
    from course_logic_upgrade import compile_course_logic_upgrade

    promoted = deepcopy(course)
    for node in promoted.get("nodes") or []:
        # These legacy fixtures model root-level teaching sections as level one.
        # The official teaching-plan contract is intentionally section-based.
        if int(node.get("node_level") or 1) == 1:
            node["node_level"] = 2
        node["node_content"] = str(node.get("node_content") or "").strip() or "\n\n".join(
            str(block.get("content") or "")
            for block in node.get("content_blocks") or []
        )
    promoted.update(compile_course_logic_upgrade(promoted)["updates"])
    return promoted


def narrative_course() -> dict:
    return {
        "course_id": "narrative-course",
        "course_name": "线性代数的结构与方法",
        "nodes": [
            {
                "node_id": "chapter-vectors",
                "parent_node_id": "root",
                "node_name": "第一章 向量空间",
                "node_level": 1,
                "content_blocks": [{
                    "block_id": "chapter-vectors-summary",
                    "title": "正文",
                    "content": "从向量空间的结构出发，建立后续线性变换的语言基础。",
                    "metadata": {"role": "orientation"},
                }],
            },
            {
                "node_id": "topic-definition",
                "parent_node_id": "chapter-vectors",
                "node_name": "1.1 向量空间的定义",
                "node_level": 2,
                "content_blocks": [{
                    "block_id": "topic-definition-summary",
                    "title": "正文",
                    "content": "先识别研究对象，再理解公理为何能够统一不同实例。",
                    "metadata": {"role": "concept"},
                }],
            },
            {
                "node_id": "detail-definition",
                "parent_node_id": "topic-definition",
                "node_name": "1.1.1 从实例到公理",
                "node_level": 3,
                "content_blocks": [{
                    "block_id": "detail-definition-body",
                    "title": "正文",
                    "content": (
                        "### 核心概念与背景\n"
                        "向量空间把看似不同的对象放进同一套运算规则中。\n\n"
                        "### 深度原理与底层机制\n"
                        "封闭性保证运算不会离开研究对象，分配律保证结构可以稳定组合。\n\n"
                        "### 技术实现与方法论\n"
                        "- 检查集合是否封闭\n"
                        "- 检查加法与数乘公理\n"
                        "- 判断子集是否构成子空间\n\n"
                        "### 实战案例与行业应用\n"
                        "多项式集合与矩阵集合都可以按同一组公理验证为向量空间。\n\n"
                        "### 思考与挑战\n"
                        "如果集合对数乘不封闭，它还可能是向量空间吗？\n\n"
                        "### 延伸阅读\n"
                        "向量空间概念还可以继续推广到更一般的代数结构。"
                    ),
                    "metadata": {"role": "concept"},
                }],
            },
            {
                "node_id": "chapter-transformations",
                "parent_node_id": "root",
                "node_name": "第二章 线性变换",
                "node_level": 1,
                "content_blocks": [{
                    "block_id": "chapter-transformations-summary",
                    "title": "正文",
                    "content": "把向量空间的结构映射到新的表示中。",
                    "metadata": {"role": "orientation"},
                }],
            },
        ],
    }


def mermaid_and_formula_course() -> dict:
    return {
        "course_id": "typed-visual-source",
        "course_name": "Thermodynamics",
        "nodes": [{
            "node_id": "chapter-system",
            "parent_node_id": "root",
            "node_name": "Systems and energy",
            "node_level": 1,
            "content_blocks": [
                {
                    "block_id": "mermaid-flow",
                    "title": "System classification",
                    "content": (
                        "A closed system is classified by its exchange with the environment."
                        "\n\n#### ### 🎨 可视化图解"
                        "\n```mermaid"
                        "\ngraph TD"
                        "\nA[Closed system] -->|cannot exchange matter| B[Environment]"
                        "\n```"
                        "\n\nThis relation distinguishes a closed system from an open system."
                    ),
                    "metadata": {"role": "concept"},
                },
                {
                    "block_id": "formula-explanation",
                    "title": "Internal energy change",
                    "content": (
                        "Internal energy change is measured between two states."
                        "\n\n- State-difference equation:"
                        "\n\n$$\\Delta U = U_2 - U_1$$"
                        "\n\nThe sign records whether the final state has more or less energy."
                    ),
                    "metadata": {"role": "reasoning"},
                },
            ],
        }],
    }


def test_fragmenter_preserves_source_headings_as_semantic_boundaries() -> None:
    document = document_from_legacy_course(narrative_course())
    fragments = fragment_course_document(document)
    headings = [
        item.text for item in fragments
        if item.block_id == "detail-definition-body" and item.kind == "heading"
    ]

    assert headings == [
        "核心概念与背景",
        "深度原理与底层机制",
        "技术实现与方法论",
        "实战案例与行业应用",
        "思考与挑战",
        "延伸阅读",
    ]


def test_fragmenter_preserves_mermaid_as_diagram_and_drops_visual_marker() -> None:
    document = document_from_legacy_course(mermaid_and_formula_course())
    fragments = [
        item
        for item in fragment_course_document(document)
        if item.block_id == "mermaid-flow"
    ]

    assert [item.kind for item in fragments] == [
        "paragraph",
        "diagram",
        "paragraph",
    ]
    assert all("可视化图解" not in item.text for item in fragments)
    assert fragments[1].language == "mermaid"
    assert fragments[1].text.startswith("graph TD")


def test_fragmenter_drops_quoted_diagram_id_authoring_metadata() -> None:
    course = source_course()
    course["nodes"][0]["content_blocks"][0]["content"] = (
        "课程正文介绍系统分类。\n\n"
        "```mermaid\n"
        "graph TD\n"
        '    A["系统"] --> B["分类"]\n'
        "```\n\n"
        '> ID: "ThermodynamicSystemClassification"\n\n'
        "#### 实战案例/行业应用\n\n"
        "空调循环用于说明开放系统。"
    )
    document = document_from_legacy_course(course)

    fragments = fragment_course_document(document)
    visible_text = "\n".join(fragment.text for fragment in fragments)

    assert 'ID: "ThermodynamicSystemClassification"' not in visible_text
    assert any(
        fragment.kind == "heading"
        and fragment.text == "实战案例/行业应用"
        for fragment in fragments
    )


@pytest.mark.parametrize(
    ("kind", "content", "expected_text"),
    [
        (
            "paragraph",
            "没有自然句号的Unity状态同步说明" * 50,
            "没有自然句号的Unity状态同步说明" * 50,
        ),
        (
            "list_item",
            "- " + "单个列表项同时解释输入状态和输出状态" * 45,
            "单个列表项同时解释输入状态和输出状态" * 45,
        ),
    ],
)
def test_fragmenter_partitions_oversized_atomic_prose_without_losing_text(
    kind: str,
    content: str,
    expected_text: str,
) -> None:
    course = source_course()
    course["nodes"][0]["content_blocks"][0]["content"] = content
    document = document_from_legacy_course(course)

    fragments = [
        fragment
        for fragment in fragment_course_document(document)
        if fragment.block_id == "block-core" and fragment.kind == kind
    ]

    assert len(fragments) > 1
    assert all(len(fragment.text) <= 230 for fragment in fragments)
    assert "".join(fragment.text for fragment in fragments) == expected_text


def test_fragmenter_partitions_oversized_display_formula_at_safe_boundaries() -> None:
    course = source_course()
    formula_body = " + ".join(f"x_{{{index}}}" for index in range(1, 90)) + " = 0"
    course["nodes"][0]["content_blocks"][0]["content"] = f"$${formula_body}$$"
    document = document_from_legacy_course(course)

    fragments = [
        fragment
        for fragment in fragment_course_document(document)
        if fragment.block_id == "block-core" and fragment.kind == "formula"
    ]

    assert len(fragments) > 1
    assert all(len(fragment.text) <= 230 for fragment in fragments)
    assert all(
        fragment.text.startswith("$$") and fragment.text.endswith("$$")
        for fragment in fragments
    )
    assert "".join(fragment.text[2:-2] for fragment in fragments) == formula_body


def test_pagination_keeps_enumeration_unit_and_next_heading_with_its_body() -> None:
    course = source_course()
    course["nodes"][0]["content_blocks"][0]["content"] = (
        "### 1.1 热力学系统的分类与描述\n\n"
        "#### 核心概念与背景\n\n"
        "系统是研究的物理对象或区域，环境是系统以外的部分，二者通过边界区分。\n\n"
        "根据系统与环境之间的交互方式，热力学将系统分为三类：\n\n"
        "- 孤立系统：既不交换物质，也不交换能量。\n"
        "- 封闭系统：不交换物质，但可以交换能量。\n"
        "- 开放系统：既可以交换物质，也可以交换能量。\n\n"
        "这些分类是后续建立模型和分析的基础。\n\n"
        "#### 深度原理/底层机制\n\n"
        "系统边界决定了系统是否能与环境交换物质和能量。"
    )
    document = document_from_legacy_course(course)
    fragments = [
        fragment
        for fragment in fragment_course_document(document)
        if fragment.block_id == "block-core"
    ]

    pages = _paginate_fragments(fragments, 1100)
    page_by_text = {
        fragment.text: page_index
        for page_index, page in enumerate(pages)
        for fragment in page
    }

    enumeration_texts = [
        "根据系统与环境之间的交互方式，热力学将系统分为三类：",
        "孤立系统：既不交换物质，也不交换能量。",
        "封闭系统：不交换物质，但可以交换能量。",
        "开放系统：既可以交换物质，也可以交换能量。",
        "这些分类是后续建立模型和分析的基础。",
    ]
    assert len({page_by_text[text] for text in enumeration_texts}) == 1
    assert (
        page_by_text["系统是研究的物理对象或区域，环境是系统以外的部分，二者通过边界区分。"]
        != page_by_text[enumeration_texts[0]]
    )
    assert (
        page_by_text["深度原理/底层机制"]
        == page_by_text["系统边界决定了系统是否能与环境交换物质和能量。"]
    )
    assert page_by_text["深度原理/底层机制"] != page_by_text[enumeration_texts[0]]
    assert all(page[-1].kind != "heading" for page in pages)


def test_formula_is_not_paginated_without_adjacent_explanation() -> None:
    document = document_from_legacy_course(mermaid_and_formula_course())
    fragments = [
        item
        for item in fragment_course_document(document)
        if item.block_id == "formula-explanation"
    ]

    # Teaching-mode promotion can end the mainline run at the formula while
    # moving the later interpretation to detail pages.
    pages = _paginate_fragments(fragments[:-1], 1200)
    formula_page = next(
        page for page in pages
        if any(fragment.kind == "formula" for fragment in page)
    )

    assert any(
        fragment.kind in {"paragraph", "list_item"}
        for fragment in formula_page
    )
    assert [fragment.kind for fragment in formula_page] in (
        ["paragraph", "formula"],
        ["list_item", "formula"],
        ["formula", "paragraph"],
        ["formula", "list_item"],
        ["paragraph", "formula", "paragraph"],
    )


def test_preflight_blocks_decks_that_require_chapter_splitting() -> None:
    document = document_from_legacy_course(source_course())
    fragments = fragment_course_document(document)
    plan = deterministic_slide_allocation(
        document,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    template = plan.pages[-1]
    while len(plan.pages) <= MAX_GENERATED_SLIDE_COUNT:
        clone = template.model_copy(deep=True)
        clone.page_id = f"slide:oversized:{len(plan.pages)}"
        plan.pages.append(clone)

    quality = slide_deck_preflight_quality(plan)

    assert quality["passed"] is False
    assert quality["blockers"] == [{
        "severity": "critical",
        "code": "deck_split_required",
        "message": "课件预计页数超过单次生成上限，请按章节拆分后生成。",
        "estimated_slide_count": len(plan.pages),
        "maximum_slide_count": MAX_GENERATED_SLIDE_COUNT,
    }]


def test_oversized_plan_splits_on_chapter_boundaries_without_losing_source() -> None:
    course = multi_chapter_course()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    plan = deterministic_slide_allocation(
        document,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )

    parts = split_slide_deck_plan_by_chapter(
        document,
        plan,
        maximum_slide_count=8,
    )

    assert len(parts) > 1
    assert all(len(part.allocation_plan.pages) <= 8 for part in parts)
    part_fragment_ids = [
        {
            fragment.fragment_id
            for fragment in fragment_course_document(part.document)
        }
        for part in parts
    ]
    assert set().union(*part_fragment_ids) == {
        fragment.fragment_id for fragment in fragments
    }
    assert sum(len(values) for values in part_fragment_ids) == len(fragments)
    assert all(
        part.allocation_plan.source_document_revision
        == part.document.document_revision
        for part in parts
    )
    assert all(
        part.document.document_revision == document.document_revision
        for part in parts
    )


def test_oversized_variant_publishes_atomic_ready_part_representations(
    tmp_path: Path,
) -> None:
    course = multi_chapter_course()
    document = document_from_legacy_course(course)
    plan = deterministic_slide_allocation(
        document,
        fragment_course_document(document),
        mode="teaching",
        theme="qizhi-classroom",
    )
    parts = split_slide_deck_plan_by_chapter(
        document,
        plan,
        maximum_slide_count=8,
    )
    repository = TeachingRepresentationRepository(tmp_path / "registry")
    events: list[dict] = []

    result = rebuild_slide_deck_variant_bundle_safely(
        document,
        course,
        repository,
        mode="teaching",
        theme="qizhi-classroom",
        parts=parts,
        progress_callback=events.append,
    )

    assert result["status"] == "synchronized"
    assert result["bundle"] is True
    assert result["part_count"] == len(parts)
    registry = repository.load(document.course_id)
    representations = [
        item for item in registry.representations
        if item.variant_key.startswith("teaching:qizhi-classroom:part:")
    ]
    assert len(representations) == len(parts)
    assert all(item.status == "ready" for item in representations)
    specs = {
        spec.spec_id: spec
        for spec in registry.specs
    }
    assert [
        specs[item.spec_id].payload["content"]["bundle_part"]["part_index"]
        for item in representations
    ] == list(range(1, len(parts) + 1))
    assert all(
        any(
            binding.source_revisions.get("course_document")
            == document.document_revision
            for binding in specs[item.spec_id].source_bindings
        )
        for item in representations
    )
    canonical_title_revision = revision_vector_for_document(document).revisions["course_title"]
    assert all(
        any(
            binding.source_revisions.get("course_title")
            == canonical_title_revision
            for binding in specs[item.spec_id].source_bindings
        )
        for item in representations
    )
    assert any(event["event"] == "bundle_plan" for event in events)
    assert sum(event["event"] == "bundle_part_complete" for event in events) == len(parts)


def test_teaching_plan_builds_a_chapter_level_learning_progression() -> None:
    course = narrative_course()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    plan = deterministic_slide_allocation(
        document,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    detail_ids = {
        item.fragment_id
        for item in fragments
        if item.block_id == "detail-definition-body"
    }
    detail_pages = [
        page for page in plan.pages
        if detail_ids & set(page.fragment_ids)
    ]
    mainline_detail_pages = [page for page in detail_pages if not page.appendix]
    appendix_detail_pages = [page for page in detail_pages if page.appendix]

    first_chapter_index = next(
        index for index, page in enumerate(plan.pages)
        if page.page_id == "slide:chapter:chapter-vectors"
    )
    first_topic_index = next(
        index for index, page in enumerate(plan.pages)
        if set(page.fragment_ids) & detail_ids and not page.appendix
    )
    first_recap_index = next(
        index for index, page in enumerate(plan.pages)
        if page.page_id == "slide:chapter-recap:chapter-vectors"
    )
    second_chapter_index = next(
        index for index, page in enumerate(plan.pages)
        if page.page_id == "slide:chapter:chapter-transformations"
    )

    assert first_chapter_index < first_topic_index < first_recap_index
    assert first_recap_index < second_chapter_index
    assert mainline_detail_pages
    assert appendix_detail_pages
    assert {
        page.narrative_role for page in mainline_detail_pages
    } >= {"concept", "reasoning", "example", "checkpoint"}
    assert len({page.layout for page in mainline_detail_pages}) >= 3
    assert {
        fragment_id
        for page in plan.pages
        for fragment_id in page.fragment_ids
    } == {item.fragment_id for item in fragments}


def test_source_heading_becomes_slide_claim_without_body_duplication() -> None:
    course = narrative_course()
    document = document_from_legacy_course(course)
    plan = deterministic_slide_allocation(
        document,
        fragment_course_document(document),
        mode="teaching",
        theme="qizhi-classroom",
    )
    content = compile_slide_deck_v3(
        document,
        course,
        mode="teaching",
        theme="qizhi-classroom",
        allocation_plan=plan,
    )
    reasoning = next(
        slide for slide in content["slides"]
        if slide["slide_purpose"] == "reasoning"
    )
    body_text = "\n".join(
        text
        for block in reasoning["blocks"]
        for text in [block.get("content", ""), *(block.get("items") or [])]
    )

    assert reasoning["title"] == "深度原理与底层机制"
    assert "深度原理与底层机制" not in body_text
    assert "封闭性保证运算不会离开研究对象" in body_text


def test_consecutive_source_headings_do_not_create_empty_teaching_pages() -> None:
    course = source_course()
    course["nodes"].append({
        "node_id": "section-reasoning-detail",
        "parent_node_id": "section-core",
        "node_name": "Reasoning detail",
        "node_level": 3,
        "content_blocks": [{
            "block_id": "block-reasoning-detail",
            "title": "正文",
            "content": (
                "### 深度原理与底层机制\n"
                "#### 向量空间的公理体系\n"
                "封闭性、结合律与分配律共同保证结构可以稳定运算。"
            ),
            "metadata": {"role": "concept"},
        }],
    })
    document = document_from_legacy_course(course)
    plan = deterministic_slide_allocation(
        document,
        fragment_course_document(document),
        mode="teaching",
        theme="qizhi-classroom",
    )
    content = compile_slide_deck_v3(
        document,
        course,
        mode="teaching",
        theme="qizhi-classroom",
        allocation_plan=plan,
    )
    reasoning_slides = [
        slide for slide in content["slides"]
        if slide["slide_purpose"] == "reasoning"
    ]

    assert reasoning_slides
    assert all(slide["blocks"] for slide in reasoning_slides)
    assert any(
        "向量空间的公理体系" in block.get("content", "")
        for slide in reasoning_slides
        for block in slide["blocks"]
    )


@pytest.mark.parametrize("mode", ["full", "teaching"])
def test_full_coverage_modes_materialize_every_source_fragment(mode: str) -> None:
    course = source_course()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    plan = deterministic_slide_allocation(
        document,
        fragments,
        mode=mode,
        theme="qizhi-classroom",
    )
    content = compile_slide_deck_v3(
        document,
        course,
        mode=mode,
        theme="qizhi-classroom",
        allocation_plan=plan,
    )

    assert content["schema_version"] == "slide_deck_v3"
    assert content["coverage_report"]["visible_coverage_ratio"] == 1.0
    assert content["coverage_report"]["hash_integrity_passed"] is True
    assert content["quality_summary"]["passed"] is True
    assert content["quality_summary"]["score"] >= 90
    assert len([
        issue
        for issue in content["quality_report"]["warnings"]
        if issue["code"] == "knowledge_binding_missing"
    ]) <= 1
    referenced = {
        fragment_id
        for page in content["allocation_plan"]["pages"]
        for fragment_id in page["fragment_ids"]
    }
    assert referenced == {item.fragment_id for item in fragments}
    if mode == "teaching":
        appendix_pages = [
            item for item in content["allocation_plan"]["pages"]
            if item["appendix"] and item["fragment_ids"]
        ]
        assert appendix_pages


def test_concise_mode_records_every_omitted_fragment() -> None:
    course = source_course()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    plan = deterministic_slide_allocation(
        document,
        fragments,
        mode="concise",
        theme="modern-geometric",
    )
    content = compile_slide_deck_v3(
        document,
        course,
        mode="concise",
        theme="modern-geometric",
        allocation_plan=plan,
    )

    included = set(content["coverage_report"]["included_fragment_ids"])
    excluded = set(content["coverage_report"]["excluded_fragment_ids"])
    assert included.isdisjoint(excluded)
    assert included | excluded == {item.fragment_id for item in fragments}
    assert content["coverage_report"]["decision_coverage_ratio"] == 1.0
    assert all(item["reason"] == "mode_concise" for item in content["exclusions"])


def test_process_content_uses_cumulative_reveal_pages_without_double_counting_coverage() -> None:
    course = source_course()
    course["nodes"][0]["content_blocks"].insert(1, {
        "block_id": "block-process",
        "title": "生成流程",
        "content": (
            "- 切分课程正文\n"
            "- 分配页面内容\n"
            "- 检查页面容量\n"
            "- 发布可用版本"
        ),
        "metadata": {"role": "reasoning"},
    })
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    plan = deterministic_slide_allocation(
        document,
        fragments,
        mode="full",
        theme="qizhi-classroom",
    )
    sequence_pages = [page for page in plan.pages if page.sequence_id]

    assert len(sequence_pages) == 4
    assert [page.step_index for page in sequence_pages] == [1, 2, 3, 4]
    final_fragment_ids = sequence_pages[-1].fragment_ids
    assert [
        page.fragment_ids for page in sequence_pages
    ] == [
        final_fragment_ids[:1],
        final_fragment_ids[:2],
        final_fragment_ids[:3],
        final_fragment_ids[:4],
    ]

    content = compile_slide_deck_v3(
        document,
        course,
        mode="full",
        theme="qizhi-classroom",
        allocation_plan=plan,
    )

    assert content["coverage_report"]["visible_coverage_ratio"] == 1.0
    assert content["coverage_report"]["included_fragment_count"] == len(fragments)
    built_sequence = [
        slide for slide in content["slides"]
        if slide["quality"].get("sequence_id") == sequence_pages[0].sequence_id
    ]
    assert [slide["quality"]["step_index"] for slide in built_sequence] == [1, 2, 3, 4]


def test_inline_math_stays_in_prose_and_visible_text_contains_no_markup() -> None:
    course = source_course()
    course["nodes"][0]["content_blocks"][0].update({
        "title": "正文",
        "content": (
            "常见向量空间包括 $\\mathbb{R}^n$ 和矩阵空间。\n\n"
            "<!-- BODY_START -->\n"
            "**向量空间**需要满足封闭性。"
        ),
    })
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)

    assert all(item.kind != "formula" for item in fragments if "常见向量空间" in item.text)

    plan = deterministic_slide_allocation(
        document,
        fragments,
        mode="full",
        theme="qizhi-classroom",
    )
    content = compile_slide_deck_v3(
        document,
        course,
        mode="full",
        theme="qizhi-classroom",
        allocation_plan=plan,
    )
    visible = "\n".join(
        value
        for slide in content["slides"]
        for block in slide["blocks"]
        for value in [block.get("content", ""), *(block.get("items") or [])]
    )

    assert "\\mathbb" not in visible
    assert "<!--" not in visible
    assert "**" not in visible
    assert "ℝⁿ" in visible


def test_placeholder_block_title_uses_semantic_section_title() -> None:
    course = source_course()
    course["nodes"][0]["content_blocks"][0]["title"] = "正文"
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    plan = deterministic_slide_allocation(
        document,
        fragments,
        mode="full",
        theme="qizhi-classroom",
    )
    content = compile_slide_deck_v3(
        document,
        course,
        mode="full",
        theme="qizhi-classroom",
        allocation_plan=plan,
    )
    teaching_slides = [
        slide for slide in content["slides"]
        if slide.get("source_block_ids") == ["block-core"]
    ]

    assert teaching_slides
    assert all(slide["title"] != "正文" for slide in teaching_slides)
    assert teaching_slides[0]["title"] == "第一章 核心内容"


def test_deterministic_pages_respect_materialized_layout_capacity() -> None:
    course = source_course()
    course["nodes"][0]["content_blocks"][0]["content"] = "\n\n".join(
        f"第 {index} 个概念用于验证页面容量。"
        for index in range(1, 8)
    )
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    plan = deterministic_slide_allocation(
        document,
        fragments,
        mode="full",
        theme="qizhi-classroom",
    )
    content = compile_slide_deck_v3(
        document,
        course,
        mode="full",
        theme="qizhi-classroom",
        allocation_plan=plan,
    )

    concept_slides = [
        slide for slide in content["slides"]
        if slide["layout"] == "concept"
    ]
    assert concept_slides
    assert all(len(slide["blocks"]) <= 3 for slide in concept_slides)
    assert not {
        issue["code"]
        for issue in content["quality_report"]["blockers"]
    } & {"slide_block_overflow", "concept_card_overflow"}


def test_sequential_prose_is_grouped_into_readable_body_blocks() -> None:
    course = source_course()
    course["nodes"][0]["content_blocks"][0]["content"] = "\n\n".join(
        f"Paragraph {index} explains one connected part of the same concept."
        for index in range(1, 9)
    )
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    plan = deterministic_slide_allocation(
        document,
        fragments,
        mode="full",
        theme="qizhi-classroom",
    )
    content = compile_slide_deck_v3(
        document,
        course,
        mode="full",
        theme="qizhi-classroom",
        allocation_plan=plan,
    )
    teaching_slides = [
        slide for slide in content["slides"]
        if slide.get("source_block_ids") == ["block-core"]
    ]

    assert len(teaching_slides) <= 3
    assert all(len(slide["blocks"]) <= 2 for slide in teaching_slides)
    assert content["quality_report"]["passed"] is True


def test_common_latex_is_translated_without_leaking_commands() -> None:
    course = source_course()
    course["nodes"][0]["content_blocks"][0]["content"] = (
        r"The eigenvalue equation is $\det(A-\lambda I)=0$, "
        r"and $\frac{1}{n}\sum_{i=1}^{n}x_i$ is an average."
    )
    document = document_from_legacy_course(course)
    plan = deterministic_slide_allocation(
        document,
        fragment_course_document(document),
        mode="full",
        theme="qizhi-classroom",
    )
    content = compile_slide_deck_v3(
        document,
        course,
        mode="full",
        theme="qizhi-classroom",
        allocation_plan=plan,
    )
    visible = "\n".join(
        value
        for slide in content["slides"]
        for block in slide["blocks"]
        for value in [block.get("content", ""), *(block.get("items") or [])]
    )

    assert "\\" not in visible
    assert "$" not in visible
    assert "λ" in visible
    assert "∑" in visible
    assert not {
        issue["code"] for issue in content["quality_report"]["blockers"]
    } & {"raw_latex_leaked"}


def test_teaching_mode_promotes_semantic_deep_content_and_keeps_remainder_in_appendix() -> None:
    course = source_course()
    course["nodes"].append({
        "node_id": "section-detail",
        "parent_node_id": "section-core",
        "node_name": "Proof details",
        "node_level": 3,
        "content_blocks": [{
            "block_id": "block-detail",
            "title": "Proof details",
            "content": (
                "### Core idea\n"
                "A definition used by the teaching mainline.\n\n"
                "### Deep derivation\n"
                "A detailed derivation retained verbatim for reference.\n\n"
                "### Further reading\n"
                "An optional historical note remains in the appendix."
            ),
            "metadata": {"role": "concept"},
        }],
    })
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    plan = deterministic_slide_allocation(
        document,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    detail_ids = {
        item.fragment_id for item in fragments if item.block_id == "block-detail"
    }
    detail_pages = [
        page for page in plan.pages if detail_ids & set(page.fragment_ids)
    ]

    assert detail_pages
    assert any(not page.appendix for page in detail_pages)
    assert any(page.appendix for page in detail_pages)


def test_teaching_appendix_uses_compact_readable_editorial_layout() -> None:
    course = source_course()
    course["nodes"].append({
        "node_id": "section-appendix",
        "parent_node_id": "section-core",
        "node_name": "Detailed appendix",
        "node_level": 3,
        "content_blocks": [{
            "block_id": "block-appendix",
            "title": "Detailed appendix",
            "content": "\n\n".join([
                *[
                    f"Supporting paragraph {index} remains verbatim and editable."
                    for index in range(1, 7)
                ],
                "- First reference condition",
                "- Second reference condition",
            ]),
            "metadata": {"role": "concept"},
        }],
    })
    document = document_from_legacy_course(course)
    plan = deterministic_slide_allocation(
        document,
        fragment_course_document(document),
        mode="teaching",
        theme="qizhi-classroom",
    )
    content = compile_slide_deck_v3(
        document,
        course,
        mode="teaching",
        theme="qizhi-classroom",
        allocation_plan=plan,
    )
    appendix_slides = [
        slide for slide in content["slides"]
        if slide.get("source_block_ids") == ["block-appendix"]
    ]

    assert len(appendix_slides) <= 2
    assert all(slide["layout"] == "appendix" for slide in appendix_slides)
    assert content["quality_report"]["passed"] is True


def test_summary_page_contains_source_derived_recap_content() -> None:
    course = source_course()
    document = document_from_legacy_course(course)
    plan = deterministic_slide_allocation(
        document,
        fragment_course_document(document),
        mode="full",
        theme="qizhi-classroom",
    )
    content = compile_slide_deck_v3(
        document,
        course,
        mode="full",
        theme="qizhi-classroom",
        allocation_plan=plan,
    )
    summary = next(
        slide for slide in content["slides"] if slide["unit_id"] == "slide:summary"
    )

    assert summary["blocks"]
    assert any(
        block.get("items") or block.get("content")
        for block in summary["blocks"]
    )
    assert summary["source_section_ids"]


@pytest.mark.asyncio
async def test_ai_planner_cannot_inject_teaching_body_text() -> None:
    course = source_course()
    document = document_from_legacy_course(course)

    async def invalid_planner(_request):
        return {
            "schema_version": "slide_allocation_plan_v2",
            "title": document.title,
            "mode": "teaching",
            "theme": "qizhi-classroom",
            "variant_key": "teaching:qizhi-classroom",
            "source_document_revision": document.document_revision,
            "pages": [{
                "page_id": "slide:title",
                "layout": "cover",
                "fragment_ids": [],
                "appendix": False,
                "sequence_id": "",
                "step_index": 0,
                "derived_text": [],
                "body": "模型擅自生成的教学正文",
            }],
            "exclusions": [],
            "planner": "ai",
            "fallback_reason": "",
            "review": {},
        }

    plan = await plan_slide_deck_v3(
        document,
        course,
        mode="teaching",
        theme="qizhi-classroom",
        ai_planner=invalid_planner,
    )

    assert isinstance(plan, SlideAllocationPlanV2)
    assert plan.planner == "deterministic_fallback"
    assert plan.fallback_reason == "invalid_or_failed_ai_plan"


def test_mode_theme_variants_are_cached_independently_and_do_not_replace_core_specs() -> None:
    course = source_course()
    document = document_from_legacy_course(course)
    with TemporaryDirectory() as temp_dir:
        repository = TeachingRepresentationRepository(temp_dir)
        compile_core_representations(document, course, repository)
        before = repository.load(document.course_id)
        core_ids = {
            item.representation_type: item.representation_id
            for item in before.representations
        }
        for mode in ("full", "teaching", "concise"):
            for theme in SLIDE_DECK_THEMES:
                fragments = fragment_course_document(document)
                plan = deterministic_slide_allocation(
                    document,
                    fragments,
                    mode=mode,
                    theme=theme,
                )
                result = rebuild_slide_deck_variant_safely(
                    document,
                    course,
                    repository,
                    mode=mode,
                    theme=theme,
                    allocation_plan=plan,
                )
                assert result["status"] == "synchronized"

        registry = repository.load(document.course_id)
        variants = [
            item for item in registry.representations
            if item.representation_type == "slide_deck" and item.variant_key
        ]
        assert len(variants) == 15
        assert len({item.representation_id for item in variants}) == 15
        for representation_type, representation_id in core_ids.items():
            assert any(
                item.representation_type == representation_type
                and item.representation_id == representation_id
                for item in registry.representations
            )


def test_failed_variant_rebuild_emits_the_structured_terminal_quality(
    monkeypatch,
    tmp_path: Path,
) -> None:
    import representation_compiler

    course = source_course()
    document = document_from_legacy_course(course)
    repository = TeachingRepresentationRepository(tmp_path / "registry")
    events: list[dict] = []

    def fail_compile(*_args, **_kwargs):
        raise ValueError("final SlideDeckContent schema mismatch")

    monkeypatch.setattr(
        representation_compiler,
        "compile_slide_deck_variant",
        fail_compile,
    )

    result = rebuild_slide_deck_variant_safely(
        document,
        course,
        repository,
        mode="teaching",
        theme="qizhi-classroom",
        progress_callback=events.append,
    )

    terminal = events[-1]
    assert terminal["event"] == "build_failed"
    assert terminal["code"] == "slide_variant_rebuild_failed"
    assert terminal["quality"] == result["quality"]
    assert terminal["quality"]["blockers"][0]["message"] == (
        "final SlideDeckContent schema mismatch"
    )


def test_requested_v5_without_story_plan_fails_closed_with_actionable_reason(
    tmp_path: Path,
) -> None:
    course = source_course()
    document = document_from_legacy_course(course)
    repository = TeachingRepresentationRepository(tmp_path / "registry")
    events: list[dict] = []

    result = rebuild_slide_deck_variant_safely(
        document,
        course,
        repository,
        mode="teaching",
        theme="qizhi-classroom",
        requested_schema="slide_deck_v5",
        story_plan=None,
        progress_callback=events.append,
    )

    assert result["candidate_status"] == "v5_failed"
    assert result["failure"] == {
        "stage": "source_preflight",
        "code": "v5_story_plan_missing",
        "message": "V5 构建缺少课程级 story plan。",
        "retryable": True,
        "source_revision": document.document_revision,
    }
    assert events[-1]["event"] == "build_failed"
    assert events[-1]["failure"] == result["failure"]
    assert not any(event.get("event") == "slide_upsert" for event in events)
    assert not repository.load(document.course_id).representations


def test_requested_v5_rejects_source_revision_change_before_atomic_commit(
    monkeypatch,
    tmp_path: Path,
) -> None:
    import representation_compiler

    course = source_course()
    document = document_from_legacy_course(course)
    repository = TeachingRepresentationRepository(tmp_path / "registry")
    events: list[dict] = []

    monkeypatch.setattr(
        representation_compiler,
        "compile_slide_deck_variant",
        lambda *_args, **_kwargs: {
            "status": "ready",
            "quality": {"passed": True, "score": 100, "blockers": []},
        },
    )

    result = rebuild_slide_deck_variant_safely(
        document,
        course,
        repository,
        mode="teaching",
        theme="qizhi-classroom",
        requested_schema="slide_deck_v5",
        story_plan={"schema_version": "slide_story_plan_v2"},
        progress_callback=events.append,
        source_revision_provider=lambda: "newer-course-revision",
    )

    assert result["candidate_status"] == "v5_failed"
    assert result["failure"]["code"] == "v5_source_revision_conflict"
    assert result["failure"]["stage"] == "source_commit"
    assert result["failure"]["retryable"] is True
    assert events[-1]["event"] == "build_failed"
    assert events[-1]["failure"] == result["failure"]
    assert not repository.load(document.course_id).representations


def test_requested_v5_rejects_a_non_v5_compiler_candidate(
    monkeypatch,
    tmp_path: Path,
) -> None:
    import representation_compiler
    from slide_deck_v5 import SlideDeckV5BuildError

    course = source_course()
    document = document_from_legacy_course(course)
    repository = TeachingRepresentationRepository(tmp_path / "registry")
    monkeypatch.setattr(
        representation_compiler,
        "slide_deck_v5_enabled",
        lambda: True,
    )
    monkeypatch.setattr(
        representation_compiler,
        "compile_slide_deck_v5",
        lambda *_args, **_kwargs: {
            "schema_version": "slide_deck_v4",
            "title": "invalid candidate",
            "slides": [],
            "quality_report": {},
            "quality_summary": {},
        },
    )
    monkeypatch.setattr(
        representation_compiler,
        "export_structured_slide_deck",
        lambda _content, path, **_kwargs: path,
    )
    monkeypatch.setattr(
        representation_compiler,
        "audit_exported_pptx",
        lambda *_args, **_kwargs: {
            "passed": True,
            "issues": [],
            "blockers": [],
        },
    )
    monkeypatch.setattr(
        representation_compiler,
        "validate_slide_deck_v5",
        lambda *_args, **_kwargs: {
            "passed": True,
            "score": 100,
            "issues": [],
            "blockers": [],
        },
    )

    with pytest.raises(SlideDeckV5BuildError) as captured:
        representation_compiler.compile_slide_deck_variant(
            document,
            course,
            repository,
            mode="teaching",
            theme="qizhi-classroom",
            requested_schema="slide_deck_v5",
            story_plan={"schema_version": "slide_story_plan_v2"},
        )

    assert captured.value.public_detail() == {
        "stage": "compiler",
        "code": "v5_schema_mismatch",
        "message": "V5 编译器返回了非 V5 候选。",
        "retryable": False,
        "source_revision": document.document_revision,
    }
    assert not repository.load(document.course_id).representations


def test_v3_export_is_editable_widescreen_and_uses_variant_theme(tmp_path: Path) -> None:
    course = source_course()
    document = document_from_legacy_course(course)
    repository = TeachingRepresentationRepository(tmp_path / "registry")
    plan = deterministic_slide_allocation(
        document,
        fragment_course_document(document),
        mode="teaching",
        theme="dark-tech",
    )
    result = rebuild_slide_deck_variant_safely(
        document,
        course,
        repository,
        mode="teaching",
        theme="dark-tech",
        allocation_plan=plan,
    )
    assert result["status"] == "synchronized"
    registry = repository.load(document.course_id)
    representation = next(
        item for item in registry.representations
        if item.variant_key == "teaching:dark-tech"
    )
    spec = next(item for item in registry.specs if item.spec_id == representation.spec_id)
    path = export_slide_deck_pptx(spec, tmp_path / "source-first.pptx")
    presentation = Presentation(path)

    assert presentation.slide_width / presentation.slide_height == pytest.approx(16 / 9, rel=0.01)
    assert len(presentation.slides) == len(spec.payload["content"]["slides"])
    assert any(
        shape.has_text_frame and shape.text.strip()
        for slide in presentation.slides
        for shape in slide.shapes
    )
    cover_sizes = [
        paragraph.font.size.pt
        for shape in presentation.slides[0].shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        if paragraph.font.size
    ]
    assert max(cover_sizes) >= 50
    non_cover_heading_sizes = [
        max(
            (
                paragraph.font.size.pt
                for shape in slide.shapes
                if shape.has_text_frame
                for paragraph in shape.text_frame.paragraphs
                if paragraph.font.size
            ),
            default=0,
        )
        for slide in list(presentation.slides)[1:]
    ]
    assert all(size >= 35 for size in non_cover_heading_sizes)
    appendix_index = next(
        index
        for index, slide in enumerate(spec.payload["content"]["slides"])
        if slide["layout"] == "appendix"
    )
    appendix_body_sizes = [
        paragraph.font.size.pt
        for shape in presentation.slides[appendix_index].shapes
        if shape.has_text_frame
        for paragraph in shape.text_frame.paragraphs
        if paragraph.font.size
    ]
    assert appendix_body_sizes
    assert 16 in appendix_body_sizes
    assert all(
        shape.left >= 0
        and shape.top >= 0
        and shape.left + shape.width <= presentation.slide_width
        and shape.top + shape.height <= presentation.slide_height
        for slide in presentation.slides
        for shape in slide.shapes
    )


def test_variant_stream_endpoint_builds_only_requested_combination(tmp_path: Path, monkeypatch) -> None:
    from routers import teaching_representations as representation_router

    course = course_with_ready_slide_story_inputs(source_course())
    document = document_from_legacy_course(course)
    canonical = {
        **course,
        "course_schema_version": COURSE_DOCUMENT_SCHEMA,
        "course_document": document.model_dump(mode="json"),
        "course_document_authoritative": True,
        "course_operation_log": [],
    }
    course_repository = CourseDocumentRepository(MemoryStorage(canonical))
    representation_repository = TeachingRepresentationRepository(tmp_path / "registry")
    monkeypatch.setattr(
        representation_router,
        "get_course_document_repository",
        lambda: course_repository,
    )
    monkeypatch.setattr(
        representation_router,
        "get_teaching_representation_repository",
        lambda: representation_repository,
    )
    monkeypatch.setattr(representation_router, "get_task_manager_optional", lambda: None)

    async def existing_course(_course_id: str):
        return course_repository.load_course_view(document.course_id)

    monkeypatch.setattr(representation_router, "get_course_or_404", existing_course)
    app = FastAPI()
    app.include_router(representation_router.router, prefix="/api")
    client = TestClient(app)

    with client.stream(
        "POST",
        f"/api/courses/{document.course_id}/teaching-representations/slide-decks/build/stream",
        headers={"X-User-Id": "teacher-1"},
        json={
            "mode": "teaching",
            "theme": "grid-notebook",
            "force_rebuild": False,
            "engine_version": "v5",
        },
    ) as response:
        stream = "".join(response.iter_text())

    assert response.status_code == 200
    assert "event: slide_upsert" in stream
    assert "event: build_complete" in stream
    registry = representation_repository.load(document.course_id)
    assert [
        (item.representation_type, item.variant_key)
        for item in registry.representations
    ] == [("slide_deck", "teaching:grid-notebook")]


def test_variant_stream_rebuilds_cached_v5_after_compiler_reliability_upgrade(
    tmp_path: Path,
    monkeypatch,
) -> None:
    import slide_deck_v5
    from course_logic_upgrade import compile_course_logic_upgrade
    from routers import teaching_representations as representation_router

    course = source_course()
    course["nodes"][0]["node_level"] = 2
    course["nodes"][0]["node_content"] = "\n\n".join(
        str(block.get("content") or "")
        for block in course["nodes"][0]["content_blocks"]
    )
    course.update(compile_course_logic_upgrade(course)["updates"])
    document = document_from_legacy_course(course)
    canonical = {
        **course,
        "course_schema_version": COURSE_DOCUMENT_SCHEMA,
        "course_document": document.model_dump(mode="json"),
        "course_document_authoritative": True,
        "course_operation_log": [],
    }
    course_repository = CourseDocumentRepository(MemoryStorage(canonical))
    representation_repository = TeachingRepresentationRepository(tmp_path / "registry")
    monkeypatch.setattr(
        representation_router,
        "get_course_document_repository",
        lambda: course_repository,
    )
    monkeypatch.setattr(
        representation_router,
        "get_teaching_representation_repository",
        lambda: representation_repository,
    )
    monkeypatch.setattr(representation_router, "get_task_manager_optional", lambda: None)

    async def existing_course(_course_id: str):
        return course_repository.load_course_view(document.course_id)

    monkeypatch.setattr(representation_router, "get_course_or_404", existing_course)
    app = FastAPI()
    app.include_router(representation_router.router, prefix="/api")
    client = TestClient(app)
    endpoint = (
        f"/api/courses/{document.course_id}"
        "/teaching-representations/slide-decks/build/stream"
    )
    request = {
        "mode": "teaching",
        "theme": "qizhi-classroom",
        "force_rebuild": False,
        "engine_version": "v5",
    }
    current_compiler_version = slide_deck_v5.SLIDE_DECK_V5_COMPILER_VERSION
    monkeypatch.setattr(
        slide_deck_v5,
        "SLIDE_DECK_V5_COMPILER_VERSION",
        "course_logic_slide_compiler_v5.5",
    )

    with client.stream(
        "POST",
        endpoint,
        headers={"X-User-Id": "teacher-1"},
        json=request,
    ) as response:
        first_stream = "".join(response.iter_text())

    assert response.status_code == 200
    assert "event: build_complete" in first_stream
    registry = representation_repository.load(document.course_id)
    spec = next(item for item in registry.specs if item.representation_type == "slide_deck")
    assert spec.payload["content"]["schema_version"] == "slide_deck_v5"

    monkeypatch.setattr(
        slide_deck_v5,
        "SLIDE_DECK_V5_COMPILER_VERSION",
        current_compiler_version,
    )
    with client.stream(
        "POST",
        endpoint,
        headers={"X-User-Id": "teacher-1"},
        json=request,
    ) as response:
        upgraded_stream = "".join(response.iter_text())

    assert response.status_code == 200
    assert '"cached": true' not in upgraded_stream
    assert "event: slide_upsert" in upgraded_stream


def test_variant_stream_publishes_and_reuses_an_atomic_v5_variant(
    tmp_path: Path,
    monkeypatch,
) -> None:
    from routers import teaching_representations as representation_router

    course = course_with_ready_slide_story_inputs(multi_chapter_course())
    document = document_from_legacy_course(course)
    canonical = {
        **course,
        "course_schema_version": COURSE_DOCUMENT_SCHEMA,
        "course_document": document.model_dump(mode="json"),
        "course_document_authoritative": True,
        "course_operation_log": [],
    }
    course_repository = CourseDocumentRepository(MemoryStorage(canonical))
    representation_repository = TeachingRepresentationRepository(tmp_path / "registry")
    monkeypatch.setattr(
        representation_router,
        "get_course_document_repository",
        lambda: course_repository,
    )
    monkeypatch.setattr(
        representation_router,
        "get_teaching_representation_repository",
        lambda: representation_repository,
    )
    monkeypatch.setattr(representation_router, "get_task_manager_optional", lambda: None)

    async def existing_course(_course_id: str):
        return course_repository.load_course_view(document.course_id)

    monkeypatch.setattr(representation_router, "get_course_or_404", existing_course)
    app = FastAPI()
    app.include_router(representation_router.router, prefix="/api")
    client = TestClient(app)
    endpoint = (
        f"/api/courses/{document.course_id}"
        "/teaching-representations/slide-decks/build/stream"
    )
    request = {
        "mode": "teaching",
        "theme": "qizhi-classroom",
        "force_rebuild": False,
        "engine_version": "v5",
    }

    with client.stream(
        "POST",
        endpoint,
        headers={"X-User-Id": "teacher-1"},
        json=request,
    ) as response:
        first_stream = "".join(response.iter_text())

    assert response.status_code == 200
    assert "event: bundle_plan" not in first_stream
    assert "event: build_complete" in first_stream
    registry = representation_repository.load(document.course_id)
    assert [
        item.variant_key
        for item in registry.representations
    ] == ["teaching:qizhi-classroom"]

    with client.stream(
        "POST",
        endpoint,
        headers={"X-User-Id": "teacher-1"},
        json=request,
    ) as response:
        cached_stream = "".join(response.iter_text())

    assert response.status_code == 200
    assert '"cached": true' in cached_stream
    assert "event: slide_upsert" not in cached_stream
