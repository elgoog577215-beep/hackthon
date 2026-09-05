from copy import deepcopy

import pytest
from pptx import Presentation
from pptx.util import Pt

from ppt_page_draft import lower_teaching_draft
from ppt_teaching_content import PageTeachingV2
from ppt_native_scene import audit_scene, render_scene
from .test_ppt_teaching_content import branch_fixture, scene_for


def draft_fixture():
    text = '先观察输入；符合条件时分类，不符合条件时继续观察。输入决定分类与观察两个分支。'
    citation = [{'block_id': 'b', 'quote': text}]
    return {'title': '观察输入后怎样选择？', 'page_goal': '理解任务分支', 'expression_kind': 'process',
        'elements': [{'key': k, 'text': t, 'sources': citation, 'show_from': 1} for k, t in [('input', '观察输入'), ('classify', '分类'), ('observe', '继续观察')]],
        'relations': [{'source_key': 'input', 'target_key': k, 'kind': 'sequence', 'sources': citation} for k in ['classify', 'observe']],
        'reveal_notes': ['比较两个分支']}, {'b': {'block_id': 'b', 'block_revision': 'r1', 'full_text': text}}


def test_compact_graph_assigns_source_and_preserves_branch_edges():
    value, source = draft_fixture()
    result = lower_teaching_draft(value, source)
    content = PageTeachingV2.model_validate(result['teaching'])
    assert [(r.source_id, r.target_id) for r in content.expression.relations] == [('input', 'classify'), ('input', 'observe')]
    assert all(e.sources[0].block_revision == 'r1' for e in content.elements)
    assert content.source_dispositions[0].element_ids == ['input', 'classify', 'observe']
    assert value['elements'][0]['key'] == 'input'


def test_compact_draft_never_drops_relation_or_fabricates_quote():
    value, source = draft_fixture()
    linear = deepcopy(value)
    linear['expression_kind'] = 'derivation'
    with pytest.raises(ValueError, match='linear_expression_cannot_discard_relations'):
        lower_teaching_draft(linear, source)
    value['elements'][0]['sources'] = [{'block_id': 'b', 'quote': '来源没有这句话'}]
    with pytest.raises(ValueError, match='source_excerpt_mismatch'):
        lower_teaching_draft(value, source)


def test_compact_exercise_pairs_multiple_answers_to_questions():
    value, source = draft_fixture()
    value['expression_kind'], value['relations'] = 'exercise', []
    citation = value['elements'][0]['sources']
    value['elements'] = [
        {'key': 'q1', 'text': '先做什么？', 'role': 'question', 'sources': citation},
        {'key': 'q2', 'text': '之后做什么？', 'role': 'question', 'sources': citation},
        {'key': 'a1', 'text': '观察输入', 'role': 'answer', 'sources': citation, 'show_from': 2, 'answers_question_id': 'q1'},
        {'key': 'a2', 'text': '分类', 'role': 'answer', 'sources': citation, 'show_from': 3, 'answers_question_id': 'q2'},
    ]
    value['reveal_notes'] = ['提出问题', '讨论第一个答案', '讨论第二个答案']
    content = lower_teaching_draft(value, source)['teaching']
    assert content['states'][0]['visible_element_ids'] == ['q1', 'q2']
    value['elements'][2]['show_from'] = 1
    with pytest.raises(ValueError, match='answer_revealed_before_question'):
        lower_teaching_draft(value, source)


def test_relation_label_geometry_font_and_anchor_are_read_back(tmp_path):
    value, _ = branch_fixture()
    value['expression']['relations'][0]['label'] = '分支'
    scene = scene_for(value, 'concept-map')
    deck = Presentation()
    deck.slide_width, deck.slide_height = Pt(960), Pt(540)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    render_scene(slide, scene)
    path = tmp_path / 'branch-label.pptx'
    deck.save(path)
    restored = Presentation(path).slides[0]
    assert audit_scene(restored, scene)['passed']
    label = next(s for s in restored.shapes if s.name == 'relation-label:ab')
    label.left += Pt(40)
    with pytest.raises(ValueError, match='export_relation_label_geometry_mismatch'):
        audit_scene(restored, scene)
    value['expression']['relations'][0]['label'] = '这是一段过长的分支关系标签' * 8
    with pytest.raises(ValueError, match='relation_label_capacity_exceeded'):
        scene_for(value, 'concept-map')


def test_narrative_repairs_missing_sources_and_keeps_failed_candidate():
    import asyncio
    from ppt_teaching_planner import plan_teaching_manuscript
    from .test_ppt_teaching_content import compiled_manuscript
    doc, graph, template, existing = compiled_manuscript()
    requests, saved = [], []
    page = existing.pages[0]
    plan = {'page_id': 'p1', 'teaching_unit_id': page.teaching_unit_id, 'source_block_ids': ['b'],
            'title': page.title, 'page_goal': page.page_goal, 'layout_id': page.layout_id}
    async def planner(request):
        requests.append(request)
        if request['teaching_request'] == 'narrative':
            planned = deepcopy(plan)
            if len(requests) == 1:
                planned['source_block_ids'] = ['missing']
            return {'narrative_brief': {'central_question': '比较执行方式'}, 'pages': [planned]}
        return {'title': page.title, 'page_goal': page.page_goal, 'teaching': page.teaching.model_dump(mode='json')}
    async def checkpoint(value, event):
        saved.append(deepcopy(value))
    result, trace = asyncio.run(plan_teaching_manuscript(doc, graph, template, planner, on_checkpoint=checkpoint))
    assert len(requests) == 3 and result.pages[0].title == page.title
    assert 'missing=' in requests[1]['validation_error']
    assert requests[1]['previous_candidate']['pages'][0]['source_block_ids'] == ['missing']
    assert saved[0]['signature'] == trace['signature']
    assert any(s.get('draft_narrative', {}).get('pages', [{}])[0].get('source_block_ids') == ['missing'] for s in saved)
    schema = requests[0]['response_contract']['$defs']['PlannedPage']['properties']
    assert 'audience_question' not in schema and 'source_block_ids' in schema


def test_exact_quote_selection_retains_duplicate_occurrence_and_rejects_conflicts():
    from ppt_source_quotes import source_excerpt_catalog
    from ppt_page_draft import bind_choices, draft_element_text, PageDraftElement
    from ppt_comparison_draft import QuoteChoice
    text = '例一：$x=1$。\n例二：$x=1$。\n```python\nx = 1\nprint(x)\n```'
    sources = {'b': {'block_id': 'b', 'block_revision': 'r1', 'full_text': text}}
    catalog = source_excerpt_catalog(sources)
    repeated = [q for q in catalog if q['quote'] == '$x=1$']
    selected = repeated[-1]
    choice = QuoteChoice(quote_id=selected['quote_id'])
    assert bind_choices([choice], sources, owner='formula')[0]['start'] == selected['start']
    for quote in [selected, next(q for q in catalog if q['quote'].startswith('```python') and q['quote'].endswith('```'))]:
        element = PageDraftElement(key='artifact', kind='formula' if quote == selected else 'code',
            use_source_text=True, sources=[QuoteChoice(quote_id=quote['quote_id'])])
        assert draft_element_text(element, sources) == quote['quote']
    for conflict in [{'block_id': 'wrong'}, {'quote': '$x=2$'}]:
        with pytest.raises(ValueError, match='source_quote_choice_conflict'):
            bind_choices([choice.model_copy(update=conflict)], sources, owner='formula')
    with pytest.raises(ValueError, match='source_quote_id_unknown'):
        bind_choices([QuoteChoice(quote_id='not-supplied')], sources, owner='formula')


def test_comparison_repair_reports_actual_invalid_endpoints():
    from .test_ppt_teaching_content import comparison_fixture
    value, _ = comparison_fixture()
    value['expression']['relations'] = [{'relation_id': 'bad-edge', 'source_id': 'condition',
        'target_id': 'serial', 'kind': 'causal', 'sources': value['elements'][0]['sources']}]
    with pytest.raises(ValueError, match='comparison_relation_endpoint_invalid:bad-edge: condition->serial'):
        PageTeachingV2.model_validate(value)
    value['expression']['relations'][0].update(source_id='a-mode', target_id='b-mode')
    with pytest.raises(ValueError, match='comparison_relation_crosses_subjects:bad-edge'):
        PageTeachingV2.model_validate(value)


def test_page_group_is_compiled_before_confirmation_and_can_resume_without_models(monkeypatch):
    import asyncio
    from ppt_teaching_planner import plan_teaching_manuscript, page_response_contract, regenerate_teaching_pages
    from .test_ppt_teaching_content import compiled_manuscript
    doc, graph, template, existing = compiled_manuscript()
    page = existing.pages[0]
    group = {'pages': [{'title': title, 'page_goal': page.page_goal, 'primary_claim': page.primary_claim,
        'teaching': page.teaching.model_dump(mode='json')} for title in ['比较执行方式', '观察执行差异']]}
    calls = []
    async def planner(request):
        calls.append(request)
        if request['teaching_request'] == 'narrative':
            return {'narrative_brief': {'central_question': '比较执行方式'}, 'pages': [{
                'page_id': 'p1', 'teaching_unit_id': page.teaching_unit_id, 'source_block_ids': ['b'],
                'title': page.title, 'page_goal': page.page_goal, 'layout_id': page.layout_id}]}
        return group
    result, trace = asyncio.run(plan_teaching_manuscript(doc, graph, template, planner))
    assert len(calls) == 2 and [p.page_id for p in result.pages] == ['p1-part-1', 'p1-part-2']
    assert result.page_count == result.story_page_count == 2
    async def forbidden(_):
        pytest.fail('accepted page group was regenerated')
    restored, _ = asyncio.run(plan_teaching_manuscript(doc, graph, template, forbidden, checkpoint=trace))
    assert restored.model_dump() == result.model_dump()
    schema = page_response_contract()['$defs']['LinearTeachingPageDraft']
    assert 'title' in schema['properties'] and 'title' in schema['required']
    assert schema['properties']['relations']['maxItems'] == 0
    monkeypatch.setattr('ppt_teaching_manuscript.template_for_manuscript', lambda _: template)
    before = existing.model_dump()
    split = asyncio.run(regenerate_teaching_pages(existing, ['p1'], planner))
    assert [p.page_id for p in split.pages] == ['p1-part-1', 'p1-part-2']
    assert existing.model_dump() == before
