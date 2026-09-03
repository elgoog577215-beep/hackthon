from __future__ import annotations

import asyncio
import json
import re
import threading
import time
from copy import deepcopy

import pytest
from fastapi import FastAPI
from fastapi.testclient import TestClient

import teacher_script as teacher_script_module
from ai_base import AIProviderUnavailable

from course_document import document_from_generation_draft
from teacher_lesson_authoring import (
    TeacherLessonAuthoringError,
    TeacherLessonAuthoringRepository,
    TeacherLessonAuthoringService,
    align_teacher_lesson_plan_to_arrangement,
    build_uploaded_ppt_review_report,
    complete_teacher_lesson_plan_fields,
    extract_uploaded_pptx_evidence,
    extract_uploaded_pptx_review,
    lesson_scope,
    normalize_teacher_lesson_plan,
    project_confirmed_teacher_scripts,
    teacher_lesson_script_revision,
    teacher_lesson_section_content,
    teacher_lesson_v6_source,
    validate_teacher_lesson_plan,
)
from teacher_script import (
    SCRIPT_PIPELINE_VERSION,
    SCRIPT_QUALITY_VERSION,
    compile_teacher_script_module_contract,
    compile_teacher_script_section,
    normalize_teacher_script_section,
    teacher_script_revision_is_publishable,
    validate_teacher_script_section,
    validate_teacher_script_revision,
)
from course_presentation_graph import compile_course_presentation_graph
from course_generation.service import CourseService
from dependencies import get_teacher_lesson_authoring_repository, require_task_manager
from teaching_design.lesson_arrangement import (
    _lesson_type,
    apply_lesson_arrangement_to_plan,
    normalize_lesson_arrangement,
    recommend_lesson_arrangement,
    validate_lesson_arrangement,
)
from routers import teacher_lesson_authoring as teacher_lesson_router
from routers import courses as courses_router


@pytest.mark.parametrize(
    ("asset_status", "parse_status", "expected_code"),
    [
        ("parsing", "", "lesson_material_source_processing"),
        ("failed", "failed", "lesson_material_source_parse_failed"),
    ],
)
def test_selected_material_must_finish_parsing_before_generation(
    monkeypatch,
    asset_status,
    parse_status,
    expected_code,
):
    class FakeCourseSpace:
        @staticmethod
        def list_owned(actor, course_id):
            assert (actor, course_id) == ("teacher-1", "course-1")
            return [{"package_id": "package-1"}]

        @staticmethod
        def load_owned(package_id, actor):
            assert (package_id, actor) == ("package-1", "teacher-1")
            return {"assets": [{"material_asset_id": "material-1"}]}

    class FakeMaterials:
        @staticmethod
        def get_asset(asset_id):
            return type("Asset", (), {"status": asset_status})()

        @staticmethod
        def load_parsed_document(asset_id):
            return type("Parsed", (), {"parse_status": parse_status})() if parse_status else None

        @staticmethod
        def load_evidence(asset_id):
            return []

    monkeypatch.setattr(
        teacher_lesson_router,
        "teacher_course_space_repository",
        FakeCourseSpace(),
    )
    monkeypatch.setattr(teacher_lesson_router, "material_repository", FakeMaterials())

    with pytest.raises(TeacherLessonAuthoringError) as raised:
        teacher_lesson_router._course_material_evidence(
            "course-1", "teacher-1", ["material-1"]
        )

    assert raised.value.code == expected_code
    assert raised.value.details == {"material_asset_ids": ["material-1"]}


def course_data():
    return {
        "course_id": "course-1",
        "nodes": [
            {"node_id": "L1-1", "parent_node_id": "root", "node_level": 1, "node_name": "第一讲"},
            {"node_id": "L2-1-1", "parent_node_id": "L1-1", "node_level": 2, "node_name": "1.1"},
            {"node_id": "L2-1-2", "parent_node_id": "L1-1", "node_level": 2, "node_name": "1.2"},
            {"node_id": "L1-2", "parent_node_id": "root", "node_level": 1, "node_name": "第二讲"},
            {"node_id": "L2-2-1", "parent_node_id": "L1-2", "node_level": 2, "node_name": "2.1"},
        ],
        "course_plan": {
            "chapters": [
                {"node_id": "L1-1", "title": "第一讲", "sections": [{"node_id": "L2-1-1"}, {"node_id": "L2-1-2"}]},
                {"node_id": "L1-2", "title": "第二讲", "sections": [{"node_id": "L2-2-1"}]},
            ]
        },
    }


def lecture_course_data():
    source = course_data()
    source["authoring_structure_version"] = "lecture_v1"
    source["course_plan"]["authoring_structure_version"] = "lecture_v1"
    for lecture_number, chapter in enumerate(
        source["course_plan"]["chapters"],
        start=1,
    ):
        chapter.pop("node_id")
        chapter["chapter_number"] = lecture_number
        chapter["lecture_number"] = lecture_number
    return source


def standard_lesson_plan():
    return {
        "schema_version": "course_teaching_plan_v3",
        "source_outline_revision_id": "outline-v1",
        "sections": [{
            "node_id": "L2-1-1",
            "learning_objective": "能独立解释核心概念并完成一次迁移应用。",
            "key_points": ["核心概念"],
            "key_difficulties": ["概念的适用边界"],
            "in_class_checks": ["完成一道情境判断并说明依据。"],
            "homework": ["用新场景复述概念并给出反例。"],
            "teaching_notes": ["保留学生产出用于课后复盘。"],
            "knowledge_structure": [{
                "knowledge_points": [{
                    "name": "核心概念",
                    "statement": "核心概念由定义、条件和边界共同构成。",
                }],
            }],
            "teaching_modules": [{
                "module_id": "core_explanation",
                "teaching_purpose": "建立概念框架",
                "knowledge_names": ["核心概念"],
                "planned_minutes": 20,
                "teacher_activity": "用正例与反例对照讲解定义和边界。",
                "student_activity": "归纳判断标准并完成情境判断。",
            }],
        }],
    }


def single_section_course_data(
    lesson_unit_id: str = "L1-1",
    section_node_id: str = "L2-1-1",
):
    source = course_data()
    source["course_plan"]["reference_books"] = [
        "张三：《核心概念导论》，高等教育出版社，2025"
    ]
    chapter = next(
        item
        for item in source["course_plan"]["chapters"]
        if item["node_id"] == lesson_unit_id
    )
    chapter["sections"] = [
        item for item in chapter["sections"]
        if item["node_id"] == section_node_id
    ]
    source["course_plan"]["chapters"] = [chapter]
    source["nodes"] = [
        item for item in source["nodes"]
        if item["node_id"] in {lesson_unit_id, section_node_id}
    ]
    return source


def test_lesson_scope_keeps_all_sections_inside_one_lesson():
    scoped = lesson_scope(course_data(), "L1-1")
    assert scoped["lesson"]["node_name"] == "第一讲"
    assert [item["node_id"] for item in scoped["sections"]] == ["L2-1-1", "L2-1-2"]
    assert scoped["chapter"]["node_id"] == "L1-1"


def test_generated_lesson_completes_formal_fields_and_specific_block_names():
    source = course_data()
    source["course_plan"]["reference_books"] = ["张三：《核心概念导论》，高等教育出版社，2025"]
    source["course_plan"]["chapters"][1]["title"] = "第二讲 迁移应用"
    plan = standard_lesson_plan()
    plan["sections"][0]["teaching_modules"][0]["module_id"] = "math_formalization"

    completed = complete_teacher_lesson_plan_fields(source, "L1-1", plan)
    section = completed["sections"][0]

    assert completed["formal_field_policy_version"] == "teacher_lesson_formal_fields_v1"
    assert section["class_summary"]
    assert section["homework_evaluation"]
    assert "第二讲 迁移应用" in section["next_lesson_connection"]
    assert section["resource_refs"] == [
        "已确认来源｜课程参考资料｜张三：《核心概念导论》，高等教育出版社，2025"
    ]
    assert section["teaching_modules"][0]["label"] == "形式化推导"
    assert "homework_submission" in section["teacher_confirmation_fields"]
    assert validate_teacher_lesson_plan(completed)["passed"] is True


def test_lecture_v1_lesson_completion_keeps_current_and_next_lecture_context():
    source = lecture_course_data()
    source["course_plan"]["chapters"][0]["extension_resources"] = [
        {
            "title": "界面设计基础",
            "resource_type": "book",
            "verification_status": "verified",
        }
    ]

    completed = complete_teacher_lesson_plan_fields(
        source,
        "L1-1",
        standard_lesson_plan(),
    )
    section = completed["sections"][0]

    assert "第二讲" in section["next_lesson_connection"]
    assert section["resource_refs"] == [
        "已确认来源｜教材/专著｜界面设计基础"
    ]


def test_material_absorption_creates_unconfirmed_linked_working_drafts(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    repository.set_outline("course-1", "outline-confirmed")
    confirmed = repository.save_plan_revision(
        "course-1",
        "L1-1",
        standard_lesson_plan(),
        source_outline_revision_id="outline-confirmed",
    )
    repository.confirm_plan_revision("course-1", "L1-1", confirmed["working_revision_id"])
    bundle = {
        "bundle_id": "bundle-1",
        "plan_id": "audit-1",
        "package_id": "package-1",
        "course_id": "course-1",
        "targets": [
            {
                "target_id": "managed:outline",
                "target_type": "outline",
                "target_scope_id": "course",
                "target_scope_label": "整门课程",
                "title": "课程大纲",
                "sources": [{"asset_id": "outline-source", "role": "primary"}],
                "structured_document": {
                    "content_hash": "outline-hash",
                    "sections": [{"section_id": "outline-s1", "title": "课程目标", "blocks": []}],
                },
            },
            {
                "target_id": "lesson-plan:L1-1",
                "target_type": "lesson_plan",
                "target_scope_id": "L1-1",
                "target_scope_label": "第一讲",
                "title": "第一讲教案",
                "sources": [{"asset_id": "plan-source", "role": "primary"}],
                "structured_document": {
                    "content_hash": "plan-hash",
                    "sections": [{"section_id": "plan-s1", "title": "教学流程", "blocks": []}],
                },
            },
        ],
    }

    receipt = repository.apply_material_absorption("course-1", bundle)
    duplicate = repository.apply_material_absorption("course-1", bundle)
    view = repository.view("course-1")

    assert receipt == duplicate
    assert receipt["status"] == "working_drafts_created"
    assert len(receipt["drafts"]) == 2
    assert view["outline_revision_id"] == "outline-confirmed"
    assert view["lessons"]["L1-1"]["confirmed_revision_id"] == confirmed["working_revision_id"]
    drafts = repository.current_material_drafts("course-1")
    assert drafts["outline"]["confirmation_required"] is True
    assert drafts["lessons"]["L1-1"]["lesson_plan"]["status"] == "working_draft"


def test_lesson_projection_recommends_current_arrangement_after_outline_change(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    old_source = course_data()
    old_source["blueprint_revision_id"] = "outline-old"
    old_source["nodes"] = [
        item for item in old_source["nodes"] if item.get("node_id") != "L2-1-2"
    ]
    old_source["course_plan"]["chapters"][0]["sections"] = [
        {"node_id": "L2-1-1"}
    ]
    old_arrangement = recommend_lesson_arrangement(
        old_source,
        "L1-1",
        source_outline_revision_id="outline-old",
    )
    repository.save_arrangement_revision(
        "course-1",
        "L1-1",
        old_arrangement,
        source_outline_revision_id="outline-old",
        confirm=True,
    )

    current_source = course_data()
    current_source["blueprint_revision_id"] = "outline-new"
    repository.set_outline("course-1", "outline-new")
    lesson = teacher_lesson_router._lesson_projection(
        current_source,
        repository,
    )[0]

    assert lesson["arrangement"]["source_state"] == "current"
    assert lesson["arrangement"]["confirmed"] is False
    assert {
        item["section_node_id"] for item in lesson["arrangement"]["blocks"]
    } == {"L2-1-1", "L2-1-2"}


def test_lesson_projection_keeps_legacy_fingerprint_out_of_formal_script_revision(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)

    lesson = teacher_lesson_router._lesson_projection(
        course_data(),
        repository,
    )[0]

    assert lesson["script"]["current_revision_id"] == ""
    assert lesson["script"]["legacy_source_fingerprint"].startswith("tlsr-")
    assert lesson["script"]["ready"] is False


def test_lesson_arrangement_projects_existing_modules_without_example_exam_collision():
    arrangement = recommend_lesson_arrangement(
        course_data(),
        "L1-2",
        source_outline_revision_id="outline-v1",
    )

    assert arrangement["lesson_type"] == "theory"
    assert arrangement["confirmed"] is False
    assert {item["section_node_id"] for item in arrangement["blocks"]} == {"L2-2-1"}
    assert validate_lesson_arrangement(
        arrangement,
        expected_section_ids=["L2-2-1"],
    ) == []

    applied = apply_lesson_arrangement_to_plan(
        course_data()["course_plan"],
        {**arrangement, "lesson_type": "case_discussion"},
    )
    section = applied["chapters"][1]["sections"][0]
    assert section["lesson_archetype"]["label"] == "案例研讨"
    assert [item["arrangement_block_id"] for item in section["module_plan"]] == [
        item["block_id"] for item in arrangement["blocks"]
    ]

    assert _lesson_type(
        ["engineering_debugging_lab", "engineering_guided_build"],
        ["core_explanation", "engineering_review", "engineering_testing"],
    ) == "theory_practice"


def test_changed_lesson_type_reorders_blocks_and_preserves_resources_and_tools():
    raw = {
        "lesson_type": "theory",
        "blocks": [
            {
                "block_id": "concept",
                "section_node_id": "L2-1-1",
                "module_id": "core_explanation",
                "name": "概念讲解",
                "role": "concept",
                "planned_minutes": 20,
                "resource_refs": ["教材第 3 章", "教材第 3 章"],
                "tools": ["函数图像工具", "函数图像工具"],
            },
            {
                "block_id": "checkpoint",
                "section_node_id": "L2-1-1",
                "module_id": "feedback_check",
                "name": "达成检查",
                "role": "checkpoint",
                "planned_minutes": 10,
            },
        ],
    }

    theory = normalize_lesson_arrangement(
        raw,
        lesson_unit_id="L1-1",
        source_outline_revision_id="outline-v1",
    )
    assessment = normalize_lesson_arrangement(
        {**raw, "lesson_type": "review_assessment"},
        lesson_unit_id="L1-1",
        source_outline_revision_id="outline-v1",
    )

    assert [item["block_id"] for item in theory["blocks"]] == ["concept", "checkpoint"]
    assert [item["block_id"] for item in assessment["blocks"]] == ["checkpoint", "concept"]
    concept = next(item for item in assessment["blocks"] if item["block_id"] == "concept")
    assert concept["resource_refs"] == ["教材第 3 章"]
    assert concept["tools"] == ["函数图像工具"]

    applied = apply_lesson_arrangement_to_plan(course_data()["course_plan"], assessment)
    module = next(
        item
        for item in applied["chapters"][0]["sections"][0]["module_plan"]
        if item["arrangement_block_id"] == "concept"
    )
    assert module["resource_refs"] == ["教材第 3 章"]
    assert module["tools"] == ["函数图像工具"]


def test_lesson_arrangement_adapts_legacy_node_outline_into_same_recommendation():
    legacy = course_data()
    legacy.pop("course_plan")

    arrangement = recommend_lesson_arrangement(
        legacy,
        "L1-1",
        source_outline_revision_id="legacy-outline-v1",
    )

    assert arrangement["blocks"]
    assert {item["section_node_id"] for item in arrangement["blocks"]} == {
        "L2-1-1",
        "L2-1-2",
    }
    assert validate_lesson_arrangement(
        arrangement,
        expected_section_ids=["L2-1-1", "L2-1-2"],
    ) == []


def test_lesson_arrangement_resolves_lecture_v1_outline_by_lecture_number():
    source = lecture_course_data()

    arrangement = recommend_lesson_arrangement(
        source,
        "L1-1",
        source_outline_revision_id="lecture-outline-v1",
    )

    assert arrangement["blocks"]
    assert {item["section_node_id"] for item in arrangement["blocks"]} == {
        "L2-1-1",
        "L2-1-2",
    }
    assert validate_lesson_arrangement(
        arrangement,
        expected_section_ids=["L2-1-1", "L2-1-2"],
    ) == []

    applied = apply_lesson_arrangement_to_plan(
        source["course_plan"],
        arrangement,
    )
    assert applied["chapters"][0]["sections"][0]["module_plan"]

    second_arrangement = recommend_lesson_arrangement(source, "L1-2")
    assert second_arrangement["lesson_phase"] == "closing"


def test_lesson_arrangement_compacts_large_legacy_chapter_to_one_row_per_section():
    legacy = course_data()
    legacy.pop("course_plan")
    legacy["nodes"] = [legacy["nodes"][0], *[
        {
            "node_id": f"legacy-section-{index}",
            "parent_node_id": "L1-1",
            "node_level": 2,
            "node_name": f"1.{index} 历史小节{index}",
        }
        for index in range(1, 7)
    ]]

    arrangement = recommend_lesson_arrangement(legacy, "L1-1")

    assert len(arrangement["blocks"]) == 6
    assert {item["section_node_id"] for item in arrangement["blocks"]} == {
        f"legacy-section-{index}" for index in range(1, 7)
    }
    assert all(item["module_id"] == "teacher_section_sequence" for item in arrangement["blocks"])
    assert all("→" in item["content_summary"] for item in arrangement["blocks"])


def test_arrangement_revision_marks_dependent_preparation_assets_stale(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    source = course_data()
    source["blueprint_revision_id"] = "outline-v1"
    repository.set_outline("course-1", "outline-v1")
    arrangement = recommend_lesson_arrangement(
        source,
        "L1-2",
        source_outline_revision_id="outline-v1",
    )
    saved_arrangement = repository.save_arrangement_revision(
        "course-1",
        "L1-2",
        arrangement,
        source_outline_revision_id="outline-v1",
        confirm=True,
    )
    first_arrangement_revision_id = saved_arrangement["arrangement"]["confirmed_revision_id"]

    plan = standard_lesson_plan()
    plan["sections"][0]["node_id"] = "L2-2-1"
    plan_asset = repository.save_plan_revision(
        "course-1",
        "L1-2",
        plan,
        source_outline_revision_id="outline-v1",
        quality_report={
            "passed": True,
            "blocking_issues": [],
            "review_issues": [],
            "metrics": {},
        },
    )
    plan_revision_id = plan_asset["working_revision_id"]
    repository.confirm_plan_revision(
        "course-1",
        "L1-2",
        plan_revision_id,
        quality_report={
            "passed": True,
            "blocking_issues": [],
            "review_issues": [],
            "metrics": {},
        },
    )
    state = repository.load("course-1")
    prepared = state["lessons"]["L1-2"]
    prepared["script_confirmation"] = {
        "confirmed_revision_id": "script-v1",
        "source_lesson_plan_revision_id": plan_revision_id,
        "source_state": "current",
    }
    prepared["ppt_assets"] = [{
        "asset_id": "ppt-v1",
        "source_lesson_plan_revision_id": plan_revision_id,
        "source_state": "current",
    }]
    repository._save(state)

    updated = deepcopy(repository.confirmed_arrangement("course-1", "L1-2"))
    assert updated is not None
    updated["blocks"][0]["feedback_strategy"] = "按教师的新备课要求调整反馈预案。"
    repository.save_arrangement_revision(
        "course-1",
        "L1-2",
        updated,
        source_outline_revision_id="outline-v1",
        confirm=True,
    )

    lesson = repository.lesson("course-1", "L1-2")
    plan_revision = next(
        item for item in lesson["revisions"]
        if item["revision_id"] == plan_revision_id
    )
    assert plan_revision["source_arrangement_revision_id"] == first_arrangement_revision_id
    assert lesson["source_state"] == "stale"
    assert lesson["source_state_reason"] == "arrangement_changed"
    assert lesson["script_confirmation"]["source_state"] == "stale"
    assert lesson["ppt_assets"][0]["source_state"] == "stale"

    repository.set_outline("course-1", "outline-v1")
    assert repository.lesson("course-1", "L1-2")["source_state"] == "stale"
    with pytest.raises(TeacherLessonAuthoringError) as exc_info:
        repository.confirm_plan_revision(
            "course-1",
            "L1-2",
            plan_revision_id,
            quality_report={"passed": True, "blocking_issues": [], "review_issues": [], "metrics": {}},
        )
    assert exc_info.value.code == "lesson_plan_arrangement_conflict"


def test_standard_lesson_plan_quality_gate_is_shared_by_draft_and_confirmation():
    report = validate_teacher_lesson_plan(
        standard_lesson_plan(),
        expected_section_ids=["L2-1-1"],
        expected_outline_revision_id="outline-v1",
        source_outline_revision_id="outline-v1",
    )
    assert report["passed"] is True
    assert report["pipeline_version"] == "standard_lesson_plan_v1"
    assert report["metrics"]["planned_minutes"] == 20

    incomplete = standard_lesson_plan()
    incomplete["sections"][0]["in_class_checks"] = []
    incomplete["sections"][0]["teaching_modules"] = []
    blocked = validate_teacher_lesson_plan(
        incomplete,
        expected_section_ids=["L2-1-1"],
    )
    assert blocked["passed"] is False
    assert {item["code"] for item in blocked["blocking_issues"]} >= {
        "lesson_plan:modules",
        "lesson_plan:checks",
    }


def test_standard_lesson_plan_blocks_internal_policy_and_abstract_activity_language():
    plan = standard_lesson_plan()
    section = plan["sections"][0]
    section["teaching_notes"] = [
        "无已确认教师资料或外部来源，不声称来自真实数据。"
    ]
    section["teaching_modules"][0]["teacher_activity"] = "建立问题、价值与任务边界"

    report = validate_teacher_lesson_plan(plan)

    assert report["passed"] is False
    assert {item["code"] for item in report["blocking_issues"]} >= {
        "lesson_plan:internal_register",
        "lesson_plan:abstract_activity",
    }
    assert report["metrics"]["teacher_language_rule_version"] == "teacher_plan_language_v2"


def test_lesson_plan_clock_and_module_order_follow_confirmed_arrangement():
    arrangement = {
        "blocks": [
            {
                "block_id": "a-1",
                "section_node_id": "L2-1-1",
                "module_id": "lesson_goal",
                "name": "目标",
                "planned_minutes": 10,
            },
            {
                "block_id": "a-2",
                "section_node_id": "L2-1-1",
                "module_id": "core_explanation",
                "name": "讲解",
                "planned_minutes": 11,
            },
            {
                "block_id": "a-3",
                "section_node_id": "L2-1-2",
                "module_id": "worked_example",
                "name": "例题",
                "planned_minutes": 11,
            },
            {
                "block_id": "a-4",
                "section_node_id": "L2-1-2",
                "module_id": "feedback_check",
                "name": "检查",
                "planned_minutes": 13,
            },
        ],
    }
    drifting = {
        "schema_version": "course_teaching_plan_v3",
        "sections": [
            {
                "node_id": "L2-1-1",
                "learning_objective": "理解第一部分。",
                "key_points": ["第一部分"],
                "teaching_modules": [{
                    "module_id": "core_explanation",
                    "planned_minutes": 65,
                    "teacher_activity": "讲解第一部分。",
                    "student_activity": "完成判断。",
                }],
            },
            {
                "node_id": "L2-1-2",
                "learning_objective": "理解第二部分。",
                "key_points": ["第二部分"],
                "teaching_modules": [],
            },
        ],
    }

    aligned = align_teacher_lesson_plan_to_arrangement(drifting, arrangement)

    assert aligned["lesson_duration_minutes"] == 45
    assert [
        module["module_id"]
        for section in aligned["sections"]
        for module in section["teaching_modules"]
    ] == ["lesson_goal", "core_explanation", "worked_example", "feedback_check"]
    assert [
        module["planned_minutes"]
        for section in aligned["sections"]
        for module in section["teaching_modules"]
    ] == [10, 11, 11, 13]
    assert validate_teacher_lesson_plan(
        aligned,
        expected_section_ids=["L2-1-1", "L2-1-2"],
        expected_total_minutes=45,
    )["passed"] is True


def test_teacher_script_inherits_confirmed_archetype_and_module_order():
    outline = {
        "node_id": "L2-1-1",
        "node_name": "1.1 核心概念",
        "lesson_archetype": {
            "archetype_id": "general_concept_building",
            "label": "概念建构",
        },
        "module_plan": [
            {"module_id": "lesson_goal", "label": "本节任务", "required": True},
            {"module_id": "general_concept_model", "label": "概念模型", "required": True},
            {"module_id": "feedback_check", "label": "检查与反馈", "required": True},
        ],
    }
    plan = {
        "node_id": "L2-1-1",
        "learning_objective": "能解释概念并划清边界。",
        "teaching_modules": [
            {"module_id": "lesson_goal", "planned_minutes": 3},
            {"module_id": "general_concept_model", "planned_minutes": 20},
            {"module_id": "feedback_check", "planned_minutes": 7},
        ],
    }

    contract = compile_teacher_script_module_contract(outline, plan)
    assert contract["lesson_archetype"]["label"] == "概念建构"
    assert [item["module_id"] for item in contract["modules"]] == [
        "lesson_goal", "general_concept_model", "feedback_check",
    ]
    assert [item["title"] for item in contract["modules"]] == [
        "本节任务", "概念模型", "检查与反馈",
    ]
    assert 650 <= contract["modules"][0]["max_characters"] <= 1000
    assert 1800 <= contract["modules"][1]["max_characters"] <= 2600

    compiled = compile_teacher_script_section(
        "## 本节任务\n\n我们先明确本节要解决的核心问题：怎样形成稳定、可检查的概念判断标准。\n\n"
        "## 概念模型\n\n带着这个问题来看，概念模型由定义、成立条件与适用边界三个部分构成，正反例共同验证这一结构。\n\n"
        "## 检查与反馈\n\n最后请大家用一个新情境逐项对照定义、条件与边界。典型错误是漏掉成立条件；核对标准是三项齐全，发现错误后要说明修正原因并再次判断。",
        contract,
    )
    assert compiled["quality_report"]["passed"] is True
    assert [item["module_id"] for item in compiled["blocks"]] == [
        "lesson_goal", "general_concept_model", "feedback_check",
    ]
    source = course_data()
    source["nodes"][1].update(outline)
    _document, view, _synthetic_id = teacher_lesson_v6_source(
        source,
        lesson_unit_id="L1-1",
        plan_revision={
            "revision_id": "plan-1",
            "plan": {
                "schema_version": "course_teaching_plan_v3",
                "sections": [plan],
            },
        },
        script_revision={
            "revision_id": "script-1",
            "sections": [
                compiled,
                {
                    "section_node_id": "L2-1-2",
                    "blocks": [{
                        "block_id": "second-section-block",
                        "module_id": "legacy_script",
                        "role": "concept",
                        "title": "第二小节讲稿",
                        "content": "这是第二小节已经确认的讲稿内容。",
                    }],
                },
            ],
        },
    )
    projected_blocks = view["nodes"][1]["content_blocks"]
    assert [item["metadata"]["module_id"] for item in projected_blocks] == [
        "lesson_goal", "general_concept_model", "feedback_check",
    ]
    assert all(
        item["metadata"]["source_kind"] == "current_teacher_script_block"
        for item in projected_blocks
    )
    assert all(
        item["metadata"]["content_perspective"] == "teacher_delivery"
        for item in projected_blocks
    )
    assert all(
        "teacher_activity" not in item["metadata"]
        and "student_activity" not in item["metadata"]
        for item in projected_blocks
    )
    assert all(item["title"] != "讲稿正文" for item in projected_blocks)

    generic = compile_teacher_script_section(
        "## 背景导入\n\n这是一段模型自行增加的通用模板内容。",
        contract,
    )
    report = validate_teacher_script_section(generic, contract)
    assert report["passed"] is False
    assert {item["code"] for item in report["blocking_issues"]} >= {
        "teacher_script:module_contract",
        "teacher_script:module_heading",
    }

    tampered = json.loads(json.dumps(compiled))
    tampered["blocks"][0]["block_id"] = "replacement-block"
    tampered["blocks"][0]["role"] = "example"
    tampered["blocks"][0]["knowledge_names"] = ["教案范围外知识"]
    tampered_report = validate_teacher_script_section(tampered, contract)
    assert {item["code"] for item in tampered_report["blocking_issues"]} >= {
        "teacher_script:block_contract",
        "teacher_script:role_contract",
        "teacher_script:knowledge_scope",
    }


def test_confirmed_script_blocks_compile_into_teaching_page_groups():
    section_ids = [f"L2-1-{index}" for index in range(1, 5)]
    source = {
        "course_id": "linear-algebra-chapter-one",
        "nodes": [
            {
                "node_id": "L1-1",
                "parent_node_id": "root",
                "node_level": 1,
                "node_name": "第1章 行列式",
            },
            *[
                {
                    "node_id": section_id,
                    "parent_node_id": "L1-1",
                    "node_level": 2,
                    "node_name": f"1.{index} 小节",
                }
                for index, section_id in enumerate(section_ids, start=1)
            ],
        ],
    }
    block_counts = [8, 7, 7, 7]
    script_sections = []
    for section_id, block_count in zip(section_ids, block_counts):
        script_sections.append({
            "section_node_id": section_id,
            "blocks": [
                {
                    "block_id": f"{section_id}-block-{index}",
                    "module_id": "core_explanation",
                    "role": "concept",
                    "title": f"{section_id} 教学点 {index}",
                    "content": f"这是第 {index} 个完整教学点的定义、条件与边界。",
                    "planned_minutes": 2 if index <= 4 else 1,
                }
                for index in range(1, block_count + 1)
            ],
        })
    document, _view, _synthetic_id = teacher_lesson_v6_source(
        source,
        lesson_unit_id="L1-1",
        plan_revision={
            "revision_id": "plan-linear-algebra",
            "plan": {
                "schema_version": "course_teaching_plan_v3",
                "sections": [
                    {"node_id": section_id, "teaching_modules": []}
                    for section_id in section_ids
                ],
            },
        },
        script_revision={
            "revision_id": "script-linear-algebra",
            "sections": script_sections,
        },
    )
    graph = compile_course_presentation_graph(document, teaching_plan={})

    assert len(document.blocks) == 29
    assert 10 <= len(graph.units) <= 14
    assert sum(len(unit.primary_block_ids) for unit in graph.units) == 29
    assert all(len(unit.primary_block_ids) <= 3 for unit in graph.units)


@pytest.mark.parametrize(
    ("module_id", "content", "missing_code"),
    [
        (
            "engineering_minimal_run",
            "```python\nprint('hello')\n```\n\n运行结果为 hello。环境版本和输出文本共同构成最小验收条件。",
            "teacher_script:required_code_artifact",
        ),
        (
            "math_formalization",
            "定义 $f(x)=x^2$。其中 $x$ 是自变量，定义域决定表达式的适用边界。",
            "teacher_script:required_math_artifact",
        ),
    ],
)
def test_teacher_script_enforces_discipline_artifacts(module_id, content, missing_code):
    outline = {
        "node_id": "L2-1-1",
        "node_name": "学科小节",
        "module_plan": [{"module_id": module_id}],
    }
    plan = {
        "node_id": "L2-1-1",
        "teaching_modules": [{"module_id": module_id}],
    }
    contract = compile_teacher_script_module_contract(outline, plan)
    title = contract["modules"][0]["title"]
    complete = compile_teacher_script_section(f"## {title}\n\n{content}", contract)
    assert complete["quality_report"]["passed"] is True

    missing = compile_teacher_script_section(
        f"## {title}\n\n这里只给出概括性说明，没有提供本模块要求的正式学科产物。",
        contract,
    )
    assert missing["quality_report"]["passed"] is False
    assert missing_code in {
        item["code"] for item in missing["quality_report"]["blocking_issues"]
    }


def test_teacher_script_repairs_unambiguous_math_but_still_rejects_code_fence():
    outline = {
        "node_id": "L2-1-1",
        "node_name": "完整性检查",
        "module_plan": [{"module_id": "core_explanation"}],
    }
    plan = {
        "node_id": "L2-1-1",
        "teaching_modules": [{"module_id": "core_explanation"}],
    }
    contract = compile_teacher_script_module_contract(outline, plan)
    title = contract["modules"][0]["title"]
    compiled = compile_teacher_script_section(
        f"## {title}\n\n```text\n未闭合代码\n\n同时出现未闭合公式 $$x+y。",
        contract,
    )
    codes = {
        item["code"] for item in compiled["quality_report"]["blocking_issues"]
    }
    assert "teacher_script:unclosed_code_fence" in codes
    assert "teacher_script:unclosed_math_delimiter" not in codes

    inline = compile_teacher_script_section(
        f"## {title}\n\n公式在返回前被截断：$F_x=6-3",
        contract,
    )
    inline_codes = {
        item["code"]
        for item in inline["quality_report"]["blocking_issues"]
    }
    assert "teacher_script:unclosed_math_delimiter" not in inline_codes
    assert inline["blocks"][0]["content"].endswith("$")
    assert inline["format_repairs"]

    cross_line = compile_teacher_script_section(
        f"## {title}\n\n$u > 0\n\n随后得到结论。\n\nx^2 < 1$",
        contract,
    )
    assert "$u > 0$" in cross_line["blocks"][0]["content"]
    assert "$x^2 < 1$" in cross_line["blocks"][0]["content"]
    assert "teacher_script:unclosed_math_delimiter" not in {
        item["code"]
        for item in cross_line["quality_report"]["blocking_issues"]
    }


def test_teacher_script_normalizes_split_matrix_shell_before_quality_gate():
    outline = {
        "node_id": "L2-1-1",
        "node_name": "矩阵表示",
        "module_plan": [{"module_id": "core_explanation"}],
    }
    plan = {
        "node_id": "L2-1-1",
        "teaching_modules": [{"module_id": "core_explanation"}],
    }
    contract = compile_teacher_script_module_contract(outline, plan)
    title = contract["modules"][0]["title"]
    compiled = compile_teacher_script_section(
        (
            f"## {title}\n\n增广矩阵写为：\n"
            "$$\n\\left[\n$$\n"
            "\\begin{array}{cc|c}\n1 & 0 & 2 \\\\n0 & 1 & 3\n\\end{array}\n"
            "$$\n\\right]\n$$\n"
            "矩阵的每一行都与原方程保持等价。"
        ),
        contract,
    )

    content = compiled["blocks"][0]["content"]
    codes = {
        item["code"]
        for item in compiled["quality_report"]["blocking_issues"]
    }
    assert content.count("$$") == 2
    assert "\\left[" in content and "\\begin{array}" in content
    assert "\\right]" in content
    assert "teacher_script:unwrapped_display_math_environment" not in codes
    assert compiled["format_repairs"] == [{
        "block_id": compiled["blocks"][0]["block_id"],
        "repairs": ["normalize:display-math-shape"],
    }]


def test_teacher_script_quality_rejects_unwrapped_matrix_environment_checkpoint():
    outline = {
        "node_id": "L2-1-1",
        "node_name": "矩阵表示",
        "module_plan": [{"module_id": "core_explanation"}],
    }
    plan = {
        "node_id": "L2-1-1",
        "teaching_modules": [{"module_id": "core_explanation"}],
    }
    contract = compile_teacher_script_module_contract(outline, plan)
    section = normalize_teacher_script_section({
        "section_node_id": contract["section_node_id"],
        "title": contract["title"],
        "blocks": [{
            **contract["modules"][0],
            "content": (
                "矩阵表示被拆开。$$\\left[$$\n"
                "\\begin{array}{cc}1 & 0 \\\\ 0 & 1\\end{array}\n"
                "$$\\right]$$，因此必须重新生成。"
            ),
        }],
    }, contract)

    codes = {
        item["code"]
        for item in validate_teacher_script_section(section, contract)[
            "blocking_issues"
        ]
    }
    assert "teacher_script:unwrapped_display_math_environment" in codes


def test_teacher_script_quality_rejects_display_formula_that_swallows_prose():
    outline = {
        "node_id": "L2-1-1",
        "node_name": "矩阵判断",
        "module_plan": [{"module_id": "core_explanation"}],
    }
    plan = {
        "node_id": "L2-1-1",
        "teaching_modules": [{"module_id": "core_explanation"}],
    }
    contract = compile_teacher_script_module_contract(outline, plan)
    section = normalize_teacher_script_section({
        "section_node_id": contract["section_node_id"],
        "title": contract["title"],
        "blocks": [{
            **contract["modules"][0],
            "content": (
                "任务示例如下。\n$$\n"
                "\\begin{bmatrix}1 & 0 \\\\ 0 & 1\\end{bmatrix}\n"
                "输出要求：圈出主元并写出判断依据。参考解法：逐行寻找首个非零元。"
                "核对标准：主元位置严格右移；若出现矛盾行则判定无解。"
                "这些解释必须位于公式之外。\n$$"
            ),
        }],
    }, contract)

    codes = {
        item["code"]
        for item in validate_teacher_script_section(section, contract)[
            "blocking_issues"
        ]
    }
    assert "teacher_script:prose_inside_display_math" in codes


def test_teacher_script_compile_moves_task_prose_outside_display_formula():
    outline = {
        "node_id": "L2-1-1",
        "node_name": "矩阵判断",
        "module_plan": [{"module_id": "learner_action"}],
    }
    plan = {
        "node_id": "L2-1-1",
        "teaching_modules": [{"module_id": "learner_action"}],
    }
    contract = compile_teacher_script_module_contract(outline, plan)
    title = contract["modules"][0]["title"]
    compiled = compile_teacher_script_section(
        (
            f"## {title}\n\n我们现在来判断矩阵的主元。任务条件如下。\n\n$$\n"
            "A=\\begin{bmatrix}1 & 0 \\\\ 0 & 1\\end{bmatrix}\n\n"
            "输出要求：圈出主元并说明依据。\n"
            "参考解法：逐行寻找首个非零元，再核对其下方是否全零。"
            "验收标准：结论、依据与边界完整。\n$$"
        ),
        contract,
    )

    content = compiled["blocks"][0]["content"]
    display_body = re.search(r"\$\$([\s\S]*?)\$\$", content).group(1)
    assert "输出要求" not in display_body
    assert "输出要求" in content and "参考解法" in content
    assert compiled["quality_report"]["passed"] is True
    assert compiled["format_repairs"] == [{
        "block_id": compiled["blocks"][0]["block_id"],
        "repairs": ["normalize:display-math-prose-boundary"],
    }]


def test_teacher_script_blocks_placeholder_repetition_and_shallow_full_lesson():
    sections = []
    for index in range(1, 5):
        outline = {
            "node_id": f"L2-1-{index}",
            "node_name": f"1.{index} 小节",
            "module_plan": [{
                "module_id": "core_explanation",
                "label": f"核心教学 {index}",
            }],
        }
        plan = {
            "node_id": f"L2-1-{index}",
            "teaching_modules": [{
                "module_id": "core_explanation",
                "planned_minutes": 10,
            }],
        }
        contract = compile_teacher_script_module_contract(outline, plan)
        content = (
            "本块内容完整。当前教学块围绕已确认的当前知识范围展开。"
            if index == 1
            else "概念需要同时说明定义、成立条件与适用边界，并通过正例和反例逐项核对判断标准。"
        )
        sections.append(compile_teacher_script_section(
            f"## 核心教学 {index}\n\n{content}",
            contract,
        ))

    report = validate_teacher_script_revision(
        sections,
        generation_source="model_block_pipeline",
    )

    assert report["passed"] is False
    assert report["publication_eligible"] is False
    assert {item["code"] for item in report["blocking_issues"]} >= {
        "teacher_script:placeholder_content",
        "teacher_script:repetitive_blocks",
        "teacher_script:lesson_too_shallow",
    }


def test_teacher_script_does_not_treat_distinct_matrices_as_repeated_prose():
    sections = []
    for index, matrix in enumerate((
        "1 & 0 & 2 \\\\ 0 & 1 & 3",
        "1 & 2 & 4 \\\\ 0 & 0 & 5",
        "2 & 1 & 0 \\\\ 0 & 3 & 6",
    ), start=1):
        sections.append({
            "schema_version": "teacher_script_v2",
            "pipeline_version": SCRIPT_PIPELINE_VERSION,
            "section_node_id": f"matrix-{index}",
            "quality_report": {
                "schema_version": SCRIPT_QUALITY_VERSION,
                "pipeline_version": SCRIPT_PIPELINE_VERSION,
                "passed": True,
                "blocking_issues": [],
                "review_issues": [],
            },
            "blocks": [{
                "block_id": f"matrix-block-{index}",
                "title": f"矩阵 {index}",
                "planned_minutes": 1,
                "content": f"$$\\begin{{bmatrix}} {matrix} \\end{{bmatrix}}$$",
            }],
        })

    report = validate_teacher_script_revision(
        sections,
        generation_source="model_block_pipeline",
    )

    assert "teacher_script:repetitive_blocks" not in {
        item["code"] for item in report["blocking_issues"]
    }


def test_teacher_script_stale_quality_contract_is_never_publishable():
    assert teacher_script_revision_is_publishable({
        "publication_eligible": True,
        "pipeline_version": SCRIPT_PIPELINE_VERSION,
        "quality_report": {
            "schema_version": "teacher_script_quality_v5",
            "pipeline_version": "neutral_course_script_v5",
            "passed": True,
            "publication_eligible": True,
        },
    }) is False
    assert SCRIPT_QUALITY_VERSION == "teacher_script_quality_v8"


def test_teacher_script_revision_blocks_repeated_canned_transitions():
    sections = [{
        "section_node_id": "L2-1-1",
        "quality_report": {
            "schema_version": SCRIPT_QUALITY_VERSION,
            "pipeline_version": SCRIPT_PIPELINE_VERSION,
            "passed": True,
            "blocking_issues": [],
            "review_issues": [],
        },
        "blocks": [
            {
                "block_id": f"block-{index}",
                "planned_minutes": 1,
                "content": f"值得注意的是，第 {index} 个问题需要结合不同条件单独判断，并写出对应依据。",
            }
            for index in range(1, 5)
        ],
    }]

    report = validate_teacher_script_revision(
        sections,
        generation_source="model_block_pipeline",
    )

    assert "teacher_script:repetitive_canned_transitions" in {
        item["code"] for item in report["blocking_issues"]
    }


def test_ppt_source_keeps_quality_report_non_blocking_for_current_script(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    plan = standard_lesson_plan()
    second_section = deepcopy(plan["sections"][0])
    second_section["node_id"] = "L2-1-2"
    plan["sections"].append(second_section)
    lesson = repository.save_plan_revision(
        "course-1",
        "L1-1",
        plan,
        source_outline_revision_id="outline-v1",
        quality_report=validate_teacher_lesson_plan(plan),
    )
    plan_revision = lesson["working_revision_id"]
    lesson = repository.save_script_revision(
        "course-1",
        "L1-1",
        [
            {"section_node_id": "L2-1-1", "title": "1.1", "content": "第一节完整讲稿包含定义、条件、边界和核对方法。"},
            {"section_node_id": "L2-1-2", "title": "1.2", "content": "第二节完整讲稿包含推理步骤、示例结果和检验标准。"},
        ],
        source_lesson_plan_revision_id=plan_revision,
        generation_source="teacher_edit",
    )
    script_revision = lesson["working_script_revision_id"]

    persisted = repository.load("course-1")
    revision = persisted["lessons"]["L1-1"]["script_revisions"][0]
    revision["quality_report"].update({
        "schema_version": "teacher_script_quality_v5",
        "pipeline_version": "neutral_course_script_v5",
        "passed": True,
        "publication_eligible": True,
    })
    revision["publication_eligible"] = True
    repository._save(persisted)

    class FakeStorage:
        @staticmethod
        def load_course(_course_id):
            source = course_data()
            source["blueprint_revision_id"] = "outline-v1"
            return source

    class FakeTaskManager:
        storage = FakeStorage()

        @staticmethod
        def get_generation_workspace_course(_course_id):
            return None

        @staticmethod
        def get_generation_preview(_course_id):
            return None

    document, _course, _synthetic_id, _lesson, _plan = (
        teacher_lesson_router._teacher_v6_source(
            FakeTaskManager(), repository, "course-1", "L1-1"
        )
    )
    assert document.document_revision


def test_teacher_script_rejects_mechanical_cues_plan_voice_and_truncation():
    outline = {
        "node_id": "L2-1-1",
        "node_name": "中性讲稿",
        "module_plan": [{"module_id": "core_explanation", "label": "核心教学"}],
    }
    plan = {
        "node_id": "L2-1-1",
        "teaching_modules": [{
            "module_id": "core_explanation",
            "teacher_activity": "请大家观察反例。",
            "student_activity": "完成判断。",
        }],
    }
    contract = compile_teacher_script_module_contract(outline, plan)
    assert contract["content_perspective"] == "teacher_delivery"
    assert contract["modules"][0]["source_plan_context"] == {
        "teacher_activity": "请大家观察反例。",
        "student_activity": "完成判断。",
    }

    compiled = compile_teacher_script_section(
        "## 核心教学\n\n【板书】教师应请大家观察。全链路验收完成，因为",
        contract,
    )
    codes = {
        item["code"]
        for item in compiled["quality_report"]["blocking_issues"]
    }
    assert codes >= {
        "teacher_script:classroom_delivery_cue",
        "teacher_script:lesson_plan_voice",
        "teacher_script:internal_process_leakage",
        "teacher_script:incomplete_block_ending",
    }


def test_teacher_script_service_generates_direct_teaching_script(monkeypatch):
    service = CourseService()
    captured = {}

    async def fake_call(user_prompt, system_prompt, **kwargs):
        captured["user_prompt"] = user_prompt
        captured["system_prompt"] = system_prompt
        captured["kwargs"] = kwargs
        return "## 核心教学\n\n我们先看核心概念怎样成立：它由定义、成立条件和适用边界共同构成。正例满足全部条件，反例则显示概念边界。"

    monkeypatch.setattr(service, "_call_llm", fake_call)
    result = asyncio.run(service.generate_teacher_script_section(
        course_id="course-1",
        outline_section={
            "node_id": "L2-1-1",
            "node_name": "1.1 核心概念",
            "lesson_archetype": {
                "archetype_id": "general_concept_building",
                "label": "概念建构",
                "purpose": "建立概念、关系和边界。",
            },
            "module_plan": [{
                "module_id": "core_explanation",
                "label": "核心教学",
                "required": True,
            }],
        },
        confirmed_plan_section={
            "node_id": "L2-1-1",
            "learning_objective": "能解释核心概念并划清边界。",
            "teaching_modules": [{
                "module_id": "core_explanation",
                "knowledge_names": ["核心概念"],
                "planned_minutes": 20,
                "teacher_activity": "用正反例讲解。",
            }],
        },
        lesson_context={"lesson_title": "第一讲"},
        requirements="贴近课堂表达",
    ))

    assert result["quality_report"]["passed"] is True
    assert "教师站在讲台上实际说的完整讲义" in captured["system_prompt"]
    assert "本节课型：概念建构" in captured["system_prompt"]
    assert "学科类型与当前教学块策略" in captured["system_prompt"]
    assert "前后小节连贯与课程总编约束" in captured["system_prompt"]
    assert "提问写出问题原话" in captured["system_prompt"]
    assert "## 核心教学" in captured["system_prompt"]
    assert "改写为教师当场会说的话" in captured["system_prompt"]
    assert "当前生成不得输出 `$$`" in captured["system_prompt"]
    assert "不得超过" in captured["system_prompt"]
    assert "自学课程的完整小节" not in captured["system_prompt"]
    assert captured["kwargs"]["use_fast_model"] is True
    assert captured["kwargs"]["enable_thinking"] is False
    assert captured["kwargs"]["max_attempts"] == 2
    assert captured["kwargs"]["reject_truncated"] is True


def test_teacher_script_service_hardens_formula_boundaries_on_retry(monkeypatch):
    service = CourseService()
    prompts = []

    async def fake_call(_user_prompt, system_prompt, **_kwargs):
        prompts.append(system_prompt)
        if len(prompts) == 1:
            return (
                "## 数学建模\n\n"
                "我们把流量写成模型。\n\n$$\nQ(t)=2t+1\n"
                "接下来请大家解释变量和单位。这个结论还要接受情境条件检查。"
                "请先判断自变量和因变量，再说明结果的实际意义。"
                "最后换一个时刻重新计算，并核对单位和适用边界。\n$$"
            )
        return (
            "## 数学建模\n\n"
            "我们把流量写成模型：\n\n\\[Q(t)=2t+1\\]\n\n"
            "接下来请大家解释变量和单位，再用新的时刻检验结论是否符合情境条件。"
        )

    monkeypatch.setattr(service, "_call_llm", fake_call)
    result = asyncio.run(service.generate_teacher_script_section(
        course_id="course-formula-boundary",
        outline_section={
            "node_id": "L2-1-1",
            "node_name": "1.1 变化率模型",
            "module_plan": [{"module_id": "math_modeling", "label": "数学建模"}],
        },
        confirmed_plan_section={
            "node_id": "L2-1-1",
            "teaching_modules": [{
                "module_id": "math_modeling",
                "planned_minutes": 6,
                "knowledge_names": ["变化率模型"],
            }],
        },
    ))

    assert len(prompts) == 2
    assert "这次禁止使用 `$$`" in prompts[1]
    assert result["quality_report"]["passed"] is True


def test_teacher_script_service_exposes_checkable_activity_structure(monkeypatch):
    service = CourseService()
    prompts = []

    async def fake_call(_user_prompt, system_prompt, **_kwargs):
        prompts.append(system_prompt)
        return (
            "## 学习者行动\n\n"
            "任务条件：给定一个包含边界情形的判断题，写出使用的定义与成立条件。"
            "参考解法：先核对对象与条件，再给出结果。验收标准：结论、依据和边界三项齐全。"
        )

    monkeypatch.setattr(service, "_call_llm", fake_call)
    result = asyncio.run(service.generate_teacher_script_section(
        course_id="course-activity-contract",
        outline_section={
            "node_id": "L2-1-1",
            "node_name": "1.1 核心概念",
            "module_plan": [{
                "module_id": "learner_action",
                "label": "学习者行动",
            }],
        },
        confirmed_plan_section={
            "node_id": "L2-1-1",
            "teaching_modules": [{
                "module_id": "learner_action",
                "planned_minutes": 4,
            }],
        },
    ))

    assert result["quality_report"]["passed"] is True
    assert "任务条件" in prompts[0]
    assert "参考解法" in prompts[0]
    assert "验收标准" in prompts[0]


def test_teacher_script_rejects_textbook_length_block():
    outline = {
        "node_id": "L2-1-1",
        "node_name": "轻量讲稿",
        "module_plan": [{"module_id": "core_explanation"}],
    }
    plan = {
        "node_id": "L2-1-1",
        "teaching_modules": [{
            "module_id": "core_explanation",
            "planned_minutes": 10,
        }],
    }
    contract = compile_teacher_script_module_contract(outline, plan)
    module = contract["modules"][0]
    oversized = "这是重复展开的教材式讲解。" * (module["max_characters"] // 10 + 20)
    compiled = compile_teacher_script_section(
        f"## {module['title']}\n\n{oversized}",
        contract,
    )
    assert compiled["quality_report"]["passed"] is False
    assert "teacher_script:block_too_long" in {
        item["code"]
        for item in compiled["quality_report"]["blocking_issues"]
    }


def test_teacher_script_service_compacts_length_only_failure(monkeypatch):
    service = CourseService()
    calls = []

    async def fake_call(user_prompt, _system_prompt, **_kwargs):
        calls.append(user_prompt)
        if len(calls) < 3:
            return "## 核心教学\n\n" + "重复讲解。" * 400
        return (
            "## 核心教学\n\n"
            "我们把重复内容收束成一个判断框架：概念的完整表达包含定义、成立条件与适用边界。"
            "接着看一个正例，用它验证三者是否被同时满足。"
        )

    monkeypatch.setattr(service, "_call_llm", fake_call)
    result = asyncio.run(service.generate_teacher_script_section(
        course_id="course-compact",
        outline_section={
            "node_id": "L2-1-1",
            "node_name": "轻量讲解",
            "module_plan": [{
                "module_id": "core_explanation",
                "label": "核心教学",
            }],
        },
        confirmed_plan_section={
            "node_id": "L2-1-1",
            "teaching_modules": [{
                "module_id": "core_explanation",
                "planned_minutes": 10,
            }],
        },
    ))

    assert len(calls) == 3
    assert "请压缩下面的教师讲义" in calls[-1]
    assert result["quality_report"]["passed"] is True


def test_teacher_script_service_uses_smart_pool_after_fast_pool_failure(monkeypatch):
    service = CourseService()
    routes = []
    stream_events = []

    async def fake_call(_user_prompt, _system_prompt, **kwargs):
        routes.append(kwargs["use_fast_model"])
        if kwargs["use_fast_model"]:
            raise AIProviderUnavailable("fast_pool_exhausted")
        await kwargs["on_content_delta"]("概念必须同时说明")
        return (
            "## 核心教学\n\n"
            "概念必须同时说明定义、成立条件和适用边界，新情境可用于核对这三项标准。"
        )

    async def on_reset():
        stream_events.append("reset")

    async def on_delta(delta):
        stream_events.append(delta)

    monkeypatch.setattr(service, "_call_llm", fake_call)
    result = asyncio.run(service.generate_teacher_script_section(
        course_id="course-smart-fallback",
        outline_section={
            "node_id": "L2-1-1",
            "node_name": "轻量讲解",
            "module_plan": [{
                "module_id": "core_explanation",
                "label": "核心教学",
            }],
        },
        confirmed_plan_section={
            "node_id": "L2-1-1",
            "teaching_modules": [{"module_id": "core_explanation"}],
        },
        on_content_reset=on_reset,
        on_content_delta=on_delta,
    ))

    assert routes == [True, False]
    assert stream_events == ["reset", "reset", "概念必须同时说明"]
    assert result["quality_report"]["passed"] is True


def test_teacher_script_requests_share_course_service_capacity(monkeypatch):
    service = CourseService()
    service._teaching_plan_semaphore = asyncio.Semaphore(1)
    active = 0
    peak = 0

    async def fake_call(_user_prompt, _system_prompt, **_kwargs):
        nonlocal active, peak
        active += 1
        peak = max(peak, active)
        await asyncio.sleep(0.01)
        active -= 1
        return (
            "## 核心教学\n\n"
            "概念必须同时说明定义、成立条件和适用边界，新情境可用于核对这三项标准。"
        )

    monkeypatch.setattr(service, "_call_llm", fake_call)
    outline = {
        "node_id": "L2-1-1",
        "node_name": "轻量讲解",
        "module_plan": [{
            "module_id": "core_explanation",
            "label": "核心教学",
        }],
    }
    plan = {
        "node_id": "L2-1-1",
        "teaching_modules": [{"module_id": "core_explanation"}],
    }

    async def scenario():
        return await asyncio.gather(*(
            service.generate_teacher_script_section(
                course_id=f"course-{index}",
                outline_section=outline,
                confirmed_plan_section=plan,
            )
            for index in range(2)
        ))

    results = asyncio.run(scenario())

    assert all(item["quality_report"]["passed"] for item in results)
    assert peak == 1


def test_script_job_keeps_completed_blocks_and_resumes_only_missing_work(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    service = TeacherLessonAuthoringService(repository)
    plan = standard_lesson_plan()
    plan["sections"][0]["teaching_modules"].append({
        "module_id": "feedback_check",
        "teaching_purpose": "检查学生是否掌握判断标准",
        "knowledge_names": ["核心概念"],
        "planned_minutes": 8,
        "teacher_activity": "给出新情境并追问判断依据。",
        "student_activity": "独立判断并说明理由。",
    })
    lesson = repository.save_plan_revision(
        "course-1",
        "L1-1",
        plan,
        source_outline_revision_id="outline-v1",
        quality_report=validate_teacher_lesson_plan(plan),
    )
    plan_revision = lesson["working_revision_id"]
    repository.confirm_plan_revision("course-1", "L1-1", plan_revision)
    outline_section = {
        "node_id": "L2-1-1",
        "node_name": "1.1 核心概念",
        "module_plan": [
            {"module_id": "core_explanation", "label": "核心教学"},
            {"module_id": "feedback_check", "label": "检查与反馈"},
        ],
    }

    first_job = repository.create_job(
        "course-1",
        "L1-1",
        job_type="teacher_lesson_script_generation",
        request_id="script-fail-on-second-block",
    )

    async def first_generator(_outline, _plan, module, _completed):
        if module["module_id"] == "feedback_check":
            raise RuntimeError("provider interrupted")
        return "我们先把判断框架立起来：核心概念由定义、成立条件与适用边界构成，正反例共同界定可检查的标准。"

    failed = asyncio.run(service.run_script_job(
        course_id="course-1",
        lesson_unit_id="L1-1",
        job_id=first_job["id"],
        source_plan_revision_id=plan_revision,
        outline_sections=[outline_section],
        plan_sections={"L2-1-1": plan["sections"][0]},
        generator=first_generator,
    ))

    assert failed["status"] == "failed"
    assert failed["completed_blocks"] == 1
    assert failed["total_blocks"] == 2
    assert failed["result_sections"][0]["blocks"][0]["module_id"] == "core_explanation"
    assert repository.lesson("course-1", "L1-1")["working_script_revision_id"] == ""

    second_job = repository.create_job(
        "course-1",
        "L1-1",
        job_type="teacher_lesson_script_generation",
        request_id="script-resume",
    )
    resumed_modules = []

    async def resume_generator(_outline, _plan, module, shard_context):
        resumed_modules.append(module["module_id"])
        assert shard_context["schema_version"] == "teacher_script_shard_context_v1"
        assert shard_context["previous_block"]["title"] == "核心教学"
        assert "content" not in shard_context["previous_block"]
        return "最后请大家判断一个新情境，逐项核对定义、成立条件和边界。典型错误是漏掉条件；核对标准是三项齐全，发现错误后要说明修正原因。"

    completed = asyncio.run(service.run_script_job(
        course_id="course-1",
        lesson_unit_id="L1-1",
        job_id=second_job["id"],
        source_plan_revision_id=plan_revision,
        outline_sections=[outline_section],
        plan_sections={"L2-1-1": plan["sections"][0]},
        generator=resume_generator,
        seed_sections=failed["result_sections"],
    ))

    assert completed["status"] == "completed"
    assert completed["stream_mode"] == "buffered_fallback"
    assert resumed_modules == ["feedback_check"]
    revision = repository.lesson("course-1", "L1-1")["script_revisions"][0]
    assert [item["module_id"] for item in revision["sections"][0]["blocks"]] == [
        "core_explanation",
        "feedback_check",
    ]


def test_script_job_runs_blocks_concurrently_and_streams_real_token_shards(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    service = TeacherLessonAuthoringService(repository)
    plan = standard_lesson_plan()
    plan["sections"][0]["teaching_modules"].append({
        "module_id": "feedback_check",
        "teaching_purpose": "检查学生是否掌握判断标准",
        "knowledge_names": ["核心概念"],
        "planned_minutes": 8,
        "teacher_activity": "给出新情境并追问判断依据。",
        "student_activity": "独立判断并说明理由。",
    })
    lesson = repository.save_plan_revision(
        "course-1",
        "L1-1",
        plan,
        source_outline_revision_id="outline-v1",
        quality_report=validate_teacher_lesson_plan(plan),
    )
    job = repository.create_job(
        "course-1",
        "L1-1",
        job_type="teacher_lesson_script_generation",
        request_id="script-parallel-blocks",
    )
    outline_section = {
        "node_id": "L2-1-1",
        "node_name": "1.1 核心概念",
        "module_plan": [
            {"module_id": "core_explanation", "label": "核心教学"},
            {"module_id": "feedback_check", "label": "检查与反馈"},
        ],
    }
    started: set[str] = set()
    observed_live_blocks: set[str] = set()
    both_started = asyncio.Event()
    release = asyncio.Event()

    async def generator(
        _outline,
        _plan,
        module,
        shard_context,
        *,
        on_content_delta,
        on_content_reset,
    ):
        started.add(module["module_id"])
        assert "previous_script_blocks" not in shard_context
        if len(started) == 2:
            both_started.set()
        await release.wait()
        if module["module_id"] == "core_explanation":
            content = "我们先把判断框架立起来：核心概念由定义、成立条件与适用边界构成，正反例共同界定可检查的标准。"
        else:
            content = "现在请大家判断一个新情境。典型错误是漏掉成立条件；修正原因是边界不完整；核对标准是定义、条件和边界三项齐全。"
        split_at = len(content) // 2
        await on_content_reset()
        await on_content_delta(content[:split_at])
        await on_content_delta(content[split_at:])
        live_job = repository.get_job("course-1", job["id"])
        assert any(
            streamed == content
            for shard_id, streamed in live_job["stream_batches"].items()
            if shard_id.endswith(f":{module['block_id']}")
        )
        observed_live_blocks.add(module["module_id"])
        return content

    async def scenario():
        task = asyncio.create_task(service.run_script_job(
            course_id="course-1",
            lesson_unit_id="L1-1",
            job_id=job["id"],
            source_plan_revision_id=lesson["working_revision_id"],
            outline_sections=[outline_section],
            plan_sections={"L2-1-1": plan["sections"][0]},
            generator=generator,
        ))
        await asyncio.wait_for(both_started.wait(), timeout=0.5)
        release.set()
        return await task

    completed = asyncio.run(scenario())

    assert completed["status"] == "completed"
    assert completed["stream_mode"] == "token_stream"
    assert completed["stream_events"] == []
    assert completed["stream_batches"] == {}
    assert observed_live_blocks == {"core_explanation", "feedback_check"}
    assert [
        item["module_id"] for item in completed["result_sections"][0]["blocks"]
    ] == ["core_explanation", "feedback_check"]
    app = FastAPI()
    app.include_router(teacher_lesson_router.router, prefix="/api")
    app.dependency_overrides[
        get_teacher_lesson_authoring_repository
    ] = lambda: repository
    with TestClient(app) as client:
        with client.stream(
            "GET",
            f"/api/teacher/courses/course-1/lesson-jobs/{job['id']}/stream",
        ) as response:
            payload_text = "".join(response.iter_text())
    data_line = next(
        line.removeprefix("data: ")
        for line in payload_text.splitlines()
        if line.startswith("data: ")
    )
    stream_payload = json.loads(data_line)
    assert stream_payload["lesson_unit_id"] == "L1-1"
    assert stream_payload["event"] == "lesson_script_complete"
    assert stream_payload["job"]["status"] == "completed"
    assert stream_payload["job"]["stream_complete"] is True
    assert stream_payload["job"]["result_sections"]


def test_script_job_runs_bounded_shards_concurrently_and_retries_only_failed_shard(
    tmp_path,
    monkeypatch,
):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    service = TeacherLessonAuthoringService(repository)
    module_ids = [f"custom_{index}" for index in range(1, 5)]
    plan = standard_lesson_plan()
    plan["sections"][0]["teaching_modules"] = [
        {
            "module_id": module_id,
            "teaching_purpose": f"完成第 {index} 个判断任务",
            "knowledge_names": ["核心概念"],
            "planned_minutes": 1,
            "teacher_activity": f"引导学生核对第 {index} 项标准。",
            "student_activity": f"完成第 {index} 个情境判断。",
        }
        for index, module_id in enumerate(module_ids, start=1)
    ]
    lesson = repository.save_plan_revision(
        "course-1",
        "L1-1",
        plan,
        source_outline_revision_id="outline-v1",
        quality_report={"passed": True},
    )
    plan_revision = lesson["working_revision_id"]
    repository.confirm_plan_revision("course-1", "L1-1", plan_revision)
    outline_section = {
        "node_id": "L2-1-1",
        "node_name": "1.1 核心概念",
        "module_plan": [
            {"module_id": module_id, "label": f"教学环节 {index}"}
            for index, module_id in enumerate(module_ids, start=1)
        ],
    }
    contract = compile_teacher_script_module_contract(
        outline_section,
        plan["sections"][0],
    )
    ordered_block_ids = [item["block_id"] for item in contract["modules"]]
    content_values = [
        "先看红色标本：我们把对象的定义、成立条件和排除边界分开。请给出一句可核对的判断并说明依据。",
        "蓝色案例带来了新问题：条件都相似，但关键变量已经越界。大家用反例找出失效位置，再改写原结论。",
        "接下来把绿色情境换成一组真实数据。先列出可观察事实，再选择匹配规则，最后检查结果是否能被数据支持。",
        "最后处理灰色任务：它故意隐去一项前提。请找到缺口，补齐必要信息，然后用同一检查方法验证修正后的答案。",
    ]
    contents = dict(zip(ordered_block_ids, content_values, strict=True))

    # Force four complete blocks into two adjacent two-block request shards.
    monkeypatch.setattr(
        teacher_script_module,
        "SCRIPT_SINGLE_REQUEST_TARGET_CHARACTERS",
        1,
    )
    monkeypatch.setattr(
        teacher_script_module,
        "SCRIPT_SINGLE_REQUEST_MAX_CHARACTERS",
        1,
    )
    monkeypatch.setattr(
        teacher_script_module,
        "SCRIPT_SHARD_TARGET_CHARACTERS",
        1800,
    )
    monkeypatch.setattr(
        teacher_script_module,
        "SCRIPT_SHARD_MAX_CHARACTERS",
        3600,
    )

    first_job = repository.create_job(
        "course-1",
        "L1-1",
        job_type="teacher_lesson_script_generation",
        request_id="script-shard-partial-failure",
    )
    started: list[str] = []
    both_started = asyncio.Event()
    release = asyncio.Event()

    async def unused_block_generator(*_args, **_kwargs):
        raise AssertionError("shard generator must own real requests")

    async def first_shard_generator(
        entries,
        shard_context,
        *,
        on_block_delta,
        on_shard_reset,
    ):
        shard_id = shard_context["shard_id"]
        started.append(shard_id)
        assert shard_context["budget_mode"] == "bounded_shards"
        assert all(
            "content" not in item
            for item in shard_context["block_directory"]
        )
        if len(started) == 2:
            both_started.set()
        await release.wait()
        if shard_context["sequence"] == 2:
            raise RuntimeError("second shard interrupted")
        result = {}
        for entry in entries:
            block_id = entry["module"]["block_id"]
            content = contents[block_id]
            split_at = len(content) // 2
            await on_shard_reset(block_id)
            await on_block_delta(block_id, content[:split_at])
            await on_block_delta(block_id, content[split_at:])
            result[block_id] = content
        return result

    async def first_scenario():
        task = asyncio.create_task(service.run_script_job(
            course_id="course-1",
            lesson_unit_id="L1-1",
            job_id=first_job["id"],
            source_plan_revision_id=plan_revision,
            outline_sections=[outline_section],
            plan_sections={"L2-1-1": plan["sections"][0]},
            generator=unused_block_generator,
            shard_generator=first_shard_generator,
        ))
        await asyncio.wait_for(both_started.wait(), timeout=0.5)
        release.set()
        return await task

    failed = asyncio.run(first_scenario())

    assert failed["status"] == "failed"
    assert len(started) == 2
    assert failed["completed_blocks"] == 2
    assert len(failed["result_sections"][0]["blocks"]) == 2
    assert repository.lesson("course-1", "L1-1")["script_revisions"] == []
    failed_shard_id = failed["error"]["failed_shards"][0]["shard_id"]
    assert failed["stream_events"] == []
    assert failed["stream_batches"] == {}

    retry_job = repository.create_job(
        "course-1",
        "L1-1",
        job_type="teacher_lesson_script_generation",
        request_id="script-shard-resume",
    )
    retry_shards: list[str] = []

    async def retry_shard_generator(
        entries,
        shard_context,
        *,
        on_block_delta,
        on_shard_reset,
    ):
        retry_shards.append(shard_context["shard_id"])
        result = {}
        for entry in entries:
            block_id = entry["module"]["block_id"]
            content = contents[block_id]
            await on_shard_reset(block_id)
            await on_block_delta(block_id, content)
            result[block_id] = content
        return result

    completed = asyncio.run(service.run_script_job(
        course_id="course-1",
        lesson_unit_id="L1-1",
        job_id=retry_job["id"],
        source_plan_revision_id=plan_revision,
        outline_sections=[outline_section],
        plan_sections={"L2-1-1": plan["sections"][0]},
        generator=unused_block_generator,
        shard_generator=retry_shard_generator,
        seed_sections=failed["result_sections"],
    ))

    assert completed["status"] == "completed", completed.get("error")
    assert retry_shards == [failed_shard_id]
    assert [
        item["block_id"] for item in completed["result_sections"][0]["blocks"]
    ] == ordered_block_ids
    assert completed["stream_events"] == []
    assert completed["stream_batches"] == {}


def test_script_resume_discards_only_invalid_checkpoint_block(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    service = TeacherLessonAuthoringService(repository)
    plan = standard_lesson_plan()
    lesson = repository.save_plan_revision(
        "course-1",
        "L1-1",
        plan,
        source_outline_revision_id="outline-v1",
        quality_report=validate_teacher_lesson_plan(plan),
    )
    plan_revision = lesson["working_revision_id"]
    repository.confirm_plan_revision("course-1", "L1-1", plan_revision)
    outline_section = {
        "node_id": "L2-1-1",
        "node_name": "1.1 核心概念",
        "module_plan": [{
            "module_id": "core_explanation",
            "label": "核心教学",
        }],
    }
    job = repository.create_job(
        "course-1",
        "L1-1",
        job_type="teacher_lesson_script_generation",
        request_id="script-invalid-checkpoint",
    )
    generated = []

    async def generator(_outline, _plan, module, _completed):
        generated.append(module["module_id"])
        return "我们重新把计算补完整：$F_x=6-3=3$。现在请大家再按方向规定核对一次结果。"

    completed = asyncio.run(service.run_script_job(
        course_id="course-1",
        lesson_unit_id="L1-1",
        job_id=job["id"],
        source_plan_revision_id=plan_revision,
        outline_sections=[outline_section],
        plan_sections={"L2-1-1": plan["sections"][0]},
        generator=generator,
        seed_sections=[{
            "section_node_id": "L2-1-1",
            "blocks": [{
                "block_id": compile_teacher_script_module_contract(
                    outline_section,
                    plan["sections"][0],
                )["modules"][0]["block_id"],
                "module_id": "core_explanation",
                "title": "核心教学",
                "content": "被截断的公式 $F_x=6-3",
            }],
        }],
    ))

    assert completed["status"] == "completed", completed.get("error")
    assert generated == ["core_explanation"]


def test_script_resume_restores_a_missing_middle_block_in_contract_order(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    service = TeacherLessonAuthoringService(repository)
    plan = standard_lesson_plan()
    plan["sections"][0]["teaching_modules"] = [
        {
            "module_id": "lesson_goal",
            "planned_minutes": 3,
            "teaching_purpose": "明确本节任务",
            "knowledge_names": ["核心概念"],
        },
        {
            "module_id": "learner_action",
            "planned_minutes": 4,
            "teaching_purpose": "完成课堂任务",
            "knowledge_names": ["核心概念"],
        },
        {
            "module_id": "feedback_check",
            "planned_minutes": 3,
            "teaching_purpose": "核对结果",
            "knowledge_names": ["核心概念"],
        },
    ]
    lesson = repository.save_plan_revision(
        "course-1",
        "L1-1",
        plan,
        source_outline_revision_id="outline-v1",
        quality_report={"passed": True},
    )
    plan_revision = lesson["working_revision_id"]
    repository.confirm_plan_revision("course-1", "L1-1", plan_revision)
    outline_section = {
        "node_id": "L2-1-1",
        "node_name": "1.1 核心概念",
        "module_plan": [
            {"module_id": "lesson_goal", "label": "本节任务"},
            {"module_id": "learner_action", "label": "学习者行动"},
            {"module_id": "feedback_check", "label": "检查与反馈"},
        ],
    }
    contract = compile_teacher_script_module_contract(
        outline_section,
        plan["sections"][0],
    )
    by_module = {
        item["module_id"]: item for item in contract["modules"]
    }
    seed_sections = [{
        "section_node_id": "L2-1-1",
        "blocks": [
            {
                **by_module["lesson_goal"],
                "content": "我们先明确本节的判断目标、适用条件和最终可检查的课堂产出。",
                "generation_source": "model",
            },
            {
                **by_module["learner_action"],
                "content": "这是旧任务检查点中保留的本地恢复内容，恢复后必须由模型重新生成并通过硬校验。",
                "generation_source": "local_recovery",
            },
            {
                **by_module["feedback_check"],
                "content": "最后请大家逐项核对结论、推理依据和边界。典型错误是只给结论；核对标准是三项齐全，发现错误后要说明修正原因并重新回答。",
                "generation_source": "model",
            },
        ],
    }]
    job = repository.create_job(
        "course-1",
        "L1-1",
        job_type="teacher_lesson_script_generation",
        request_id="script-resume-middle",
    )
    generated = []

    async def generator(_outline, _plan, module, _completed):
        generated.append(module["module_id"])
        return "接下来请大家独立分析一个新情境，写出判断结论、使用依据和结果检查，再与同伴比较差异。"

    completed = asyncio.run(service.run_script_job(
        course_id="course-1",
        lesson_unit_id="L1-1",
        job_id=job["id"],
        source_plan_revision_id=plan_revision,
        outline_sections=[outline_section],
        plan_sections={"L2-1-1": plan["sections"][0]},
        generator=generator,
        seed_sections=seed_sections,
    ))

    assert completed["status"] == "completed", completed.get("error")
    assert generated == ["learner_action"]
    revision = repository.lesson("course-1", "L1-1")["script_revisions"][0]
    blocks = revision["sections"][0]["blocks"]
    assert [item["module_id"] for item in blocks] == [
        "lesson_goal",
        "learner_action",
        "feedback_check",
    ]
    assert [item["generation_source"] for item in blocks] == [
        "model",
        "model",
        "model",
    ]


def test_script_resume_regenerates_repetitive_checkpoint_blocks(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    service = TeacherLessonAuthoringService(repository)
    plan = standard_lesson_plan()
    plan["sections"][0]["teaching_modules"] = [
        {
            "module_id": "lesson_goal",
            "planned_minutes": 2,
            "teaching_purpose": "明确目标",
            "knowledge_names": ["核心概念"],
        },
        {
            "module_id": "core_explanation",
            "planned_minutes": 3,
            "teaching_purpose": "解释概念",
            "knowledge_names": ["核心概念"],
        },
        {
            "module_id": "feedback_check",
            "planned_minutes": 2,
            "teaching_purpose": "核对结论",
            "knowledge_names": ["核心概念"],
        },
    ]
    lesson = repository.save_plan_revision(
        "course-1",
        "L1-1",
        plan,
        source_outline_revision_id="outline-v1",
        quality_report={"passed": True},
    )
    plan_revision = lesson["working_revision_id"]
    repository.confirm_plan_revision("course-1", "L1-1", plan_revision)
    outline_section = {
        "node_id": "L2-1-1",
        "node_name": "1.1 核心概念",
        "module_plan": [
            {"module_id": "lesson_goal", "label": "本节任务"},
            {"module_id": "core_explanation", "label": "核心教学"},
            {"module_id": "feedback_check", "label": "检查与反馈"},
        ],
    }
    contract = compile_teacher_script_module_contract(
        outline_section,
        plan["sections"][0],
    )
    repeated = "概念需要同时说明定义、成立条件与适用边界，并通过正例和反例逐项核对判断标准。"
    seed_sections = [{
        "section_node_id": "L2-1-1",
        "title": "1.1 核心概念",
        "blocks": [
            {**module, "content": repeated, "generation_source": "model"}
            for module in contract["modules"]
        ],
    }]
    job = repository.create_job(
        "course-1",
        "L1-1",
        job_type="teacher_lesson_script_generation",
        request_id="script-resume-repetition",
    )
    generated = []
    replacements = {
        "core_explanation": "接着来看核心概念：它由对象、成立条件和适用边界共同界定；正例验证条件，反例负责暴露边界。",
        "feedback_check": "最后请大家先检查对象是否满足条件，再比较结论与边界；若判断错误，必须指出具体违反哪一项。",
    }

    async def generator(_outline, _plan, module, shard_context):
        generated.append(module["module_id"])
        assert shard_context["schema_version"] == "teacher_script_shard_context_v1"
        assert shard_context["previous_block"] is not None
        assert "content" not in shard_context["previous_block"]
        return replacements[module["module_id"]]

    completed = asyncio.run(service.run_script_job(
        course_id="course-1",
        lesson_unit_id="L1-1",
        job_id=job["id"],
        source_plan_revision_id=plan_revision,
        outline_sections=[outline_section],
        plan_sections={"L2-1-1": plan["sections"][0]},
        generator=generator,
        seed_sections=seed_sections,
    ))

    assert completed["status"] == "completed", completed.get("error")
    assert set(generated) == {"core_explanation", "feedback_check"}
    revision = repository.lesson("course-1", "L1-1")["script_revisions"][0]
    assert revision["quality_report"]["passed"] is True


def test_script_fallback_content_is_rejected_without_formal_revision(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    service = TeacherLessonAuthoringService(repository)
    plan = standard_lesson_plan()
    lesson = repository.save_plan_revision(
        "course-1",
        "L1-1",
        plan,
        source_outline_revision_id="outline-v1",
        quality_report=validate_teacher_lesson_plan(plan),
    )
    plan_revision = lesson["working_revision_id"]
    repository.confirm_plan_revision("course-1", "L1-1", plan_revision)
    outline_section = {
        "node_id": "L2-1-1",
        "node_name": "1.1 核心概念",
        "module_plan": [{
            "module_id": "core_explanation",
            "label": "核心教学",
        }],
    }
    job = repository.create_job(
        "course-1",
        "L1-1",
        job_type="teacher_lesson_script_generation",
        request_id="script-provider-fallback",
    )
    async def fallback_generator(_outline, _plan, module, _completed):
        raise TeacherLessonAuthoringError(
            "lesson_script_provider_failed",
            f"{module.get('title') or '教学块'}生成失败，请重试。",
        )

    completed = asyncio.run(service.run_script_job(
        course_id="course-1",
        lesson_unit_id="L1-1",
        job_id=job["id"],
        source_plan_revision_id=plan_revision,
        outline_sections=[outline_section],
        plan_sections={"L2-1-1": plan["sections"][0]},
        generator=fallback_generator,
    ))

    assert completed["status"] == "failed"
    assert completed["error"]["code"] == "lesson_script_provider_failed"
    assert completed["completed_blocks"] == 0
    assert repository.lesson("course-1", "L1-1")["script_revisions"] == []


def test_legacy_script_adapter_keeps_the_original_body_as_one_compatibility_block():
    original = "## 老师原有标题\n\n老师原本的一整篇课堂讲稿。"
    migrated = normalize_teacher_script_section({
        "section_node_id": "L2-1-1",
        "title": "1.1 核心概念",
        "content": original,
    })

    assert len(migrated["blocks"]) == 1
    assert migrated["blocks"][0]["module_id"] == "legacy_script"
    assert migrated["blocks"][0]["content"] == original


def test_teacher_lesson_v6_source_accepts_confirmed_legacy_script_sections():
    source = course_data()
    source["blueprint_revision_id"] = "outline-v1"
    document, view, _synthetic_id = teacher_lesson_v6_source(
        source,
        lesson_unit_id="L1-1",
        plan_revision={
            "revision_id": "plan-1",
            "plan": standard_lesson_plan(),
        },
        script_revision={
            "revision_id": "script-legacy",
            "sections": [
                {
                    "section_node_id": "L2-1-1",
                    "title": "1.1 核心概念",
                    "content": "这是老版已确认讲稿的正文。",
                },
                {
                    "section_node_id": "L2-1-2",
                    "title": "1.2 能力迁移",
                    "content": "这是第二小节的已确认讲稿。",
                },
            ],
        },
    )

    assert document.blocks
    assert view["nodes"][1]["content_blocks"][0]["metadata"]["module_id"] == "legacy_script"
    assert view["nodes"][1]["content_blocks"][0]["content"] == "这是老版已确认讲稿的正文。"


def test_canonical_outline_revision_recovers_current_state_and_blocks_weak_plan(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    repository.set_outline("course-1", "blueprint-v1")
    quality = validate_teacher_lesson_plan(
        standard_lesson_plan(),
        expected_section_ids=["L2-1-1"],
        expected_outline_revision_id="knowledge-scope-v1",
        source_outline_revision_id="knowledge-scope-v1",
    )
    lesson = repository.save_plan_revision(
        "course-1",
        "L1-1",
        standard_lesson_plan(),
        source_outline_revision_id="knowledge-scope-v1",
        quality_report=quality,
    )
    assert lesson["source_state"] == "stale"

    repository.set_outline("course-1", "knowledge-scope-v1")
    recovered = repository.lesson("course-1", "L1-1")
    assert recovered["source_state"] == "current"
    confirmed = repository.confirm_plan_revision(
        "course-1",
        "L1-1",
        recovered["working_revision_id"],
    )
    assert confirmed["confirmed_revision_id"] == recovered["working_revision_id"]

    weak = repository.save_plan_revision(
        "course-1",
        "L1-1",
        {"schema_version": "course_teaching_plan_v3", "sections": [{"node_id": "L2-1-1"}]},
        source_outline_revision_id="knowledge-scope-v1",
    )
    try:
        repository.confirm_plan_revision(
            "course-1",
            "L1-1",
            weak["working_revision_id"],
        )
    except Exception as exc:
        assert getattr(exc, "code", "") == "lesson_plan_quality_blocked"
    else:
        raise AssertionError("不完整教案不应该被确认")


def test_outline_only_lesson_scope_reuses_existing_pedagogy_compiler(monkeypatch):
    service = CourseService()
    captured = {}

    async def fake_prepare(*, course_data, plan, **_kwargs):
        planned_sections = plan["chapters"][0]["sections"]
        captured["module_ids"] = [
            [item["module_id"] for item in section.get("module_plan") or []]
            for section in planned_sections
        ]
        course_data["course_teaching_plan"] = {
            "schema_version": "course_teaching_plan_v3",
            "sections": [
                {"node_id": section["node_id"], "teaching_modules": []}
                for section in planned_sections
            ],
        }
        course_data.setdefault("generation_stage_artifacts", {})["course_teaching_plan"] = {
            "source_outline_revision_id": "outline-v1",
            "fallback_units": [],
        }
        return plan

    monkeypatch.setattr(service, "_prepare_course_teaching_plan", fake_prepare)

    result = asyncio.run(service.prepare_teacher_lesson_plan(
        course_data={**lecture_course_data(), "blueprint_revision_id": "outline-v1"},
        lesson_unit_id="L1-1",
    ))

    assert all("core_explanation" in module_ids for module_ids in captured["module_ids"])
    assert [item["node_id"] for item in result["plan"]["sections"]] == ["L2-1-1", "L2-1-2"]


def test_teacher_lesson_plan_resume_reuses_planner_checkpoint(monkeypatch):
    service = CourseService()
    emitted = []

    async def fake_prepare(*, course_data, plan, on_checkpoint, **_kwargs):
        stage = course_data["generation_stage_artifacts"]["course_teaching_plan"]
        assert stage["batches"]["TP-B01"]["status"] == "completed"
        course_data["course_teaching_plan"] = {
            "schema_version": "course_teaching_plan_v3",
            "sections": [
                {"node_id": section["node_id"], "teaching_modules": []}
                for section in plan["chapters"][0]["sections"]
            ],
        }
        stage["source_outline_revision_id"] = "outline-v1"
        stage["fallback_units"] = []
        await on_checkpoint(course_data)
        return plan

    monkeypatch.setattr(service, "_prepare_course_teaching_plan", fake_prepare)
    result = asyncio.run(service.prepare_teacher_lesson_plan(
        course_data={**course_data(), "blueprint_revision_id": "outline-v1"},
        lesson_unit_id="L1-1",
        resume_checkpoint={
            "planner_course_data": {
                "generation_stage_artifacts": {
                    "course_teaching_plan": {
                        "batches": {"TP-B01": {"status": "completed"}},
                    },
                },
            },
        },
        on_checkpoint=lambda checkpoint: emitted.append(checkpoint),
    ))

    assert result["plan"]["sections"]
    assert emitted[0]["schema_version"] == "teacher_lesson_plan_checkpoint_v1"
    assert emitted[0]["planner_course_data"]["generation_stage_artifacts"]["course_teaching_plan"]["batches"]["TP-B01"]["status"] == "completed"


def test_uploaded_pptx_is_extracted_as_immutable_lesson_evidence(tmp_path):
    from pptx import Presentation

    source = tmp_path / "旧课件.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "DeepSeek 4"
    slide.placeholders[1].text = "从模型能力变化到课堂案例"
    presentation.save(source)

    evidence = extract_uploaded_pptx_evidence(source, asset_id="asset-1")

    assert len(evidence) == 1
    assert evidence[0]["evidence_id"] == "uploaded-ppt-asset-1-slide-1"
    assert "DeepSeek 4" in evidence[0]["source_text"]
    assert source.is_file()


def test_uploaded_pptx_review_indexes_editable_blocks_and_keeps_original(tmp_path):
    from pptx import Presentation

    source = tmp_path / "老师原稿.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "核心概念"
    slide.placeholders[1].text = "定义、条件与边界"
    presentation.save(source)
    original = source.read_bytes()

    parsed = extract_uploaded_pptx_review(
        source,
        asset_id="asset-1",
        filename="老师原稿.pptx",
    )

    assert parsed["source_filename"] == "老师原稿.pptx"
    assert parsed["slides"][0]["title"] == "核心概念"
    assert all("shape_index" in block for block in parsed["slides"][0]["blocks"])
    assert any(block["editable"] for block in parsed["slides"][0]["blocks"])
    assert source.read_bytes() == original


def test_uploaded_ppt_review_report_names_sources_without_fake_score():
    slides = [{
        "slide_id": "slide-1",
        "slide_number": 1,
        "title": "另一个话题",
        "blocks": [{"block_id": "b1", "kind": "title", "text": "另一个话题", "editable": True}],
    }]
    report = build_uploaded_ppt_review_report(
        slides,
        sources=[
            {"kind": "lesson_plan", "label": "已确认教案", "revision_id": "plan-1", "status": "confirmed"},
            {"kind": "script", "label": "已确认讲稿", "revision_id": "script-1", "status": "confirmed"},
        ],
        reference_units=[{
            "kind": "script",
            "label": "核心概念",
            "revision_id": "script-1",
            "text": "核心概念的定义、条件与适用边界",
        }],
    )

    assert report["summary"]["finding_count"] >= 1
    assert report["findings"][0]["confidence"] == "high"
    assert "score" not in report["summary"]
    assert {item["label"] for item in report["sources"]} == {"已确认教案", "已确认讲稿"}


def test_imported_ppt_review_requires_revision_and_confirmation(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    repository.set_outline("course-1", "outline-v1")
    review = repository.save_imported_ppt_review(
        "course-1",
        "L1-1",
        package_id="package-1",
        source_asset_id="asset-1",
        source_filename="老师原稿.pptx",
        slides=[{"slide_id": "slide-1", "slide_number": 1, "title": "原标题", "blocks": []}],
        report={"findings": [], "sources": []},
        source_outline_revision_id="outline-v1",
        source_lesson_plan_revision_id="",
        source_script_revision_id="",
        actor="teacher-1",
    )
    updated = repository.replace_imported_ppt_review(
        "course-1",
        "L1-1",
        review_id=review["review_id"],
        base_revision_id=review["revision_id"],
        slides=[{"slide_id": "slide-1", "slide_number": 1, "title": "新标题", "blocks": []}],
        report={"findings": [], "sources": []},
        actor="teacher-1",
    )

    with pytest.raises(TeacherLessonAuthoringError, match="已更新"):
        repository.confirm_imported_ppt_review(
            "course-1",
            "L1-1",
            review_id=review["review_id"],
            revision_id=review["revision_id"],
        )
    confirmed = repository.confirm_imported_ppt_review(
        "course-1",
        "L1-1",
        review_id=review["review_id"],
        revision_id=updated["revision_id"],
    )
    lesson = repository.lesson("course-1", "L1-1")
    assert confirmed["status"] == "confirmed"
    assert lesson["ppt_assets"][-1]["engine"] == "uploaded_pptx"
    assert len(confirmed["revision_history"]) == 2


def test_repository_keeps_sibling_lesson_assets_independent(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    repository.set_outline("course-1", "outline-v1")
    first = repository.save_plan_revision(
        "course-1",
        "L1-1",
        {"sections": [{"node_id": "L2-1-1"}]},
        source_outline_revision_id="outline-v1",
    )
    second = repository.save_plan_revision(
        "course-1",
        "L1-2",
        {"sections": [{"node_id": "L2-2-1"}]},
        source_outline_revision_id="outline-v1",
    )

    assert first["working_revision_id"] != second["working_revision_id"]
    view = repository.view("course-1")
    assert set(view["lessons"]) == {"L1-1", "L1-2"}
    assert len(view["lessons"]["L1-1"]["revisions"]) == 1


def test_plan_history_restore_creates_a_new_working_revision(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    repository.set_outline("course-1", "outline-v1")
    first = repository.save_plan_revision(
        "course-1",
        "L1-1",
        {"sections": [{"node_id": "L2-1-1", "learning_objective": "解释极限"}]},
        source_outline_revision_id="outline-v1",
    )
    first_revision = first["working_revision_id"]
    second = repository.save_plan_revision(
        "course-1",
        "L1-1",
        {"sections": [{"node_id": "L2-1-1", "learning_objective": "计算极限"}]},
        source_outline_revision_id="outline-v1",
    )
    second_revision = second["working_revision_id"]

    restored = repository.restore_plan_revision(
        "course-1",
        "L1-1",
        first_revision,
        expected_working_revision_id=second_revision,
        actor="teacher-1",
    )

    restored_revision = restored["working_revision_id"]
    assert restored_revision not in {first_revision, second_revision}
    saved = next(item for item in restored["revisions"] if item["revision_id"] == restored_revision)
    assert saved["generation_source"] == "history_restore"
    assert saved["restored_from_revision_id"] == first_revision
    assert saved["plan"]["sections"][0]["learning_objective"] == "解释极限"


def test_plan_fallback_fails_without_formal_revision(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    service = TeacherLessonAuthoringService(repository)
    job = repository.create_job(
        "course-1",
        "L1-1",
        request_id="request-1",
        source_outline_revision_id="outline-v1",
    )

    async def planner(_course, _lesson_id, on_progress):
        await on_progress("lesson_plan_validation", 70, "正在校验")
        return {
            "plan": {"sections": [{"node_id": "L2-1-1", "teaching_modules": []}]},
            "warnings": [{"code": "model_output_failed_validation"}],
            "generation_source": "deterministic_local_fallback",
            "source_refs": [{"source_kind": "uploaded_ppt", "asset_id": "asset-1", "slide": 1}],
            "source_outline_revision_id": "outline-v1",
        }

    completed = asyncio.run(service.run_plan_job(
        course_id="course-1",
        lesson_unit_id="L1-1",
        job_id=job["id"],
        course_data=course_data(),
        planner=planner,
    ))

    assert completed["status"] == "failed"
    assert completed["error"]["code"] == "lesson_plan_generation_incomplete"
    assert repository.view("course-1")["lessons"] == {}
    assert completed["message"] == "本讲教案生成失败"


def test_plan_job_progress_never_moves_backwards(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    repository.set_outline("course-1", "outline-v1")
    service = TeacherLessonAuthoringService(repository)
    job = repository.create_job(
        "course-1",
        "L1-1",
        request_id="request-monotonic-progress",
        source_outline_revision_id="outline-v1",
    )
    observed_progress = []

    async def planner(_course, _lesson_id, on_progress):
        await on_progress("lesson_plan_generation", 36, "正在生成")
        observed_progress.append(repository.get_job("course-1", job["id"])["progress"])
        await on_progress("lesson_plan_generation", 35, "仍在生成")
        observed_progress.append(repository.get_job("course-1", job["id"])["progress"])
        return {
            "plan": standard_lesson_plan(),
            "warnings": [],
            "generation_source": "model",
            "source_outline_revision_id": "outline-v1",
        }

    completed = asyncio.run(service.run_plan_job(
        course_id="course-1",
        lesson_unit_id="L1-1",
        job_id=job["id"],
        course_data=single_section_course_data(),
        planner=planner,
    ))

    assert observed_progress == [36, 36]
    assert completed["progress"] == 100


def test_plan_job_stream_updates_do_not_block_event_loop(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    repository.set_outline("course-1", "outline-v1")
    service = TeacherLessonAuthoringService(repository)
    job = repository.create_job(
        "course-1",
        "L1-1",
        request_id="request-responsive-progress",
        source_outline_revision_id="outline-v1",
    )
    stream_write_started = threading.Event()
    release_stream_write = threading.Event()
    original_update_job_stream = repository.update_job_stream

    def blocking_update_job_stream(*args, **kwargs):
        stream_write_started.set()
        release_stream_write.wait(timeout=0.5)
        return original_update_job_stream(*args, **kwargs)

    repository.update_job_stream = blocking_update_job_stream  # type: ignore[method-assign]
    progress_elapsed: list[float] = []

    async def planner(_course, _lesson_id, on_progress):
        async def release_from_event_loop() -> None:
            while not stream_write_started.is_set():
                await asyncio.sleep(0)
            await asyncio.sleep(0.01)
            release_stream_write.set()

        started_at = time.perf_counter()
        await asyncio.gather(
            on_progress(
                "course_teaching_plan_batch",
                40,
                "正在写入流式进度",
                40,
                {
                    "stream_event": "delta",
                    "stream_batch_id": "TP-B01",
                    "stream_delta": '{"title":"第一批"}',
                },
            ),
            release_from_event_loop(),
        )
        progress_elapsed.append(time.perf_counter() - started_at)
        return {
            "plan": standard_lesson_plan(),
            "warnings": [],
            "generation_source": "model",
            "source_outline_revision_id": "outline-v1",
        }

    completed = asyncio.run(service.run_plan_job(
        course_id="course-1",
        lesson_unit_id="L1-1",
        job_id=job["id"],
        course_data=single_section_course_data(),
        planner=planner,
    ))

    assert completed["status"] in {"completed", "completed_with_warnings"}
    assert progress_elapsed[0] < 0.2


def test_plan_job_keeps_formal_outline_and_planner_scope_revisions_separate(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    repository.set_outline("course-1", "outline-v1")
    service = TeacherLessonAuthoringService(repository)
    job = repository.create_job(
        "course-1",
        "L1-2",
        request_id="request-scope-revision",
        source_outline_revision_id="outline-v1",
    )
    plan = standard_lesson_plan()
    plan["sections"][0]["node_id"] = "L2-2-1"

    async def planner(_course, _lesson_id, _on_progress):
        return {
            "plan": plan,
            "warnings": [],
            "generation_source": "model",
            "source_outline_revision_id": "knowledge-scope-v2",
        }

    completed = asyncio.run(service.run_plan_job(
        course_id="course-1",
        lesson_unit_id="L1-2",
        job_id=job["id"],
        course_data=single_section_course_data("L1-2", "L2-2-1"),
        planner=planner,
    ))

    assert completed["status"] == "completed"
    revision = repository.lesson("course-1", "L1-2")["revisions"][0]
    assert revision["source_outline_revision_id"] == "outline-v1"
    assert revision["source_knowledge_scope_revision_id"] == "knowledge-scope-v2"
    assert revision["quality_report"]["passed"] is True
    assert completed["warnings"] == []


def test_failed_plan_job_discards_unvalidated_streamed_working_copy(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    service = TeacherLessonAuthoringService(repository)
    job = repository.create_job(
        "course-1",
        "L1-1",
        request_id="request-stream-failure",
        source_outline_revision_id="outline-v1",
    )

    async def planner(_course, _lesson_id, on_progress):
        await asyncio.gather(
            on_progress(
                "course_teaching_plan_batch", 45, "生成第一批", 45,
                {
                    "stream_event": "delta",
                    "stream_batch_id": "TP-B01",
                    "stream_delta": '{"title":"第一批已生成"',
                },
            ),
            on_progress(
                "course_teaching_plan_batch", 45, "生成第二批", 45,
                {
                    "stream_event": "delta",
                    "stream_batch_id": "TP-B02",
                    "stream_delta": '{"title":"第二批已生成"',
                },
            ),
        )
        raise RuntimeError("模型连接中断")

    failed = asyncio.run(service.run_plan_job(
        course_id="course-1",
        lesson_unit_id="L1-1",
        job_id=job["id"],
        course_data=course_data(),
        planner=planner,
    ))

    assert failed["status"] == "failed"
    assert failed["stream_complete"] is True
    assert failed["stream_batches"] == {}
    assert failed["stream_events"] == []


def test_orphaned_lesson_job_recovers_only_last_semantic_checkpoint(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    job = repository.create_job(
        "course-1",
        "L1-1",
        request_id="request-orphaned",
        source_outline_revision_id="outline-v1",
    )
    repository.update_job_stream(
        "course-1",
        job["id"],
        phase="course_teaching_plan_batch",
        progress=42,
        message="正在生成",
        batch_id="TP-B01",
        event="delta",
        delta='{"learning_objective":"已生成部分目标"',
    )
    live = repository.get_job("course-1", job["id"])
    assert "已生成部分目标" in live["stream_batches"]["TP-B01"]

    before_checkpoint = TeacherLessonAuthoringRepository(tmp_path).get_job(
        "course-1",
        job["id"],
    )
    assert before_checkpoint["stream_batches"] == {}

    repository.update_job(
        "course-1",
        job["id"],
        phase="course_teaching_plan_batch_saved",
        progress=60,
        checkpoint={
            "schema_version": "teacher_lesson_plan_checkpoint_v1",
            "validated_batch_ids": ["TP-B01"],
        },
    )
    path = tmp_path / "course-1.json"
    stored = json.loads(path.read_text(encoding="utf-8"))
    stored["jobs"][job["id"]]["updated_at"] = "2020-01-01T00:00:00+00:00"
    path.write_text(json.dumps(stored, ensure_ascii=False), encoding="utf-8")

    reloaded = TeacherLessonAuthoringRepository(tmp_path)
    expired = reloaded.expire_stale_job("course-1", job["id"])

    assert expired["status"] == "failed"
    assert expired["phase"] == "lesson_plan_interrupted"
    assert expired["error"]["retryable"] is True
    assert expired["stream_batches"] == {}
    assert expired["checkpoint"]["validated_batch_ids"] == ["TP-B01"]


def test_plan_v3_projection_is_editable_and_never_serializes_module_json():
    section = {
        "node_id": "L2-1-1",
        "key_points": ["二进制转换"],
        "knowledge_structure": [{
            "knowledge_points": [{
                "name": "二进制转换",
                "statement": "完成二进制与十进制之间的相互转换。",
                "boundaries": ["仅处理无符号整数"],
                "capability_points": [{"observable_behavior": "能够独立完成一次进制转换并核对结果。"}],
            }],
        }],
        "teaching_modules": [
            {
                "module_id": "core_explanation",
                "teaching_purpose": "按模板完成「核心讲解」",
                "teaching_guidance": "使用位权展开演示转换过程",
                "knowledge_names": ["二进制转换"],
            },
            {
                "module_id": "learner_action",
                "teaching_guidance": "学生独立完成一道转换题",
                "knowledge_names": ["二进制转换"],
            },
        ],
    }

    view = teacher_lesson_section_content(section)
    assert view["learning_objective"] == "能够独立完成一次进制转换并核对结果。"
    assert "仅处理无符号整数" in view["key_difficulties"]
    assert view["teacher_activities"] == ["核心讲解：围绕二进制转换；使用位权展开演示转换过程"]
    assert view["student_activities"] == ["学习者行动：围绕二进制转换；学生独立完成一道转换题"]

    normalized = normalize_teacher_lesson_plan({"sections": [section]})
    projected = normalized["sections"][0]
    assert projected["learning_objective"] == view["learning_objective"]
    assert projected["teacher_activities"] == view["teacher_activities"]
    assert "{" not in "\n".join(projected["teacher_activities"])


def test_v6_source_rejects_plan_only_fallback_without_confirmed_script():
    source = course_data()
    revision = {
        "revision_id": "fallback-v1",
        "plan": {
            "sections": [{
                "node_id": "L2-1-1",
                "knowledge_structure": [{
                    "knowledge_points": [{
                        "name": "位权展开",
                        "statement": "二进制数可按位权展开并转换为十进制。",
                    }],
                }],
                "teaching_modules": [{
                    "module_id": "core_explanation",
                    "teaching_purpose": "按模板完成「正式定义」",
                    "teaching_guidance": "逐位演示位权展开",
                    "knowledge_names": ["位权展开"],
                }],
            }],
        },
    }
    with pytest.raises(TeacherLessonAuthoringError) as exc_info:
        teacher_lesson_v6_source(
            source,
            lesson_unit_id="L1-1",
            plan_revision=revision,
            script_revision={"revision_id": "", "sections": []},
        )
    assert exc_info.value.code == "lesson_script_source_missing"


def test_request_id_is_idempotent(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    first = repository.create_job("course-1", "L1-1", request_id="same")
    second = repository.create_job("course-1", "L1-1", request_id="same")
    assert first["id"] == second["id"]


def test_teacher_job_contract_preserves_checkpoint_when_cancelled(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    job = repository.create_job(
        "course-1",
        "L1-1",
        job_type="teacher_lesson_script_generation",
        request_id="cancel-me",
    )
    repository.update_job(
        "course-1",
        job["id"],
        status="running",
        completed_blocks=1,
        current_block_id="block-2",
        result_sections=[{
            "section_node_id": "L2-1-1",
            "content": "已完成内容",
        }],
    )

    cancelled = repository.cancel_job("course-1", job["id"])

    assert cancelled["schema_version"] == "teacher_asset_job_v1"
    assert cancelled["asset_type"] == "script"
    assert cancelled["status"] == "cancelled"
    assert cancelled["stream_complete"] is True
    assert cancelled["checkpoint"]["completed_blocks"] == 1
    assert cancelled["checkpoint"]["result_sections"][0]["content"] == "已完成内容"
    assert cancelled["error"]["retryable"] is True


@pytest.mark.parametrize(
    "frozen_status",
    ["paused", "failed", "cancelled", "completed_with_warnings", "completed"],
)
def test_teacher_job_rejects_late_status_progress_and_stream_updates(
    tmp_path,
    frozen_status,
):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    job = repository.create_job(
        "course-1",
        "L1-1",
        job_type="teacher_lesson_script_generation",
        request_id=f"frozen-{frozen_status}",
    )
    repository.update_job(
        "course-1",
        job["id"],
        status="running",
        phase="lesson_script_generation",
        progress=40,
    )
    frozen = repository.update_job(
        "course-1",
        job["id"],
        status=frozen_status,
        phase="frozen_phase",
        progress=45,
        message="用户或终态已确立",
    )

    late_status = repository.update_job(
        "course-1",
        job["id"],
        status="completed",
        phase="late_complete",
        progress=100,
        message="迟到完成",
    )
    late_stream = repository.update_job_stream(
        "course-1",
        job["id"],
        phase="late_stream",
        progress=90,
        message="迟到增量",
        batch_id="block-1",
        event="delta",
        delta="不应写入",
        lesson_unit_id="L1-1",
        block_id="block-1",
    )

    assert late_status == frozen
    assert late_stream == frozen
    assert late_status["status"] == frozen_status
    assert late_status["phase"] == "frozen_phase"
    assert late_status["progress"] == 45
    assert late_status.get("stream_batches") == {}


def test_teacher_job_cannot_save_a_plan_revision_after_pause(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    job = repository.create_job(
        "course-1",
        "L1-1",
        request_id="paused-before-save",
    )
    repository.update_job("course-1", job["id"], status="running")
    repository.pause_job("course-1", job["id"])

    with pytest.raises(TeacherLessonAuthoringError) as exc_info:
        repository.save_plan_revision(
            "course-1",
            "L1-1",
            standard_lesson_plan(),
            source_outline_revision_id="outline-v1",
            active_job_id=job["id"],
        )

    assert exc_info.value.code == "teacher_job_not_active"
    assert repository.view("course-1")["lessons"] == {}


def test_teacher_only_course_is_hidden_from_student_list(monkeypatch):
    courses = [
        {"course_id": "student-course", "is_published": True},
        {
            "course_id": "teacher-course",
            "generation_job_id": "teacher-job",
        },
    ]
    monkeypatch.setattr(courses_router.storage, "list_courses", lambda: courses)
    monkeypatch.setattr(courses_router.learning_snapshot_repository, "load", lambda *_args: None)
    student = courses_router._list_courses_with_resume(
        "learner",
        {"teacher-job"},
        {"teacher-course"},
    )
    teacher = courses_router._list_teacher_courses({"teacher-job"})
    assert [item["course_id"] for item in student] == ["student-course"]
    assert [item["course_id"] for item in teacher] == ["student-course", "teacher-course"]


def test_ai_candidate_acceptance_creates_new_working_revision(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    lesson = repository.save_plan_revision(
        "course-1",
        "L1-1",
        {"sections": [{"node_id": "L2-1-1", "learning_objective": "before"}]},
        source_outline_revision_id="outline-v1",
    )
    candidate = repository.save_ai_candidate(
        "course-1",
        "L1-1",
        base_revision_id=lesson["working_revision_id"],
        instruction="优化目标",
        section_node_id="L2-1-1",
        plan={"sections": [{"node_id": "L2-1-1", "learning_objective": "after"}]},
    )
    accepted = repository.resolve_ai_candidate(
        "course-1",
        "L1-1",
        candidate["candidate_id"],
        accept=True,
    )

    assert len(accepted["revisions"]) == 2
    assert accepted["working_revision_id"] != lesson["working_revision_id"]
    assert accepted["revisions"][-1]["plan"]["sections"][0]["learning_objective"] == "after"
    assert accepted["ai_candidates"][0]["status"] == "accepted"


def test_ai_optimizer_uses_compact_editable_contract_and_merges_one_section():
    plan = {
        "schema_version": "course_teaching_plan_v3",
        "sections": [
            {
                "node_id": "L2-1-1",
                "learning_objective": "原目标",
                "key_points": ["进制转换"],
                "key_difficulties": ["原难点"],
                "in_class_checks": ["完成一道转换题"],
                "homework": ["原作业"],
                "teaching_notes": ["原备注"],
                "teaching_modules": [{
                    "module_id": "core_explanation",
                    "teaching_purpose": "讲清位权展开",
                    "planned_minutes": 15,
                    "teacher_activity": "原教师活动",
                    "student_activity": "原学生活动",
                    "knowledge_names": ["进制转换"],
                }],
                "knowledge_structure": [{"knowledge_points": [{"statement": "不可改写的事实"}]}],
            },
            {"node_id": "L2-1-2", "learning_objective": "兄弟小节"},
        ],
    }

    class FakeOptimizer:
        captured_prompt = ""

        async def _call_llm(self, prompt, **_kwargs):
            self.captured_prompt = prompt
            return json.dumps({
                "sections": [{
                    "node_id": "L2-1-1",
                    "learning_objective": "学生能够独立完成一次进制转换并解释步骤。",
                    "key_points": ["进制转换", "位权展开"],
                    "key_difficulties": ["位权展开"],
                    "in_class_checks": ["独立完成一道转换题并说明每一步。"],
                    "homework": ["完成两道相邻进制转换题。"],
                    "teaching_notes": ["先检查位权表，再处理转换步骤。"],
                    "teaching_modules": [{
                        "module_id": "core_explanation",
                        "teaching_purpose": "用演示和练习建立位权展开方法",
                        "planned_minutes": 18,
                        "teacher_activity": "演示十进制转二进制并逐步核对余数。",
                        "student_activity": "独立完成一道转换题并说明每一步。",
                    }],
                }],
            }, ensure_ascii=False)

        @staticmethod
        def _extract_json(value):
            return json.loads(value)

    fake = FakeOptimizer()
    result = asyncio.run(CourseService.optimize_teacher_lesson_plan(
        fake,
        plan=plan,
        instruction="增加可观察目标和课堂练习",
        section_node_id="L2-1-1",
    ))

    optimized = result["plan"]["sections"]
    assert optimized[0]["learning_objective"].startswith("学生能够独立")
    assert optimized[0]["teaching_modules"][0]["teacher_activity"].startswith("演示十进制")
    assert optimized[0]["teacher_activities"] == ["演示十进制转二进制并逐步核对余数。"]
    assert optimized[0]["knowledge_structure"] == plan["sections"][0]["knowledge_structure"]
    assert optimized[1]["node_id"] == plan["sections"][1]["node_id"]
    assert optimized[1]["learning_objective"] == plan["sections"][1]["learning_objective"]
    assert '"schema_version"' not in fake.captured_prompt
    assert '"knowledge_context"' in fake.captured_prompt
    assert "只修改实现要求所必需的字段" in fake.captured_prompt
    assert "保持原有总时长" in fake.captured_prompt


def test_v6_ppt_binds_exact_plan_and_script_revisions_and_becomes_stale(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    lesson = repository.save_plan_revision(
        "course-1",
        "L1-1",
        standard_lesson_plan(),
        source_outline_revision_id="outline-v1",
    )
    source_revision = lesson["working_revision_id"]
    repository.confirm_plan_revision("course-1", "L1-1", source_revision)
    lesson = repository.save_script_revision(
        "course-1",
        "L1-1",
        [{
            "section_node_id": "L2-1-1",
            "title": "1.1",
            "content": "这是已经确认的讲稿正文。",
        }],
        source_lesson_plan_revision_id=source_revision,
        generation_source="teacher_edit",
    )
    script_revision = lesson["working_script_revision_id"]
    repository.confirm_script_revision(
        "course-1", "L1-1", script_revision
    )
    asset = repository.bind_v6_ppt_revision(
        "course-1",
        "L1-1",
        source_lesson_plan_revision_id=source_revision,
        source_script_revision_id=script_revision,
        synthetic_course_id="teacher-lesson-1",
        representation_id="rep-1",
        spec_id="spec-1",
        candidate_status="ready",
    )
    assert asset["source_state"] == "current"
    assert asset["v6_revisions"][0]["source_lesson_plan_revision_id"] == source_revision
    assert asset["v6_revisions"][0]["source_script_revision_id"] == script_revision

    first_binding_id = asset["working_v6_revision_id"]
    updated_asset = repository.bind_v6_ppt_revision(
        "course-1",
        "L1-1",
        source_lesson_plan_revision_id=source_revision,
        source_script_revision_id=script_revision,
        synthetic_course_id="teacher-lesson-1",
        representation_id="rep-1",
        spec_id="spec-2",
        candidate_status="ready",
    )
    restored_asset = repository.restore_v6_ppt_revision(
        "course-1",
        "L1-1",
        first_binding_id,
        expected_working_revision_id=updated_asset["working_v6_revision_id"],
    )
    assert restored_asset["working_v6_revision_id"] == first_binding_id
    assert restored_asset["source_lesson_plan_revision_id"] == source_revision
    assert len(restored_asset["v6_revisions"]) == 2

    repository.save_plan_revision(
        "course-1",
        "L1-1",
        {"lesson_title": "第一讲 v2", "sections": [{"node_id": "L2-1-1", "title": "1.1"}]},
        source_outline_revision_id="outline-v1",
        actor="teacher",
    )
    stale = repository.lesson("course-1", "L1-1")["ppt_assets"][0]
    assert stale["source_state"] == "stale"
    assert stale["v6_revisions"][0]["representation_id"] == "rep-1"


def test_teacher_lesson_v6_source_is_synthetic_and_covers_only_one_lesson():
    source = course_data()
    source["nodes"][1]["node_content"] = "这是一段已确认的一手讲稿正文。"
    source["nodes"][2]["node_content"] = "这是第二小节已确认的一手讲稿正文。"
    source_before = str(source)
    revision = {
        "revision_id": "plan-v1",
        "plan": {
            "revision_id": "plan-v1",
            "sections": [
                {
                    "node_id": "L2-1-1",
                    "learning_objective": "理解第一节",
                    "key_points": ["概念一"],
                    "teaching_modules": [{
                        "module_id": "core_explanation",
                        "teaching_purpose": "讲清概念一",
                        "knowledge_names": ["概念一"],
                    }],
                },
                {
                    "node_id": "L2-1-2",
                    "learning_objective": "理解第二节",
                    "key_points": ["概念二"],
                    "teaching_modules": [{
                        "module_id": "learner_action",
                        "teaching_purpose": "完成概念二练习",
                        "knowledge_names": ["概念二"],
                    }],
                },
            ],
        },
    }
    document, view, synthetic_id = teacher_lesson_v6_source(
        source,
        lesson_unit_id="L1-1",
        plan_revision=revision,
        script_revision={
            "revision_id": "legacy-script-v1",
            "sections": [
                {
                    "section_node_id": "L2-1-1",
                    "blocks": [{
                        "block_id": "legacy-1",
                        "module_id": "legacy_script",
                        "role": "concept",
                        "title": "讲稿正文",
                        "content": "这是一段已确认的一手讲稿正文。",
                    }],
                },
                {
                    "section_node_id": "L2-1-2",
                    "blocks": [{
                        "block_id": "legacy-2",
                        "module_id": "legacy_script",
                        "role": "concept",
                        "title": "讲稿正文",
                        "content": "这是第二小节已确认的一手讲稿正文。",
                    }],
                },
            ],
        },
    )
    graph = compile_course_presentation_graph(
        document,
        teaching_plan=view["course_teaching_plan"],
    )
    assert synthetic_id.startswith("teacher-lesson-")
    assert synthetic_id != source["course_id"]
    assert {section.section_id for section in document.sections} == {"L1-1", "L2-1-1", "L2-1-2"}
    assert {block.section_id for block in document.blocks} == {"L2-1-1", "L2-1-2"}
    assert graph.primary_block_coverage == 1.0
    assert graph.diagnostics == []
    assert view["teacher_lesson_source"]["script_revision_id"] == "legacy-script-v1"
    assert view["nodes"][1]["content_blocks"][0]["content"] == "这是一段已确认的一手讲稿正文。"
    assert str(source) == source_before


def test_teacher_preview_projects_current_script_without_confirmation_gate():
    plan = standard_lesson_plan()
    preview = {
        "course_id": "course-1",
        "course_name": "测试课程",
        "nodes": [
            {
                "node_id": "L1-1",
                "parent_node_id": "root",
                "node_level": 1,
                "node_name": "第一讲",
                "node_content": "",
            },
            {
                "node_id": "L2-1-1",
                "parent_node_id": "L1-1",
                "node_level": 2,
                "node_name": "1.1 核心概念",
                "node_content": "",
            },
        ],
    }
    quality = {
        "schema_version": SCRIPT_QUALITY_VERSION,
        "pipeline_version": SCRIPT_PIPELINE_VERSION,
        "passed": True,
        "publication_eligible": True,
    }
    authoring_state = {
        "course_id": "course-1",
        "outline_revision_id": "outline-v1",
        "lessons": {
            "L1-1": {
                "source_state": "current",
                "working_revision_id": "plan-v1",
                "confirmed_revision_id": "plan-v1",
                "revisions": [{"revision_id": "plan-v1", "plan": plan}],
                "working_script_revision_id": "script-v8",
                "script_confirmation": {
                    "confirmed_revision_id": "script-v8",
                    "source_lesson_plan_revision_id": "plan-v1",
                    "source_state": "current",
                },
                "script_revisions": [{
                    "revision_id": "script-v8",
                    "source_lesson_plan_revision_id": "plan-v1",
                    "publication_eligible": True,
                    "quality_report": quality,
                    "sections": [{
                        "section_node_id": "L2-1-1",
                        "title": "1.1 核心概念",
                        "content": "这是一段教师可以直接讲授的正式内容。",
                        "blocks": [{
                            "block_id": "block-1",
                            "module_id": "core_explanation",
                            "role": "concept",
                            "title": "核心教学",
                            "content": "这是一段教师可以直接讲授的正式内容。",
                            "planned_minutes": 8,
                        }],
                    }],
                }],
            },
        },
    }

    projected = project_confirmed_teacher_scripts(preview, authoring_state)

    section = next(item for item in projected["nodes"] if item["node_id"] == "L2-1-1")
    assert projected["projection"] == "teacher_lesson_authoring"
    assert projected["teacher_lesson_projection"]["covered_lesson_unit_ids"] == ["L1-1"]
    assert projected["teacher_lesson_projection"]["covered_section_count"] == 1
    assert section["node_content"] == "## 核心教学\n\n这是一段教师可以直接讲授的正式内容。"
    assert section["content_blocks"][0]["block_id"] == "block-1"
    assert section["generation_status"] == "completed"


def test_v6_ppt_manuscript_requires_matching_confirmation_before_formal_export(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    lesson = repository.save_plan_revision(
        "course-1",
        "L1-1",
        standard_lesson_plan(),
        source_outline_revision_id="outline-v1",
        quality_report=validate_teacher_lesson_plan(standard_lesson_plan()),
    )
    plan_revision = lesson["working_revision_id"]
    repository.confirm_plan_revision("course-1", "L1-1", plan_revision)
    lesson = repository.save_script_revision(
        "course-1",
        "L1-1",
        [{"section_node_id": "L2-1-1", "title": "1.1", "content": "中性讲稿正文。"}],
        source_lesson_plan_revision_id=plan_revision,
    )
    script_revision = lesson["working_script_revision_id"]
    repository.confirm_script_revision("course-1", "L1-1", script_revision)
    asset = repository.bind_v6_ppt_revision(
        "course-1",
        "L1-1",
        source_lesson_plan_revision_id=plan_revision,
        source_script_revision_id=script_revision,
        synthetic_course_id="teacher-lesson-1",
        representation_id="representation-1",
        spec_id="spec-1",
        candidate_status="v6_ready",
        ppt_manuscript_revision="pptman-1",
        ppt_manuscript_status="draft",
    )
    assert asset["ppt_manuscript_status"] == "draft"

    with pytest.raises(TeacherLessonAuthoringError) as conflict:
        repository.confirm_v6_ppt_manuscript(
            "course-1",
            "L1-1",
            representation_id="representation-1",
            manuscript_revision="pptman-old",
        )
    assert conflict.value.code == "lesson_ppt_manuscript_revision_conflict"

    confirmed = repository.confirm_v6_ppt_manuscript(
        "course-1",
        "L1-1",
        representation_id="representation-1",
        manuscript_revision="pptman-1",
    )
    assert confirmed["ppt_manuscript_status"] == "confirmed"
    assert confirmed["v6_revisions"][0]["ppt_manuscript_status"] == "confirmed"


def test_v6_ppt_manuscript_quality_gate_blocks_confirmation():
    with pytest.raises(TeacherLessonAuthoringError) as blocked:
        teacher_lesson_router._assert_ppt_manuscript_confirmable({
            "quality_status": "blocked",
            "quality_issues": ["page-4: 标题仍是原始 LaTeX 表达式"],
        })

    assert blocked.value.code == "lesson_ppt_manuscript_quality_blocked"
    assert blocked.value.details["quality_issues"] == [
        "page-4: 标题仍是原始 LaTeX 表达式"
    ]

    teacher_lesson_router._assert_ppt_manuscript_confirmable({
        "quality_status": "passed",
        "quality_issues": [],
    })


def test_ppt_manuscript_state_separates_original_branch_and_stale_materials():
    original = teacher_lesson_router._ppt_manuscript_state_payload(
        None, generation_branch="original_ppt_review"
    )
    assert original["generation_branch"] == "original_ppt_review"
    assert original["can_generate_ppt"] is False

    stale = teacher_lesson_router._ppt_manuscript_state_payload(
        {
            "revision": "pptman-1",
            "status": "confirmed",
            "source_state": "current",
            "source_material_revision": "pptrefs-old",
            "manuscript": {"quality_status": "passed"},
        },
        generation_branch="manuscript_first",
        current_material_revision="pptrefs-new",
    )
    assert stale["source_state"] == "stale"
    assert stale["confirmable"] is False
    assert stale["can_generate_ppt"] is False


def test_independent_ppt_manuscript_must_be_confirmed_before_binding_deck(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    lesson = repository.save_plan_revision(
        "course-1",
        "L1-1",
        standard_lesson_plan(),
        source_outline_revision_id="outline-v1",
        quality_report=validate_teacher_lesson_plan(standard_lesson_plan()),
    )
    plan_revision = lesson["working_revision_id"]
    repository.confirm_plan_revision("course-1", "L1-1", plan_revision)
    lesson = repository.save_script_revision(
        "course-1",
        "L1-1",
        [{"section_node_id": "L2-1-1", "title": "1.1", "content": "第一版讲稿。"}],
        source_lesson_plan_revision_id=plan_revision,
    )
    script_revision = lesson["working_script_revision_id"]
    repository.confirm_script_revision(
        "course-1", "L1-1", script_revision
    )
    state = repository.save_v6_ppt_manuscript(
        "course-1",
        "L1-1",
        source_lesson_plan_revision_id=plan_revision,
        source_script_revision_id=script_revision,
        source_material_revision="pptrefs-1",
        task_id="task-manuscript-1",
        mode="teaching",
        theme="qizhi-classroom",
        template_id="pptp-demo",
        template_version="3",
        template_digest="tmpl-demo",
        template_pack_id="pptp-demo",
        manuscript={
            "schema_version": "ppt_manuscript_v1",
            "manuscript_revision": "pptman-1",
            "quality_status": "passed",
            "pages": [{"page_id": "page-1", "title": "第一页"}],
        },
    )
    assert state["status"] == "draft"
    assert state["template_id"] == "pptp-demo"
    assert state["template_version"] == "3"
    assert state["template_digest"] == "tmpl-demo"
    assert state["template_pack_id"] == "pptp-demo"
    assert repository.lesson("course-1", "L1-1")["ppt_assets"] == []

    with pytest.raises(TeacherLessonAuthoringError) as blocked:
        repository.bind_v6_ppt_manuscript_result(
            "course-1",
            "L1-1",
            manuscript_revision="pptman-1",
            representation_id="representation-1",
        )
    assert blocked.value.code == "lesson_ppt_manuscript_not_confirmed"

    confirmed = repository.confirm_v6_ppt_manuscript_draft(
        "course-1", "L1-1", manuscript_revision="pptman-1"
    )
    assert confirmed["status"] == "confirmed"
    bound = repository.bind_v6_ppt_manuscript_result(
        "course-1",
        "L1-1",
        manuscript_revision="pptman-1",
        representation_id="representation-1",
    )
    assert bound["generated_representation_id"] == "representation-1"

    repository.save_script_revision(
        "course-1",
        "L1-1",
        [{"section_node_id": "L2-1-1", "title": "1.1", "content": "第二版讲稿。"}],
        source_lesson_plan_revision_id=plan_revision,
    )
    assert (
        repository.current_v6_ppt_manuscript("course-1", "L1-1")[
            "source_state"
        ]
        == "stale"
    )


def test_script_confirmation_versions_the_body_and_stales_bound_v6_ppt(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    quality = validate_teacher_lesson_plan(standard_lesson_plan())
    lesson = repository.save_plan_revision(
        "course-1",
        "L1-1",
        standard_lesson_plan(),
        source_outline_revision_id="outline-v1",
        quality_report=quality,
    )
    plan_revision = lesson["working_revision_id"]
    repository.confirm_plan_revision("course-1", "L1-1", plan_revision)

    source = course_data()
    source["nodes"][1]["node_content"] = "第一节正式讲稿"
    source["nodes"][2]["node_content"] = "第二节正式讲稿"
    first_script_revision = teacher_lesson_script_revision(source, "L1-1")
    repository.save_script_revision(
        "course-1",
        "L1-1",
        [
            {"section_node_id": "L2-1-1", "title": "1.1", "content": "第一节正式讲稿"},
            {"section_node_id": "L2-1-2", "title": "1.2", "content": "第二节正式讲稿"},
        ],
        source_lesson_plan_revision_id=plan_revision,
    )
    repository.confirm_script_revision("course-1", "L1-1", first_script_revision)
    asset = repository.bind_v6_ppt_revision(
        "course-1",
        "L1-1",
        source_lesson_plan_revision_id=plan_revision,
        source_script_revision_id=first_script_revision,
        synthetic_course_id="teacher-lesson-1",
        representation_id="representation-1",
        spec_id="spec-1",
        candidate_status="v6_ready",
    )
    assert asset["source_state"] == "current"
    assert asset["source_script_revision_id"] == first_script_revision

    source["nodes"][1]["node_content"] = "第一节修改后的正式讲稿"
    second_script_revision = teacher_lesson_script_revision(source, "L1-1")
    assert second_script_revision != first_script_revision
    repository.save_script_revision(
        "course-1",
        "L1-1",
        [
            {"section_node_id": "L2-1-1", "title": "1.1", "content": "第一节修改后的正式讲稿"},
            {"section_node_id": "L2-1-2", "title": "1.2", "content": "第二节正式讲稿"},
        ],
        source_lesson_plan_revision_id=plan_revision,
        generation_source="teacher_edit",
    )
    repository.confirm_script_revision("course-1", "L1-1", second_script_revision)
    stale = repository.lesson("course-1", "L1-1")["ppt_assets"][0]
    assert stale["source_state"] == "stale"


def test_current_complete_script_unlocks_the_v6_ppt_api_without_confirmation(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    source = course_data()
    source["blueprint_revision_id"] = "outline-v1"
    source["nodes"][1]["node_content"] = "第一节正式讲稿"
    source["nodes"][2]["node_content"] = "第二节正式讲稿"
    plan = standard_lesson_plan()
    second_section = deepcopy(plan["sections"][0])
    second_section["node_id"] = "L2-1-2"
    plan["sections"].append(second_section)
    lesson = repository.save_plan_revision(
        "course-1",
        "L1-1",
        plan,
        source_outline_revision_id="outline-v1",
        quality_report=validate_teacher_lesson_plan(plan),
    )
    class FakeStorage:
        @staticmethod
        def load_course(course_id):
            assert course_id == "course-1"
            return source

    class FakeTaskManager:
        storage = FakeStorage()

        @staticmethod
        def get_generation_workspace_course(_course_id):
            return None

        @staticmethod
        def get_generation_preview(_course_id):
            return None

    app = FastAPI()
    app.include_router(teacher_lesson_router.router, prefix="/api")
    app.dependency_overrides[require_task_manager] = lambda: FakeTaskManager()
    app.dependency_overrides[get_teacher_lesson_authoring_repository] = lambda: repository
    saved_script = repository.save_script_revision(
        "course-1",
        "L1-1",
        [
            {"section_node_id": "L2-1-1", "title": "1.1", "content": "第一节正式讲稿，包含完整定义、成立条件与适用边界。"},
            {"section_node_id": "L2-1-2", "title": "1.2", "content": "第二节正式讲稿，包含推理过程、示例结果与核对方法。"},
        ],
        source_lesson_plan_revision_id=lesson["working_revision_id"],
        generation_source="teacher_edit",
    )
    script_revision = saved_script["working_script_revision_id"]

    with TestClient(app) as client:
        source_response = client.get(
            "/api/teacher/courses/course-1/lessons/L1-1/ppt-v6/source"
        )
        assert source_response.status_code == 200
        assert source_response.json()["document"]["document_revision"]
        assert repository.lesson("course-1", "L1-1")["script_confirmation"] == {}

        legacy_generation = client.post(
            "/api/teacher/courses/course-1/lessons/L1-1/ppt/generate",
            json={"source_revision_id": lesson["working_revision_id"]},
        )
        assert legacy_generation.status_code == 404


def test_script_generation_edit_candidate_and_confirmation_share_one_asset_chain(tmp_path, monkeypatch):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    source = {**course_data(), "blueprint_revision_id": "outline-v1"}
    for section in (source["nodes"][1], source["nodes"][2]):
        section["module_plan"] = [{
            "module_id": "core_explanation",
            "label": "核心教学",
            "block_role": "concept",
            "required": True,
            "lesson_archetype_id": "general_concept_building",
            "lesson_archetype_label": "概念建构",
        }]
    lesson_plan = standard_lesson_plan()
    second_section = json.loads(json.dumps(lesson_plan["sections"][0]))
    second_section["node_id"] = "L2-1-2"
    lesson_plan["sections"].append(second_section)
    lesson = repository.save_plan_revision(
        "course-1",
        "L1-1",
        lesson_plan,
        source_outline_revision_id="outline-v1",
        quality_report=validate_teacher_lesson_plan(
            lesson_plan,
            expected_section_ids=["L2-1-1", "L2-1-2"],
        ),
    )
    plan_revision = lesson["working_revision_id"]
    repository.confirm_plan_revision("course-1", "L1-1", plan_revision)

    class FakeCourseService:
        registered = False
        script_calls = []
        rewrite_calls = []

        @classmethod
        def register_course_generation_metadata(cls, course_id, course):
            assert course_id == "course-1"
            assert course["nodes"]
            cls.registered = True

        @staticmethod
        async def generate_teacher_script_section(**kwargs):
            FakeCourseService.script_calls.append(kwargs)
            contract = compile_teacher_script_module_contract(
                kwargs["outline_section"], kwargs["confirmed_plan_section"]
            )
            return compile_teacher_script_section(
                "\n\n".join(
                    f"## {module['title']}\n\n"
                    f"这是第 {index} 个严格遵循已确认教案的正式讲义块，"
                    "定义、成立条件、适用边界与课堂核对方法均已说明。"
                    for index, module in enumerate(contract["modules"], start=1)
                ),
                contract,
            )

        @staticmethod
        async def rewrite_selection(**kwargs):
            FakeCourseService.rewrite_calls.append(kwargs)
            assert "增加一个真实案例" in kwargs["user_requirement"]
            return {"replacement_text": "## 核心教学\n\nAI 候选讲稿，已增加真实案例并保持原教学模块。"}

    class FakeStorage:
        @staticmethod
        def load_course(course_id):
            assert course_id == "course-1"
            return source

    class FakeTaskManager:
        storage = FakeStorage()
        course_service = FakeCourseService()

        @staticmethod
        def get_generation_workspace_course(_course_id):
            return None

        @staticmethod
        def get_generation_preview(_course_id):
            return None

    app = FastAPI()
    app.include_router(teacher_lesson_router.router, prefix="/api")
    app.dependency_overrides[require_task_manager] = lambda: FakeTaskManager()
    app.dependency_overrides[get_teacher_lesson_authoring_repository] = lambda: repository
    monkeypatch.setattr(
        teacher_lesson_router,
        "_course_material_evidence",
        lambda _course_id, _actor, material_ids: (
            material_ids,
            [{"asset_id": "material-1", "unit_id": "evidence-1", "text": "资料中的可靠案例"}],
        ),
    )

    with TestClient(app) as client:
        generated = client.post(
            "/api/teacher/courses/course-1/lessons/L1-1/script/generate",
            json={
                "request_id": "script-first",
                "requirements": "增加案例",
                "material_asset_ids": ["material-1"],
            },
            headers={"X-User-Id": "teacher-1"},
        )
        assert generated.status_code == 202
        job_id = generated.json()["job"]["id"]
        for _ in range(100):
            job = client.get(
                f"/api/teacher/courses/course-1/lesson-jobs/{job_id}"
            ).json()["job"]
            if job["status"] in {"completed", "completed_with_warnings", "failed"}:
                break
            time.sleep(0.01)
        assert job["status"] == "completed"
        assert job["completed_blocks"] == job["total_blocks"] == 2
        view = client.get(
            "/api/teacher/courses/course-1/lesson-authoring"
        ).json()
        first_script = next(
            item for item in view["lessons"]
            if item["lesson_unit_id"] == "L1-1"
        )["script"]
        assert first_script["ready"] is True
        assert len(first_script["sections"]) == 2
        first_revision = first_script["current_revision_id"]

        candidate = client.post(
            "/api/teacher/courses/course-1/lessons/L1-1/script/rewrite-candidate",
            json={
                "base_revision_id": first_revision,
                "section_node_id": "L2-1-1",
                "instruction": "增加一个真实案例",
            },
            headers={"X-User-Id": "teacher-1"},
        )
        assert candidate.status_code == 200
        assert repository.lesson("course-1", "L1-1")["working_script_revision_id"] == first_revision

        edited_sections = first_script["sections"]
        edited_sections[0]["content"] = candidate.json()["candidate"]["replacement_text"]
        edited_sections[0].pop("blocks", None)
        saved = client.put(
            "/api/teacher/courses/course-1/lessons/L1-1/script/draft",
            json={"base_revision_id": first_revision, "sections": edited_sections},
            headers={"X-User-Id": "teacher-1"},
        )
        assert saved.status_code == 200
        second_revision = saved.json()["lesson"]["script"]["current_revision_id"]
        assert second_revision != first_revision
        stored_second_revision = next(
            item for item in repository.lesson("course-1", "L1-1")["script_revisions"]
            if item["revision_id"] == second_revision
        )
        assert stored_second_revision["requirements"] == "增加案例"
        assert stored_second_revision["material_asset_ids"] == ["material-1"]
        assert stored_second_revision["publication_eligible"] is True
        assert stored_second_revision["sections"][0]["blocks"][0]["generation_source"] == "teacher_edit"
        assert stored_second_revision["sections"][1]["blocks"][0]["generation_source"] == "model"

        confirmed = client.post(
            "/api/teacher/courses/course-1/lessons/L1-1/script/confirm",
            json={"revision_id": second_revision},
        )
        assert confirmed.status_code == 200
        assert confirmed.json()["lesson"]["script"]["confirmed"] is True

        ppt_source = client.get(
            "/api/teacher/courses/course-1/lessons/L1-1/ppt-v6/source"
        )
        assert ppt_source.status_code == 200
        payload = json.dumps(ppt_source.json(), ensure_ascii=False)
        assert "AI 候选讲稿" in payload

        restored = client.post(
            f"/api/teacher/courses/course-1/lessons/L1-1/script/revisions/{first_revision}/restore",
            json={"expected_current_revision_id": second_revision},
            headers={"X-User-Id": "teacher-1"},
        )
        assert restored.status_code == 200
        restored_script = restored.json()["lesson"]["script"]
        assert restored_script["current_revision_id"] not in {first_revision, second_revision}
        assert restored_script["revisions"][0]["restored_from_revision_id"] == first_revision

    assert FakeCourseService.registered is True
    assert len(FakeCourseService.script_calls) == 1
    assert len(
        FakeCourseService.script_calls[0]["confirmed_plan_section"]["teaching_modules"]
    ) == 2
    assert (
        FakeCourseService.script_calls[0]["lesson_context"]
        ["script_shard_context"]["budget_mode"]
        == "single_request"
    )
    assert all(call["requirements"] == "增加案例" for call in FakeCourseService.script_calls)
    generation_context = FakeCourseService.script_calls[0]["lesson_context"]
    assert generation_context["selected_material_evidence"][0]["text"] == "资料中的可靠案例"
    rewrite_context = json.loads(FakeCourseService.rewrite_calls[0]["course_context"])
    assert rewrite_context["confirmed_lesson_plan"]["node_id"] == "L2-1-1"
    assert rewrite_context["teacher_requirements"] == "增加案例"


def test_teacher_can_save_first_script_draft_without_model_revision(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    source = {**course_data(), "blueprint_revision_id": "outline-v1"}
    lesson_plan = standard_lesson_plan()
    second_section = json.loads(json.dumps(lesson_plan["sections"][0]))
    second_section["node_id"] = "L2-1-2"
    lesson_plan["sections"].append(second_section)
    for section in (source["nodes"][1], source["nodes"][2]):
        section["module_plan"] = [{
            "module_id": "core_explanation",
            "label": "核心教学",
            "block_role": "concept",
            "required": True,
        }]
    lesson = repository.save_plan_revision(
        "course-1",
        "L1-1",
        lesson_plan,
        source_outline_revision_id="outline-v1",
        quality_report=validate_teacher_lesson_plan(
            lesson_plan,
            expected_section_ids=["L2-1-1", "L2-1-2"],
        ),
    )
    repository.confirm_plan_revision(
        "course-1", "L1-1", lesson["working_revision_id"]
    )

    class FakeStorage:
        @staticmethod
        def load_course(_course_id):
            return source

    class FakeTaskManager:
        storage = FakeStorage()

        @staticmethod
        def get_generation_workspace_course(_course_id):
            return None

        @staticmethod
        def get_generation_preview(_course_id):
            return None

    sections = []
    for outline_section, plan_section in zip(
        (source["nodes"][1], source["nodes"][2]),
        lesson_plan["sections"],
    ):
        contract = compile_teacher_script_module_contract(
            outline_section, plan_section
        )
        sections.append(compile_teacher_script_section(
            "## 核心教学\n\n教师手工补全的可编辑讲稿，包含概念、例子与核对结论。",
            contract,
        ))

    app = FastAPI()
    app.include_router(teacher_lesson_router.router, prefix="/api")
    app.dependency_overrides[require_task_manager] = lambda: FakeTaskManager()
    app.dependency_overrides[get_teacher_lesson_authoring_repository] = lambda: repository
    with TestClient(app) as client:
        saved = client.put(
            "/api/teacher/courses/course-1/lessons/L1-1/script/draft",
            json={"base_revision_id": "", "sections": sections},
            headers={"X-User-Id": "teacher-1"},
        )

    assert saved.status_code == 200
    stored = repository.lesson("course-1", "L1-1")
    assert stored["working_script_revision_id"]
    assert stored["script_revisions"][-1]["generation_source"] == "teacher_edit"
    assert all(
        block["generation_source"] == "teacher_edit"
        for section in stored["script_revisions"][-1]["sections"]
        for block in section["blocks"]
    )


def test_teacher_v6_route_uses_shared_ai_planner_factories():
    assert teacher_lesson_router.build_ai_base_story_planner_v6.__module__ == (
        "slide_ai_planning_v6"
    )
    assert teacher_lesson_router.build_ai_base_visual_planner_v2.__module__ == (
        "slide_ai_planning_v6"
    )
    assert not hasattr(teacher_lesson_router, "_teacher_v6_story_planner")
    assert not hasattr(teacher_lesson_router, "_teacher_v6_visual_planner")


def test_legacy_lightweight_ppt_write_chain_is_retired(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    service = CourseService()
    assert not hasattr(repository, "save_ppt_revision")
    assert not hasattr(repository, "save_ppt_ai_candidate")
    assert not hasattr(repository, "resolve_ppt_ai_candidate")
    assert not hasattr(service, "generate_teacher_lesson_ppt")
    assert not hasattr(service, "optimize_teacher_lesson_ppt")


def test_teacher_lesson_api_projects_sessions_from_canonical_course_document(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    source = {**course_data(), "course_name": "数据结构"}
    document = document_from_generation_draft(source)
    canonical = {
        "course_id": "course-1",
        "course_name": "数据结构",
        "course_schema_version": "course_document_v1",
        "blueprint_revision_id": "outline-v2",
        "course_document": document.model_dump(mode="json"),
    }

    class FakeStorage:
        @staticmethod
        def load_course(course_id):
            assert course_id == "course-1"
            return canonical

    class FakeTaskManager:
        storage = FakeStorage()

        @staticmethod
        def get_generation_workspace_course(_course_id):
            return None

        @staticmethod
        def get_generation_preview(_course_id):
            return None

    app = FastAPI()
    app.include_router(teacher_lesson_router.router, prefix="/api")
    app.dependency_overrides[require_task_manager] = lambda: FakeTaskManager()
    app.dependency_overrides[get_teacher_lesson_authoring_repository] = lambda: repository

    with TestClient(app) as client:
        view = client.get("/api/teacher/courses/course-1/lesson-authoring")

    assert view.status_code == 200
    assert view.json()["outline_revision_id"] == "outline-v2"
    assert [item["lesson_unit_id"] for item in view.json()["lessons"]] == ["L1-1", "L1-2"]
    assert [item["section_node_id"] for item in view.json()["lessons"][0]["sections"]] == ["L2-1-1", "L2-1-2"]


def test_teacher_lesson_view_expires_orphaned_jobs_before_frontend_recovery(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    job = repository.create_job(
        "course-1",
        "L1-1",
        request_id="request-orphaned-view",
        source_outline_revision_id="outline-v1",
        job_type="teacher_lesson_script_generation",
    )
    stored_path = tmp_path / "course-1.json"
    stored = json.loads(stored_path.read_text(encoding="utf-8"))
    stored["jobs"][job["id"]]["status"] = "running"
    stored["jobs"][job["id"]]["updated_at"] = "2020-01-01T00:00:00+00:00"
    stored_path.write_text(json.dumps(stored, ensure_ascii=False), encoding="utf-8")

    class FakeStorage:
        @staticmethod
        def load_course(_course_id):
            return course_data()

    class FakeTaskManager:
        storage = FakeStorage()

        @staticmethod
        def get_generation_workspace_course(_course_id):
            return None

        @staticmethod
        def get_generation_preview(_course_id):
            return None

    app = FastAPI()
    app.include_router(teacher_lesson_router.router, prefix="/api")
    app.dependency_overrides[require_task_manager] = lambda: FakeTaskManager()
    app.dependency_overrides[get_teacher_lesson_authoring_repository] = lambda: repository

    with TestClient(app) as client:
        response = client.get("/api/teacher/courses/course-1/lesson-authoring")

    returned = next(item for item in response.json()["jobs"] if item["id"] == job["id"])
    assert returned["status"] == "failed"
    assert returned["error"]["code"] == "lesson_script_generation_interrupted"


def test_teacher_lesson_view_does_not_rewrite_unchanged_authoring_state(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    repository.set_outline("course-1", "outline-v1")
    before = repository.view("course-1")["revision"]

    class FakeStorage:
        @staticmethod
        def load_course(_course_id):
            return course_data()

    class FakeTaskManager:
        storage = FakeStorage()

        @staticmethod
        def get_generation_workspace_course(_course_id):
            return None

        @staticmethod
        def get_generation_preview(_course_id):
            return None

    app = FastAPI()
    app.include_router(teacher_lesson_router.router, prefix="/api")
    app.dependency_overrides[require_task_manager] = lambda: FakeTaskManager()
    app.dependency_overrides[get_teacher_lesson_authoring_repository] = lambda: repository

    with TestClient(app) as client:
        first = client.get("/api/teacher/courses/course-1/lesson-authoring")
        second = client.get("/api/teacher/courses/course-1/lesson-authoring")

    assert first.status_code == 200
    assert second.status_code == 200
    assert repository.view("course-1")["revision"] == before


def test_teacher_lesson_view_treats_an_empty_teacher_draft_as_ready_for_setup(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    empty_draft = {
        "course_id": "course-1",
        "course_name": "物理",
        "nodes": [],
        "course_document": {
            "schema_version": "course_document_v1",
            "course_id": "course-1",
            "title": "物理",
            "sections": [],
            "blocks": [],
        },
    }

    class FakeStorage:
        @staticmethod
        def load_course(course_id):
            assert course_id == "course-1"
            return empty_draft

    class FakeTaskManager:
        storage = FakeStorage()

        @staticmethod
        def get_generation_workspace_course(_course_id):
            return None

        @staticmethod
        def get_generation_preview(_course_id):
            return None

    app = FastAPI()
    app.include_router(teacher_lesson_router.router, prefix="/api")
    app.dependency_overrides[require_task_manager] = lambda: FakeTaskManager()
    app.dependency_overrides[get_teacher_lesson_authoring_repository] = lambda: repository

    with TestClient(app) as client:
        view = client.get("/api/teacher/courses/course-1/lesson-authoring")

    assert view.status_code == 200
    assert view.json()["course_id"] == "course-1"
    assert view.json()["outline_revision_id"] == ""
    assert view.json()["lessons"] == []
    assert view.json()["jobs"] == []


def test_teacher_lesson_api_ignores_empty_persisted_shell_and_uses_workspace(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    empty_shell = {
        "course_id": "course-1",
        "nodes": [],
        "course_document": {
            "schema_version": "course_document_v1",
            "sections": [],
            "blocks": [],
        },
    }
    workspace = {**course_data(), "blueprint_revision_id": "outline-workspace-v1"}

    class FakeStorage:
        @staticmethod
        def load_course(course_id):
            assert course_id == "course-1"
            return empty_shell

    class FakeTaskManager:
        storage = FakeStorage()

        @staticmethod
        def get_generation_workspace_course(course_id):
            assert course_id == "course-1"
            return workspace

        @staticmethod
        def get_generation_preview(_course_id):
            return None

    app = FastAPI()
    app.include_router(teacher_lesson_router.router, prefix="/api")
    app.dependency_overrides[require_task_manager] = lambda: FakeTaskManager()
    app.dependency_overrides[get_teacher_lesson_authoring_repository] = lambda: repository

    with TestClient(app) as client:
        view = client.get("/api/teacher/courses/course-1/lesson-authoring")

    assert view.status_code == 200
    assert view.json()["outline_revision_id"] == "outline-workspace-v1"
    assert [item["lesson_unit_id"] for item in view.json()["lessons"]] == ["L1-1", "L1-2"]


def test_lesson_type_is_adjusted_from_outline_without_confirming_teaching_blocks(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    workspace = {**course_data(), "blueprint_revision_id": "outline-v1"}
    repository.save_arrangement_revision(
        "course-1",
        "L1-1",
        {"lesson_type": "theory", "blocks": []},
        source_outline_revision_id="outline-v1",
        confirm=False,
    )

    class FakeTaskManager:
        storage = None

        @staticmethod
        def get_generation_workspace_course(course_id):
            assert course_id == "course-1"
            return workspace

        @staticmethod
        def get_generation_preview(_course_id):
            return None

    app = FastAPI()
    app.include_router(teacher_lesson_router.router, prefix="/api")
    app.dependency_overrides[require_task_manager] = lambda: FakeTaskManager()
    app.dependency_overrides[get_teacher_lesson_authoring_repository] = lambda: repository

    with TestClient(app) as client:
        response = client.put(
            "/api/teacher/courses/course-1/lessons/L1-1/arrangement/type",
            json={"lesson_type": "project_workshop"},
        )

    assert response.status_code == 200
    arrangement = response.json()["lesson"]["arrangement"]
    assert arrangement["lesson_type"] == "project_workshop"
    assert arrangement["lesson_type_label"] == "项目工作坊"
    assert arrangement["status"] == "draft"
    assert arrangement["confirmed"] is False
    assert arrangement["blocks"] == []


def test_teacher_lesson_api_generates_only_requested_lesson(tmp_path):
    repository = TeacherLessonAuthoringRepository(tmp_path)

    class FakeCourseService:
        calls = []

        async def prepare_teacher_lesson_plan(
            self,
            *,
            course_data,
            lesson_unit_id,
            on_phase,
            source_evidence=None,
            lesson_arrangement=None,
            resume_checkpoint=None,
            on_checkpoint=None,
        ):
            self.calls.append(lesson_unit_id)
            assert source_evidence == []
            assert resume_checkpoint == {}
            assert callable(on_checkpoint)
            assert lesson_arrangement["lesson_type"] == "theory"
            assert lesson_arrangement["status"] == "draft"
            assert {item["section_node_id"] for item in lesson_arrangement["blocks"]} == {"L2-2-1"}
            assert course_data["requirements"] == "突出课堂讨论与案例分析"
            selected = next(
                item
                for item in course_data["course_plan"]["chapters"]
                if item["lecture_number"] == int(lesson_unit_id.rsplit("-", 1)[1])
            )
            assert selected["teacher_requirements"] == "突出课堂讨论与案例分析"
            await on_phase(
                "course_teaching_plan_batch", 60, "生成中", 60,
                {"stream_event": "reset", "stream_batch_id": "TP-B01", "stream_delta": ""},
            )
            await on_phase(
                "course_teaching_plan_batch", 60, "生成中", 60,
                {
                    "stream_event": "delta",
                    "stream_batch_id": "TP-B01",
                    "stream_delta": '{"sections":[{"learning_objective":"正在生成',
                },
            )
            await on_phase(
                "course_teaching_plan_batch", 60, "生成中", 60,
                {
                    "stream_event": "delta",
                    "stream_batch_id": "TP-B01",
                    "stream_delta": '专业教学目标"}]}',
                },
            )
            scope = lesson_scope(course_data, lesson_unit_id)
            plan = standard_lesson_plan()
            plan["sections"][0]["node_id"] = scope["sections"][0]["node_id"]
            return {
                "plan": plan,
                "warnings": [],
                "source_outline_revision_id": "outline-v1",
                "generation_source": "model",
            }

    class FakeTaskManager:
        storage = None
        course_service = FakeCourseService()

        @staticmethod
        def get_generation_workspace_course(course_id):
            assert course_id == "course-1"
            source = lecture_course_data()
            source["course_plan"]["reference_books"] = [
                "张三：《核心概念导论》，高等教育出版社，2025"
            ]
            return {**source, "blueprint_revision_id": "outline-v1"}

        @staticmethod
        def get_generation_preview(_course_id):
            return None

    app = FastAPI()
    app.include_router(teacher_lesson_router.router, prefix="/api")
    app.dependency_overrides[require_task_manager] = lambda: FakeTaskManager()
    app.dependency_overrides[get_teacher_lesson_authoring_repository] = lambda: repository

    with TestClient(app) as client:
        view = client.get("/api/teacher/courses/course-1/lesson-authoring")
        assert view.status_code == 200
        assert [item["lesson_unit_id"] for item in view.json()["lessons"]] == ["L1-1", "L1-2"]

        response = client.post(
            "/api/teacher/courses/course-1/lessons/L1-2/plan/generate",
            json={
                "request_id": "lesson-two",
                "requirements": "突出课堂讨论与案例分析",
            },
        )
        assert response.status_code == 202
        job_id = response.json()["job"]["id"]
        for _ in range(50):
            job = client.get(f"/api/teacher/courses/course-1/lesson-jobs/{job_id}").json()["job"]
            if job["status"] in {"completed", "completed_with_warnings", "failed"}:
                break
            time.sleep(0.01)
        with client.stream(
            "GET",
            f"/api/teacher/courses/course-1/lesson-jobs/{job_id}/stream",
        ) as stream_response:
            stream_payload = "".join(stream_response.iter_text())

    assert job["status"] == "completed"
    assert job["stream_complete"] is True
    assert job["stream_batches"] == {}
    assert "lesson_plan_complete" in stream_payload
    assert "正在生成专业教学目标" not in stream_payload
    assert FakeTaskManager.course_service.calls == ["L1-2"]
    assets = repository.view("course-1")["lessons"]
    assert set(assets) == {"L1-2"}
    generated_section = assets["L1-2"]["revisions"][0]["plan"]["sections"][0]
    assert generated_section["node_id"] == "L2-2-1"
    assert generated_section["teaching_modules"]
    assert all(
        item["arrangement_block_id"]
        for item in generated_section["teaching_modules"]
    )
    current_arrangement = repository.current_arrangement("course-1", "L1-2")
    assert sum(
        item["planned_minutes"]
        for item in generated_section["teaching_modules"]
    ) == sum(
        item["planned_minutes"]
        for item in current_arrangement["blocks"]
    )
    assert generated_section["teacher_activities"]
    generated_revision = assets["L1-2"]["revisions"][0]
    assert generated_revision["status"] == "draft"
    assert generated_revision["quality_report"]["passed"] is True


def test_generate_all_lesson_plans_queues_lecture_v1_lessons(
    tmp_path,
    monkeypatch,
):
    repository = TeacherLessonAuthoringRepository(tmp_path)

    class FakeTaskManager:
        storage = None

        @staticmethod
        def get_generation_workspace_course(course_id):
            assert course_id == "course-1"
            return {
                **lecture_course_data(),
                "blueprint_revision_id": "outline-v1",
            }

        @staticmethod
        def get_generation_preview(_course_id):
            return None

    requested_children = []

    async def fake_generate_lesson_plan(
        course_id,
        lesson_unit_id,
        body,
        _request,
        _tm,
        child_repository,
    ):
        requested_children.append((lesson_unit_id, body))
        job = child_repository.create_job(
            course_id,
            lesson_unit_id,
            request_id=body.request_id,
            source_outline_revision_id="outline-v1",
        )
        return {"job": job}

    monkeypatch.setattr(
        teacher_lesson_router,
        "generate_lesson_plan",
        fake_generate_lesson_plan,
    )
    monkeypatch.setattr(
        teacher_lesson_router,
        "_lesson_plan_material_scope",
        lambda _course_id, _actor, lesson_unit_id: {
            "source_package_id": "",
            "source_asset_id": "",
            "material_asset_ids": [f"material-{lesson_unit_id}"],
        },
    )
    app = FastAPI()
    app.include_router(teacher_lesson_router.router, prefix="/api")
    app.dependency_overrides[require_task_manager] = lambda: FakeTaskManager()
    app.dependency_overrides[
        get_teacher_lesson_authoring_repository
    ] = lambda: repository

    with TestClient(app) as client:
        response = client.post(
            "/api/teacher/courses/course-1/lesson-plans/generate-all",
            json={"request_id": "all-plans", "requirements": "重视案例"},
        )

    assert response.status_code == 202
    payload = response.json()
    assert len(payload["parent_job"]["child_job_ids"]) == 2
    assert payload["parent_job"]["skipped_lesson_ids"] == []
    assert [item[0] for item in requested_children] == ["L1-1", "L1-2"]
    assert all(item[1].batch_size == 2 for item in requested_children)
    assert repository.current_arrangement("course-1", "L1-1")["blocks"]
    assert repository.current_arrangement("course-1", "L1-2")["blocks"]


def test_generate_all_lesson_plans_skips_lessons_without_teaching_structure(
    tmp_path,
    monkeypatch,
):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    source = lecture_course_data()
    source["course_plan"]["chapters"] = source["course_plan"]["chapters"][:1]

    class FakeTaskManager:
        storage = None

        @staticmethod
        def get_generation_workspace_course(course_id):
            assert course_id == "course-1"
            return {**source, "blueprint_revision_id": "outline-v1"}

        @staticmethod
        def get_generation_preview(_course_id):
            return None

    requested_children = []

    async def fake_generate_lesson_plan(
        course_id,
        lesson_unit_id,
        body,
        _request,
        _tm,
        child_repository,
    ):
        requested_children.append((lesson_unit_id, body))
        job = child_repository.create_job(
            course_id,
            lesson_unit_id,
            request_id=body.request_id,
            source_outline_revision_id="outline-v1",
        )
        return {"job": job}

    monkeypatch.setattr(
        teacher_lesson_router,
        "generate_lesson_plan",
        fake_generate_lesson_plan,
    )
    monkeypatch.setattr(
        teacher_lesson_router,
        "_lesson_plan_material_scope",
        lambda _course_id, _actor, _lesson_unit_id: {
            "source_package_id": "",
            "source_asset_id": "",
            "material_asset_ids": [],
        },
    )
    app = FastAPI()
    app.include_router(teacher_lesson_router.router, prefix="/api")
    app.dependency_overrides[require_task_manager] = lambda: FakeTaskManager()
    app.dependency_overrides[
        get_teacher_lesson_authoring_repository
    ] = lambda: repository

    with TestClient(app) as client:
        response = client.post(
            "/api/teacher/courses/course-1/lesson-plans/generate-all",
            json={"request_id": "ready-plans-only"},
        )

    assert response.status_code == 202
    payload = response.json()
    assert len(payload["parent_job"]["child_job_ids"]) == 1
    assert payload["skipped_lesson_ids"] == ["L1-2"]
    assert payload["skipped_lessons"] == [{
        "lesson_unit_id": "L1-2",
        "reason": "lesson_arrangement:blocks_empty",
    }]
    assert [item[0] for item in requested_children] == ["L1-1"]
    assert requested_children[0][1].batch_size == 1


def test_generate_all_lesson_scripts_queues_every_lesson_with_one_parent(
    tmp_path,
    monkeypatch,
):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    for lesson_unit_id, section_node_ids in (
        ("L1-1", ["L2-1-1", "L2-1-2"]),
        ("L1-2", ["L2-2-1"]),
    ):
        plan = standard_lesson_plan()
        template_section = plan["sections"][0]
        plan["sections"] = [
            {**deepcopy(template_section), "node_id": section_node_id}
            for section_node_id in section_node_ids
        ]
        repository.save_plan_revision(
            "course-1",
            lesson_unit_id,
            plan,
            source_outline_revision_id="outline-v1",
            quality_report=validate_teacher_lesson_plan(plan),
        )

    class FakeTaskManager:
        storage = None
        course_service = object()

        @staticmethod
        def get_generation_workspace_course(course_id):
            assert course_id == "course-1"
            return {**course_data(), "blueprint_revision_id": "outline-v1"}

        @staticmethod
        def get_generation_preview(_course_id):
            return None

    requested_children = []

    async def fake_generate_lesson_script(
        course_id,
        lesson_unit_id,
        body,
        _request,
        _tm,
        _repository,
    ):
        requested_children.append((lesson_unit_id, body))
        return {
            "job": {
                "id": f"job-{lesson_unit_id}",
                "course_id": course_id,
                "lesson_unit_id": lesson_unit_id,
                "lesson_id": lesson_unit_id,
                "status": "pending",
                "phase": "queued",
                "message": "已入队",
                "parent_job_id": body.batch_parent_job_id,
                "batch_position": body.batch_position,
                "batch_size": body.batch_size,
            }
        }

    monkeypatch.setattr(
        teacher_lesson_router,
        "generate_lesson_script",
        fake_generate_lesson_script,
    )
    monkeypatch.setattr(
        teacher_lesson_router,
        "_lesson_script_material_scope",
        lambda _course_id, _actor, lesson_unit_id: {
            "source_package_id": "",
            "source_asset_id": "",
            "material_asset_ids": [f"material-{lesson_unit_id}"],
        },
    )
    app = FastAPI()
    app.include_router(teacher_lesson_router.router, prefix="/api")
    app.dependency_overrides[require_task_manager] = lambda: FakeTaskManager()
    app.dependency_overrides[
        get_teacher_lesson_authoring_repository
    ] = lambda: repository

    with TestClient(app) as client:
        response = client.post(
            "/api/teacher/courses/course-1/lesson-scripts/generate-all",
            json={"request_id": "all-scripts", "requirements": "重视案例"},
        )

    assert response.status_code == 202
    payload = response.json()
    assert payload["child_job_ids"] == ["job-L1-1", "job-L1-2"]
    assert payload["skipped_lesson_ids"] == []
    assert [item[0] for item in requested_children] == ["L1-1", "L1-2"]
    parent_ids = {item[1].batch_parent_job_id for item in requested_children}
    assert len(parent_ids) == 1
    assert [item[1].batch_position for item in requested_children] == [1, 2]
    assert all(item[1].batch_size == 2 for item in requested_children)
    assert requested_children[0][1].material_asset_ids == ["material-L1-1"]


def test_generate_all_lesson_scripts_skips_lessons_without_ready_plan(
    tmp_path,
    monkeypatch,
):
    repository = TeacherLessonAuthoringRepository(tmp_path)
    plan = standard_lesson_plan()
    template_section = plan["sections"][0]
    plan["sections"] = [
        {**deepcopy(template_section), "node_id": section_node_id}
        for section_node_id in ("L2-1-1", "L2-1-2")
    ]
    repository.save_plan_revision(
        "course-1",
        "L1-1",
        plan,
        source_outline_revision_id="outline-v1",
    )

    class FakeTaskManager:
        storage = None
        course_service = object()

        @staticmethod
        def get_generation_workspace_course(course_id):
            assert course_id == "course-1"
            return {**course_data(), "blueprint_revision_id": "outline-v1"}

        @staticmethod
        def get_generation_preview(_course_id):
            return None

    requested_children = []

    async def fake_generate_lesson_script(
        course_id,
        lesson_unit_id,
        body,
        _request,
        _tm,
        _repository,
    ):
        requested_children.append((lesson_unit_id, body))
        return {
            "job": {
                "id": f"job-{lesson_unit_id}",
                "course_id": course_id,
                "lesson_unit_id": lesson_unit_id,
                "lesson_id": lesson_unit_id,
                "status": "pending",
                "phase": "queued",
                "message": "已入队",
                "parent_job_id": body.batch_parent_job_id,
                "batch_position": body.batch_position,
                "batch_size": body.batch_size,
            }
        }

    monkeypatch.setattr(
        teacher_lesson_router,
        "generate_lesson_script",
        fake_generate_lesson_script,
    )
    monkeypatch.setattr(
        teacher_lesson_router,
        "_lesson_script_material_scope",
        lambda _course_id, _actor, _lesson_unit_id: {
            "source_package_id": "",
            "source_asset_id": "",
            "material_asset_ids": [],
        },
    )
    app = FastAPI()
    app.include_router(teacher_lesson_router.router, prefix="/api")
    app.dependency_overrides[require_task_manager] = lambda: FakeTaskManager()
    app.dependency_overrides[
        get_teacher_lesson_authoring_repository
    ] = lambda: repository

    with TestClient(app) as client:
        response = client.post(
            "/api/teacher/courses/course-1/lesson-scripts/generate-all",
            json={"request_id": "ready-scripts-only"},
        )

    assert response.status_code == 202
    payload = response.json()
    assert payload["child_job_ids"] == ["job-L1-1"]
    assert payload["skipped_lesson_ids"] == ["L1-2"]
    assert payload["skipped_lessons"] == [{
        "lesson_unit_id": "L1-2",
        "reason": "revision_missing",
    }]
    assert [item[0] for item in requested_children] == ["L1-1"]
    assert requested_children[0][1].batch_size == 1
