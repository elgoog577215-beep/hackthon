from copy import deepcopy

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Pt

from content_blocks import content_fingerprint
from ppt_adopted_visuals import bind_adopted_assets, current_visual_catalog
from ppt_layout_execution import file_digest
from ppt_native_scene import audit_scene, render_scene
from ppt_teaching_content import PageTeachingV2
from slide_asset_repository import SlideAssetRepository
from .test_ppt_teaching_content import comparison_fixture, scene_for


def adopted_fixture(tmp_path):
    assets = SlideAssetRepository(tmp_path / 'assets')
    path = tmp_path / 'image.png'
    Image.new('RGB', (180, 90), '#3157c7').save(path)
    asset = assets.promote(assets.stage_image(path, course_id='course', source_fragment_ids=['b'], alt_text='执行方式示意', purpose='comparison'))
    content, notes = comparison_fixture()
    text = notes['b'][1]
    item = {'status': 'accepted', 'stale_reasons': [], 'representation_id': 'adopted-image', 'revision': 'representation-r1',
        'representation_type': 'image', 'source': {'block_id': 'b', 'script_revision_id': 'script-r1', 'block_content_fingerprint': content_fingerprint(text)},
        'content': {'visual_asset_manifest': [asset.model_dump(mode='json')]}}
    content['elements'][4].update(kind='image', asset_id=asset.asset_id, asset_digest=asset.sha256)
    return assets, asset, item, content, {'b': text}


def test_only_current_accepted_images_are_bound_to_their_source(tmp_path):
    assets, asset, item, value, sources = adopted_fixture(tmp_path)
    variants = [deepcopy(item) for _ in range(4)]
    variants[0]['status'] = 'candidate'
    variants[1]['source']['script_revision_id'] = 'old'
    variants[2]['source']['block_content_fingerprint'] = 'changed'
    variants[3]['source']['block_id'] = 'other'
    catalog = current_visual_catalog([*variants, item], course_id='course', script_revision_id='script-r1', sources=sources, asset_repository=assets)
    assert len(catalog) == 1
    content = bind_adopted_assets(PageTeachingV2.model_validate(value), catalog, {'b'})
    binding = content.adopted_assets[0]
    assert binding.asset_id == asset.asset_id and binding.sha256 == file_digest(assets.resolve(asset.asset_id))
    assert binding.representation_id == 'adopted-image' and binding.source_block_id == 'b'
    with pytest.raises(ValueError, match='teaching_asset_not_adopted'):
        bind_adopted_assets(content, catalog, {'other'})


def test_image_preview_identity_and_native_crop_are_audited(tmp_path):
    assets, asset, item, value, sources = adopted_fixture(tmp_path)
    catalog = current_visual_catalog([item], course_id='course', script_revision_id='script-r1', sources=sources, asset_repository=assets)
    content = bind_adopted_assets(PageTeachingV2.model_validate(value), catalog, {'b'})
    scene = scene_for(content.model_dump(mode='json'), 'compare-visual')
    image = next(o for o in scene.objects if o.kind == 'image')
    assert (image.asset_course_id, image.asset_representation_id) == ('course', 'adopted-image')
    deck = Presentation()
    deck.slide_width, deck.slide_height = Pt(960), Pt(540)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    render_scene(slide, scene, assets=assets)
    assert audit_scene(slide, scene)['passed']
    picture = next(s for s in slide.shapes if s.name == 'teaching:a-mode')
    picture.crop_right = 0.25
    with pytest.raises(ValueError, match='export_asset_geometry_mismatch'):
        audit_scene(slide, scene)
    tampered = deepcopy(item)
    tampered['content']['visual_asset_manifest'][0]['sha256'] = 'wrong'
    with pytest.raises(ValueError, match='teaching_adopted_asset_identity_mismatch'):
        current_visual_catalog([tampered], course_id='course', script_revision_id='script-r1', sources=sources, asset_repository=assets)


def test_adopted_diagram_copies_branch_edges_and_detects_semantic_edits(tmp_path):
    from ppt_adopted_visuals import lower_adopted_diagram, validate_adopted_diagram, diagram_semantics
    assets = SlideAssetRepository(tmp_path / 'assets')
    text = '观察结果为判断路径和继续记录提供依据。'
    spec = {'title': '观察任务', 'units': [{'unit_id': 'unit', 'section_id': 'section', 'title': '观察任务',
        'source_block_ids': ['b'], 'nodes': [{'node_id': k, 'label': label, 'kind': 'course_block', 'source_ref': 'b'}
            for k, label in [('a', '观察结果'), ('b', '判断路径'), ('c', '继续记录')]],
        'edges': [{'edge_id': f'a-{k}', 'source_node_id': 'a', 'target_node_id': k, 'relation': 'supports'} for k in ['b', 'c']]}]}
    item = {'status': 'accepted', 'representation_id': 'diagram', 'revision': 'r1', 'representation_type': 'diagram',
        'source': {'block_id': 'b', 'script_revision_id': 'script', 'block_content_fingerprint': content_fingerprint(text)}, 'content': spec}
    catalog = current_visual_catalog([item], course_id='course', script_revision_id='script', sources={'b': text}, asset_repository=assets)
    raw = lower_adopted_diagram({'adopted_diagram_id': 'diagram', 'diagram_unit_id': 'unit',
        'title': '观察任务', 'page_goal': '理解分支', 'teaching_note': '观察两条分支'},
        {'b': {'block_id': 'b', 'block_revision': 'source-r1', 'full_text': text}}, catalog)
    content = bind_adopted_assets(PageTeachingV2.model_validate(raw['teaching']), catalog, {'b'})
    assert [(r.source_id, r.target_id, r.kind) for r in content.expression.relations] == [('a', 'b', 'supports'), ('a', 'c', 'supports')]
    scene = scene_for(content.model_dump(mode='json'), 'concept-map')
    deck = Presentation()
    deck.slide_width, deck.slide_height = Pt(960), Pt(540)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    render_scene(slide, scene)
    assert audit_scene(slide, scene)['passed']
    content.expression.relations[1].source_id = 'b'
    with pytest.raises(ValueError, match='teaching_adopted_diagram_changed'):
        validate_adopted_diagram(content)
    content.adopted_diagram.semantic_digest = diagram_semantics(content)
    with pytest.raises(ValueError, match='teaching_adopted_diagram_changed'):
        bind_adopted_assets(content, catalog, {'b'})
