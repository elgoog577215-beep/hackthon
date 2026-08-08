from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches

from slide_deck import SlideSpec
from slide_deck_renderer import (
    _heading,
    _heading_excerpt,
    _render_claim_only,
    _render_classification_three,
    _render_code,
    _render_editorial_body,
    _render_practice_feedback,
    _render_process,
    _uses_visual_directed_renderer,
    _worked_example_labels,
    validate_theme,
)
from slide_deck_v5 import (
    DeckChapterV5,
    _assign_heading_modes_v5,
    _chapter_recap_slide,
    _consolidate_task_activity_pages_v5,
    _disambiguate_duplicate_titles_v5,
    _enrich_practice_feedback_slides_v5,
    _normalize_concept_definition_slide_v5,
    _promote_sparse_single_claim_v5,
    _restore_chapter_entry_mainlines_v5,
    _strip_instructional_scaffolding_v5,
    _structure_labeled_reasoning_pairs_v5,
    _structure_long_editorial_prose_v5,
    apply_page_contract_v5,
    compile_page_title_v5,
    finalize_v5_quality_report,
    split_mixed_intent_slides_v5,
    v5_contract_issues,
)


def test_standalone_micro_transition_is_folded_into_previous_page_metadata() -> None:
    slides = split_mixed_intent_slides_v5([
        {
            "unit_id": "practice",
            "position": 0,
            "scene_kind": "practice_feedback",
            "beat_role": "prompt",
            "title": "判断系统类型",
            "blocks": [{
                "block_id": "question",
                "type": "exercise",
                "content": "水壶盖子关闭时属于哪类系统？",
            }],
            "quality": {},
        },
        {
            "unit_id": "slide:v4:0003:transition",
            "position": 1,
            "scene_kind": "transition",
            "beat_role": "transition",
            "title": "下一节：热力学第一定律",
            "blocks": [{
                "block_id": "transition",
                "type": "statement",
                "content": "下一节将深入探讨热力学第一定律。",
            }],
            "quality": {},
        },
        {
            "unit_id": "next-concept",
            "position": 2,
            "scene_kind": "concept",
            "title": "热力学第一定律描述能量守恒",
            "blocks": [],
            "quality": {},
        },
    ])

    assert [slide["unit_id"] for slide in slides] == [
        "practice",
        "next-concept",
    ]
    assert slides[0]["quality"]["removed_redundant_transition"] is True
    assert slides[0]["quality"]["removed_transition_unit_ids"] == [
        "slide:v4:0003:transition"
    ]
    assert slides[0]["quality"]["next_topic"] == "热力学第一定律描述能量守恒"


def test_mixed_question_and_transition_drops_redundant_navigation_page() -> None:
    slides = split_mixed_intent_slides_v5([{
        "unit_id": "mixed-question-transition",
        "position": 5,
        "layout": "practice",
        "slide_purpose": "practice_feedback",
        "scene_kind": "practice_feedback",
        "beat_role": "prompt",
        "title": "水壶盖子没有打开，那么这个系统应该归类为什么类型",
        "key_message": "",
        "teaching_job": "用来源问题检查理解",
        "takeaway": "本节内容为后续学习打下基础",
        "narrative_role": "checkpoint",
        "composition": "exercise",
        "blocks": [
            {
                "block_id": "question",
                "type": "exercise",
                "content": "水壶盖子没有打开，这个系统属于哪种类型？",
            },
            {
                "block_id": "transition",
                "type": "statement",
                "content": (
                    "本节介绍了系统分类，并比较了孤立系统、封闭系统和开放系统在"
                    "物质交换、能量交换、边界条件及典型工程情境中的差异。"
                    "这些内容帮助我们建立后续分析所需要的对象边界和状态描述。"
                    "下一节将深入探讨热力学第一定律。"
                ),
            },
        ],
        "quality": {"requested_layout": "two-column"},
    }, {
        "unit_id": "actual-next-page",
        "position": 6,
        "layout": "concept",
        "slide_purpose": "concept",
        "scene_kind": "concept",
        "beat_role": "formal_explanation",
        "title": "1.2 状态变量与过程量",
        "key_message": "",
        "blocks": [{
            "block_id": "definition",
            "type": "rich_text",
            "content": "状态变量只依赖于系统当前状态。",
        }],
    }])

    assert len(slides) == 2
    assert slides[0]["quality"]["requested_layout"] == "question-prompt"
    assert [block["block_id"] for block in slides[0]["blocks"]] == ["question"]
    assert slides[0]["quality"]["removed_redundant_transition"] is True
    assert slides[0]["quality"]["next_topic"] == "状态变量与过程量"
    assert slides[1]["unit_id"] == "actual-next-page"
    assert all(
        slide["quality"].get("requested_layout") != "hero-claim"
        for slide in slides
    )


def test_question_and_transition_inside_one_block_are_split_at_sentence_level() -> None:
    slides = split_mixed_intent_slides_v5([{
        "unit_id": "single-mixed-block",
        "position": 7,
        "layout": "practice",
        "slide_purpose": "practice_feedback",
        "scene_kind": "practice_feedback",
        "beat_role": "prompt",
        "title": "判断系统类型并衔接下一节",
        "key_message": "",
        "blocks": [{
            "block_id": "mixed",
            "type": "exercise",
            "content": (
                "水壶盖子没有打开，这个系统属于哪种类型？"
                "本节介绍了系统分类。"
                "下一节将深入探讨热力学第一定律。"
            ),
        }],
        "quality": {"requested_layout": "two-column"},
    }])

    assert len(slides) == 1
    assert slides[0]["blocks"][0]["content"] == (
        "水壶盖子没有打开，这个系统属于哪种类型？"
    )
    assert slides[0]["quality"]["removed_redundant_transition"] is True


def test_derived_chapter_recap_compacts_long_claims_before_quality_gate() -> None:
    chapter = DeckChapterV5(
        chapter_id="chapter-density",
        agenda_id="agenda-density",
        position=0,
        eyebrow="第一章",
        title="系统与边界",
        driving_question="系统边界如何决定分析方式？",
        learning_objective="能够依据交换关系判断系统类型",
    )
    source_slides = [
        {
            "unit_id": f"source-{index}",
            "title": f"来源页 {index}",
            "takeaway": (
                f"判断{index}：当研究对象与环境发生物质和能量交换时，"
                "必须同时识别边界条件、交换方向、状态变量及过程约束，"
                "才能选择合适的热力学模型并解释实际现象，同时还要验证"
                "所选假设是否适用于稳态、瞬态和不可逆过程。"
            ),
            "blocks": [],
        }
        for index in range(1, 5)
    ]

    recap = apply_page_contract_v5(_chapter_recap_slide(chapter, source_slides))

    assert "body_density_overflow" not in {
        issue["code"] for issue in v5_contract_issues([recap])
    }
    assert all(len(item) <= 64 for item in recap["blocks"][0]["items"])
    assert not {
        "recap_item_incomplete",
        "recap_retrieval_prompt_missing",
    } & {issue["code"] for issue in v5_contract_issues([recap])}


def test_concept_definition_is_promoted_and_generic_label_is_removed() -> None:
    normalized = _normalize_concept_definition_slide_v5({
        "unit_id": "state-variable",
        "scene_kind": "concept",
        "title": "状态变量是指只依赖于系统当前状态的",
        "blocks": [
            {
                "block_id": "context",
                "type": "rich_text",
                "title": "核心概念与背景",
                "content": "在热力学中，宏观性质通常分为两类物理量。",
            },
            {
                "block_id": "definition-source",
                "type": "rich_text",
                "content": "状态变量是指只依赖于系统当前状态的物理量。",
            },
        ],
        "quality": {},
    })
    contracted = apply_page_contract_v5(normalized)

    assert contracted["blocks"][0]["title"] == "定义"
    assert contracted["blocks"][0]["content"] == (
        "状态变量是指只依赖于系统当前状态的物理量。"
    )
    assert contracted["blocks"][0]["metadata"]["semantic_role"] == "definition"
    assert all(
        block.get("title") != "核心概念与背景"
        for block in contracted["blocks"]
    )
    assert contracted["title"] == "状态变量只取决于系统当前状态"
    assert "concept_definition_missing" not in {
        issue["code"] for issue in v5_contract_issues([contracted])
    }


def test_long_foundation_claim_is_compressed_without_mid_word_truncation() -> None:
    title = compile_page_title_v5(
        explicit_title=(
            "热力学第零定律是热力学四大定律中最基础的一条，"
            "它是温度这一物理量的定义基础。"
        ),
        primary_claim="",
        body_text="",
        fallback_context="",
    )

    assert title == "热力学第零定律奠定温度定义基础"
    assert not title.endswith("最基")


def test_title_compiler_rejects_markdown_internal_production_labels() -> None:
    title = compile_page_title_v5(
        explicit_title="本节知识规范名称",
        primary_claim=(
            "**知识规范名称：MonoBehaviour 脚本命名规范与生命周期回调执行顺序**"
        ),
        body_text=(
            "MonoBehaviour 脚本命名规范与生命周期回调执行顺序。"
            "本节通过创建符合规范的脚本验证初始化时序。"
        ),
        fallback_context="脚本生命周期",
    )

    assert title
    assert "知识规范名称" not in title
    assert "生命周期" in title or "MonoBehaviour" in title


def test_title_compiler_prefers_source_heading_over_generic_teaching_job() -> None:
    title = compile_page_title_v5(
        explicit_title="说明结论如何从条件推出",
        primary_claim="1. Inspector 参数实时热更机制",
        body_text=(
            "Unity 的 Inspector 面板不仅是属性查看器，"
            "也是运行时参数控制台。"
        ),
        fallback_context="说明结论如何从条件推出",
    )

    assert title == "Inspector 参数实时热更机制"


def test_page_contract_preserves_visible_sequence_for_continuation_titles() -> None:
    slide = apply_page_contract_v5({
        "unit_id": "slide:v5:continuation:002",
        "layout": "concept",
        "title": "Collider Is Trigger 属性与碰撞事件路由机制",
        "takeaway": "Collider Is Trigger 属性与碰撞事件路由机制",
        "blocks": [{
            "block_id": "conditions",
            "type": "statement",
            "title": "",
            "content": "至少有一个物体必须挂载非运动学的 Rigidbody 组件。",
            "items": [],
        }],
        "quality": {
            "requested_layout": "editorial-body",
            "continuation_of": "slide:v5:continuation:001",
            "continuation_index": 2,
            "continuation_total": 4,
            "title_character_budget": 18,
        },
    })

    assert slide["title"].startswith("属性与碰撞事件路由机制")
    assert slide["title"].endswith("（续2/4）")


def test_page_contract_keeps_the_only_source_claim_for_hero_quality_audit() -> None:
    slide = apply_page_contract_v5({
        "unit_id": "slide:v5:source-hero",
        "layout": "concept",
        "title": "建立本节核心概念与边界",
        "key_message": "第 1 章概念",
        "takeaway": "",
        "blocks": [{
            "block_id": "source-claim",
            "type": "statement",
            "title": "",
            "content": "这是第 1 章的课程正文与方法说明。",
            "items": [],
            "metadata": {"source_fragment_ids": ["fragment-1"]},
        }],
        "visuals": [],
        "quality": {"requested_layout": "editorial-body"},
    })

    assert slide["title"] == "这是第 1 章的课程正文与方法说明"
    assert slide["blocks"][0]["content"].startswith("这是第 1 章")
    assert slide["quality"]["resolved_layout"] == "hero-claim"
    assert slide["quality"]["suppress_redundant_body"] is True


def test_numbered_semantic_heading_recovers_an_incomplete_body_excerpt() -> None:
    title = compile_page_title_v5(
        explicit_title="卡诺循环（Carnot",
        primary_claim="3.2 卡诺定理与最大效率",
        body_text=(
            "在热力学中，卡诺循环（Carnot Cycle）是最早提出的一个"
            "理想化热机模型。"
        ),
        fallback_context="热力学第二定律",
        prefer_body_claim=True,
    )

    assert title == "卡诺定理与最大效率"


def test_incomplete_body_claim_does_not_replace_a_complete_semantic_title() -> None:
    title = compile_page_title_v5(
        explicit_title="内能的本质",
        primary_claim="内能的本质",
        body_text="从微观角度来看，内能 U 包括：",
        fallback_context="热力学第一定律",
        prefer_body_claim=True,
    )

    assert title == "内能的本质"


def test_complete_question_without_terminal_punctuation_is_not_blocked() -> None:
    title = compile_page_title_v5(
        explicit_title="从热力学角度看，这是由什么驱动的",
        primary_claim="思考与挑战",
        body_text="",
        fallback_context="溶液热力学",
    )
    issues = v5_contract_issues([{
        "unit_id": "practice-question",
        "title": title,
        "blocks": [],
        "quality": {},
    }])

    assert title == "从热力学角度看，这是由什么驱动的"
    assert "incomplete_title_claim" not in {
        issue["code"] for issue in issues
    }


def test_title_ending_in_a_dependent_conjunction_is_blocked() -> None:
    issues = v5_contract_issues([{
        "unit_id": "dependent-title",
        "title": "碰撞回调事件的封装与分层处理模式以及",
        "blocks": [],
        "quality": {},
    }])

    assert "incomplete_title_claim" in {
        issue["code"] for issue in issues
    }


def test_title_ending_in_a_single_character_after_a_connector_is_blocked() -> None:
    issues = v5_contract_issues([{
        "unit_id": "hard-cut-title",
        "title": "精准定位运行时性能瓶颈与逻",
        "blocks": [],
        "quality": {},
    }])

    assert "incomplete_title_claim" in {
        issue["code"] for issue in issues
    }


def test_bare_instructional_scaffold_is_not_a_publishable_title() -> None:
    issues = v5_contract_issues([{
        "unit_id": "scaffold-title",
        "title": "本节课的目标是",
        "blocks": [],
        "quality": {},
    }])

    assert "incomplete_title_claim" in {
        issue["code"] for issue in issues
    }


def test_recap_excludes_question_morphology_and_instructional_prompts() -> None:
    chapter = DeckChapterV5(
        chapter_id="chapter-declarative",
        agenda_id="agenda-declarative",
        position=0,
        eyebrow="第一章",
        title="系统分类",
        driving_question="如何判断系统类型？",
        learning_objective="能够根据交换关系判断系统类型。",
    )
    recap = _chapter_recap_slide(chapter, [
        {
            "unit_id": "question-source",
            "title": "判断系统类型",
            "takeaway": "水壶盖子关闭时应该归类为什么类型",
            "blocks": [],
        },
        {
            "unit_id": "instruction-source",
            "title": "理解检查",
            "takeaway": "考虑空调制冷的过程",
            "blocks": [],
        },
        {
            "unit_id": "declarative-source",
            "title": "封闭系统不交换物质",
            "takeaway": "封闭系统不与外界交换物质，但可以交换能量。",
            "blocks": [],
        },
    ])

    assert recap["blocks"][0]["items"] == [
        "封闭系统不与外界交换物质，但可以交换能量。"
    ]


def test_optional_visual_does_not_override_v5_practice_composition() -> None:
    unit = SlideSpec.model_validate({
        "unit_id": "practice-with-visual",
        "position": 0,
        "layout": "practice",
        "slide_purpose": "practice_feedback",
        "title": "判断系统类型",
        "blocks": [],
        "visuals": [{
            "visual_id": "optional-relation",
            "kind": "relational_diagram",
            "purpose": "structure",
            "alt_text": "可选关系图",
        }],
        "quality": {"resolved_layout": "practice-feedback"},
    })

    assert _uses_visual_directed_renderer(unit, "practice-feedback") is False
    assert _uses_visual_directed_renderer(unit, "figure-text") is True


def test_worked_example_labels_must_be_explicit_or_neutral() -> None:
    assert _worked_example_labels({}, 3) == ("步骤 1", "步骤 2", "步骤 3")
    assert _worked_example_labels(
        {"worked_step_labels": ["条件", "推演", "验证"]},
        3,
    ) == ("条件", "推演", "验证")


def test_prompt_only_practice_is_enriched_with_grounded_answer_evidence() -> None:
    slides = [
        {
            "unit_id": "concept-system-types",
            "position": 0,
            "chapter_id": "chapter-1",
            "knowledge_refs": ["system-types"],
            "layout": "concept",
            "scene_kind": "explanation",
            "title": "三类热力学系统",
            "blocks": [{
                "block_id": "definitions",
                "type": "bullets",
                "items": [
                    "孤立系统：不交换物质，也不交换能量。",
                    "封闭系统：不交换物质，但可以交换能量。",
                    "开放系统：既可以交换物质，也可以交换能量。",
                ],
            }],
            "quality": {"requested_layout": "classification-3"},
        },
        {
            "unit_id": "practice-system-types",
            "position": 1,
            "chapter_id": "chapter-1",
            "knowledge_refs": ["system-types"],
            "layout": "practice",
            "scene_kind": "practice_feedback",
            "beat_role": "prompt",
            "title": "判断水壶属于哪类系统",
            "blocks": [{
                "block_id": "question",
                "type": "exercise",
                "items": [
                    "水壶盖子没有打开时属于哪类系统？",
                    "盖子打开并有蒸汽逸出时呢？",
                ],
            }],
            "quality": {"requested_layout": "question-prompt"},
        },
    ]

    enriched = _enrich_practice_feedback_slides_v5(slides)
    practice = apply_page_contract_v5(enriched[1])

    assert practice["quality"]["requested_layout"] == "practice-feedback"
    assert practice["quality"]["resolved_layout"] == "practice-feedback"
    assert practice["quality"]["grounded_feedback_source_ids"] == [
        "concept-system-types"
    ]
    assert practice["blocks"][0]["items"] == [
        "水壶盖子没有打开时属于哪类系统？",
        "盖子打开并有蒸汽逸出时呢？",
    ]
    assert practice["blocks"][1]["title"] == "判断依据"
    assert practice["blocks"][1]["items"] == [
        "孤立系统：不交换物质，也不交换能量。",
        "封闭系统：不交换物质，但可以交换能量。",
    ]
    assert practice["blocks"][1]["metadata"]["direct_answer"] is False
    assert practice["quality"]["feedback_mode"] == "shared_evidence"
    assert practice["quality"]["feedback_pair_count"] == 0
    assert practice["quality"]["feedback_evidence_count"] == 2
    assert "practice_feedback_missing_answer" not in {
        issue["code"] for issue in v5_contract_issues([practice])
    }


def test_prompt_only_practice_fails_the_v5_contract_without_feedback() -> None:
    practice = apply_page_contract_v5({
        "unit_id": "prompt-only",
        "position": 0,
        "layout": "practice",
        "scene_kind": "practice_feedback",
        "beat_role": "prompt",
        "title": "先判断再说明",
        "blocks": [{
            "block_id": "question",
            "type": "exercise",
            "content": "这个系统属于哪一类？",
        }],
        "quality": {"requested_layout": "question-prompt"},
    })

    assert "practice_feedback_missing_answer" in {
        issue["code"] for issue in v5_contract_issues([practice])
    }


def test_action_task_does_not_receive_unrelated_generic_feedback() -> None:
    slides = [
        {
            "unit_id": "cpu-concept",
            "chapter_id": "chapter-1",
            "knowledge_refs": ["cpu"],
            "scene_kind": "concept",
            "blocks": [{
                "block_id": "cpu-fact",
                "type": "statement",
                "content": "CPU 调度是基于指令执行的。",
            }],
            "quality": {},
        },
        {
            "unit_id": "profiler-task",
            "chapter_id": "chapter-1",
            "source_section_ids": ["lesson-profiler"],
            "knowledge_refs": ["cpu"],
            "layout": "practice",
            "scene_kind": "practice_feedback",
            "beat_role": "prompt",
            "title": "创建性能诊断场景",
            "blocks": [{
                "block_id": "task",
                "type": "exercise",
                "items": [
                    "切换到 CPU 标签并勾选 Deep Profile。",
                    "点击 Play 运行场景并录制至少 5 秒。",
                ],
            }],
            "quality": {"requested_layout": "question-prompt"},
        },
    ]

    enriched = _enrich_practice_feedback_slides_v5(slides)
    task = enriched[-1]

    assert task["quality"]["feedback_mode"] == "task_only"
    assert task["quality"]["task_prompt_mode"] == "action"
    assert task["quality"]["requested_layout"] == "process-sequence"
    assert task["quality"]["prompt_label"] == "执行步骤"
    assert len(task["blocks"]) == 1
    assert "CPU 调度是基于指令执行的" not in str(task["blocks"])
    assert "practice_feedback_missing_answer" not in {
        issue["code"] for issue in v5_contract_issues([task])
    }


def test_task_activity_pages_consolidate_by_visual_grammar_and_reindex() -> None:
    def page(
        index: int,
        *,
        mode: str,
        layout: str,
        items: list[str],
    ) -> dict:
        return {
            "unit_id": f"task-page-{index}",
            "position": index,
            "chapter_id": "chapter-1",
            "source_section_ids": ["lesson-profiler"],
            "layout": "practice",
            "scene_kind": "practice_feedback",
            "beat_role": "prompt",
            "title": "创建一个场景",
            "blocks": [{
                "block_id": f"task-block-{index}",
                "type": "process" if layout == "process-sequence" else "exercise",
                "items": items,
                "metadata": {
                    "semantic_role": (
                        "process_step" if layout == "process-sequence" else "prompt"
                    ),
                    "question_mode": "task",
                },
            }],
            "quality": {
                "requested_layout": layout,
                "feedback_mode": "task_only",
                "task_prompt_mode": mode,
                "fragment_ids": [f"fragment-{index}"],
            },
        }

    source = [
        page(0, mode="action", layout="question-prompt", items=[
            "场景构建：创建一个场景，制造可复现的性能问题。",
        ]),
        page(1, mode="action", layout="process-sequence", items=["切换 CPU 标签。", "开始录制。"]),
        page(2, mode="action", layout="process-sequence", items=["观察曲线。", "定位热点。", "保存结果。"]),
        page(3, mode="action", layout="process-sequence", items=["切换 Memory 标签。"]),
        page(4, mode="action", layout="process-sequence", items=["定位持续增长对象。", "完成修复验证。"]),
        page(5, mode="verification", layout="question-prompt", items=["是否截取了 CPU 视图？"]),
        page(6, mode="verification", layout="question-prompt", items=["是否指出了持续增长对象？"]),
    ]

    consolidated = _consolidate_task_activity_pages_v5(source)

    assert len(consolidated) == 4
    assert [
        item
        for slide in consolidated
        for block in slide["blocks"]
        for item in block.get("items") or []
    ] == [
        item
        for slide in source
        for block in slide["blocks"]
        for item in block.get("items") or []
    ]
    assert [slide["quality"]["practice_page_index"] for slide in consolidated] == [
        1, 2, 3, 4,
    ]
    assert all(
        slide["quality"]["practice_page_count"] == 4
        for slide in consolidated
    )
    assert [slide["quality"]["task_prompt_phase"] for slide in consolidated] == [
        "overview", "procedure", "procedure", "verification",
    ]
    assert [
        len(slide["blocks"][0].get("items") or [])
        for slide in consolidated
    ] == [1, 4, 4, 2]
    assert consolidated[2]["quality"]["requested_layout"] == "process-sequence"
    assert consolidated[2]["blocks"][0]["type"] == "process"
    assert consolidated[0]["title"] == "场景构建：创建一个场景"
    assert consolidated[-1]["title"].endswith("（续4/4）")


def test_task_activity_over_four_pages_is_a_critical_contract_failure() -> None:
    slides = [
        {
            "unit_id": f"task-page-{index}",
            "position": index,
            "chapter_id": "chapter-1",
            "layout": "practice",
            "scene_kind": "practice_feedback",
            "beat_role": "prompt",
            "title": f"任务（续{index + 1}/5）",
            "blocks": [{
                "block_id": f"task-block-{index}",
                "type": "exercise",
                "items": [f"执行任务 {index + 1}"],
                "metadata": {"question_mode": "task"},
            }],
            "quality": {
                "feedback_mode": "task_only",
                "task_prompt_mode": "action",
                "task_activity_id": "task-activity-1",
            },
        }
        for index in range(5)
    ]

    issues = v5_contract_issues(slides)

    assert any(
        issue["code"] == "task_activity_page_limit_exceeded"
        and issue["severity"] == "critical"
        and issue["page_count"] == 5
        for issue in issues
    )


def test_generated_practice_answers_are_bound_to_stable_question_ids() -> None:
    enriched = _enrich_practice_feedback_slides_v5([{
        "unit_id": "generated-practice",
        "position": 0,
        "chapter_id": "chapter-1",
        "scene_kind": "practice_feedback",
        "beat_role": "prompt",
        "title": "判断系统类型",
        "blocks": [{
            "block_id": "questions",
            "type": "exercise",
            "items": [
                "水壶盖子关闭时属于哪类系统？",
                "水壶盖子打开并有蒸汽逸出时呢？",
            ],
        }],
        "quality": {
            "question_ids": ["pptq_closed", "pptq_open"],
            "generated_practice_answers": [
                {
                    "question_index": 0,
                    "question_id": "pptq_closed",
                    "answer_source": "llm_generated",
                    "answer_text": "属于封闭系统，因为没有物质穿过边界。",
                    "supporting_fragment_ids": ["definition-closed"],
                },
                {
                    "question_index": 1,
                    "question_id": "pptq_open",
                    "answer_source": "llm_generated",
                    "answer_text": "属于开放系统，因为蒸汽会穿过边界。",
                    "supporting_fragment_ids": ["definition-open"],
                },
            ],
        },
    }])
    practice = apply_page_contract_v5(enriched[0])
    prompt = practice["blocks"][0]
    answer = practice["blocks"][1]

    assert practice["quality"]["answer_generation_mode"] == "llm_generated"
    assert answer["metadata"]["answer_for_question_ids"] == (
        prompt["metadata"]["question_ids"]
    )
    assert answer["metadata"]["source_fragment_ids"] == [
        "definition-closed",
        "definition-open",
    ]
    assert not {
        "practice_direct_answer_unbound",
        "practice_direct_answer_count_mismatch",
    } & {issue["code"] for issue in v5_contract_issues([practice])}


def test_compound_visible_prompt_coalesces_fragment_level_generated_answers() -> None:
    enriched = _enrich_practice_feedback_slides_v5([{
        "unit_id": "compound-generated-practice",
        "position": 0,
        "chapter_id": "chapter-1",
        "scene_kind": "practice_feedback",
        "beat_role": "prompt",
        "title": "判断并说明理由",
        "blocks": [{
            "block_id": "compound-question",
            "type": "exercise",
            "content": "盖子关闭时属于哪类系统？盖子打开时呢？",
        }],
        "quality": {
            "generated_practice_answers": [
                {
                    "question_index": 0,
                    "answer_text": "盖子关闭时属于封闭系统。",
                    "supporting_fragment_ids": ["definition-closed"],
                },
                {
                    "question_index": 1,
                    "answer_text": "盖子打开时属于开放系统。",
                    "supporting_fragment_ids": ["definition-open"],
                },
            ],
        },
    }])
    practice = apply_page_contract_v5(enriched[0])
    prompt = practice["blocks"][0]
    answer = practice["blocks"][1]

    assert practice["quality"]["answer_generation_mode"] == "llm_generated"
    assert len(answer["items"]) == 1
    assert "封闭系统" in answer["items"][0]
    assert "开放系统" in answer["items"][0]
    assert answer["metadata"]["answer_for_question_ids"] == (
        prompt["metadata"]["question_ids"]
    )
    assert not {
        "practice_direct_answer_unbound",
        "practice_direct_answer_count_mismatch",
    } & {issue["code"] for issue in v5_contract_issues([practice])}


def test_direct_answers_with_missing_question_identity_fail_the_contract() -> None:
    practice = apply_page_contract_v5({
        "unit_id": "unbound-practice",
        "position": 0,
        "scene_kind": "practice_feedback",
        "beat_role": "prompt",
        "title": "判断系统类型",
        "blocks": [
            {
                "block_id": "questions",
                "type": "exercise",
                "items": ["问题一？", "问题二？"],
                "metadata": {
                    "semantic_role": "prompt",
                    "question_ids": ["question-1", "question-2"],
                },
            },
            {
                "block_id": "answers",
                "type": "callout",
                "items": ["答案一。", "答案二。"],
                "metadata": {
                    "semantic_role": "answer",
                    "direct_answer": True,
                    "answer_for_question_ids": ["question-1"],
                },
            },
        ],
        "quality": {
            "requested_layout": "practice-feedback",
            "feedback_mode": "paired",
        },
    })

    assert "practice_direct_answer_unbound" in {
        issue["code"] for issue in v5_contract_issues([practice])
    }


def test_repeated_episode_pages_keep_their_distinct_visible_heading() -> None:
    slides = _assign_heading_modes_v5([
        {
            "unit_id": "concept-1",
            "position": 0,
            "layout": "concept",
            "slide_purpose": "concept",
            "episode_id": "episode-state",
            "section_id": "section-state",
            "eyebrow": "核心概念",
            "title": "状态变量只由当前状态决定",
            "key_message": "1.2 状态变量与过程量",
            "blocks": [{"block_id": "a", "type": "rich_text", "content": "定义。"}],
            "quality": {},
        },
        {
            "unit_id": "concept-2",
            "position": 1,
            "layout": "concept",
            "slide_purpose": "concept",
            "episode_id": "episode-state",
            "section_id": "section-state",
            "eyebrow": "核心概念",
            "title": "温度和压力都是状态变量",
            "key_message": "",
            "blocks": [{"block_id": "b", "type": "rich_text", "content": "举例。"}],
            "quality": {},
        },
    ])

    assert slides[0]["quality"]["heading_mode"] == "full"
    assert slides[0]["quality"]["section_label"] == "1.2 状态变量与过程量"
    assert slides[1]["quality"]["heading_mode"] == "full"
    assert slides[1]["quality"]["section_label"] == "1.2 状态变量与过程量"
    assert slides[1]["title"] == "温度和压力都是状态变量"


def test_export_keeps_the_page_claim_visible_even_if_legacy_metadata_hides_it() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    unit = SlideSpec.model_validate({
        "unit_id": "continuation",
        "position": 0,
        "layout": "concept",
        "slide_purpose": "concept",
        "eyebrow": "核心概念",
        "title": "温度和压力都是状态变量",
        "blocks": [],
        "quality": {
            "heading_mode": "hidden",
            "section_label": "1.2 状态变量与过程量",
        },
    })

    _heading(slide, unit, validate_theme("qizhi-classroom"))

    visible_text = "\n".join(
        shape.text
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert "1.2 状态变量与过程量" in visible_text
    assert "温度和压力都是状态变量" in visible_text


def test_renderer_does_not_hard_cut_a_complete_compiled_heading() -> None:
    title = "碰撞回调事件的封装与分层处理模式以及性能裁剪策略"

    assert _heading_excerpt(title) == title


def test_long_editorial_objective_is_structured_into_three_visible_points() -> None:
    slide = {
        "unit_id": "slide:v5:long-objective",
        "layout": "concept",
        "slide_purpose": "concept",
        "scene_kind": "concept",
        "title": "Input System connects actions to movement",
        "blocks": [{
            "block_id": "objective",
            "type": "statement",
            "title": "",
            "content": (
                "This lesson configures the Input System and defines an Action Map. "
                "Learners create Move and Jump actions with explicit bindings. "
                "The final script reads action data and drives movement in the scene."
            ),
            "items": [],
            "metadata": {"fragment_ids": ["fragment-objective"]},
        }],
        "visuals": [],
        "quality": {
            "requested_layout": "editorial-body",
            "resolved_layout": "editorial-body",
        },
    }

    structured = _structure_long_editorial_prose_v5(slide)

    assert structured["quality"]["requested_layout"] == "classification-3"
    assert structured["blocks"][0]["type"] == "bullets"
    assert len(structured["blocks"][0]["items"]) == 3
    assert structured["blocks"][0]["metadata"]["fragment_ids"] == [
        "fragment-objective"
    ]


def test_enumeration_scaffolding_is_removed_before_rendering_cards() -> None:
    slide = {
        "unit_id": "slide:v5:objectives",
        "title": "本页核心判断",
        "scene_kind": "concept",
        "blocks": [{
            "block_id": "objectives",
            "type": "bullets",
            "content": "本节聚焦于构建与发布配置。完成本节后，你将能够：",
            "items": ["切换目标平台", "校验依赖", "输出发布清单"],
            "metadata": {"fragment_ids": ["fragment-objectives"]},
        }],
        "quality": {"requested_layout": "classification-3"},
    }

    cleaned = _strip_instructional_scaffolding_v5(slide)

    assert cleaned["blocks"][0]["content"] == ""
    assert cleaned["blocks"][0]["metadata"]["fragment_ids"] == [
        "fragment-objectives"
    ]
    assert cleaned["quality"]["instructional_scaffolding_suppressed"] is True


def test_long_method_paragraph_becomes_two_source_bound_regions() -> None:
    slide = {
        "unit_id": "slide:v5:method-prose",
        "title": "线程安全单例",
        "scene_kind": "method",
        "blocks": [{
            "block_id": "method",
            "type": "statement",
            "content": (
                "单例负责管理跨场景共享的全局资源。"
                "多线程访问时必须通过同步机制保护初始化过程，避免创建重复实例。"
            ),
            "items": [],
            "metadata": {"fragment_ids": ["fragment-method"]},
        }],
        "visuals": [],
        "quality": {
            "requested_layout": "editorial-body",
            "resolved_layout": "editorial-body",
        },
    }

    structured = _structure_long_editorial_prose_v5(slide)

    assert structured["quality"]["requested_layout"] == "balanced-two-column"
    assert len(structured["blocks"][0]["items"]) == 2


def test_sparse_single_concept_claim_becomes_an_intentional_hero_page() -> None:
    slide = {
        "unit_id": "slide:v5:single-claim",
        "title": "适配非标准屏幕比例",
        "scene_kind": "concept",
        "blocks": [{
            "block_id": "claim",
            "type": "bullets",
            "content": "",
            "items": ["脚本动态调整 UI 适配参数，以应对非标准屏幕比例。"],
            "metadata": {"fragment_ids": ["fragment-claim"]},
        }],
        "visuals": [],
        "quality": {
            "requested_layout": "classification-3",
            "resolved_layout": "editorial-body",
        },
    }

    promoted = _promote_sparse_single_claim_v5(slide)

    assert promoted["quality"]["requested_layout"] == "hero-claim"
    assert promoted["quality"]["suppress_redundant_body"] is True
    assert promoted["key_message"] == (
        "脚本动态调整 UI 适配参数，以应对非标准屏幕比例。"
    )
    assert promoted["blocks"][0]["type"] == "statement"
    assert promoted["blocks"][0]["items"] == []
    assert promoted["blocks"][0]["metadata"]["fragment_ids"] == [
        "fragment-claim"
    ]


def test_promoted_hero_renderer_shows_the_source_claim() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    unit = SlideSpec.model_validate({
        "unit_id": "slide:v5:promoted-hero",
        "position": 0,
        "layout": "concept",
        "slide_purpose": "concept",
        "title": "适配非标准屏幕比例",
        "key_message": "脚本动态调整 UI 适配参数，以应对非标准屏幕比例。",
        "teaching_job": "建立本节核心概念与边界",
        "blocks": [{
            "block_id": "claim",
            "type": "statement",
            "content": "脚本动态调整 UI 适配参数，以应对非标准屏幕比例。",
            "items": [],
            "metadata": {},
        }],
        "quality": {"resolved_layout": "hero-claim"},
    })

    _render_claim_only(slide, unit, validate_theme("qizhi-classroom"))

    visible_text = "\n".join(
        shape.text
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert "脚本动态调整 UI 适配参数" in visible_text
    assert "建立本节核心概念与边界" not in visible_text
    claim_shapes = [
        shape
        for shape in slide.shapes
        if (
            getattr(shape, "has_text_frame", False)
            and "脚本动态调整 UI 适配参数" in shape.text
        )
    ]
    assert len(claim_shapes) == 1
    assert claim_shapes[0].height >= Inches(1.5)


def test_code_renderer_uses_full_width_when_no_real_annotation_exists() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    unit = SlideSpec.model_validate({
        "unit_id": "slide:v5:code-only",
        "position": 0,
        "layout": "code",
        "slide_purpose": "method",
        "title": "生命周期回调顺序",
        "key_message": "",
        "blocks": [{
            "block_id": "code",
            "type": "code",
            "content": "void Awake() {}\nvoid Start() {}\nvoid Update() {}",
            "items": [],
            "metadata": {"language": "csharp"},
        }],
        "quality": {"resolved_layout": "code"},
    })

    _render_code(slide, unit, validate_theme("qizhi-classroom"))

    visible_text = "\n".join(
        shape.text
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    code_shapes = [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and "void Awake" in shape.text
    ]
    assert "阅读线索" not in visible_text
    assert len(code_shapes) == 1
    assert code_shapes[0].width >= Inches(10.5)


def test_task_process_renderer_uses_vertical_numbered_steps() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    items = [
        "切换到 CPU 标签并开始录制。",
        "定位最耗时的方法并记录调用路径。",
        "修改脚本后重新运行场景，比较前后曲线。",
    ]
    unit = SlideSpec.model_validate({
        "unit_id": "slide:v5:task-procedure",
        "position": 0,
        "layout": "practice",
        "slide_purpose": "practice",
        "title": "性能诊断任务（续2/4）",
        "blocks": [{
            "block_id": "procedure",
            "type": "process",
            "items": items,
            "metadata": {"semantic_role": "process_step"},
        }],
        "quality": {
            "resolved_layout": "process-sequence",
            "task_prompt_mode": "action",
            "task_prompt_phase": "procedure",
            "prompt_label": "执行步骤",
        },
    })

    _render_process(slide, unit, validate_theme("qizhi-classroom"))

    visible_text = "\n".join(
        shape.text
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert "执行步骤" in visible_text
    assert all(item in visible_text for item in items)
    assert all(f"{index:02d}" in visible_text for index in range(1, 4))


def test_code_renderer_keeps_annotation_column_when_source_annotation_exists() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    unit = SlideSpec.model_validate({
        "unit_id": "slide:v5:annotated-code",
        "position": 0,
        "layout": "code",
        "slide_purpose": "method",
        "title": "生命周期回调顺序",
        "key_message": "",
        "blocks": [
            {
                "block_id": "code",
                "type": "code",
                "content": "void Awake() {}\nvoid Start() {}",
                "items": [],
                "metadata": {"language": "csharp"},
            },
            {
                "block_id": "annotation",
                "type": "bullets",
                "content": "",
                "items": ["Awake 先于 Start 执行。"],
                "metadata": {},
            },
        ],
        "quality": {"resolved_layout": "code"},
    })

    _render_code(slide, unit, validate_theme("qizhi-classroom"))

    visible_text = "\n".join(
        shape.text
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert "阅读线索" in visible_text
    assert "Awake 先于 Start 执行" in visible_text


def test_empty_chapter_entry_restores_its_outline_learning_objective() -> None:
    chapter = DeckChapterV5(
        chapter_id="chapter-input",
        agenda_id="agenda-input",
        position=0,
        eyebrow="第03章",
        title="交互逻辑与物理系统基础",
        driving_question="配置输入系统并完成角色移动控制。",
        learning_objective="配置输入系统并完成角色移动控制。",
    )
    slide = {
        "unit_id": "slide:v5:chapter:chapter-input",
        "chapter_id": "chapter-input",
        "scene_kind": "chapter_entry",
        "title": chapter.title,
        "key_message": "",
        "quality": {"requested_layout": "chapter-entry"},
    }

    restored = _restore_chapter_entry_mainlines_v5([slide], [chapter])

    assert restored[0]["key_message"] == chapter.learning_objective


def test_restored_chapter_entry_uses_navigation_copy_when_objective_repeats_next_page() -> None:
    chapter = DeckChapterV5(
        chapter_id="chapter-input",
        agenda_id="agenda-input",
        position=0,
        eyebrow="第三章",
        title="交互逻辑与物理系统基础",
        driving_question="配置输入系统并完成角色移动控制。",
        learning_objective="配置输入系统并完成角色移动控制。",
    )
    entry = {
        "unit_id": "slide:v5:chapter:chapter-input",
        "chapter_id": "chapter-input",
        "scene_kind": "chapter_entry",
        "title": chapter.title,
        "key_message": "",
        "blocks": [],
        "quality": {"requested_layout": "chapter-entry"},
    }
    first_content_page = {
        "unit_id": "slide:v5:first-content",
        "chapter_id": "chapter-input",
        "scene_kind": "concept",
        "title": "输入系统配置",
        "key_message": chapter.learning_objective,
        "blocks": [],
        "quality": {"requested_layout": "hero-claim"},
    }

    restored = _restore_chapter_entry_mainlines_v5(
        [entry, first_content_page],
        [chapter],
    )

    assert restored[0]["key_message"] != chapter.learning_objective
    assert chapter.title in restored[0]["key_message"]
    assert restored[0]["quality"]["chapter_entry_mainline_restored"] is True


def test_labeled_error_and_inference_sequence_becomes_three_peer_regions() -> None:
    slide = {
        "unit_id": "slide:v5:reasoning-errors",
        "scene_kind": "reasoning",
        "title": "典型错误与推导依据",
        "blocks": [{
            "block_id": "reasoning-errors",
            "type": "process",
            "content": "",
            "items": [
                "错误 1：只关注 CPU 总耗时，忽略了 GC Pause。",
                "推导：需在 Memory 标签页验证垃圾回收。",
                "错误 2：把初始化代码误判为每帧瓶颈。",
                "推导：检查时间轴，区分一次性开销与每帧开销。",
                "错误 3：修复后未做回归测试，导致旧功能失效。",
            ],
            "metadata": {"source_fragment_ids": ["fragment-errors"]},
        }],
        "visuals": [{
            "visual_id": "visual-errors",
            "kind": "relational_diagram",
            "alt_text": "错误与推导关系图",
        }],
        "quality": {
            "requested_layout": "editorial-body",
            "resolved_layout": "figure-text",
        },
    }

    structured = _structure_labeled_reasoning_pairs_v5(slide)
    contracted = apply_page_contract_v5(structured)

    assert structured["blocks"][0]["type"] == "bullets"
    assert len(structured["blocks"][0]["items"]) == 3
    assert structured["quality"]["requested_layout"] == "classification-3"
    assert contracted["quality"]["resolved_layout"] == "classification-3"
    assert all(
        text in " ".join(structured["blocks"][0]["items"])
        for text in ("错误 1", "推导", "错误 3")
    )


def test_editorial_fallback_does_not_render_the_heading_twice() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    unit = SlideSpec.model_validate({
        "unit_id": "slide:v5:editorial-fallback",
        "position": 0,
        "layout": "concept",
        "slide_purpose": "teach",
        "eyebrow": "核心概念",
        "title": "生命周期回调执行顺序",
        "blocks": [{
            "block_id": "body",
            "type": "statement",
            "content": "Awake 完成初始化，Start 在首次 Update 前执行。",
        }],
        "quality": {"resolved_layout": "editorial-body"},
    })
    theme = validate_theme("qingfeng-classroom")

    _heading(slide, unit, theme)
    _render_editorial_body(
        slide,
        unit,
        theme,
        heading_already_rendered=True,
    )

    visible = [
        shape.text
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    ]
    assert visible.count(unit.title) == 1


def test_required_programming_code_is_a_final_publication_gate() -> None:
    report = finalize_v5_quality_report(
        previous_quality={"passed": True, "issues": [], "blockers": []},
        slides=[{
            "unit_id": "concept-only",
            "position": 0,
            "layout": "concept",
            "slide_purpose": "concept",
            "scene_kind": "concept",
            "title": "MonoBehaviour 生命周期",
            "blocks": [{
                "block_id": "body",
                "type": "statement",
                "content": "Awake、Start 和 Update 按生命周期顺序执行。",
                "metadata": {},
            }],
            "visuals": [],
            "quality": {
                "passed": True,
                "resolved_layout": "editorial-body",
            },
        }],
        planner="ai",
        fallback_reason="",
        planning_diagnostics={
            "subject_presentation_contract": {
                "schema_version": "subject_presentation_contract_v1",
                "profile_id": "engineering_programming",
                "primary_mode": "programming_engineering",
                "required_representation_kinds": ["code"],
                "optional_representation_kinds": ["output", "debugging"],
                "characteristic_fragment_ids": {"code": ["code-fragment"]},
                "chapter_requirements": [{
                    "chapter_id": "chapter-unity",
                    "required_representation_kinds": ["code"],
                    "minimum_artifact_count": 1,
                }],
                "classification_confidence": 0.96,
                "classification_source": "course_generation_v16",
                "evidence_conflicts": [],
            },
        },
    )

    issue = next(
        item for item in report["issues"]
        if item["code"] == "required_subject_representation_missing"
    )
    assert issue["severity"] == "critical"
    assert issue["representation_kind"] == "code"
    assert report["passed"] is False


def test_presentation_grammar_mismatch_requires_manual_review() -> None:
    contracted = apply_page_contract_v5({
        "unit_id": "process-as-columns",
        "position": 0,
        "layout": "concept",
        "slide_purpose": "method",
        "scene_kind": "method",
        "title": "回调按生命周期顺序执行",
        "blocks": [
            {
                "block_id": "left",
                "type": "statement",
                "content": "Awake 完成初始化。",
                "metadata": {},
            },
            {
                "block_id": "right",
                "type": "statement",
                "content": "Update 每帧执行。",
                "metadata": {},
            },
        ],
        "visuals": [],
        "quality": {
            "requested_layout": "balanced-two-column",
            "presentation_grammar": {
                "presentation_intent": "process",
                "copy_voice": "ordered_instructional",
                "information_structure": "sequence",
                "visual_grammar": "control_flow",
                "allowed_layouts": ["process-sequence", "figure-text"],
                "forbidden_fallbacks": ["editorial-body", "balanced-two-column"],
            },
        },
    })

    issue = next(
        item for item in contracted["quality"]["issues"]
        if item["code"] == "presentation_grammar_mismatch"
    )
    assert issue["expected_grammar"] == "control_flow"
    assert issue["observed_layout"] == "balanced-two-column"
    assert contracted["quality"]["manual_edit_required"] is True


def test_classification_renderer_uses_content_and_items_as_three_regions() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    unit = SlideSpec.model_validate({
        "unit_id": "slide:v5:mixed-classification",
        "position": 0,
        "layout": "concept",
        "slide_purpose": "teach",
        "eyebrow": "核心概念",
        "title": "三项判断",
        "blocks": [
            {
                "block_id": "definition",
                "type": "statement",
                "content": "第一项判断",
            },
            {
                "block_id": "details",
                "type": "bullets",
                "items": ["第二项判断", "第三项判断"],
            },
        ],
        "quality": {"resolved_layout": "classification-3"},
    })

    _render_classification_three(
        slide,
        unit,
        validate_theme("qingfeng-classroom"),
    )

    visible_shapes = [
        shape.text
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    ]
    visible = "\n".join(visible_shapes)
    assert visible_shapes.count(unit.title) == 1
    assert all(value in visible for value in ("第一项判断", "第二项判断", "第三项判断"))


def test_duplicate_source_heading_uses_the_next_pages_source_question() -> None:
    first = {
        "unit_id": "slide:v5:first",
        "title": "预制体源文件与实例化对象",
        "scene_kind": "concept",
        "blocks": [{"content": "学习者需要完成三个可观察目标。", "items": []}],
        "quality": {"resolved_layout": "classification-3"},
    }
    second = {
        "unit_id": "slide:v5:second",
        "title": "预制体源文件与实例化对象",
        "scene_kind": "concept",
        "blocks": [{
            "content": "为什么变体能保持独立性又同步更新？变体保存父级引用与局部覆盖表。",
            "items": [],
        }],
        "quality": {"resolved_layout": "editorial-body"},
    }

    resolved = _disambiguate_duplicate_titles_v5([first, second])

    assert resolved[0]["title"] == "预制体源文件与实例化对象"
    assert resolved[1]["title"] == "为什么变体能保持独立性又同步更新"
    assert resolved[1]["blocks"][0]["content"] == (
        "变体保存父级引用与局部覆盖表。"
    )


def test_export_pairs_each_practice_question_with_its_answer() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    unit = SlideSpec.model_validate({
        "unit_id": "paired-practice",
        "position": 0,
        "layout": "practice",
        "slide_purpose": "practice_feedback",
        "eyebrow": "理解检查",
        "title": "判断系统类型",
        "blocks": [
            {
                "block_id": "questions",
                "type": "exercise",
                "items": ["盖子关闭时属于哪类系统？", "盖子打开时呢？"],
            },
            {
                "block_id": "answers",
                "type": "callout",
                "items": ["封闭系统。", "开放系统。"],
            },
        ],
        "quality": {"resolved_layout": "practice-feedback"},
    })

    _render_practice_feedback(slide, unit, validate_theme("qizhi-classroom"))

    text_shapes = {
        shape.text: shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text
    }
    assert text_shapes["盖子关闭时属于哪类系统？"].top == text_shapes["封闭系统。"].top
    assert text_shapes["盖子打开时呢？"].top == text_shapes["开放系统。"].top


def test_export_pairs_answers_by_identity_even_when_answer_order_differs() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    unit = SlideSpec.model_validate({
        "unit_id": "identity-paired-practice",
        "position": 0,
        "layout": "practice",
        "slide_purpose": "practice_feedback",
        "eyebrow": "理解检查",
        "title": "判断系统类型",
        "blocks": [
            {
                "block_id": "questions",
                "type": "exercise",
                "items": ["盖子关闭时属于哪类系统？", "盖子打开时呢？"],
                "metadata": {
                    "question_ids": ["closed", "open"],
                },
            },
            {
                "block_id": "answers",
                "type": "callout",
                "items": ["开放系统。", "封闭系统。"],
                "metadata": {
                    "semantic_role": "answer",
                    "answer_for_question_ids": ["open", "closed"],
                },
            },
        ],
        "quality": {
            "resolved_layout": "practice-feedback",
            "feedback_mode": "paired",
        },
    })

    _render_practice_feedback(slide, unit, validate_theme("qizhi-classroom"))

    text_shapes = {
        shape.text: shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text
    }
    assert text_shapes["盖子关闭时属于哪类系统？"].top == text_shapes["封闭系统。"].top
    assert text_shapes["盖子打开时呢？"].top == text_shapes["开放系统。"].top


def test_export_does_not_present_related_evidence_as_direct_answers() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    unit = SlideSpec.model_validate({
        "unit_id": "shared-evidence",
        "position": 0,
        "layout": "practice",
        "slide_purpose": "practice_feedback",
        "eyebrow": "理解检查",
        "title": "判断系统类型",
        "blocks": [
            {
                "block_id": "questions",
                "type": "exercise",
                "items": ["盖子关闭时属于哪类系统？", "盖子打开时呢？"],
            },
            {
                "block_id": "evidence",
                "type": "callout",
                "items": ["封闭系统不交换物质。", "开放系统可以交换物质。"],
                "metadata": {"direct_answer": False},
            },
        ],
        "quality": {
            "resolved_layout": "practice-feedback",
            "feedback_mode": "shared_evidence",
        },
    })

    _render_practice_feedback(slide, unit, validate_theme("qizhi-classroom"))

    labels = [
        shape.text
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text
    ]
    assert labels.count("判断依据") == 1
    assert not any("回答与判断依据" in label for label in labels)
