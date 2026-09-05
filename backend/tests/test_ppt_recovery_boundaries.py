import asyncio
import json
from copy import deepcopy
from types import SimpleNamespace

import pytest

from .test_ppt_teaching_content import compiled_manuscript


def test_provider_failure_persists_accepted_pages_for_a_new_planner_instance(tmp_path):
    from ppt_teaching_planner import plan_teaching_manuscript
    from slide_deck_v6_models import V6BuildError
    doc, graph, template, manuscript = compiled_manuscript()
    page = manuscript.pages[0]
    trace_path = tmp_path / 'planning.json'
    revision = {'title': page.title, 'page_goal': page.page_goal, 'primary_claim': page.primary_claim,
                'teaching': page.teaching.model_dump(mode='json')}
    async def save(trace, event):
        trace_path.write_text(json.dumps(trace))
    async def first(request):
        if request['teaching_request'] == 'narrative':
            return {'narrative_brief': {'central_question': '比较执行方式'}, 'pages': [
                {'page_id': key, 'teaching_unit_id': page.teaching_unit_id, 'source_block_ids': ['b'],
                 'title': page.title, 'page_goal': page.page_goal, 'layout_id': page.layout_id} for key in ('p1', 'p2')]}
        if request['page']['page_id'] == 'p2':
            raise RuntimeError('provider unavailable')
        return deepcopy(revision)
    with pytest.raises(V6BuildError, match='teaching_provider_failed'):
        asyncio.run(plan_teaching_manuscript(doc, graph, template, first, on_checkpoint=save))
    saved = json.loads(trace_path.read_text())
    assert set(saved['pages']) == {'p1'}
    calls = []
    async def recovered(request):
        calls.append(request['page']['page_id'])
        return deepcopy(revision)
    restored, _ = asyncio.run(plan_teaching_manuscript(doc, graph, template, recovered, checkpoint=saved))
    assert calls == ['p2'] and restored.story_page_count == 2
    changed = template.model_copy(deep=True, update={'template_digest': 'changed'})
    with pytest.raises(V6BuildError, match='teaching_planning_checkpoint_mismatch'):
        asyncio.run(plan_teaching_manuscript(doc, graph, changed, recovered, checkpoint=saved))
    assert calls == ['p2']


def test_missing_font_and_tool_change_reject_export_without_overwriting(tmp_path, monkeypatch):
    from ppt_native_scene import render_teaching_deck
    from ppt_teaching_manuscript import physical_pages
    from ppt_runtime_identity import tool_identity
    _, _, _, manuscript = compiled_manuscript()
    deck = SimpleNamespace(pages=physical_pages(manuscript))
    output = tmp_path / 'last.pptx'
    output.write_bytes(b'last-good')
    tools = tool_identity()
    with monkeypatch.context() as m:
        m.setattr('ppt_runtime_identity.FONT_PATH', tmp_path / 'missing.otf')
        with pytest.raises(ValueError, match='teaching_font_missing'):
            render_teaching_deck(deck, output)
    with monkeypatch.context() as m:
        m.setattr('ppt_page_scene.tool_identity', lambda: {**tools, 'renderer': 'changed'})
        with pytest.raises(ValueError, match='teaching_execution_environment_changed'):
            render_teaching_deck(deck, output)
    assert output.read_bytes() == b'last-good'
    assert not list(tmp_path.glob('.last-*'))


def test_missing_adopted_asset_rejects_scene_without_placeholder_success(tmp_path):
    from ppt_native_scene import render_scene
    from ppt_adopted_visuals import current_visual_catalog, bind_adopted_assets
    from ppt_teaching_content import PageTeachingV2
    from .test_ppt_adopted_visuals import adopted_fixture
    from .test_ppt_teaching_content import scene_for
    from pptx import Presentation
    assets, asset, item, value, sources = adopted_fixture(tmp_path)
    catalog = current_visual_catalog([item], course_id='course', script_revision_id='script-r1', sources=sources, asset_repository=assets)
    content = bind_adopted_assets(PageTeachingV2.model_validate(value), catalog, {'b'})
    scene = scene_for(content.model_dump(mode='json'), 'compare-visual')
    assets.resolve(asset.asset_id).unlink()
    p = Presentation()
    with pytest.raises(FileNotFoundError):
        render_scene(p.slides.add_slide(p.slide_layouts[6]), scene, assets=assets)
