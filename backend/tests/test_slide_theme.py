from slide_layout_registry import select_layout_v2
from slide_theme import (
    REPOSITORY_ROOT,
    load_slide_theme_pack,
    slide_theme,
    slide_theme_asset_path,
)
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


COMPLETE_THEME_NAMES = (
    "qizhi-classroom",
    "academic-editorial",
    "grid-notebook",
    "modern-geometric",
    "dark-tech",
)
BACKGROUND_ASSET_NAMES = (
    "cover",
    "chapter",
    "recap",
    "interior_content",
    "interior_reasoning",
    "interior_practice",
    "interior_evidence",
)
TEXT_BOX_STYLE_NAMES = (
    "standard",
    "message",
    "definition",
    "boundary",
    "reasoning",
    "practice",
    "feedback",
    "misconception",
    "evidence",
    "note",
)


def test_qizhi_classroom_theme_bundles_authored_visual_assets() -> None:
    theme = slide_theme("qizhi-classroom")

    assert theme["template"]["template_id"] == "qizhi-classroom-v2"
    for asset_name in (
        "cover",
        "chapter",
        "recap",
        "interior_content",
        "interior_reasoning",
        "interior_practice",
        "interior_evidence",
    ):
        asset_path = slide_theme_asset_path(theme, asset_name)
        assert asset_path is not None
        assert asset_path.is_file()
    assert theme["background_profiles"]["practice"]["asset"] == "interior_practice"
    assert theme["text_box_styles"]["misconception"]["accent"] == "C45443"
    assert theme["text_box_styles"]["misconception"]["depth"] == "E9BEB5"


def test_qizhi_classroom_layout_preferences_participate_in_selection() -> None:
    selection = select_layout_v2(
        scene_kind="chapter_entry",
        evidence_kinds=["text"],
        character_count=120,
        item_count=1,
        theme="qizhi-classroom",
    )

    assert selection.layout_id == "chapter-question"
    assert selection.theme_score == 1.3
    assert "theme=qizhi-classroom:1.30" in selection.reason


def test_every_selectable_theme_is_a_complete_template_pack() -> None:
    pack = load_slide_theme_pack()

    assert pack["schema_version"] == "slide_theme_pack_v1"
    for theme_name in COMPLETE_THEME_NAMES:
        theme = slide_theme(theme_name)
        assert theme["template"]["template_id"]
        assert theme["template"]["reference_deck"].endswith(".pptx")
        assert theme["typography"]["slide_title_pt"] >= 35
        assert theme["geometry"]["card_radius_in"] > 0
        assert set(theme["background_profiles"]) == {
            "content", "reasoning", "practice", "evidence",
        }
        assert set(theme["text_box_styles"]) == set(TEXT_BOX_STYLE_NAMES)
        assert len(theme["semantic_layout_weights"]) >= 18
        for asset_name in BACKGROUND_ASSET_NAMES:
            assert slide_theme_asset_path(theme, asset_name) is not None


def test_every_complete_theme_reference_deck_is_bundled() -> None:
    total_pages = 0
    for theme_name in COMPLETE_THEME_NAMES:
        theme = slide_theme(theme_name)
        reference = str(theme["template"]["reference_deck"]).lstrip("/")
        path = REPOSITORY_ROOT / "frontend" / "public" / reference
        assert path.is_file(), f"missing reference deck for {theme_name}"
        assert path.read_bytes().startswith(b"PK")
        deck = Presentation(path)
        assert len(deck.slides) == 18
        assert round(deck.slide_width / deck.slide_height, 4) == 1.7778
        total_pages += len(deck.slides)
        empty_text_boxes = [
            shape
            for slide in deck.slides
            for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
            and not shape.text.strip()
        ]
        assert not empty_text_boxes, f"empty text boxes in {theme_name}"
    assert total_pages == 90
