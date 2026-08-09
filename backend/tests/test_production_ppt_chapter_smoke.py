import json

from pptx import Presentation
from pptx.util import Inches

from course_document import CourseBlock, CourseDocument, CourseSection
from production_ppt_chapter_smoke import (
    SmokeFailure,
    _planner_failure_reason_code,
    _pptx_presentation_mode_audit,
    _planned_scene_requirements,
    _source_disposition,
    build_chapter_document,
    extract_source_code_lines,
    finalize_deferred_render,
    rank_programming_chapter_candidates,
)
from slide_deck import SlideSpec
from slide_deck_renderer import _render_claim_only, _render_code, validate_theme


def test_smoke_uses_final_source_bound_story_scenes() -> None:
    story = {
        "chapters": [{
            "episodes": [
                {"scene_kind": "chapter_entry", "beats": []},
                {
                    "scene_kind": "concept",
                    "beats": [{"fragment_ids": ["concept-1"]}],
                },
                {
                    "scene_kind": "worked_example",
                    "beats": [{"fragment_ids": []}],
                },
                {
                    "scene_kind": "practice_feedback",
                    "beats": [{"fragment_ids": ["feedback-1"]}],
                },
                {"scene_kind": "chapter_recap", "beats": []},
            ],
        }],
    }

    assert _planned_scene_requirements(story) == {
        "concept",
        "practice_feedback",
    }


def test_smoke_reports_provider_balance_failure_without_raw_request_data() -> None:
    reason = _planner_failure_reason_code([{
        "error_type": "RuntimeError",
        "message": "Error code: 429 - insufficient balance; request_id=private",
    }])

    assert reason == "ai_provider_balance_exhausted"


def test_smoke_failure_classifier_ignores_opaque_identifier_digits() -> None:
    reason = _planner_failure_reason_code([{
        "batch_index": 1,
        "chapter_ids": ["chapter-429-observability"],
        "page_ids": ["page-401-lifecycle"],
        "failure_category": "ValueError",
        "message": "Visual batch contained no valid requested pages",
    }])

    assert reason == "ai_planner_failed"


def test_smoke_reports_authentication_only_from_explicit_failure_diagnostics() -> None:
    reason = _planner_failure_reason_code([{
        "page_ids": ["page-safe"],
        "failure_category": "AuthenticationError",
        "message": "Error code: 401 - invalid api key",
    }])

    assert reason == "ai_provider_authentication_failed"


def test_smoke_accepts_only_explicit_code_source_disposition() -> None:
    allocation = {
        "pages": [{"fragment_ids": ["code-1", "prose-1"]}],
        "exclusions": [{
            "fragment_id": "code-2",
            "reason": "subject_artifact_redundant_after_chapter_coverage",
        }],
    }

    allocated, excluded = _source_disposition(
        allocation,
        {"code-1", "code-2"},
    )

    assert allocated == {"code-1"}
    assert excluded == {
        "code-2": "subject_artifact_redundant_after_chapter_coverage",
    }


def test_deferred_render_finalizer_promotes_only_matching_page_count(
    tmp_path,
    monkeypatch,
) -> None:
    (tmp_path / "chapter-smoke.pptx").write_bytes(b"pptx")
    (tmp_path / "report.json").write_text(
        json.dumps({
            "status": "passed_pending_render",
            "chain": {"candidate_status": "v5_ready"},
            "deck": {"slide_count": 3},
            "gates": {"quality_gate": True},
            "export": {"pptx_bytes": 4},
        }),
        encoding="utf-8",
    )
    monkeypatch.setattr(
        "production_ppt_chapter_smoke._render_artifacts",
        lambda _pptx, _output: {
            "pdf_bytes": 100,
            "rendered_page_count": 3,
            "contact_sheet_bytes": 50,
        },
    )

    report = finalize_deferred_render(tmp_path)

    assert report["status"] == "passed"
    assert report["gates"]["rendered_page_count_matches"] is True
    assert report["render_verification"] == {
        "status": "passed",
        "executor": "isolated_ci_runner",
    }


def test_deferred_render_finalizer_blocks_page_count_mismatch(
    tmp_path,
    monkeypatch,
) -> None:
    (tmp_path / "chapter-smoke.pptx").write_bytes(b"pptx")
    (tmp_path / "report.json").write_text(
        json.dumps({
            "status": "passed_pending_render",
            "chain": {"candidate_status": "v5_needs_manual_edit"},
            "deck": {"slide_count": 3},
            "gates": {"quality_gate": True},
            "export": {"pptx_bytes": 4},
        }),
        encoding="utf-8",
    )
    monkeypatch.setattr(
        "production_ppt_chapter_smoke._render_artifacts",
        lambda _pptx, _output: {
            "pdf_bytes": 100,
            "rendered_page_count": 2,
            "contact_sheet_bytes": 50,
        },
    )

    report = finalize_deferred_render(tmp_path)

    assert report["status"] == "failed"
    assert report["failure"]["failed_gates"] == [
        "rendered_page_count_matches",
    ]


def _document() -> CourseDocument:
    return CourseDocument(
        course_id="course-online",
        title="Online programming course",
        document_revision="cdr_online",
        sections=[
            CourseSection(
                section_id="chapter-code",
                title="Programming chapter",
                position=0,
                level=1,
            ),
            CourseSection(
                section_id="lesson-code",
                parent_section_id="chapter-code",
                title="Lifecycle methods",
                position=1,
                level=2,
            ),
            CourseSection(
                section_id="chapter-prose",
                title="Background chapter",
                position=2,
                level=1,
            ),
            CourseSection(
                section_id="lesson-prose",
                parent_section_id="chapter-prose",
                title="Background reading",
                position=3,
                level=2,
            ),
        ],
        blocks=[
            CourseBlock(
                block_id="concept",
                section_id="lesson-code",
                position=0,
                kind="rich_text",
                role="concept",
                payload={"markdown": "Lifecycle callbacks have a defined order."},
            ),
            CourseBlock(
                block_id="code",
                section_id="lesson-code",
                position=1,
                kind="code",
                role="example",
                payload={
                    "markdown": (
                        "```csharp\n"
                        "void Tick1() { Debug.Log(1); }\n"
                        "void Tick32() { Debug.Log(32); }\n"
                        "```"
                    )
                },
            ),
            CourseBlock(
                block_id="practice",
                section_id="lesson-code",
                position=2,
                kind="practice_ref",
                role="checkpoint",
                payload={"markdown": "Explain the callback order."},
            ),
            CourseBlock(
                block_id="prose",
                section_id="lesson-prose",
                position=0,
                kind="rich_text",
                role="concept",
                payload={"markdown": "Only prose is available here."},
            ),
        ],
    )


def test_programming_smoke_selects_one_source_chapter_with_code_and_loop() -> None:
    document = _document()

    candidates = rank_programming_chapter_candidates(document)

    assert [item.chapter_id for item in candidates] == ["chapter-code"]
    assert candidates[0].source_role_count == 3
    assert candidates[0].code_character_count > 40

    chapter_document = build_chapter_document(
        document,
        candidates[0].chapter_id,
    )
    assert [item.section_id for item in chapter_document.sections] == [
        "chapter-code",
        "lesson-code",
    ]
    assert [item.block_id for item in chapter_document.blocks] == [
        "concept",
        "code",
        "practice",
    ]
    assert extract_source_code_lines(chapter_document) == [
        "void Tick1() { Debug.Log(1); }",
        "void Tick32() { Debug.Log(32); }",
    ]


def test_programming_smoke_rejects_requested_chapter_without_code() -> None:
    try:
        rank_programming_chapter_candidates(
            _document(),
            requested_chapter_id="chapter-prose",
        )
    except SmokeFailure as exc:
        assert exc.code == "production_chapter_has_no_code_source"
    else:
        raise AssertionError("A programming smoke must not use a prose-only chapter")


def test_smoke_audits_dominant_claim_and_full_width_code_export(tmp_path) -> None:
    theme = validate_theme("qizhi-classroom")
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    claim_model = {
        "unit_id": "claim",
        "position": 0,
        "layout": "concept",
        "slide_purpose": "concept",
        "title": "空间换时间",
        "blocks": [{
            "block_id": "claim-block",
            "type": "statement",
            "content": "对象池通过复用实例，以空间换取稳定的运行时间。",
            "items": [],
            "metadata": {},
        }],
        "quality": {
            "resolved_layout": "hero-claim",
            "hero_claim_display_mode": "dominant_canvas",
        },
    }
    code_model = {
        "unit_id": "code",
        "position": 1,
        "layout": "code",
        "slide_purpose": "method",
        "title": "生命周期回调顺序",
        "blocks": [{
            "block_id": "code-block",
            "type": "code",
            "content": "void Awake() {}\nvoid Start() {}",
            "items": [],
            "metadata": {"language": "csharp"},
        }],
        "quality": {
            "resolved_layout": "code",
            "code_region_mode": "full_width",
        },
    }
    claim_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _render_claim_only(claim_slide, SlideSpec.model_validate(claim_model), theme)
    code_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _render_code(code_slide, SlideSpec.model_validate(code_model), theme)
    output = tmp_path / "presentation-modes.pptx"
    presentation.save(output)

    report = _pptx_presentation_mode_audit(
        output,
        [claim_model, code_model],
    )

    assert report == {
        "passed": True,
        "issues": [],
        "hero_claim_page_count": 1,
        "full_width_code_page_count": 1,
    }
