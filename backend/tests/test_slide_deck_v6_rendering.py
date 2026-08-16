from copy import deepcopy
from pathlib import Path
from types import SimpleNamespace

import pytest
from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR
from pydantic import ValidationError

import slide_deck_renderer
from course_document import CourseBlock, CourseDocument, CourseSection, refresh_document_revision
from course_presentation_graph import compile_course_presentation_graph
from representation_compiler import export_slide_deck_pptx
from slide_deck_renderer import audit_exported_pptx
from slide_deck_v6 import (
    SlideStoryBatchV3,
    SlideStoryPageV3,
    SlideStoryPlanV3,
    SlideVisualDecisionV2,
    SlideVisualPlanV2,
    V6BuildError,
    _compile_course_cover_page,
    compile_slide_deck_v6,
)
from slide_deck_v6_renderer import adapt_v6_page_to_slide_spec, export_slide_deck_v6_pptx
from slide_layout_geometry import balanced_two_column_body_metrics
from template_layout_contract import compile_builtin_template_layout_contract_v1


def _code_deck(code_source: str = "function onEvent(value) {\n  return validate(value);\n}"):
    document = refresh_document_revision(
        CourseDocument(
            course_id="generic-code-render-fixture",
            title="Event-driven interaction",
            sections=[CourseSection(section_id="chapter-1", title="Callbacks", position=0)],
            blocks=[
                CourseBlock(
                    block_id="condition",
                    section_id="chapter-1",
                    position=0,
                    role="concept",
                    payload={"markdown": "The handler runs only after the event is emitted."},
                ),
                CourseBlock(
                    block_id="implementation",
                    section_id="chapter-1",
                    position=1,
                    role="example",
                    kind="code",
                    payload={"markdown": code_source},
                ),
                CourseBlock(
                    block_id="result",
                    section_id="chapter-1",
                    position=2,
                    role="feedback",
                    payload={"markdown": "A rejected value remains visible with its validation reason."},
                ),
            ],
        )
    )
    graph = compile_course_presentation_graph(document, teaching_plan={})
    template = compile_builtin_template_layout_contract_v1("qizhi-classroom")
    unit = graph.units[0]
    layout_id = template.layout_id("evidence-code")
    story = SlideStoryPlanV3(
        source_document_revision=document.document_revision,
        template_digest=template.template_digest,
        batches=[SlideStoryBatchV3(
            batch_id="story-1",
            chapter_id="chapter-1",
            provider="fixture-pool",
            model="fixture-story",
            duration_ms=1,
            attempts=1,
            validation_status="passed",
            pages=[SlideStoryPageV3(
                page_id="page-code",
                teaching_unit_id=unit.teaching_unit_id,
                template_layout_id=layout_id,
                title="Connect the event to observable feedback",
                source_block_ids=unit.primary_block_ids,
                page_ordinal=0,
            )],
        )],
    )
    visual = SlideVisualPlanV2(
        source_document_revision=document.document_revision,
        template_digest=template.template_digest,
        decisions=[SlideVisualDecisionV2(
            page_id="page-code",
            decision="code",
            source_block_ids=unit.primary_block_ids,
            resolved_template_layout_id=layout_id,
        )],
    )
    return document, compile_slide_deck_v6(document, graph, story, visual, template)


def _practice_code_deck():
    document = refresh_document_revision(CourseDocument(
        course_id="generic-practice-code-render-fixture",
        title="Sensor verification",
        sections=[CourseSection(section_id="practice", title="Verify", position=0)],
        blocks=[CourseBlock(
            block_id="verification-task",
            section_id="practice",
            position=0,
            role="activity",
            payload={"markdown": (
                "1. Capture the sensor reading.\n"
                "2. Compare it with the declared threshold.\n"
                "3. Record the decision and evidence.\n\n"
                "4. Repeat the check with a boundary value.\n"
                "5. Preserve the original reading.\n"
                "6. Explain any rejected result.\n"
                "7. Submit the signed verification record.\n\n"
                "```python\n"
                "def accept(reading, threshold):\n"
                "    return reading <= threshold\n"
                "```"
            )},
        )],
    ))
    graph = compile_course_presentation_graph(document, teaching_plan={})
    template = compile_builtin_template_layout_contract_v1("qizhi-classroom")
    unit = graph.units[0]
    layout_id = template.layout_id("practice-code")
    story = SlideStoryPlanV3(
        source_document_revision=document.document_revision,
        template_digest=template.template_digest,
        batches=[SlideStoryBatchV3(
            batch_id="story-practice-code",
            chapter_id="practice",
            provider="fixture-pool",
            model="fixture-story",
            duration_ms=1,
            attempts=1,
            validation_status="passed",
            pages=[SlideStoryPageV3(
                page_id="practice-code-page",
                teaching_unit_id=unit.teaching_unit_id,
                template_layout_id=layout_id,
                title="Verify the reading before accepting it",
                source_block_ids=unit.primary_block_ids,
                page_ordinal=0,
            )],
        )],
    )
    visual = SlideVisualPlanV2(
        source_document_revision=document.document_revision,
        template_digest=template.template_digest,
        decisions=[SlideVisualDecisionV2(
            page_id="practice-code-page",
            decision="code",
            source_block_ids=unit.primary_block_ids,
            resolved_template_layout_id=layout_id,
        )],
    )
    return compile_slide_deck_v6(document, graph, story, visual, template)


def _classification_three_deck(items: list[str]):
    document = refresh_document_revision(CourseDocument(
        course_id="generic-classification-capacity-fixture",
        title="Runtime allocation trade-offs",
        sections=[CourseSection(
            section_id="classification",
            title="Classify the trade-offs",
            position=0,
        )],
        blocks=[CourseBlock(
            block_id="classification-source",
            section_id="classification",
            position=0,
            role="concept",
            kind="rich_text",
            payload={"markdown": "\n".join(f"- {item}" for item in items)},
        )],
    ))
    graph = compile_course_presentation_graph(document, teaching_plan={})
    template = compile_builtin_template_layout_contract_v1("qizhi-classroom")
    unit = graph.units[0]
    layout_id = template.layout_id("classification-three")
    story = SlideStoryPlanV3(
        source_document_revision=document.document_revision,
        template_digest=template.template_digest,
        batches=[SlideStoryBatchV3(
            batch_id="story-classification",
            chapter_id="classification",
            provider="fixture-pool",
            model="fixture-story",
            duration_ms=1,
            attempts=1,
            validation_status="passed",
            pages=[SlideStoryPageV3(
                page_id="classification-page",
                teaching_unit_id=unit.teaching_unit_id,
                template_layout_id=layout_id,
                title="以空间换时间",
                source_block_ids=unit.primary_block_ids,
                page_ordinal=0,
            )],
        )],
    )
    visual = SlideVisualPlanV2(
        source_document_revision=document.document_revision,
        template_digest=template.template_digest,
        decisions=[SlideVisualDecisionV2(
            page_id="classification-page",
            decision="text_native",
            source_block_ids=unit.primary_block_ids,
            resolved_template_layout_id=layout_id,
        )],
    )
    return compile_slide_deck_v6(document, graph, story, visual, template)


def _practice_long_table_deck():
    table = "\n".join([
        "| Check | Input | Expected result | Evidence | Repair |",
        "| --- | --- | --- | --- | --- |",
        (
            "| Verify the integration boundary | Preserve the original editor and "
            "runtime settings before the check begins | The external editor opens "
            "the exact project without changing its source state | Record the project, "
            "editor version, runtime log, and reviewer identity in one traceable result | "
            "Restore every missing field from the signed source record before approval |"
        ),
    ])
    document = refresh_document_revision(CourseDocument(
        course_id="generic-practice-table-capacity-fixture",
        title="Integration verification",
        sections=[CourseSection(
            section_id="practice-table",
            title="Verify the integration",
            position=0,
        )],
        blocks=[
            CourseBlock(
                block_id="practice-action",
                section_id="practice-table",
                position=0,
                role="activity",
                payload={"markdown": "1. Run the integration check and preserve the original log."},
            ),
            CourseBlock(
                block_id="practice-evidence",
                section_id="practice-table",
                position=1,
                role="feedback",
                kind="review_checkpoint",
                payload={"markdown": table},
            ),
        ],
    ))
    graph = compile_course_presentation_graph(document, teaching_plan={})
    template = compile_builtin_template_layout_contract_v1("qizhi-classroom")
    unit = graph.units[0]
    layout_id = template.layout_id("practice-table")
    story = SlideStoryPlanV3(
        source_document_revision=document.document_revision,
        template_digest=template.template_digest,
        batches=[SlideStoryBatchV3(
            batch_id="story-practice-table",
            chapter_id="practice-table",
            provider="fixture-pool",
            model="fixture-story",
            duration_ms=1,
            attempts=1,
            validation_status="passed",
            pages=[SlideStoryPageV3(
                page_id="practice-table-page",
                teaching_unit_id=unit.teaching_unit_id,
                template_layout_id=layout_id,
                title="Verify the integration boundary",
                source_block_ids=unit.primary_block_ids,
                page_ordinal=0,
            )],
        )],
    )
    visual = SlideVisualPlanV2(
        source_document_revision=document.document_revision,
        template_digest=template.template_digest,
        decisions=[SlideVisualDecisionV2(
            page_id="practice-table-page",
            decision="table",
            source_block_ids=unit.primary_block_ids,
            resolved_template_layout_id=layout_id,
        )],
    )
    return compile_slide_deck_v6(document, graph, story, visual, template)


def _dense_table_deck():
    table = "\n".join([
        "| Check | Required evidence | Result |",
        "| --- | --- | --- |",
        *(
            f"| Field item {index} | Record the observation condition and supporting evidence {index} | Verified |"
            for index in range(1, 9)
        ),
    ])
    document = refresh_document_revision(
        CourseDocument(
            course_id="generic-field-table-render-fixture",
            title="Field evidence review",
            sections=[
                CourseSection(
                    section_id="field-review",
                    title="Review the evidence",
                    position=0,
                )
            ],
            blocks=[
                CourseBlock(
                    block_id="interpretation",
                    section_id="field-review",
                    position=0,
                    role="activity",
                    payload={
                        "markdown": (
                            "Compare every recorded condition with the required evidence, "
                            "then identify the first unsupported observation before publishing "
                            "the field result. "
                        ) * 3,
                    },
                ),
                CourseBlock(
                    block_id="evidence-table",
                    section_id="field-review",
                    position=1,
                    kind="review_checkpoint",
                    role="feedback",
                    payload={"markdown": table},
                ),
            ],
        )
    )
    graph = compile_course_presentation_graph(document, teaching_plan={})
    template = compile_builtin_template_layout_contract_v1("qizhi-classroom")
    unit = graph.units[0]
    layout_id = template.layout_id("evidence-table")
    story = SlideStoryPlanV3(
        source_document_revision=document.document_revision,
        template_digest=template.template_digest,
        batches=[
            SlideStoryBatchV3(
                batch_id="story-table",
                chapter_id="field-review",
                provider="fixture-pool",
                model="fixture-story",
                duration_ms=1,
                attempts=1,
                validation_status="passed",
                pages=[
                    SlideStoryPageV3(
                        page_id="field-table-page",
                        teaching_unit_id=unit.teaching_unit_id,
                        template_layout_id=layout_id,
                        title="Review the evidence",
                        source_block_ids=unit.primary_block_ids,
                        page_ordinal=0,
                    )
                ],
            )
        ],
    )
    visual = SlideVisualPlanV2(
        source_document_revision=document.document_revision,
        template_digest=template.template_digest,
        decisions=[
            SlideVisualDecisionV2(
                page_id="field-table-page",
                decision="table",
                source_block_ids=unit.primary_block_ids,
                resolved_template_layout_id=layout_id,
            )
        ],
    )
    return compile_slide_deck_v6(document, graph, story, visual, template)


def _wide_markdown_table_deck():
    table = "\n".join([
        "| Stage | Standard | Evidence | Basis | Repair |",
        "| :--- | :--- | :--- | :--- :--- |",
        (
            "| **Observe** | Record the `site` \\| time, weather, and observer before "
            "sampling begins | The log preserves the original field condition | "
            "A stable context makes later comparisons meaningful | **Error**: context "
            "is missing.<br>**Repair**: restore it from the signed field note |"
        ),
        (
            "| **Compare** | Check every observation against the declared criterion | "
            "The first mismatch remains visible | Evidence must precede interpretation | "
            "**Error**: a conclusion replaces the observation.<br>**Repair**: separate them |"
        ),
    ])
    support = (
        "The full audit record remains available for traceability and later review. "
        "Compare each field observation with its declared evidence before publishing."
    )
    summary = "Compare each field observation with its declared evidence before publishing."
    document = refresh_document_revision(CourseDocument(
        course_id="generic-field-audit",
        title="Field audit",
        sections=[CourseSection(section_id="audit", title="Audit", position=0)],
        blocks=[
            CourseBlock(
                block_id="audit-table",
                section_id="audit",
                position=0,
                role="feedback",
                kind="table",
                payload={"markdown": table},
            ),
            CourseBlock(
                block_id="audit-interpretation",
                section_id="audit",
                position=1,
                role="reasoning",
                payload={"markdown": support},
            ),
        ],
    ))
    graph = compile_course_presentation_graph(document, teaching_plan={})
    template = compile_builtin_template_layout_contract_v1("qizhi-classroom")
    layout_id = template.layout_id("evidence-table")
    unit = graph.units[0]
    story = SlideStoryPlanV3(
        source_document_revision=document.document_revision,
        template_digest=template.template_digest,
        batches=[SlideStoryBatchV3(
            batch_id="story-wide-table",
            chapter_id="audit",
            provider="fixture-pool",
            model="fixture-story",
            duration_ms=1,
            attempts=1,
            validation_status="passed",
            pages=[SlideStoryPageV3(
                page_id="wide-table-page",
                teaching_unit_id=unit.teaching_unit_id,
                template_layout_id=layout_id,
                title="Compare the field evidence",
                summary=summary,
                source_block_ids=unit.primary_block_ids,
                page_ordinal=0,
            )],
        )],
    )
    visual = SlideVisualPlanV2(
        source_document_revision=document.document_revision,
        template_digest=template.template_digest,
        decisions=[SlideVisualDecisionV2(
            page_id="wide-table-page",
            decision="table",
            source_block_ids=unit.primary_block_ids,
            resolved_template_layout_id=layout_id,
        )],
    )
    return (
        compile_slide_deck_v6(document, graph, story, visual, template),
        template,
        summary,
    )


def _long_cell_table_deck():
    table = "\n".join([
        "| 观察现象 | 判断依据 | 修正动作 |",
        "| --- | --- | --- |",
        (
            "| The declared field observation and its acceptance criterion could not be reconciled | 原始记录必须同时保留地点、时间、天气、观察者和采样批次，"
            "才能支持后续复核 | 从签字确认的现场记录恢复全部环境字段后再发布结论 |"
        ),
        (
            "| 对照结论缺少证据 | 每项结论都必须指向对应观察条件、验收标准和原始测量值，"
            "不能用解释代替证据 | 分离观察事实与研究者解释并重新执行逐项对照 |"
        ),
        (
            "| 审核过程无法追溯 | 审核者身份、证据修订号和最终决定必须在同一记录中保持可见，"
            "以便复查 | 重新打开审核并补齐缺失的来源信息后才能批准 |"
        ),
        (
            "| 重复测量结果冲突 | 多次测量必须使用相同单位、采样窗口和校准基准，"
            "否则不能直接比较 | 统一记录口径并在审核轨迹中解释每个排除值 |"
        ),
    ])
    document = refresh_document_revision(CourseDocument(
        course_id="generic-field-review-with-long-table-cells",
        title="Field evidence review",
        sections=[CourseSection(section_id="review", title="Review", position=0)],
        blocks=[CourseBlock(
            block_id="review-table",
            section_id="review",
            position=0,
            role="feedback",
            kind="review_checkpoint",
            payload={"markdown": table},
        )],
    ))
    graph = compile_course_presentation_graph(document, teaching_plan={})
    template = compile_builtin_template_layout_contract_v1("qizhi-classroom")
    layout_id = template.layout_id("evidence-table")
    unit = graph.units[0]
    story = SlideStoryPlanV3(
        source_document_revision=document.document_revision,
        template_digest=template.template_digest,
        batches=[SlideStoryBatchV3(
            batch_id="story-long-table",
            chapter_id="review",
            provider="fixture-pool",
            model="fixture-story",
            duration_ms=1,
            attempts=1,
            validation_status="passed",
            pages=[SlideStoryPageV3(
                page_id="long-table-page",
                teaching_unit_id=unit.teaching_unit_id,
                template_layout_id=layout_id,
                title="对照结论缺少证据",
                summary=(
                    "核对每项观察的环境条件、原始证据和验收标准；若结论缺少来源、审核轨迹"
                    "或统一测量口径，应先补齐记录并重新执行逐项对照，再决定是否发布。"
                ),
                source_block_ids=unit.primary_block_ids,
                page_ordinal=0,
            )],
        )],
    )
    visual = SlideVisualPlanV2(
        source_document_revision=document.document_revision,
        template_digest=template.template_digest,
        decisions=[SlideVisualDecisionV2(
            page_id="long-table-page",
            decision="table",
            source_block_ids=unit.primary_block_ids,
            resolved_template_layout_id=layout_id,
        )],
    )
    return compile_slide_deck_v6(document, graph, story, visual, template)


def _chapter_entry_at_contract_capacity_deck():
    driving_question = (
        "开展实地观察前，怎样确认地点、时间、天气、观察者与采样批次均已记录，"
        "并且每项结论都能追溯到原始证据、验收标准和审核修订，从而避免用解释替代事实？"
        "若任一字段缺失，应先停止发布并返回现场记录补齐来源，之后重新核对。"
    )
    document = refresh_document_revision(CourseDocument(
        course_id="generic-field-orientation",
        title="Field observation",
        sections=[CourseSection(section_id="orientation", title="Orientation", position=0)],
        blocks=[CourseBlock(
            block_id="orientation-question",
            section_id="orientation",
            position=0,
            role="objective",
            payload={"markdown": driving_question},
        )],
    ))
    graph = compile_course_presentation_graph(document, teaching_plan={})
    template = compile_builtin_template_layout_contract_v1("qizhi-classroom")
    layout_id = template.layout_id("chapter-entry")
    unit = graph.units[0]
    story = SlideStoryPlanV3(
        source_document_revision=document.document_revision,
        template_digest=template.template_digest,
        batches=[SlideStoryBatchV3(
            batch_id="story-orientation",
            chapter_id="orientation",
            provider="fixture-pool",
            model="fixture-story",
            duration_ms=1,
            attempts=1,
            validation_status="passed",
            pages=[SlideStoryPageV3(
                page_id="orientation-page",
                teaching_unit_id=unit.teaching_unit_id,
                template_layout_id=layout_id,
                title="开展实地观察前",
                source_block_ids=unit.primary_block_ids,
                page_ordinal=0,
            )],
        )],
    )
    visual = SlideVisualPlanV2(
        source_document_revision=document.document_revision,
        template_digest=template.template_digest,
        decisions=[SlideVisualDecisionV2(
            page_id="orientation-page",
            decision="text_native",
            source_block_ids=unit.primary_block_ids,
            resolved_template_layout_id=layout_id,
            degraded=True,
            degradation_reason="visual_text_native",
        )],
    )
    return compile_slide_deck_v6(document, graph, story, visual, template)


def _ordered_step_deck():
    document = refresh_document_revision(CourseDocument(
        course_id="generic-lab-transfer",
        title="Laboratory transfer",
        sections=[CourseSection(
            section_id="transfer",
            title="Transfer the specimen",
            position=0,
        )],
        blocks=[CourseBlock(
            block_id="transfer-steps",
            section_id="transfer",
            position=0,
            role="activity",
            payload={"markdown": (
                "Complete the transfer in order:\n\n"
                "1. **Verify the specimen**\n"
                "   - Match the identifier to the record.\n"
                "2. **Close the container**\n"
                "   - Confirm the seal is intact.\n"
                "3. **Record the handoff**\n"
                "   - Capture the receiver name."
            )},
        )],
    ))
    graph = compile_course_presentation_graph(document, teaching_plan={})
    template = compile_builtin_template_layout_contract_v1("qizhi-classroom")
    unit = graph.units[0]
    layout_id = template.layout_id("practice-prompt")
    story = SlideStoryPlanV3(
        source_document_revision=document.document_revision,
        template_digest=template.template_digest,
        batches=[SlideStoryBatchV3(
            batch_id="story-transfer",
            chapter_id="transfer",
            provider="fixture-pool",
            model="fixture-story",
            duration_ms=1,
            attempts=1,
            validation_status="passed",
            pages=[SlideStoryPageV3(
                page_id="transfer-page",
                teaching_unit_id=unit.teaching_unit_id,
                template_layout_id=layout_id,
                title="Transfer the specimen",
                source_block_ids=unit.primary_block_ids,
                page_ordinal=0,
            )],
        )],
    )
    visual = SlideVisualPlanV2(
        source_document_revision=document.document_revision,
        template_digest=template.template_digest,
        decisions=[SlideVisualDecisionV2(
            page_id="transfer-page",
            decision="text_native",
            source_block_ids=unit.primary_block_ids,
            resolved_template_layout_id=layout_id,
        )],
    )
    return compile_slide_deck_v6(document, graph, story, visual, template)


def test_v6_ordered_steps_render_as_numbered_lines_in_pptx(tmp_path: Path) -> None:
    deck = _ordered_step_deck()
    spec = adapt_v6_page_to_slide_spec(deck.pages[0])
    process = next(block for block in spec.blocks if block.type == "process")

    assert spec.quality["resolved_layout"] == "practice-sequence"
    assert process.items == [
        "Verify the specimen: Match the identifier to the record.",
        "Close the container: Confirm the seal is intact.",
        "Record the handoff: Capture the receiver name.",
    ]

    output = export_slide_deck_v6_pptx(deck, tmp_path / "ordered-steps.pptx")
    presentation = Presentation(output)
    visible_text = "\n".join(
        shape.text
        for shape in presentation.slides[0].shapes
        if hasattr(shape, "text_frame") and shape.text.strip()
    )

    assert "01\nVerify the specimen" in visible_text
    assert "02\nClose the container" in visible_text
    assert "03\nRecord the handoff" in visible_text


def test_v6_ordered_steps_separate_action_titles_from_explanations(tmp_path: Path) -> None:
    deck = _ordered_step_deck()

    output = export_slide_deck_v6_pptx(deck, tmp_path / "ordered-step-path.pptx")
    presentation = Presentation(output)
    text_shapes = [
        shape.text.strip()
        for shape in presentation.slides[0].shapes
        if hasattr(shape, "text_frame") and shape.text.strip()
    ]

    assert "Verify the specimen" in text_shapes
    assert "Match the identifier to the record." in text_shapes
    assert "Close the container" in text_shapes
    assert "Confirm the seal is intact." in text_shapes
    assert "Record the handoff" in text_shapes
    assert "Capture the receiver name." in text_shapes


def test_v6_dense_ordered_steps_pass_the_export_frame_audit(tmp_path: Path) -> None:
    deck = _ordered_step_deck()
    task = next(region for region in deck.pages[0].regions if region.slot_id == "task")
    task.content = "\n".join([
        "采集样本：在指定样区的观测窗口内记录采集时间、地点、观察者和完整的仪器编号。",
        "密封容器：检查盖体是否牢固，并确认防拆标记在整个交接过程里保持清晰可见。",
        "标注证据：逐字复制完整样本标识、采样窗口和现场记录编号，不能省略来源字段。",
        "转移包裹：在放行前记录接收人、交接时间、运输路线和当前保管条件。",
        "异常证据模拟与交接记录修正（可选核验）：把签字回执与原始记录逐项比较，并在发布结论前补齐所有缺失字段。",
    ])

    output = export_slide_deck_v6_pptx(deck, tmp_path / "ordered-step-dense.pptx")
    report = audit_exported_pptx(output, expected_slide_count=1)

    assert report["passed"], report["blockers"]


def test_v6_ordered_steps_are_safe_across_provider_font_metrics(
    tmp_path: Path,
    monkeypatch,
) -> None:
    import slide_deck_renderer as renderer

    class NarrowProviderFont:
        def getlength(self, _character: str) -> float:
            return 0.1

    deck = _ordered_step_deck()
    task = next(
        region for region in deck.pages[0].regions if region.slot_id == "task"
    )
    task.content = "\n".join([
        "Build scene: \u5728 Hierarchy \u7a97\u53e3\u521b\u5efa\u4e00\u4e2a\u540d\u4e3a SpaceStation \u7684\u7a7a\u7269\u4f53\u4f5c\u4e3a\u7236\u8282\u70b9\uff1b"
        "\u5728\u5176\u4e0b\u521b\u5efa ModuleA\u3001ModuleB\u3001ModuleC\uff0c\u5e76\u8bbe\u7f6e\u4f4d\u7f6e\u3001\u65cb\u8f6c\u89d2\u5ea6\u548c\u5c40\u90e8\u5750\u6807\u3002",
        "Write script: \u65b0\u5efa RelativeMover.cs \u5e76\u6302\u8f7d\u5230 ModuleC\uff1b\u8ba9 ModuleC \u5728 SpaceStation "
        "\u7684\u5c40\u90e8\u5750\u6807\u7cfb\u4e2d\u6cbf Y \u8f74\u4e0a\u4e0b\u6d6e\u52a8\uff0c\u5e76\u8bb0\u5f55 Mathf.Sin \u8fd0\u884c\u72b6\u6001\u3002",
        "Validate: \u8fd0\u884c\u573a\u666f\u5e76\u5728 Inspector \u4e2d\u65cb\u8f6c SpaceStation \u7684 Z \u8f74\uff1b\u89c2\u5bdf ModuleC "
        "\u8f68\u8ff9\u662f\u5426\u59cb\u7ec8\u4fdd\u6301\u5728\u7236\u8282\u70b9\u5b9a\u4e49\u7684\u5e73\u9762\u5185\uff0c\u4e14\u6ca1\u6709\u79bb\u5fc3\u504f\u79fb\u3002",
        "Check standard",
        "Complete the extension task",
    ])
    original_font = renderer._audit_font
    monkeypatch.setattr(
        renderer,
        "_audit_font",
        lambda _font_size: NarrowProviderFont(),
    )
    output = export_slide_deck_v6_pptx(
        deck,
        tmp_path / "ordered-step-cross-font.pptx",
    )
    monkeypatch.setattr(renderer, "_audit_font", original_font)

    report = audit_exported_pptx(output, expected_slide_count=1)

    assert report["passed"], report["blockers"]


def test_table_renderer_rejects_row_compression_that_would_clip_content() -> None:
    import slide_deck_renderer as renderer

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])

    with pytest.raises(ValueError, match="table_render_capacity_exceeded"):
        renderer._table(
            slide,
            ["Observation", "Evidence", "Decision"],
            [
                [
                    f"Field sample {index}",
                    "Preserve the source condition, recorded value, and reviewer note.",
                    "Review required",
                ]
                for index in range(1, 7)
            ],
            0.78,
            1.92,
            11.78,
            0.72,
            renderer.THEMES["qizhi-classroom"],
        )


def test_v6_summary_band_yields_safe_space_to_a_tall_source_table(
    tmp_path: Path,
) -> None:
    deck, _template, _summary = _wide_markdown_table_deck()
    first_page = deck.pages[0]
    table_region = next(
        region for region in first_page.regions if region.slot_id == "table"
    )
    table_region.content = "\n".join([
        "| Task | Acceptance | Evidence | Rationale | Errors and correction |",
        "| --- | --- | --- | --- | --- |",
        (
            "| Field review | Preserve the signed observation and declared threshold | "
            "The record contains time, location, measurement, and reviewer | "
            "Source evidence must remain visible before interpretation | "
            "Error: a conclusion replaces the original evidence. Correction: restore "
            "the signed record, compare every declared condition, document each mismatch, "
            "and repeat the review before publishing the result. |"
        ),
    ])

    output = export_slide_deck_v6_pptx(
        deck,
        tmp_path / "adaptive-table-summary-band.pptx",
    )
    report = audit_exported_pptx(output, expected_slide_count=len(deck.pages))

    assert report["passed"], report["blockers"]


def test_v6_materializes_typed_template_slots_without_mixing_code_and_explanation() -> None:
    _document, deck = _code_deck()
    page = deck.pages[0]
    regions = {region.slot_id: region for region in page.regions}

    assert page.resolved_layout.endswith("/evidence-code")
    assert regions["code"].content_kind == "code"
    assert "function onEvent" in regions["code"].content
    assert regions["code"].source_block_ids == ["implementation"]
    assert regions["annotation"].content_kind == "body"
    assert "event is emitted" in regions["annotation"].content
    assert "validation reason" in regions["annotation"].content
    assert set(regions["annotation"].source_block_ids) == {"condition", "result"}


def test_v6_web_and_pptx_adapters_resolve_the_same_template_page(tmp_path: Path) -> None:
    document, deck = _code_deck()
    page = deck.pages[0]
    renderer_slide = adapt_v6_page_to_slide_spec(page)

    assert renderer_slide.quality["v6_template_layout_id"] == page.resolved_layout
    assert renderer_slide.quality["v6_layout_slug"] == "evidence-code"
    assert renderer_slide.quality["resolved_layout"] == "code"
    assert renderer_slide.source_block_ids == page.source_block_ids

    output = export_slide_deck_v6_pptx(
        deck.model_dump(mode="json"),
        tmp_path / "v6-code.pptx",
    )
    presentation = Presentation(output)
    assert len(presentation.slides) == 1
    assert "function onEvent" in "\n".join(
        shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text")
    )
    notes = presentation.slides[0].notes_slide.notes_text_frame.text
    assert document.document_revision in notes
    assert "The handler runs only after the event is emitted." in notes
    assert "A rejected value remains visible" in notes


def test_published_continuation_titles_hide_legacy_pagination_suffixes() -> None:
    deck = _dense_table_deck()
    page = deck.pages[0].model_copy(update={
        "title": "Review the evidence (1/3)",
        "continuation_index": 1,
        "continuation_count": 3,
    })

    adapted = adapt_v6_page_to_slide_spec(page)

    assert adapted.title == "Review the evidence"


def test_course_cover_adapter_keeps_the_title_page_minimal_and_internal_labels_hidden() -> None:
    document, _deck = _code_deck()
    template = compile_builtin_template_layout_contract_v1("qizhi-classroom")
    page = _compile_course_cover_page(document, template)

    adapted = adapt_v6_page_to_slide_spec(page)

    assert adapted.eyebrow == ""
    assert adapted.subtitle == ""
    assert adapted.quality["audience_label_policy"] == "source_only"
    assert [(region.slot_id, region.content) for region in page.regions] == [
        ("title", document.title),
    ]


def test_course_cover_title_only_contract_survives_pptx_frame_audit(tmp_path: Path) -> None:
    document, deck = _code_deck()
    template = compile_builtin_template_layout_contract_v1("qizhi-classroom")
    long_objective = (
        "记录湿地观察地点时间天气仪器校准采样窗口证据负责人验收标准异常现象"
        "复核结论推导依据修正方案跟进责任人并保留签字原始记录供独立审查追溯"
        "同时区分事实观察与解释判断确保每项结论都能回到冻结来源"
    )
    cover = _compile_course_cover_page(
        document.model_copy(update={
            "title": "A field guide to reproducible environmental evidence",
            "sections": [
                section.model_copy(update={
                    "learning_objective": long_objective,
                })
                for section in document.sections
            ],
        }),
        template,
    )
    cover_deck = deck.model_copy(update={"pages": [cover]})
    output = export_slide_deck_v6_pptx(
        cover_deck.model_dump(mode="json"),
        tmp_path / "v6-cover-capacity.pptx",
    )

    assert [(region.slot_id, region.content) for region in cover.regions] == [
        ("title", cover.title),
    ]
    audit = audit_exported_pptx(output, expected_slide_count=1)
    presentation = Presentation(output)
    visible_text = [
        shape.text.strip()
        for shape in presentation.slides[0].shapes
        if hasattr(shape, "text") and shape.text.strip()
    ]

    assert "COURSE DECK" not in visible_text
    assert "课堂演示" not in visible_text
    assert long_objective not in visible_text
    assert not [
        issue
        for issue in audit["issues"]
        if issue.get("code") in {
            "exported_text_frame_overflow",
            "exported_title_unexpected_wrap",
        }
    ]


def test_evidence_code_contract_capacity_survives_pptx_frame_audit(tmp_path: Path) -> None:
    template = compile_builtin_template_layout_contract_v1("qizhi-classroom")
    layout = template.get_layout(template.layout_id("evidence-code"))
    assert layout is not None
    slots = {slot.slot_id: slot for slot in layout.slots}
    code_slot = slots["code"]
    annotation_slot = slots["annotation"]
    code_line_width = max(1, code_slot.max_chars // code_slot.max_lines)
    wrapped_line_capacity = code_line_width * code_slot.max_lines
    wide_prefix = 'const label = "'
    code_samples = {
        "logical-lines": "\n".join(
            f"stage_{index:02d}: " + "validate(input);".ljust(code_line_width - 10, " ")
            for index in range(code_slot.max_lines)
        )[: code_slot.max_chars],
        "wide-literal": wide_prefix + "状态" * (
            (wrapped_line_capacity - len(wide_prefix) - 2) // 4
        ) + '";',
    }
    for sample_name, code in code_samples.items():
        _document, deck = _code_deck()
        regions = {region.slot_id: region for region in deck.pages[0].regions}
        regions["code"].content = code
        regions["annotation"].content = "验证输入事件、保留结果并说明失败边界。" * 20
        regions["annotation"].content = regions["annotation"].content[: annotation_slot.max_chars]

        output = export_slide_deck_v6_pptx(
            deck,
            tmp_path / f"v6-code-capacity-{sample_name}.pptx",
        )
        report = audit_exported_pptx(output, expected_slide_count=1)

        assert report["passed"], (sample_name, report["blockers"])


def test_content_stack_contract_capacity_survives_pptx_frame_audit(
    tmp_path: Path,
) -> None:
    deck = _dense_table_deck()
    template = compile_builtin_template_layout_contract_v1("qizhi-classroom")
    layout = template.get_layout(template.layout_id("content-stack"))
    assert layout is not None
    body_slot = next(slot for slot in layout.slots if slot.slot_id == "body")
    support_page = next(
        page for page in deck.pages if page.resolved_layout.endswith("/content-stack")
    )
    body = next(region for region in support_page.regions if region.slot_id == "body")
    english = (
        "Trace each observation to its signed source record, preserve the declared "
        "condition, separate evidence from interpretation, and document every repair. "
    ) * 10
    visible_chars = body_slot.max_chars - 6
    english = english[:visible_chars]
    split_size = visible_chars // 4
    english = "\n\n".join(
        english[index : index + split_size]
        for index in range(0, len(english), split_size)
    )[: body_slot.max_chars]
    chinese = (
        "核对观察条件、原始记录、推导依据与异常修订，确认每个结论都能回溯到签字证据；"
        "同时区分事实观察与解释判断，确保发布前保留完整来源。"
    ) * 20
    samples = {
        "character-limit": english,
        "wrapped-line-limit": chinese[: body_slot.max_chars],
    }
    for sample_name, sample in samples.items():
        body.content = sample
        one_page_deck = deck.model_copy(update={"pages": [support_page]})

        output = export_slide_deck_v6_pptx(
            one_page_deck,
            tmp_path / f"v6-content-stack-contract-capacity-{sample_name}.pptx",
        )
        report = audit_exported_pptx(output, expected_slide_count=1)

        assert len(body.content) == body_slot.max_chars
        assert report["passed"], (sample_name, report["blockers"])


def test_content_stack_balances_uneven_source_paragraphs_before_export(
    tmp_path: Path,
) -> None:
    deck = _dense_table_deck()
    support_page = next(
        page for page in deck.pages if page.resolved_layout.endswith("/content-stack")
    )
    body = next(region for region in support_page.regions if region.slot_id == "body")
    body.content = "\n\n".join([
        "先确认输入与原始状态。",
        (
            "逐项核对观察条件、签字记录、推导依据与异常修订，确保结论能够回溯到来源；"
            "同时区分事实观察和解释判断，发布前记录每项差异及其修复结果。"
        ) * 8,
        "最后由另一位复核者确认结果。",
    ])[:620]
    one_page_deck = deck.model_copy(update={"pages": [support_page]})

    output = export_slide_deck_v6_pptx(
        one_page_deck,
        tmp_path / "v6-content-stack-uneven-paragraphs.pptx",
    )
    report = audit_exported_pptx(output, expected_slide_count=1)
    presentation = Presentation(output)
    capacity_shapes = [
        shape
        for shape in presentation.slides[0].shapes
        if "[v6-body-capacity=balanced-two-column-body-v1]"
        in str(shape.name or "")
    ]

    assert report["passed"], report["blockers"]
    assert capacity_shapes
    assert all(
        "[v6-body-max-lines=15]" in str(shape.name or "")
        for shape in capacity_shapes
    )


def test_content_stack_renderer_uses_the_shared_balanced_split() -> None:
    source = "\n\n".join([
        "先确认输入、原始状态与冻结来源。",
        (
            "逐项核对中英混排 identifier_with_a_long_name、$P_world = R × P_local$、"
            "硬换行和异常修订，确保所有结论都能回溯到完整证据。"
        ) * 7,
        "最后由另一位复核者确认结果。",
    ])
    metrics = balanced_two_column_body_metrics(source)

    assert metrics["mode"] == "two-column"
    assert slide_deck_renderer._balanced_two_column_body(source) == metrics["segments"]
    assert "".join("".join(metrics["segments"]).split()) == "".join(source.split())


def test_content_stack_renderer_rejects_a_body_outside_the_shared_profile(
    tmp_path: Path,
) -> None:
    deck = _dense_table_deck()
    support_page = next(
        page for page in deck.pages if page.resolved_layout.endswith("/content-stack")
    )
    body = next(region for region in support_page.regions if region.slot_id == "body")
    body.content = "\n".join(
        f"核对项 {index:02d}：保留来源、判断依据与复验结果。"
        for index in range(1, 33)
    )
    metrics = balanced_two_column_body_metrics(body.content)
    one_page_deck = deck.model_copy(update={"pages": [support_page]})

    assert metrics["wrapped_lines"] == [16, 16]
    assert not metrics["fits"]
    with pytest.raises(ValueError, match="template_slot_capacity_exceeded"):
        export_slide_deck_v6_pptx(
            one_page_deck,
            tmp_path / "v6-content-stack-profile-overflow.pptx",
        )


def test_content_stack_splits_newline_dense_short_body_before_export(
    tmp_path: Path,
) -> None:
    deck = _dense_table_deck()
    support_page = next(
        page for page in deck.pages if page.resolved_layout.endswith("/content-stack")
    )
    body = next(region for region in support_page.regions if region.slot_id == "body")
    body.content = "\n".join(
        f"核对项 {index:02d}：保留来源。"
        for index in range(1, 11)
    )
    assert len(body.content) < 180
    one_page_deck = deck.model_copy(update={"pages": [support_page]})

    output = export_slide_deck_v6_pptx(
        one_page_deck,
        tmp_path / "v6-content-stack-newline-dense.pptx",
    )
    report = audit_exported_pptx(output, expected_slide_count=1)

    assert report["passed"], report["blockers"]


def test_unbalanced_classification_items_use_lossless_safe_continuation(
    tmp_path: Path,
) -> None:
    items = [
        "本节的核心理念是以空间换时间，通过预先分配内存来消除运行时的堆内存分配峰值。",
        "栈式复用与内存预分配",
        (
            "在标准 Unity 开发中，频繁调用 GameObject.Instantiate() 会在堆（Heap）上创建新对象，"
            "而 Object.Destroy() 虽然标记对象为待销毁，但实际内存回收由 .NET 垃圾回收器（GC）"
            "在后续周期执行。这种按需分配、延迟回收的模式会导致 GC 压力波动，引发帧率卡顿"
            "（Stuttering）。"
        ),
    ]

    deck = _classification_three_deck(items)
    visible_source = "\n".join(
        region.content
        for page in deck.pages
        for region in page.regions
        if region.content_kind in {"body", "items"}
    )

    assert all(page.resolved_layout.endswith("/content-stack") for page in deck.pages)
    assert deck.quality.source_prose_visible_fidelity == 1.0
    assert all(item in visible_source for item in items)
    output = export_slide_deck_v6_pptx(
        deck,
        tmp_path / "unbalanced-classification-safe-continuation.pptx",
    )
    report = audit_exported_pptx(output, expected_slide_count=len(deck.pages))

    assert report["passed"], report["blockers"]


def test_practice_artifact_allocates_step_height_by_wrapped_line_cost(
    tmp_path: Path,
) -> None:
    deck = _practice_code_deck()
    page = deck.pages[0]
    task = next(region for region in page.regions if region.slot_id == "task")
    task.content = "\n".join([
        "采集输入并记录初始值。",
        "运行一次完整流程。",
        "保存原始日志。",
        (
            "复核关键步骤：逐项比较输入、状态转换、输出、异常记录和签字证据，"
            "确认所有观察都能回到冻结来源，并在发布结论前完成独立复验。"
        ),
        "记录最终通过条件。",
    ])
    one_page_deck = deck.model_copy(update={"pages": [page]})

    output = export_slide_deck_v6_pptx(
        one_page_deck,
        tmp_path / "v6-practice-artifact-weighted-steps.pptx",
    )
    report = audit_exported_pptx(output, expected_slide_count=1)

    assert report["passed"], report["blockers"]


def test_practice_table_separates_an_oversized_row_from_required_steps(
    tmp_path: Path,
) -> None:
    deck = _practice_long_table_deck()

    assert [page.resolved_layout.rsplit("/", 1)[-1] for page in deck.pages] == [
        "practice-prompt",
        "evidence-table",
    ]
    adapted = [adapt_v6_page_to_slide_spec(page) for page in deck.pages]
    assert adapted[1].quality["v6_layout_variant"] == "table-row-detail"

    output = export_slide_deck_v6_pptx(
        deck,
        tmp_path / "v6-practice-table-row-detail.pptx",
    )
    report = audit_exported_pptx(output, expected_slide_count=2)

    assert report["passed"], report["blockers"]


@pytest.mark.parametrize(
    "course_title",
    [
        "Unity 6 环境证据与 reproducible workflow 实战指南",
        "面向复杂现场证据复核与持续改进的完整教学实践指南",
    ],
)
def test_course_cover_preserves_complete_breakable_long_titles(
    course_title: str,
    tmp_path: Path,
) -> None:
    document, deck = _code_deck()
    template = compile_builtin_template_layout_contract_v1("qizhi-classroom")
    cover = _compile_course_cover_page(
        document.model_copy(update={"title": course_title}),
        template,
    )

    assert cover.title == course_title
    assert cover.regions[0].content == course_title
    output = export_slide_deck_v6_pptx(
        deck.model_copy(update={"pages": [cover]}),
        tmp_path / "v6-breakable-long-cover.pptx",
    )
    audit = audit_exported_pptx(output, expected_slide_count=1)
    assert not [
        issue
        for issue in audit["issues"]
        if issue.get("code") in {
            "exported_text_frame_overflow",
            "exported_title_unexpected_wrap",
        }
    ]


def test_course_cover_rejects_a_single_unbreakable_token_that_cannot_fit() -> None:
    document, _deck = _code_deck()
    template = compile_builtin_template_layout_contract_v1("qizhi-classroom")
    unbreakable_title = "BuildArtifactIdentity_" + "x" * 80

    with pytest.raises(V6BuildError, match="template_title_capacity_exceeded"):
        _compile_course_cover_page(
            document.model_copy(update={"title": unbreakable_title}),
            template,
        )


def test_evidence_table_paginates_complete_interpretation_before_table_rows(
    tmp_path: Path,
) -> None:
    deck = _dense_table_deck()
    assert len(deck.pages) == 4
    support_page = deck.pages[0]
    assert support_page.resolved_layout.endswith("/content-stack")
    support_body = next(
        region.content
        for region in support_page.regions
        if region.content_kind == "body"
    )
    assert "…" not in support_body
    assert support_body.count("Compare every recorded condition") == 3
    table_pages = deck.pages[1:]
    split_slide = adapt_v6_page_to_slide_spec(table_pages[0])

    table_row_counts = []
    for page in table_pages:
        table_region = next(region for region in page.regions if region.content_kind == "table")
        table_row_counts.append(
            len([line for line in table_region.content.splitlines() if line.strip()]) - 2
        )

    assert split_slide.quality["v6_layout_variant"] == "table-continuation"
    assert split_slide.quality["v6_artifact_support_mode"] == "full"
    assert split_slide.eyebrow == ""
    assert all(block.title == "" for block in split_slide.blocks)
    for page in table_pages[1:]:
        continuation_slide = adapt_v6_page_to_slide_spec(page)
        assert continuation_slide.quality["v6_layout_variant"] == "table-continuation"
        assert continuation_slide.quality["v6_artifact_support_mode"] == "full"
    assert sum(table_row_counts) == 8

    output = export_slide_deck_v6_pptx(
        deck,
        tmp_path / "v6-dense-table.pptx",
    )
    report = audit_exported_pptx(output, expected_slide_count=len(deck.pages))

    assert report["passed"], report["blockers"]
    presentation = Presentation(output)
    for slide_index, slide in enumerate(presentation.slides):
        visible_text = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        assert "EVIDENCE TABLE" not in visible_text
        assert "INTERPRETATION" not in visible_text
        assert "SUMMARY" not in visible_text
        assert "COURSE" not in visible_text
        assert not any("(1/" in value or "(2/" in value for value in visible_text)
        if slide_index == 0:
            assert not any(shape.has_table for shape in slide.shapes)
            continue
        table = next(shape.table for shape in slide.shapes if shape.has_table)
        assert all(
            cell.vertical_anchor == MSO_ANCHOR.MIDDLE
            for row in table.rows
            for cell in row.cells
        )


def test_wide_table_summary_band_fits_three_lines_at_readable_size(tmp_path: Path) -> None:
    deck, _template, _summary = _wide_markdown_table_deck()
    interpretation = next(
        region
        for region in deck.pages[0].regions
        if region.slot_id == "interpretation"
    )
    interpretation.content = ((
        "核对观察条件、原始记录、推导依据与异常修订，确认每个结论都能回溯到签字证据；"
        "同时保留地点、时间、观察者、仪器与采样窗口，避免把解释误写成事实。"
    ) * 2)[:120]

    output = export_slide_deck_v6_pptx(deck, tmp_path / "wide-table-summary-band.pptx")
    report = audit_exported_pptx(output, expected_slide_count=len(deck.pages))

    assert report["passed"], report["blockers"]


def test_wide_table_summary_band_fits_declared_multilingual_capacity(
    tmp_path: Path,
) -> None:
    deck, _template, _summary = _wide_markdown_table_deck()
    template = compile_builtin_template_layout_contract_v1("qizhi-classroom")
    layout = template.get_layout(template.layout_id("evidence-table"))
    assert layout is not None
    interpretation_slot = next(
        slot for slot in layout.slots if slot.slot_id == "interpretation"
    )
    interpretation = next(
        region
        for region in deck.pages[0].regions
        if region.slot_id == "interpretation"
    )
    interpretation.content = (
        "正确的观察结果应呈现如下特征：初始状态下所有字段使用默认值，并显示稳定结果。\n\n"
        "请对照以下标准自查练习成果：地点、时间、观察者、仪器和采样窗口均已保留。\n\n"
        "场景描述：学习者发现修改采样窗口后结果没有变化，应核对条件、日志与签字证据。"
    )

    output = export_slide_deck_v6_pptx(
        deck,
        tmp_path / "wide-table-multilingual-summary-capacity.pptx",
    )
    report = audit_exported_pptx(output, expected_slide_count=len(deck.pages))

    assert len(interpretation.content) <= interpretation_slot.max_chars
    assert report["passed"], report["blockers"]


def test_wide_markdown_table_uses_llm_summary_and_exports_template_safe_cells(
    tmp_path: Path,
) -> None:
    deck, template, summary = _wide_markdown_table_deck()
    page = deck.pages[0]
    regions = {region.slot_id: region for region in page.regions}

    assert regions["interpretation"].content == summary
    adapted = adapt_v6_page_to_slide_spec(page)
    assert adapted.quality["v6_layout_variant"] == "table-wide-with-summary"
    assert adapted.quality["v6_artifact_support_mode"] == "band"
    assert len(deck.pages) == 3
    continuation = adapt_v6_page_to_slide_spec(deck.pages[1])
    assert continuation.quality["v6_layout_variant"] == "table-row-detail"
    support_page = deck.pages[2]
    assert support_page.resolved_layout.endswith("/content-stack")
    assert any(
        "The full audit record remains available" in region.content
        for region in support_page.regions
    )

    compiled_table_text = "\n".join(
        region.content
        for compiled_page in deck.pages
        for region in compiled_page.regions
        if region.content_kind == "table"
    )
    assert (
        "Record the site \\| time, weather, and observer before sampling begins"
        in compiled_table_text
    )
    assert "restore it from the signed field note" in compiled_table_text
    assert "separate them" in compiled_table_text

    output = export_slide_deck_v6_pptx(deck, tmp_path / "wide-table.pptx")
    report = audit_exported_pptx(output, expected_slide_count=len(deck.pages))
    assert report["passed"], report["blockers"]

    presentation = Presentation(output)
    table_shape = next(shape for shape in presentation.slides[0].shapes if shape.has_table)
    rendered_rows = [[cell.text for cell in row.cells] for row in table_shape.table.rows]
    rendered_text = "\n".join(
        shape.text
        for slide in presentation.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    rendered_text += "\n" + "\n".join(cell for row in rendered_rows for cell in row)
    assert len(rendered_rows) == 2
    assert len(rendered_rows[0]) == 5
    assert "site | time" in rendered_text
    assert not any(marker in rendered_text for marker in ("**", "`", "<br>"))
    assert "Record the site | time, weather, and observer before sampling begins" in rendered_text
    assert "restore it from the signed field note" in rendered_text
    assert "…" not in rendered_text
    assert not any(set(cell.replace(" ", "")) <= {":", "-"} for cell in rendered_rows[1])


def test_three_column_table_with_long_cells_is_split_before_export_overflow(
    tmp_path: Path,
) -> None:
    deck = _long_cell_table_deck()

    assert len(deck.pages) == 2
    assert deck.pages[1].continuation_of_page_id == deck.pages[0].page_id
    adapted = [adapt_v6_page_to_slide_spec(page) for page in deck.pages]
    assert adapted[0].quality["v6_layout_variant"] == "table-wide-with-summary"
    assert adapted[1].quality["v6_layout_variant"] == "table-row-detail"
    table_content = "\n".join(
        region.content
        for page in deck.pages
        for region in page.regions
        if region.content_kind == "table"
    )
    assert "The declared field observation and its acceptance criterion could not be reconciled" in table_content
    assert "原始记录必须同时保留地点、时间、天气、观察者和采样批次，才能支持后续复核" in table_content
    assert "重新打开审核并补齐缺失的来源信息后才能批准" in table_content
    assert "…" not in table_content

    output = export_slide_deck_v6_pptx(deck, tmp_path / "long-cell-table.pptx")
    report = audit_exported_pptx(output, expected_slide_count=len(deck.pages))

    assert report["passed"], report["blockers"]


def test_chapter_entry_renderer_honors_declared_driving_question_capacity(
    tmp_path: Path,
) -> None:
    deck = _chapter_entry_at_contract_capacity_deck()

    output = export_slide_deck_v6_pptx(deck, tmp_path / "chapter-entry-capacity.pptx")
    report = audit_exported_pptx(output, expected_slide_count=1)

    assert report["passed"], report["blockers"]


def test_oversized_code_requires_lossless_pagination_instead_of_an_excerpt() -> None:
    source = "\n".join([
        "using Example.Runtime;",
        "",
        "public class AuditRunner",
        "{",
        "    void FirstCheck()",
        "    {",
        "        VerifyContext();",
        "    }",
        "",
        "    // The second check runs after the first result is visible",
        "    void SecondCheck()",
        "    {",
        "        VerifyEvidence();",
        "    }",
        "}",
    ])
    block = CourseBlock(
        block_id="generic-code",
        section_id="audit",
        position=0,
        role="example",
        kind="code",
        payload={"markdown": source},
    )

    from slide_deck_v6 import _bounded_slot_content

    with pytest.raises(ValueError, match="template_slot_capacity_exceeded"):
        _bounded_slot_content(
            [block],
            slot_kind="code",
            max_chars=300,
            max_items=0,
            max_lines=10,
            max_rows=0,
        )


def test_oversized_code_with_comment_braces_still_requires_lossless_pagination() -> None:
    source = "\n".join([
        "public class Formatter",
        "{",
        "    void Render()",
        "    {",
        '        var label = "{pending}";',
        "        Publish(label); // unmatched } belongs to the comment",
        "    }",
        "",
        "    void Archive()",
        "    {",
        "        Save();",
        "    }",
        "}",
    ])
    block = CourseBlock(
        block_id="generic-brace-code",
        section_id="generic-section",
        position=0,
        role="example",
        kind="code",
        payload={"markdown": source},
    )

    from slide_deck_v6 import _bounded_slot_content

    with pytest.raises(ValueError, match="template_slot_capacity_exceeded"):
        _bounded_slot_content(
            [block],
            slot_kind="code",
            max_chars=220,
            max_items=0,
            max_lines=8,
            max_rows=0,
        )


def test_chapter_entry_title_contract_allows_only_declared_safe_wrapping(
    tmp_path: Path,
) -> None:
    _document, deck = _code_deck()
    template = compile_builtin_template_layout_contract_v1("qizhi-classroom")
    layout = template.get_layout(template.layout_id("chapter-entry"))
    assert layout is not None
    title_slot = next(slot for slot in layout.slots if slot.slot_kind == "title")
    page = deck.pages[0]
    page.title = ("教学单元逻辑顺序与来源证据" * 4)[: title_slot.max_chars]
    page.title_max_lines = title_slot.max_lines
    page.resolved_layout = layout.template_layout_id
    page.visual_decision.resolved_template_layout_id = layout.template_layout_id

    output = export_slide_deck_v6_pptx(deck, tmp_path / "v6-chapter-title.pptx")
    report = audit_exported_pptx(output, expected_slide_count=1)

    assert report["passed"], report["blockers"]
    title_shapes = [
        shape
        for shape in Presentation(output).slides[0].shapes
        if getattr(shape, "has_text_frame", False)
        and str(shape.text or "").strip() == page.title
    ]
    assert len(title_shapes) == 1
    assert f"[v6-title-max-lines={title_slot.max_lines}]" in title_shapes[0].name


def test_official_representation_export_dispatches_v6_without_legacy_schema_coercion(
    tmp_path: Path,
) -> None:
    _document, deck = _code_deck()
    spec = SimpleNamespace(
        representation_type="slide_deck",
        payload={"content": deck.model_dump(mode="json")},
    )

    output = export_slide_deck_pptx(spec, tmp_path / "official-v6.pptx")

    assert output.is_file()
    assert len(Presentation(output).slides) == len(deck.pages)


def test_official_v6_export_accepts_declared_publication_metadata(
    tmp_path: Path,
) -> None:
    _document, deck = _code_deck()
    published = {
        **deck.model_dump(mode="json"),
        "build_signature": {"signature": "slidebuildv6_fixture"},
        "source_contract": {"schema_version": "ppt_source_contract_v2"},
        "course_presentation_graph": {
            "schema_version": "course_presentation_graph_v1",
        },
        "story_plan": {"schema_version": "slide_story_plan_v3"},
        "visual_plan": {"schema_version": "slide_visual_plan_v2"},
        "template_contract": {
            "schema_version": "template_layout_pack_contract_v1",
        },
        "ai_batch_diagnostics": [{
            "schema_version": "ai_batch_diagnostic_v1",
        }],
        "planning_status": {
            "story_ai": {"status": "completed"},
            "visual_ai": {"status": "completed"},
        },
    }
    frozen_publication = deepcopy(published)
    spec = SimpleNamespace(
        representation_type="slide_deck",
        payload={"content": published},
    )

    output = export_slide_deck_pptx(
        spec,
        tmp_path / "published-official-v6.pptx",
    )

    assert output.is_file()
    assert len(Presentation(output).slides) == len(deck.pages)
    assert published == frozen_publication
    assert published["build_signature"] == frozen_publication["build_signature"]
    assert published["source_contract"] == frozen_publication["source_contract"]
    assert published["template_contract"] == frozen_publication["template_contract"]


@pytest.mark.parametrize("unknown_field", ["undeclared_payload", "pagse"])
def test_official_v6_export_still_rejects_unknown_publication_fields(
    tmp_path: Path,
    unknown_field: str,
) -> None:
    _document, deck = _code_deck()
    published = {
        **deck.model_dump(mode="json"),
        unknown_field: {"should": "not pass"},
    }
    spec = SimpleNamespace(
        representation_type="slide_deck",
        payload={"content": published},
    )

    with pytest.raises(ValidationError, match=unknown_field):
        export_slide_deck_pptx(
            spec,
            tmp_path / "invalid-published-v6.pptx",
        )


@pytest.mark.parametrize("required_field", ["pages", "quality"])
def test_official_v6_export_rejects_missing_required_deck_fields(
    tmp_path: Path,
    required_field: str,
) -> None:
    _document, deck = _code_deck()
    published = {
        **deck.model_dump(mode="json"),
        "template_contract": {
            "schema_version": "template_layout_pack_contract_v1",
        },
    }
    published.pop(required_field)
    spec = SimpleNamespace(
        representation_type="slide_deck",
        payload={"content": published},
    )

    with pytest.raises(ValidationError, match=required_field):
        export_slide_deck_pptx(
            spec,
            tmp_path / "incomplete-published-v6.pptx",
        )


@pytest.mark.parametrize(
    ("quality_field", "rejected_value"),
    [
        ("source_artifact_visible_fidelity", 0.0),
        ("source_prose_visible_fidelity", 0.0),
        ("ordered_step_visible_fidelity", 0.0),
        ("generated_ellipsis_free", False),
        ("pagination_within_dynamic_bound", False),
    ],
)
def test_official_v6_export_rejects_each_current_quality_gate(
    tmp_path: Path,
    quality_field: str,
    rejected_value: object,
) -> None:
    _document, deck = _code_deck()
    setattr(deck.quality, quality_field, rejected_value)

    with pytest.raises(slide_deck_renderer.SlideDeckQualityError) as captured:
        export_slide_deck_v6_pptx(
            deck.model_dump(mode="json"),
            tmp_path / f"rejected-{quality_field}.pptx",
        )

    assert f"v6_{quality_field}_failed" in {
        blocker["code"] for blocker in captured.value.report["blockers"]
    }


def test_official_v6_export_preserves_literal_csharp_interpolation(
    tmp_path: Path,
) -> None:
    code_source = 'Debug.Log($"frame={Time.frameCount} phase={phase}");'
    _document, deck = _code_deck(code_source)

    output = export_slide_deck_v6_pptx(
        deck,
        tmp_path / "literal-csharp-interpolation.pptx",
    )
    visible = "\n".join(
        str(shape.text or "")
        for slide in Presentation(output).slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )

    assert code_source in visible


def test_practice_code_layout_exports_numbered_steps_and_readable_code(
    tmp_path: Path,
) -> None:
    deck = _practice_code_deck()
    page = deck.pages[0]

    assert page.resolved_layout.endswith("/practice-code")
    assert {region.content_kind for region in page.regions} == {"steps", "code"}
    renderer_slide = adapt_v6_page_to_slide_spec(page)
    assert renderer_slide.quality["resolved_layout"] == "practice-artifact"
    assert [block.type for block in renderer_slide.blocks] == ["code", "process"]

    output = export_slide_deck_v6_pptx(deck, tmp_path / "practice-code.pptx")
    report = audit_exported_pptx(output, expected_slide_count=1)

    assert report["passed"], report["blockers"]
    visible_text = "\n".join(
        str(shape.text or "")
        for shape in Presentation(output).slides[0].shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert "01" in visible_text
    assert "02" in visible_text
    assert "07" in visible_text
    assert "def accept" in visible_text


def test_long_code_exports_every_source_line_across_content_driven_pages(
    tmp_path: Path,
) -> None:
    code = "\n".join(
        f"step_{index} = observe(source_value_{index})"
        for index in range(70)
    )
    _document, deck = _code_deck(code)

    assert len(deck.pages) > 3
    assert deck.quality.source_artifact_visible_fidelity == 1.0
    output = export_slide_deck_v6_pptx(deck, tmp_path / "long-code.pptx")
    report = audit_exported_pptx(output, expected_slide_count=len(deck.pages))

    assert report["passed"], report["blockers"]
    presentation = Presentation(output)
    visible_text = "\n".join(
        str(shape.text or "")
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )
    assert all(line in visible_text for line in code.splitlines())


def test_pptx_renderer_applies_the_frozen_template_theme_overrides(
    tmp_path: Path,
    monkeypatch,
) -> None:
    import slide_deck_v6_renderer as renderer

    _document, deck = _code_deck()
    deck.template_theme_overrides = {
        "accent": "315E7D",
        "title_font": "Noto Serif SC",
    }
    observed: dict[str, str] = {}
    original = renderer._render_slide

    def capture(slide, unit, page_number, page_count, theme, assets):
        observed.update(theme)
        return original(slide, unit, page_number, page_count, theme, assets)

    monkeypatch.setattr(renderer, "_render_slide", capture)
    renderer.export_slide_deck_v6_pptx(deck, tmp_path / "personal-theme.pptx")

    assert observed["accent"] == "315E7D"
    assert observed["title_font"] == "Noto Serif SC"
