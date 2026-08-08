from __future__ import annotations

import asyncio
from copy import deepcopy

import pytest
from pptx import Presentation

from course_document import document_from_legacy_course
from course_revisions import revision_vector_for_course
from representation_compiler import (
    compile_core_representations,
    rebuild_slide_deck_variant_bundle_safely,
    rebuild_slide_deck_variant_safely,
)
from slide_deck_renderer import export_structured_slide_deck
from slide_deck_v3 import (
    V3_LAYOUTS,
    fragment_course_document,
    split_slide_deck_plan_by_chapter,
)
from slide_deck_v4 import (
    _presentation_quality,
    allocation_from_story_plan_v2,
    build_signature_v4,
    compile_slide_deck_v4,
)
from slide_deck_v5 import allocation_from_story_plan_v5, compact_story_plan_v5
from slide_layout_registry import registry_summary_v2, select_layout_v2
from slide_story_plan import (
    STORY_BEAT_TEXT_CAPACITY,
    SlideStoryPlanPrerequisiteError,
    _claim_for_scene,
    compile_slide_story_plan_v2,
    plan_slide_story_v2,
    resolve_slide_deck_schema,
)
from teaching_representations import TeachingRepresentationRepository


def _course_with_teaching_plan() -> dict:
    return {
        "course_id": "linear-map-v4",
        "course_name": "线性映射的定义与矩阵表示",
        "course_revision": "course-rev-1",
        "knowledge_base_revision": "kb-rev-1",
        "coherence_contract_revision": "coherence-rev-1",
        "course_knowledge_base": {
            "schema_version": "course_knowledge_base_v2",
            "revision_id": "kb-rev-1",
            "lifecycle_status": "active",
        },
        "course_coherence_contract": {
            "schema_version": "course_coherence_v2",
            "revision_id": "coherence-rev-1",
            "status": "active",
            "quality_report": {"passed": True},
        },
        "generation_stage_artifacts": {
            "course_teaching_plan": {
                "status": "completed",
                "section_count": 1,
            },
        },
        "course_teaching_plan": {
            "revision_id": "teaching-plan-rev-1",
            "sections": [
                {
                    "node_id": "chapter-linear-map",
                    "key_points": ["线性映射同时保持加法与数乘"],
                    "reused_knowledge_names": ["向量空间"],
                    "knowledge_structure": [
                        {
                            "concept_group": "线性映射",
                            "knowledge_points": [
                                {
                                    "knowledge_id": "kp-linear-map",
                                    "name": "线性映射",
                                    "statement": "线性映射同时保持向量加法与数乘。",
                                    "knowledge_type": "concept",
                                    "conditions": ["定义域和值域均为向量空间"],
                                    "boundaries": ["必须同时满足两条保持性质"],
                                    "capability": "判断一个映射是否为线性映射",
                                    "capability_points": [
                                        {
                                            "capability_id": "cap-check-linear",
                                            "statement": "能按两条保持性质完成判断",
                                        },
                                    ],
                                    "misconceptions": [
                                        {
                                            "misconception_id": "mis-one-rule",
                                            "statement": "只验证加法保持就判定线性",
                                            "repair": "必须同时验证数乘保持",
                                        },
                                    ],
                                    "mastery_criteria": [
                                        {
                                            "criterion_id": "mastery-linear",
                                            "statement": "能解释并验证两条保持性质",
                                        },
                                    ],
                                    "prerequisite_names": ["向量空间"],
                                },
                            ],
                        },
                    ],
                    "knowledge_relations": [
                        {
                            "source_name": "线性映射",
                            "target_name": "矩阵表示",
                            "relation_type": "representation",
                            "derivation_steps": ["选定基", "计算基向量的像", "按列组成矩阵"],
                        },
                    ],
                    "teaching_modules": [
                        {
                            "module_id": "module-concept",
                            "teaching_purpose": "建立线性映射的正式定义与条件边界",
                            "knowledge_names": ["线性映射"],
                        },
                        {
                            "module_id": "module-example",
                            "teaching_purpose": "通过完整例题练习判断方法",
                            "knowledge_names": ["线性映射"],
                        },
                        {
                            "module_id": "module-method",
                            "teaching_purpose": "把线性映射判定转化为可执行的方法步骤",
                            "knowledge_names": ["线性映射"],
                        },
                        {
                            "module_id": "module-practice",
                            "teaching_purpose": "检查是否达到掌握标准",
                            "knowledge_names": ["线性映射"],
                        },
                    ],
                },
            ],
        },
        "nodes": [
            {
                "node_id": "chapter-linear-map",
                "parent_node_id": "root",
                "node_name": "线性映射的定义与矩阵表示",
                "node_level": 1,
                "learning_objective": "理解线性映射的定义，并能判断给定映射是否线性。",
                "objective_id": "objective-linear-map",
                "content_blocks": [
                    {
                        "block_id": "block-entry",
                        "title": "进入问题",
                        "content": "什么样的映射不会破坏向量空间中的线性结构？",
                        "metadata": {
                            "role": "orientation",
                            "module_id": "module-concept",
                        },
                    },
                    {
                        "block_id": "block-concept",
                        "title": "正式定义",
                        "content": (
                            "设 V 与 W 为向量空间。映射 T:V→W 若对任意 u,v∈V 与标量 c，"
                            "满足 T(u+v)=T(u)+T(v) 且 T(cu)=cT(u)，则称 T 为线性映射。"
                        ),
                        "metadata": {
                            "role": "concept",
                            "module_id": "module-concept",
                        },
                    },
                    {
                        "block_id": "block-method",
                        "title": "判断方法",
                        "content": "先验证加法保持，再验证数乘保持；两项都成立才能判定为线性映射。",
                        "metadata": {
                            "role": "method",
                            "module_id": "module-method",
                        },
                    },
                    {
                        "block_id": "block-example-prompt",
                        "title": "例题",
                        "content": "判断 T(x,y)=(x+y,y) 是否为线性映射。",
                        "metadata": {
                            "role": "example",
                            "module_id": "module-example",
                        },
                    },
                    {
                        "block_id": "block-example-solution",
                        "title": "例题解答",
                        "content": "依次验证加法保持与数乘保持，两项均成立，因此该映射是线性映射。",
                        "metadata": {
                            "role": "answer",
                            "module_id": "module-example",
                        },
                    },
                    {
                        "block_id": "block-practice",
                        "title": "课堂练习",
                        "content": "判断 S(x,y)=(x+1,y) 是否为线性映射。",
                        "metadata": {
                            "role": "checkpoint",
                            "module_id": "module-practice",
                        },
                    },
                    {
                        "block_id": "block-feedback",
                        "title": "练习反馈",
                        "content": "S(0,0)≠(0,0)，所以 S 不是线性映射。",
                        "metadata": {
                            "role": "feedback",
                            "module_id": "module-practice",
                        },
                    },
                    {
                        "block_id": "block-misconception",
                        "title": "常见误区",
                        "content": "只验证加法保持不够，还必须验证数乘保持。",
                        "metadata": {
                            "role": "misconception",
                            "module_id": "module-concept",
                        },
                    },
                    {
                        "block_id": "block-recap",
                        "title": "本章总结",
                        "content": "线性映射的判断必须同时检查加法保持与数乘保持。",
                        "metadata": {
                            "role": "summary",
                            "module_id": "module-practice",
                        },
                    },
                ],
            },
        ],
    }


def test_story_plan_uses_official_course_logic_and_closes_the_chapter() -> None:
    course = _course_with_teaching_plan()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)

    plan = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )

    assert plan.schema_version == "slide_story_plan_v2"
    assert plan.source_revisions.teaching_plan_revision == "teaching-plan-rev-1"
    assert plan.chapters[0].driving_question == (
        "理解线性映射的定义，并能判断给定映射是否线性。"
    )
    scene_kinds = [episode.scene_kind for episode in plan.chapters[0].episodes]
    assert scene_kinds[0] == "chapter_entry"
    assert "concept" in scene_kinds
    assert "method" in scene_kinds
    assert "worked_example" in scene_kinds
    assert "practice_feedback" in scene_kinds
    assert "misconception" in scene_kinds
    assert scene_kinds[-1] == "chapter_recap"
    assert plan.chapters[0].owned_knowledge_ids == ["kp-linear-map"]
    assert plan.chapters[0].prerequisite_knowledge_names == ["向量空间"]


def test_example_role_in_an_application_module_stays_a_parallel_application() -> None:
    course = _course_with_teaching_plan()
    section_plan = course["course_teaching_plan"]["sections"][0]
    section_plan["teaching_modules"].append({
        "module_id": "module-application",
        "teaching_purpose": "比较多个行业应用情境",
        "knowledge_names": ["线性映射"],
    })
    course["nodes"][0]["content_blocks"].append({
        "block_id": "block-application-examples",
        "title": "行业应用",
        "content": "计算机图形、信号处理和数据降维都使用线性映射。",
        "metadata": {
            "role": "example",
            "module_id": "module-application",
        },
    })
    document = document_from_legacy_course(course)

    plan = compile_slide_story_plan_v2(
        document,
        course,
        fragment_course_document(document),
        mode="teaching",
        theme="qizhi-classroom",
    )

    application = next(
        episode for episode in plan.chapters[0].episodes
        if episode.scene_kind == "application"
    )
    assert application.beats[0].beat_role == "mapping"


def test_single_family_scene_remains_selectable_after_rhythm_limit() -> None:
    selection = select_layout_v2(
        scene_kind="chapter_entry",
        evidence_kinds=["text"],
        character_count=0,
        item_count=0,
        theme="qizhi-classroom",
        recent_layout_families=["hero", "hero"],
    )

    assert selection.layout_id == "chapter-question"
    assert selection.capacity_passed is True
    assert selection.rhythm_score < 1


def test_content_scene_claim_is_derived_from_the_local_beat() -> None:
    course = _course_with_teaching_plan()
    document = document_from_legacy_course(course)
    source_fragment = next(
        fragment
        for fragment in fragment_course_document(document)
        if fragment.block_id == "block-concept"
    )
    local_heading = source_fragment.model_copy(update={
        "kind": "heading",
        "text": "Local continuation: scalar multiplication",
    })

    claim = _claim_for_scene(
        scene="concept",
        chapter=document.sections[0],
        section_plan={
            "knowledge_structure": [{
                "knowledge_points": [{
                    "knowledge_id": "kp-global",
                    "statement": "Global chapter claim",
                }],
            }],
        },
        fragments=[local_heading],
        module={"module_id": "module-concept", "teaching_purpose": "Global purpose"},
    )

    assert claim.kind == "source_heading"
    assert claim.text == "Local continuation: scalar multiplication"
    assert claim.fragment_id == local_heading.fragment_id


def test_example_and_practice_answers_are_revealed_after_the_prompt() -> None:
    course = _course_with_teaching_plan()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    fragment_by_block = {
        fragment.block_id: fragment.fragment_id
        for fragment in fragments
        if fragment.kind != "title"
    }

    plan = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    episodes = plan.chapters[0].episodes
    example = next(item for item in episodes if item.scene_kind == "worked_example")
    practice = next(item for item in episodes if item.scene_kind == "practice_feedback")

    assert example.beats[0].beat_role == "prompt"
    assert fragment_by_block["block-example-solution"] not in example.beats[0].fragment_ids
    assert any(
        fragment_by_block["block-example-solution"] in beat.fragment_ids
        for beat in example.beats[1:]
    )
    assert practice.beats[0].beat_role == "prompt"
    assert fragment_by_block["block-feedback"] not in practice.beats[0].fragment_ids
    assert any(
        fragment_by_block["block-feedback"] in beat.fragment_ids
        for beat in practice.beats[1:]
    )


def test_layout_selection_is_scene_aware_capacity_safe_and_deterministic() -> None:
    first = select_layout_v2(
        scene_kind="worked_example",
        evidence_kinds=["formula"],
        character_count=520,
        item_count=4,
        theme="qizhi-classroom",
        recent_layout_families=["split"],
    )
    second = select_layout_v2(
        scene_kind="worked_example",
        evidence_kinds=["formula"],
        character_count=520,
        item_count=4,
        theme="qizhi-classroom",
        recent_layout_families=["split"],
    )

    assert first.layout_id == second.layout_id
    assert first.scene_match_score > 0
    assert first.capacity_passed is True
    assert first.layout_family != "split"


def test_new_programming_course_long_code_is_partitioned_before_story_layout_selection(
) -> None:
    course = deepcopy(_course_with_teaching_plan())
    course["subject_pedagogy_profile"] = {
        "primary_mode": "programming_engineering",
        "confidence": 0.96,
        "classification_source": "course_generation_v16",
    }
    code_lines = [
        (
            f"public void Tick{index}(GameObject player) {{ "
            f"player.transform.position += velocity{index} * Time.deltaTime; }}"
        )
        for index in range(1, 33)
    ]
    concept = next(
        block
        for block in course["nodes"][0]["content_blocks"]
        if block["block_id"] == "block-concept"
    )
    concept["content"] = "```csharp\n" + "\n".join(code_lines) + "\n```"
    document = document_from_legacy_course(course)

    fragments = fragment_course_document(document)
    code_fragments = [
        fragment
        for fragment in fragments
        if fragment.block_id == "block-concept" and fragment.kind == "code"
    ]

    assert len(code_fragments) > 1
    assert all(
        len(fragment.text) <= STORY_BEAT_TEXT_CAPACITY
        for fragment in code_fragments
    )
    assert "\n".join(fragment.text for fragment in code_fragments) == (
        "\n".join(code_lines)
    )

    story = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )

    subject_contract = story.planning_diagnostics[
        "subject_presentation_contract"
    ]
    assert subject_contract["schema_version"] == (
        "subject_presentation_contract_v1"
    )
    assert subject_contract["profile_id"] == "engineering_programming"
    assert subject_contract["required_representation_kinds"] == ["code"]
    assert {
        fragment.fragment_id for fragment in code_fragments
    } <= set(subject_contract["characteristic_fragment_ids"]["code"])

    assert any(
        code_fragment.fragment_id in beat.fragment_ids
        for chapter in story.chapters
        for episode in chapter.episodes
        for beat in episode.beats
        for code_fragment in code_fragments
    )

    compact_story = compact_story_plan_v5(document, story, fragments)
    allocation, _ = allocation_from_story_plan_v5(
        document,
        fragments,
        compact_story,
    )
    decided_fragment_ids = {
        fragment_id
        for page in allocation.pages
        for fragment_id in page.fragment_ids
    } | {
        exclusion.fragment_id for exclusion in allocation.exclusions
    }

    assert decided_fragment_ids == {
        fragment.fragment_id for fragment in fragments
    }
    allocated_code_ids = {
        fragment_id
        for page in allocation.pages
        if page.layout == "code"
        for fragment_id in page.fragment_ids
    }
    assert allocated_code_ids == {
        fragment.fragment_id for fragment in code_fragments
    }
    assert not allocated_code_ids & {
        exclusion.fragment_id for exclusion in allocation.exclusions
    }
    code_pages = [page for page in allocation.pages if page.layout == "code"]
    assert 1 <= len(code_pages) <= 3


def test_layout_registry_only_exposes_renderer_layouts_accepted_by_allocation() -> None:
    renderer_layouts = {
        str(layout["renderer_layout"])
        for layout in registry_summary_v2()
    }

    assert renderer_layouts <= set(V3_LAYOUTS)


def test_layout_rhythm_resets_at_chapter_boundaries() -> None:
    course = _course_with_teaching_plan()
    document = document_from_legacy_course(course)
    story = compile_slide_story_plan_v2(
        document,
        course,
        fragment_course_document(document),
        mode="teaching",
        theme="qizhi-classroom",
    )
    first_chapter = story.chapters[0].model_copy(deep=True)
    first_concept = next(
        episode
        for episode in first_chapter.episodes
        if episode.scene_kind == "concept"
    )
    template = first_concept.beats[0]
    first = template.model_copy(update={
        "beat_id": "beat-chapter-1-a",
        "layout_family": "comparison",
    })
    second = template.model_copy(update={
        "beat_id": "beat-chapter-1-b",
        "layout_family": "comparison",
    })
    first_concept.beats = [first, second]
    second_chapter = first_chapter.model_copy(
        deep=True,
        update={"chapter_id": "chapter-two"},
    )
    second_concept = next(
        episode
        for episode in second_chapter.episodes
        if episode.scene_kind == "concept"
    )
    third = template.model_copy(update={
        "beat_id": "beat-chapter-2-a",
        "layout_family": "comparison",
    })
    second_concept.beats = [third]
    two_chapter_story = story.model_copy(
        deep=True,
        update={"chapters": [first_chapter, second_chapter]},
    )

    quality = _presentation_quality(
        two_chapter_story,
        {
            "page-1": first,
            "page-2": second,
            "page-3": third,
        },
    )

    assert quality["passed"] is True
    assert not any(
        issue["code"] == "layout_family_repeated_more_than_twice"
        for issue in quality["issues"]
    )


def test_dense_mixed_concept_scene_is_split_before_layout_selection() -> None:
    course = _course_with_teaching_plan()
    concept_block = next(
        block
        for block in course["nodes"][0]["content_blocks"]
        if block["block_id"] == "block-concept"
    )
    concept_block["content"] = "\n\n".join([
        (
            "热力学系统由大量微观粒子构成，宏观状态需要用统计量描述。"
            "当系统从一个宏观状态演化到另一个宏观状态时，必须同时区分状态函数、"
            "过程量和约束条件，才能解释熵变、热量与功之间的关系。"
        ) * 4,
        "\n".join(
            f"- 判断要点 {index}：核对系统边界、状态变量与适用条件。"
            for index in range(1, 11)
        ),
        "$$\\Delta S = \\int \\frac{\\delta Q_{rev}}{T}$$",
        "```python\n" + "\n".join(
            f"state_{index} = energy_{index} / temperature_{index}"
            for index in range(1, 13)
        ) + "\n```",
    ])
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)

    plan = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )

    concept = next(
        episode
        for episode in plan.chapters[0].episodes
        if episode.scene_kind == "concept"
    )
    concept_source_ids = {
        fragment.fragment_id
        for fragment in fragments
        if fragment.block_id == "block-concept"
    }
    allocated_ids = [
        fragment_id
        for beat in concept.beats
        for fragment_id in beat.fragment_ids
    ]
    layouts = {
        layout["layout_id"]: layout
        for layout in registry_summary_v2()
    }
    fragment_by_id = {
        fragment.fragment_id: fragment
        for fragment in fragments
    }

    assert len(concept.beats) > 1
    assert set(allocated_ids) == concept_source_ids
    assert len(allocated_ids) == len(set(allocated_ids))
    for beat in concept.beats:
        layout = layouts[beat.layout_intent]
        beat_fragments = [
            fragment_by_id[fragment_id]
            for fragment_id in beat.fragment_ids
        ]
        assert sum(len(fragment.text) for fragment in beat_fragments) <= layout["density_budget"]
        assert sum(
            fragment.kind == "list_item"
            for fragment in beat_fragments
        ) <= layout["item_budget"]
        if sum(len(fragment.text) for fragment in beat_fragments) > 230:
            assert len(beat_fragments) == 1
            assert beat_fragments[0].kind in {"code", "formula", "table"}

    allocation, _ = allocation_from_story_plan_v2(
        document,
        fragments,
        plan,
    )
    content = compile_slide_deck_v4(
        document,
        course,
        story_plan=plan,
        allocation_plan=allocation,
    )
    blocker_codes = {
        item["code"]
        for item in content["quality_report"]["blockers"]
    }
    assert not blocker_codes & {
        "slide_block_overflow",
        "slide_item_overflow",
        "slide_text_overflow",
    }


def test_v4_allocation_preserves_source_order_across_semantic_scenes() -> None:
    course = _course_with_teaching_plan()
    original_section = course["nodes"][0]
    original_section["parent_node_id"] = "chapter-root"
    original_section["node_level"] = 2
    course["nodes"].insert(0, {
        "node_id": "chapter-root",
        "parent_node_id": "root",
        "node_name": "第一章 线性映射",
        "node_level": 1,
        "content_blocks": [{
            "block_id": "block-root-concept",
            "title": "章节概览",
            "content": "本章从保持线性结构的问题出发。",
            "metadata": {
                "role": "concept",
                "module_id": "module-concept",
            },
        }],
    })
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    story = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    fragment_by_id = {
        fragment.fragment_id: fragment
        for fragment in fragments
    }
    for chapter in story.chapters:
        for episode in chapter.episodes:
            for beat in episode.beats:
                beat_ordinals = [
                    fragment_by_id[fragment_id].ordinal
                    for fragment_id in beat.fragment_ids
                ]
                assert beat_ordinals == sorted(beat_ordinals)

    allocation, _ = allocation_from_story_plan_v2(
        document,
        fragments,
        story,
    )

    allocated_ordinals = [
        fragment_by_id[fragment_id].ordinal
        for page in allocation.pages
        if not page.appendix
        for fragment_id in page.fragment_ids
    ]
    assert allocated_ordinals == sorted(allocated_ordinals)


def test_story_plan_requires_completed_official_teaching_plan() -> None:
    course = _course_with_teaching_plan()
    course["generation_stage_artifacts"]["course_teaching_plan"]["status"] = "pending"
    document = document_from_legacy_course(course)

    with pytest.raises(SlideStoryPlanPrerequisiteError):
        compile_slide_story_plan_v2(
            document,
            course,
            fragment_course_document(document),
            mode="teaching",
            theme="qizhi-classroom",
        )


def test_enabled_story_engine_rejects_silent_v3_fallback() -> None:
    course = _course_with_teaching_plan()
    course.pop("course_knowledge_base")

    with pytest.raises(
        SlideStoryPlanPrerequisiteError,
        match="active official course knowledge base",
    ):
        resolve_slide_deck_schema(
            course,
            story_engine_enabled=True,
        )

    assert resolve_slide_deck_schema(
        course,
        story_engine_enabled=False,
    ) == "slide_deck_v3"


def test_enabled_story_engine_targets_v5_for_a_ready_course() -> None:
    assert resolve_slide_deck_schema(
        _course_with_teaching_plan(),
        story_engine_enabled=True,
    ) == "slide_deck_v5"


def test_v5_can_be_explicitly_rolled_back_to_v4() -> None:
    assert resolve_slide_deck_schema(
        _course_with_teaching_plan(),
        story_engine_enabled=True,
        v5_enabled=False,
    ) == "slide_deck_v4"


def test_v4_signature_binds_teaching_logic_revisions() -> None:
    course = _course_with_teaching_plan()
    document = document_from_legacy_course(course)
    original = build_signature_v4(
        document=document,
        course_data=course,
        mode="teaching",
        theme="qizhi-classroom",
    )
    changed_course = deepcopy(course)
    changed_course["course_teaching_plan"]["revision_id"] = "teaching-plan-rev-2"
    changed = build_signature_v4(
        document=document,
        course_data=changed_course,
        mode="teaching",
        theme="qizhi-classroom",
    )

    assert original["signature"] != changed["signature"]
    assert original["teaching_plan_revision"] == "teaching-plan-rev-1"


def test_v4_compilation_persists_scene_layout_and_source_bound_claims() -> None:
    course = _course_with_teaching_plan()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    story = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    allocation, _ = allocation_from_story_plan_v2(document, fragments, story)

    content = compile_slide_deck_v4(
        document,
        course,
        story_plan=story,
        allocation_plan=allocation,
    )

    assert content["schema_version"] == "slide_deck_v4"
    assert content["story_plan"]["schema_version"] == "slide_story_plan_v2"
    assert content["layout_plan"]["schema_version"] == "slide_layout_plan_v2"
    teaching_slides = [slide for slide in content["slides"] if slide["scene_kind"]]
    assert teaching_slides
    assert all(slide["teaching_job"] for slide in teaching_slides)
    assert all(slide["primary_claim_source"]["kind"] for slide in teaching_slides)
    prompt_index = next(
        index
        for index, slide in enumerate(content["slides"])
        if slide["scene_kind"] == "worked_example" and slide["beat_role"] == "prompt"
    )
    solution_index = next(
        index
        for index, slide in enumerate(content["slides"])
        if slide["scene_kind"] == "worked_example" and slide["beat_role"] == "solution"
    )
    assert prompt_index < solution_index


def test_illegal_ai_claim_is_rejected_and_uses_deterministic_story() -> None:
    course = _course_with_teaching_plan()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    baseline = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    ).model_dump(mode="json")
    baseline["chapters"][0]["episodes"][1]["beats"][0][
        "primary_claim_source"
    ]["text"] = "模型补造的课程观点"

    async def illegal_planner(_request: dict) -> dict:
        return baseline

    planned = asyncio.run(plan_slide_story_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
        ai_planner=illegal_planner,
    ))

    assert planned.planner == "deterministic_fallback"
    assert planned.fallback_reason == "invalid_or_failed_ai_story_plan"
    assert all(
        beat.primary_claim_source.text != "模型补造的课程观点"
        for chapter in planned.chapters
        for episode in chapter.episodes
        for beat in episode.beats
    )


def test_ai_story_planner_receives_bounded_source_text_for_semantic_decisions() -> None:
    course = _course_with_teaching_plan()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    baseline = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    ).model_dump(mode="json")
    captured: dict = {}

    async def planner(request: dict) -> dict:
        captured.update(request)
        return baseline

    planned = asyncio.run(plan_slide_story_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
        ai_planner=planner,
    ))

    assert planned.planner == "ai"
    assert captured["rules"]["structured_headlines_required"] is True
    assert captured["rules"]["body_text_forbidden"] is False
    assert captured["rules"]["copy_policy"] == "source_faithful_rewrite"
    assert captured["rules"]["unsupported_new_facts_forbidden"] is True
    assert all(
        item["source_text"]
        and len(item["source_text"]) <= 400
        for item in captured["fragments"]
    )
    assert all(item["semantic_unit_id"] for item in captured["fragments"])
    assert all("presentation_intent" in item for item in captured["fragments"])
    assert all("evidence_refs" in item for item in captured["fragments"])
    assert all(
        "semantic_unit_ids" in beat
        for beat in captured["beat_catalog"]
    )


def test_ai_story_planner_batches_large_decks_by_chapter(monkeypatch) -> None:
    course = _course_with_teaching_plan()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    fallback = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    second = fallback.chapters[0].model_copy(
        update={
            "chapter_id": "chapter-second",
            "next_chapter_id": "",
            "episodes": [
                episode.model_copy(update={
                    "beats": [
                        beat.model_copy(update={"fragment_ids": []})
                        for beat in episode.beats
                    ],
                })
                for episode in fallback.chapters[0].episodes
            ],
        },
    )
    batched_fallback = fallback.model_copy(update={
        "chapters": [
            fallback.chapters[0].model_copy(
                update={"next_chapter_id": "chapter-second"},
            ),
            second,
        ],
    })
    monkeypatch.setattr(
        "slide_story_plan.compile_slide_story_plan_v2",
        lambda *_args, **_kwargs: batched_fallback,
    )
    requests: list[dict] = []

    async def planner(request: dict) -> dict:
        requests.append(request)
        beat = request["beat_catalog"][0]
        return {
            "slide_story_chapter_directives_v2": {
                "chapter_id": request["scope"]["chapter_id"],
                "beat_directives": [{
                    "beat_id": beat["beat_id"],
                    "layout_id": beat["current_layout_id"],
                }],
            },
        }

    planned = asyncio.run(plan_slide_story_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
        ai_planner=planner,
    ))

    assert planned.planner == "ai"
    assert [item.chapter_id for item in planned.chapters] == [
        fallback.chapters[0].chapter_id,
        "chapter-second",
    ]
    assert len(requests) == 2
    assert all("deterministic_baseline" not in request for request in requests)
    assert all(request["chapter_contract"]["chapter_id"] for request in requests)
    assert [request["scope"]["chapter_index"] for request in requests] == [0, 1]


def test_ai_story_planner_retries_one_invalid_chapter_with_validation_errors() -> None:
    course = _course_with_teaching_plan()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    requests: list[dict] = []

    async def planner(request: dict) -> dict:
        requests.append(request)
        if len(requests) == 1:
            return {
                "schema_version": "slide_story_chapter_directives_v2",
                "chapter_id": request["scope"]["chapter_id"],
                "beat_directives": [{
                    "beat_id": "unknown-beat",
                }],
            }
        beat = request["beat_catalog"][0]
        return {
            "schema_version": "slide_story_chapter_directives_v2",
            "chapter_id": request["scope"]["chapter_id"],
            "beat_directives": [{
                "beat_id": beat["beat_id"],
                "layout_id": beat["current_layout_id"],
            }],
        }

    planned = asyncio.run(plan_slide_story_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
        ai_planner=planner,
    ))

    assert planned.planner == "ai"
    assert len(requests) == 2
    retry = requests[1]["validation_retry"]
    assert retry["attempt"] == 1
    assert retry["errors"][0]["code"] == "invalid_structure"
    assert "unknown beat" in retry["errors"][0]["message"].lower()


def test_ai_story_planner_keeps_valid_chapters_when_one_times_out(
    monkeypatch,
) -> None:
    course = _course_with_teaching_plan()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    fallback = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    second = fallback.chapters[0].model_copy(update={
        "chapter_id": "chapter-timeout",
        "next_chapter_id": "",
        "episodes": [
            episode.model_copy(update={
                "beats": [
                    beat.model_copy(update={"fragment_ids": []})
                    for beat in episode.beats
                ],
            })
            for episode in fallback.chapters[0].episodes
        ],
    })
    batched_fallback = fallback.model_copy(update={
        "chapters": [
            fallback.chapters[0].model_copy(
                update={"next_chapter_id": "chapter-timeout"},
            ),
            second,
        ],
    })
    monkeypatch.setattr(
        "slide_story_plan.compile_slide_story_plan_v2",
        lambda *_args, **_kwargs: batched_fallback,
    )

    async def planner(request: dict) -> dict:
        if request["scope"]["chapter_id"] == "chapter-timeout":
            raise asyncio.TimeoutError
        beat = next(
            item
            for item in request["beat_catalog"]
            if item["headline_candidates"]
        )
        return {
            "schema_version": "slide_story_chapter_directives_v2",
            "chapter_id": request["scope"]["chapter_id"],
            "beat_directives": [{
                "beat_id": beat["beat_id"],
                "layout_id": beat["current_layout_id"],
            }],
        }

    planned = asyncio.run(plan_slide_story_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
        ai_planner=planner,
    ))

    assert planned.planner == "ai"
    assert planned.fallback_reason == "partial_ai_story_plan"
    assert planned.chapters[1] == second
    assert planned.planning_diagnostics["failed_chapter_count"] == 1
    assert planned.planning_diagnostics["chapter_failures"] == [{
        "chapter_id": "chapter-timeout",
        "code": "timeout",
        "error_type": "TimeoutError",
    }]


def test_ai_story_planner_applies_compact_source_bound_directives() -> None:
    course = _course_with_teaching_plan()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    fallback = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    target = next(
        beat
        for chapter in fallback.chapters
        for episode in chapter.episodes
        for beat in episode.beats
        if beat.fragment_ids
    )
    headline_fragment_id = target.fragment_ids[-1]

    async def planner(request: dict) -> dict:
        return {
            "schema_version": "slide_story_chapter_directives_v2",
            "chapter_id": request["scope"]["chapter_id"],
            "episode_directives": [{
                "episode_id": next(
                    episode.episode_id
                    for episode in fallback.chapters[0].episodes
                    if any(
                        beat.beat_id == target.beat_id
                        for beat in episode.beats
                    )
                ),
                "beat_directives": [{
                    "beat_id": target.beat_id,
                    "headline_fragment_id": headline_fragment_id,
                    "layout_id": target.layout_intent,
                    "copy_mode": "source_faithful_rewrite",
                }],
            }],
        }

    planned = asyncio.run(plan_slide_story_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
        baseline=fallback,
        ai_planner=planner,
    ))
    planned_target = next(
        beat
        for chapter in planned.chapters
        for episode in chapter.episodes
        for beat in episode.beats
        if beat.beat_id == target.beat_id
    )

    assert planned.planner == "ai"
    assert planned_target.primary_claim_source.fragment_id == headline_fragment_id
    assert planned_target.primary_claim_source.text == next(
        item.text
        for item in fragments
        if item.fragment_id == headline_fragment_id
    )
    assert planned_target.copy_mode == "source_exact"


def test_ai_story_planner_accepts_grounded_audience_facing_copy() -> None:
    course = _course_with_teaching_plan()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    fallback = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    target = next(
        beat
        for chapter in fallback.chapters
        for episode in chapter.episodes
        for beat in episode.beats
        if beat.fragment_ids
    )

    async def planner(request: dict) -> dict:
        return {
            "schema_version": "slide_story_chapter_directives_v2",
            "chapter_id": request["scope"]["chapter_id"],
            "beat_directives": [{
                "beat_id": target.beat_id,
                "layout_id": target.layout_intent,
                "copy_mode": "source_faithful_rewrite",
                "audience_facing_title": "先抓住定义条件，再判断概念是否成立",
                "audience_facing_summary": "把来源中的定义、条件和边界压缩成一条可讲授的判断路径。",
                "supporting_fragment_ids": target.fragment_ids,
            }],
        }

    planned = asyncio.run(plan_slide_story_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
        baseline=fallback,
        ai_planner=planner,
    ))
    planned_target = next(
        beat
        for chapter in planned.chapters
        for episode in chapter.episodes
        for beat in episode.beats
        if beat.beat_id == target.beat_id
    )

    assert planned.planner == "ai"
    assert planned_target.primary_claim_source == target.primary_claim_source
    assert planned_target.audience_facing_title == "先抓住定义条件，再判断概念是否成立"
    assert planned_target.audience_facing_summary == (
        "把来源中的定义、条件和边界压缩成一条可讲授的判断路径。"
    )
    assert planned_target.copy_mode == "source_faithful_rewrite"
    assert planned_target.copy_source_fragment_ids == target.fragment_ids


def test_ai_story_planner_generates_answers_only_when_source_answer_is_missing() -> None:
    course = _course_with_teaching_plan()
    course["nodes"][0]["content_blocks"] = [
        block
        for block in course["nodes"][0]["content_blocks"]
        if block["block_id"] != "block-feedback"
    ]
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    fallback = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    support = next(
        fragment
        for fragment in fragments
        if fragment.block_id == "block-method" and fragment.kind != "heading"
    )
    captured_prompt: dict = {}

    async def planner(request: dict) -> dict:
        prompt = next(
            beat
            for beat in request["beat_catalog"]
            if beat["needs_generated_answers"]
        )
        captured_prompt.update(prompt)
        return {
            "schema_version": "slide_story_chapter_directives_v2",
            "chapter_id": request["scope"]["chapter_id"],
            "beat_directives": [{
                "beat_id": prompt["beat_id"],
                "layout_id": prompt["current_layout_id"],
                "generated_practice_answers": [{
                    "question_index": 0,
                    "question_id": prompt["question_ids"][0],
                    "answer_text": (
                        "不是线性映射，因为线性映射必须同时保持加法与数乘。"
                    ),
                    "supporting_fragment_ids": [support.fragment_id],
                }],
            }],
        }

    planned = asyncio.run(plan_slide_story_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
        baseline=fallback,
        ai_planner=planner,
    ))
    prompt = next(
        beat
        for chapter in planned.chapters
        for episode in chapter.episodes
        for beat in episode.beats
        if beat.generated_practice_answers
    )

    assert captured_prompt["needs_generated_answers"] is True
    assert len(captured_prompt["question_ids"]) == 1
    assert captured_prompt["prompt_questions"] == [
        "判断 S(x,y)=(x+1,y) 是否为线性映射。"
    ]
    assert prompt.generated_practice_answers[0].question_id == (
        captured_prompt["question_ids"][0]
    )
    assert prompt.generated_practice_answers[0].answer_source == "llm_generated"
    assert prompt.generated_practice_answers[0].answer_text.startswith(
        "不是线性映射"
    )
    assert prompt.generated_practice_answers[0].supporting_fragment_ids == [
        support.fragment_id
    ]


def test_generated_answer_may_reuse_numeric_premises_from_its_question() -> None:
    course = _course_with_teaching_plan()
    course["nodes"][0]["content_blocks"] = [
        block
        for block in course["nodes"][0]["content_blocks"]
        if block["block_id"] != "block-feedback"
    ]
    prompt_block = next(
        block
        for block in course["nodes"][0]["content_blocks"]
        if block["block_id"] == "block-practice"
    )
    prompt_block["content"] = (
        "两条路径分别吸收 100 J 和 120 J，哪条路径做功更多？为什么？"
    )
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    fallback = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    support = next(
        fragment
        for fragment in fragments
        if fragment.block_id == "block-method" and fragment.kind != "heading"
    )

    async def planner(request: dict) -> dict:
        prompt = next(
            beat
            for beat in request["beat_catalog"]
            if beat["needs_generated_answers"]
        )
        prompt_fragment_id = prompt["headline_candidates"][0]["fragment_id"]
        assert prompt["prompt_questions"] == [prompt_block["content"]]
        return {
            "schema_version": "slide_story_chapter_directives_v2",
            "chapter_id": request["scope"]["chapter_id"],
            "beat_directives": [{
                "beat_id": prompt["beat_id"],
                "layout_id": prompt["current_layout_id"],
                "audience_facing_title": "数值题判断",
                "supporting_fragment_ids": [prompt_fragment_id],
                "generated_practice_answers": [{
                    "question_index": 0,
                    "answer_text": (
                        "从 A 到 B 时，吸收 120 J 的路径做功更多，"
                        "因为题干给出的该路径吸热更多。"
                        + "这一判断依据题干给出的能量信息。" * 10
                        + "依据：sfg_hidden_1, sfg_hidden_2。"
                    ),
                    "supporting_fragment_ids": [support.fragment_id],
                }],
            }],
        }

    planned = asyncio.run(plan_slide_story_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
        baseline=fallback,
        ai_planner=planner,
    ))
    generated = [
        answer
        for chapter in planned.chapters
        for episode in chapter.episodes
        for beat in episode.beats
        for answer in beat.generated_practice_answers
    ]
    assert generated, planned.planning_diagnostics
    generated_beat = next(
        beat
        for chapter in planned.chapters
        for episode in chapter.episodes
        for beat in episode.beats
        if beat.generated_practice_answers
    )

    assert planned.planner == "ai"
    assert len(generated) == 1
    assert "120 J" in generated[0].answer_text
    assert "A 到 B" in generated[0].answer_text
    assert "sfg_" not in generated[0].answer_text
    assert len(generated[0].answer_text) <= 140
    assert generated_beat.copy_mode == "source_faithful_rewrite"


def test_ai_story_planner_drops_unsafe_optional_copy_without_losing_chapter() -> None:
    course = _course_with_teaching_plan()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    fallback = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    target = next(
        beat
        for chapter in fallback.chapters
        for episode in chapter.episodes
        for beat in episode.beats
        if beat.fragment_ids
    )

    async def planner(request: dict) -> dict:
        return {
            "schema_version": "slide_story_chapter_directives_v2",
            "chapter_id": request["scope"]["chapter_id"],
            "beat_directives": [{
                "beat_id": target.beat_id,
                "copy_mode": "source_faithful_rewrite",
                "audience_facing_title": "这种方法能提升50%效率",
                "supporting_fragment_ids": target.fragment_ids,
            }],
        }

    planned = asyncio.run(plan_slide_story_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
        baseline=fallback,
        ai_planner=planner,
    ))

    planned_target = next(
        beat
        for chapter in planned.chapters
        for episode in chapter.episodes
        for beat in episode.beats
        if beat.beat_id == target.beat_id
    )

    assert planned.planner == "ai"
    assert planned_target.copy_mode == "source_exact"
    assert planned_target.audience_facing_title == ""


def test_ai_story_planner_drops_incompatible_optional_fields_per_beat() -> None:
    course = _course_with_teaching_plan()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    fallback = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    target = next(
        beat
        for chapter in fallback.chapters
        for episode in chapter.episodes
        for beat in episode.beats
        if beat.fragment_ids
    )
    outside_fragment = next(
        fragment.fragment_id
        for fragment in fragments
        if fragment.fragment_id not in target.fragment_ids
    )

    async def planner(request: dict) -> dict:
        return {
            "schema_version": "slide_story_chapter_directives_v2",
            "chapter_id": request["scope"]["chapter_id"],
            "beat_directives": [{
                "beat_id": target.beat_id,
                "headline_fragment_id": outside_fragment,
                "layout_id": "not-a-compatible-layout",
                "copy_mode": "source_faithful_rewrite",
                "audience_facing_title": "可选标题",
                "supporting_fragment_ids": [outside_fragment],
            }],
        }

    planned = asyncio.run(plan_slide_story_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
        baseline=fallback,
        ai_planner=planner,
    ))
    planned_target = next(
        beat
        for chapter in planned.chapters
        for episode in chapter.episodes
        for beat in episode.beats
        if beat.beat_id == target.beat_id
    )

    assert planned.planner == "ai"
    assert planned_target.primary_claim_source == target.primary_claim_source
    assert planned_target.layout_intent == target.layout_intent
    assert planned_target.copy_mode == "source_exact"


def test_invalid_generated_answer_does_not_discard_other_chapter_directives() -> None:
    course = _course_with_teaching_plan()
    course["nodes"][0]["content_blocks"] = [
        block
        for block in course["nodes"][0]["content_blocks"]
        if block["block_id"] != "block-feedback"
    ]
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    fallback = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    prompt = next(
        beat
        for chapter in fallback.chapters
        for episode in chapter.episodes
        for beat in episode.beats
        if beat.beat_role == "prompt"
        and episode.scene_kind == "practice_feedback"
    )
    support = next(
        fragment.fragment_id
        for fragment in fragments
        if fragment.fragment_id not in prompt.fragment_ids
    )

    async def planner(request: dict) -> dict:
        return {
            "schema_version": "slide_story_chapter_directives_v2",
            "chapter_id": request["scope"]["chapter_id"],
            "beat_directives": [{
                "beat_id": prompt.beat_id,
                "layout_id": prompt.layout_intent,
                "generated_practice_answers": [{
                    "question_index": 0,
                    "answer_text": "温度会无依据地提升到 999 K。",
                    "supporting_fragment_ids": [support],
                }],
            }],
        }

    planned = asyncio.run(plan_slide_story_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
        baseline=fallback,
        ai_planner=planner,
    ))
    planned_prompt = next(
        beat
        for chapter in planned.chapters
        for episode in chapter.episodes
        for beat in episode.beats
        if beat.beat_id == prompt.beat_id
    )

    assert planned.planner == "ai"
    assert planned_prompt.layout_intent == prompt.layout_intent
    assert planned_prompt.generated_practice_answers == []


def test_surplus_optional_answers_do_not_discard_the_ai_chapter() -> None:
    course = _course_with_teaching_plan()
    course["nodes"][0]["content_blocks"] = [
        block
        for block in course["nodes"][0]["content_blocks"]
        if block["block_id"] != "block-feedback"
    ]
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    fallback = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    prompt = next(
        beat
        for chapter in fallback.chapters
        for episode in chapter.episodes
        for beat in episode.beats
        if beat.beat_role == "prompt"
        and episode.scene_kind == "practice_feedback"
    )
    support = next(
        fragment.fragment_id
        for fragment in fragments
        if fragment.fragment_id not in prompt.fragment_ids
    )

    async def planner(request: dict) -> dict:
        return {
            "schema_version": "slide_story_chapter_directives_v2",
            "chapter_id": request["scope"]["chapter_id"],
            "beat_directives": [{
                "beat_id": prompt.beat_id,
                "layout_id": prompt.layout_intent,
                "generated_practice_answers": [
                    {
                        "question_index": index,
                        "answer_text": f"Optional answer {index}",
                        "supporting_fragment_ids": [support],
                    }
                    for index in range(6)
                ],
            }],
        }

    planned = asyncio.run(plan_slide_story_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
        baseline=fallback,
        ai_planner=planner,
    ))
    planned_prompt = next(
        beat
        for chapter in planned.chapters
        for episode in chapter.episodes
        for beat in episode.beats
        if beat.beat_id == prompt.beat_id
    )

    assert planned.planner == "ai"
    assert planned.fallback_reason == ""
    assert planned_prompt.generated_practice_answers == []


def test_ai_practice_prompt_layout_compiles_to_a_supported_allocation() -> None:
    course = _course_with_teaching_plan()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    fallback = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    practice_episode = next(
        episode
        for episode in fallback.chapters[0].episodes
        if episode.scene_kind == "practice_feedback"
    )
    prompt = practice_episode.beats[0]

    async def planner(request: dict) -> dict:
        return {
            "schema_version": "slide_story_chapter_directives_v2",
            "chapter_id": request["scope"]["chapter_id"],
            "episode_directives": [{
                "episode_id": practice_episode.episode_id,
                "beat_directives": [{
                    "beat_id": prompt.beat_id,
                    "layout_id": "practice-prompt",
                }],
            }],
        }

    planned = asyncio.run(plan_slide_story_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
        baseline=fallback,
        ai_planner=planner,
    ))
    allocation, _ = allocation_from_story_plan_v2(
        document,
        fragments,
        planned,
    )
    planned_prompt = next(
        beat
        for chapter in planned.chapters
        for episode in chapter.episodes
        for beat in episode.beats
        if beat.beat_id == prompt.beat_id
    )

    assert planned.planner == "ai"
    assert planned_prompt.renderer_layout == "question"
    assert all(page.layout in V3_LAYOUTS for page in allocation.pages)


def test_ai_story_with_unknown_renderer_layout_uses_deterministic_fallback() -> None:
    course = _course_with_teaching_plan()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    fallback = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    raw_candidate = fallback.model_dump(mode="json")
    raw_candidate["chapters"][0]["episodes"][0]["beats"][0][
        "renderer_layout"
    ] = "unknown-renderer-layout"
    raw_candidate["planner"] = "ai"
    raw_candidate["fallback_reason"] = ""

    async def planner(_request: dict) -> dict:
        return raw_candidate

    planned = asyncio.run(plan_slide_story_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
        baseline=fallback,
        ai_planner=planner,
    ))

    assert planned.planner == "deterministic_fallback"
    assert planned.fallback_reason == "invalid_or_failed_ai_story_plan"
    assert all(
        beat.renderer_layout in V3_LAYOUTS
        for chapter in planned.chapters
        for episode in chapter.episodes
        for beat in episode.beats
    )


def test_v4_exports_editable_widescreen_pptx(tmp_path) -> None:
    course = _course_with_teaching_plan()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    story = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    allocation, _ = allocation_from_story_plan_v2(document, fragments, story)
    content = compile_slide_deck_v4(
        document,
        course,
        story_plan=story,
        allocation_plan=allocation,
    )
    path = export_structured_slide_deck(
        content,
        tmp_path / "linear-map-v4.pptx",
        theme="qizhi-classroom",
    )
    presentation = Presentation(path)

    assert len(presentation.slides) == len(content["slides"])
    assert round(presentation.slide_width / presentation.slide_height, 3) == round(16 / 9, 3)
    assert any(
        shape.has_text_frame and shape.text.strip()
        for slide in presentation.slides
        for shape in slide.shapes
    )


def test_v5_variant_is_atomically_published_under_existing_variant_key(tmp_path) -> None:
    course = _course_with_teaching_plan()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    story = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    allocation, _ = allocation_from_story_plan_v2(document, fragments, story)
    repository = TeachingRepresentationRepository(tmp_path / "registry")
    compile_core_representations(document, course, repository)
    legacy = next(
        item for item in repository.load(document.course_id).representations
        if item.representation_type == "slide_deck" and not item.variant_key
    )
    assert legacy.status == "ready"

    result = rebuild_slide_deck_variant_safely(
        document,
        course,
        repository,
        mode="teaching",
        theme="qizhi-classroom",
        allocation_plan=allocation,
        story_plan=story,
    )
    registry = repository.load(document.course_id)
    representation = next(
        item for item in registry.representations
        if item.variant_key == "teaching:qizhi-classroom"
    )
    spec = next(item for item in registry.specs if item.spec_id == representation.spec_id)

    assert result["status"] == "synchronized"
    assert representation.status == "ready"
    archived_legacy = next(
        item for item in registry.representations
        if item.representation_id == legacy.representation_id
    )
    assert archived_legacy.status == "archived"
    assert spec.payload["content"]["schema_version"] == "slide_deck_v5"
    assert spec.payload["content"]["deck_outline"]["schema_version"] == "deck_outline_v5"
    resolved_layouts = [
        (slide.get("quality") or {}).get("resolved_layout")
        for slide in spec.payload["content"]["slides"]
    ]
    assert resolved_layouts[:2] == ["cover-editorial", "agenda-linear"]
    assert resolved_layouts.count("chapter-entry") == len(story.chapters)


def test_v4_variant_becomes_stale_when_course_logic_revision_changes(tmp_path) -> None:
    course = _course_with_teaching_plan()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    story = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    allocation, _ = allocation_from_story_plan_v2(document, fragments, story)
    repository = TeachingRepresentationRepository(tmp_path / "registry")
    rebuild_slide_deck_variant_safely(
        document,
        course,
        repository,
        mode="teaching",
        theme="qizhi-classroom",
        allocation_plan=allocation,
        story_plan=story,
    )

    registry = repository.load(document.course_id)
    representation = next(
        item for item in registry.representations
        if item.variant_key == "teaching:qizhi-classroom"
    )
    assert representation.source_revision_vector["course_teaching_plan"] == (
        "teaching-plan-rev-1"
    )
    assert representation.source_revision_vector["course_knowledge_base"] == (
        "kb-rev-1"
    )
    assert representation.source_revision_vector["course_coherence_contract"] == (
        "coherence-rev-1"
    )

    changed_course = deepcopy(course)
    changed_course["course_knowledge_base"]["revision_id"] = "kb-rev-2"
    reconciled = repository.reconcile_source_revision_vector(
        document.course_id,
        revision_vector_for_course(document, changed_course),
    )
    changed_representation = next(
        item for item in reconciled.representations
        if item.representation_id == representation.representation_id
    )

    assert changed_representation.status == "stale"
    assert "source_revision_changed:course_knowledge_base" in (
        changed_representation.stale_reasons
    )
    assert set(changed_representation.stale_unit_ids) == set(
        next(
            item for item in reconciled.specs
            if item.spec_id == changed_representation.spec_id
        ).unit_bindings
    )


def test_v5_bundle_parts_keep_the_latest_story_engine(tmp_path) -> None:
    course = _course_with_teaching_plan()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    story = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="teaching",
        theme="qizhi-classroom",
    )
    allocation, _ = allocation_from_story_plan_v2(document, fragments, story)
    parts = split_slide_deck_plan_by_chapter(document, allocation)
    repository = TeachingRepresentationRepository(tmp_path / "registry")

    result = rebuild_slide_deck_variant_bundle_safely(
        document,
        course,
        repository,
        mode="teaching",
        theme="qizhi-classroom",
        parts=parts,
        story_plan=story,
    )

    assert result["status"] == "synchronized"
    registry = repository.load(document.course_id)
    schemas = {
        spec.payload["content"]["schema_version"]
        for spec in registry.specs
        if spec.variant_key.startswith("teaching:qizhi-classroom:part:")
    }
    assert schemas == {"slide_deck_v5"}


def test_concise_mode_keeps_a_minimum_loop_and_records_every_omission() -> None:
    course = _course_with_teaching_plan()
    document = document_from_legacy_course(course)
    fragments = fragment_course_document(document)
    story = compile_slide_story_plan_v2(
        document,
        course,
        fragments,
        mode="concise",
        theme="qizhi-classroom",
    )
    allocation, _ = allocation_from_story_plan_v2(document, fragments, story)
    included = {
        fragment_id
        for page in allocation.pages
        for fragment_id in page.fragment_ids
    }
    excluded = {item.fragment_id for item in allocation.exclusions}

    assert [item.scene_kind for item in story.chapters[0].episodes][0] == "chapter_entry"
    assert [item.scene_kind for item in story.chapters[0].episodes][-1] == "chapter_recap"
    assert "practice_feedback" in {
        item.scene_kind for item in story.chapters[0].episodes
    }
    assert included | excluded == {item.fragment_id for item in fragments}
    assert included & excluded == set()


def test_chapter_story_aggregates_official_plans_bound_to_child_sections() -> None:
    course = _course_with_teaching_plan()
    original_section = course["nodes"][0]
    original_section["parent_node_id"] = "chapter-root"
    original_section["node_level"] = 2
    course["nodes"].insert(0, {
        "node_id": "chapter-root",
        "parent_node_id": "root",
        "node_name": "第一章 线性映射",
        "node_level": 1,
        "content_blocks": [],
    })
    document = document_from_legacy_course(course)

    story = compile_slide_story_plan_v2(
        document,
        course,
        fragment_course_document(document),
        mode="teaching",
        theme="qizhi-classroom",
    )

    assert [chapter.chapter_id for chapter in story.chapters] == ["chapter-root"]
    assert story.chapters[0].owned_knowledge_ids == ["kp-linear-map"]
    assert story.chapters[0].learning_objective
    assert "concept" in {
        episode.scene_kind for episode in story.chapters[0].episodes
    }
