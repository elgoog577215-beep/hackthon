from copy import deepcopy

import pytest
from pptx import Presentation
from pptx.util import Pt

from ppt_layout_execution import NativeTarget, compile_teaching_template
from ppt_native_scene import audit_scene, inspect_native_bindings, native_target_geometry, render_scene
from ppt_page_scene import resolve_page_scenes
from ppt_teaching_content import PageTeachingV2
from .test_ppt_teaching_content import comparison_fixture, scene_for


def test_native_branch_connections_keep_original_targets_and_hide_with_nodes(tmp_path):
    from ppt_layout_schema import NativeConnection
    from ppt_native_scene import clone_native_slide, inspect_native_connections
    from .test_ppt_teaching_content import branch_fixture
    value, _ = branch_fixture()
    complete = value['states'][0]
    value['states'] = [{**deepcopy(complete), 'state_id': 'first', 'visible_element_ids': ['a']}, complete]
    template = compile_teaching_template('academic-editorial', certification_required=False)
    layout = template.get_layout(template.layout_id('concept-map'))
    scenes = resolve_page_scenes(page_id='branch', title='分支', content=PageTeachingV2.model_validate(value),
        layout=layout, template=template, source_document_revision='doc')
    deck = Presentation()
    deck.slide_width, deck.slide_height = Pt(960), Pt(540)
    source = deck.slides.add_slide(deck.slide_layouts[6])
    render_scene(source, scenes[-1])
    names = {s.name: s for s in source.shapes}
    execution = layout.execution.model_copy(deep=True)
    execution.mode, execution.source_slide_number, execution.source_sha256 = 'native_fill', 1, 'fixture'
    execution.targets = {o.slot_id: NativeTarget(shape_id=names[f'teaching:{o.object_id}'].shape_id,
        geometry_pt=(o.x, o.y, o.width, o.height)) for o in scenes[-1].objects}
    objects = {o.object_id: o for o in scenes[-1].objects}
    execution.connections = {f'{objects[e.source_id].slot_id}->{objects[e.target_id].slot_id}': NativeConnection(
        shape_id=names[f'relation:{e.relation_id}'].shape_id, source_slot=objects[e.source_id].slot_id,
        target_slot=objects[e.target_id].slot_id, start_site=e.start_site, end_site=e.end_site,
        geometry_pt=(e.x1, e.y1, e.x2, e.y2), directed=False) for e in scenes[-1].edges}
    inspect_native_connections(source, execution)
    layout.execution = execution
    native = resolve_page_scenes(page_id='branch', title='分支', content=PageTeachingV2.model_validate(value),
        layout=layout, template=template, source_document_revision='doc')
    for scene in native:
        slide = clone_native_slide(deck, source, execution)
        render_scene(slide, scene)
        assert audit_scene(slide, scene)['passed']
        assert len([s for s in slide.shapes if s.name.startswith('relation:')]) == len(scene.edges)
    execution.connections[next(iter(execution.connections))].directed = True
    with pytest.raises(ValueError, match='native_connection_direction_mismatch'):
        inspect_native_connections(source, execution)


def test_native_table_cells_fill_and_read_back_without_new_text_boxes(tmp_path):
    content, _ = comparison_fixture()
    component = scene_for(content)
    deck = Presentation()
    deck.slide_width, deck.slide_height = Pt(960), Pt(540)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    table_shape = slide.shapes.add_table(2, 3, Pt(42), Pt(160), Pt(876), Pt(294))
    table_shape.table.columns[0].width = Pt(150)
    for index in (1, 2):
        table_shape.table.columns[index].width = Pt(363)
    cells = {"serial": (0, 1), "parallel": (0, 2), "dimension": (1, 0), "a-mode": (1, 1), "b-mode": (1, 2)}
    targets = {}
    for obj in component.objects:
        if obj.element_id in cells:
            row, column = cells[obj.element_id]
            targets[obj.slot_id] = NativeTarget(shape_id=table_shape.shape_id, row=row, column=column)
        else:
            shape = slide.shapes.add_textbox(Pt(obj.x), Pt(obj.y), Pt(obj.width), Pt(obj.height))
            shape.text = "待填充"
            targets[obj.slot_id] = NativeTarget(shape_id=shape.shape_id)
    bindings = inspect_native_bindings(slide, targets)
    template = compile_teaching_template("academic-editorial", certification_required=False)
    layout = template.get_layout(template.layout_id("compare-matrix")).model_copy(deep=True)
    layout.execution.mode = "native_fill"
    layout.execution.targets = bindings
    layout.execution.source_sha256 = "fixture-source"
    layout.execution.source_slide_number = 1
    scene = resolve_page_scenes(page_id="p", title="执行方式", content=PageTeachingV2.model_validate(content),
        layout=layout, template=template, source_document_revision="doc")[0]
    before_shapes = len(slide.shapes)
    render_scene(slide, scene)
    path = tmp_path / "native-table.pptx"
    deck.save(path)
    restored = Presentation(path).slides[0]
    assert len(restored.shapes) == before_shapes
    assert audit_scene(restored, scene)["passed"]
    restored.shapes[0].table.cell(1, 1).text = "同时执行"
    with pytest.raises(ValueError, match="export_element_text_mismatch:a-mode"):
        audit_scene(restored, scene)


def test_native_target_missing_ambiguous_duplicate_and_group_geometry():
    p = Presentation()
    slide = p.slides.add_slide(p.slide_layouts[6])
    group = slide.shapes.add_group_shape()
    shape = group.shapes.add_textbox(Pt(100), Pt(100), Pt(200), Pt(60))
    group.left, group.top = Pt(200), Pt(150)
    group.width = Pt(400)
    target = NativeTarget(shape_id=shape.shape_id, group_path=[group.shape_id])
    assert native_target_geometry(slide, target) == (200, 150, 400, 60)
    with pytest.raises(ValueError, match="native_scaled_group_unsupported"):
        inspect_native_bindings(slide, {"body": target})
    group.width = Pt(200)
    assert inspect_native_bindings(slide, {"body": target})["body"].geometry_pt == (200, 150, 200, 60)
    with pytest.raises(ValueError, match="bound_twice"):
        inspect_native_bindings(slide, {"a": target, "b": target})
    with pytest.raises(ValueError, match="missing_or_ambiguous"):
        inspect_native_bindings(slide, {"bad": NativeTarget(shape_id=999)})
    group.shapes._spTree.insert_element_before(deepcopy(shape._element), "p:extLst")
    with pytest.raises(ValueError, match="missing_or_ambiguous"):
        inspect_native_bindings(slide, {"body": target})
