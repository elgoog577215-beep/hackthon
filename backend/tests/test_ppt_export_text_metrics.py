from pptx import Presentation
from pptx.util import Pt

from slide_deck_renderer import _text_frame_audit


def frame(text, *, width=80, height=26):
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    shape = slide.shapes.add_textbox(Pt(50), Pt(200), Pt(width), Pt(height))
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Pt(0)
    tf.word_wrap = False
    paragraph = tf.paragraphs[0]
    paragraph.text = text
    paragraph.font.name = "Noto Sans CJK SC"
    paragraph.font.size = Pt(20)
    paragraph.line_spacing = Pt(26)
    return shape


def test_fixed_chinese_line_does_not_get_wrapped_again():
    report = _text_frame_audit(frame("触发问题"))
    assert not report["overflow"]
    assert report["maximum_wrapped_lines"] == 1
    assert report["required_height_pt"] == 26


def test_no_wrap_still_rejects_real_horizontal_overflow():
    report = _text_frame_audit(frame("触发问题后"))
    assert report["overflow"] and report["horizontal_overflow"]


def test_fixed_lines_use_the_written_line_spacing():
    shape = frame("触发\v问题", height=40)
    report = _text_frame_audit(shape)
    assert report["overflow"] and report["required_height_pt"] == 52
    shape.height = Pt(52)
    assert not _text_frame_audit(shape)["overflow"]


def test_title_height_tolerance_cannot_hide_horizontal_overflow(tmp_path):
    from slide_deck_renderer import audit_exported_pptx
    deck = Presentation()
    deck.slide_width, deck.slide_height = Pt(960), Pt(540)
    slide = deck.slides.add_slide(deck.slide_layouts[6])
    shape = slide.shapes.add_textbox(Pt(48), Pt(60), Pt(120), Pt(60))
    shape.text_frame.word_wrap = False
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.font.name, paragraph.font.size = "Noto Sans CJK SC", Pt(30)
    paragraph.text = "标题不能横向溢出"
    path = tmp_path / "title.pptx"
    deck.save(path)
    report = audit_exported_pptx(path, require_pixel_audit=False)
    assert any(b["code"] == "exported_text_frame_overflow" for b in report["blockers"])


def test_pdf_aliases_are_identical_in_the_pinned_font():
    from PIL import ImageFont
    from ppt_layout_execution import FONT_PATH
    from ppt_render_audit import _glyphs
    font = ImageFont.truetype(str(FONT_PATH), 80)
    for left, right in (("…", "⋯"), ("⎢", "⎥")):
        assert font.getbbox(left) == font.getbbox(right)
        assert font.getlength(left) == font.getlength(right)
        assert bytes(font.getmask(left)) == bytes(font.getmask(right))
        assert _glyphs(left) == _glyphs(right)
    assert _glyphs("≠") != _glyphs("=")
