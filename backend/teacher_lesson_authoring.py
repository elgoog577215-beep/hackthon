"""Teacher-only lesson plan assets and jobs.

This module deliberately does not write ``CourseDocument``.  It is the
authoring boundary for a teacher lesson (one L1 node plus all direct L2
sections) while the existing learner course-generation pipeline remains
unchanged.
"""

from __future__ import annotations

import asyncio
import inspect
import json
import os
import re
import tempfile
import threading
import time
import uuid
import hashlib
from copy import deepcopy
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Awaitable, Callable

from course_document import document_from_generation_draft
from course_pedagogy import module_block_role
from lesson_identity import lesson_chapter_index, resolve_lesson_chapter
from teaching_design import normalize_lesson_arrangement
from teacher_script import (
    SCRIPT_PIPELINE_VERSION,
    SCRIPT_QUALITY_VERSION,
    compile_teacher_script_generation_shards,
    compile_teacher_script_module_contract,
    compile_teacher_script_shard_context,
    normalize_teacher_script_section,
    teacher_script_revision_is_publishable,
    validate_teacher_script_section,
    validate_teacher_script_revision,
)
from teacher_visible_language import has_unnatural_system_language


SCHEMA_VERSION = "teacher_lesson_authoring_v1"
LESSON_PLAN_PIPELINE_VERSION = "standard_lesson_plan_v1"
LESSON_PLAN_FORMAL_FIELD_POLICY_VERSION = "teacher_lesson_formal_fields_v1"
TEACHER_ASSET_JOB_SCHEMA_VERSION = "teacher_asset_job_v1"
LESSON_JOB_STALE_SECONDS = 300
LESSON_BATCH_QUEUED_STALE_SECONDS = 14400
JOB_TYPES = {
    "teacher_lesson_plan_generation",
    "teacher_lesson_script_generation",
}
_PLAN_INTERNAL_REGISTER_PATTERN = re.compile(
    r"全课知识地图|先修链定位|学习路径角色|可观察成果证据|证据闭环|"
    r"输入对象|输出对象|系统(?:策略|将会|将|会自动)|模型(?:生成|输出)|质量门|"
    r"资料不足|不得编造来源|不声称来自真实数据|无已确认教师资料或外部来源"
)
_PLAN_ABSTRACT_ACTIVITY_PATTERN = re.compile(
    r"建立问题、价值与任务边界|调取经验并作出初始判断|"
    r"依据[“\"].{0,60}[”\"]检查是否服务本讲目标"
)


class TeacherLessonAuthoringError(RuntimeError):
    def __init__(self, code: str, message: str, *, details: dict[str, Any] | None = None):
        super().__init__(message)
        self.code = code
        self.details = details or {}


def _now() -> str:
    return datetime.now(timezone.utc).isoformat()


def _empty_lesson_asset(lesson_unit_id: str) -> dict[str, Any]:
    return {
        "lesson_unit_id": lesson_unit_id,
        "arrangement": {
            "working_revision_id": "",
            "confirmed_revision_id": "",
            "source_state": "current",
            "revisions": [],
        },
        "working_revision_id": "",
        "confirmed_revision_id": "",
        "source_state": "current",
        "revisions": [],
        "ai_candidates": [],
        "working_script_revision_id": "",
        "script_revisions": [],
        "script_confirmation": {},
        "ppt_manuscript": {},
        "ppt_assets": [],
        "imported_ppt_reviews": [],
        "material_drafts": {},
        "current_material_draft_ids": {},
    }


def _mark_lesson_dependents_stale(
    lesson: dict[str, Any],
    *,
    reason: str,
) -> None:
    """Keep last usable preparation assets while marking their source as changed."""
    if lesson.get("working_revision_id"):
        lesson["source_state"] = "stale"
        lesson["source_state_reason"] = reason
    confirmation = lesson.get("script_confirmation")
    if isinstance(confirmation, dict) and confirmation.get("confirmed_revision_id"):
        confirmation["source_state"] = "stale"
    for asset in lesson.get("ppt_assets") or []:
        if isinstance(asset, dict):
            asset["source_state"] = "stale"
    manuscript = lesson.get("ppt_manuscript")
    if isinstance(manuscript, dict) and manuscript:
        manuscript["source_state"] = "stale"
    for review in lesson.get("imported_ppt_reviews") or []:
        if isinstance(review, dict):
            review["source_state"] = "stale"


def _text_list(value: Any) -> list[str]:
    if isinstance(value, str):
        return [item.strip() for item in value.splitlines() if item.strip()]
    if isinstance(value, list):
        return [str(item).strip() for item in value if isinstance(item, str) and item.strip()]
    return []


def _unique_text(values: list[str]) -> list[str]:
    return list(dict.fromkeys(item.strip() for item in values if item.strip()))


_MODULE_LABELS = {
    "lesson_goal": "明确本讲目标",
    "core_explanation": "核心概念讲解",
    "learner_action": "学生解释与练习",
    "explained_example": "例题示范",
    "guided_practice": "带支架练习",
    "feedback_check": "检查与反馈",
    "general_concept_model": "概念模型建构",
    "general_comparison": "对比辨析",
    "general_explained_example": "例题推演",
    "general_application": "迁移应用",
    "general_checklist": "课堂小结",
    "summary_and_transfer": "总结与迁移",
    "assessment": "学习检查",
    "math_problem_strategy": "解题策略选择",
    "math_worked_example": "例题推演",
    "math_intuition": "直觉与问题导入",
    "math_representation": "多重表征转换",
    "math_formalization": "形式化推导",
    "math_variation": "变式练习",
}
_GENERIC_MODULE_LABEL = re.compile(r"^(?:环节|教学块|模块)\s*\d*$")


def _specific_module_label(module: dict[str, Any], index: int) -> str:
    module_id = str(module.get("module_id") or "").strip()
    label = str(module.get("label") or "").strip()
    if label and label != module_id and not _GENERIC_MODULE_LABEL.fullmatch(label):
        return label
    known = _MODULE_LABELS.get(module_id)
    if known:
        return known
    purpose = str(
        module.get("teaching_purpose")
        or module.get("teaching_guidance")
        or module.get("teacher_activity")
        or ""
    ).strip()
    purpose = re.split(r"[，。；：]", purpose, maxsplit=1)[0].strip()
    if purpose:
        return purpose[:24]
    return f"教学环节 {index + 1}"


def _format_recommended_reading(value: Any) -> str:
    if isinstance(value, str):
        return value.strip()
    if not isinstance(value, dict):
        return ""
    title = str(value.get("title") or value.get("name") or "").strip()
    source_ref = str(value.get("source_ref") or "").strip()
    if not title and not source_ref:
        return ""
    resource_type = {
        "book": "教材/专著",
        "article": "论文/文章",
        "standard": "标准",
        "regulation": "法规",
        "dataset": "数据集",
        "video": "视频",
        "website": "网站",
        "other": "资料",
    }.get(str(value.get("resource_type") or value.get("type") or "").lower(), "资料")
    citation = source_ref or title
    extras = _unique_text([
        str(value.get("edition") or "").strip(),
        str(value.get("locator") or value.get("chapter") or "").strip(),
    ])
    status = str(value.get("verification_status") or "").strip().lower()
    prefix = "已确认来源" if status == "verified" else "AI 推荐（待教师确认）"
    return "｜".join([
        prefix,
        resource_type,
        citation,
        *(extras[:2]),
    ])


def complete_teacher_lesson_plan_fields(
    course_data: dict[str, Any],
    lesson_unit_id: str,
    plan: dict[str, Any],
) -> dict[str, Any]:
    """Complete inferable formal fields without inventing teacher-owned facts.

    The teaching-plan model owns substantive lesson content. This boundary
    carries confirmed outline readings into the lesson and deterministically
    completes fields that can be derived from that content. Submission channel,
    deadline and real post-class activity records remain teacher decisions.
    """
    completed = normalize_teacher_lesson_plan(plan)
    source_plan = course_data.get("course_plan") or course_data.get("course_outline") or {}
    chapters = [
        item for item in source_plan.get("chapters") or []
        if isinstance(item, dict)
    ]
    resolved_chapter_index = lesson_chapter_index(source_plan, lesson_unit_id)
    chapter_index = (
        resolved_chapter_index if resolved_chapter_index is not None else -1
    )
    chapter = chapters[chapter_index] if chapter_index >= 0 else {}
    next_chapter = (
        chapters[chapter_index + 1]
        if 0 <= chapter_index < len(chapters) - 1
        else None
    )
    source_sections = {
        str(item.get("node_id") or ""): item
        for item in chapter.get("sections") or []
        if isinstance(item, dict) and str(item.get("node_id") or "")
    }
    course_references = _unique_text([
        *_text_list(source_plan.get("reference_books")),
        *_text_list(source_plan.get("reference_websites")),
    ])
    chapter_readings = [
        _format_recommended_reading(item)
        for item in chapter.get("extension_resources") or []
    ]
    chapter_readings = _unique_text(chapter_readings)
    next_title = str(
        (next_chapter or {}).get("title")
        or (next_chapter or {}).get("learning_focus")
        or ""
    ).strip()

    for section in completed.get("sections") or []:
        if not isinstance(section, dict):
            continue
        source_section = source_sections.get(str(section.get("node_id") or "")) or {}
        key_points = _text_list(section.get("key_points"))
        objectives = _text_list(section.get("knowledge_objectives")) or _text_list(
            section.get("learning_objective")
        )
        checks = _text_list(section.get("in_class_checks"))
        difficulties = _text_list(section.get("key_difficulties"))
        anchor_items = key_points[:3] or objectives[:2]
        anchor = "、".join(anchor_items) or str(
            source_section.get("title") or chapter.get("title") or "本讲内容"
        ).strip()

        for index, module in enumerate(section.get("teaching_modules") or []):
            if isinstance(module, dict):
                module["label"] = _specific_module_label(module, index)

        if not _text_list(section.get("class_summary")):
            summary = [f"回到本讲目标，梳理{anchor}之间的关系、成立条件与适用边界。"]
            if checks:
                summary.append(f"通过“{checks[0]}”确认学生能否独立说明关键依据。")
            section["class_summary"] = summary

        if not str(section.get("homework_evaluation") or "").strip():
            section["homework_evaluation"] = (
                f"重点检查{anchor}相关结论是否准确、关键步骤与依据是否完整，"
                "以及能否在任务给定的情境中正确应用。"
            )

        if not str(section.get("next_lesson_connection") or "").strip():
            section["next_lesson_connection"] = (
                f"本讲形成的“{anchor}”将作为下一讲“{next_title}”的前置基础，"
                "课后任务用于检查是否具备进入条件。"
                if next_title
                else f"本讲是课程收束，围绕“{anchor}”完成整课知识连接与迁移复盘。"
            )

        if not _text_list(section.get("teaching_notes")):
            difficulty = difficulties[0] if difficulties else anchor
            section["teaching_notes"] = [
                f"重点观察学生在“{difficulty}”处的判断依据；出现偏差时先追问条件，"
                "再用反例、变式或重新示范完成就近纠正。"
            ]

        if not _text_list(section.get("resource_refs")):
            section_readings = [
                _format_recommended_reading(item)
                for item in source_section.get("extension_resources") or []
            ]
            readings = _unique_text([
                *section_readings,
                *chapter_readings,
                *[f"已确认来源｜课程参考资料｜{item}" for item in course_references],
            ])
            if readings:
                section["resource_refs"] = readings[:3]

        teacher_confirmation_fields = _text_list(
            section.get("teacher_confirmation_fields")
        )
        if not str(section.get("homework_submission") or "").strip():
            teacher_confirmation_fields.append("homework_submission")
        if not _text_list(section.get("teaching_activity_photos")):
            teacher_confirmation_fields.append("teaching_activity_photos")
        section["teacher_confirmation_fields"] = _unique_text(
            teacher_confirmation_fields
        )

    completed["formal_field_policy_version"] = (
        LESSON_PLAN_FORMAL_FIELD_POLICY_VERSION
    )
    return completed


def teacher_lesson_section_content(section: dict[str, Any]) -> dict[str, Any]:
    """Project plan-v3 knowledge/modules into concrete teacher-editable fields.

    Real provider fallbacks keep useful knowledge facts but often omit the old
    ``learning_objective`` / ``teacher_activities`` convenience fields.  This
    projection is the compatibility boundary shared by preview, editing and
    PPT source compilation; it never serializes module dictionaries as prose.
    """
    points = [
        point
        for group in section.get("knowledge_structure") or []
        if isinstance(group, dict)
        for point in group.get("knowledge_points") or []
        if isinstance(point, dict)
    ]
    modules = [
        item for item in section.get("teaching_modules") or []
        if isinstance(item, dict)
    ]
    capability_objectives = [
        str(capability.get("observable_behavior") or capability.get("name") or "").strip()
        for point in points
        for capability in point.get("capability_points") or []
        if isinstance(capability, dict)
        and str(capability.get("observable_behavior") or capability.get("name") or "").strip()
    ]
    statements = [
        str(point.get("statement") or point.get("description") or "").strip()
        for point in points
        if str(point.get("statement") or point.get("description") or "").strip()
    ]
    boundaries = [
        item for point in points for item in _text_list(point.get("boundaries"))
    ]
    misconceptions = [
        str(item.get("discrimination") or item.get("name") or "").strip()
        for point in points
        for item in point.get("misconceptions") or []
        if isinstance(item, dict)
        and str(item.get("discrimination") or item.get("name") or "").strip()
    ]

    def module_line(module: dict[str, Any]) -> str:
        labels = {
            "lesson_goal": "本节目标",
            "core_explanation": "核心讲解",
            "math_problem_strategy": "策略选择",
            "math_worked_example": "例题推演",
            "math_intuition": "直觉导入",
            "math_representation": "多重表征",
            "math_formalization": "正式定义",
            "math_variation": "变式练习",
            "learner_action": "学习者行动",
            "feedback_check": "检查与反馈",
        }
        module_id = str(module.get("module_id") or "")
        label = str(module.get("label") or labels.get(module_id) or "教学活动")
        knowledge = _text_list(module.get("knowledge_names"))
        guidance = str(module.get("teaching_guidance") or "").strip()
        if guidance.startswith("按模板完成"):
            guidance = ""
        details = [f"围绕{'、'.join(knowledge)}" if knowledge else "", guidance]
        suffix = "；".join(item for item in details if item)
        return f"{label}：{suffix}" if suffix else label

    explicit_objective = str(
        section.get("learning_objective") or section.get("objective") or ""
    ).strip()
    knowledge_objectives = _text_list(
        section.get("knowledge_objectives") or section.get("knowledge_objective")
    ) or ([explicit_objective] if explicit_objective else statements[:3])
    ability_objectives = _text_list(
        section.get("ability_objectives") or section.get("ability_objective")
    ) or capability_objectives[:3] or ([explicit_objective] if explicit_objective else [])
    education_objectives = _text_list(
        section.get("education_objectives")
        or section.get("education_objective")
        or section.get("education_goal")
        or section.get("ideology_goal")
    )
    learning_objective = (
        explicit_objective
        or "；".join(ability_objectives)
        or "；".join(knowledge_objectives)
        or "；".join(statements)
        or "；".join(_text_list(section.get("key_points")))
    )
    key_difficulties = _text_list(section.get("key_difficulties")) or _unique_text([
        *_text_list(section.get("key_points")),
        *boundaries,
        *misconceptions,
    ])
    module_teacher_activities = _unique_text([
        str(module.get("teacher_activity") or "").strip()
        for module in modules
        if str(module.get("teacher_activity") or "").strip()
    ])
    module_student_activities = _unique_text([
        str(module.get("student_activity") or "").strip()
        for module in modules
        if str(module.get("student_activity") or "").strip()
    ])
    teacher_activities = module_teacher_activities or _text_list(
        section.get("teacher_activities")
    ) or _unique_text([
        module_line(module)
        for module in modules
        if str(module.get("module_id") or "") not in {
            "learner_action", "math_variation", "feedback_check",
        }
    ])
    student_activities = module_student_activities or _text_list(
        section.get("student_activities")
    ) or _unique_text([
        module_line(module)
        for module in modules
        if str(module.get("module_id") or "") in {"learner_action", "math_variation"}
    ])
    homework = _text_list(section.get("homework")) or _unique_text([
        module_line(module)
        for module in modules
        if str(module.get("module_id") or "") == "feedback_check"
    ])
    return {
        "learning_objective": learning_objective,
        "knowledge_objectives": knowledge_objectives,
        "ability_objectives": ability_objectives,
        "education_objectives": education_objectives,
        "key_difficulties": key_difficulties,
        "teacher_activities": teacher_activities,
        "student_activities": student_activities,
        "homework": homework,
        "knowledge_statements": statements,
        "knowledge_boundaries": boundaries,
        "misconceptions": misconceptions,
    }


def normalize_teacher_lesson_plan(plan: dict[str, Any]) -> dict[str, Any]:
    normalized = deepcopy(plan)
    normalized["sections"] = []
    for section in plan.get("sections") or []:
        if not isinstance(section, dict):
            continue
        next_section = deepcopy(section)
        content = teacher_lesson_section_content(next_section)
        for key, value in content.items():
            if key in {"knowledge_statements", "knowledge_boundaries", "misconceptions"}:
                continue
            if key in {"teacher_activities", "student_activities"} and next_section.get("teaching_modules"):
                next_section[key] = deepcopy(value)
            elif not next_section.get(key):
                next_section[key] = deepcopy(value)
        checks = _text_list(next_section.get("in_class_checks"))
        key_points = _text_list(next_section.get("key_points"))
        normalized_modules: list[dict[str, Any]] = []
        for index, module in enumerate(next_section.get("teaching_modules") or []):
            if not isinstance(module, dict):
                continue
            normalized_module = deepcopy(module)
            label = str(
                normalized_module.get("label")
                or normalized_module.get("teaching_purpose")
                or normalized_module.get("module_id")
                or f"教学块 {index + 1}"
            ).strip()
            expected_output = str(
                normalized_module.get("expected_output")
                or (checks[0] if checks else "")
                or normalized_module.get("student_activity")
                or "完成本教学块的当堂任务并说明依据。"
            ).strip()
            if not str(normalized_module.get("teaching_purpose") or "").strip():
                normalized_module["teaching_purpose"] = label
            if not str(normalized_module.get("expected_output") or "").strip():
                normalized_module["expected_output"] = expected_output
            if not str(normalized_module.get("check_method") or "").strip():
                normalized_module["check_method"] = checks[0] if checks else expected_output
            if not str(normalized_module.get("feedback_strategy") or "").strip():
                normalized_module["feedback_strategy"] = "根据当堂表现给出具体反馈，并决定直接推进、补充支架或重新示范。"
            if not _text_list(normalized_module.get("adaptation_options")):
                normalized_module["adaptation_options"] = [
                    "达到目标时进入下一教学块。",
                    "部分达到时补充提示或示例后再次检查。",
                    "未达到时缩小任务、重新示范并安排复查。",
                ]
            if not str(normalized_module.get("transition") or "").strip():
                normalized_module["transition"] = f"根据“{label}”的当堂证据决定是否进入下一教学块。"
            if not _text_list(normalized_module.get("resource_refs")):
                normalized_module["resource_refs"] = _text_list(next_section.get("resource_refs"))
            normalized_module["tools"] = _text_list(normalized_module.get("tools"))
            mapping_anchor = key_points[0] if key_points else label
            if not str(normalized_module.get("handout_ppt_mapping") or "").strip():
                normalized_module["handout_ppt_mapping"] = f"讲义对应“{mapping_anchor}”内容；PPT 对应同名教学页面。"
            normalized_modules.append(normalized_module)
        next_section["teaching_modules"] = normalized_modules
        normalized["sections"].append(next_section)
    return normalized


def align_teacher_lesson_plan_to_arrangement(
    plan: dict[str, Any],
    arrangement: dict[str, Any] | None,
) -> dict[str, Any]:
    """Make the current lesson arrangement authoritative for order and time.

    The model owns the substantive explanation.  It does not own the lesson
    clock: real-model responses have assigned the full lesson duration to each
    section, omitted early sections, and silently changed module order.  This
    boundary preserves usable model content while deterministically restoring
    the teacher-visible block identities, sequence and minute allocation.
    Missing classroom fields are completed only from the same current block
    and the section's already-grounded knowledge, never from a second model.
    """
    normalized = normalize_teacher_lesson_plan(plan)
    blocks = [
        deepcopy(item)
        for item in (arrangement or {}).get("blocks") or []
        if isinstance(item, dict)
        and str(item.get("section_node_id") or "").strip()
    ]
    if not blocks:
        return normalized

    blocks_by_section: dict[str, list[dict[str, Any]]] = {}
    for block in blocks:
        blocks_by_section.setdefault(
            str(block.get("section_node_id") or ""), []
        ).append(block)

    for section in normalized.get("sections") or []:
        if not isinstance(section, dict):
            continue
        section_id = str(section.get("node_id") or "")
        arranged = blocks_by_section.get(section_id) or []
        if not arranged:
            continue
        existing_modules = [
            deepcopy(item)
            for item in section.get("teaching_modules") or []
            if isinstance(item, dict)
        ]
        used_indexes: set[int] = set()
        aligned_modules: list[dict[str, Any]] = []

        for block in arranged:
            module_id = str(block.get("module_id") or "core_explanation")
            matched_index = next(
                (
                    index
                    for index, item in enumerate(existing_modules)
                    if index not in used_indexes
                    and str(item.get("module_id") or "") == module_id
                ),
                None,
            )
            source_module = (
                existing_modules[matched_index]
                if matched_index is not None
                else {}
            )
            if matched_index is not None:
                used_indexes.add(matched_index)
            knowledge_names = _text_list(source_module.get("knowledge_names"))
            if not knowledge_names:
                knowledge_names = _text_list(section.get("key_points"))
            block_name = str(block.get("name") or module_id).strip()
            purpose = str(
                block.get("purpose")
                or block.get("content_summary")
                or source_module.get("teaching_purpose")
                or f"完成“{block_name}”并形成可检查的课堂产出。"
            ).strip()
            expected_output = str(
                block.get("expected_output")
                or f"完成“{block_name}”对应的课堂任务并说明依据。"
            ).strip()
            teacher_activity = str(
                source_module.get("teacher_activity")
                or block.get("teacher_activity")
                or (f"围绕“{block_name}”完成讲解、示范与要点收束。" if block_name else "")
            ).strip()
            student_activity = str(
                source_module.get("student_activity")
                or block.get("student_activity")
                or (f"完成“{expected_output}”，并说明判断依据。" if expected_output else "记录关键步骤并完成当堂自检。")
            ).strip()
            aligned_modules.append({
                **source_module,
                "module_id": module_id,
                "label": str(source_module.get("label") or block_name),
                "teaching_purpose": str(
                    source_module.get("teaching_purpose") or purpose
                ),
                "teaching_guidance": str(
                    source_module.get("teaching_guidance")
                    or block.get("content_summary")
                    or purpose
                ),
                "knowledge_names": knowledge_names,
                "planned_minutes": max(
                    1, int(block.get("planned_minutes") or 1)
                ),
                "arrangement_block_id": str(block.get("block_id") or ""),
                "teacher_activity": teacher_activity,
                "student_activity": student_activity,
                "expected_output": str(
                    source_module.get("expected_output") or expected_output
                ).strip(),
                "check_method": str(
                    source_module.get("check_method")
                    or block.get("check_method")
                    or expected_output
                ).strip(),
                "feedback_strategy": str(
                    source_module.get("feedback_strategy")
                    or block.get("feedback_strategy")
                    or "根据当堂表现给出具体反馈，并决定直接推进、补充支架或重新示范。"
                ).strip(),
                "adaptation_options": _text_list(
                    source_module.get("adaptation_options")
                    or block.get("adaptation_options")
                ) or [
                    "达到目标时进入下一教学块。",
                    "部分达到时补充提示或示例后再次检查。",
                    "未达到时缩小任务、重新示范并安排复查。",
                ],
                "resource_refs": _text_list(
                    source_module.get("resource_refs")
                    or block.get("resource_refs")
                    or section.get("resource_refs")
                ),
                "tools": _text_list(
                    source_module.get("tools")
                    or block.get("tools")
                ),
                "transition": str(
                    source_module.get("transition")
                    or block.get("transition")
                    or f"根据“{block_name}”的当堂证据决定是否进入下一教学块。"
                ).strip(),
                "handout_ppt_mapping": str(
                    source_module.get("handout_ppt_mapping")
                    or block.get("handout_ppt_mapping")
                    or f"讲义对应“{block_name}”内容；PPT 对应同名教学页面。"
                ).strip(),
            })

        section["teaching_modules"] = aligned_modules
        section["planned_minutes"] = sum(
            int(item.get("planned_minutes") or 0)
            for item in aligned_modules
        )
        section["teacher_activities"] = _unique_text([
            str(item.get("teacher_activity") or "")
            for item in aligned_modules
        ])
        section["student_activities"] = _unique_text([
            str(item.get("student_activity") or "")
            for item in aligned_modules
        ])
        if not _text_list(section.get("key_difficulties")):
            section["key_difficulties"] = _text_list(
                section.get("key_points")
            )
        if not _text_list(section.get("in_class_checks")):
            section["in_class_checks"] = _unique_text([
                str(item.get("expected_output") or "")
                for item in arranged
            ]) or ["根据本节目标完成一次当堂自检并说明依据。"]
        if not _text_list(section.get("homework")):
            anchor = next(
                iter(_text_list(section.get("key_points"))),
                str(section.get("learning_objective") or "本节目标").strip(),
            )
            section["homework"] = [
                f"围绕“{anchor}”完成一次课后巩固与迁移练习。"
            ]
        section["timing_source"] = "current_lesson_arrangement"

    normalized["lesson_duration_minutes"] = sum(
        max(1, int(item.get("planned_minutes") or 1))
        for item in blocks
    )
    normalized["timing_source"] = "current_lesson_arrangement"
    return normalize_teacher_lesson_plan(normalized)


def validate_teacher_lesson_plan(
    plan: dict[str, Any],
    *,
    expected_section_ids: list[str] | None = None,
    expected_outline_revision_id: str = "",
    source_outline_revision_id: str = "",
    expected_total_minutes: int | None = None,
) -> dict[str, Any]:
    """Validate the one teacher-facing standard lesson-plan contract.

    The V3 planner remains responsible for grounded knowledge and pedagogy.
    This report checks the classroom document: complete section coverage,
    observable objectives, executable activities, checks, timing and homework.
    Findings remain available for teacher review; downstream authoring uses the
    current structurally complete revision instead of treating this score as a gate.
    """
    normalized = normalize_teacher_lesson_plan(plan)
    sections = [
        item for item in normalized.get("sections") or []
        if isinstance(item, dict)
    ]
    actual_ids = [str(item.get("node_id") or "") for item in sections]
    blocking: list[dict[str, str]] = []
    review: list[dict[str, str]] = []

    def issue(
        target: list[dict[str, str]],
        code: str,
        message: str,
        section_id: str = "",
    ) -> None:
        payload = {"code": code, "message": message}
        if section_id:
            payload["section_id"] = section_id
        target.append(payload)

    if str(normalized.get("schema_version") or "") != "course_teaching_plan_v3":
        issue(blocking, "lesson_plan:schema", "教案必须使用统一的 CourseTeachingPlanV3 结构。")
    if not sections:
        issue(blocking, "lesson_plan:sections_empty", "教案没有可用的教学小节。")
    if any(not section_id for section_id in actual_ids):
        issue(blocking, "lesson_plan:section_identity", "教案小节缺少稳定标识。")
    if len(actual_ids) != len(set(actual_ids)):
        issue(blocking, "lesson_plan:duplicate_section", "教案包含重复小节。")
    if expected_section_ids is not None and actual_ids != expected_section_ids:
        issue(blocking, "lesson_plan:section_scope", "教案必须按大纲顺序完整覆盖当前讲次的所有小节。")
    if (
        expected_outline_revision_id
        and source_outline_revision_id
        and expected_outline_revision_id != source_outline_revision_id
    ):
        issue(blocking, "lesson_plan:stale_outline", "教案对应的大纲版本已经过期。")

    total_modules = 0
    total_minutes = 0
    knowledge_point_count = 0
    enforce_formal_fields = (
        str(normalized.get("formal_field_policy_version") or "")
        == LESSON_PLAN_FORMAL_FIELD_POLICY_VERSION
    )
    for section in sections:
        section_id = str(section.get("node_id") or "")
        objective = str(section.get("learning_objective") or "").strip()
        knowledge_objectives = _text_list(section.get("knowledge_objectives"))
        ability_objectives = _text_list(section.get("ability_objectives"))
        key_points = _text_list(section.get("key_points"))
        difficulties = _text_list(section.get("key_difficulties"))
        checks = _text_list(section.get("in_class_checks"))
        homework = _text_list(section.get("homework"))
        modules = [
            item for item in section.get("teaching_modules") or []
            if isinstance(item, dict)
        ]
        total_modules += len(modules)
        section_minutes = sum(
            max(0, int(item.get("planned_minutes") or 0))
            for item in modules
            if str(item.get("planned_minutes") or "").isdigit()
        )
        total_minutes += section_minutes

        if not objective or not knowledge_objectives:
            issue(blocking, "lesson_plan:knowledge_objective", "本讲缺少明确的知识目标。", section_id)
        if not ability_objectives:
            issue(blocking, "lesson_plan:ability_objective", "本讲缺少可观察的能力目标。", section_id)
        if not key_points:
            issue(blocking, "lesson_plan:key_points", "小节缺少教学重点。", section_id)
        if not difficulties:
            issue(blocking, "lesson_plan:difficulties", "小节缺少教学难点。", section_id)
        if not modules:
            issue(blocking, "lesson_plan:modules", "小节缺少可执行的教学流程。", section_id)
        if modules and not any(str(item.get("teacher_activity") or "").strip() for item in modules):
            issue(blocking, "lesson_plan:teacher_activity", "小节缺少具体的教师活动。", section_id)
        if modules and not any(str(item.get("student_activity") or "").strip() for item in modules):
            issue(blocking, "lesson_plan:student_activity", "小节缺少具体的学生活动。", section_id)
        if modules and section_minutes <= 0:
            issue(blocking, "lesson_plan:timing", "小节缺少有效的时间分配。", section_id)
        for module in modules:
            missing = [
                field for field in (
                    "teaching_purpose", "teacher_activity", "student_activity",
                    "expected_output", "check_method", "feedback_strategy",
                    "transition", "handout_ppt_mapping",
                )
                if not str(module.get(field) or "").strip()
            ]
            if not _text_list(module.get("adaptation_options")):
                missing.append("adaptation_options")
            if missing:
                issue(
                    blocking,
                    "lesson_plan:block_contract",
                    f"教学块缺少可执行字段：{'、'.join(missing)}。",
                    section_id,
                )
        if not checks:
            issue(blocking, "lesson_plan:checks", "小节缺少课堂检查或可观察产出。", section_id)
        if not homework:
            issue(blocking, "lesson_plan:homework", "小节缺少课后巩固或迁移任务。", section_id)
        if enforce_formal_fields:
            if not _text_list(section.get("class_summary")):
                issue(blocking, "lesson_plan:class_summary", "本讲缺少课程总结。", section_id)
            if not _text_list(section.get("resource_refs")):
                issue(blocking, "lesson_plan:recommended_reading", "本讲缺少可识别的推荐阅读来源。", section_id)
            if not str(section.get("homework_evaluation") or "").strip():
                issue(blocking, "lesson_plan:homework_evaluation", "课后作业缺少评价标准。", section_id)
            if not str(section.get("next_lesson_connection") or "").strip():
                issue(blocking, "lesson_plan:next_lesson_connection", "本讲缺少后续衔接或课程收束说明。", section_id)
            if not _text_list(section.get("teaching_notes")):
                issue(blocking, "lesson_plan:teaching_notes", "本讲缺少可执行的教学提醒。", section_id)
            if any(
                not str(item.get("label") or "").strip()
                or _GENERIC_MODULE_LABEL.fullmatch(str(item.get("label") or "").strip())
                for item in modules
            ):
                issue(blocking, "lesson_plan:block_name", "每个教学环节都需要与内容对应的具体名称。", section_id)
        public_copy = [
            objective,
            *key_points,
            *difficulties,
            *checks,
            *homework,
            *_text_list(section.get("teacher_activities")),
            *_text_list(section.get("student_activities")),
            *_text_list(section.get("teaching_notes")),
        ]
        for module in modules:
            public_copy.extend(
                str(module.get(field) or "").strip()
                for field in (
                    "teaching_purpose",
                    "teacher_activity",
                    "student_activity",
                    "expected_output",
                    "check_method",
                    "feedback_strategy",
                    "transition",
                )
            )
        visible_text = "\n".join(item for item in public_copy if item)
        if (
            _PLAN_INTERNAL_REGISTER_PATTERN.search(visible_text)
            or has_unnatural_system_language(visible_text)
        ):
            issue(
                blocking,
                "lesson_plan:internal_register",
                "教案夹带了资料、模型或内部规划说明，没有写成教师实际备课时会使用的语言。",
                section_id,
            )
        if _PLAN_ABSTRACT_ACTIVITY_PATTERN.search(visible_text):
            issue(
                blocking,
                "lesson_plan:abstract_activity",
                "教案仍使用抽象流程套话，师生活动需要改成针对本节内容的具体课堂动作。",
                section_id,
            )

        points = [
            point
            for group in section.get("knowledge_structure") or []
            if isinstance(group, dict)
            for point in group.get("knowledge_points") or []
            if isinstance(point, dict)
        ]
        knowledge_point_count += len(points)
        for point in points:
            name = str(point.get("name") or "").strip()
            statement = str(point.get("statement") or point.get("description") or "").strip()
            if name and not statement:
                issue(blocking, "lesson_plan:knowledge_statement", f"知识点「{name}」缺少准确陈述。", section_id)
            if point.get("conflict") or point.get("needs_manual_review"):
                issue(blocking, "lesson_plan:knowledge_conflict", f"知识点「{name or '未命名'}」存在待核实冲突。", section_id)

    if (
        expected_total_minutes is not None
        and total_minutes != int(expected_total_minutes)
    ):
        issue(
            blocking,
            "lesson_plan:total_timing",
            (
                f"本讲教案共 {total_minutes} 分钟，与已确认的 "
                f"{int(expected_total_minutes)} 分钟编排不一致。"
            ),
        )

    return {
        "schema_version": "teacher_lesson_plan_quality_v1",
        "pipeline_version": LESSON_PLAN_PIPELINE_VERSION,
        "passed": not blocking,
        "blocking_issues": blocking,
        "review_issues": review,
        "metrics": {
            "section_count": len(sections),
            "teaching_module_count": total_modules,
            "knowledge_point_count": knowledge_point_count,
            "planned_minutes": total_minutes,
            "teacher_language_rule_version": "teacher_plan_language_v2",
        },
    }


def extract_uploaded_pptx_evidence(
    path: Path,
    *,
    asset_id: str,
) -> list[dict[str, Any]]:
    """Extract visible slide text without mutating the teacher's source deck."""
    if path.suffix.lower() != ".pptx":
        raise TeacherLessonAuthoringError(
            "uploaded_ppt_format_unsupported",
            "旧课件同源解析目前仅支持 PPTX 文件。",
        )
    try:
        from pptx import Presentation

        presentation = Presentation(path)
    except Exception as exc:
        raise TeacherLessonAuthoringError(
            "uploaded_ppt_parse_failed",
            "旧课件无法解析，请确认文件未损坏。",
        ) from exc

    evidence: list[dict[str, Any]] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            text = ""
            if getattr(shape, "has_text_frame", False):
                text = str(getattr(shape, "text", "") or "")
            elif getattr(shape, "has_table", False):
                text = "\n".join(
                    str(cell.text or "")
                    for row in shape.table.rows
                    for cell in row.cells
                )
            normalized = " ".join(text.split()).strip()
            if normalized and normalized not in parts:
                parts.append(normalized)
        source_text = "\n".join(parts).strip()
        if not source_text:
            continue
        evidence.append({
            "evidence_id": f"uploaded-ppt-{asset_id}-slide-{slide_number}",
            "kind": "uploaded_ppt_slide",
            "summary": source_text[:1200],
            "source_text": source_text[:2400],
            "slide": slide_number,
            "asset_id": asset_id,
        })
    if not evidence:
        raise TeacherLessonAuthoringError(
            "uploaded_ppt_empty",
            "旧课件中没有可用于生成教案的文字内容。",
        )
    return evidence


def extract_uploaded_pptx_review(
    path: Path,
    *,
    asset_id: str,
    filename: str = "",
) -> dict[str, Any]:
    """Build an editable text index while keeping the uploaded PPTX immutable."""
    if path.suffix.lower() != ".pptx":
        raise TeacherLessonAuthoringError(
            "uploaded_ppt_format_unsupported",
            "PPT 审阅目前仅支持 PPTX 文件。",
        )
    try:
        from pptx import Presentation

        presentation = Presentation(path)
    except Exception as exc:
        raise TeacherLessonAuthoringError(
            "uploaded_ppt_parse_failed",
            "PPT 无法解析，请确认文件未损坏。",
        ) from exc

    slides: list[dict[str, Any]] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        blocks: list[dict[str, Any]] = []
        for shape_index, shape in enumerate(slide.shapes):
            raw_text = ""
            kind = "text"
            editable = False
            if getattr(shape, "has_text_frame", False):
                raw_text = str(getattr(shape, "text", "") or "")
                editable = True
                shape_name = str(getattr(shape, "name", "") or "").lower()
                if "title" in shape_name or "标题" in shape_name:
                    kind = "title"
            elif getattr(shape, "has_table", False):
                raw_text = "\n".join(
                    str(cell.text or "")
                    for row in shape.table.rows
                    for cell in row.cells
                )
                kind = "table"
            text = "\n".join(line.strip() for line in raw_text.splitlines() if line.strip())
            if not text:
                continue
            blocks.append({
                "block_id": f"uploaded-ppt-{asset_id}-s{slide_number}-b{shape_index}",
                "shape_index": shape_index,
                "kind": kind,
                "text": text,
                "original_text": text,
                "editable": editable,
            })
        if blocks and not any(item.get("kind") == "title" for item in blocks):
            blocks[0]["kind"] = "title"
        title_block = next(
            (item for item in blocks if item.get("kind") == "title"),
            None,
        )
        content_hash = hashlib.sha256(
            json.dumps(blocks, ensure_ascii=False, sort_keys=True).encode("utf-8")
        ).hexdigest()[:24]
        slides.append({
            "slide_id": f"uploaded-ppt-{asset_id}-slide-{slide_number}",
            "slide_number": slide_number,
            "title": str((title_block or {}).get("text") or ""),
            "blocks": blocks,
            "content_hash": content_hash,
        })
    if not slides:
        raise TeacherLessonAuthoringError(
            "uploaded_ppt_empty",
            "PPT 中没有可审阅的页面。",
        )
    return {
        "source_asset_id": asset_id,
        "source_filename": filename or path.name,
        "slides": slides,
    }


def _review_terms(text: str) -> set[str]:
    words = re.findall(r"[a-z0-9_]{3,}|[\u4e00-\u9fff]{2,}", text.lower())
    terms: set[str] = set()
    for word in words:
        terms.add(word)
        if re.fullmatch(r"[\u4e00-\u9fff]{2,}", word):
            terms.update(word[index:index + 2] for index in range(len(word) - 1))
    return terms


def build_uploaded_ppt_review_report(
    slides: list[dict[str, Any]],
    *,
    sources: list[dict[str, Any]],
    reference_units: list[dict[str, Any]],
) -> dict[str, Any]:
    """Return explainable review findings; indexes assist but never publish edits."""
    findings: list[dict[str, Any]] = []
    confirmed_sources = [item for item in sources if item.get("status") == "confirmed"]
    confidence = "high" if len(confirmed_sources) >= 2 else "medium" if sources else "low"

    def add_finding(
        code: str,
        title: str,
        detail: str,
        *,
        slide_id: str = "",
        slide_number: int | None = None,
        severity: str = "suggestion",
        evidence: list[dict[str, Any]] | None = None,
    ) -> None:
        target = slide_id or "deck"
        digest = hashlib.sha256(f"{code}:{target}".encode("utf-8")).hexdigest()[:12]
        findings.append({
            "finding_id": f"ppt-review-{digest}",
            "code": code,
            "title": title,
            "detail": detail,
            "severity": severity,
            "confidence": confidence,
            "slide_id": slide_id,
            "slide_number": slide_number,
            "status": "open",
            "evidence": deepcopy(evidence or []),
        })

    reference_terms = set().union(*(
        _review_terms(str(item.get("text") or "")) for item in reference_units
    )) if reference_units else set()
    slide_terms: dict[str, set[str]] = {}
    for slide in slides:
        slide_id = str(slide.get("slide_id") or "")
        blocks = [item for item in slide.get("blocks") or [] if isinstance(item, dict)]
        text = "\n".join(str(item.get("text") or "") for item in blocks).strip()
        terms = _review_terms(text)
        slide_terms[slide_id] = terms
        slide_number = int(slide.get("slide_number") or 0)
        title = str(slide.get("title") or "").strip()
        if not blocks:
            add_finding(
                "visual_only_slide",
                "该页未识别到可审阅文字",
                "当前只检查到视觉内容，请确认该页是否需要讲解文字或备注。",
                slide_id=slide_id,
                slide_number=slide_number,
            )
        elif not title:
            add_finding(
                "slide_title_missing",
                "该页缺少明确标题",
                "补充页面标题可以让教学进度和学习目标更容易被识别。",
                slide_id=slide_id,
                slide_number=slide_number,
            )
        if len(text) > 260 or len(blocks) > 8:
            add_finding(
                "slide_content_dense",
                "该页文字较密",
                f"已识别 {len(text)} 个字符、{len(blocks)} 个文字块，建议拆分或保留一个主结论。",
                slide_id=slide_id,
                slide_number=slide_number,
            )
        if reference_terms and terms and len(reference_terms & terms) < 2:
            add_finding(
                "slide_alignment_unresolved",
                "与已确认教学内容的对应关系不明确",
                "索引未找到该页与当前教案或讲义的明确对应，请确认是补充材料还是需要调整。",
                slide_id=slide_id,
                slide_number=slide_number,
                evidence=[{"kind": item.get("kind"), "label": item.get("label"), "revision_id": item.get("revision_id")} for item in confirmed_sources],
            )

    for unit in reference_units:
        unit_terms = _review_terms(str(unit.get("text") or ""))
        if not unit_terms or any(len(unit_terms & terms) >= 2 for terms in slide_terms.values()):
            continue
        label = str(unit.get("label") or "教学内容")
        add_finding(
            "source_unit_not_covered",
            f"未找到“{label}”的明确对应页",
            "已确认的教案或讲义中包含该内容，但 PPT 索引未找到足够相关的页面。",
            evidence=[{
                "kind": unit.get("kind"),
                "label": label,
                "revision_id": unit.get("revision_id"),
            }],
        )

    return {
        "schema_version": "uploaded_ppt_review_report_v1",
        "generated_at": _now(),
        "sources": deepcopy(sources),
        "findings": findings,
        "summary": {
            "slide_count": len(slides),
            "finding_count": len(findings),
            "high_confidence_count": sum(item.get("confidence") == "high" for item in findings),
        },
    }


def _default_root() -> Path:
    configured = os.getenv("TEACHER_LESSON_AUTHORING_DIR", "").strip()
    if configured:
        return Path(configured)
    return Path(__file__).resolve().parent / "data" / "teacher_lesson_authoring"


def lesson_scope(course_data: dict[str, Any], lesson_unit_id: str) -> dict[str, Any]:
    """Resolve one stable teacher lesson and its direct ordered sections."""
    nodes = [item for item in course_data.get("nodes") or [] if isinstance(item, dict)]
    lesson = next(
        (
            item for item in nodes
            if str(item.get("node_id") or "") == lesson_unit_id
            and int(item.get("node_level") or 0) == 1
        ),
        None,
    )
    if not lesson:
        raise TeacherLessonAuthoringError(
            "lesson_unit_not_found",
            "当前课程中不存在该讲次。",
            details={"lesson_unit_id": lesson_unit_id},
        )
    sections = [
        deepcopy(item) for item in nodes
        if str(item.get("parent_node_id") or "") == lesson_unit_id
    ]
    if not sections:
        raise TeacherLessonAuthoringError(
            "lesson_sections_empty",
            "当前讲次没有可生成教案的小节。",
            details={"lesson_unit_id": lesson_unit_id},
        )
    plan = course_data.get("course_plan") or course_data.get("course_outline") or {}
    chapters = [item for item in plan.get("chapters") or [] if isinstance(item, dict)]
    chapter = resolve_lesson_chapter(plan, lesson_unit_id)
    if chapter is None:
        section_ids = {str(item.get("node_id") or "") for item in sections}
        chapter = next(
            (
                item for item in chapters
                if section_ids
                and section_ids.issubset({
                    str(section.get("node_id") or "")
                    for section in item.get("sections") or []
                    if isinstance(section, dict)
                })
            ),
            None,
        )
    return {
        "lesson": deepcopy(lesson),
        "sections": sections,
        "chapter": deepcopy(chapter) if isinstance(chapter, dict) else None,
    }


def teacher_lesson_script_revision(
    course_data: dict[str, Any],
    lesson_unit_id: str,
) -> str:
    """Fingerprint the saved course body used as this lesson's single script source."""
    scope = lesson_scope(course_data, lesson_unit_id)
    payload = [
        {
            "section_node_id": str(section.get("node_id") or ""),
            "title": str(section.get("node_name") or ""),
            "content": str(section.get("node_content") or ""),
        }
        for section in scope["sections"]
    ]
    digest = hashlib.sha256(
        json.dumps(payload, ensure_ascii=False, sort_keys=True).encode("utf-8")
    ).hexdigest()[:24]
    return f"tlsr-{digest}"


def teacher_lesson_script_sections_revision(
    sections: list[dict[str, Any]],
) -> str:
    """Fingerprint one teacher-script asset without depending on course storage."""
    payload = []
    for section in sections:
        if not isinstance(section, dict):
            continue
        blocks = [
            block for block in section.get("blocks") or [] if isinstance(block, dict)
        ]
        # Preserve the historic content fingerprint for one-way migrated v1
        # scripts so existing confirmation links remain valid.
        legacy_only = bool(blocks) and all(
            str(block.get("module_id") or "") == "legacy_script" for block in blocks
        )
        item = {
            "section_node_id": str(section.get("section_node_id") or ""),
            "title": str(section.get("title") or ""),
        }
        if blocks and not legacy_only:
            item["blocks"] = [
                {
                    "block_id": str(block.get("block_id") or ""),
                    "module_id": str(block.get("module_id") or ""),
                    "role": str(block.get("role") or ""),
                    "title": str(block.get("title") or ""),
                    "content": str(block.get("content") or ""),
                    "knowledge_names": list(block.get("knowledge_names") or []),
                    "planned_minutes": block.get("planned_minutes"),
                }
                for block in blocks
            ]
        else:
            item["content"] = str(
                (blocks[0].get("content") if legacy_only else section.get("content"))
                or ""
            )
        payload.append(item)
    digest = hashlib.sha256(
        json.dumps(payload, ensure_ascii=False, sort_keys=True).encode("utf-8")
    ).hexdigest()[:24]
    return f"tlsr-{digest}"


def teacher_lesson_deck_to_structured_slide_deck(
    deck: dict[str, Any],
    *,
    source_revision_id: str,
) -> dict[str, Any]:
    """Adapt the lightweight teacher deck to the shared editable PPTX renderer."""
    raw_slides = [item for item in deck.get("slides") or [] if isinstance(item, dict)]
    if not raw_slides:
        raise TeacherLessonAuthoringError("lesson_ppt_empty", "本讲 PPT 没有可导出的页面。")
    slides = []
    for index, slide in enumerate(raw_slides):
        body = slide.get("body") or []
        if isinstance(body, str):
            body = [body]
        layout = "cover" if index == 0 else "recap" if index == len(raw_slides) - 1 else "concept"
        slides.append({
            "unit_id": str(slide.get("slide_id") or f"slide-{index + 1}"),
            "position": index,
            "layout": layout,
            "slide_purpose": "teacher_lesson_presentation",
            "title": str(slide.get("title") or f"第 {index + 1} 页"),
            "key_message": str(body[0] if body else slide.get("title") or ""),
            "takeaway": str(body[-1] if body else ""),
            "blocks": [{
                "block_id": f"{slide.get('slide_id') or index}-bullets",
                "type": "bullets",
                "items": [str(item) for item in body if str(item).strip()][:8],
            }],
            "speaker_notes": str(slide.get("speaker_notes") or ""),
        })
    return {
        "schema_version": "slide_deck_v2",
        "title": str(deck.get("title") or "本讲课件"),
        "theme": "qingfeng-classroom",
        "source_document_revision": source_revision_id,
        "slides": slides,
    }


def teacher_lesson_v6_source(
    course_data: dict[str, Any],
    *,
    lesson_unit_id: str,
    plan_revision: dict[str, Any],
    script_revision: dict[str, Any],
) -> tuple[Any, dict[str, Any], str]:
    """Adapt one current, complete teacher script revision to the V6 contracts.

    The returned synthetic course id is stable for one real course + lesson,
    while the CourseDocument revision changes with the teacher plan. Nothing is
    persisted to the learner CourseDocument repository.
    """
    scope = lesson_scope(course_data, lesson_unit_id)
    plan = normalize_teacher_lesson_plan(plan_revision.get("plan") or {})
    plan_sections = {
        str(item.get("node_id") or ""): deepcopy(item)
        for item in plan.get("sections") or []
        if isinstance(item, dict) and item.get("node_id")
    }
    digest = hashlib.sha256(
        f"{course_data.get('course_id')}:{lesson_unit_id}".encode("utf-8")
    ).hexdigest()[:20]
    synthetic_course_id = f"teacher-lesson-{digest}"
    lesson_title = str(scope["lesson"].get("node_name") or "本讲课件")
    script_revision_id = str(script_revision.get("revision_id") or "")
    if not script_revision_id:
        raise TeacherLessonAuthoringError(
            "lesson_script_source_missing",
            "PPT 必须使用当前已生成讲义，当前讲义修订标识缺失。",
        )
    script_sections = {
        str(item.get("section_node_id") or ""): normalize_teacher_script_section(item)
        for item in script_revision.get("sections") or []
        if isinstance(item, dict) and item.get("section_node_id")
    }
    required_section_ids = [
        str(item.get("node_id") or "") for item in scope["sections"]
    ]
    missing_sections = [
        section_id
        for section_id in required_section_ids
        if section_id not in script_sections
    ]
    if missing_sections:
        raise TeacherLessonAuthoringError(
            "lesson_script_source_incomplete",
            "当前讲义没有完整覆盖本讲全部小节，不能生成 PPT。",
            details={"missing_section_node_ids": missing_sections},
        )
    lesson_node = deepcopy(scope["lesson"])
    lesson_node.update({
        "node_id": lesson_unit_id,
        "parent_node_id": "root",
        "node_level": 1,
        "node_name": lesson_title,
        "node_content": "",
        "content_blocks": [],
    })
    nodes = [lesson_node]
    for section_index, outline_section in enumerate(scope["sections"], start=1):
        section_id = str(outline_section.get("node_id") or "")
        planned = plan_sections.get(section_id) or {}
        section_content = teacher_lesson_section_content(planned)
        blocks: list[dict[str, Any]] = []
        script_section = script_sections.get(section_id) or {}
        all_script_blocks = [
            item for item in script_section.get("blocks") or []
            if isinstance(item, dict) and str(item.get("content") or "").strip()
        ]
        if not all_script_blocks:
            raise TeacherLessonAuthoringError(
                "lesson_script_source_incomplete",
                "当前讲义仍有空白小节，不能生成 PPT。",
                details={"section_node_id": section_id},
            )
        ppt_group_index = 1
        ppt_group_minutes = 0
        ppt_group_size = 0
        for script_block in all_script_blocks:
            module_id = str(script_block.get("module_id") or "core_explanation")
            role = str(script_block.get("role") or module_block_role(module_id))
            if role not in {
                "orientation", "prerequisite", "objective", "concept", "reasoning",
                "example", "counterexample", "application", "activity", "feedback",
                "misconception", "checkpoint", "remediation", "summary", "transfer",
            }:
                role = "concept"
            knowledge_names = [
                str(item) for item in script_block.get("knowledge_names") or []
                if str(item).strip()
            ]
            try:
                block_minutes = max(
                    1, int(script_block.get("planned_minutes") or 2)
                )
            except (TypeError, ValueError):
                block_minutes = 2
            if ppt_group_size and (
                ppt_group_minutes >= 4 or ppt_group_size >= 3
            ):
                ppt_group_index += 1
                ppt_group_minutes = 0
                ppt_group_size = 0
            ppt_page_group_id = (
                f"{section_id}:ppt-page-group:{ppt_group_index}"
            )
            blocks.append({
                "block_id": str(script_block.get("block_id") or f"{section_id}-{module_id}"),
                "parent_block_id": ppt_page_group_id,
                "type": role,
                "title": str(script_block.get("title") or module_id),
                "content": str(script_block.get("content") or "").strip(),
                "metadata": {
                    "role": role,
                    "module_id": module_id,
                    "module_instance_id": str(
                        script_block.get("block_id") or f"{section_id}:{module_id}"
                    ),
                    "concept_refs": knowledge_names,
                    "planned_minutes": script_block.get("planned_minutes"),
                    "ppt_page_group_id": ppt_page_group_id,
                    "content_perspective": "teacher_delivery",
                    "source_kind": "current_teacher_script_block",
                    "legacy_adapter": module_id == "legacy_script",
                },
            })
            ppt_group_minutes += block_minutes
            ppt_group_size += 1
        node = deepcopy(outline_section)
        node.update({
            "node_id": section_id,
            "parent_node_id": lesson_unit_id,
            "node_level": 2,
            "node_name": str(outline_section.get("node_name") or f"第{section_index}节"),
            "learning_objective": str(
                section_content.get("learning_objective")
                or outline_section.get("learning_objective")
                or ""
            ),
            "knowledge_structure": deepcopy(
                planned.get("knowledge_structure")
                or outline_section.get("knowledge_structure")
                or []
            ),
            "key_points": deepcopy(planned.get("key_points") or []),
            "content_blocks": blocks,
            "node_content": "\n\n".join(
                f"## {block['title']}\n\n{block['content']}" for block in blocks
            ),
        })
        nodes.append(node)
    synthetic = {
        "course_id": synthetic_course_id,
        "course_name": lesson_title,
        "language": str(course_data.get("language") or "zh-CN"),
        "nodes": nodes,
        "course_teaching_plan": plan,
        "course_knowledge_base": deepcopy(course_data.get("course_knowledge_base") or {}),
        "course_coherence_contract": deepcopy(course_data.get("course_coherence_contract") or {}),
        "generation_request": deepcopy(course_data.get("generation_request") or {}),
        "teacher_lesson_source": {
            "real_course_id": str(course_data.get("course_id") or ""),
            "lesson_unit_id": lesson_unit_id,
            "lesson_plan_revision_id": str(plan_revision.get("revision_id") or ""),
            "script_revision_id": script_revision_id,
            "script_content_perspective": "teacher_delivery",
        },
    }
    document = document_from_generation_draft(synthetic)
    return document, synthetic, synthetic_course_id


def project_current_teacher_scripts(
    preview: dict[str, Any],
    authoring_state: dict[str, Any],
) -> dict[str, Any]:
    """Overlay current complete lesson scripts onto the teacher's current outline.

    The teacher preview is a read model, not another publication channel.  It
    keeps the generated outline workspace as the structural truth and exposes
    scripts that still match the current plan. Quality findings stay available
    for review but do not hide an otherwise complete working revision.
    """
    projected = deepcopy(preview)
    nodes = [
        deepcopy(item)
        for item in projected.get("nodes") or []
        if isinstance(item, dict)
    ]
    projected["nodes"] = nodes
    node_ids = {
        str(item.get("node_id") or "")
        for item in nodes
        if item.get("node_id")
    }
    overlay_by_id: dict[str, dict[str, Any]] = {}
    covered_lessons: list[str] = []
    skipped_lessons: list[dict[str, str]] = []
    course_data = {
        **deepcopy(projected),
        "course_id": str(projected.get("course_id") or authoring_state.get("course_id") or ""),
        "nodes": nodes,
    }

    for lesson_id, lesson in (authoring_state.get("lessons") or {}).items():
        if not isinstance(lesson, dict):
            continue
        lesson_id = str(lesson_id or "")
        plan_revision_id = str(lesson.get("working_revision_id") or "")
        script_revision_id = str(lesson.get("working_script_revision_id") or "")
        if (
            not lesson_id
            or lesson_id not in node_ids
            or str(lesson.get("source_state") or "") != "current"
            or not plan_revision_id
            or not script_revision_id
        ):
            skipped_lessons.append({"lesson_unit_id": lesson_id, "reason": "source_not_current"})
            continue
        plan_revision = next(
            (
                item for item in lesson.get("revisions") or []
                if isinstance(item, dict)
                and str(item.get("revision_id") or "") == plan_revision_id
            ),
            None,
        )
        script_revision = next(
            (
                item for item in lesson.get("script_revisions") or []
                if isinstance(item, dict)
                and str(item.get("revision_id") or "") == script_revision_id
            ),
            None,
        )
        if (
            not isinstance(plan_revision, dict)
            or not isinstance(script_revision, dict)
            or str(script_revision.get("source_lesson_plan_revision_id") or "")
            != plan_revision_id
        ):
            skipped_lessons.append({"lesson_unit_id": lesson_id, "reason": "source_not_current"})
            continue
        try:
            _document, lesson_view, _synthetic_id = teacher_lesson_v6_source(
                course_data,
                lesson_unit_id=lesson_id,
                plan_revision=plan_revision,
                script_revision=script_revision,
            )
        except TeacherLessonAuthoringError as exc:
            skipped_lessons.append({"lesson_unit_id": lesson_id, "reason": exc.code})
            continue
        lesson_nodes = [
            item for item in lesson_view.get("nodes") or []
            if isinstance(item, dict) and int(item.get("node_level") or 0) == 2
        ]
        if not lesson_nodes or any(
            str(item.get("node_id") or "") not in node_ids
            for item in lesson_nodes
        ):
            skipped_lessons.append({"lesson_unit_id": lesson_id, "reason": "outline_scope_mismatch"})
            continue
        for item in lesson_nodes:
            overlay_by_id[str(item.get("node_id") or "")] = item
        covered_lessons.append(lesson_id)

    for index, node in enumerate(nodes):
        overlay = overlay_by_id.get(str(node.get("node_id") or ""))
        if not overlay:
            continue
        nodes[index] = {
            **node,
            "learning_objective": str(overlay.get("learning_objective") or node.get("learning_objective") or ""),
            "knowledge_structure": deepcopy(overlay.get("knowledge_structure") or []),
            "key_points": deepcopy(overlay.get("key_points") or []),
            "node_content": str(overlay.get("node_content") or ""),
            "content_blocks": deepcopy(overlay.get("content_blocks") or []),
            "generation_status": "completed",
            "content_state": "finalized",
            "generated_chars": len(str(overlay.get("node_content") or "")),
            "error_summary": None,
        }

    projected["projection"] = "teacher_lesson_authoring"
    projected["teacher_lesson_projection"] = {
        "schema_version": "teacher_lesson_preview_v1",
        "outline_revision_id": str(authoring_state.get("outline_revision_id") or ""),
        "covered_lesson_unit_ids": covered_lessons,
        "covered_section_count": len(overlay_by_id),
        "skipped_lessons": skipped_lessons,
    }
    return projected


def project_confirmed_teacher_scripts(
    preview: dict[str, Any],
    authoring_state: dict[str, Any],
) -> dict[str, Any]:
    """Backward-compatible name for the former confirmation-based projection."""
    return project_current_teacher_scripts(preview, authoring_state)


TEACHER_JOB_ACTIVE_STATUSES = frozenset({"pending", "running"})
TEACHER_JOB_FROZEN_STATUSES = frozenset({
    "paused",
    "failed",
    "cancelled",
    "completed_with_warnings",
    "completed",
})


class TeacherLessonAuthoringRepository:
    def __init__(self, root: str | Path | None = None):
        self.root = Path(root) if root is not None else _default_root()
        self.root.mkdir(parents=True, exist_ok=True)
        self._lock = threading.RLock()
        self._live_stream_jobs: dict[str, dict[str, dict[str, Any]]] = {}
        self._live_stream_touched_at: dict[tuple[str, str], float] = {}

    def _path(self, course_id: str) -> Path:
        safe = "".join(char for char in course_id if char.isalnum() or char in {"-", "_"})
        if not safe or safe != course_id:
            raise TeacherLessonAuthoringError("invalid_course_id", "课程标识无效。")
        return self.root / f"{safe}.json"

    def _empty(self, course_id: str) -> dict[str, Any]:
        return {
            "schema_version": SCHEMA_VERSION,
            "course_id": course_id,
            "revision": 0,
            "outline_revision_id": "",
            "outline_material_drafts": [],
            "current_outline_material_draft_id": "",
            "material_absorptions": [],
            "lessons": {},
            "jobs": {},
            "updated_at": _now(),
        }

    def _overlay_live_stream_jobs_locked(
        self,
        course_id: str,
        value: dict[str, Any],
    ) -> dict[str, Any]:
        live_jobs = self._live_stream_jobs.get(course_id) or {}
        if not live_jobs:
            return value
        jobs = value.get("jobs")
        if not isinstance(jobs, dict):
            jobs = {}
            value["jobs"] = jobs
        for job_id, live_job in live_jobs.items():
            jobs[job_id] = deepcopy(live_job)
        return value

    def _drop_live_stream_job_locked(self, course_id: str, job_id: str) -> None:
        live_jobs = self._live_stream_jobs.get(course_id)
        if live_jobs is not None:
            live_jobs.pop(job_id, None)
            if not live_jobs:
                self._live_stream_jobs.pop(course_id, None)
        self._live_stream_touched_at.pop((course_id, job_id), None)
    def load(self, course_id: str) -> dict[str, Any]:
        with self._lock:
            path = self._path(course_id)
            if not path.exists():
                return self._overlay_live_stream_jobs_locked(
                    course_id,
                    self._empty(course_id),
                )
            try:
                data = json.loads(path.read_text(encoding="utf-8"))
            except (OSError, json.JSONDecodeError) as exc:
                raise TeacherLessonAuthoringError(
                    "authoring_repository_corrupt",
                    "教师讲次资产读取失败。",
                ) from exc
            value = data if isinstance(data, dict) else self._empty(course_id)
            return self._overlay_live_stream_jobs_locked(course_id, value)

    def _save(self, value: dict[str, Any]) -> dict[str, Any]:
        course_id = str(value.get("course_id") or "")
        path = self._path(course_id)
        payload = deepcopy(value)
        for job in (payload.get("jobs") or {}).values():
            if not isinstance(job, dict):
                continue
            # Raw model deltas are only a same-process SSE projection. Durable
            # recovery is based on validated plan batches and teaching blocks.
            job["stream_batches"] = {}
            job["stream_events"] = []
            job["last_stream_event"] = {}
        payload["schema_version"] = SCHEMA_VERSION
        payload["revision"] = int(payload.get("revision") or 0) + 1
        payload["updated_at"] = _now()
        fd, temp_name = tempfile.mkstemp(prefix=f".{course_id}.", suffix=".tmp", dir=str(self.root))
        try:
            with os.fdopen(fd, "w", encoding="utf-8") as handle:
                json.dump(payload, handle, ensure_ascii=False, indent=2)
                handle.flush()
                os.fsync(handle.fileno())
            os.replace(temp_name, path)
        finally:
            if os.path.exists(temp_name):
                os.unlink(temp_name)
        return payload

    def set_outline(self, course_id: str, outline_revision_id: str) -> dict[str, Any]:
        with self._lock:
            value = self.load(course_id)
            original = deepcopy(value)
            value["outline_revision_id"] = outline_revision_id
            for lesson in (value.get("lessons") or {}).values():
                if not isinstance(lesson, dict):
                    continue
                arrangement = lesson.get("arrangement") or {}
                arrangement_revision = next(
                    (
                        item for item in arrangement.get("revisions") or []
                        if isinstance(item, dict)
                        and item.get("revision_id") == arrangement.get("working_revision_id")
                    ),
                    None,
                )
                arrangement["source_state"] = (
                    "current"
                    if not arrangement_revision
                    or str(arrangement_revision.get("source_outline_revision_id") or "") == outline_revision_id
                    else "stale"
                )
                lesson["arrangement"] = arrangement
                if not lesson.get("working_revision_id"):
                    continue
                working = next(
                    (
                        item for item in lesson.get("revisions") or []
                        if isinstance(item, dict)
                        and item.get("revision_id") == lesson.get("working_revision_id")
                    ),
                    None,
                )
                source_arrangement_revision_id = str(
                    (working or {}).get("source_arrangement_revision_id") or ""
                )
                current_arrangement_revision_id = str(
                    arrangement.get("working_revision_id") or ""
                )
                arrangement_matches = (
                    not source_arrangement_revision_id
                    or source_arrangement_revision_id == current_arrangement_revision_id
                )
                current = bool(
                    isinstance(working, dict)
                    and str(working.get("source_outline_revision_id") or "") == outline_revision_id
                    and arrangement_matches
                    and lesson.get("source_state_reason") != "arrangement_changed"
                )
                lesson["source_state"] = "current" if current else "stale"
                if current:
                    lesson.pop("source_state_reason", None)
                else:
                    reason = str(lesson.get("source_state_reason") or "outline_changed")
                    _mark_lesson_dependents_stale(lesson, reason=reason)
            if value == original:
                return deepcopy(value)
            return self._save(value)

    def save_arrangement_revision(
        self,
        course_id: str,
        lesson_unit_id: str,
        arrangement: dict[str, Any],
        *,
        source_outline_revision_id: str,
        actor: str = "teacher",
        confirm: bool = True,
    ) -> dict[str, Any]:
        with self._lock:
            value = self.load(course_id)
            lesson = value.setdefault("lessons", {}).setdefault(
                lesson_unit_id,
                _empty_lesson_asset(lesson_unit_id),
            )
            state = lesson.setdefault("arrangement", {
                "working_revision_id": "",
                "confirmed_revision_id": "",
                "source_state": "current",
                "revisions": [],
            })
            previous_working_revision_id = str(
                state.get("working_revision_id") or ""
            )
            normalized = normalize_lesson_arrangement(
                arrangement,
                lesson_unit_id=lesson_unit_id,
                source_outline_revision_id=source_outline_revision_id,
            )
            revision_id = f"tlar-{uuid.uuid4().hex}"
            revision = {
                **normalized,
                "revision_id": revision_id,
                "status": "confirmed" if confirm else "draft",
                "actor": actor,
                "created_at": _now(),
                **({"confirmed_at": _now()} if confirm else {}),
            }
            state.setdefault("revisions", []).append(revision)
            state["working_revision_id"] = revision_id
            state["source_state"] = (
                "current"
                if not value.get("outline_revision_id")
                or str(value.get("outline_revision_id") or "") == source_outline_revision_id
                else "stale"
            )
            if confirm:
                state["confirmed_revision_id"] = revision_id
            lesson["arrangement"] = state
            if previous_working_revision_id and previous_working_revision_id != revision_id:
                _mark_lesson_dependents_stale(
                    lesson,
                    reason="arrangement_changed",
                )
            if source_outline_revision_id and not value.get("outline_revision_id"):
                value["outline_revision_id"] = source_outline_revision_id
            saved = self._save(value)
            return deepcopy(saved["lessons"][lesson_unit_id])

    def current_arrangement(
        self,
        course_id: str,
        lesson_unit_id: str,
    ) -> dict[str, Any] | None:
        """Return the current structurally usable arrangement working revision."""
        lesson = self.lesson(course_id, lesson_unit_id)
        state = lesson.get("arrangement") or {}
        revision_id = str(state.get("working_revision_id") or "")
        if not revision_id or state.get("source_state", "current") != "current":
            return None
        revision = next(
            (
                item for item in state.get("revisions") or []
                if isinstance(item, dict) and item.get("revision_id") == revision_id
            ),
            None,
        )
        if not isinstance(revision, dict) or not list(revision.get("blocks") or []):
            return None
        return deepcopy(revision)

    def confirmed_arrangement(
        self,
        course_id: str,
        lesson_unit_id: str,
    ) -> dict[str, Any] | None:
        lesson = self.lesson(course_id, lesson_unit_id)
        state = lesson.get("arrangement") or {}
        revision_id = str(state.get("confirmed_revision_id") or "")
        if not revision_id or state.get("source_state", "current") != "current":
            return None
        return deepcopy(next(
            (
                item for item in state.get("revisions") or []
                if isinstance(item, dict) and item.get("revision_id") == revision_id
            ),
            None,
        ))

    def create_job(
        self,
        course_id: str,
        lesson_unit_id: str,
        *,
        job_type: str = "teacher_lesson_plan_generation",
        request_id: str = "",
        source_outline_revision_id: str = "",
    ) -> dict[str, Any]:
        if job_type not in JOB_TYPES:
            raise TeacherLessonAuthoringError("unsupported_teacher_job", "不支持的教师讲次任务。")
        with self._lock:
            value = self.load(course_id)
            if request_id:
                existing = next(
                    (
                        job for job in (value.get("jobs") or {}).values()
                        if isinstance(job, dict)
                        and job.get("request_id") == request_id
                        and job.get("lesson_unit_id") == lesson_unit_id
                        and job.get("type") == job_type
                    ),
                    None,
                )
                if existing:
                    return deepcopy(existing)
            job_id = f"tlj-{uuid.uuid4().hex}"
            initial_message = (
                "等待生成本讲讲义"
                if job_type == "teacher_lesson_script_generation"
                else "等待生成本讲教案"
            )
            job = {
                "schema_version": TEACHER_ASSET_JOB_SCHEMA_VERSION,
                "id": job_id,
                "course_id": course_id,
                "lesson_unit_id": lesson_unit_id,
                "lesson_id": lesson_unit_id,
                "type": job_type,
                "asset_type": (
                    "script"
                    if job_type == "teacher_lesson_script_generation"
                    else "lesson_plan"
                ),
                "state_owner": "teacher_lesson_authoring",
                "request_id": request_id,
                "idempotency_key": request_id,
                "source_outline_revision_id": source_outline_revision_id,
                "status": "pending",
                "progress": 0,
                "phase": "queued",
                "stage": "queued",
                "message": initial_message,
                "stream_sequence": 0,
                "stream_batches": {},
                "stream_mode": "",
                "stream_events": [],
                "last_stream_event": {},
                "stream_complete": False,
                "checkpoint": {},
                "cancel_requested": False,
                "retryable": True,
                "warnings": [],
                "error": None,
                "created_at": _now(),
                "heartbeat_at": _now(),
                "updated_at": _now(),
            }
            value.setdefault("jobs", {})[job_id] = job
            self._save(value)
            return deepcopy(job)

    def _apply_job_changes_locked(
        self,
        job: dict[str, Any],
        changes: dict[str, Any],
    ) -> dict[str, Any]:
        job.update(deepcopy(changes))
        if "phase" in changes:
            job["stage"] = str(changes.get("phase") or "")
        timestamp = _now()
        job["updated_at"] = timestamp
        if str(job.get("status") or "") in {"pending", "running"}:
            job["heartbeat_at"] = timestamp
        if "result_sections" in changes:
            job["checkpoint"] = {
                "result_sections": deepcopy(changes.get("result_sections") or []),
                "completed_blocks": int(changes.get("completed_blocks") or 0),
                "current_block_id": str(changes.get("current_block_id") or ""),
                "saved_at": timestamp,
            }
        if str(job.get("status") or "") in {
            "completed", "completed_with_warnings", "failed", "cancelled"
        }:
            job["completed_at"] = timestamp
            job["retryable"] = bool(
                (job.get("error") or {}).get("retryable")
                if isinstance(job.get("error"), dict)
                else False
            )
        return job

    def update_job(self, course_id: str, job_id: str, **changes: Any) -> dict[str, Any]:
        """Persist a semantic checkpoint or a lifecycle state transition."""
        with self._lock:
            value = self.load(course_id)
            job = (value.get("jobs") or {}).get(job_id)
            if not isinstance(job, dict):
                raise TeacherLessonAuthoringError("teacher_job_not_found", "教师讲次任务不存在。")
            if str(job.get("status") or "") in TEACHER_JOB_FROZEN_STATUSES:
                return deepcopy(job)
            job = self._apply_job_changes_locked(job, changes)
            value["jobs"][job_id] = job
            saved = self._save(value)
            if str(job.get("status") or "") in TEACHER_JOB_FROZEN_STATUSES:
                self._drop_live_stream_job_locked(course_id, job_id)
            elif job_id in (self._live_stream_jobs.get(course_id) or {}):
                # A validated checkpoint must not erase token streams that are
                # still arriving concurrently for other plan/script shards.
                self._live_stream_jobs[course_id][job_id] = job
            return deepcopy(saved["jobs"][job_id])

    def update_job_live(
        self,
        course_id: str,
        job_id: str,
        **changes: Any,
    ) -> dict[str, Any]:
        """Publish progress and heartbeat changes without rewriting course JSON."""
        with self._lock:
            live_jobs = self._live_stream_jobs.setdefault(course_id, {})
            job = live_jobs.get(job_id)
            if not isinstance(job, dict):
                value = self.load(course_id)
                job = (value.get("jobs") or {}).get(job_id)
            if not isinstance(job, dict):
                raise TeacherLessonAuthoringError("teacher_job_not_found", "教师讲次任务不存在。")
            if str(job.get("status") or "") not in TEACHER_JOB_ACTIVE_STATUSES:
                return deepcopy(job)
            job = self._apply_job_changes_locked(job, changes)
            live_jobs[job_id] = job
            self._live_stream_touched_at[(course_id, job_id)] = time.monotonic()
            return deepcopy(job)

    def update_job_stream(
        self,
        course_id: str,
        job_id: str,
        *,
        phase: str,
        progress: int,
        message: str,
        batch_id: str,
        event: str,
        delta: str = "",
        lesson_unit_id: str = "",
        block_id: str = "",
        shard_id: str = "",
        stream_mode: str = "",
    ) -> dict[str, Any]:
        """Publish one model delta to the same-process SSE projection only."""
        with self._lock:
            live_jobs = self._live_stream_jobs.setdefault(course_id, {})
            job = live_jobs.get(job_id)
            if not isinstance(job, dict):
                value = self.load(course_id)
                job = (value.get("jobs") or {}).get(job_id)
            if not isinstance(job, dict):
                raise TeacherLessonAuthoringError("teacher_job_not_found", "教师讲次任务不存在。")
            if str(job.get("status") or "") not in TEACHER_JOB_ACTIVE_STATUSES:
                return deepcopy(job)
            batches = deepcopy(job.get("stream_batches") or {})
            if event == "reset":
                batches[batch_id] = ""
            elif event == "delta":
                batches[batch_id] = (
                    str(batches.get(batch_id) or "") + str(delta or "")
                )[-200_000:]
            sequence = int(job.get("stream_sequence") or 0) + 1
            stream_event = {
                "event": event,
                "lesson_unit_id": str(
                    lesson_unit_id or job.get("lesson_unit_id") or ""
                ),
                "block_id": str(block_id or ""),
                "shard_id": str(shard_id or batch_id or ""),
                "sequence": sequence,
                "delta": str(delta or ""),
            }
            stream_events = [
                item for item in job.get("stream_events") or []
                if isinstance(item, dict)
            ]
            stream_events.append(stream_event)
            stream_events = stream_events[-500:]
            timestamp = _now()
            job.update({
                "phase": phase,
                "stage": phase,
                "progress": progress,
                "message": message,
                "stream_sequence": sequence,
                "stream_batches": batches,
                "stream_mode": str(stream_mode or job.get("stream_mode") or ""),
                "stream_events": stream_events,
                "last_stream_event": stream_event,
                "stream_complete": False,
                "heartbeat_at": timestamp,
                "updated_at": timestamp,
            })
            live_jobs[job_id] = job
            self._live_stream_touched_at[(course_id, job_id)] = time.monotonic()
            return deepcopy(job)

    def cancel_job(self, course_id: str, job_id: str) -> dict[str, Any]:
        """Cancel one durable teacher job while preserving its last checkpoint."""
        with self._lock:
            value = self.load(course_id)
            job = (value.get("jobs") or {}).get(job_id)
            if not isinstance(job, dict):
                raise TeacherLessonAuthoringError(
                    "teacher_job_not_found",
                    "教师讲次任务不存在。",
                )
            if str(job.get("status") or "") in TEACHER_JOB_FROZEN_STATUSES:
                return deepcopy(job)
            timestamp = _now()
            job.update({
                "status": "cancelled",
                "phase": "teacher_asset_job_cancelled",
                "message": "已停止生成，已完成的内容仍然保留",
                "cancel_requested": True,
                "stream_sequence": int(job.get("stream_sequence") or 0) + 1,
                "stream_complete": True,
                "retryable": True,
                "error": {
                    "code": "teacher_asset_job_cancelled",
                    "message": "生成已由教师停止，可以从已保存进度继续。",
                    "retryable": True,
                },
                "completed_at": timestamp,
                "updated_at": timestamp,
            })
            value["jobs"][job_id] = job
            saved = self._save(value)
            self._drop_live_stream_job_locked(course_id, job_id)
            return deepcopy(saved["jobs"][job_id])

    def pause_job(self, course_id: str, job_id: str) -> dict[str, Any]:
        """Pause one job at its next safe checkpoint without discarding progress."""
        with self._lock:
            value = self.load(course_id)
            job = (value.get("jobs") or {}).get(job_id)
            if not isinstance(job, dict):
                raise TeacherLessonAuthoringError(
                    "teacher_job_not_found",
                    "教师讲次任务不存在。",
                )
            if str(job.get("status") or "") not in {"pending", "running"}:
                return deepcopy(job)
            timestamp = _now()
            job.update({
                "status": "paused",
                "phase": "teacher_asset_job_paused",
                "message": "生成已暂停，继续时将沿用已保存进度",
                "pause_requested": True,
                "cancel_requested": True,
                "stream_sequence": int(job.get("stream_sequence") or 0) + 1,
                "stream_complete": True,
                "retryable": True,
                "error": None,
                "updated_at": timestamp,
            })
            value["jobs"][job_id] = job
            saved = self._save(value)
            self._drop_live_stream_job_locked(course_id, job_id)
            return deepcopy(saved["jobs"][job_id])

    def save_plan_revision(
        self,
        course_id: str,
        lesson_unit_id: str,
        plan: dict[str, Any],
        *,
        source_outline_revision_id: str,
        source_knowledge_scope_revision_id: str = "",
        generation_source: str = "model",
        warnings: list[dict[str, Any]] | None = None,
        source_refs: list[dict[str, Any]] | None = None,
        quality_report: dict[str, Any] | None = None,
        actor: str = "teacher",
        restored_from_revision_id: str = "",
        active_job_id: str = "",
    ) -> dict[str, Any]:
        with self._lock:
            value = self.load(course_id)
            if active_job_id:
                active_job = (value.get("jobs") or {}).get(active_job_id)
                if (
                    not isinstance(active_job, dict)
                    or str(active_job.get("status") or "") not in TEACHER_JOB_ACTIVE_STATUSES
                ):
                    raise TeacherLessonAuthoringError(
                        "teacher_job_not_active",
                        "任务已停止，迟到的教案结果未保存。",
                    )
            normalized_plan = normalize_teacher_lesson_plan(plan)
            effective_quality = deepcopy(
                quality_report or validate_teacher_lesson_plan(normalized_plan)
            )
            lesson = value.setdefault("lessons", {}).setdefault(
                lesson_unit_id,
                _empty_lesson_asset(lesson_unit_id),
            )
            arrangement_state = lesson.get("arrangement") or {}
            source_arrangement_revision_id = str(
                arrangement_state.get("working_revision_id") or ""
            )
            revision_id = f"tlpr-{uuid.uuid4().hex}"
            revision = {
                "revision_id": revision_id,
                "lesson_unit_id": lesson_unit_id,
                "source_outline_revision_id": source_outline_revision_id,
                "source_knowledge_scope_revision_id": source_knowledge_scope_revision_id,
                "source_arrangement_revision_id": source_arrangement_revision_id,
                "generation_source": generation_source,
                "status": "draft",
                "warnings": deepcopy(warnings or []),
                "source_refs": deepcopy(source_refs or []),
                "pipeline_version": LESSON_PLAN_PIPELINE_VERSION,
                "quality_report": effective_quality,
                "plan": normalized_plan,
                "actor": actor,
                "created_at": _now(),
            }
            if restored_from_revision_id:
                revision["restored_from_revision_id"] = restored_from_revision_id
            lesson.setdefault("revisions", []).append(revision)
            lesson["working_revision_id"] = revision_id
            lesson["source_state"] = (
                "current"
                if not value.get("outline_revision_id")
                or str(value.get("outline_revision_id") or "") == source_outline_revision_id
                else "stale"
            )
            if arrangement_state.get("source_state", "current") != "current":
                lesson["source_state"] = "stale"
                lesson["source_state_reason"] = "arrangement_stale"
            else:
                lesson.pop("source_state_reason", None)
            for asset in lesson.get("ppt_assets") or []:
                if not isinstance(asset, dict):
                    continue
                source_revision = str(asset.get("source_lesson_plan_revision_id") or "")
                if source_revision and source_revision != revision_id:
                    asset["source_state"] = "stale"
            manuscript = lesson.get("ppt_manuscript")
            if (
                isinstance(manuscript, dict)
                and manuscript.get("source_lesson_plan_revision_id")
                != revision_id
            ):
                manuscript["source_state"] = "stale"
            for review in lesson.get("imported_ppt_reviews") or []:
                if not isinstance(review, dict):
                    continue
                source_revision = str(review.get("source_lesson_plan_revision_id") or "")
                if source_revision and source_revision != revision_id:
                    review["source_state"] = "stale"
            if source_outline_revision_id and not value.get("outline_revision_id"):
                value["outline_revision_id"] = source_outline_revision_id
            saved = self._save(value)
            return deepcopy(saved["lessons"][lesson_unit_id])

    def restore_plan_revision(
        self,
        course_id: str,
        lesson_unit_id: str,
        revision_id: str,
        *,
        expected_working_revision_id: str,
        actor: str = "teacher",
    ) -> dict[str, Any]:
        with self._lock:
            lesson = self.lesson(course_id, lesson_unit_id)
            current_revision_id = str(lesson.get("working_revision_id") or "")
            if current_revision_id != expected_working_revision_id:
                raise TeacherLessonAuthoringError(
                    "lesson_plan_revision_conflict",
                    "教案已在其他页面修改，请重新载入后再恢复。",
                )
            source = next(
                (
                    item for item in lesson.get("revisions") or []
                    if isinstance(item, dict) and item.get("revision_id") == revision_id
                ),
                None,
            )
            if not isinstance(source, dict):
                raise TeacherLessonAuthoringError(
                    "lesson_plan_revision_not_found",
                    "教案历史版本不存在。",
                )
            return self.save_plan_revision(
                course_id,
                lesson_unit_id,
                deepcopy(source.get("plan") or {}),
                source_outline_revision_id=str(source.get("source_outline_revision_id") or ""),
                source_knowledge_scope_revision_id=str(
                    source.get("source_knowledge_scope_revision_id") or ""
                ),
                generation_source="history_restore",
                warnings=deepcopy(source.get("warnings") or []),
                source_refs=deepcopy(source.get("source_refs") or []),
                actor=actor,
                restored_from_revision_id=revision_id,
            )

    def bind_v6_ppt_revision(
        self,
        course_id: str,
        lesson_unit_id: str,
        *,
        source_lesson_plan_revision_id: str,
        source_script_revision_id: str,
        synthetic_course_id: str,
        representation_id: str,
        spec_id: str,
        candidate_status: str,
        ppt_manuscript_revision: str = "",
        ppt_manuscript_status: str = "draft",
    ) -> dict[str, Any]:
        """Register one real V6 representation without copying it into student data."""
        with self._lock:
            value = self.load(course_id)
            lesson = (value.get("lessons") or {}).get(lesson_unit_id)
            if not isinstance(lesson, dict):
                raise TeacherLessonAuthoringError("lesson_plan_not_found", "请先生成本讲教案。")
            if (
                lesson.get("working_revision_id") != source_lesson_plan_revision_id
                or lesson.get("source_state", "current") != "current"
            ):
                raise TeacherLessonAuthoringError(
                    "lesson_plan_revision_conflict",
                    "当前教案已经变化，V6 结果未登记。",
                )
            assets = lesson.setdefault("ppt_assets", [])
            asset = next(
                (item for item in assets if isinstance(item, dict) and item.get("role") == "primary"),
                None,
            )
            if asset is None:
                asset = {
                    "asset_id": f"tlpa-{uuid.uuid4().hex}",
                    "lesson_unit_id": lesson_unit_id,
                    "role": "primary",
                    "working_revision_id": "",
                    "source_lesson_plan_revision_id": source_lesson_plan_revision_id,
                    "source_state": "current",
                    "revisions": [],
                    "ai_candidates": [],
                }
                assets.append(asset)
            binding = {
                "revision_id": f"tlv6r-{uuid.uuid4().hex}",
                "engine": "slide_deck_v6",
                "synthetic_course_id": synthetic_course_id,
                "representation_id": representation_id,
                "spec_id": spec_id,
                "source_lesson_plan_revision_id": source_lesson_plan_revision_id,
                "source_script_revision_id": source_script_revision_id,
                "ppt_manuscript_revision": ppt_manuscript_revision,
                "ppt_manuscript_status": ppt_manuscript_status,
                "candidate_status": candidate_status,
                "created_at": _now(),
            }
            asset.setdefault("v6_revisions", []).append(binding)
            asset["engine"] = "slide_deck_v6"
            asset["working_v6_revision_id"] = binding["revision_id"]
            asset["working_representation_id"] = representation_id
            asset["synthetic_course_id"] = synthetic_course_id
            asset["source_lesson_plan_revision_id"] = source_lesson_plan_revision_id
            asset["source_script_revision_id"] = source_script_revision_id
            if ppt_manuscript_revision:
                asset["ppt_manuscript_revision"] = ppt_manuscript_revision
                asset["ppt_manuscript_status"] = ppt_manuscript_status
            asset["source_state"] = "current"
            saved = self._save(value)
            return deepcopy(next(item for item in saved["lessons"][lesson_unit_id]["ppt_assets"] if item["asset_id"] == asset["asset_id"]))

    def restore_v6_ppt_revision(
        self,
        course_id: str,
        lesson_unit_id: str,
        revision_id: str,
        *,
        expected_working_revision_id: str,
    ) -> dict[str, Any]:
        """Move the V6 working pointer back to an existing binding.

        A restore must not create a new binding: the historical binding keeps
        the exact plan/script/spec relationship that was active before the
        course-wide change.  The expected pointer prevents an undo from
        overwriting a later teacher edit.
        """

        with self._lock:
            value = self.load(course_id)
            lesson = (value.get("lessons") or {}).get(lesson_unit_id)
            if not isinstance(lesson, dict):
                raise TeacherLessonAuthoringError("lesson_plan_not_found", "请先生成本讲教案。")
            asset = next(
                (
                    item for item in lesson.get("ppt_assets") or []
                    if isinstance(item, dict)
                    and item.get("role") == "primary"
                    and item.get("engine") == "slide_deck_v6"
                ),
                None,
            )
            if not isinstance(asset, dict):
                raise TeacherLessonAuthoringError("ppt_revision_not_found", "PPT 历史版本不存在。")
            if str(asset.get("working_v6_revision_id") or "") != expected_working_revision_id:
                raise TeacherLessonAuthoringError(
                    "ppt_revision_conflict",
                    "PPT 已在其他页面修改，请重新载入后再恢复。",
                )
            binding = next(
                (
                    item for item in asset.get("v6_revisions") or []
                    if isinstance(item, dict) and str(item.get("revision_id") or "") == revision_id
                ),
                None,
            )
            if not isinstance(binding, dict):
                raise TeacherLessonAuthoringError("ppt_revision_not_found", "PPT 历史版本不存在。")
            asset["working_v6_revision_id"] = revision_id
            asset["working_representation_id"] = str(binding.get("representation_id") or "")
            asset["synthetic_course_id"] = str(binding.get("synthetic_course_id") or "")
            asset["source_lesson_plan_revision_id"] = str(
                binding.get("source_lesson_plan_revision_id") or ""
            )
            asset["source_script_revision_id"] = str(binding.get("source_script_revision_id") or "")
            asset["ppt_manuscript_revision"] = str(binding.get("ppt_manuscript_revision") or "")
            asset["ppt_manuscript_status"] = str(binding.get("ppt_manuscript_status") or "draft")
            asset["source_state"] = (
                "current"
                if str(lesson.get("working_revision_id") or "")
                == str(binding.get("source_lesson_plan_revision_id") or "")
                else "stale"
            )
            saved = self._save(value)
            return deepcopy(next(
                item
                for item in saved["lessons"][lesson_unit_id]["ppt_assets"]
                if item["asset_id"] == asset["asset_id"]
            ))

    def save_v6_ppt_manuscript(
        self,
        course_id: str,
        lesson_unit_id: str,
        *,
        source_lesson_plan_revision_id: str,
        source_script_revision_id: str,
        source_material_revision: str,
        task_id: str,
        mode: str,
        theme: str,
        manuscript: dict[str, Any],
        template_id: str = "",
        template_version: str = "",
        template_digest: str = "",
        template_pack_id: str = "",
    ) -> dict[str, Any]:
        """保存无原版 PPT 分支的独立页面内容稿工作稿，不提前创建 PPT 资产。"""
        with self._lock:
            value = self.load(course_id)
            lesson = (value.get("lessons") or {}).get(lesson_unit_id)
            if not isinstance(lesson, dict):
                raise TeacherLessonAuthoringError(
                    "lesson_plan_not_found", "请先生成本讲教案。"
                )
            if (
                lesson.get("working_revision_id")
                != source_lesson_plan_revision_id
                or lesson.get("source_state", "current") != "current"
                or lesson.get("working_script_revision_id")
                != source_script_revision_id
            ):
                raise TeacherLessonAuthoringError(
                    "lesson_ppt_source_stale",
                    "教案或讲义已经变化，请基于最新内容重新生成 页面内容稿。",
                )
            state = {
                "revision": str(manuscript.get("manuscript_revision") or ""),
                "status": "draft",
                "source_state": "current",
                "source_lesson_plan_revision_id": source_lesson_plan_revision_id,
                "source_script_revision_id": source_script_revision_id,
                "source_material_revision": source_material_revision,
                "task_id": task_id,
                "mode": mode,
                "theme": theme,
                "template_id": template_id,
                "template_version": template_version,
                "template_digest": template_digest,
                "template_pack_id": template_pack_id,
                "manuscript": deepcopy(manuscript),
                "created_at": _now(),
                "confirmed_at": "",
                "generated_representation_id": "",
            }
            lesson["ppt_manuscript"] = state
            saved = self._save(value)
            return deepcopy(
                saved["lessons"][lesson_unit_id]["ppt_manuscript"]
            )

    def current_v6_ppt_manuscript(
        self, course_id: str, lesson_unit_id: str
    ) -> dict[str, Any] | None:
        lesson = self.lesson(course_id, lesson_unit_id)
        state = lesson.get("ppt_manuscript")
        return deepcopy(state) if isinstance(state, dict) and state else None

    def confirm_v6_ppt_manuscript_draft(
        self,
        course_id: str,
        lesson_unit_id: str,
        *,
        manuscript_revision: str,
    ) -> dict[str, Any]:
        """确认独立 页面内容稿；确认后才可进入 PPT 编译。"""
        with self._lock:
            value = self.load(course_id)
            lesson = (value.get("lessons") or {}).get(lesson_unit_id)
            state = (lesson or {}).get("ppt_manuscript")
            if not isinstance(state, dict) or not state:
                raise TeacherLessonAuthoringError(
                    "lesson_ppt_manuscript_not_found", "请先生成 页面内容稿。"
                )
            if state.get("source_state") != "current":
                raise TeacherLessonAuthoringError(
                    "lesson_ppt_source_stale",
                    "教案或讲义已经变化，请重新生成 页面内容稿。",
                )
            if str(state.get("revision") or "") != manuscript_revision:
                raise TeacherLessonAuthoringError(
                    "lesson_ppt_manuscript_revision_conflict",
                    "页面内容稿已更新，请刷新后再确认。",
                )
            state["status"] = "confirmed"
            state["confirmed_at"] = _now()
            saved = self._save(value)
            return deepcopy(
                saved["lessons"][lesson_unit_id]["ppt_manuscript"]
            )

    def bind_v6_ppt_manuscript_result(
        self,
        course_id: str,
        lesson_unit_id: str,
        *,
        manuscript_revision: str,
        representation_id: str,
    ) -> dict[str, Any]:
        with self._lock:
            value = self.load(course_id)
            lesson = (value.get("lessons") or {}).get(lesson_unit_id)
            state = (lesson or {}).get("ppt_manuscript")
            if (
                not isinstance(state, dict)
                or state.get("status") != "confirmed"
                or state.get("source_state") != "current"
                or state.get("revision") != manuscript_revision
            ):
                raise TeacherLessonAuthoringError(
                    "lesson_ppt_manuscript_not_confirmed",
                    "页面内容稿尚未确认，不能登记生成结果。",
                )
            state["generated_representation_id"] = representation_id
            state["generated_at"] = _now()
            saved = self._save(value)
            return deepcopy(
                saved["lessons"][lesson_unit_id]["ppt_manuscript"]
            )

    def confirm_v6_ppt_manuscript(
        self,
        course_id: str,
        lesson_unit_id: str,
        *,
        representation_id: str,
        manuscript_revision: str,
    ) -> dict[str, Any]:
        """确认逐页 页面内容稿，作为正式导出的显式门。"""
        with self._lock:
            value = self.load(course_id)
            lesson = (value.get("lessons") or {}).get(lesson_unit_id)
            if not isinstance(lesson, dict):
                raise TeacherLessonAuthoringError(
                    "lesson_plan_not_found", "请先生成本讲教案。"
                )
            asset = next(
                (
                    item
                    for item in lesson.get("ppt_assets") or []
                    if isinstance(item, dict)
                    and item.get("working_representation_id") == representation_id
                ),
                None,
            )
            if not isinstance(asset, dict):
                raise TeacherLessonAuthoringError(
                    "lesson_ppt_not_found", "本讲 页面内容稿不存在。"
                )
            if asset.get("source_state") != "current":
                raise TeacherLessonAuthoringError(
                    "lesson_ppt_source_stale", "讲义或教案已更新，请先重新生成 页面内容稿。"
                )
            if str(asset.get("ppt_manuscript_revision") or "") != manuscript_revision:
                raise TeacherLessonAuthoringError(
                    "lesson_ppt_manuscript_revision_conflict",
                    "页面内容稿已更新，请刷新后再确认。",
                )
            asset["ppt_manuscript_status"] = "confirmed"
            for revision in asset.get("v6_revisions") or []:
                if (
                    isinstance(revision, dict)
                    and revision.get("representation_id") == representation_id
                    and revision.get("ppt_manuscript_revision") == manuscript_revision
                ):
                    revision["ppt_manuscript_status"] = "confirmed"
                    revision["ppt_manuscript_confirmed_at"] = _now()
            saved = self._save(value)
            return deepcopy(next(
                item
                for item in saved["lessons"][lesson_unit_id]["ppt_assets"]
                if item["asset_id"] == asset["asset_id"]
            ))

    def save_imported_ppt_review(
        self,
        course_id: str,
        lesson_unit_id: str,
        *,
        package_id: str,
        source_asset_id: str,
        source_filename: str,
        slides: list[dict[str, Any]],
        report: dict[str, Any],
        source_outline_revision_id: str,
        source_lesson_plan_revision_id: str,
        source_script_revision_id: str,
        actor: str,
    ) -> dict[str, Any]:
        with self._lock:
            value = self.load(course_id)
            lesson = value.setdefault("lessons", {}).setdefault(
                lesson_unit_id, _empty_lesson_asset(lesson_unit_id)
            )
            for item in lesson.setdefault("imported_ppt_reviews", []):
                if isinstance(item, dict) and item.get("status") == "reviewing":
                    item["status"] = "superseded"
            review_id = f"tlpir-{uuid.uuid4().hex}"
            revision_id = f"tlpivr-{uuid.uuid4().hex}"
            review = {
                "review_id": review_id,
                "lesson_unit_id": lesson_unit_id,
                "package_id": package_id,
                "source_asset_id": source_asset_id,
                "source_filename": source_filename,
                "status": "reviewing",
                "source_state": "current",
                "revision_id": revision_id,
                "source_outline_revision_id": source_outline_revision_id,
                "source_lesson_plan_revision_id": source_lesson_plan_revision_id,
                "source_script_revision_id": source_script_revision_id,
                "slides": deepcopy(slides),
                "report": deepcopy(report),
                "ai_candidates": [],
                "revision_history": [{
                    "revision_id": revision_id,
                    "slides": deepcopy(slides),
                    "actor": actor,
                    "created_at": _now(),
                }],
                "created_at": _now(),
                "updated_at": _now(),
            }
            lesson["imported_ppt_reviews"].append(review)
            saved = self._save(value)
            return deepcopy(saved["lessons"][lesson_unit_id]["imported_ppt_reviews"][-1])

    def current_imported_ppt_review(
        self, course_id: str, lesson_unit_id: str
    ) -> dict[str, Any] | None:
        lesson = self.lesson(course_id, lesson_unit_id)
        reviews = [
            item for item in lesson.get("imported_ppt_reviews") or []
            if isinstance(item, dict) and item.get("status") in {"reviewing", "confirmed"}
        ]
        return deepcopy(reviews[-1]) if reviews else None

    def replace_imported_ppt_review(
        self,
        course_id: str,
        lesson_unit_id: str,
        *,
        review_id: str,
        base_revision_id: str,
        slides: list[dict[str, Any]],
        report: dict[str, Any],
        actor: str,
    ) -> dict[str, Any]:
        with self._lock:
            value = self.load(course_id)
            lesson = (value.get("lessons") or {}).get(lesson_unit_id)
            review = next(
                (item for item in (lesson or {}).get("imported_ppt_reviews") or []
                 if isinstance(item, dict) and item.get("review_id") == review_id),
                None,
            )
            if not isinstance(review, dict):
                raise TeacherLessonAuthoringError("uploaded_ppt_review_not_found", "PPT 审阅记录不存在。")
            if review.get("revision_id") != base_revision_id:
                raise TeacherLessonAuthoringError("uploaded_ppt_revision_conflict", "PPT 工作稿已更新，请刷新后再修改。")
            revision_id = f"tlpivr-{uuid.uuid4().hex}"
            review["slides"] = deepcopy(slides)
            review["report"] = deepcopy(report)
            review["revision_id"] = revision_id
            review["status"] = "reviewing"
            review["updated_at"] = _now()
            review.setdefault("revision_history", []).append({
                "revision_id": revision_id,
                "slides": deepcopy(slides),
                "actor": actor,
                "created_at": _now(),
            })
            saved = self._save(value)
            return deepcopy(next(
                item for item in saved["lessons"][lesson_unit_id]["imported_ppt_reviews"]
                if item.get("review_id") == review_id
            ))

    def save_imported_ppt_ai_candidate(
        self,
        course_id: str,
        lesson_unit_id: str,
        *,
        review_id: str,
        base_revision_id: str,
        slide_id: str,
        instruction: str,
        proposed_blocks: list[dict[str, Any]],
    ) -> dict[str, Any]:
        with self._lock:
            value = self.load(course_id)
            lesson = (value.get("lessons") or {}).get(lesson_unit_id)
            review = next((item for item in (lesson or {}).get("imported_ppt_reviews") or [] if isinstance(item, dict) and item.get("review_id") == review_id), None)
            if not isinstance(review, dict):
                raise TeacherLessonAuthoringError("uploaded_ppt_review_not_found", "PPT 审阅记录不存在。")
            if review.get("revision_id") != base_revision_id:
                raise TeacherLessonAuthoringError("uploaded_ppt_revision_conflict", "PPT 工作稿已更新，请重新生成 AI 候选。")
            for item in review.get("ai_candidates") or []:
                if isinstance(item, dict) and item.get("status") == "pending":
                    item["status"] = "superseded"
            candidate = {
                "candidate_id": f"tlpiac-{uuid.uuid4().hex}",
                "base_revision_id": base_revision_id,
                "slide_id": slide_id,
                "instruction": instruction,
                "proposed_blocks": deepcopy(proposed_blocks),
                "status": "pending",
                "created_at": _now(),
            }
            review.setdefault("ai_candidates", []).append(candidate)
            self._save(value)
            return deepcopy(candidate)

    def mark_imported_ppt_ai_candidate(
        self,
        course_id: str,
        lesson_unit_id: str,
        *,
        review_id: str,
        candidate_id: str,
        status: str,
    ) -> dict[str, Any]:
        if status not in {"accepted", "rejected", "superseded"}:
            raise ValueError("unsupported imported PPT candidate status")
        with self._lock:
            value = self.load(course_id)
            lesson = (value.get("lessons") or {}).get(lesson_unit_id)
            review = next((item for item in (lesson or {}).get("imported_ppt_reviews") or [] if isinstance(item, dict) and item.get("review_id") == review_id), None)
            candidate = next((item for item in (review or {}).get("ai_candidates") or [] if isinstance(item, dict) and item.get("candidate_id") == candidate_id), None)
            if not isinstance(candidate, dict):
                raise TeacherLessonAuthoringError("uploaded_ppt_candidate_not_found", "AI PPT 修改候选不存在。")
            if candidate.get("status") == "pending":
                candidate["status"] = status
                candidate["resolved_at"] = _now()
                self._save(value)
            return deepcopy(candidate)

    def confirm_imported_ppt_review(
        self,
        course_id: str,
        lesson_unit_id: str,
        *,
        review_id: str,
        revision_id: str,
    ) -> dict[str, Any]:
        with self._lock:
            value = self.load(course_id)
            lesson = (value.get("lessons") or {}).get(lesson_unit_id)
            review = next((item for item in (lesson or {}).get("imported_ppt_reviews") or [] if isinstance(item, dict) and item.get("review_id") == review_id), None)
            if not isinstance(review, dict):
                raise TeacherLessonAuthoringError("uploaded_ppt_review_not_found", "PPT 审阅记录不存在。")
            if review.get("revision_id") != revision_id:
                raise TeacherLessonAuthoringError("uploaded_ppt_revision_conflict", "PPT 工作稿已更新，请刷新后再确认。")
            if review.get("source_state") != "current":
                raise TeacherLessonAuthoringError("uploaded_ppt_source_stale", "上游教学内容已更新，请先重新审阅。")
            review["status"] = "confirmed"
            review["confirmed_revision_id"] = revision_id
            review["confirmed_at"] = _now()
            assets = lesson.setdefault("ppt_assets", [])
            for asset in assets:
                if isinstance(asset, dict) and asset.get("role") == "primary":
                    asset["role"] = "supplemental"
            assets.append({
                "asset_id": f"tlpa-{uuid.uuid4().hex}",
                "lesson_unit_id": lesson_unit_id,
                "role": "primary",
                "engine": "uploaded_pptx",
                "working_revision_id": revision_id,
                "source_lesson_plan_revision_id": str(review.get("source_lesson_plan_revision_id") or ""),
                "source_script_revision_id": str(review.get("source_script_revision_id") or ""),
                "source_state": "current",
                "package_id": str(review.get("package_id") or ""),
                "source_asset_id": str(review.get("source_asset_id") or ""),
                "review_id": review_id,
                "confirmed_at": review["confirmed_at"],
                "revisions": [],
                "ai_candidates": [],
            })
            saved = self._save(value)
            return deepcopy(next(item for item in saved["lessons"][lesson_unit_id]["imported_ppt_reviews"] if item.get("review_id") == review_id))

    def save_v6_ppt_ai_candidate(
        self,
        course_id: str,
        lesson_unit_id: str,
        *,
        representation_id: str,
        base_spec_id: str,
        base_spec_revision: str,
        page_id: str,
        instruction: str,
        candidate_page: dict[str, Any],
        changed_fields: list[str],
        page_changes: list[dict[str, Any]] | None = None,
    ) -> dict[str, Any]:
        with self._lock:
            value = self.load(course_id)
            lesson = (value.get("lessons") or {}).get(lesson_unit_id)
            asset = next(
                (
                    item for item in (lesson or {}).get("ppt_assets") or []
                    if isinstance(item, dict)
                    and item.get("role") == "primary"
                    and item.get("engine") == "slide_deck_v6"
                ),
                None,
            )
            if not isinstance(asset, dict):
                raise TeacherLessonAuthoringError(
                    "lesson_ppt_not_found", "本讲还没有可优化的 V6 PPT。"
                )
            if (
                asset.get("working_representation_id") != representation_id
                or not any(
                    isinstance(item, dict)
                    and item.get("spec_id") == base_spec_id
                    and item.get("representation_id") == representation_id
                    for item in asset.get("v6_revisions") or []
                )
            ):
                raise TeacherLessonAuthoringError(
                    "lesson_ppt_revision_conflict",
                    "PPT 工作稿已经变化，请重新生成 AI 候选。",
                )
            for item in asset.get("v6_ai_candidates") or []:
                if isinstance(item, dict) and item.get("status") == "pending":
                    item["status"] = "superseded"
                    item["resolved_at"] = _now()
            candidate = {
                "candidate_id": f"tlv6ac-{uuid.uuid4().hex}",
                "representation_id": representation_id,
                "base_spec_id": base_spec_id,
                "base_spec_revision": base_spec_revision,
                "page_id": page_id,
                "instruction": instruction,
                "candidate_page": deepcopy(candidate_page),
                "changed_fields": list(changed_fields),
                "page_changes": deepcopy(page_changes or []),
                "status": "pending",
                "created_at": _now(),
            }
            asset.setdefault("v6_ai_candidates", []).append(candidate)
            self._save(value)
            return deepcopy(candidate)

    def pending_v6_ppt_ai_candidate(
        self,
        course_id: str,
        lesson_unit_id: str,
        *,
        representation_id: str,
        spec_id: str,
        spec_revision: str,
    ) -> dict[str, Any] | None:
        lesson = self.lesson(course_id, lesson_unit_id)
        for asset in lesson.get("ppt_assets") or []:
            if not isinstance(asset, dict) or asset.get("role") != "primary":
                continue
            for candidate in reversed(asset.get("v6_ai_candidates") or []):
                if (
                    isinstance(candidate, dict)
                    and candidate.get("status") == "pending"
                    and candidate.get("representation_id") == representation_id
                    and candidate.get("base_spec_id") == spec_id
                    and candidate.get("base_spec_revision") == spec_revision
                ):
                    return deepcopy(candidate)
        return None

    def mark_v6_ppt_ai_candidate(
        self,
        course_id: str,
        lesson_unit_id: str,
        candidate_id: str,
        *,
        status: str,
        result_spec_id: str = "",
    ) -> dict[str, Any]:
        if status not in {"accepted", "rejected", "superseded"}:
            raise ValueError("unsupported V6 candidate status")
        with self._lock:
            value = self.load(course_id)
            lesson = (value.get("lessons") or {}).get(lesson_unit_id)
            candidate = next(
                (
                    candidate
                    for asset in (lesson or {}).get("ppt_assets") or []
                    if isinstance(asset, dict)
                    for candidate in asset.get("v6_ai_candidates") or []
                    if isinstance(candidate, dict)
                    and candidate.get("candidate_id") == candidate_id
                ),
                None,
            )
            if not isinstance(candidate, dict):
                raise TeacherLessonAuthoringError(
                    "lesson_ppt_candidate_not_found", "AI PPT 候选不存在。"
                )
            if candidate.get("status") == "pending":
                candidate["status"] = status
                candidate["resolved_at"] = _now()
                if result_spec_id:
                    candidate["result_spec_id"] = result_spec_id
                self._save(value)
            return deepcopy(candidate)

    def get_job(self, course_id: str, job_id: str) -> dict[str, Any]:
        with self._lock:
            live_job = (self._live_stream_jobs.get(course_id) or {}).get(job_id)
            if isinstance(live_job, dict):
                return deepcopy(live_job)
        value = self.load(course_id)
        job = (value.get("jobs") or {}).get(job_id)
        if not isinstance(job, dict):
            raise TeacherLessonAuthoringError("teacher_job_not_found", "教师讲次任务不存在。")
        return deepcopy(job)

    def expire_stale_job(
        self,
        course_id: str,
        job_id: str,
        *,
        stale_after_seconds: int = LESSON_JOB_STALE_SECONDS,
    ) -> dict[str, Any]:
        """Close an orphaned generation job left behind by a process reload."""
        with self._lock:
            live_job = (self._live_stream_jobs.get(course_id) or {}).get(job_id)
            touched_at = self._live_stream_touched_at.get((course_id, job_id))
            if isinstance(live_job, dict) and touched_at is not None:
                effective_stale_seconds = stale_after_seconds
                if (
                    stale_after_seconds == LESSON_JOB_STALE_SECONDS
                    and str(live_job.get("status") or "") == "pending"
                    and str(live_job.get("parent_job_id") or "")
                ):
                    effective_stale_seconds = LESSON_BATCH_QUEUED_STALE_SECONDS
                if time.monotonic() - touched_at < max(1, int(effective_stale_seconds)):
                    return deepcopy(live_job)
            value = self.load(course_id)
            job = (value.get("jobs") or {}).get(job_id)
            if not isinstance(job, dict):
                raise TeacherLessonAuthoringError("teacher_job_not_found", "教师讲次任务不存在。")
            if str(job.get("status") or "") not in {"pending", "running"}:
                return deepcopy(job)
            try:
                updated_at = datetime.fromisoformat(
                    str(job.get("updated_at") or "").replace("Z", "+00:00")
                )
                if updated_at.tzinfo is None:
                    updated_at = updated_at.replace(tzinfo=timezone.utc)
            except ValueError:
                updated_at = datetime.fromtimestamp(0, tz=timezone.utc)
            age_seconds = (datetime.now(timezone.utc) - updated_at).total_seconds()
            effective_stale_seconds = stale_after_seconds
            if (
                stale_after_seconds == LESSON_JOB_STALE_SECONDS
                and str(job.get("status") or "") == "pending"
                and str(job.get("parent_job_id") or "")
            ):
                effective_stale_seconds = LESSON_BATCH_QUEUED_STALE_SECONDS
            if age_seconds < max(1, int(effective_stale_seconds)):
                return deepcopy(job)
            script_job = str(job.get("type") or "") == "teacher_lesson_script_generation"
            job.update({
                "status": "failed",
                "phase": (
                    "lesson_script_interrupted"
                    if script_job
                    else "lesson_plan_interrupted"
                ),
                "message": (
                    "讲义生成进程已中断"
                    if script_job
                    else "教案生成进程已中断"
                ),
                "stream_sequence": int(job.get("stream_sequence") or 0) + 1,
                "stream_complete": True,
                "error": {
                    "code": (
                        "lesson_script_generation_interrupted"
                        if script_job
                        else "lesson_plan_generation_interrupted"
                    ),
                    "message": (
                        "生成进程已中断，已完成的讲义块仍然保留，可以继续生成。"
                        if script_job
                        else "生成进程已中断，请重新生成本讲教案。"
                    ),
                    "retryable": True,
                },
                "updated_at": _now(),
            })
            value["jobs"][job_id] = job
            saved = self._save(value)
            self._drop_live_stream_job_locked(course_id, job_id)
            return deepcopy(saved["jobs"][job_id])

    def expire_stale_jobs(
        self,
        course_id: str,
        *,
        stale_after_seconds: int = LESSON_JOB_STALE_SECONDS,
    ) -> dict[str, Any]:
        """Expire all orphaned jobs with one repository read and at most one write."""
        with self._lock:
            value = self.load(course_id)
            changed = False
            expired_job_ids: list[str] = []
            now = datetime.now(timezone.utc)
            for job_id, job in (value.get("jobs") or {}).items():
                if not isinstance(job, dict) or str(job.get("status") or "") not in {"pending", "running"}:
                    continue
                try:
                    updated_at = datetime.fromisoformat(
                        str(job.get("updated_at") or "").replace("Z", "+00:00")
                    )
                    if updated_at.tzinfo is None:
                        updated_at = updated_at.replace(tzinfo=timezone.utc)
                except ValueError:
                    updated_at = datetime.fromtimestamp(0, tz=timezone.utc)
                effective_stale_seconds = stale_after_seconds
                if (
                    stale_after_seconds == LESSON_JOB_STALE_SECONDS
                    and str(job.get("status") or "") == "pending"
                    and str(job.get("parent_job_id") or "")
                ):
                    effective_stale_seconds = LESSON_BATCH_QUEUED_STALE_SECONDS
                if (now - updated_at).total_seconds() < max(1, int(effective_stale_seconds)):
                    continue
                script_job = str(job.get("type") or "") == "teacher_lesson_script_generation"
                job.update({
                    "status": "failed",
                    "phase": "lesson_script_interrupted" if script_job else "lesson_plan_interrupted",
                    "message": "讲义生成进程已中断" if script_job else "教案生成进程已中断",
                    "stream_sequence": int(job.get("stream_sequence") or 0) + 1,
                    "stream_complete": True,
                    "error": {
                        "code": "lesson_script_generation_interrupted" if script_job else "lesson_plan_generation_interrupted",
                        "message": (
                            "生成进程已中断，已完成的讲义块仍然保留，可以继续生成。"
                            if script_job
                            else "生成进程已中断，请重新生成本讲教案。"
                        ),
                        "retryable": True,
                    },
                    "updated_at": _now(),
                })
                changed = True
                expired_job_ids.append(str(job_id))
            if not changed:
                return deepcopy(value)
            saved = self._save(value)
            for job_id in expired_job_ids:
                self._drop_live_stream_job_locked(course_id, job_id)
            return saved

    def lesson(self, course_id: str, lesson_unit_id: str) -> dict[str, Any]:
        value = self.load(course_id)
        lesson = (value.get("lessons") or {}).get(lesson_unit_id)
        if not isinstance(lesson, dict):
            return _empty_lesson_asset(lesson_unit_id)
        return deepcopy(lesson)

    def confirm_plan_revision(
        self,
        course_id: str,
        lesson_unit_id: str,
        revision_id: str,
        *,
        quality_report: dict[str, Any] | None = None,
    ) -> dict[str, Any]:
        with self._lock:
            value = self.load(course_id)
            lesson = (value.get("lessons") or {}).get(lesson_unit_id)
            if not isinstance(lesson, dict):
                raise TeacherLessonAuthoringError("lesson_plan_not_found", "本讲还没有可确认的教案。")
            revision = next(
                (
                    item for item in lesson.get("revisions") or []
                    if isinstance(item, dict) and item.get("revision_id") == revision_id
                ),
                None,
            )
            if revision is None:
                raise TeacherLessonAuthoringError("lesson_plan_revision_not_found", "教案修订不存在。")
            if lesson.get("working_revision_id") != revision_id:
                raise TeacherLessonAuthoringError(
                    "lesson_plan_revision_conflict",
                    "只能确认当前教案工作稿。",
                )
            effective_quality = deepcopy(
                quality_report
                or revision.get("quality_report")
                or validate_teacher_lesson_plan(revision.get("plan") or {})
            )
            revision["pipeline_version"] = LESSON_PLAN_PIPELINE_VERSION
            revision["quality_report"] = effective_quality
            if not effective_quality.get("passed"):
                raise TeacherLessonAuthoringError(
                    "lesson_plan_quality_blocked",
                    "教案尚未通过专业性与完整性检查。",
                    details={
                        "quality_report": deepcopy(effective_quality),
                        "blocking_issues": deepcopy(
                            effective_quality.get("blocking_issues") or []
                        ),
                    },
                )
            if lesson.get("source_state") != "current":
                arrangement_changed = lesson.get("source_state_reason") in {
                    "arrangement_changed", "arrangement_stale",
                }
                raise TeacherLessonAuthoringError(
                    (
                        "lesson_plan_arrangement_conflict"
                        if arrangement_changed
                        else "lesson_plan_outline_conflict"
                    ),
                    (
                        "本讲教学结构已经变化，请先更新教案。"
                        if arrangement_changed
                        else "教案对应的大纲已经变化，请先更新教案。"
                    ),
                )
            lesson["confirmed_revision_id"] = revision_id
            revision["status"] = "confirmed"
            revision["confirmed_at"] = _now()
            script_confirmation = lesson.get("script_confirmation")
            if isinstance(script_confirmation, dict) and script_confirmation.get("confirmed_revision_id"):
                if script_confirmation.get("source_lesson_plan_revision_id") != revision_id:
                    script_confirmation["source_state"] = "stale"
            saved = self._save(value)
            return deepcopy(saved["lessons"][lesson_unit_id])

    def save_script_revision(
        self,
        course_id: str,
        lesson_unit_id: str,
        sections: list[dict[str, Any]],
        *,
        source_lesson_plan_revision_id: str,
        generation_source: str = "model",
        requirements: str = "",
        material_asset_ids: list[str] | None = None,
        actor: str = "teacher",
        expected_working_revision_id: str | None = None,
        revision_id_override: str = "",
        restored_from_revision_id: str = "",
        active_job_id: str = "",
    ) -> dict[str, Any]:
        normalized_sections = []
        for item in sections:
            if not isinstance(item, dict):
                continue
            normalized = normalize_teacher_script_section(item)
            quality_report = deepcopy(item.get("quality_report") or {})
            if not quality_report:
                compatibility_modules = [
                    {
                        **deepcopy(block),
                        "artifact_contract": {},
                        "target_characters": 0,
                        "max_characters": 0,
                    }
                    for block in normalized.get("blocks") or []
                    if isinstance(block, dict)
                ]
                quality_report = validate_teacher_script_section(
                    normalized,
                    {
                        "section_node_id": normalized.get("section_node_id"),
                        "title": normalized.get("title"),
                        "modules": compatibility_modules,
                    },
                )
            normalized["quality_report"] = quality_report
            normalized["pipeline_version"] = str(
                item.get("pipeline_version")
                or quality_report.get("pipeline_version")
                or SCRIPT_PIPELINE_VERSION
            )
            normalized_sections.append(normalized)
        if not normalized_sections or any(
            not item["section_node_id"]
            or not item["content"]
            or not item.get("blocks")
            for item in normalized_sections
        ):
            raise TeacherLessonAuthoringError(
                "lesson_script_incomplete",
                "本讲仍有小节没有讲义内容，暂时不能保存。",
            )
        revision_quality = validate_teacher_script_revision(
            normalized_sections,
            generation_source=generation_source,
        )
        publication_eligible = bool(revision_quality.get("publication_eligible"))
        revision_id = revision_id_override or teacher_lesson_script_sections_revision(normalized_sections)
        with self._lock:
            value = self.load(course_id)
            if active_job_id:
                active_job = (value.get("jobs") or {}).get(active_job_id)
                if (
                    not isinstance(active_job, dict)
                    or str(active_job.get("status") or "") not in TEACHER_JOB_ACTIVE_STATUSES
                ):
                    raise TeacherLessonAuthoringError(
                        "teacher_job_not_active",
                        "任务已停止，迟到的讲义结果未保存。",
                    )
            lesson = (value.get("lessons") or {}).get(lesson_unit_id)
            if not isinstance(lesson, dict):
                raise TeacherLessonAuthoringError(
                    "lesson_plan_not_found",
                    "请先生成本讲教案。",
                )
            if (
                expected_working_revision_id is not None
                and lesson.get("working_script_revision_id")
                != expected_working_revision_id
            ):
                raise TeacherLessonAuthoringError(
                    "lesson_script_revision_conflict",
                    "讲义工作稿已经变化，请基于当前版本重新修改。",
                )
            if (
                lesson.get("working_revision_id") != source_lesson_plan_revision_id
                or lesson.get("source_state", "current") != "current"
            ):
                raise TeacherLessonAuthoringError(
                    "lesson_plan_revision_conflict",
                    "当前教案已经变化，请基于最新教案生成讲义。",
                )
            revisions = lesson.setdefault("script_revisions", [])
            existing = next(
                (
                    item for item in revisions
                    if isinstance(item, dict) and item.get("revision_id") == revision_id
                ),
                None,
            )
            if existing is None:
                revision = {
                    "revision_id": revision_id,
                    "lesson_unit_id": lesson_unit_id,
                    "source_lesson_plan_revision_id": source_lesson_plan_revision_id,
                    "generation_source": generation_source,
                    "requirements": requirements,
                    "material_asset_ids": list(dict.fromkeys(
                        str(value or "").strip()
                        for value in material_asset_ids or []
                        if str(value or "").strip()
                    )),
                    "sections": normalized_sections,
                    "quality_report": revision_quality,
                    "publication_eligible": publication_eligible,
                    "quality_contract_version": SCRIPT_QUALITY_VERSION,
                    "pipeline_version": SCRIPT_PIPELINE_VERSION,
                    "actor": actor,
                    "created_at": _now(),
                }
                if restored_from_revision_id:
                    revision["restored_from_revision_id"] = restored_from_revision_id
                revisions.append(revision)
            else:
                # The content digest may stay unchanged while the quality
                # contract is tightened. Refresh the same revision in place so
                # stale v5 reports can never remain publishable by accident.
                existing.update({
                    "generation_source": generation_source,
                    "requirements": requirements,
                    "material_asset_ids": list(dict.fromkeys(
                        str(item or "").strip()
                        for item in material_asset_ids or []
                        if str(item or "").strip()
                    )),
                    "sections": normalized_sections,
                    "quality_report": revision_quality,
                    "publication_eligible": publication_eligible,
                    "quality_contract_version": SCRIPT_QUALITY_VERSION,
                    "pipeline_version": SCRIPT_PIPELINE_VERSION,
                    "actor": actor,
                    "updated_at": _now(),
                })
            lesson["working_script_revision_id"] = revision_id
            confirmation = lesson.get("script_confirmation")
            if (
                isinstance(confirmation, dict)
                and confirmation.get("confirmed_revision_id")
                and confirmation.get("confirmed_revision_id") != revision_id
            ):
                confirmation["source_state"] = "stale"
            for asset in lesson.get("ppt_assets") or []:
                if not isinstance(asset, dict) or asset.get("engine") != "slide_deck_v6":
                    continue
                if asset.get("source_script_revision_id") != revision_id:
                    asset["source_state"] = "stale"
            manuscript = lesson.get("ppt_manuscript")
            if (
                isinstance(manuscript, dict)
                and manuscript.get("source_script_revision_id") != revision_id
            ):
                manuscript["source_state"] = "stale"
            for review in lesson.get("imported_ppt_reviews") or []:
                if not isinstance(review, dict):
                    continue
                source_revision = str(review.get("source_script_revision_id") or "")
                if source_revision and source_revision != revision_id:
                    review["source_state"] = "stale"
            saved = self._save(value)
            return deepcopy(saved["lessons"][lesson_unit_id])

    def restore_script_revision(
        self,
        course_id: str,
        lesson_unit_id: str,
        revision_id: str,
        *,
        expected_working_revision_id: str,
        actor: str = "teacher",
    ) -> dict[str, Any]:
        with self._lock:
            lesson = self.lesson(course_id, lesson_unit_id)
            current_revision_id = str(lesson.get("working_script_revision_id") or "")
            if current_revision_id != expected_working_revision_id:
                raise TeacherLessonAuthoringError(
                    "lesson_script_revision_conflict",
                    "讲义已在其他页面修改，请重新载入后再恢复。",
                )
            source = next(
                (
                    item for item in lesson.get("script_revisions") or []
                    if isinstance(item, dict) and item.get("revision_id") == revision_id
                ),
                None,
            )
            if not isinstance(source, dict):
                raise TeacherLessonAuthoringError(
                    "lesson_script_revision_not_found",
                    "讲义历史版本不存在。",
                )
            return self.save_script_revision(
                course_id,
                lesson_unit_id,
                deepcopy(source.get("sections") or []),
                source_lesson_plan_revision_id=str(
                    source.get("source_lesson_plan_revision_id") or ""
                ),
                generation_source="history_restore",
                requirements=str(source.get("requirements") or ""),
                material_asset_ids=list(source.get("material_asset_ids") or []),
                actor=actor,
                expected_working_revision_id=expected_working_revision_id,
                revision_id_override=f"tlsr-restore-{uuid.uuid4().hex}",
                restored_from_revision_id=revision_id,
            )

    def save_script_ai_candidate(
        self,
        course_id: str,
        lesson_unit_id: str,
        *,
        base_revision_id: str,
        section_node_id: str,
        instruction: str,
        replacement_text: str,
        source_lesson_plan_revision_id: str,
        material_asset_ids: list[str] | None = None,
        section_replacements: dict[str, str] | None = None,
    ) -> dict[str, Any]:
        with self._lock:
            value = self.load(course_id)
            lesson = (value.get("lessons") or {}).get(lesson_unit_id)
            if not isinstance(lesson, dict):
                raise TeacherLessonAuthoringError(
                    "lesson_plan_not_found",
                    "请先生成并确认本讲教案。",
                )
            if lesson.get("working_script_revision_id") != base_revision_id:
                raise TeacherLessonAuthoringError(
                    "lesson_script_revision_conflict",
                    "讲义工作稿已经变化，请重新生成 AI 候选。",
                )
            for item in lesson.get("script_ai_candidates") or []:
                if isinstance(item, dict) and item.get("status") == "pending":
                    item["status"] = "superseded"
                    item["resolved_at"] = _now()
            candidate = {
                "candidate_id": f"tlsac-{uuid.uuid4().hex}",
                "lesson_unit_id": lesson_unit_id,
                "base_revision_id": base_revision_id,
                "source_lesson_plan_revision_id": source_lesson_plan_revision_id,
                "section_node_id": section_node_id,
                "instruction": instruction,
                "replacement_text": replacement_text,
                "section_replacements": {
                    str(key): str(value).strip()
                    for key, value in (section_replacements or {}).items()
                    if str(key).strip() and str(value).strip()
                },
                "material_asset_ids": list(dict.fromkeys(
                    str(item).strip()
                    for item in material_asset_ids or []
                    if str(item).strip()
                )),
                "status": "pending",
                "created_at": _now(),
            }
            lesson.setdefault("script_ai_candidates", []).append(candidate)
            self._save(value)
            return deepcopy(candidate)

    def script_ai_candidate(
        self,
        course_id: str,
        lesson_unit_id: str,
        candidate_id: str,
    ) -> dict[str, Any]:
        lesson = self.lesson(course_id, lesson_unit_id)
        candidate = next(
            (
                item for item in lesson.get("script_ai_candidates") or []
                if isinstance(item, dict) and item.get("candidate_id") == candidate_id
            ),
            None,
        )
        if not isinstance(candidate, dict):
            raise TeacherLessonAuthoringError(
                "lesson_script_candidate_not_found",
                "AI 讲义候选不存在。",
            )
        return deepcopy(candidate)

    def mark_script_ai_candidate(
        self,
        course_id: str,
        lesson_unit_id: str,
        candidate_id: str,
        *,
        status: str,
        result_revision_id: str = "",
    ) -> dict[str, Any]:
        if status not in {"accepted", "rejected", "superseded"}:
            raise ValueError("unsupported script candidate status")
        with self._lock:
            value = self.load(course_id)
            lesson = (value.get("lessons") or {}).get(lesson_unit_id)
            candidate = next(
                (
                    item for item in (lesson or {}).get("script_ai_candidates") or []
                    if isinstance(item, dict) and item.get("candidate_id") == candidate_id
                ),
                None,
            )
            if not isinstance(candidate, dict):
                raise TeacherLessonAuthoringError(
                    "lesson_script_candidate_not_found",
                    "AI 讲义候选不存在。",
                )
            if candidate.get("status") == "pending":
                candidate["status"] = status
                candidate["resolved_at"] = _now()
                if result_revision_id:
                    candidate["result_revision_id"] = result_revision_id
                self._save(value)
            return deepcopy(candidate)

    def confirm_script_revision(
        self,
        course_id: str,
        lesson_unit_id: str,
        revision_id: str,
    ) -> dict[str, Any]:
        with self._lock:
            value = self.load(course_id)
            lesson = (value.get("lessons") or {}).get(lesson_unit_id)
            if not isinstance(lesson, dict):
                raise TeacherLessonAuthoringError(
                    "lesson_plan_not_found",
                    "请先生成并确认本讲教案。",
                )
            source_plan_revision = str(lesson.get("confirmed_revision_id") or "")
            if not source_plan_revision:
                raise TeacherLessonAuthoringError(
                    "lesson_plan_not_confirmed",
                    "请先确认本讲教案，再确认讲义。",
                )
            if lesson.get("working_script_revision_id") != revision_id:
                raise TeacherLessonAuthoringError(
                    "lesson_script_revision_conflict",
                    "只能确认当前讲义工作稿。",
                )
            revision = next(
                (
                    item for item in lesson.get("script_revisions") or []
                    if isinstance(item, dict) and item.get("revision_id") == revision_id
                ),
                None,
            )
            if not isinstance(revision, dict):
                raise TeacherLessonAuthoringError(
                    "lesson_script_revision_not_found",
                    "讲义修订不存在。",
                )
            if revision.get("source_lesson_plan_revision_id") != source_plan_revision:
                raise TeacherLessonAuthoringError(
                    "lesson_plan_revision_conflict",
                    "讲义对应的教案已经变化，请重新生成讲义。",
                )
            quality_report = revision.get("quality_report") or {}
            if not teacher_script_revision_is_publishable(revision):
                raise TeacherLessonAuthoringError(
                    "lesson_script_quality_blocked",
                    "讲义尚未通过当前教学质量与来源检查，请修正或重新生成后再确认。",
                    details={
                        "quality_report": deepcopy(quality_report),
                    },
                )
            lesson["script_confirmation"] = {
                "confirmed_revision_id": revision_id,
                "source_lesson_plan_revision_id": source_plan_revision,
                "source_state": "current",
                "confirmed_at": _now(),
            }
            for asset in lesson.get("ppt_assets") or []:
                if not isinstance(asset, dict):
                    continue
                source_revision = str(asset.get("source_script_revision_id") or "")
                if asset.get("engine") == "slide_deck_v6" and source_revision != revision_id:
                    asset["source_state"] = "stale"
            manuscript = lesson.get("ppt_manuscript")
            if (
                isinstance(manuscript, dict)
                and manuscript.get("source_script_revision_id") != revision_id
            ):
                manuscript["source_state"] = "stale"
            saved = self._save(value)
            return deepcopy(saved["lessons"][lesson_unit_id])

    def save_ai_candidate(
        self,
        course_id: str,
        lesson_unit_id: str,
        *,
        base_revision_id: str,
        instruction: str,
        plan: dict[str, Any],
        section_node_id: str = "",
        material_asset_ids: list[str] | None = None,
    ) -> dict[str, Any]:
        with self._lock:
            value = self.load(course_id)
            lesson = (value.get("lessons") or {}).get(lesson_unit_id)
            if not isinstance(lesson, dict):
                raise TeacherLessonAuthoringError("lesson_plan_not_found", "本讲还没有可优化的教案。")
            if lesson.get("working_revision_id") != base_revision_id:
                raise TeacherLessonAuthoringError("lesson_plan_revision_conflict", "教案草稿已经变化，请重新生成 AI 候选。")
            candidate = {
                "candidate_id": f"tlpc-{uuid.uuid4().hex}",
                "lesson_unit_id": lesson_unit_id,
                "base_revision_id": base_revision_id,
                "instruction": instruction,
                "section_node_id": section_node_id,
                "material_asset_ids": sorted({
                    str(value).strip()
                    for value in (material_asset_ids or [])
                    if str(value).strip()
                }),
                "plan": deepcopy(plan),
                "status": "pending",
                "created_at": _now(),
            }
            lesson.setdefault("ai_candidates", []).append(candidate)
            self._save(value)
            return deepcopy(candidate)

    def resolve_ai_candidate(
        self,
        course_id: str,
        lesson_unit_id: str,
        candidate_id: str,
        *,
        accept: bool,
        quality_report: dict[str, Any] | None = None,
        accepted_plan: dict[str, Any] | None = None,
        actor: str = "teacher",
    ) -> dict[str, Any]:
        with self._lock:
            value = self.load(course_id)
            lesson = (value.get("lessons") or {}).get(lesson_unit_id)
            if not isinstance(lesson, dict):
                raise TeacherLessonAuthoringError("lesson_plan_not_found", "本讲还没有可优化的教案。")
            candidate = next(
                (
                    item for item in lesson.get("ai_candidates") or []
                    if isinstance(item, dict) and item.get("candidate_id") == candidate_id
                ),
                None,
            )
            if candidate is None:
                raise TeacherLessonAuthoringError("lesson_plan_candidate_not_found", "AI 教案候选不存在。")
            if candidate.get("status") != "pending":
                return deepcopy(lesson)
            if lesson.get("working_revision_id") != candidate.get("base_revision_id"):
                raise TeacherLessonAuthoringError("lesson_plan_revision_conflict", "教案草稿已经变化，不能覆盖新修改。")
            candidate["status"] = "accepted" if accept else "rejected"
            candidate["resolved_at"] = _now()
            if not accept:
                saved = self._save(value)
                return deepcopy(saved["lessons"][lesson_unit_id])
            source_outline_revision_id = str(value.get("outline_revision_id") or "")
            plan = deepcopy(accepted_plan or candidate.get("plan") or {})
            candidate["plan"] = deepcopy(plan)
            self._save(value)
        return self.save_plan_revision(
            course_id,
            lesson_unit_id,
            plan,
            source_outline_revision_id=source_outline_revision_id,
            generation_source="ai_optimization",
            quality_report=quality_report,
            actor=actor,
        )

    def apply_material_absorption(
        self,
        course_id: str,
        bundle: dict[str, Any],
    ) -> dict[str, Any]:
        """Create linked structured working drafts without confirming them.

        The whole bundle is written in one authoring-state replacement and is
        idempotent by bundle_id.  Existing confirmed revisions are never
        changed or hidden.
        """
        bundle_id = str(bundle.get("bundle_id") or "")
        if not bundle_id or str(bundle.get("course_id") or "") != course_id:
            raise TeacherLessonAuthoringError(
                "material_absorption_bundle_invalid",
                "材料吸收执行包不属于当前课程。",
            )
        with self._lock:
            value = self.load(course_id)
            existing = next(
                (
                    item for item in value.get("material_absorptions") or []
                    if isinstance(item, dict) and item.get("bundle_id") == bundle_id
                ),
                None,
            )
            if isinstance(existing, dict):
                return deepcopy(existing)

            created: list[dict[str, Any]] = []
            for target in bundle.get("targets") or []:
                if not isinstance(target, dict):
                    continue
                target_id = str(target.get("target_id") or "")
                target_type = str(target.get("target_type") or "")
                scope_id = str(target.get("target_scope_id") or "")
                structured = deepcopy(target.get("structured_document") or {})
                if target_type not in {"outline", "lesson_plan", "script", "ppt"} or not target_id or not structured:
                    raise TeacherLessonAuthoringError(
                        "material_absorption_target_invalid",
                        "材料吸收包中存在不完整的正式文件工作稿。",
                    )
                revision_id = "tmad-" + hashlib.sha256(
                    f"{bundle_id}:{target_id}:{structured.get('content_hash', '')}".encode("utf-8")
                ).hexdigest()[:24]
                draft = {
                    "revision_id": revision_id,
                    "bundle_id": bundle_id,
                    "plan_id": str(bundle.get("plan_id") or ""),
                    "package_id": str(bundle.get("package_id") or ""),
                    "target_id": target_id,
                    "target_type": target_type,
                    "target_scope_id": scope_id,
                    "target_scope_label": str(target.get("target_scope_label") or ""),
                    "title": str(target.get("title") or target_id),
                    "status": "working_draft",
                    "source_state": "current",
                    "confirmation_required": True,
                    "structured_document": structured,
                    "source_refs": deepcopy(target.get("sources") or []),
                    "created_at": _now(),
                }
                if target_type == "outline":
                    for prior in value.setdefault("outline_material_drafts", []):
                        if isinstance(prior, dict) and prior.get("status") == "working_draft":
                            prior["status"] = "superseded"
                            prior["superseded_at"] = _now()
                    value["outline_material_drafts"].append(draft)
                    value["current_outline_material_draft_id"] = revision_id
                else:
                    lesson = value.setdefault("lessons", {}).setdefault(
                        scope_id,
                        _empty_lesson_asset(scope_id),
                    )
                    drafts = lesson.setdefault("material_drafts", {}).setdefault(target_type, [])
                    for prior in drafts:
                        if isinstance(prior, dict) and prior.get("status") == "working_draft":
                            prior["status"] = "superseded"
                            prior["superseded_at"] = _now()
                    drafts.append(draft)
                    lesson.setdefault("current_material_draft_ids", {})[target_type] = revision_id
                created.append({
                    "revision_id": revision_id,
                    "target_id": target_id,
                    "target_type": target_type,
                    "target_scope_id": scope_id,
                    "status": "working_draft",
                })

            receipt = {
                "schema_version": "teacher_material_absorption_receipt_v1",
                "bundle_id": bundle_id,
                "plan_id": str(bundle.get("plan_id") or ""),
                "status": "working_drafts_created",
                "created_at": _now(),
                "drafts": created,
            }
            value.setdefault("material_absorptions", []).append(receipt)
            self._save(value)
            return deepcopy(receipt)

    def current_material_drafts(self, course_id: str) -> dict[str, Any]:
        value = self.load(course_id)
        outline_id = str(value.get("current_outline_material_draft_id") or "")
        outline = next(
            (
                deepcopy(item) for item in reversed(value.get("outline_material_drafts") or [])
                if isinstance(item, dict) and item.get("revision_id") == outline_id
            ),
            None,
        )
        lessons: dict[str, dict[str, Any]] = {}
        for lesson_id, lesson in (value.get("lessons") or {}).items():
            if not isinstance(lesson, dict):
                continue
            selected: dict[str, Any] = {}
            for target_type, revision_id in (lesson.get("current_material_draft_ids") or {}).items():
                draft = next(
                    (
                        deepcopy(item) for item in reversed((lesson.get("material_drafts") or {}).get(target_type) or [])
                        if isinstance(item, dict) and item.get("revision_id") == revision_id
                    ),
                    None,
                )
                if draft:
                    selected[str(target_type)] = draft
            if selected:
                lessons[str(lesson_id)] = selected
        return {"outline": outline, "lessons": lessons}

    def view(self, course_id: str) -> dict[str, Any]:
        return self.load(course_id)


Planner = Callable[[dict[str, Any], str, Callable[..., Awaitable[None]]], Awaitable[dict[str, Any]]]


class TeacherLessonAuthoringService:
    def __init__(self, repository: TeacherLessonAuthoringRepository):
        self.repository = repository

    def _quality_report(
        self,
        course_data: dict[str, Any],
        lesson_unit_id: str,
        plan: dict[str, Any],
        *,
        expected_outline_revision_id: str,
        source_outline_revision_id: str,
    ) -> dict[str, Any]:
        scope = lesson_scope(course_data, lesson_unit_id)
        arrangement = self.repository.current_arrangement(
            str(course_data.get("course_id") or ""),
            lesson_unit_id,
        )
        expected_total_minutes = (
            sum(
                max(1, int(item.get("planned_minutes") or 1))
                for item in arrangement.get("blocks") or []
                if isinstance(item, dict)
            )
            if arrangement
            else None
        )
        return validate_teacher_lesson_plan(
            plan,
            expected_section_ids=[
                str(item.get("node_id") or "")
                for item in scope["sections"]
            ],
            expected_outline_revision_id=expected_outline_revision_id,
            source_outline_revision_id=source_outline_revision_id,
            expected_total_minutes=expected_total_minutes,
        )

    def save_plan_draft(
        self,
        *,
        course_id: str,
        lesson_unit_id: str,
        course_data: dict[str, Any],
        plan: dict[str, Any],
        source_outline_revision_id: str,
        actor: str,
    ) -> dict[str, Any]:
        canonical_outline_revision = str(
            self.repository.view(course_id).get("outline_revision_id") or ""
        )
        effective_source_revision = (
            source_outline_revision_id or canonical_outline_revision
        )
        plan = align_teacher_lesson_plan_to_arrangement(
            plan,
            self.repository.current_arrangement(course_id, lesson_unit_id),
        )
        quality_report = self._quality_report(
            course_data,
            lesson_unit_id,
            plan,
            expected_outline_revision_id=canonical_outline_revision,
            source_outline_revision_id=effective_source_revision,
        )
        return self.repository.save_plan_revision(
            course_id,
            lesson_unit_id,
            plan,
            source_outline_revision_id=effective_source_revision,
            generation_source="teacher_edit",
            quality_report=quality_report,
            actor=actor,
        )

    def confirm_plan(
        self,
        *,
        course_id: str,
        lesson_unit_id: str,
        course_data: dict[str, Any],
        revision_id: str,
    ) -> dict[str, Any]:
        lesson = self.repository.lesson(course_id, lesson_unit_id)
        revision = next(
            (
                item for item in lesson.get("revisions") or []
                if isinstance(item, dict) and item.get("revision_id") == revision_id
            ),
            None,
        )
        if not isinstance(revision, dict):
            raise TeacherLessonAuthoringError(
                "lesson_plan_revision_not_found",
                "教案修订不存在。",
            )
        canonical_outline_revision = str(
            self.repository.view(course_id).get("outline_revision_id") or ""
        )
        quality_report = self._quality_report(
            course_data,
            lesson_unit_id,
            revision.get("plan") or {},
            expected_outline_revision_id=canonical_outline_revision,
            source_outline_revision_id=str(
                revision.get("source_outline_revision_id") or ""
            ),
        )
        return self.repository.confirm_plan_revision(
            course_id,
            lesson_unit_id,
            revision_id,
            quality_report=quality_report,
        )

    def resolve_ai_candidate(
        self,
        *,
        course_id: str,
        lesson_unit_id: str,
        course_data: dict[str, Any],
        candidate_id: str,
        accept: bool,
        actor: str,
    ) -> dict[str, Any]:
        lesson = self.repository.lesson(course_id, lesson_unit_id)
        candidate = next(
            (
                item for item in lesson.get("ai_candidates") or []
                if isinstance(item, dict)
                and item.get("candidate_id") == candidate_id
            ),
            None,
        )
        if not isinstance(candidate, dict) or not accept:
            return self.repository.resolve_ai_candidate(
                course_id,
                lesson_unit_id,
                candidate_id,
                accept=accept,
                actor=actor,
            )
        canonical_outline_revision = str(
            self.repository.view(course_id).get("outline_revision_id") or ""
        )
        aligned_plan = align_teacher_lesson_plan_to_arrangement(
            candidate.get("plan") or {},
            self.repository.current_arrangement(course_id, lesson_unit_id),
        )
        quality_report = self._quality_report(
            course_data,
            lesson_unit_id,
            aligned_plan,
            expected_outline_revision_id=canonical_outline_revision,
            source_outline_revision_id=canonical_outline_revision,
        )
        return self.repository.resolve_ai_candidate(
            course_id,
            lesson_unit_id,
            candidate_id,
            accept=True,
            quality_report=quality_report,
            accepted_plan=aligned_plan,
            actor=actor,
        )

    async def run_plan_job(
        self,
        *,
        course_id: str,
        lesson_unit_id: str,
        job_id: str,
        course_data: dict[str, Any],
        planner: Planner,
    ) -> dict[str, Any]:
        started_job = await asyncio.to_thread(
            self.repository.update_job,
            course_id,
            job_id,
            status="running",
            phase="lesson_plan_generation",
            progress=5,
            message="正在生成本讲全部小节教案",
        )
        if str(started_job.get("status") or "") not in TEACHER_JOB_ACTIVE_STATUSES:
            return started_job

        async def on_progress(
            phase: str,
            progress: int,
            message: str,
            _phase_progress: int = 0,
            phase_detail: dict[str, Any] | None = None,
        ) -> None:
            current_job = await asyncio.to_thread(
                self.repository.get_job,
                course_id,
                job_id,
            )
            if current_job.get("cancel_requested"):
                raise asyncio.CancelledError
            current_progress = int(current_job.get("progress") or 0)
            next_progress = max(
                current_progress,
                max(5, min(95, int(progress))),
            )
            changes: dict[str, Any] = {
                "phase": phase,
                "progress": next_progress,
                "message": message,
            }
            detail = phase_detail or {}
            stream_event = str(detail.get("stream_event") or "")
            batch_id = str(detail.get("stream_batch_id") or "")
            if stream_event in {"reset", "delta"} and batch_id:
                await asyncio.to_thread(
                    self.repository.update_job_stream,
                    course_id,
                    job_id,
                    phase=phase,
                    progress=changes["progress"],
                    message=message,
                    batch_id=batch_id,
                    event=stream_event,
                    delta=str(detail.get("stream_delta") or ""),
                )
                return
            await asyncio.to_thread(
                self.repository.update_job_live,
                course_id,
                job_id,
                **changes,
            )

        try:
            result = await planner(course_data, lesson_unit_id, on_progress)
            current_job = await asyncio.to_thread(
                self.repository.get_job,
                course_id,
                job_id,
            )
            if current_job.get("cancel_requested"):
                raise asyncio.CancelledError
            plan = result.get("plan") if isinstance(result, dict) else None
            if not isinstance(plan, dict) or not plan.get("sections"):
                raise TeacherLessonAuthoringError(
                    "lesson_plan_empty",
                    "本讲教案生成结果为空。",
                )
            arrangement = await asyncio.to_thread(
                self.repository.current_arrangement,
                course_id,
                lesson_unit_id,
            )
            plan = align_teacher_lesson_plan_to_arrangement(plan, arrangement)
            plan = complete_teacher_lesson_plan_fields(
                course_data,
                lesson_unit_id,
                plan,
            )
            generation_warnings = list(result.get("warnings") or [])
            warnings = list(generation_warnings)
            generation_source = str(
                result.get("generation_source")
                or (
                    "deterministic_local_fallback"
                    if generation_warnings
                    else "model"
                )
            )
            if "fallback" in generation_source:
                raise TeacherLessonAuthoringError(
                    "lesson_plan_generation_incomplete",
                    "本讲教案的模型生成未完整完成，请从任务检查点重试。",
                )
            source_refs = [
                deepcopy(item)
                for item in result.get("source_refs") or []
                if isinstance(item, dict)
            ]
            current_job = await asyncio.to_thread(
                self.repository.get_job,
                course_id,
                job_id,
            )
            job_source_revision = str(current_job.get("source_outline_revision_id") or "")
            knowledge_scope_revision = str(
                result.get("source_outline_revision_id") or ""
            )
            outline_revision = job_source_revision or knowledge_scope_revision
            course_view = await asyncio.to_thread(self.repository.view, course_id)
            quality_report = self._quality_report(
                course_data,
                lesson_unit_id,
                plan,
                expected_outline_revision_id=str(
                    course_view.get("outline_revision_id")
                    or outline_revision
                ),
                source_outline_revision_id=outline_revision,
            )
            if not quality_report.get("passed"):
                messages = "；".join(
                    str(item.get("message") or "未知教案错误")
                    for item in quality_report.get("blocking_issues") or []
                    if isinstance(item, dict)
                )
                raise TeacherLessonAuthoringError(
                    "lesson_plan_quality_failed",
                    f"本讲教案未通过硬校验：{messages or '请重试'}",
                    details={"quality_report": deepcopy(quality_report)},
                )
            lesson = await asyncio.to_thread(
                self.repository.save_plan_revision,
                course_id,
                lesson_unit_id,
                plan,
                source_outline_revision_id=outline_revision,
                source_knowledge_scope_revision_id=knowledge_scope_revision,
                generation_source=generation_source,
                warnings=warnings,
                source_refs=source_refs,
                quality_report=quality_report,
                active_job_id=job_id,
            )
            current_job = await asyncio.to_thread(
                self.repository.get_job,
                course_id,
                job_id,
            )
            return await asyncio.to_thread(
                self.repository.update_job,
                course_id,
                job_id,
                status="completed",
                phase="lesson_plan_ready",
                progress=100,
                message="本讲教案已生成",
                warnings=warnings,
                result_revision_id=lesson.get("working_revision_id"),
                stream_sequence=int(current_job.get("stream_sequence") or 0) + 1,
                stream_complete=True,
                error=None,
            )
        except Exception as exc:
            if isinstance(exc, asyncio.CancelledError):
                raise
            code = exc.code if isinstance(exc, TeacherLessonAuthoringError) else "lesson_plan_generation_failed"
            current_job = await asyncio.to_thread(
                self.repository.get_job,
                course_id,
                job_id,
            )
            return await asyncio.to_thread(
                self.repository.update_job,
                course_id,
                job_id,
                status="failed",
                phase="lesson_plan_failed",
                message="本讲教案生成失败",
                stream_sequence=int(current_job.get("stream_sequence") or 0) + 1,
                stream_complete=True,
                error={"code": code, "message": str(exc), "retryable": True},
            )

    async def run_script_job(
        self,
        *,
        course_id: str,
        lesson_unit_id: str,
        job_id: str,
        source_plan_revision_id: str,
        outline_sections: list[dict[str, Any]],
        plan_sections: dict[str, dict[str, Any]],
        generator: Callable[..., Awaitable[str]],
        shard_generator: Callable[..., Awaitable[dict[str, str]]] | None = None,
        seed_sections: list[dict[str, Any]] | None = None,
        requirements: str = "",
        material_asset_ids: list[str] | None = None,
        actor: str = "teacher",
    ) -> dict[str, Any]:
        """Generate one lesson's stable teaching blocks in a bounded wave.

        Partial blocks live in the durable job until every current-plan block is
        complete.  Each block reads only its deterministic plan-derived shard
        context, so it can run independently; local code owns checkpointing,
        ordering, validation and final assembly.
        """
        contracts: list[tuple[dict[str, Any], dict[str, Any], dict[str, Any]]] = []
        for outline_section in outline_sections:
            section_id = str(outline_section.get("node_id") or "")
            plan_section = plan_sections.get(section_id) or {}
            contract = compile_teacher_script_module_contract(
                outline_section,
                plan_section,
            )
            if not contract.get("modules"):
                raise TeacherLessonAuthoringError(
                    "lesson_script_contract_empty",
                    f"{contract.get('title') or section_id} 没有可生成的教学块。",
                )
            contracts.append((outline_section, plan_section, contract))

        total_blocks = sum(
            len(contract.get("modules") or [])
            for _outline, _plan, contract in contracts
        )
        global_block_order = {
            str(module.get("block_id") or ""): index
            for index, module in enumerate(
                [
                    module
                    for _outline, _plan, contract in contracts
                    for module in contract.get("modules") or []
                    if isinstance(module, dict)
                ]
            )
        }
        invalid_seed_ids: set[str] = set()
        if seed_sections:
            seed_quality = validate_teacher_script_revision(
                [
                    item for item in seed_sections
                    if isinstance(item, dict)
                ],
                generation_source="model_block_pipeline",
            )
            for issue in seed_quality.get("blocking_issues") or []:
                if not isinstance(issue, dict) or str(issue.get("code") or "") != (
                    "teacher_script:repetitive_blocks"
                ):
                    continue
                for group in issue.get("repeated_clause_groups") or []:
                    ordered_group = sorted(
                        [str(block_id) for block_id in group if str(block_id)],
                        key=lambda block_id: global_block_order.get(
                            block_id,
                            len(global_block_order),
                        ),
                    )
                    invalid_seed_ids.update(ordered_group[1:])
        seed_by_section = {
            str(item.get("section_node_id") or ""): item
            for item in seed_sections or []
            if isinstance(item, dict) and item.get("section_node_id")
        }
        completed_by_section: dict[str, list[dict[str, Any]]] = {}
        block_states: dict[str, str] = {}
        for _outline, _plan, contract in contracts:
            section_id = str(contract.get("section_node_id") or "")
            expected = [
                item for item in contract.get("modules") or [] if isinstance(item, dict)
            ]
            existing = {
                str(item.get("block_id") or ""): item
                for item in (seed_by_section.get(section_id) or {}).get("blocks") or []
                if isinstance(item, dict)
                and item.get("block_id")
                and str(item.get("content") or "").strip()
                and str(item.get("generation_source") or "") != "local_recovery"
                and str(item.get("block_id") or "") not in invalid_seed_ids
            }
            completed: list[dict[str, Any]] = []
            for module in expected:
                block_id = str(module.get("block_id") or "")
                previous = existing.get(block_id)
                if previous:
                    candidate = {
                        **deepcopy(module),
                        "content": str(previous.get("content") or "").strip(),
                        "generation_source": str(
                            previous.get("generation_source") or "model"
                        ),
                    }
                    single_contract = {
                        **deepcopy(contract),
                        "modules": [deepcopy(module)],
                    }
                    seed_report = validate_teacher_script_section(
                        {
                            "section_node_id": section_id,
                            "title": contract.get("title"),
                            "blocks": [candidate],
                        },
                        single_contract,
                    )
                    if seed_report.get("passed"):
                        completed.append(candidate)
                        block_states[block_id] = "completed"
                    else:
                        # Checkpoints are reusable evidence, not trusted final
                        # output. A truncated formula or newly tightened
                        # quality rule must regenerate only the affected block.
                        block_states[block_id] = "pending"
                else:
                    block_states[block_id] = "pending"
            completed_by_section[section_id] = completed

        def checkpoint_sections() -> list[dict[str, Any]]:
            result: list[dict[str, Any]] = []
            for _outline, _plan, contract in contracts:
                section_id = str(contract.get("section_node_id") or "")
                blocks = completed_by_section.get(section_id) or []
                if not blocks:
                    continue
                result.append(normalize_teacher_script_section({
                    "section_node_id": section_id,
                    "title": contract.get("title"),
                    "blocks": blocks,
                }, contract))
            return result

        completed_count = sum(
            1 for state in block_states.values() if state == "completed"
        )
        started_job = self.repository.update_job(
            course_id,
            job_id,
            status="running",
            phase="lesson_script_generation",
            progress=max(5, int(90 * completed_count / max(1, total_blocks))),
            message=(
                f"继续生成本讲讲义，已保留 {completed_count}/{total_blocks} 个教学块"
                if completed_count
                else "正在按当前教案生成本讲讲义"
            ),
            total_blocks=total_blocks,
            completed_blocks=completed_count,
            block_states=block_states,
            result_sections=checkpoint_sections(),
            stream_mode="",
            stream_complete=False,
            error=None,
        )
        if str(started_job.get("status") or "") not in TEACHER_JOB_ACTIVE_STATUSES:
            return started_job

        current_block_id = ""
        current_block_title = ""
        try:
            block_order = {
                str(module.get("block_id") or ""): index
                for index, module in enumerate(
                    [
                        module
                        for _outline, _plan, contract in contracts
                        for module in contract.get("modules") or []
                        if isinstance(module, dict)
                    ]
                )
            }
            block_entries: dict[str, dict[str, Any]] = {}
            for outline_section, plan_section, contract in contracts:
                for module in contract.get("modules") or []:
                    if not isinstance(module, dict):
                        continue
                    block_id = str(module.get("block_id") or "")
                    block_entries[block_id] = {
                        "outline_section": outline_section,
                        "plan_section": plan_section,
                        "contract": contract,
                        "module": module,
                    }
            shards: list[dict[str, Any]] = []
            if shard_generator:
                for planned_shard in compile_teacher_script_generation_shards(
                    lesson_unit_id,
                    [contract for _outline, _plan, contract in contracts],
                ):
                    entries = [
                        block_entries[block_id]
                        for block_id in planned_shard.get("block_ids") or []
                        if block_id in block_entries
                        and block_states.get(block_id) != "completed"
                    ]
                    if not entries:
                        continue
                    shards.append({
                        "entries": entries,
                        "context": deepcopy(planned_shard.get("context") or {}),
                        "shard_id": str(planned_shard.get("shard_id") or ""),
                    })
                    for entry in entries:
                        block_states[str(entry["module"].get("block_id") or "")] = (
                            "running"
                        )
            else:
                for block_id, entry in block_entries.items():
                    if block_states.get(block_id) == "completed":
                        continue
                    shards.append({
                        "entries": [entry],
                        "context": compile_teacher_script_shard_context(
                            entry["contract"],
                            entry["module"],
                        ),
                        "shard_id": f"{block_id}:shard:1",
                    })
                    block_states[block_id] = "running"

            if shards:
                self.repository.update_job_live(
                    course_id,
                    job_id,
                    phase="lesson_script_block_generation",
                    message=f"已将 {len(shards)} 个讲义分片并发入队",
                    current_block_id="",
                    current_block_title="",
                    block_states=block_states,
                )

            semaphore = asyncio.Semaphore(4)

            async def generate_shard(shard: dict[str, Any]) -> dict[str, Any]:
                entries = list(shard.get("entries") or [])
                modules = [entry["module"] for entry in entries]
                context = shard["context"]
                block_ids = [str(module.get("block_id") or "") for module in modules]
                block_titles = [
                    str(module.get("title") or "教学块") for module in modules
                ]
                block_title = " / ".join(block_titles)
                shard_id = str(shard.get("shard_id") or context.get("shard_id") or "")
                stream_state = {
                    "reset_blocks": set(),
                    "delta_blocks": set(),
                }

                async def persist_stream_reset(block_id: str = "") -> None:
                    targets = [block_id] if block_id else block_ids
                    for target_block_id in targets:
                        stream_key = f"{shard_id}:{target_block_id}"
                        stream_state["reset_blocks"].add(target_block_id)
                        stream_state["delta_blocks"].discard(target_block_id)
                        await asyncio.to_thread(
                            self.repository.update_job_stream,
                            course_id,
                            job_id,
                            phase="lesson_script_block_generation",
                            progress=max(
                                5,
                                int(90 * completed_count / max(1, total_blocks)),
                            ),
                            message=f"正在生成：{block_title}",
                            batch_id=stream_key,
                            event="reset",
                            lesson_unit_id=lesson_unit_id,
                            block_id=target_block_id,
                            shard_id=stream_key,
                            stream_mode="token_stream",
                        )

                async def persist_stream_delta(block_id: str, delta: str) -> None:
                    if not str(delta or ""):
                        return
                    if block_id not in stream_state["reset_blocks"]:
                        await persist_stream_reset(block_id)
                    stream_state["delta_blocks"].add(block_id)
                    stream_key = f"{shard_id}:{block_id}"
                    await asyncio.to_thread(
                        self.repository.update_job_stream,
                        course_id,
                        job_id,
                        phase="lesson_script_block_generation",
                        progress=max(
                            5,
                            int(90 * completed_count / max(1, total_blocks)),
                        ),
                        message=f"正在生成：{block_title}",
                        batch_id=stream_key,
                        event="delta",
                        delta=str(delta),
                        lesson_unit_id=lesson_unit_id,
                        block_id=block_id,
                        shard_id=stream_key,
                        stream_mode="token_stream",
                    )

                try:
                    async with semaphore:
                        current = await asyncio.to_thread(
                            self.repository.get_job,
                            course_id,
                            job_id,
                        )
                        if current.get("cancel_requested"):
                            raise asyncio.CancelledError
                        if shard_generator:
                            generated_map = await shard_generator(
                                deepcopy(entries),
                                deepcopy(context),
                                on_block_delta=persist_stream_delta,
                                on_shard_reset=persist_stream_reset,
                            )
                        else:
                            entry = entries[0]
                            module = entry["module"]
                            block_id = str(module.get("block_id") or "")
                            parameters = inspect.signature(generator).parameters
                            supports_stream_callbacks = (
                                "on_content_delta" in parameters
                                and "on_content_reset" in parameters
                            ) or any(
                                parameter.kind == inspect.Parameter.VAR_KEYWORD
                                for parameter in parameters.values()
                            )
                            if supports_stream_callbacks:
                                generated = await generator(
                                    entry["outline_section"],
                                    entry["plan_section"],
                                    module,
                                    deepcopy(context),
                                    on_content_delta=lambda delta: persist_stream_delta(
                                        block_id,
                                        delta,
                                    ),
                                    on_content_reset=lambda: persist_stream_reset(block_id),
                                )
                            else:
                                generated = await generator(
                                    entry["outline_section"],
                                    entry["plan_section"],
                                    module,
                                    deepcopy(context),
                                )
                            generated_map = {block_id: str(generated or "").strip()}
                    current = await asyncio.to_thread(
                        self.repository.get_job,
                        course_id,
                        job_id,
                    )
                    if current.get("cancel_requested"):
                        raise asyncio.CancelledError
                    if str(current.get("status") or "") not in {"pending", "running"}:
                        return {"terminal_job": current}
                    if not isinstance(generated_map, dict):
                        raise TeacherLessonAuthoringError(
                            "lesson_script_shard_invalid",
                            f"{block_title} 没有返回可定位的教学块。",
                        )
                    candidates: list[dict[str, Any]] = []
                    candidate_failures: list[dict[str, Any]] = []
                    for entry in entries:
                        module = entry["module"]
                        contract = entry["contract"]
                        block_id = str(module.get("block_id") or "")
                        content = str(generated_map.get(block_id) or "").strip()
                        if not content:
                            candidate_failures.append({
                                "block_id": block_id,
                                "title": str(module.get("title") or block_id),
                                "code": "lesson_script_block_empty",
                                "message": (
                                    f"{module.get('title') or block_id} "
                                    "没有生成有效内容。"
                                ),
                            })
                            continue
                        candidate = {
                            **deepcopy(module),
                            "content": content,
                            "generation_source": "model",
                        }
                        candidate_report = validate_teacher_script_section(
                            {
                                "section_node_id": str(
                                    contract.get("section_node_id") or ""
                                ),
                                "title": contract.get("title"),
                                "blocks": [candidate],
                            },
                            {**deepcopy(contract), "modules": [deepcopy(module)]},
                        )
                        if not candidate_report.get("passed"):
                            messages = "；".join(
                                str(item.get("message") or "未知讲义错误")
                                for item in candidate_report.get("blocking_issues") or []
                                if isinstance(item, dict)
                            )
                            candidate_failures.append({
                                "block_id": block_id,
                                "title": str(module.get("title") or block_id),
                                "code": "lesson_script_block_quality_failed",
                                "message": (
                                    f"{module.get('title') or block_id}未通过硬校验："
                                    f"{messages or '请重试'}"
                                ),
                                "quality_report": deepcopy(candidate_report),
                            })
                            continue
                        candidates.append({
                            "section_id": str(
                                contract.get("section_node_id") or ""
                            ),
                            "candidate": candidate,
                        })
                    return {
                        "block_ids": block_ids,
                        "block_title": block_title,
                        "shard_id": shard_id,
                        "candidates": candidates,
                        "failures": candidate_failures,
                        "streamed_block_ids": list(stream_state["delta_blocks"]),
                    }
                except asyncio.CancelledError:
                    raise
                except Exception as exc:
                    return {
                        "block_ids": block_ids,
                        "block_title": block_title,
                        "shard_id": shard_id,
                        "error": exc,
                    }

            tasks = [asyncio.create_task(generate_shard(shard)) for shard in shards]
            failed_shards: list[dict[str, Any]] = []
            failed_blocks: list[dict[str, Any]] = []
            try:
                for completed_task in asyncio.as_completed(tasks):
                    result = await completed_task
                    if result.get("terminal_job"):
                        for task in tasks:
                            task.cancel()
                        await asyncio.gather(*tasks, return_exceptions=True)
                        return result["terminal_job"]
                    block_ids = [
                        str(block_id) for block_id in result.get("block_ids") or []
                        if str(block_id)
                    ]
                    block_id = block_ids[0] if block_ids else ""
                    block_title = str(result.get("block_title") or "教学块")
                    shard_id = str(result.get("shard_id") or block_id)
                    if result.get("error") is not None:
                        error = result["error"]
                        for failed_block_id in block_ids:
                            block_states[failed_block_id] = "failed"
                        shard_failure = {
                            "block_id": block_id,
                            "block_ids": block_ids,
                            "shard_id": shard_id,
                            "title": block_title,
                            "code": (
                                error.code
                                if isinstance(error, TeacherLessonAuthoringError)
                                else "lesson_script_generation_failed"
                            ),
                            "message": str(error),
                        }
                        failed_shards.append(shard_failure)
                        for failed_block_id in block_ids:
                            failed_module = (
                                block_entries.get(failed_block_id) or {}
                            ).get("module") or {}
                            failed_blocks.append({
                                "block_id": failed_block_id,
                                "shard_id": shard_id,
                                "title": str(
                                    failed_module.get("title") or failed_block_id
                                ),
                                "code": shard_failure["code"],
                                "message": str(error),
                            })
                        self.repository.update_job_live(
                            course_id,
                            job_id,
                            phase="lesson_script_block_failed",
                            message=f"{block_title}生成失败，其他教学块继续",
                            current_block_id=block_id,
                            current_block_title=block_title,
                            block_states=block_states,
                        )
                        continue

                    shard_block_failures = [
                        item for item in result.get("failures") or []
                        if isinstance(item, dict) and item.get("block_id")
                    ]
                    if shard_block_failures:
                        failed_ids = [
                            str(item.get("block_id") or "")
                            for item in shard_block_failures
                        ]
                        for failure in shard_block_failures:
                            failed_block_id = str(failure.get("block_id") or "")
                            block_states[failed_block_id] = "failed"
                            failed_blocks.append({
                                **deepcopy(failure),
                                "block_id": failed_block_id,
                                "shard_id": shard_id,
                            })
                        failed_shards.append({
                            "block_id": failed_ids[0],
                            "block_ids": failed_ids,
                            "shard_id": shard_id,
                            "title": " / ".join(
                                str(item.get("title") or item.get("block_id") or "教学块")
                                for item in shard_block_failures
                            ),
                            "code": str(
                                shard_block_failures[0].get("code")
                                or "lesson_script_generation_failed"
                            ),
                            "message": str(
                                shard_block_failures[0].get("message")
                                or "讲义教学块未通过校验"
                            ),
                        })
                        self.repository.update_job_live(
                            course_id,
                            job_id,
                            phase="lesson_script_block_failed",
                            message=(
                                f"{len(shard_block_failures)} 个教学块未通过校验，"
                                "其他教学块继续"
                            ),
                            current_block_id=failed_ids[0],
                            current_block_title=str(
                                shard_block_failures[0].get("title") or "教学块"
                            ),
                            block_states=block_states,
                        )

                    streamed_block_ids = set(result.get("streamed_block_ids") or [])
                    for generated_item in result.get("candidates") or []:
                        if not isinstance(generated_item, dict):
                            continue
                        candidate = generated_item["candidate"]
                        section_id = str(generated_item.get("section_id") or "")
                        generated_block_id = str(candidate.get("block_id") or "")
                        completed_by_section[section_id].append(candidate)
                        completed_by_section[section_id].sort(
                            key=lambda item: block_order.get(
                                str(item.get("block_id") or ""),
                                len(block_order),
                            )
                        )
                        block_states[generated_block_id] = "completed"
                        completed_count += 1
                        if generated_block_id not in streamed_block_ids:
                            # A non-streaming provider exposes one honest whole
                            # block, never a timer-sliced imitation.
                            for event, delta in (
                                ("reset", ""),
                                ("delta", str(candidate.get("content") or "")),
                            ):
                                self.repository.update_job_stream(
                                    course_id,
                                    job_id,
                                    phase="lesson_script_block_saved",
                                    progress=max(
                                        5,
                                        min(
                                            95,
                                            int(95 * completed_count / max(1, total_blocks)),
                                        ),
                                    ),
                                    message=f"已生成 {completed_count}/{total_blocks} 个教学块",
                                    batch_id=f"{shard_id}:{generated_block_id}",
                                    event=event,
                                    delta=delta,
                                    lesson_unit_id=lesson_unit_id,
                                    block_id=generated_block_id,
                                    shard_id=f"{shard_id}:{generated_block_id}",
                                    stream_mode="buffered_fallback",
                                )
                    self.repository.update_job(
                        course_id,
                        job_id,
                        phase="lesson_script_block_saved",
                        progress=max(
                            5,
                            min(95, int(95 * completed_count / max(1, total_blocks))),
                        ),
                        message=f"已生成 {completed_count}/{total_blocks} 个教学块",
                        completed_blocks=completed_count,
                        current_block_id="",
                        current_block_title="",
                        block_states=block_states,
                        result_sections=checkpoint_sections(),
                    )
            except asyncio.CancelledError:
                for task in tasks:
                    task.cancel()
                await asyncio.gather(*tasks, return_exceptions=True)
                raise

            if failed_blocks:
                first = failed_blocks[0]
                current_block_id = str(first.get("block_id") or "")
                current_block_title = str(first.get("title") or "教学块")
                raise TeacherLessonAuthoringError(
                    str(first.get("code") or "lesson_script_generation_failed"),
                    f"{len(failed_blocks)} 个教学块生成失败，已保留其他成功结果。",
                    details={
                        "failed_shards": failed_shards,
                        "failed_blocks": failed_blocks,
                    },
                )

            final_sections: list[dict[str, Any]] = []
            for _outline, _plan, contract in contracts:
                section_id = str(contract.get("section_node_id") or "")
                section = normalize_teacher_script_section({
                    "section_node_id": section_id,
                    "title": contract.get("title"),
                    "blocks": completed_by_section.get(section_id) or [],
                }, contract)
                section["quality_report"] = validate_teacher_script_section(
                    section,
                    contract,
                )
                section["pipeline_version"] = SCRIPT_PIPELINE_VERSION
                if not section["quality_report"].get("passed"):
                    messages = "；".join(
                        str(item.get("message") or "未知讲义错误")
                        for item in section["quality_report"].get("blocking_issues") or []
                        if isinstance(item, dict)
                    )
                    raise TeacherLessonAuthoringError(
                        "lesson_script_section_quality_failed",
                        f"{section.get('title') or section_id}未通过硬校验：{messages or '请重试'}",
                        details={"quality_report": deepcopy(section["quality_report"])},
                    )
                final_sections.append(section)

            revision_quality = validate_teacher_script_revision(
                final_sections,
                generation_source="model_block_pipeline",
            )
            if not revision_quality.get("passed"):
                messages = "；".join(
                    str(item.get("message") or "未知讲义错误")
                    for item in revision_quality.get("blocking_issues") or []
                    if isinstance(item, dict)
                )
                raise TeacherLessonAuthoringError(
                    "lesson_script_quality_failed",
                    f"本讲讲义未通过硬校验：{messages or '请重试'}",
                    details={"quality_report": deepcopy(revision_quality)},
                )
            lesson = self.repository.save_script_revision(
                course_id,
                lesson_unit_id,
                final_sections,
                source_lesson_plan_revision_id=source_plan_revision_id,
                generation_source="model_block_pipeline",
                requirements=requirements,
                material_asset_ids=material_asset_ids or [],
                actor=actor,
                active_job_id=job_id,
            )
            current_job = self.repository.get_job(course_id, job_id)
            return self.repository.update_job(
                course_id,
                job_id,
                status="completed",
                phase="lesson_script_ready",
                progress=100,
                message="本讲讲义已生成",
                completed_blocks=total_blocks,
                result_sections=final_sections,
                result_revision_id=str(lesson.get("working_script_revision_id") or ""),
                warnings=[],
                current_block_id="",
                current_block_title="",
                stream_sequence=int(current_job.get("stream_sequence") or 0) + 1,
                stream_complete=True,
                error=None,
            )
        except Exception as exc:
            if isinstance(exc, asyncio.CancelledError):
                raise
            if current_block_id:
                block_states[current_block_id] = "failed"
            code = (
                exc.code
                if isinstance(exc, TeacherLessonAuthoringError)
                else "lesson_script_generation_failed"
            )
            current_job = self.repository.get_job(course_id, job_id)
            return self.repository.update_job(
                course_id,
                job_id,
                status="failed",
                phase="lesson_script_failed",
                progress=max(5, int(95 * completed_count / max(1, total_blocks))),
                message=f"讲义生成暂停，已保留 {completed_count}/{total_blocks} 个教学块",
                completed_blocks=completed_count,
                current_block_id=current_block_id,
                current_block_title=current_block_title,
                block_states=block_states,
                result_sections=checkpoint_sections(),
                stream_sequence=int(current_job.get("stream_sequence") or 0) + 1,
                stream_complete=True,
                error={
                    "code": code,
                    "message": str(exc),
                    "retryable": True,
                    **(
                        {"failed_shards": deepcopy(exc.details["failed_shards"])}
                        if isinstance(exc, TeacherLessonAuthoringError)
                        and isinstance(exc.details.get("failed_shards"), list)
                        else {}
                    ),
                    **(
                        {"failed_blocks": deepcopy(exc.details["failed_blocks"])}
                        if isinstance(exc, TeacherLessonAuthoringError)
                        and isinstance(exc.details.get("failed_blocks"), list)
                        else {}
                    ),
                },
            )
