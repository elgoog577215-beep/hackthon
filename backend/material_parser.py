"""成熟解析器适配与统一 ParsedDocument 归一化。"""

from __future__ import annotations

import asyncio
import hashlib
import importlib.metadata
import mimetypes
import os
import re
import shutil
import subprocess
import tempfile
import uuid
import zipfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Protocol

from material_models import DocumentBlock, DocumentLocator, MaterialAsset, ParsedDocument
from material_storage import (
    IMAGE_EXTENSIONS,
    LEGACY_OFFICE_EXTENSIONS,
    TEXT_EXTENSIONS,
    MaterialRepository,
)

PARSE_OPTIONS_VERSION = "material_parse_v1"


class DocumentParser(Protocol):
    name: str

    def supports(self, extension: str) -> bool: ...

    def parse(self, asset: MaterialAsset, source_path: Path) -> ParsedDocument: ...


def _now() -> str:
    return datetime.now(timezone.utc).isoformat()


def _compact(value: Any, limit: int) -> str:
    text = re.sub(r"\s+", " ", str(value or "")).strip()
    return text if len(text) <= limit else f"{text[: limit - 1].rstrip()}…"


def _options_hash(parser_name: str) -> str:
    return hashlib.sha256(f"{PARSE_OPTIONS_VERSION}:{parser_name}".encode()).hexdigest()[:16]


def _quality(blocks: list[DocumentBlock]) -> dict[str, Any]:
    text_chars = sum(len(block.text) for block in blocks)
    located = sum(
        1
        for block in blocks
        if block.locator.page is not None
        or block.locator.slide is not None
        or block.locator.section_path
    )
    pages = [block.locator.page for block in blocks if block.locator.page]
    slides = [block.locator.slide for block in blocks if block.locator.slide]
    return {
        "block_count": len(blocks),
        "text_chars": text_chars,
        "located_blocks": located,
        "location_coverage": round(located / max(1, len(blocks)), 3),
        "page_count": max(pages, default=0),
        "slide_count": max(slides, default=0),
    }


class TextDocumentParser:
    name = "builtin_text"
    version = "1"

    def supports(self, extension: str) -> bool:
        return extension in TEXT_EXTENSIONS

    def parse(self, asset: MaterialAsset, source_path: Path) -> ParsedDocument:
        text = source_path.read_text(encoding="utf-8")
        blocks = self._to_blocks(text)
        return ParsedDocument(
            document_id=f"doc-{uuid.uuid4().hex}",
            asset_id=asset.asset_id,
            source_sha256=asset.sha256,
            parse_status="parsed" if blocks else "metadata_only",
            parser_name=self.name,
            parser_version=self.version,
            parse_options_hash=_options_hash(self.name),
            blocks=blocks,
            quality=_quality(blocks),
            warnings=[] if blocks else ["文本文件没有可用正文"],
            created_at=_now(),
        )

    @staticmethod
    def _to_blocks(text: str) -> list[DocumentBlock]:
        blocks: list[DocumentBlock] = []
        section_path: list[str] = []
        paragraph: list[str] = []

        def flush() -> None:
            if not paragraph:
                return
            value = "\n".join(paragraph).strip()
            paragraph.clear()
            if not value:
                return
            blocks.append(DocumentBlock(
                block_id=f"blk-{len(blocks) + 1}",
                kind=_detect_block_kind(value),
                text=value,
                order=len(blocks),
                locator=DocumentLocator(section_path=list(section_path)),
            ))

        for raw in text.splitlines():
            line = raw.rstrip()
            heading = re.match(r"^(#{1,6})\s+(.+)$", line)
            if heading:
                flush()
                level = len(heading.group(1))
                title = heading.group(2).strip()
                section_path[:] = section_path[: level - 1]
                section_path.append(title)
                blocks.append(DocumentBlock(
                    block_id=f"blk-{len(blocks) + 1}",
                    kind="title" if level == 1 else "heading",
                    text=title,
                    order=len(blocks),
                    locator=DocumentLocator(section_path=list(section_path)),
                    metadata={"heading_level": level},
                ))
            elif not line.strip():
                flush()
            else:
                paragraph.append(line)
        flush()
        return blocks


class ImageOcrParser:
    name = "rapidocr"
    version = "1"

    def supports(self, extension: str) -> bool:
        return extension in IMAGE_EXTENSIONS

    def parse(self, asset: MaterialAsset, source_path: Path) -> ParsedDocument:
        segments = _ocr_image(source_path)
        blocks: list[DocumentBlock] = []
        confidences: list[float] = []
        for segment in segments:
            text = str(segment.get("text") or "").strip()
            if not text:
                continue
            confidence = max(0.0, min(1.0, float(segment.get("confidence") or 0)))
            confidences.append(confidence)
            blocks.append(DocumentBlock(
                block_id=f"blk-{len(blocks) + 1}",
                kind=_detect_block_kind(text),
                text=text,
                order=len(blocks),
                locator=DocumentLocator(
                    page=max(1, int(segment.get("page") or 1)),
                    bbox=_normalized_bbox(segment.get("bbox")),
                ),
                metadata={
                    "ocr_engine": self.name,
                    "ocr_confidence": round(confidence, 4),
                },
            ))
        if not blocks:
            raise RuntimeError("OCR 没有从图片中提取到可用文字")
        average_confidence = round(sum(confidences) / max(1, len(confidences)), 4)
        quality = {
            **_quality(blocks),
            "ocr_confidence": average_confidence,
            "ocr_engine": self.name,
        }
        degraded = average_confidence < 0.85
        return ParsedDocument(
            document_id=f"doc-{uuid.uuid4().hex}",
            asset_id=asset.asset_id,
            source_sha256=asset.sha256,
            parse_status="degraded" if degraded else "parsed",
            parser_name=self.name,
            parser_version=self.version,
            parse_options_hash=_options_hash(self.name),
            blocks=blocks,
            quality=quality,
            warnings=(
                ["OCR 平均置信度低于 0.85，相关题目必须进入教师审核"]
                if degraded
                else []
            ),
            created_at=_now(),
        )


class DoclingDocumentParser:
    name = "docling"

    @property
    def version(self) -> str:
        for package in ("docling-slim", "docling"):
            try:
                return importlib.metadata.version(package)
            except importlib.metadata.PackageNotFoundError:
                continue
        return "unavailable"

    def supports(self, extension: str) -> bool:
        return extension in {".pdf", ".docx", ".pptx", ".xlsx"}

    def parse(self, asset: MaterialAsset, source_path: Path) -> ParsedDocument:
        try:
            from docling.backend.msexcel_backend import MsExcelDocumentBackend
            from docling.backend.mspowerpoint_backend import MsPowerpointDocumentBackend
            from docling.backend.msword_backend import MsWordDocumentBackend
            from docling.backend.pypdfium2_backend import PyPdfiumDocumentBackend
            from docling.datamodel.base_models import InputFormat
            from docling.datamodel.document import InputDocument
        except ImportError as exc:
            raise RuntimeError("Docling 未安装") from exc
        formats: dict[str, tuple[Any, Any]] = {
            ".docx": (InputFormat.DOCX, MsWordDocumentBackend),
            ".pptx": (InputFormat.PPTX, MsPowerpointDocumentBackend),
            ".xlsx": (InputFormat.XLSX, MsExcelDocumentBackend),
            ".pdf": (InputFormat.PDF, PyPdfiumDocumentBackend),
        }
        input_format, backend_class = formats[asset.extension]
        input_document = InputDocument(source_path, format=input_format, backend=backend_class)
        if not input_document.valid:
            raise RuntimeError("Docling 文件后端无法打开该资料")
        backend = input_document._backend
        warnings: list[str] = []
        parse_status = "parsed"
        visual_quality: dict[str, Any] = {}
        if asset.extension == ".pdf":
            blocks = _blocks_from_pdf_backend(backend)
        else:
            document = backend.convert()
            blocks = _blocks_from_docling(document.export_to_dict(), asset.extension)
            if asset.extension == ".docx":
                blocks, docx_quality, docx_warnings = _enrich_docx_structure(
                    source_path,
                    blocks,
                )
                visual_quality.update(docx_quality)
                warnings.extend(docx_warnings)
                if docx_quality.get("tracked_change_count") or docx_quality.get("comment_part_present"):
                    parse_status = "degraded"
            elif asset.extension == ".pptx":
                visual_blocks, visual_quality, visual_warnings = _pptx_visual_evidence(
                    source_path,
                    starting_order=len(blocks),
                )
                blocks.extend(visual_blocks)
                warnings.extend(visual_warnings)
                if visual_quality.get("animation_slide_count") or visual_quality.get("unread_picture_count"):
                    parse_status = "degraded"
        if not blocks:
            raise RuntimeError("Docling 没有提取到可用文本")
        return ParsedDocument(
            document_id=f"doc-{uuid.uuid4().hex}",
            asset_id=asset.asset_id,
            source_sha256=asset.sha256,
            parse_status=parse_status,
            parser_name=self.name,
            parser_version=self.version,
            parse_options_hash=_options_hash(self.name),
            blocks=blocks,
            quality={**_quality(blocks), **visual_quality},
            warnings=warnings,
            created_at=_now(),
        )


class ScannedPdfOcrParser:
    """Local page OCR for PDFs whose text layer is empty.

    This adapter runs only after the lightweight PDF text backend fails.  It
    renders pages in memory, keeps page locators and confidence, and never sends
    teacher material to an external service.
    """

    name = "rapidocr_pdf"

    @property
    def version(self) -> str:
        for package in ("rapidocr", "rapidocr-onnxruntime"):
            try:
                return importlib.metadata.version(package)
            except importlib.metadata.PackageNotFoundError:
                continue
        return "unavailable"

    def supports(self, extension: str) -> bool:
        return extension == ".pdf"

    def parse(self, asset: MaterialAsset, source_path: Path) -> ParsedDocument:
        try:
            import pypdfium2 as pdfium
        except ImportError as exc:
            raise RuntimeError("扫描 PDF OCR 组件未安装") from exc
        RapidOCR = _rapidocr_class()

        maximum_pages = max(1, int(os.getenv("MATERIAL_OCR_MAX_PAGES", "250")))
        document = pdfium.PdfDocument(str(source_path))
        page_count = len(document)
        if page_count > maximum_pages:
            document.close()
            raise RuntimeError(f"扫描 PDF 共 {page_count} 页，超过 OCR 上限 {maximum_pages} 页")

        engine = RapidOCR()
        blocks: list[DocumentBlock] = []
        confidences: list[float] = []
        try:
            with tempfile.TemporaryDirectory(prefix="lingzhi-pdf-ocr-") as temp_value:
                temp_dir = Path(temp_value)
                for page_index in range(page_count):
                    page = document[page_index]
                    image_path = temp_dir / f"page-{page_index + 1}.png"
                    try:
                        page.render(scale=2.0).to_pil().save(image_path, format="PNG")
                    finally:
                        page.close()
                    for segment in _ocr_image(
                        image_path,
                        engine=engine,
                        page_number=page_index + 1,
                    ):
                        text = str(segment.get("text") or "").strip()
                        if not text:
                            continue
                        confidence = max(0.0, min(1.0, float(segment.get("confidence") or 0)))
                        confidences.append(confidence)
                        blocks.append(DocumentBlock(
                            block_id=f"blk-{len(blocks) + 1}",
                            kind=_detect_block_kind(text),
                            text=text,
                            order=len(blocks),
                            locator=DocumentLocator(
                                page=page_index + 1,
                                bbox=_normalized_bbox(segment.get("bbox")),
                            ),
                            metadata={
                                "ocr_engine": self.name,
                                "ocr_confidence": round(confidence, 4),
                            },
                        ))
        finally:
            document.close()
        if not blocks:
            raise RuntimeError("OCR 没有从扫描 PDF 中提取到可用文字")
        average_confidence = round(sum(confidences) / max(1, len(confidences)), 4)
        warnings = ["扫描 PDF 由本地 OCR 提取，版式、手写内容和图片语义需要教师复核"]
        if average_confidence < 0.85:
            warnings.append("OCR 平均置信度低于 0.85，相关内容必须进入教师审核")
        return ParsedDocument(
            document_id=f"doc-{uuid.uuid4().hex}",
            asset_id=asset.asset_id,
            source_sha256=asset.sha256,
            parse_status="degraded",
            parser_name=self.name,
            parser_version=self.version,
            parse_options_hash=_options_hash(self.name),
            blocks=blocks,
            quality={
                **_quality(blocks),
                "ocr_confidence": average_confidence,
                "ocr_engine": self.name,
                "source_page_count": page_count,
            },
            warnings=warnings,
            created_at=_now(),
        )


class LegacyOfficeConversionParser:
    """Convert legacy binary Office files in a disposable directory."""

    name = "libreoffice_legacy"
    version = "1"
    targets = {".doc": ".docx", ".ppt": ".pptx", ".xls": ".xlsx"}

    def supports(self, extension: str) -> bool:
        return extension in LEGACY_OFFICE_EXTENSIONS

    def parse(self, asset: MaterialAsset, source_path: Path) -> ParsedDocument:
        executable = shutil.which("soffice") or shutil.which("libreoffice")
        if not executable:
            raise RuntimeError("旧版 Office 转换组件不可用")
        target_extension = self.targets[asset.extension]
        with tempfile.TemporaryDirectory(prefix="lingzhi-office-") as temp_value:
            temp_dir = Path(temp_value)
            input_path = temp_dir / f"source{asset.extension}"
            shutil.copy2(source_path, input_path)
            profile_dir = temp_dir / "profile"
            result = subprocess.run(
                [
                    executable,
                    "--headless",
                    f"-env:UserInstallation={profile_dir.as_uri()}",
                    "--convert-to",
                    target_extension.lstrip("."),
                    "--outdir",
                    str(temp_dir),
                    str(input_path),
                ],
                capture_output=True,
                check=False,
                text=True,
                timeout=120,
            )
            converted_path = temp_dir / f"source{target_extension}"
            if result.returncode != 0 or not converted_path.is_file():
                detail = _compact(result.stderr or result.stdout or "转换没有生成目标文件", 500)
                raise RuntimeError(f"旧版 Office 转换失败：{detail}")
            converted_asset = asset.model_copy(update={"extension": target_extension})
            errors: list[str] = []
            for parser in (DoclingDocumentParser(), MarkItDownFallbackParser()):
                try:
                    document = parser.parse(converted_asset, converted_path)
                    document.parser_name = f"{self.name}+{document.parser_name}"
                    document.parser_version = f"{self.version}+{document.parser_version}"
                    document.parse_options_hash = _options_hash(document.parser_name)
                    document.warnings = [
                        "旧版 Office 原件已在临时目录转换；正文可用，但复杂版式、宏和动画需复核",
                        *document.warnings,
                    ]
                    if asset.extension == ".ppt" and document.parse_status == "parsed":
                        document.parse_status = "degraded"
                    return document
                except Exception as exc:
                    errors.append(f"{parser.name}: {exc}")
            raise RuntimeError("；".join(errors) or "转换后的 Office 文件无法解析")


class MarkItDownFallbackParser:
    name = "markitdown"

    @property
    def version(self) -> str:
        try:
            return importlib.metadata.version("markitdown")
        except importlib.metadata.PackageNotFoundError:
            return "unavailable"

    def supports(self, extension: str) -> bool:
        return extension in {".pdf", ".docx", ".pptx", ".xlsx"}

    def parse(self, asset: MaterialAsset, source_path: Path) -> ParsedDocument:
        try:
            from markitdown import MarkItDown
        except ImportError as exc:
            raise RuntimeError("MarkItDown 未安装") from exc
        result = MarkItDown().convert(str(source_path))
        text = str(getattr(result, "text_content", "") or "").strip()
        blocks = TextDocumentParser._to_blocks(text)
        if not blocks:
            raise RuntimeError("降级解析器没有提取到可用内容")
        return ParsedDocument(
            document_id=f"doc-{uuid.uuid4().hex}",
            asset_id=asset.asset_id,
            source_sha256=asset.sha256,
            parse_status="degraded",
            parser_name=self.name,
            parser_version=self.version,
            parse_options_hash=_options_hash(self.name),
            blocks=blocks,
            quality=_quality(blocks),
            warnings=["当前资料仅完成文本降级提取，页码、布局或 OCR 来源可能不完整"],
            created_at=_now(),
        )


async def parse_material_asset(
    repository: MaterialRepository,
    asset: MaterialAsset,
) -> ParsedDocument:
    cached = repository.load_parsed_document(asset.asset_id)
    if cached and cached.source_sha256 == asset.sha256 and cached.parse_status in {"parsed", "degraded"}:
        return cached

    repository.update_status(asset.asset_id, "parsing")
    source = repository.source_path(asset)
    parsers: list[DocumentParser]
    if asset.extension in TEXT_EXTENSIONS:
        parsers = [TextDocumentParser()]
    elif asset.extension in IMAGE_EXTENSIONS:
        parsers = [ImageOcrParser()]
    elif asset.extension in LEGACY_OFFICE_EXTENSIONS:
        parsers = [LegacyOfficeConversionParser()]
    elif asset.extension == ".pdf":
        parsers = [DoclingDocumentParser(), ScannedPdfOcrParser(), MarkItDownFallbackParser()]
    else:
        parsers = [DoclingDocumentParser(), MarkItDownFallbackParser()]

    errors: list[str] = []
    for parser in parsers:
        if not parser.supports(asset.extension):
            continue
        try:
            document = await asyncio.to_thread(parser.parse, asset, source)
            repository.save_parsed_document(document)
            repository.update_status(
                asset.asset_id,
                document.parse_status,
                warnings=document.warnings,
                parser_name=document.parser_name,
                parser_version=document.parser_version,
                parse_options_hash=document.parse_options_hash,
                parse_quality=document.quality,
            )
            return document
        except Exception as exc:
            errors.append(f"{parser.name}: {exc}")

    message = "；".join(errors) or "没有可用解析器"
    failed = ParsedDocument(
        document_id=f"doc-{uuid.uuid4().hex}",
        asset_id=asset.asset_id,
        source_sha256=asset.sha256,
        parse_status="failed",
        parser_name="none",
        parser_version="",
        parse_options_hash=_options_hash("none"),
        blocks=[],
        quality=_quality([]),
        error=message,
        created_at=_now(),
    )
    repository.save_parsed_document(failed)
    repository.update_status(asset.asset_id, "failed", error=message)
    return failed


async def parse_document_path(
    source_path: Path,
    *,
    asset_id: str,
    filename: str,
) -> ParsedDocument:
    """Parse a course-space source without creating a second material asset.

    Course-space uploads may be stored either as a reference to ``mat-*`` or as
    package-owned bytes.  The latter still needs the same mature parser chain;
    this adapter deliberately returns a transient ``ParsedDocument`` instead
    of copying the file into ``material_storage`` and creating a parallel
    source of truth.
    """
    path = Path(source_path)
    extension = path.suffix.lower()
    detected_mime = mimetypes.guess_type(filename)[0] or "application/octet-stream"
    asset = MaterialAsset(
        asset_id=asset_id,
        filename=filename or path.name,
        extension=extension,
        mime_type=detected_mime,
        detected_mime=detected_mime,
        size_bytes=path.stat().st_size,
        sha256=hashlib.sha256(path.read_bytes()).hexdigest(),
        source_name=path.name,
        uploaded_at=_now(),
        updated_at=_now(),
    )
    if extension in TEXT_EXTENSIONS:
        parsers: list[DocumentParser] = [TextDocumentParser()]
    elif extension in IMAGE_EXTENSIONS:
        parsers = [ImageOcrParser()]
    elif extension in LEGACY_OFFICE_EXTENSIONS:
        parsers = [LegacyOfficeConversionParser()]
    elif extension == ".pdf":
        parsers = [DoclingDocumentParser(), ScannedPdfOcrParser(), MarkItDownFallbackParser()]
    else:
        parsers = [DoclingDocumentParser(), MarkItDownFallbackParser()]
    errors: list[str] = []
    for parser in parsers:
        if not parser.supports(extension):
            continue
        try:
            return await asyncio.to_thread(parser.parse, asset, path)
        except Exception as exc:
            errors.append(f"{parser.name}: {exc}")
    raise RuntimeError("；".join(errors) or f"不支持解析 {extension or '未知格式'}")


def _blocks_from_docling(data: dict[str, Any], extension: str) -> list[DocumentBlock]:
    blocks: list[DocumentBlock] = []
    visited: set[str] = set()
    section_path: list[str] = []

    def resolve(ref: str) -> dict[str, Any] | None:
        if not ref.startswith("#/"):
            return None
        value: Any = data
        try:
            for part in ref[2:].split("/"):
                value = value[int(part)] if isinstance(value, list) else value[part]
            return value if isinstance(value, dict) else None
        except (KeyError, IndexError, ValueError, TypeError):
            return None

    def visit(item: dict[str, Any]) -> None:
        ref = str(item.get("self_ref") or "")
        if ref and ref in visited:
            return
        if ref:
            visited.add(ref)
        label = str(item.get("label") or "text")
        text = str(item.get("text") or item.get("orig") or "").strip()
        if not text and label == "table":
            cells = (item.get("data") or {}).get("table_cells") or []
            text = " | ".join(str(cell.get("text") or "").strip() for cell in cells if str(cell.get("text") or "").strip())
        if text:
            kind = _docling_kind(label, text)
            if kind in {"title", "heading"}:
                level = int(item.get("level") or (1 if kind == "title" else 2))
                section_path[:] = section_path[: max(0, level - 1)]
                section_path.append(text[:200])
            provenance = item.get("prov") or []
            first = provenance[0] if provenance else {}
            page_no = first.get("page_no")
            locator = DocumentLocator(
                page=int(page_no) if page_no else None,
                slide=int(page_no) if page_no and extension == ".pptx" else None,
                section_path=list(section_path),
                bbox=_normalized_bbox(first.get("bbox")),
            )
            blocks.append(DocumentBlock(
                block_id=f"blk-{len(blocks) + 1}",
                kind=kind,
                text=text,
                order=len(blocks),
                locator=locator,
                metadata={"docling_label": label, "source_ref": ref},
            ))
        for child in item.get("children") or []:
            child_item = resolve(str(child.get("$ref") or "")) if isinstance(child, dict) else None
            if child_item:
                visit(child_item)

    root = data.get("body") or {}
    visit(root)
    if not blocks:
        for collection in ("texts", "tables", "pictures"):
            for item in data.get(collection) or []:
                visit(item)
    return blocks


def _blocks_from_pdf_backend(backend: Any) -> list[DocumentBlock]:
    blocks: list[DocumentBlock] = []
    try:
        for page_number, page in enumerate(backend.iter_pages(), start=1):
            try:
                cells = [
                    cell for cell in page.get_text_cells()
                    if re.sub(r"\s+", " ", str(getattr(cell, "text", "") or "")).strip()
                ]
                for cell_index, cell in enumerate(cells, start=1):
                    text = re.sub(r"\s+", " ", str(getattr(cell, "text", "") or "")).strip()
                    blocks.append(DocumentBlock(
                        block_id=f"blk-{len(blocks) + 1}",
                        kind=_detect_block_kind(text),
                        text=text,
                        order=len(blocks),
                        locator=DocumentLocator(
                            page=page_number,
                            bbox=_pdf_cell_bbox(cell),
                        ),
                        metadata={
                            "evidence_kind": "pdf_text_cell",
                            "page_cell_index": cell_index,
                        },
                    ))
            finally:
                page.unload()
    finally:
        backend.unload()
    return blocks


def _pdf_cell_bbox(cell: Any) -> dict[str, float] | None:
    raw = getattr(cell, "rect", None) or getattr(cell, "bbox", None)
    if isinstance(raw, dict):
        return _normalized_bbox(raw)
    values: dict[str, float] = {}
    aliases = {
        "l": ("l", "left", "x0"),
        "t": ("t", "top", "y0"),
        "r": ("r", "right", "x1"),
        "b": ("b", "bottom", "y1"),
    }
    for target, names in aliases.items():
        for name in names:
            value = getattr(raw, name, None)
            if isinstance(value, (int, float)):
                values[target] = float(value)
                break
    return values or None


def _enrich_docx_structure(
    source_path: Path,
    blocks: list[DocumentBlock],
) -> tuple[list[DocumentBlock], dict[str, Any], list[str]]:
    """Preserve Word structure and explicit review boundaries without duplicating prose."""
    try:
        from docx import Document
    except ImportError as exc:
        raise RuntimeError("DOCX 结构识别组件未安装") from exc

    document = Document(str(source_path))
    normalized_blocks = [re.sub(r"\s+", " ", item.text).strip() for item in blocks]
    used_block_indexes: set[int] = set()
    section_path: list[str] = []
    heading_count = 0
    list_item_count = 0
    for paragraph_index, paragraph in enumerate(document.paragraphs, start=1):
        text = re.sub(r"\s+", " ", paragraph.text or "").strip()
        if not text:
            continue
        style_name = str(getattr(paragraph.style, "name", "") or "")
        heading_match = re.search(r"(?:Heading|Title|标题)\s*(\d+)?", style_name, re.I)
        heading_level = int(heading_match.group(1) or 1) if heading_match else 0
        if heading_level:
            heading_count += 1
            section_path[:] = section_path[: max(0, heading_level - 1)]
            section_path.append(text[:200])
        has_numbering = bool(paragraph._p.xpath("./w:pPr/w:numPr")) or bool(
            re.search(r"(?:^|\s)List(?:\s|$)|列表", style_name, re.I)
        )
        if has_numbering:
            list_item_count += 1
        block_index = next((
            index for index, value in enumerate(normalized_blocks)
            if index not in used_block_indexes and value == text
        ), None)
        if block_index is None:
            continue
        used_block_indexes.add(block_index)
        block = blocks[block_index]
        block.metadata = {
            **block.metadata,
            "docx_paragraph_index": paragraph_index,
            "docx_style": style_name,
            "heading_level": heading_level,
            "numbered_or_bulleted": has_numbering,
        }
        if heading_level:
            block.kind = "title" if heading_level == 1 else "heading"
        elif has_numbering:
            block.kind = "list_item"
        block.locator.section_path = list(section_path)

    table_cell_count = 0
    for table_index, table in enumerate(document.tables, start=1):
        matrix = [
            [re.sub(r"\s+", " ", cell.text or "").strip() for cell in row.cells]
            for row in table.rows
        ]
        table_cell_count += sum(len(row) for row in matrix)
        table_block = next((
            item for item in blocks
            if item.kind == "table" and "docx_table_index" not in item.metadata
        ), None)
        if table_block is not None:
            table_block.metadata = {
                **table_block.metadata,
                "docx_table_index": table_index,
                "rows": matrix,
                "row_count": len(matrix),
                "column_count": max((len(row) for row in matrix), default=0),
            }

    header_footer_texts: list[tuple[str, str]] = []
    for section_index, section in enumerate(document.sections, start=1):
        for location, container in (("header", section.header), ("footer", section.footer)):
            text = "\n".join(
                paragraph.text.strip() for paragraph in container.paragraphs
                if paragraph.text.strip()
            )
            if text and (location, text) not in header_footer_texts:
                header_footer_texts.append((location, text))
                blocks.append(DocumentBlock(
                    block_id=f"blk-{len(blocks) + 1}",
                    kind="other",
                    text=text,
                    order=len(blocks),
                    metadata={
                        "evidence_kind": f"docx_{location}",
                        "docx_section_index": section_index,
                    },
                ))

    tracked_change_count = 0
    comment_part_present = False
    formula_count = 0
    media_count = 0
    try:
        with zipfile.ZipFile(source_path) as archive:
            names = set(archive.namelist())
            document_xml = archive.read("word/document.xml") if "word/document.xml" in names else b""
            tracked_change_count = document_xml.count(b"<w:ins") + document_xml.count(b"<w:del")
            formula_count = document_xml.count(b"<m:oMath")
            comment_part_present = "word/comments.xml" in names
            media_count = sum(1 for name in names if name.startswith("word/media/"))
    except (OSError, KeyError, zipfile.BadZipFile):
        pass

    warnings: list[str] = []
    if media_count:
        warnings.append(f"DOCX 含 {media_count} 个图片或媒体对象，未将图像内容冒充为可引用事实")
    if tracked_change_count:
        warnings.append("DOCX 含修订痕迹，当前保留检测结果，采用哪个修订版本需要教师确认")
    if comment_part_present:
        warnings.append("DOCX 含批注，批注不会自动合并到正文")
    return blocks, {
        "heading_count": heading_count,
        "list_item_count": list_item_count,
        "table_count": len(document.tables),
        "table_cell_count": table_cell_count,
        "header_footer_block_count": len(header_footer_texts),
        "picture_or_media_count": media_count,
        "formula_count": formula_count,
        "tracked_change_count": tracked_change_count,
        "comment_part_present": comment_part_present,
    }, warnings


def _pptx_visual_evidence(
    source_path: Path,
    *,
    starting_order: int,
) -> tuple[list[DocumentBlock], dict[str, Any], list[str]]:
    """Extract compact visual evidence without inventing picture semantics."""
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:
        raise RuntimeError("PPTX 视觉证据组件未安装") from exc

    presentation = Presentation(str(source_path))
    slide_width = max(1, int(presentation.slide_width or 1))
    slide_height = max(1, int(presentation.slide_height or 1))
    blocks: list[DocumentBlock] = []
    chart_count = 0
    picture_count = 0
    unread_picture_count = 0
    note_slide_count = 0
    animation_slide_count = 0
    for slide_number, slide in enumerate(presentation.slides, start=1):
        shape_evidence: list[dict[str, Any]] = []
        chart_summaries: list[str] = []
        slide_picture_count = 0
        for shape in list(slide.shapes)[:160]:
            shape_type = getattr(shape, "shape_type", None)
            shape_name = _compact(getattr(shape, "name", "") or "未命名对象", 120)
            evidence = {
                "name": shape_name,
                "shape_type": str(shape_type),
                "bbox": {
                    "x": round(float(getattr(shape, "left", 0) or 0) / slide_width, 6),
                    "y": round(float(getattr(shape, "top", 0) or 0) / slide_height, 6),
                    "width": round(float(getattr(shape, "width", 0) or 0) / slide_width, 6),
                    "height": round(float(getattr(shape, "height", 0) or 0) / slide_height, 6),
                },
            }
            if bool(getattr(shape, "has_chart", False)):
                chart_count += 1
                chart = shape.chart
                series_names = [
                    _compact(getattr(series, "name", "") or "未命名系列", 80)
                    for series in list(chart.series)[:20]
                ]
                chart_title = ""
                if bool(getattr(chart, "has_title", False)):
                    chart_title = _compact(chart.chart_title.text_frame.text, 140)
                evidence["chart"] = {
                    "title": chart_title,
                    "series_names": series_names,
                    "series_count": len(chart.series),
                }
                chart_summaries.append(
                    f"图表“{chart_title or shape_name}”含 {len(chart.series)} 个系列"
                    + (f"（{'、'.join(series_names)}）" if series_names else "")
                )
            if shape_type == MSO_SHAPE_TYPE.PICTURE:
                picture_count += 1
                slide_picture_count += 1
                alt_text = _compact(
                    getattr(shape, "alternative_text", "")
                    or getattr(shape, "description", ""),
                    240,
                )
                if alt_text:
                    evidence["alternative_text"] = alt_text
                else:
                    unread_picture_count += 1
            shape_evidence.append(evidence)

        has_animation = bool(slide._element.xpath(".//p:timing"))
        if has_animation:
            animation_slide_count += 1
        notes_text = ""
        try:
            notes_text = _compact(slide.notes_slide.notes_text_frame.text, 1600)
        except (AttributeError, KeyError, ValueError):
            notes_text = ""
        if notes_text:
            note_slide_count += 1
        summary_parts = [
            f"幻灯片 {slide_number} 含 {len(shape_evidence)} 个可定位对象",
        ]
        if chart_summaries:
            summary_parts.extend(chart_summaries)
        if slide_picture_count:
            summary_parts.append(f"含 {slide_picture_count} 张图片，未对无替代文本图片臆测语义")
        if has_animation:
            summary_parts.append("检测到动画时间线，未解释播放顺序和触发语义")
        blocks.append(DocumentBlock(
            block_id=f"blk-{starting_order + len(blocks) + 1}",
            kind="other",
            text="；".join(summary_parts),
            order=starting_order + len(blocks),
            locator=DocumentLocator(page=slide_number, slide=slide_number),
            metadata={
                "evidence_kind": "pptx_visual_structure",
                "shape_count": len(shape_evidence),
                "shapes": shape_evidence,
                "has_animation": has_animation,
            },
        ))
        if notes_text:
            blocks.append(DocumentBlock(
                block_id=f"blk-{starting_order + len(blocks) + 1}",
                kind="paragraph",
                text=notes_text,
                order=starting_order + len(blocks),
                locator=DocumentLocator(page=slide_number, slide=slide_number),
                metadata={"evidence_kind": "speaker_notes"},
            ))

    warnings = ["PPTX 已补充对象布局、图表系列、讲者备注与动画存在性证据"]
    if unread_picture_count:
        warnings.append(
            f"{unread_picture_count} 张图片没有替代文本，未将图片内容冒充为可引用事实"
        )
    if animation_slide_count:
        warnings.append(
            f"{animation_slide_count} 页检测到动画；当前仅保留存在性，不解释播放顺序、触发和转场语义"
        )
    return blocks, {
        "visual_evidence_slide_count": len(presentation.slides),
        "chart_count": chart_count,
        "picture_count": picture_count,
        "unread_picture_count": unread_picture_count,
        "speaker_note_slide_count": note_slide_count,
        "animation_slide_count": animation_slide_count,
    }, warnings


def _docling_kind(label: str, text: str) -> str:
    mapping = {
        "title": "title",
        "section_header": "heading",
        "list_item": "list_item",
        "table": "table",
        "formula": "formula",
        "code": "code",
        "picture": "picture",
    }
    return mapping.get(label, _detect_block_kind(text))


def _detect_block_kind(text: str) -> str:
    stripped = text.strip()
    if re.search(r"(^|\n)(题目|问题|练习|思考)[：:]", stripped) or stripped.endswith(("?", "？")):
        return "question"
    if "```" in stripped:
        return "code"
    if re.search(r"\$[^$]+\$|\\\([^)]*\\\)|\\\[[^]]*\\\]", stripped):
        return "formula"
    if stripped.startswith(("- ", "* ", "1. ")):
        return "list_item"
    return "paragraph"


def _normalized_bbox(raw: Any) -> dict[str, float] | None:
    if not isinstance(raw, dict):
        return None
    result: dict[str, float] = {}
    for key in ("l", "t", "r", "b", "x", "y", "width", "height"):
        if isinstance(raw.get(key), (int, float)):
            result[key] = float(raw[key])
    return result or None


def _ocr_image(
    path: Path,
    *,
    engine: Any | None = None,
    page_number: int = 1,
) -> list[dict[str, Any]]:
    """Run optional local OCR without sending course material to a third party."""
    try:
        from PIL import Image
    except ImportError as exc:
        raise RuntimeError(
            "图片 OCR 图像组件未安装"
        ) from exc
    RapidOCR = _rapidocr_class()

    ocr_engine = engine or RapidOCR()
    output = ocr_engine(str(path))
    if hasattr(output, "boxes") and hasattr(output, "txts"):
        boxes = getattr(output, "boxes", None)
        texts = getattr(output, "txts", None)
        scores = getattr(output, "scores", None)
        raw_result = [
            [points, text, confidence]
            for points, text, confidence in zip(
                boxes if boxes is not None else [],
                texts if texts is not None else [],
                scores if scores is not None else [],
                strict=False,
            )
        ]
    else:
        raw_result = output[0] if isinstance(output, tuple) and output else output
    if not raw_result:
        return []
    with Image.open(path) as image:
        width, height = image.size
    result: list[dict[str, Any]] = []
    for raw in raw_result:
        if not isinstance(raw, (list, tuple)) or len(raw) < 3:
            continue
        points, text, confidence = raw[0], raw[1], raw[2]
        point_values = points if points is not None else []
        xs = [float(point[0]) for point in point_values if len(point) >= 2]
        ys = [float(point[1]) for point in point_values if len(point) >= 2]
        bbox = None
        if xs and ys and width > 0 and height > 0:
            left, right = min(xs), max(xs)
            top, bottom = min(ys), max(ys)
            bbox = {
                "x": round(left / width, 6),
                "y": round(top / height, 6),
                "width": round((right - left) / width, 6),
                "height": round((bottom - top) / height, 6),
            }
        result.append({
            "text": str(text or ""),
            "confidence": float(confidence or 0),
            "bbox": bbox,
            "page": page_number,
        })
    return result


def _rapidocr_class() -> Any:
    """Load either the current or legacy RapidOCR package behind one adapter."""
    try:
        from rapidocr import RapidOCR

        return RapidOCR
    except ImportError:
        try:
            from rapidocr_onnxruntime import RapidOCR

            return RapidOCR
        except ImportError as exc:
            raise RuntimeError(
                "OCR 组件未安装；Python 3.13+ 请安装 rapidocr，旧环境请安装 rapidocr-onnxruntime"
            ) from exc


__all__ = [
    "DoclingDocumentParser",
    "DocumentParser",
    "LegacyOfficeConversionParser",
    "MarkItDownFallbackParser",
    "ImageOcrParser",
    "ScannedPdfOcrParser",
    "TextDocumentParser",
    "parse_document_path",
    "parse_material_asset",
]
