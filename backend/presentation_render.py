"""Render one validated deck schema to self-contained HTML and editable PPTX."""

from __future__ import annotations

import hashlib
import html
import json
from html.parser import HTMLParser
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from presentation_models import DeckRevision, Slide
from presentation_templates import (
    LAYOUT_REGISTRY_VERSION,
    TEMPLATE_VERSION,
    get_layout,
    get_template,
    validate_slide_capacity,
)


SLIDE_WIDTH = Inches(13.333333)
SLIDE_HEIGHT = Inches(7.5)


def sha256_file(path: str | Path) -> str:
    digest = hashlib.sha256(Path(path).read_bytes()).hexdigest()
    return f"sha256:{digest}"


def _text_digest(values: Iterable[str]) -> str:
    payload = json.dumps(list(values), ensure_ascii=False, separators=(",", ":"))
    return f"sha256:{hashlib.sha256(payload.encode('utf-8')).hexdigest()}"


def title_digest(slides: Iterable[Slide]) -> str:
    return _text_digest(slide.title for slide in slides)


def speaker_notes_digest(slides: Iterable[Slide]) -> str:
    return _text_digest(slide.speaker_notes for slide in slides)


class _RenderedHtmlParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.page_count = 0
        self.titles: list[str] = []
        self._in_slide = False
        self._in_title = False
        self._title_parts: list[str] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        attributes = dict(attrs)
        if tag == "section" and "slide" in str(attributes.get("class") or "").split():
            self.page_count += 1
            self._in_slide = True
        elif tag == "h1" and self._in_slide:
            self._in_title = True
            self._title_parts = []

    def handle_data(self, data: str) -> None:
        if self._in_title:
            self._title_parts.append(data)

    def handle_endtag(self, tag: str) -> None:
        if tag == "h1" and self._in_title:
            self.titles.append("".join(self._title_parts))
            self._in_title = False
        elif tag == "section" and self._in_slide:
            self._in_slide = False


def validate_renderable(revision: DeckRevision) -> None:
    for slide in revision.slides:
        get_layout(slide.layout_id)
        text = "\n".join(
            [slide.title, slide.subtitle, slide.key_message]
            + [block.title + block.content + "".join(block.items) for block in slide.blocks]
        )
        errors = validate_slide_capacity(
            slide.layout_id,
            [block.model_dump() for block in slide.blocks],
            text,
        )
        if errors:
            raise ValueError("; ".join(errors))


def _block_html(slide: Slide) -> str:
    cards: list[str] = []
    for block in slide.blocks:
        heading = f"<h3>{html.escape(block.title)}</h3>" if block.title else ""
        if block.type in {"bullets", "exercise"} and block.items:
            body = "<ul>" + "".join(f"<li>{html.escape(item)}</li>" for item in block.items) + "</ul>"
        elif block.type == "code":
            body = f"<pre><code>{html.escape(block.content)}</code></pre>"
        elif block.type == "comparison":
            left = html.escape(str(block.metadata.get("left", block.content)))
            right = html.escape(str(block.metadata.get("right", "")))
            body = f'<div class="compare"><div>{left}</div><div>{right}</div></div>'
        else:
            body = f"<p>{html.escape(block.content)}</p>"
        cards.append(f'<article class="block block-{block.type}">{heading}{body}</article>')
    return "".join(cards)


def render_html(
    revision: DeckRevision,
    template_id: str,
    output_path: str | Path,
) -> Path:
    validate_renderable(revision)
    template = get_template(template_id)
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    slide_markup = []
    for index, slide in enumerate(revision.slides):
        slide_markup.append(
            f'''<section class="slide layout-{slide.layout_id}" data-slide-id="{html.escape(slide.slide_id)}" data-index="{index}" aria-label="第 {index + 1} 页">
              <div class="accent"></div>
              <header><span class="eyebrow">{html.escape(slide.layout_id)}</span><h1>{html.escape(slide.title)}</h1><p>{html.escape(slide.subtitle or slide.key_message)}</p></header>
              <div class="blocks">{_block_html(slide)}</div>
              <footer><span>{index + 1} / {len(revision.slides)}</span><strong>灵知课件</strong></footer>
            </section>'''
        )
    metadata = json.dumps(
        {"deckId": revision.deck_id, "revisionId": revision.revision_id, "slideIds": revision.slide_order},
        ensure_ascii=False,
    ).replace("<", "\\u003c")
    document = f'''<!doctype html>
<html lang="zh-CN"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>{html.escape(revision.slides[0].title if revision.slides else "灵知课件")}</title>
<style>
:root{{--surface:{template.surface};--ink:{template.ink};--muted:{template.muted};--accent:{template.accent};--soft:{template.soft_accent};}}
*{{box-sizing:border-box}}html,body{{margin:0;width:100%;height:100%;font-family:{template.font_family};background:#f6f7fb;color:var(--ink)}}
body{{display:grid;place-items:center;overflow:hidden}}.deck{{width:100%;height:100%;display:grid;place-items:center;padding:2.5vh 3vw}}
.slide{{display:none;position:relative;aspect-ratio:16/9;width:min(100%,calc(94vh * 16 / 9));max-height:94vh;background:var(--surface);border:1px solid #e4e7ec;box-shadow:0 18px 50px rgba(27,36,64,.10);padding:5.4% 5.2% 4.2%;overflow:hidden}}
.slide.active{{display:flex;flex-direction:column}}.accent{{position:absolute;top:0;left:0;right:0;height:7px;background:var(--accent)}}
header h1{{font-size:clamp(28px,3.2vw,54px);line-height:1.14;margin:.2rem 0 1rem;letter-spacing:-.02em}}header p{{font-size:clamp(15px,1.4vw,24px);color:var(--muted);margin:0;line-height:1.55}}
.eyebrow{{font-size:12px;color:var(--accent);font-weight:700;letter-spacing:.12em}}.blocks{{flex:1;display:grid;grid-template-columns:repeat(2,minmax(0,1fr));gap:2.2%;align-content:center;margin-top:3%}}
.block{{min-width:0;padding:3.5%;border:1px solid #e8eaf0;border-radius:12px;background:#fbfcff}}.block h3{{font-size:clamp(15px,1.4vw,23px);margin:0 0 .7em}}.block p,.block li{{font-size:clamp(13px,1.12vw,19px);line-height:1.6;color:#465166}}.block ul{{margin:0;padding-left:1.2em}}
.block-code{{background:#f8f9fc}}pre{{white-space:pre-wrap;margin:0;font-family:{template.code_font_family};font-size:clamp(12px,1.05vw,18px);line-height:1.55;color:#26324a}}.compare{{display:grid;grid-template-columns:1fr 1fr;gap:1px;background:#e3e6ee}}.compare>div{{background:white;padding:1rem;line-height:1.55}}
.layout-L01,.layout-L03{{justify-content:center}}.layout-L01 .blocks,.layout-L03 .blocks{{display:none}}.layout-L07 .blocks,.layout-L08 .blocks{{grid-template-columns:1.25fr .75fr}}footer{{display:flex;justify-content:space-between;color:#98a2b3;font-size:12px;margin-top:2%}}footer strong{{color:var(--accent)}}
.controls{{position:fixed;left:50%;bottom:12px;transform:translateX(-50%);display:flex;gap:8px;padding:6px;background:rgba(255,255,255,.92);border:1px solid #e5e7eb;border-radius:10px;box-shadow:0 8px 24px rgba(15,23,42,.12)}}button{{border:0;background:transparent;padding:8px 12px;border-radius:7px;color:#475467;cursor:pointer}}button:hover{{background:var(--soft);color:var(--accent)}}
</style></head><body><main class="deck">{''.join(slide_markup)}</main><nav class="controls" aria-label="翻页"><button id="prev">上一页</button><button id="next">下一页</button></nav>
<script>const meta={metadata};let current=0;const slides=[...document.querySelectorAll('.slide')];function emit(type,extra={{}}){{parent.postMessage({{version:'presentation-preview/v1',type,deck_id:meta.deckId,revision_id:meta.revisionId,...extra,payload:extra.payload||{{}}}},location.origin)}}function show(i){{current=Math.max(0,Math.min(slides.length-1,i));slides.forEach((s,n)=>s.classList.toggle('active',n===current));const slideId=slides[current]?.dataset.slideId||null;emit('slide:selected',{{slide_id:slideId}})}}document.getElementById('prev').onclick=()=>show(current-1);document.getElementById('next').onclick=()=>show(current+1);addEventListener('keydown',e=>{{if(e.key==='ArrowLeft')show(current-1);if(e.key==='ArrowRight')show(current+1)}});show(0);emit('preview:ready');document.fonts.ready.then(()=>emit('render:measured',{{payload:{{overflow:slides.some(s=>s.scrollHeight>s.clientHeight||s.scrollWidth>s.clientWidth),slide_count:slides.length}}}}));</script></body></html>'''
    output.write_text(document, encoding="utf-8")
    return output


def _rgb(hex_color: str) -> RGBColor:
    value = hex_color.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _add_textbox(slide, left, top, width, height, text: str, *, size: int, color: RGBColor, bold: bool = False, font: str = "Microsoft YaHei"):
    shape = slide.shapes.add_textbox(left, top, width, height)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def render_pptx(revision: DeckRevision, template_id: str, output_path: str | Path) -> Path:
    validate_renderable(revision)
    template = get_template(template_id)
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[6]
    for index, item in enumerate(revision.slides):
        slide = prs.slides.add_slide(blank)
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = _rgb(template.surface)
        accent = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.08))
        accent.fill.solid(); accent.fill.fore_color.rgb = _rgb(template.accent); accent.line.fill.background()
        title_size = 32 if item.layout_id not in {"L01", "L03"} else 38
        title_shape = _add_textbox(slide, Inches(.72), Inches(.58), Inches(11.9), Inches(.72), item.title, size=title_size, color=_rgb(template.ink), bold=True)
        title_shape.name = f"LingzhiTitle:{item.slide_id}"
        subtitle = item.subtitle or item.key_message
        if subtitle:
            _add_textbox(slide, Inches(.72), Inches(1.33), Inches(11.7), Inches(.56), subtitle, size=15, color=_rgb(template.muted))
        if item.layout_id not in {"L01", "L03"}:
            columns = 2 if len(item.blocks) > 1 else 1
            col_width = Inches(5.75 if columns == 2 else 11.7)
            for block_index, block in enumerate(item.blocks[:4]):
                col = block_index % columns
                row = block_index // columns
                left = Inches(.72) + col * Inches(6.0)
                top = Inches(2.08) + row * Inches(2.05)
                body = block.content
                if block.items:
                    body = "\n".join(f"• {value}" for value in block.items)
                if block.title:
                    body = f"{block.title}\n{body}".strip()
                font_name = "Cascadia Code" if block.type == "code" else "Microsoft YaHei"
                _add_textbox(slide, left, top, col_width, Inches(1.72), body, size=15 if block.type != "code" else 13, color=_rgb(template.ink), font=font_name)
        _add_textbox(slide, Inches(.72), Inches(7.02), Inches(2), Inches(.25), f"{index + 1} / {len(revision.slides)}", size=9, color=_rgb(template.muted))
        if item.speaker_notes:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_frame.text = item.speaker_notes
    prs.save(output)
    return output


def render_artifacts(revision: DeckRevision, template_id: str, artifact_dir: str | Path) -> dict[str, str | int]:
    target = Path(artifact_dir)
    html_path = render_html(revision, template_id, target / "deck.html")
    pptx_path = render_pptx(revision, template_id, target / "deck.pptx")
    expected_titles = [slide.title for slide in revision.slides]
    expected_notes = [slide.speaker_notes for slide in revision.slides]

    html_parser = _RenderedHtmlParser()
    html_parser.feed(html_path.read_text(encoding="utf-8"))
    if html_parser.page_count != len(revision.slides):
        raise ValueError("html_page_count_mismatch")
    if html_parser.titles != expected_titles:
        raise ValueError("html_title_mismatch")

    reloaded = Presentation(pptx_path)
    if len(reloaded.slides) != len(revision.slides):
        raise ValueError("pptx_page_count_mismatch")
    reloaded_titles: list[str] = []
    reloaded_notes: list[str] = []
    for expected, rendered in zip(revision.slides, reloaded.slides, strict=True):
        title_shapes = [shape for shape in rendered.shapes if shape.name == f"LingzhiTitle:{expected.slide_id}"]
        if len(title_shapes) != 1 or not hasattr(title_shapes[0], "text"):
            raise ValueError("pptx_title_shape_missing")
        reloaded_titles.append(str(title_shapes[0].text))
        reloaded_notes.append(rendered.notes_slide.notes_text_frame.text or "")
    if reloaded_titles != expected_titles:
        raise ValueError("pptx_title_mismatch")
    if reloaded_notes != expected_notes:
        raise ValueError("pptx_speaker_notes_mismatch")
    return {
        "html_path": str(html_path),
        "html_sha256": sha256_file(html_path),
        "pptx_path": str(pptx_path),
        "pptx_sha256": sha256_file(pptx_path),
        "page_count": len(revision.slides),
        "title_digest": _text_digest(reloaded_titles),
        "speaker_notes_digest": _text_digest(reloaded_notes),
        "template_version": TEMPLATE_VERSION,
        "layout_registry_version": LAYOUT_REGISTRY_VERSION,
    }
