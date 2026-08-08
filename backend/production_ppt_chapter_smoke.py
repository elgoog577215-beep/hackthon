"""Read-only production smoke for one source-backed V5 chapter deck.

The smoke deliberately uses the deployed compiler, provider configuration, and
published course data.  It never persists a representation or modifies course
content; the only writes are diagnostic artifacts under the requested output
directory.
"""

from __future__ import annotations

import argparse
import asyncio
import hashlib
import json
import os
import re
import shutil
import subprocess
import sys
import traceback
from collections import Counter
from copy import deepcopy
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from course_document import (
    CourseDocument,
    course_view_from_document,
    refresh_document_revision,
)

PROGRAMMING_MODES = frozenset({
    "programming_engineering",
    "engineering_programming",
})
CODE_BLOCK_KINDS = frozenset({"code", "code_lab"})
SOURCE_LOOP_ROLES = frozenset({
    "concept",
    "reasoning",
    "example",
    "application",
    "activity",
    "checkpoint",
    "feedback",
    "misconception",
    "summary",
    "transfer",
})


class SmokeFailure(RuntimeError):
    def __init__(
        self,
        code: str,
        message: str,
        *,
        details: dict[str, Any] | None = None,
    ) -> None:
        super().__init__(message)
        self.code = code
        self.details = details or {}


@dataclass(frozen=True)
class ChapterCandidate:
    chapter_id: str
    section_ids: tuple[str, ...]
    source_role_count: int
    code_character_count: int
    source_block_count: int
    subject_artifact_kinds: tuple[str, ...] = ()
    subject_artifact_fragment_count: int = 0
    subject_profile_id: str = ""

    @property
    def rank(self) -> tuple[int, int, int, int]:
        return (
            self.source_role_count,
            min(self.code_character_count, 6000),
            len(self.section_ids),
            -self.source_block_count,
        )

    @property
    def subject_rank(self) -> tuple[int, int, int, int, int]:
        return (
            len(self.subject_artifact_kinds),
            self.subject_artifact_fragment_count,
            self.source_role_count,
            len(self.section_ids),
            -self.source_block_count,
        )


@dataclass
class SelectedProductionChapter:
    course_id: str
    chapter: ChapterCandidate
    document: CourseDocument
    course_view: dict[str, Any]
    source_digest: str
    baseline_story: Any
    is_published: bool = True
    fragments: list[Any] = field(default_factory=list)
    rejected_candidate_codes: list[str] = field(default_factory=list)


def _chapter_root_id(
    section_id: str,
    parent_by_id: dict[str, str | None],
) -> str:
    current = section_id
    seen: set[str] = set()
    while current not in seen:
        seen.add(current)
        parent = parent_by_id.get(current)
        if not parent or parent not in parent_by_id:
            return current
        current = parent
    raise SmokeFailure(
        "production_course_section_cycle",
        "The selected production course contains a cyclic section hierarchy.",
    )


def _code_lines_from_markdown(markdown: str) -> list[str]:
    fenced = re.findall(r"```[^\n]*\n(.*?)```", markdown, flags=re.DOTALL)
    sources = fenced or [markdown]
    return [
        line.rstrip()
        for source in sources
        for line in source.splitlines()
        if line.strip() and not line.strip().startswith("```")
    ]


def extract_source_code_lines(document: CourseDocument) -> list[str]:
    lines: list[str] = []
    for block in sorted(document.blocks, key=lambda item: item.position):
        markdown = str(
            block.payload.get("markdown")
            or block.payload.get("text")
            or block.payload.get("content")
            or ""
        )
        if block.kind not in CODE_BLOCK_KINDS and "```" not in markdown:
            continue
        lines.extend(_code_lines_from_markdown(markdown))
    return lines


def rank_programming_chapter_candidates(
    document: CourseDocument,
    *,
    requested_chapter_id: str = "",
) -> list[ChapterCandidate]:
    parent_by_id = {
        section.section_id: section.parent_section_id
        for section in document.sections
    }
    root_by_section = {
        section_id: _chapter_root_id(section_id, parent_by_id)
        for section_id in parent_by_id
    }
    root_ids = {
        root_by_section[section_id]
        for section_id in parent_by_id
    }
    if requested_chapter_id and requested_chapter_id not in root_ids:
        raise SmokeFailure(
            "production_chapter_not_found",
            "The requested production chapter is not a chapter root.",
        )

    candidates: list[ChapterCandidate] = []
    roots = [requested_chapter_id] if requested_chapter_id else sorted(root_ids)
    for root_id in roots:
        section_ids = tuple(
            section.section_id
            for section in sorted(document.sections, key=lambda item: item.position)
            if root_by_section.get(section.section_id) == root_id
        )
        section_id_set = set(section_ids)
        blocks = [
            block
            for block in document.blocks
            if block.section_id in section_id_set and block.status != "retired"
        ]
        chapter_document = document.model_copy(
            deep=True,
            update={
                "sections": [
                    section
                    for section in document.sections
                    if section.section_id in section_id_set
                ],
                "blocks": blocks,
            },
        )
        code_lines = extract_source_code_lines(chapter_document)
        if not code_lines:
            continue
        roles = {
            str(block.role)
            for block in blocks
            if str(block.role) in SOURCE_LOOP_ROLES
        }
        candidates.append(ChapterCandidate(
            chapter_id=root_id,
            section_ids=section_ids,
            source_role_count=len(roles),
            code_character_count=sum(len(line) for line in code_lines),
            source_block_count=len(blocks),
        ))

    if not candidates:
        if requested_chapter_id:
            raise SmokeFailure(
                "production_chapter_has_no_code_source",
                "The requested production chapter has no source-backed code artifact.",
            )
        raise SmokeFailure(
            "production_programming_chapter_not_found",
            "No production programming chapter with source-backed code was found.",
        )
    return sorted(candidates, key=lambda item: item.rank, reverse=True)


def rank_subject_chapter_candidates(
    document: CourseDocument,
    course_view: dict[str, Any],
    *,
    requested_chapter_id: str = "",
) -> list[ChapterCandidate]:
    """Rank non-programming chapters by source-backed subject artifacts."""

    from slide_deck_v3 import fragment_course_document
    from slide_semantics import (
        compile_ppt_semantic_units,
        compile_subject_presentation_contract_v1,
    )

    parent_by_id = {
        section.section_id: section.parent_section_id
        for section in document.sections
    }
    root_ids = sorted({
        _chapter_root_id(section_id, parent_by_id)
        for section_id in parent_by_id
    })
    if requested_chapter_id and requested_chapter_id not in root_ids:
        raise SmokeFailure(
            "production_chapter_not_found",
            "The requested production chapter is not a chapter root.",
        )

    candidates: list[ChapterCandidate] = []
    roots = [requested_chapter_id] if requested_chapter_id else root_ids
    for root_id in roots:
        chapter_document = build_chapter_document(document, root_id)
        fragments = fragment_course_document(chapter_document)
        semantic_units = compile_ppt_semantic_units(chapter_document, fragments)
        contract = compile_subject_presentation_contract_v1(
            chapter_document,
            course_view,
            semantic_units,
            fragments,
        )
        required_kinds = tuple(sorted(set(
            contract.required_representation_kinds
        )))
        if not required_kinds:
            continue
        characteristic_ids = {
            fragment_id
            for kind in required_kinds
            for fragment_id in contract.characteristic_fragment_ids.get(kind, [])
        }
        blocks = list(chapter_document.blocks)
        roles = {
            str(block.role)
            for block in blocks
            if str(block.role) in SOURCE_LOOP_ROLES
        }
        candidates.append(ChapterCandidate(
            chapter_id=root_id,
            section_ids=tuple(
                section.section_id for section in chapter_document.sections
            ),
            source_role_count=len(roles),
            code_character_count=sum(
                len(line) for line in extract_source_code_lines(chapter_document)
            ),
            source_block_count=len(blocks),
            subject_artifact_kinds=required_kinds,
            subject_artifact_fragment_count=len(characteristic_ids),
            subject_profile_id=str(contract.profile_id or "generic"),
        ))

    if not candidates:
        if requested_chapter_id:
            raise SmokeFailure(
                "production_chapter_has_no_subject_artifact",
                "The requested production chapter has no source-backed subject artifact.",
            )
        raise SmokeFailure(
            "production_subject_chapter_not_found",
            "No production chapter with a source-backed subject artifact was found.",
        )
    return sorted(candidates, key=lambda item: item.subject_rank, reverse=True)


def build_subject_artifact_gate_summary(
    *,
    required_kinds: set[str],
    characteristic_fragment_ids: dict[str, set[str]],
    slide_artifact_kinds: set[str],
    allocated_fragment_ids: set[str],
    excluded_fragment_reasons: dict[str, str],
    editorial_fallback_kinds: set[str],
) -> dict[str, Any]:
    """Build subject-neutral evidence gates without inventing course content."""

    required_source_ids = {
        fragment_id
        for kind in required_kinds
        for fragment_id in characteristic_fragment_ids.get(kind, set())
    }
    missing_slide_kinds = sorted(required_kinds - slide_artifact_kinds)
    missing_disposition_ids = sorted(
        required_source_ids
        - (allocated_fragment_ids | set(excluded_fragment_reasons))
    )
    invalid_exclusions = sorted(
        fragment_id
        for fragment_id, reason in excluded_fragment_reasons.items()
        if reason != "subject_artifact_redundant_after_chapter_coverage"
    )
    editorial_required_kinds = sorted(
        required_kinds & editorial_fallback_kinds
    )
    return {
        "gates": {
            "subject_contract_has_required_artifact": bool(required_kinds),
            "required_subject_artifacts_present": not missing_slide_kinds,
            "subject_source_disposition_complete": not missing_disposition_ids,
            "subject_exclusions_are_explicit": not invalid_exclusions,
            "subject_artifacts_not_editorial_fallback": not editorial_required_kinds,
        },
        "missing_slide_artifact_kinds": missing_slide_kinds,
        "missing_disposition_fragment_ids": missing_disposition_ids,
        "invalid_exclusion_fragment_ids": invalid_exclusions,
        "editorial_fallback_artifact_kinds": editorial_required_kinds,
    }


def choose_cross_domain_sample(
    accepted: list[SelectedProductionChapter],
    *,
    sample_index: int,
) -> SelectedProductionChapter:
    """Choose one strongest sample per primary mode, then select by index."""

    strongest_by_mode: dict[str, SelectedProductionChapter] = {}
    for item in accepted:
        primary_mode = str(
            (item.course_view.get("subject_pedagogy_profile") or {}).get(
                "primary_mode"
            )
            or ""
        )
        subject_mode = str(item.chapter.subject_profile_id or primary_mode)
        if subject_mode == "generic":
            subject_mode = "generic:" + ",".join(
                item.chapter.subject_artifact_kinds
            )
        current = strongest_by_mode.get(subject_mode)
        if current is None or item.chapter.subject_rank > current.chapter.subject_rank:
            strongest_by_mode[subject_mode] = item
    ranked = sorted(
        strongest_by_mode.values(),
        key=lambda item: (
            item.chapter.subject_rank,
            _private_id(item.course_id),
            _private_id(item.chapter.chapter_id),
        ),
        reverse=True,
    )
    if sample_index < 0 or sample_index >= len(ranked):
        raise SmokeFailure(
            "production_cross_domain_sample_unavailable",
            "Production does not have enough distinct non-programming subject samples.",
            details={
                "requested_sample_index": sample_index,
                "available_distinct_primary_modes": len(ranked),
                "available_primary_modes": sorted(strongest_by_mode),
            },
        )
    return ranked[sample_index]


def course_is_eligible_for_sample(
    *,
    is_published: bool,
    sample_profile: str,
) -> bool:
    """Keep programming regression published-only; allow read-only draft diversity."""

    return bool(is_published) or sample_profile == "cross_domain"


def course_is_cross_domain_candidate(primary_mode: str) -> bool:
    """Allow general courses to be classified later from source artifacts."""

    return str(primary_mode or "") not in PROGRAMMING_MODES


def build_chapter_document(
    document: CourseDocument,
    chapter_id: str,
) -> CourseDocument:
    parent_by_id = {
        section.section_id: section.parent_section_id
        for section in document.sections
    }
    section_ids = {
        section_id
        for section_id in parent_by_id
        if _chapter_root_id(section_id, parent_by_id) == chapter_id
    }
    if chapter_id not in section_ids:
        raise SmokeFailure(
            "production_chapter_not_found",
            "The requested production chapter is not present in the course document.",
        )
    chapter_document = CourseDocument.model_validate({
        **document.model_dump(mode="json"),
        "document_revision": "",
        "sections": [
            section.model_dump(mode="json")
            for section in sorted(document.sections, key=lambda item: item.position)
            if section.section_id in section_ids
        ],
        "blocks": [
            block.model_dump(mode="json")
            for block in sorted(
                document.blocks,
                key=lambda item: (item.section_id, item.position),
            )
            if block.section_id in section_ids and block.status != "retired"
        ],
    })
    return refresh_document_revision(chapter_document)


def _stable_digest(value: Any) -> str:
    serialized = json.dumps(
        value,
        ensure_ascii=False,
        sort_keys=True,
        separators=(",", ":"),
        default=str,
    )
    return hashlib.sha256(serialized.encode("utf-8")).hexdigest()


def _private_id(value: str) -> str:
    return hashlib.sha256(value.encode("utf-8")).hexdigest()[:16]


def _planner_failure_reason_code(value: Any) -> str:
    """Classify provider failures without copying credentials or request IDs."""
    text = json.dumps(value, ensure_ascii=False, default=str).lower()
    if not text or text in {"{}", "[]", "null"}:
        return ""
    if "insufficient balance" in text or "insufficient_balance" in text:
        return "ai_provider_balance_exhausted"
    if "429" in text or "rate limit" in text or "too many requests" in text:
        return "ai_provider_rate_limited"
    if "401" in text or "403" in text or "authentication" in text:
        return "ai_provider_authentication_failed"
    if "timeout" in text or "timed out" in text:
        return "ai_provider_timeout"
    if "validation" in text or "invalid_response" in text:
        return "ai_planner_response_invalid"
    return "ai_planner_failed"


def _release_commit(application_root: Path) -> str:
    release_file = application_root / ".release-commit"
    if not release_file.is_file():
        raise SmokeFailure(
            "production_release_identity_missing",
            "The active production release does not expose .release-commit.",
        )
    return release_file.read_text(encoding="utf-8").strip()


def _issue_summary(items: list[Any]) -> list[dict[str, Any]]:
    result: list[dict[str, Any]] = []
    for item in items:
        if not isinstance(item, dict):
            continue
        summary = {
            "severity": str(item.get("severity") or ""),
            "code": str(item.get("code") or "unknown"),
        }
        page_id = str(item.get("page_id") or "")
        if page_id:
            summary["page_id_hash"] = _private_id(page_id)
        if isinstance(item.get("page"), int):
            summary["page"] = int(item["page"])
        dimension = str(item.get("dimension") or "")
        if dimension:
            summary["dimension"] = dimension
        for key in (
            "body_character_count",
            "body_character_budget",
            "visible_item_count",
            "visible_item_budget",
        ):
            if isinstance(item.get(key), (int, float)):
                summary[key] = item[key]
        result.append(summary)
    return result


def _planned_scene_requirements(story: dict[str, Any]) -> set[str]:
    """Return source-bound interior scenes selected by the final V5 planner."""
    return {
        str(episode.get("scene_kind") or "")
        for chapter in story.get("chapters") or []
        if isinstance(chapter, dict)
        for episode in chapter.get("episodes") or []
        if isinstance(episode, dict)
        and str(episode.get("scene_kind") or "")
        not in {"", "chapter_entry", "chapter_recap", "transition"}
        and any(
            isinstance(beat, dict) and bool(beat.get("fragment_ids"))
            for beat in episode.get("beats") or []
        )
    }


def _source_disposition(
    allocation: dict[str, Any],
    source_fragment_ids: set[str],
) -> tuple[set[str], dict[str, str]]:
    allocated = {
        str(fragment_id)
        for page in allocation.get("pages") or []
        if isinstance(page, dict)
        for fragment_id in page.get("fragment_ids") or []
        if str(fragment_id) in source_fragment_ids
    }
    excluded = {
        str(item.get("fragment_id") or ""): str(item.get("reason") or "")
        for item in allocation.get("exclusions") or []
        if isinstance(item, dict)
        and str(item.get("fragment_id") or "") in source_fragment_ids
    }
    return allocated, excluded


def _visible_code_text(slides: list[dict[str, Any]]) -> str:
    values: list[str] = []
    for slide in slides:
        artifact_kinds = {
            str(item)
            for item in (slide.get("quality") or {}).get(
                "subject_artifact_kinds",
                [],
            )
        }
        for block in slide.get("blocks") or []:
            if str(block.get("type") or "") != "code" and "code" not in artifact_kinds:
                continue
            values.extend([
                str(block.get("content") or ""),
                *[str(item) for item in block.get("items") or []],
            ])
    return "\n".join(values)


def _pptx_text(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(path)
    return "\n".join(
        str(shape.text or "")
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
    )


def _pptx_presentation_mode_audit(
    path: Path,
    slides: list[dict[str, Any]],
) -> dict[str, Any]:
    """Verify that semantic single-region layouts stay single-region in PPTX."""
    from pptx import Presentation

    presentation = Presentation(path)
    issues: list[dict[str, Any]] = []
    for page, (model, rendered) in enumerate(
        zip(slides, presentation.slides),
        start=1,
    ):
        quality = dict(model.get("quality") or {})
        resolved_layout = str(quality.get("resolved_layout") or "")
        text_shapes = [
            shape
            for shape in rendered.shapes
            if getattr(shape, "has_text_frame", False)
            and str(shape.text or "").strip()
        ]
        if resolved_layout == "hero-claim":
            claim = next(
                (
                    str(value).strip()
                    for block in model.get("blocks") or []
                    for value in (block.get("items") or [block.get("content")])
                    if str(value or "").strip()
                ),
                str(
                    model.get("key_message")
                    or model.get("takeaway")
                    or model.get("title")
                    or ""
                ).strip(),
            )
            anchor = re.sub(r"\s+", "", claim)[:12]
            dominant = any(
                anchor
                and anchor in re.sub(r"\s+", "", str(shape.text or ""))
                and int(shape.height) / 914400 >= 1.5
                for shape in text_shapes
            )
            if not dominant:
                issues.append({
                    "severity": "critical",
                    "code": "exported_hero_claim_not_dominant",
                    "page": page,
                })
        if resolved_layout != "code":
            continue
        mode = str(quality.get("code_region_mode") or "")
        code = next(
            (
                str(block.get("content") or "")
                for block in model.get("blocks") or []
                if str(block.get("type") or "") == "code"
            ),
            "",
        )
        anchor = next(
            (line.strip() for line in code.splitlines() if line.strip()),
            "",
        )
        code_shapes = [
            shape
            for shape in text_shapes
            if anchor and anchor in str(shape.text or "")
        ]
        visible_text = "\n".join(str(shape.text or "") for shape in text_shapes)
        if mode == "full_width" and (
            not code_shapes
            or max(int(shape.width) / 914400 for shape in code_shapes) < 10.5
            or "阅读线索" in visible_text
        ):
            issues.append({
                "severity": "critical",
                "code": "exported_code_single_region_not_full_width",
                "page": page,
            })
        if mode == "annotated_split" and "阅读线索" not in visible_text:
            issues.append({
                "severity": "critical",
                "code": "exported_code_annotation_region_missing",
                "page": page,
            })
    return {
        "passed": not issues,
        "issues": issues,
        "hero_claim_page_count": sum(
            str((slide.get("quality") or {}).get("resolved_layout") or "")
            == "hero-claim"
            for slide in slides
        ),
        "full_width_code_page_count": sum(
            str((slide.get("quality") or {}).get("code_region_mode") or "")
            == "full_width"
            for slide in slides
        ),
    }


def _render_artifacts(pptx_path: Path, output_dir: Path) -> dict[str, Any]:
    from PIL import Image, ImageDraw

    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    pdftoppm = shutil.which("pdftoppm")
    missing = [
        name
        for name, executable in (
            ("libreoffice", soffice),
            ("pdftoppm", pdftoppm),
        )
        if not executable
    ]
    if missing:
        raise SmokeFailure(
            "production_render_tools_missing",
            "Production render QA tools are unavailable.",
            details={"missing": missing},
        )
    render_dir = output_dir / "rendered"
    render_dir.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        [
            str(soffice),
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            str(pptx_path),
        ],
        check=True,
        capture_output=True,
        timeout=120,
    )
    pdf_path = output_dir / f"{pptx_path.stem}.pdf"
    if not pdf_path.is_file():
        raise SmokeFailure(
            "production_pdf_render_missing",
            "LibreOffice did not produce the expected PDF artifact.",
        )
    subprocess.run(
        [
            str(pdftoppm),
            "-png",
            "-r",
            "120",
            str(pdf_path),
            str(render_dir / "page"),
        ],
        check=True,
        capture_output=True,
        timeout=120,
    )
    images = sorted(
        render_dir.glob("page-*.png"),
        key=lambda item: int(item.stem.rsplit("-", 1)[-1]),
    )
    if not images:
        raise SmokeFailure(
            "production_slide_images_missing",
            "The deployed export could not be rasterized into slide images.",
        )

    columns = 3
    thumb_width = 480
    margin = 24
    thumbs: list[Image.Image] = []
    for image_path in images:
        with Image.open(image_path) as source:
            thumb = source.convert("RGB")
            thumb.thumbnail((thumb_width, round(thumb_width * 9 / 16)))
            thumbs.append(thumb.copy())
    cell_height = max(image.height for image in thumbs) + 42
    rows = (len(thumbs) + columns - 1) // columns
    sheet = Image.new(
        "RGB",
        (
            columns * thumb_width + (columns + 1) * margin,
            rows * cell_height + (rows + 1) * margin,
        ),
        "white",
    )
    draw = ImageDraw.Draw(sheet)
    for index, thumb in enumerate(thumbs):
        row, column = divmod(index, columns)
        x = margin + column * (thumb_width + margin)
        y = margin + row * (cell_height + margin)
        sheet.paste(thumb, (x, y))
        draw.text((x, y + thumb.height + 8), str(index + 1), fill="black")
    contact_sheet = output_dir / "contact-sheet.png"
    sheet.save(contact_sheet)
    return {
        "pdf_bytes": pdf_path.stat().st_size,
        "rendered_page_count": len(images),
        "contact_sheet_bytes": contact_sheet.stat().st_size,
    }


def _select_production_chapter(
    *,
    requested_course_id: str,
    requested_chapter_id: str,
    sample_profile: str = "programming",
    sample_index: int = 0,
) -> SelectedProductionChapter:
    from course_repository import CourseDocumentRepository
    from slide_deck_v3 import fragment_course_document
    from slide_deck_v5 import compact_story_plan_v5
    from slide_semantics import (
        compile_ppt_semantic_units,
        compile_subject_presentation_contract_v1,
    )
    from slide_story_plan import (
        compile_slide_story_plan_v2,
        resolve_slide_deck_schema,
    )
    from storage import storage

    repository = CourseDocumentRepository(storage)
    summaries = {
        str(item.get("course_id") or ""): item
        for item in storage.list_courses()
        if item.get("course_id")
    }
    if requested_course_id:
        course_ids = [requested_course_id]
    else:
        course_ids = sorted(summaries)
    if not course_ids:
        raise SmokeFailure(
            "production_course_catalog_empty",
            "Production has no courses available for a chapter smoke.",
        )

    accepted: list[SelectedProductionChapter] = []
    rejected_codes: list[str] = []
    for course_id in course_ids:
        try:
            summary = summaries.get(course_id) or {}
            is_published = bool(summary.get("is_published"))
            if not summary or not course_is_eligible_for_sample(
                is_published=is_published,
                sample_profile=sample_profile,
            ):
                raise SmokeFailure(
                    "production_course_not_published",
                    "The production course is not published.",
                )
            raw_before = repository.load_raw(course_id)
            document, canonical = repository.load_document(course_id)
            if not canonical:
                raise SmokeFailure(
                    "production_course_not_canonical",
                    "The production course is not on the canonical course-document chain.",
                )
            course_view = repository.load_course_view(course_id)
            primary_mode = str(
                (course_view.get("subject_pedagogy_profile") or {}).get(
                    "primary_mode",
                )
                or ""
            )
            if sample_profile == "programming":
                if primary_mode not in PROGRAMMING_MODES:
                    raise SmokeFailure(
                        "production_course_not_programming",
                        "The production course is not classified as programming engineering.",
                    )
            elif sample_profile == "cross_domain":
                if not course_is_cross_domain_candidate(primary_mode):
                    raise SmokeFailure(
                        "production_course_not_cross_domain_candidate",
                        "The production course is not a classified non-programming sample.",
                    )
            else:
                raise SmokeFailure(
                    "production_sample_profile_invalid",
                    "The production sample profile is not supported.",
                    details={"sample_profile": sample_profile},
                )
            if resolve_slide_deck_schema(
                course_view,
                story_engine_enabled=True,
                v5_enabled=True,
            ) != "slide_deck_v5":
                raise SmokeFailure(
                    "production_course_not_v5_eligible",
                    "The production course did not resolve to the V5 chain.",
                )
            candidates = (
                rank_programming_chapter_candidates(
                    document,
                    requested_chapter_id=requested_chapter_id,
                )
                if sample_profile == "programming"
                else rank_subject_chapter_candidates(
                    document,
                    course_view,
                    requested_chapter_id=requested_chapter_id,
                )
            )
            if sample_profile == "cross_domain":
                candidates = [
                    candidate
                    for candidate in candidates
                    if candidate.subject_profile_id != "engineering_programming"
                    and "code" not in candidate.subject_artifact_kinds
                ]
                if not candidates:
                    raise SmokeFailure(
                        "production_course_only_has_programming_artifacts",
                        "The production course only exposes programming subject artifacts.",
                    )
            for candidate in candidates:
                chapter_document = build_chapter_document(
                    document,
                    candidate.chapter_id,
                )
                chapter_view = course_view_from_document(
                    course_view,
                    chapter_document,
                )
                if resolve_slide_deck_schema(
                    chapter_view,
                    story_engine_enabled=True,
                    v5_enabled=True,
                ) != "slide_deck_v5":
                    raise SmokeFailure(
                        "production_chapter_not_v5_eligible",
                        "The selected production chapter did not resolve to V5.",
                    )
                fragments = fragment_course_document(chapter_document)
                semantic_units = compile_ppt_semantic_units(
                    chapter_document,
                    fragments,
                )
                subject_contract = compile_subject_presentation_contract_v1(
                    chapter_document,
                    chapter_view,
                    semantic_units,
                    fragments,
                )
                if (
                    sample_profile == "programming"
                    and "code" not in subject_contract.required_representation_kinds
                ):
                    raise SmokeFailure(
                        "production_programming_contract_missing_code",
                        "The programming chapter did not compile a required code contract.",
                    )
                if (
                    sample_profile == "cross_domain"
                    and not subject_contract.required_representation_kinds
                ):
                    raise SmokeFailure(
                        "production_cross_domain_contract_missing_artifact",
                        "The cross-domain chapter did not compile a subject artifact contract.",
                    )
                baseline = compact_story_plan_v5(
                    chapter_document,
                    compile_slide_story_plan_v2(
                        chapter_document,
                        chapter_view,
                        fragments,
                        mode="teaching",
                        theme="qizhi-classroom",
                    ),
                    fragments,
                )
                accepted.append(SelectedProductionChapter(
                    course_id=course_id,
                    chapter=candidate,
                    document=chapter_document,
                    course_view=chapter_view,
                    source_digest=_stable_digest(raw_before),
                    baseline_story=baseline,
                    is_published=is_published,
                    fragments=fragments,
                ))
        except SmokeFailure as exc:
            rejected_codes.append(exc.code)
            if requested_course_id:
                raise
        except Exception as exc:
            rejected_codes.append(type(exc).__name__)
            if requested_course_id:
                raise SmokeFailure(
                    "production_course_selection_failed",
                    "The requested production course could not compile a V5 chapter baseline.",
                    details={"exception_type": type(exc).__name__},
                ) from exc

    if not accepted:
        raise SmokeFailure(
            (
                "production_programming_sample_unavailable"
                if sample_profile == "programming"
                else "production_cross_domain_sample_unavailable"
            ),
            (
                "No published canonical programming chapter passed V5 prerequisites."
                if sample_profile == "programming"
                else "No published canonical non-programming chapter passed V5 prerequisites."
            ),
            details={
                "rejection_counts": dict(Counter(rejected_codes)),
                "course_count": len(course_ids),
            },
        )
    selected = (
        max(accepted, key=lambda item: item.chapter.rank)
        if sample_profile == "programming"
        else choose_cross_domain_sample(accepted, sample_index=sample_index)
    )
    selected.rejected_candidate_codes = rejected_codes
    return selected


async def run_production_smoke(
    *,
    application_root: Path,
    output_dir: Path,
    expected_release_commit: str,
    requested_course_id: str = "",
    requested_chapter_id: str = "",
    sample_profile: str = "programming",
    sample_index: int = 0,
    defer_render: bool = False,
) -> dict[str, Any]:
    from course_repository import CourseDocumentRepository
    from slide_deck_renderer import audit_exported_pptx, export_structured_slide_deck
    from slide_deck_v5 import (
        allocation_from_story_plan_v5,
        compile_slide_deck_v5,
    )
    from slide_story_plan import plan_slide_story_v2
    from slide_visuals import plan_slide_visuals
    from storage import storage
    from task_manager import (
        _source_first_slide_visual_ai_worker,
        _source_first_story_ai_worker,
    )

    output_dir.mkdir(parents=True, exist_ok=True)
    release_commit = _release_commit(application_root)
    if expected_release_commit and release_commit != expected_release_commit:
        raise SmokeFailure(
            "production_release_commit_mismatch",
            "The active production release is not the workflow commit under test.",
            details={
                "expected": expected_release_commit,
                "actual": release_commit,
            },
        )

    selected = _select_production_chapter(
        requested_course_id=requested_course_id,
        requested_chapter_id=requested_chapter_id,
        sample_profile=sample_profile,
        sample_index=sample_index,
    )
    os.environ["SLIDE_WEB_IMAGE_RETRIEVAL_ENABLED"] = "false"
    os.environ["SLIDE_GENERATED_ILLUSTRATIONS_ENABLED"] = "false"
    os.environ["SLIDE_LIBREOFFICE_AUDIT_ENABLED"] = (
        "false" if defer_render else "true"
    )

    story_worker = _source_first_story_ai_worker()
    visual_worker = _source_first_slide_visual_ai_worker()
    story = await plan_slide_story_v2(
        selected.document,
        selected.course_view,
        selected.fragments,
        mode="teaching",
        theme="qizhi-classroom",
        baseline=selected.baseline_story,
        ai_planner=story_worker,
    )
    allocation, _ = allocation_from_story_plan_v5(
        selected.document,
        selected.fragments,
        story,
    )
    visual_plan = await plan_slide_visuals(
        selected.document,
        allocation,
        selected.fragments,
        ai_planner=visual_worker,
    )
    story_failure_reason_code = _planner_failure_reason_code(
        (story.planning_diagnostics or {}).get("chapter_failures")
        or (story.planning_diagnostics or {}).get("deck_failure")
    )
    visual_failure_reason_code = _planner_failure_reason_code(
        (visual_plan.deck_brief or {}).get("failed_visual_batches")
    )
    planner_failure_reason_codes = list(dict.fromkeys(
        code
        for code in (
            story_failure_reason_code,
            visual_failure_reason_code,
        )
        if code
    ))
    content = compile_slide_deck_v5(
        selected.document,
        selected.course_view,
        story_plan=story,
        allocation_plan=allocation,
        visual_plan=visual_plan,
    )
    slides = list(content.get("slides") or [])
    quality = dict(content.get("quality_report") or {})
    issue_codes = {
        str(item.get("code") or "")
        for item in quality.get("issues") or []
        if isinstance(item, dict)
    }
    code_slides = [
        slide
        for slide in slides
        if "code" in {
            str(item)
            for item in (slide.get("quality") or {}).get(
                "subject_artifact_kinds",
                [],
            )
        }
        or any(
            str(block.get("type") or "") == "code"
            for block in slide.get("blocks") or []
        )
    ]
    hero_claim_slides = [
        slide
        for slide in slides
        if str((slide.get("quality") or {}).get("resolved_layout") or "")
        == "hero-claim"
    ]
    task_activity_page_counts = Counter(
        str((slide.get("quality") or {}).get("task_activity_id") or "")
        for slide in slides
        if str((slide.get("quality") or {}).get("task_activity_id") or "")
    )
    task_activity_phases: dict[str, list[str]] = {}
    for slide in slides:
        slide_quality = slide.get("quality") or {}
        activity_id = str(slide_quality.get("task_activity_id") or "")
        if not activity_id:
            continue
        task_activity_phases.setdefault(activity_id, []).append(
            str(slide_quality.get("task_prompt_phase") or "")
        )
    phase_order = {"overview": 0, "procedure": 1, "verification": 2}
    artifact_editorial_fallbacks = [
        str(slide.get("unit_id") or "")
        for slide in code_slides
        if str((slide.get("quality") or {}).get("resolved_layout") or "")
        == "editorial-body"
    ]
    normalized_titles = [
        " ".join(str(slide.get("title") or "").lower().split())
        for slide in slides
        if str(slide.get("title") or "").strip()
    ]
    duplicate_titles = sorted(
        title
        for title, count in Counter(normalized_titles).items()
        if count > 1
    )
    source_code_lines = extract_source_code_lines(selected.document)
    source_code_anchors = [
        source_code_lines[0] if source_code_lines else "",
        source_code_lines[-1] if source_code_lines else "",
    ]
    visible_code = _visible_code_text(slides)
    final_story = dict(content.get("story_plan") or {})
    final_allocation = dict(content.get("allocation_plan") or {})
    subject_contract = dict(
        (final_story.get("planning_diagnostics") or {}).get(
            "subject_presentation_contract",
        )
        or {}
    )
    required_subject_kinds = set(
        subject_contract.get("required_representation_kinds") or []
    )
    characteristic_fragment_ids = {
        str(kind): {
            str(fragment_id)
            for fragment_id in fragment_ids or []
            if str(fragment_id)
        }
        for kind, fragment_ids in (
            subject_contract.get("characteristic_fragment_ids") or {}
        ).items()
    }
    required_subject_fragment_ids = {
        fragment_id
        for kind in required_subject_kinds
        for fragment_id in characteristic_fragment_ids.get(kind, set())
    }
    allocated_subject_ids, excluded_subject = _source_disposition(
        final_allocation,
        required_subject_fragment_ids,
    )
    slide_artifact_kinds = {
        str(kind)
        for slide in slides
        for kind in (
            (slide.get("quality") or {}).get("subject_artifact_kinds") or []
        )
        if str(kind)
    }
    editorial_fallback_kinds = {
        str(kind)
        for slide in slides
        if str((slide.get("quality") or {}).get("resolved_layout") or "")
        == "editorial-body"
        for kind in (
            (slide.get("quality") or {}).get("subject_artifact_kinds") or []
        )
        if str(kind)
    }
    subject_gate_summary = build_subject_artifact_gate_summary(
        required_kinds=required_subject_kinds,
        characteristic_fragment_ids=characteristic_fragment_ids,
        slide_artifact_kinds=slide_artifact_kinds,
        allocated_fragment_ids=allocated_subject_ids,
        excluded_fragment_reasons=excluded_subject,
        editorial_fallback_kinds=editorial_fallback_kinds,
    )
    subject_evidence_gate_names = set(subject_gate_summary["gates"])
    code_source_fragment_ids = {
        str(fragment_id)
        for fragment_id in (
            subject_contract.get("characteristic_fragment_ids") or {}
        ).get("code")
        or []
        if str(fragment_id)
    }
    allocated_code_ids, excluded_code = _source_disposition(
        final_allocation,
        code_source_fragment_ids,
    )
    explicit_code_exclusions = bool(excluded_code) and all(
        reason == "subject_artifact_redundant_after_chapter_coverage"
        for reason in excluded_code.values()
    )
    code_disposition_complete = code_source_fragment_ids == (
        allocated_code_ids | set(excluded_code)
    )
    required_code_anchors = source_code_anchors[:1]
    if not excluded_code:
        required_code_anchors = source_code_anchors
    scene_kinds = {
        str(slide.get("scene_kind") or "")
        for slide in slides
        if str(slide.get("scene_kind") or "")
    }
    required_scenes = _planned_scene_requirements(final_story)
    gates: dict[str, bool] = {
        "release_commit_matches": (
            not expected_release_commit
            or release_commit == expected_release_commit
        ),
        "v5_schema": content.get("schema_version") == "slide_deck_v5",
        "v5_candidate_status": content.get("candidate_status") in {
            "v5_ready",
            "v5_needs_manual_edit",
        },
        "quality_gate": bool(quality.get("passed")),
        "section_flow_preserved": (
            {"chapter_entry", "chapter_recap"} <= scene_kinds
            and required_scenes <= scene_kinds
        ),
        "titles_unique": not duplicate_titles,
        "no_sparse_page_blocker": "sparse_non_exempt_page" not in issue_codes,
        "hero_claims_use_dominant_canvas": all(
            str((slide.get("quality") or {}).get("hero_claim_display_mode") or "")
            == "dominant_canvas"
            for slide in hero_claim_slides
        ),
        "task_activities_are_bounded": all(
            count <= 4 for count in task_activity_page_counts.values()
        ),
        "task_activity_phases_are_ordered": all(
            bool(phases)
            and all(phase in phase_order for phase in phases)
            and [phase_order[phase] for phase in phases]
            == sorted(phase_order[phase] for phase in phases)
            and phases.count("overview") <= 1
            for phases in task_activity_phases.values()
        ),
        "no_density_overflow": not bool(
            issue_codes
            & {
                "body_density_overflow",
                "visible_item_overflow",
                "slide_title_overflow",
            }
        ),
    }
    if sample_profile == "programming":
        gates.update({
            "subject_contract_requires_code": "code" in required_subject_kinds,
            "code_page_count": 1 <= len(code_slides) <= 3,
            "code_not_editorial_fallback": not artifact_editorial_fallbacks,
            "code_source_excerpt_anchored": all(
                anchor and anchor in visible_code
                for anchor in required_code_anchors
            ),
            "code_source_disposition_complete": code_disposition_complete,
            "code_exclusions_are_explicit": (
                not excluded_code or explicit_code_exclusions
            ),
            "code_regions_adapt_to_source_content": all(
                str((slide.get("quality") or {}).get("code_region_mode") or "")
                in {"full_width", "annotated_split"}
                for slide in code_slides
            ),
        })
    else:
        gates.update(subject_gate_summary["gates"])

    report: dict[str, Any] = {
        "schema_version": "production_ppt_chapter_smoke_v1",
        "status": "running",
        "sample_profile": sample_profile,
        "sample_index": sample_index,
        "subject_evidence_gate_names": sorted(subject_evidence_gate_names),
        "release_commit": release_commit,
        "source": {
            "course_id_hash": _private_id(selected.course_id),
            "chapter_id_hash": _private_id(selected.chapter.chapter_id),
            "primary_mode": str(
                (selected.course_view.get("subject_pedagogy_profile") or {}).get(
                    "primary_mode",
                )
                or ""
            ),
            "is_published": selected.is_published,
            "subject_profile_id": selected.chapter.subject_profile_id,
            "section_count": len(selected.document.sections),
            "block_count": len(selected.document.blocks),
            "source_role_count": selected.chapter.source_role_count,
            "source_code_line_count": len(source_code_lines),
            "source_code_character_count": selected.chapter.code_character_count,
            "required_subject_artifact_kinds": sorted(required_subject_kinds),
            "source_subject_artifact_fragment_count": len(
                required_subject_fragment_ids
            ),
            "read_only_digest": selected.source_digest,
        },
        "chain": {
            "schema": content.get("schema_version"),
            "candidate_status": content.get("candidate_status"),
            "story_planner_available": story_worker is not None,
            "story_planner": story.planner,
            "story_fallback_reason": story.fallback_reason,
            "story_failure_reason_code": story_failure_reason_code,
            "visual_planner_available": visual_worker is not None,
            "visual_planner": str(
                (visual_plan.deck_brief or {}).get("planner") or ""
            ),
            "visual_fallback_reason": str(
                (visual_plan.deck_brief or {}).get("fallback_reason") or ""
            ),
            "visual_failure_reason_code": visual_failure_reason_code,
        },
        "deck": {
            "slide_count": len(slides),
            "code_page_count": len(code_slides),
            "hero_claim_page_count": len(hero_claim_slides),
            "task_activity_count": len(task_activity_page_counts),
            "maximum_task_activity_page_count": max(
                task_activity_page_counts.values(),
                default=0,
            ),
            "task_activity_phase_sequences": sorted(
                task_activity_phases.values()
            ),
            "code_region_mode_counts": dict(Counter(
                str((slide.get("quality") or {}).get("code_region_mode") or "")
                for slide in code_slides
            )),
            "resolved_layout_counts": dict(Counter(
                str((slide.get("quality") or {}).get("resolved_layout") or "")
                for slide in slides
            )),
            "scene_counts": dict(Counter(
                str(slide.get("scene_kind") or "")
                for slide in slides
            )),
            "required_source_scenes": sorted(required_scenes),
            "code_source_fragment_count": len(code_source_fragment_ids),
            "code_allocated_fragment_count": len(allocated_code_ids),
            "code_excluded_fragment_count": len(excluded_code),
            "code_excerpted": bool(excluded_code),
            "code_exclusion_reason_counts": dict(Counter(
                excluded_code.values()
            )),
            "duplicate_title_count": len(duplicate_titles),
            "artifact_editorial_fallback_count": len(
                artifact_editorial_fallbacks
            ),
            "present_subject_artifact_kinds": sorted(slide_artifact_kinds),
            "missing_subject_artifact_kinds": list(
                subject_gate_summary["missing_slide_artifact_kinds"]
            ),
            "subject_allocated_fragment_count": len(allocated_subject_ids),
            "subject_excluded_fragment_count": len(excluded_subject),
            "subject_editorial_fallback_kinds": list(
                subject_gate_summary["editorial_fallback_artifact_kinds"]
            ),
        },
        "quality": {
            "passed": bool(quality.get("passed")),
            "score": int(quality.get("score") or 0),
            "blockers": _issue_summary(list(quality.get("blockers") or [])),
            "warnings": _issue_summary(list(quality.get("warnings") or [])),
        },
        "slide_diagnostics": [
            {
                "page_id_hash": _private_id(str(slide.get("unit_id") or "")),
                "position": index,
                "scene_kind": str(slide.get("scene_kind") or ""),
                "beat_role": str(slide.get("beat_role") or ""),
                "resolved_layout": str(
                    (slide.get("quality") or {}).get("resolved_layout") or ""
                ),
                "artifact_kinds": sorted({
                    str(item)
                    for item in (slide.get("quality") or {}).get(
                        "subject_artifact_kinds",
                        [],
                    )
                }),
                "fragment_count": len(
                    (slide.get("quality") or {}).get("fragment_ids") or []
                ),
                "body_character_count": int(
                    (slide.get("quality") or {}).get(
                        "body_character_count",
                    )
                    or 0
                ),
                "body_character_budget": int(
                    (slide.get("quality") or {}).get(
                        "body_character_budget",
                    )
                    or 0
                ),
                "visible_item_count": int(
                    (slide.get("quality") or {}).get("visible_item_count")
                    or 0
                ),
                "visible_item_budget": int(
                    (slide.get("quality") or {}).get("visible_item_budget")
                    or 0
                ),
                "density_band": str(
                    (slide.get("quality") or {}).get("density_band") or ""
                ),
                "hero_claim_display_mode": str(
                    (slide.get("quality") or {}).get(
                        "hero_claim_display_mode",
                    )
                    or ""
                ),
                "code_region_mode": str(
                    (slide.get("quality") or {}).get("code_region_mode") or ""
                ),
                "task_activity_id_hash": (
                    _private_id(str(
                        (slide.get("quality") or {}).get("task_activity_id")
                        or ""
                    ))
                    if (slide.get("quality") or {}).get("task_activity_id")
                    else ""
                ),
                "task_prompt_mode": str(
                    (slide.get("quality") or {}).get("task_prompt_mode") or ""
                ),
                "task_prompt_phase": str(
                    (slide.get("quality") or {}).get("task_prompt_phase") or ""
                ),
                "issue_codes": sorted({
                    str(item.get("code") or "unknown")
                    for item in quality.get("issues") or []
                    if isinstance(item, dict)
                    and str(item.get("page_id") or "")
                    == str(slide.get("unit_id") or "")
                }),
            }
            for index, slide in enumerate(slides, start=1)
        ],
        "gates": gates,
    }
    failed_before_export = [
        name
        for name, passed in gates.items()
        if not passed
        and (
            sample_profile == "programming"
            or name not in subject_evidence_gate_names
        )
    ]
    if failed_before_export:
        report["status"] = "failed"
        report["failure"] = {
            "code": "production_v5_semantic_gate_failed",
            "failed_gates": failed_before_export,
            **(
                {"external_reason_codes": planner_failure_reason_codes}
                if planner_failure_reason_codes
                else {}
            ),
        }
        return report

    pptx_path = output_dir / "chapter-smoke.pptx"
    export_structured_slide_deck(
        content,
        pptx_path,
        require_quality=True,
        theme=str(content.get("theme") or "qizhi-classroom"),
        course_data=selected.course_view,
    )
    export_audit = audit_exported_pptx(
        pptx_path,
        expected_slide_count=len(slides),
    )
    presentation_mode_audit = _pptx_presentation_mode_audit(pptx_path, slides)
    pptx_text = _pptx_text(pptx_path)
    export_anchors_preserved = all(
        anchor and anchor in pptx_text
        for anchor in required_code_anchors
    )
    repository = CourseDocumentRepository(storage)
    source_after = _stable_digest(repository.load_raw(selected.course_id))
    export_gates = {
        "pptx_created": pptx_path.is_file() and pptx_path.stat().st_size > 0,
        "pptx_audit": bool(export_audit.get("passed")),
        "pptx_presentation_modes": bool(presentation_mode_audit.get("passed")),
        "production_course_unchanged": source_after == selected.source_digest,
    }
    if sample_profile == "programming":
        export_gates["pptx_code_anchors_preserved"] = export_anchors_preserved
    report["gates"].update(export_gates)
    report["export"] = {
        "pptx_bytes": pptx_path.stat().st_size,
        "audit_schema": export_audit.get("schema_version"),
        "audit_reviewer": export_audit.get("reviewer"),
        "passed": bool(export_audit.get("passed")),
        "issues": _issue_summary(list(export_audit.get("issues") or [])),
        "presentation_modes": presentation_mode_audit,
    }
    failed = [
        name
        for name, passed in report["gates"].items()
        if not passed
        and (
            sample_profile == "programming"
            or name not in subject_evidence_gate_names
        )
    ]
    if failed:
        report["status"] = "failed"
        report["failure"] = {
            "code": "production_v5_export_gate_failed",
            "failed_gates": failed,
        }
        return report
    if defer_render:
        report["status"] = "passed_pending_render"
        report["render_verification"] = {
            "status": "pending",
            "executor": "isolated_ci_runner",
        }
        return report

    render_artifacts = _render_artifacts(pptx_path, output_dir)
    report["gates"]["rendered_page_count_matches"] = (
        int(render_artifacts["rendered_page_count"]) == len(slides)
    )
    report["export"].update(render_artifacts)
    if not report["gates"]["rendered_page_count_matches"]:
        report["status"] = "failed"
        report["failure"] = {
            "code": "production_v5_render_gate_failed",
            "failed_gates": ["rendered_page_count_matches"],
        }
        return report
    failed_subject_evidence = [
        name
        for name in sorted(subject_evidence_gate_names)
        if not report["gates"].get(name, False)
    ]
    if sample_profile == "cross_domain" and failed_subject_evidence:
        report["status"] = "failed"
        report["failure"] = {
            "code": "production_v5_subject_evidence_gate_failed",
            "failed_gates": failed_subject_evidence,
        }
        return report
    report["status"] = (
        "passed_with_manual_edit"
        if content.get("candidate_status") == "v5_needs_manual_edit"
        else "passed"
    )
    return report


def finalize_deferred_render(output_dir: Path) -> dict[str, Any]:
    report_path = output_dir / "report.json"
    pptx_path = output_dir / "chapter-smoke.pptx"
    if not report_path.is_file() or not pptx_path.is_file():
        raise SmokeFailure(
            "deferred_render_artifacts_missing",
            "Deferred render QA requires both report.json and chapter-smoke.pptx.",
        )
    report = json.loads(report_path.read_text(encoding="utf-8"))
    if str(report.get("status") or "") != "passed_pending_render":
        raise SmokeFailure(
            "deferred_render_status_invalid",
            "Deferred render QA can finalize only a pending render report.",
            details={"status": str(report.get("status") or "")},
        )
    expected_slide_count = int(
        (report.get("deck") or {}).get("slide_count") or 0
    )
    if expected_slide_count <= 0:
        raise SmokeFailure(
            "deferred_render_slide_count_missing",
            "The pending report does not contain a valid expected slide count.",
        )

    render_artifacts = _render_artifacts(pptx_path, output_dir)
    rendered_count_matches = (
        int(render_artifacts["rendered_page_count"])
        == expected_slide_count
    )
    report.setdefault("gates", {})["rendered_page_count_matches"] = (
        rendered_count_matches
    )
    report.setdefault("export", {}).update(render_artifacts)
    report["render_verification"] = {
        "status": "passed" if rendered_count_matches else "failed",
        "executor": "isolated_ci_runner",
    }
    if not rendered_count_matches:
        report["status"] = "failed"
        report["failure"] = {
            "code": "production_v5_render_gate_failed",
            "failed_gates": ["rendered_page_count_matches"],
        }
        return report
    failed_subject_evidence = [
        name
        for name in report.get("subject_evidence_gate_names") or []
        if not (report.get("gates") or {}).get(str(name), False)
    ]
    if (
        str(report.get("sample_profile") or "") == "cross_domain"
        and failed_subject_evidence
    ):
        report["status"] = "failed"
        report["failure"] = {
            "code": "production_v5_subject_evidence_gate_failed",
            "failed_gates": sorted(str(name) for name in failed_subject_evidence),
        }
        return report
    report.pop("failure", None)
    report["status"] = (
        "passed_with_manual_edit"
        if (report.get("chain") or {}).get("candidate_status")
        == "v5_needs_manual_edit"
        else "passed"
    )
    return report


def _write_report(output_dir: Path, report: dict[str, Any]) -> None:
    output_dir.mkdir(parents=True, exist_ok=True)
    (output_dir / "report.json").write_text(
        json.dumps(report, ensure_ascii=False, indent=2, sort_keys=True),
        encoding="utf-8",
    )


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--output-dir", required=True)
    parser.add_argument("--application-root", default="")
    parser.add_argument("--expected-release-commit", default="")
    parser.add_argument("--course-id", default="")
    parser.add_argument("--chapter-id", default="")
    parser.add_argument(
        "--sample-profile",
        choices=("programming", "cross_domain"),
        default="programming",
    )
    parser.add_argument("--sample-index", type=int, default=0)
    parser.add_argument("--defer-render", action="store_true")
    parser.add_argument("--finalize-deferred-render", action="store_true")
    args = parser.parse_args(argv)

    output_dir = Path(args.output_dir).resolve()
    application_root = (
        Path(args.application_root).resolve()
        if args.application_root
        else Path(__file__).resolve().parent.parent
    )
    try:
        if args.finalize_deferred_render:
            report = finalize_deferred_render(output_dir)
        else:
            report = asyncio.run(run_production_smoke(
                application_root=application_root,
                output_dir=output_dir,
                expected_release_commit=str(args.expected_release_commit or ""),
                requested_course_id=str(args.course_id or ""),
                requested_chapter_id=str(args.chapter_id or ""),
                sample_profile=str(args.sample_profile or "programming"),
                sample_index=int(args.sample_index or 0),
                defer_render=bool(args.defer_render),
            ))
    except SmokeFailure as exc:
        report = {
            "schema_version": "production_ppt_chapter_smoke_v1",
            "status": "failed",
            "failure": {
                "code": exc.code,
                "message": str(exc),
                "details": deepcopy(exc.details),
            },
        }
    except Exception as exc:
        trace = [
            {
                "file": Path(frame.filename).name,
                "line": frame.lineno,
                "function": frame.name,
            }
            for frame in traceback.extract_tb(exc.__traceback__)[-8:]
        ]
        report = {
            "schema_version": "production_ppt_chapter_smoke_v1",
            "status": "failed",
            "failure": {
                "code": "production_ppt_chapter_smoke_unhandled",
                "message": "The production chapter smoke raised an unhandled error.",
                "exception_type": type(exc).__name__,
                "details": {
                    "error": str(exc)[:500],
                    "trace": trace,
                },
            },
        }
    _write_report(output_dir, report)
    print(json.dumps(report, ensure_ascii=False, indent=2, sort_keys=True))
    return 0 if str(report.get("status") or "").startswith("passed") else 1


if __name__ == "__main__":
    sys.exit(main())
