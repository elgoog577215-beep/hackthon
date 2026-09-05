"""Native-source renderer for published personal PPT template packs.

The personal-template contract is bidirectional: the planner consumes the
source-derived slot geometry, and this renderer fills the same immutable PPTX
instead of falling back to a built-in theme on a blank presentation.
"""

from __future__ import annotations

import re
from copy import deepcopy
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.util import Inches, Pt

from markdown_table import parse_markdown_table
from slide_deck_renderer import (
    _configure_font,
    _display_text,
    _format_formula_text,
    _shape,
    _table,
)
from slide_deck_v6 import SlideDeckV6, SlidePageV6
from template_layout_contract import (
    TemplateFrameContractV1,
    TemplateLayoutContractV1,
    TemplateLayoutPackContractV1,
)


_TITLE_FONT = "Songti SC"
_BODY_FONT = "Hiragino Sans GB"
_MATH_FONT = "Times New Roman"
_TITLE = "1F497D"
_INK = "29445F"
_MUTED = "718096"
_ACCENT = "2F84D7"
_ACCENT_SOFT = "EAF3FB"
_PANEL = "F6F9FC"
_WHITE = "FFFFFF"
_LINE = "D6E4F0"


def _normalized_rgb(value: Any, fallback: str) -> str:
    candidate = str(value or "").strip().lstrip("#").upper()
    if len(candidate) == 6 and all(character in "0123456789ABCDEF" for character in candidate):
        return candidate
    return fallback


def _palette(contract: TemplateLayoutPackContractV1) -> dict[str, str]:
    theme = contract.render_theme_overrides
    return {
        "title": _normalized_rgb(theme.get("ink") or theme.get("title"), _TITLE),
        "ink": _normalized_rgb(theme.get("ink"), _INK),
        "muted": _normalized_rgb(theme.get("muted"), _MUTED),
        "accent": _normalized_rgb(theme.get("accent"), _ACCENT),
        "accent_soft": _normalized_rgb(theme.get("accent_soft"), _ACCENT_SOFT),
        "panel": _PANEL,
        "line": _LINE,
        "white": _WHITE,
    }


def _replace_relationship_ids(element: Any, replacements: dict[str, str]) -> None:
    for node in element.iter():
        for name, value in list(node.attrib.items()):
            if value in replacements:
                node.set(name, replacements[value])


def _strip_source_text(element: Any) -> None:
    """Keep source artwork while removing topic-specific template copy."""

    for node in element.iter():
        if str(node.tag).endswith("}t"):
            node.text = ""


def _copy_relationships(source_slide: Any, target_slide: Any) -> dict[str, str]:
    replacements: dict[str, str] = {}
    for relationship in source_slide.part.rels.values():
        if relationship.reltype in {RT.SLIDE_LAYOUT, RT.NOTES_SLIDE}:
            continue
        if relationship.is_external:
            new_rid = target_slide.part.rels.get_or_add_ext_rel(
                relationship.reltype,
                relationship.target_ref,
            )
        else:
            new_rid = target_slide.part.rels.get_or_add(
                relationship.reltype,
                relationship.target_part,
            )
        replacements[relationship.rId] = new_rid
    return replacements


def _is_full_bleed(shape: Any, presentation: Presentation) -> bool:
    try:
        return (
            int(shape.width) >= int(presentation.slide_width * 0.94)
            and int(shape.height) >= int(presentation.slide_height * 0.92)
            and int(shape.left) <= int(presentation.slide_width * 0.04)
            and int(shape.top) <= int(presentation.slide_height * 0.04)
        )
    except (AttributeError, TypeError):
        return False


def _clone_native_background(
    presentation: Presentation,
    source_slide: Any,
    *,
    preserve_artwork: bool,
) -> Any:
    target = presentation.slides.add_slide(source_slide.slide_layout)
    # Finished source decks often place topic-specific vector labels on the
    # master (for example RNN/CNN headings).  The selected slide artwork is
    # cloned explicitly below; inherited master shapes must stay hidden so
    # those labels do not leak into a different course.
    target._element.set("showMasterSp", "0")
    target_tree = target.shapes._spTree
    for shape in list(target.shapes):
        target_tree.remove(shape.element)
    source_bg = getattr(source_slide._element.cSld, "bg", None)
    if preserve_artwork and source_bg is not None:
        target_bg = getattr(target._element.cSld, "bg", None)
        if target_bg is not None:
            target._element.cSld.remove(target_bg)
        target._element.cSld.insert(0, deepcopy(source_bg))
    replacements = _copy_relationships(source_slide, target)
    copied = 0
    if preserve_artwork:
        for shape in source_slide.shapes:
            if not _is_full_bleed(shape, presentation):
                continue
            element = deepcopy(shape.element)
            _replace_relationship_ids(element, replacements)
            _strip_source_text(element)
            target_tree.insert_element_before(element, "p:extLst")
            copied += 1
    if copied == 0:
        background = target.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor.from_string(_WHITE)
        # Reconstruct the template's restrained technical-paper grid without
        # retaining any source-topic labels baked into an image or vector group.
        slide_width = presentation.slide_width / 914400
        slide_height = presentation.slide_height / 914400
        for grid_x in [index * 0.42 for index in range(1, int(slide_width / 0.42))]:
            _shape(target, grid_x, 0, 0.004, slide_height, "F1F5F8", radius=False)
        for grid_y in [index * 0.42 for index in range(1, int(slide_height / 0.42))]:
            _shape(target, 0, grid_y, slide_width, 0.004, "F1F5F8", radius=False)
    return target


def _remove_original_slides(presentation: Presentation, original_count: int) -> None:
    # python-pptx can materialize a fresh proxy for the same slide part, so
    # object identity is not stable here.  Generated slides are appended; the
    # immutable source slides are therefore exactly the first N slide ids.
    for slide_id in list(presentation.slides._sldIdLst)[:original_count]:
        presentation.part.drop_rel(slide_id.rId)
        presentation.slides._sldIdLst.remove(slide_id)


def _frame_inches(
    frame: TemplateFrameContractV1 | None,
    presentation: Presentation,
    fallback: tuple[float, float, float, float],
) -> tuple[float, float, float, float]:
    if frame is None:
        return fallback
    width = presentation.slide_width / 914400
    height = presentation.slide_height / 914400
    return (
        frame.x * width,
        frame.y * height,
        frame.width * width,
        frame.height * height,
    )


def _audience_text(value: str, *, formula: bool = False) -> str:
    source = str(value or "")
    if formula or "\\" in source or "$" in source or "_" in source or "^" in source:
        return _format_formula_text(source)
    return _display_text(source)


def _text_box(
    slide: Any,
    value: str,
    x: float,
    y: float,
    width: float,
    height: float,
    *,
    size: float,
    color: str,
    bold: bool = False,
    align: str = "left",
    font: str = _BODY_FONT,
    formula: bool = False,
    valign: str = "top",
    margin: float = 0.03,
) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(width), Inches(height))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(margin)
    frame.margin_top = frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }[valign]
    prepared = (
        _audience_text(str(value or ""), formula=True)
        if formula
        else str(value or "")
    )
    values = prepared.splitlines() or [""]
    for index, line in enumerate(values):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = line if formula else _audience_text(line)
        _configure_font(
            paragraph.font,
            _MATH_FONT if formula else font,
            _BODY_FONT,
        )
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = RGBColor.from_string(color)
        paragraph.alignment = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
        }[align]
        paragraph.space_after = Pt(0 if formula else max(2, size * 0.20))
        paragraph.line_spacing = 1.12 if formula else 1.22
    return box


def _page_title(
    slide: Any,
    page: SlidePageV6,
    presentation: Presentation,
    layout: TemplateLayoutContractV1,
    palette: dict[str, str],
) -> Any:
    frame = layout.slot_frames.get("title")
    x, y, width, height = _frame_inches(
        frame,
        presentation,
        (0.72, 0.48, 10.8, 0.70),
    )
    if layout.construction_role not in {"cover", "recap"}:
        # The source title group may include its subtitle/underline.  Use its
        # exact top edge but guarantee a real title line for generated CJK.
        y = min(max(y, 0.43), 0.72)
        x = min(max(x, 0.62), 0.90)
        width = max(width, 8.6)
        height = max(height, 0.62)
    box = _text_box(
        slide,
        page.title,
        x,
        y,
        width,
        height,
        size=27 if len(page.title) > 19 else 30,
        color=palette["title"],
        bold=True,
        font=_TITLE_FONT,
        valign="middle",
    )
    box.name = f"v6-native-title [v6-title-max-lines={max(1, page.title_max_lines)}]"
    _shape(slide, x, y + height + 0.05, min(width, 11.25), 0.025, palette["accent"], radius=False)
    return box


def _footer(slide: Any, page_number: int, page_count: int, palette: dict[str, str]) -> None:
    _text_box(
        slide,
        f"{page_number:02d} / {page_count:02d}",
        11.75,
        7.08,
        0.90,
        0.22,
        size=9,
        color=palette["muted"],
        align="right",
    )


def _region(page: SlidePageV6, *, kind: str = "", slot: str = "") -> Any | None:
    return next(
        (
            item
            for item in page.regions
            if (not kind or item.content_kind == kind)
            and (not slot or item.slot_id == slot)
        ),
        None,
    )


def _render_cover(
    slide: Any,
    page: SlidePageV6,
    palette: dict[str, str],
) -> None:
    _text_box(
        slide,
        "大学线性代数",
        0.82,
        1.18,
        3.5,
        0.34,
        size=13,
        color=palette["accent"],
        bold=True,
    )
    cover_title = page.title
    match = re.match(r"^(第[^ ]+章)\s+(.+)$", cover_title)
    if match:
        cover_title = f"{match.group(1)}\n{match.group(2)}"
    title = _text_box(
        slide,
        cover_title,
        0.82,
        2.08,
        5.15,
        2.00,
        size=35,
        color=palette["title"],
        bold=True,
        font=_TITLE_FONT,
        valign="middle",
    )
    title.name = f"v6-native-title [v6-title-max-lines={max(1, page.title_max_lines)}]"
    _shape(slide, 0.84, 4.35, 1.20, 0.055, palette["accent"], radius=False)
    _text_box(
        slide,
        "从向量表达到方程组求解",
        0.82,
        4.62,
        4.40,
        0.46,
        size=17,
        color=palette["muted"],
    )


def _render_agenda(slide: Any, page: SlidePageV6, palette: dict[str, str]) -> None:
    region = _region(page, slot="agenda_items") or _region(page, kind="items")
    entries = list((region.metadata or {}).get("agenda_entries") or []) if region else []
    if not entries and region:
        entries = [
            {"index": index, "title": line, "description": ""}
            for index, line in enumerate(region.content.splitlines(), start=1)
            if line.strip()
        ]
    y = 1.72
    for index, entry in enumerate(entries[:4], start=1):
        height = 1.42 if len(entries) <= 3 else 1.05
        _shape(slide, 0.84, y, 11.60, height, palette["white"], radius=True, line=palette["line"])
        _shape(slide, 0.84, y, 0.13, height, palette["accent"], radius=False)
        _text_box(
            slide,
            f"{int(entry.get('index') or index):02d}",
            1.15,
            y + 0.16,
            0.70,
            0.52,
            size=20,
            color=palette["accent"],
            bold=True,
        )
        _text_box(
            slide,
            str(entry.get("title") or ""),
            2.02,
            y + 0.16,
            4.35,
            0.58,
            size=19,
            color=palette["title"],
            bold=True,
        )
        _text_box(
            slide,
            str(entry.get("description") or ""),
            6.45,
            y + 0.14,
            5.35,
            height - 0.28,
            size=16,
            color=palette["ink"],
            valign="middle",
        )
        y += height + 0.22


def _render_formula_page(
    slide: Any,
    page: SlidePageV6,
    palette: dict[str, str],
) -> None:
    formula = _region(page, kind="formula")
    supports = [item for item in page.regions if item.content_kind not in {"formula", "visual"}]
    formula_height = 2.03 if supports else 4.65
    _shape(slide, 0.84, 1.62, 11.60, formula_height, palette["panel"], radius=True, line=palette["line"])
    if formula:
        rendered = _audience_text(formula.content, formula=True)
        line_count = max(1, len(rendered.splitlines()))
        size = 27 if line_count <= 3 and len(rendered) <= 120 else 20 if line_count <= 5 else 16
        _text_box(
            slide,
            formula.content,
            1.20,
            1.82,
            10.88,
            formula_height - 0.38,
            size=size,
            color=palette["title"],
            align="center",
            font=_MATH_FONT,
            formula=True,
            valign="middle",
        )
    if not supports:
        return
    support_y = 3.94
    support_height = 2.54
    _shape(slide, 0.84, support_y, 11.60, support_height, palette["white"], radius=True, line=palette["line"])
    _shape(slide, 0.84, support_y, 0.10, support_height, palette["accent"], radius=False)
    combined = "\n".join(item.content for item in supports)
    size = 17 if len(combined) <= 180 else 16
    _text_box(
        slide,
        combined,
        1.20,
        support_y + 0.25,
        10.78,
        support_height - 0.48,
        size=size,
        color=palette["ink"],
        valign="middle",
    )


def _render_table_page(slide: Any, page: SlidePageV6, palette: dict[str, str]) -> None:
    table_region = _region(page, kind="table")
    support = [item for item in page.regions if item.content_kind not in {"table", "visual"}]
    if table_region is None:
        _render_body_page(slide, page, palette)
        return
    headers, rows = parse_markdown_table(table_region.content)
    table_height = 3.55 if support else 4.85
    _shape(slide, 0.84, 1.60, 11.60, table_height + 0.18, palette["white"], radius=True, line=palette["line"])
    before = len(slide.shapes)
    _table(slide, headers, rows, 1.02, 1.78, 11.24, table_height, {
        "surface": palette["white"],
        "canvas": palette["accent_soft"],
        "chart_bg": palette["accent_soft"],
        "title": palette["title"],
        "ink": palette["ink"],
        "muted": palette["muted"],
        "accent": palette["accent"],
        "accent_soft": palette["accent_soft"],
    })
    for shape in list(slide.shapes)[before:]:
        if not getattr(shape, "has_table", False):
            continue
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    _configure_font(paragraph.font, _BODY_FONT, _BODY_FONT)
                    for run in paragraph.runs:
                        _configure_font(run.font, _BODY_FONT, _BODY_FONT)
    if support:
        _shape(slide, 0.84, 5.62, 11.60, 1.02, palette["accent_soft"], radius=True)
        _text_box(
            slide,
            "\n".join(item.content for item in support),
            1.12,
            5.78,
            11.05,
            0.70,
            size=16,
            color=palette["ink"],
            valign="middle",
        )


def _render_diagram_page(slide: Any, page: SlidePageV6, palette: dict[str, str]) -> None:
    body = [item for item in page.regions if item.content_kind != "visual"]
    if body:
        _shape(slide, 0.84, 1.60, 11.60, 1.20, palette["accent_soft"], radius=True)
        _text_box(
            slide,
            "\n".join(item.content for item in body),
            1.16,
            1.80,
            10.90,
            0.82,
            size=16,
            color=palette["ink"],
            valign="middle",
        )
    from types import SimpleNamespace
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.oxml.xmlchemy import OxmlElement
    from ppt_page_scene import _graph_positions, relation_anchors
    payload = page.visual_decision.visual_payload or {}
    nodes, edges = list(payload.get("nodes") or []), list(payload.get("edges") or [])
    ids = [str(n.get("node_id") or "") for n in nodes]
    if len(nodes) < 2 or not edges or "" in ids or len(ids) != len(set(ids)):
        raise ValueError("v6_visual_diagram_payload_missing")
    relations = [SimpleNamespace(source_id=str(e.get("source") or ""), target_id=str(e.get("target") or "")) for e in edges]
    if any(r.source_id not in ids or r.target_id not in ids for r in relations):
        raise ValueError("v6_visual_diagram_endpoint_missing")
    positions = _graph_positions(ids, relations, (92, 218, 770, 240))
    shapes = {}
    for node in nodes:
        key = str(node["node_id"])
        x, y, w, h = positions[key]
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Pt(x), Pt(y), Pt(w), Pt(h))
        shape.name = f"diagram-node:{key}"
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(palette["white"])
        shape.line.color.rgb = RGBColor.from_string(palette["line"])
        shape.text = str(node.get("label") or "")
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.size, paragraph.font.name = Pt(16), _BODY_FONT
            paragraph.font.color.rgb = RGBColor.from_string(palette["title"])
        shapes[key] = shape
    for index, (edge, relation) in enumerate(zip(edges, relations, strict=True)):
        x1, y1, x2, y2, start, end = relation_anchors(positions[relation.source_id], positions[relation.target_id])
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Pt(x1), Pt(y1), Pt(x2), Pt(y2))
        connector.name = f"diagram-edge:{index}"
        connector.begin_connect(shapes[relation.source_id], start)
        connector.end_connect(shapes[relation.target_id], end)
        connector.line.color.rgb = RGBColor.from_string(palette["accent"])
        connector.line.width = Pt(1.5)
        if edge.get("relation") != "association":
            arrow = OxmlElement("a:tailEnd")
            arrow.set("type", "triangle")
            connector.line._get_or_add_ln().append(arrow)
        if edge.get("label"):
            _text_box(slide, str(edge["label"]), (x1 + x2) / 144 - .5, (y1 + y2) / 144 - .35,
                      1, .3, size=14, color=palette["ink"], align="center")


def _render_body_page(slide: Any, page: SlidePageV6, palette: dict[str, str]) -> None:
    regions = [item for item in page.regions if item.content_kind != "visual"]
    if not regions:
        return
    if any(item.content_kind == "code" for item in regions):
        code = next(item for item in regions if item.content_kind == "code")
        supporting = [item for item in regions if item is not code]
        _shape(slide, 0.84, 1.60, 7.55, 4.94, "172B40", radius=True)
        _text_box(
            slide,
            code.content,
            1.10,
            1.86,
            7.05,
            4.42,
            size=16,
            color="F3F7FB",
            font="Menlo",
        )
        if supporting:
            _shape(slide, 8.68, 1.60, 3.76, 4.94, palette["white"], radius=True, line=palette["line"])
            _text_box(slide, "\n".join(item.content for item in supporting), 8.98, 1.90, 3.18, 4.34, size=16, color=palette["ink"])
        return
    if any(item.content_kind == "steps" for item in regions):
        steps = next(item for item in regions if item.content_kind == "steps")
        values = [line for line in steps.content.splitlines() if line.strip()]
        y = 1.68
        for index, value in enumerate(values[:7], start=1):
            height = min(0.66, 4.68 / max(1, len(values)))
            _shape(slide, 0.96, y, 11.30, height, palette["white"], radius=True, line=palette["line"])
            _text_box(slide, f"{index:02d}", 1.16, y + 0.10, 0.48, height - 0.18, size=12, color=palette["accent"], bold=True, valign="middle")
            _text_box(slide, value, 1.82, y + 0.08, 10.05, height - 0.14, size=16, color=palette["ink"], valign="middle")
            y += height + 0.12
        for extra in [item for item in regions if item is not steps]:
            _text_box(slide, extra.content, 1.05, 6.06, 11.15, 0.52, size=16, color=palette["muted"], valign="middle")
        return
    content = "\n".join(item.content for item in regions)
    _shape(slide, 0.84, 1.66, 11.60, 4.88, palette["white"], radius=True, line=palette["line"])
    _shape(slide, 0.84, 1.66, 0.12, 4.88, palette["accent"], radius=False)
    size = 21 if len(content) <= 90 else 18 if len(content) <= 180 else 16
    _text_box(
        slide,
        content,
        1.28,
        2.06,
        10.56,
        4.04,
        size=size,
        color=palette["ink"],
        valign="middle",
    )


def _render_recap(slide: Any, page: SlidePageV6, palette: dict[str, str]) -> None:
    _text_box(
        slide,
        "知识结构",
        5.44,
        0.43,
        2.45,
        0.55,
        size=16,
        color=palette["accent"],
        bold=True,
        align="center",
        valign="middle",
    )
    title = _text_box(
        slide,
        page.title,
        2.75,
        2.22,
        7.80,
        0.86,
        size=30 if len(page.title) <= 20 else 25,
        color=palette["title"],
        bold=True,
        align="center",
        font=_TITLE_FONT,
        valign="middle",
    )
    title.name = f"v6-native-title [v6-title-max-lines={max(1, page.title_max_lines)}]"
    takeaways = _region(page, kind="items")
    values = [line for line in (takeaways.content if takeaways else "").splitlines() if line.strip()]
    y = 3.42
    card_width = 3.62
    for index, value in enumerate(values[:3]):
        x = 0.92 + index * 4.02
        _shape(slide, x, y, card_width, 1.50, palette["white"], radius=True, line=palette["line"])
        _text_box(slide, f"{index + 1:02d}", x + 0.20, y + 0.12, 0.48, 0.42, size=16, color=palette["accent"], bold=True)
        _text_box(slide, value, x + 0.20, y + 0.60, card_width - 0.40, 0.72, size=16, color=palette["title"], bold=True, align="center", valign="middle")


def _render_native_page(
    slide: Any,
    page: SlidePageV6,
    presentation: Presentation,
    layout: TemplateLayoutContractV1,
    page_number: int,
    page_count: int,
    palette: dict[str, str],
) -> None:
    role = layout.construction_role
    slug = layout.layout_slug
    if role == "cover":
        _render_cover(slide, page, palette)
    elif role == "recap":
        _render_recap(slide, page, palette)
    else:
        _page_title(slide, page, presentation, layout, palette)
        if slug == "agenda-path":
            _render_agenda(slide, page, palette)
        elif _region(page, kind="table") is not None:
            _render_table_page(slide, page, palette)
        elif _region(page, kind="formula") is not None:
            _render_formula_page(slide, page, palette)
        elif page.visual_decision.decision == "diagram":
            _render_diagram_page(slide, page, palette)
        else:
            _render_body_page(slide, page, palette)
    _footer(slide, page_number, page_count, palette)


def render_personal_template_presentation(
    deck: SlideDeckV6,
    source_path: str | Path,
    contract: TemplateLayoutPackContractV1,
) -> Presentation:
    presentation = Presentation(str(source_path))
    originals = list(presentation.slides)
    if not originals:
        raise ValueError("personal_template_source_has_no_slides")
    palette = _palette(contract)
    page_count = len(deck.pages)
    for page_number, page in enumerate(deck.pages, start=1):
        layout = contract.get_layout(page.resolved_layout)
        if layout is None:
            raise ValueError(f"personal_template_layout_missing:{page.resolved_layout}")
        source_number = int(layout.source_slide_number or 0)
        if not 1 <= source_number <= len(originals):
            raise ValueError(
                f"personal_template_source_slide_missing:{layout.template_layout_id}:{source_number}"
            )
        slide = _clone_native_background(
            presentation,
            originals[source_number - 1],
            preserve_artwork=layout.construction_role in {"cover", "recap"},
        )
        _render_native_page(
            slide,
            page,
            presentation,
            layout,
            page_number,
            page_count,
            palette,
        )
        slide.notes_slide.notes_text_frame.text = (
            "source_document_revision: " + page.speaker_notes.source_document_revision
            + "\n\nteaching_unit_id: " + page.speaker_notes.teaching_unit_id
            + "\n\n" + "\n\n".join(block.full_text for block in page.speaker_notes.source_blocks)
        )
    _remove_original_slides(presentation, len(originals))
    return presentation


__all__ = ["render_personal_template_presentation"]
