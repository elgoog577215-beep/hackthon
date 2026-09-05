from __future__ import annotations

import asyncio
import json
import re
import uuid
from datetime import datetime, timezone
from copy import deepcopy
from typing import Any, Awaitable, Callable, Literal

from fastapi import APIRouter, Depends, HTTPException, Request
from fastapi.concurrency import run_in_threadpool
from fastapi.responses import FileResponse, StreamingResponse
from pydantic import BaseModel, Field

from ai_base import AIProviderRequestError, AIProviderUnavailable
from course_generation_budget import TeacherScriptGenerationTimeout
from teacher_asset_readiness import (
    teacher_lesson_plan_covers_sections as _plan_revision_covers_sections,
    teacher_lesson_script_can_generate,
)
from teacher_outline_source import has_teaching_structure, matches_course_shell, read_teacher_outline_source
from course_production_state import (
    read_course_production_state,
    teacher_asset_job_can_resume,
)
from dependencies import (
    get_teacher_lesson_authoring_repository,
    require_task_manager,
)
from learner_context import resolve_user_id
from lesson_identity import chapter_matches_lesson
from generation_streaming import structured_generation_stream
from course_schedule import lecture_duration_minutes
from teaching_design import (
    LESSON_TYPES,
    normalize_lesson_arrangement,
    recommend_lesson_arrangement,
    validate_lesson_arrangement,
)
from material_storage import MaterialStorageError, material_repository
from material_parser import parse_document_path, parse_material_asset
from jobs.manager import TaskManager
from teacher_lesson_authoring import (
    LESSON_PLAN_PIPELINE_VERSION,
    TeacherLessonAuthoringError,
    TeacherLessonAuthoringRepository,
    TeacherLessonAuthoringService,
    build_uploaded_ppt_review_report,
    extract_uploaded_pptx_evidence,
    extract_uploaded_pptx_review,
    lesson_scope,
    teacher_lesson_deck_to_structured_slide_deck,
    teacher_lesson_v6_source,
)
from teacher_asset_readiness import (
    teacher_lesson_plan_readiness,
    teacher_lesson_plan_revision_has_content,
    teacher_lesson_ppt_asset_readiness,
    teacher_lesson_script_readiness,
    teacher_lesson_script_revision_has_content,
)
from question_bank import approved_formal_tasks, question_bank_repository
from teacher_script import (
    compile_teacher_script_module_contract,
    normalize_teacher_script_section,
    validate_teacher_script_section,
)
from teacher_script_visuals import (
    TeacherScriptVisualService,
    script_animation_runtime_enabled,
    teacher_script_visual_service,
)
from teacher_course_space import teacher_course_space_repository
from teacher_lesson_source import compile_original_lesson_plan_evidence
from slide_deck_renderer import export_structured_slide_deck
from representation_compiler import export_slide_deck_pptx
from representation_edits import (
    classify_representation_edit,
    representation_edit_impact,
)
from slide_deck_v6_orchestrator import (
    SlideDeckV6CandidateRepository,
    SlideDeckV6Orchestrator,
    V6BuildError,
)
from slide_ai_planning_v6 import (
    build_ai_base_story_planner_v6,
    build_ai_base_visual_planner_v2,
    regenerate_ppt_manuscript_pages_v1,
)
from ppt_template_packs import (
    TemplatePackError,
    ppt_template_pack_repository,
)
from teaching_representations import (
    RepresentationConflict,
    TeachingRepresentationSpec,
    teaching_representation_repository,
)
from template_layout_contract import (
    TemplateLayoutPackContractV1,
    compile_builtin_template_layout_contract_v1,
)
from course_document import (
    CourseDocument,
    refresh_document_revision,
    stable_hash,
)
from course_presentation_graph import compile_course_presentation_graph
from course_knowledge_base import course_knowledge_base_prompt_context
from slide_deck_v6 import (
    PptManuscriptV1,
    SlideDeckV6,
    compile_slide_deck_v6_from_manuscript,
    project_ppt_manuscript_from_deck_v1,
    rebase_ppt_manuscript_source_blocks_v1,
    revise_ppt_manuscript_v1,
)


router = APIRouter(prefix="/teacher", tags=["teacher-lesson-authoring"])
_background_jobs: set[asyncio.Task] = set()


def get_teacher_script_visual_service() -> TeacherScriptVisualService:
    return teacher_script_visual_service


async def _run_lesson_plan_job(
    *,
    course_id: str,
    job_id: str,
    repository: TeacherLessonAuthoringRepository,
    run: Callable[[], Awaitable[None]],
) -> None:
    """Run one queued lesson; the shared service semaphore owns concurrency."""
    job = await run_in_threadpool(repository.get_job, course_id, job_id)
    if str(job.get("status") or "") not in {"pending", "running"}:
        return
    await run()


def _latest_teacher_asset_job(
    jobs: list[dict[str, Any]],
    lesson_unit_id: str,
    job_type: str,
) -> dict[str, Any] | None:
    """Select the same latest asset attempt that the production projection sees."""

    candidates = [
        (index, item)
        for index, item in enumerate(jobs)
        if isinstance(item, dict)
        and str(item.get("lesson_unit_id") or "") == lesson_unit_id
        and str(item.get("type") or "") == job_type
    ]
    if not candidates:
        return None
    return max(
        candidates,
        key=lambda pair: (
            str(pair[1].get("updated_at") or pair[1].get("created_at") or ""),
            pair[0],
        ),
    )[1]


def _teacher_asset_job_can_resume(job: dict[str, Any] | None) -> bool:
    return teacher_asset_job_can_resume(job)


def _validated_teacher_asset_resume_job(
    repository: TeacherLessonAuthoringRepository,
    *,
    course_id: str,
    resume_job_id: str,
    lesson_unit_id: str,
    job_type: str,
    source_revision_field: str,
    source_revision_id: str,
) -> dict[str, Any]:
    try:
        candidate = repository.get_job(course_id, resume_job_id)
    except TeacherLessonAuthoringError as exc:
        raise TeacherLessonAuthoringError(
            "teacher_asset_resume_conflict",
            "要恢复的生成任务不存在或已不属于当前课程，请刷新状态后重试。",
            details={"resume_job_id": resume_job_id, "reason": "job_not_found"},
        ) from exc
    reason = ""
    if str(candidate.get("course_id") or "") != course_id:
        reason = "course_mismatch"
    elif str(candidate.get("lesson_unit_id") or "") != lesson_unit_id:
        reason = "lesson_mismatch"
    elif str(candidate.get("type") or "") != job_type:
        reason = "job_type_mismatch"
    elif str(candidate.get(source_revision_field) or "") != source_revision_id:
        reason = "source_revision_changed"
    elif not _teacher_asset_job_can_resume(candidate):
        reason = "resume_not_allowed"
    if reason:
        raise TeacherLessonAuthoringError(
            "teacher_asset_resume_conflict",
            "原任务的来源或状态已变化，不能继续原输入；请基于当前内容重新生成。",
            details={"resume_job_id": resume_job_id, "reason": reason},
        )
    return candidate


def _validate_new_attempt(repository: TeacherLessonAuthoringRepository, course_id: str, lesson_id: str, job_type: str, body: Any) -> None:
    if not body.retry_of_job_id:
        return
    if body.resume_job_id:
        raise TeacherLessonAuthoringError("teacher_asset_retry_mode_conflict", "原输入恢复与新输入尝试不能同时提交。")
    previous = repository.get_job(course_id, body.retry_of_job_id)
    if previous.get("course_id") != course_id or previous.get("lesson_unit_id") != lesson_id or previous.get("type") != job_type or previous.get("status") not in {"failed", "cancelled", "paused"}:
        raise TeacherLessonAuthoringError("teacher_asset_retry_conflict", "只能从当前讲次已经停止的任务发起新尝试。")


class GenerateLessonPlanRequest(BaseModel):
    request_id: str = Field(default="", max_length=160)
    resume_job_id: str = Field(default="", max_length=160)
    retry_of_job_id: str = Field(default="", max_length=160)
    source_package_id: str = Field(default="", max_length=160)
    source_asset_id: str = Field(default="", max_length=160)
    requirements: str = Field(default="", max_length=4000)
    material_asset_ids: list[str] = Field(default_factory=list, max_length=24)
    batch_parent_job_id: str = Field(default="", max_length=160)
    batch_position: int = Field(default=0, ge=0, le=1000)
    batch_size: int = Field(default=0, ge=0, le=1000)
    batch_source_revision_id: str = Field(default="", max_length=160)


class GenerateAllLessonPlansRequest(BaseModel):
    request_id: str = Field(default="", max_length=160)
    source_package_id: str = Field(default="", max_length=160)
    source_asset_id: str = Field(default="", max_length=160)
    requirements: str = Field(default="", max_length=4000)
    material_asset_ids: list[str] = Field(default_factory=list, max_length=24)
    regenerate_ready: bool = False
    resume_job_ids: list[str] = Field(default_factory=list, max_length=1000)


class UpdateLessonTypeRequest(BaseModel):
    lesson_type: str


class SaveLessonPlanDraftRequest(BaseModel):
    plan: dict[str, Any]
    source_outline_revision_id: str = ""
    expected_current_revision_id: str = ""


class GenerateLessonScriptRequest(BaseModel):
    request_id: str = Field(default="", max_length=160)
    resume_job_id: str = Field(default="", max_length=160)
    retry_of_job_id: str = Field(default="", max_length=160)
    requirements: str = Field(default="", max_length=4000)
    material_asset_ids: list[str] = Field(default_factory=list, max_length=24)
    batch_parent_job_id: str = Field(default="", max_length=160)
    batch_position: int = Field(default=0, ge=0, le=1000)
    batch_size: int = Field(default=0, ge=0, le=1000)
    batch_source_revision_id: str = Field(default="", max_length=160)


class GenerateAllLessonScriptsRequest(BaseModel):
    request_id: str = Field(default="", max_length=160)
    requirements: str = Field(default="", max_length=4000)
    regenerate_ready: bool = False
    resume_job_ids: list[str] = Field(default_factory=list, max_length=1000)


class SaveLessonScriptDraftRequest(BaseModel):
    base_revision_id: str = ""
    sections: list[dict[str, Any]]


class RewriteLessonScriptRequest(BaseModel):
    base_revision_id: str
    section_node_id: str
    instruction: str = Field(min_length=1, max_length=2000)
    material_asset_ids: list[str] = Field(default_factory=list, max_length=24)


class ResolveLessonScriptCandidateRequest(BaseModel):
    accept: bool


class CreateLessonScriptVisualRequest(BaseModel):
    script_revision_id: str = Field(min_length=1, max_length=200)
    section_node_id: str = Field(min_length=1, max_length=200)
    block_id: str = Field(min_length=1, max_length=240)
    expression_type: Literal["diagram", "image", "animation"]
    instruction: str = Field(default="", max_length=1200)


class ResolveLessonScriptVisualRequest(BaseModel):
    script_revision_id: str = Field(min_length=1, max_length=200)
    accept: bool


class CreateLessonPlanCandidateRequest(BaseModel):
    instruction: str = Field(min_length=1, max_length=2000)
    section_node_id: str = ""
    target_field: str = Field(default="", max_length=80)
    target_item_id: str = Field(default="", max_length=200)
    selected_text: str = Field(default="", max_length=1200)
    base_revision_id: str
    material_asset_ids: list[str] = Field(default_factory=list, max_length=24)


class ResolveLessonPlanCandidateRequest(BaseModel):
    accept: bool


class TeacherLessonV6BuildRequest(BaseModel):
    mode: str = "teaching"
    theme: str = "academic-editorial"
    template_pack_id: str = Field(default="", max_length=200)
    template_version: int | None = Field(default=None, ge=1)
    template_pack_version: int | None = Field(default=None, ge=1)
    force_rebuild: bool = False
    resume_task_id: str = Field(default="", max_length=200)


class _TeacherPptV6JobStopped(RuntimeError):
    def __init__(self, job: dict[str, Any]) -> None:
        self.job = deepcopy(job)
        super().__init__(str(job.get("status") or "stopped"))


def _teacher_ppt_resume_job_id(
    repository: TeacherLessonAuthoringRepository,
    course_id: str,
    lesson_unit_id: str,
    resume_task_id: str,
    *,
    job_type: str,
    source_lesson_plan_revision_id: str,
    source_lesson_script_revision_id: str,
    source_material_revision: str,
) -> str:
    """Validate the exact formal PPT attempt before reusing its checkpoint."""

    candidate_id = str(resume_task_id or "").strip()
    if not candidate_id:
        return ""
    try:
        candidate = repository.get_job(course_id, candidate_id)
    except TeacherLessonAuthoringError as exc:
        if exc.code != "teacher_job_not_found":
            raise
        raise TeacherLessonAuthoringError(
            "teacher_lesson_ppt_resume_job_mismatch",
            "要恢复的 PPT 任务不存在，请刷新状态后重试。",
            details={"resume_job_id": candidate_id, "reason": "job_not_found"},
        ) from exc
    reason = ""
    if str(candidate.get("course_id") or "") != course_id:
        reason = "course_mismatch"
    elif str(candidate.get("lesson_unit_id") or "") != lesson_unit_id:
        reason = "lesson_mismatch"
    elif str(candidate.get("type") or "") != job_type:
        reason = "job_type_mismatch"
    elif not _teacher_asset_job_can_resume(candidate):
        reason = "resume_not_allowed"
    elif (
        str(candidate.get("source_lesson_plan_revision_id") or "")
        != source_lesson_plan_revision_id
    ):
        reason = "lesson_plan_revision_changed"
    elif (
        str(candidate.get("source_script_revision_id") or "")
        != source_lesson_script_revision_id
    ):
        reason = "lesson_script_revision_changed"
    elif (
        str(candidate.get("source_material_revision") or "")
        != source_material_revision
    ):
        reason = "source_material_revision_changed"
    if reason:
        raise TeacherLessonAuthoringError(
            "teacher_lesson_ppt_resume_job_mismatch",
            "要恢复的 PPT 任务与当前内容或恢复状态不一致，请刷新状态后重试。",
            details={"resume_job_id": candidate_id, "reason": reason},
        )
    return candidate_id


def _teacher_ppt_job_must_be_active(
    repository: TeacherLessonAuthoringRepository,
    course_id: str,
    job_id: str,
) -> dict[str, Any]:
    job = repository.get_job(course_id, job_id)
    if str(job.get("status") or "") not in {"pending", "running"}:
        raise _TeacherPptV6JobStopped(job)
    return job


def _teacher_ppt_stopped_event(job: dict[str, Any]) -> dict[str, Any]:
    status = str(job.get("status") or "cancelled")
    return {
        "event": "build_paused" if status == "paused" else "build_cancelled",
        "task_id": str(job.get("id") or ""),
        "progress": int(job.get("progress") or 0),
        "stage": str(job.get("phase") or status),
        "status": status,
        "message": str(job.get("message") or ""),
    }


def _teacher_ppt_queued_event(
    job: dict[str, Any],
    *,
    target_schema: str,
) -> dict[str, Any]:
    return {
        "event": "build_queued",
        "task_id": str(job.get("id") or ""),
        "job": deepcopy(job),
        "progress": int(job.get("progress") or 0),
        "stage": str(job.get("phase") or "queued"),
        "target_schema": target_schema,
    }


def _fail_teacher_ppt_job(
    repository: TeacherLessonAuthoringRepository,
    course_id: str,
    job_id: str,
    *,
    code: str,
    message: str,
    retryable: bool,
    phase: str,
) -> dict[str, Any]:
    current = repository.get_job(course_id, job_id)
    if str(current.get("status") or "") in {"paused", "cancelled"}:
        return current
    return repository.update_job(
        course_id,
        job_id,
        status="failed",
        phase=phase,
        message=message,
        stream_complete=True,
        error={"code": code, "message": message, "retryable": retryable},
    )


def _requested_template_version(body: TeacherLessonV6BuildRequest) -> int | None:
    if (
        body.template_version is not None
        and body.template_pack_version is not None
        and body.template_version != body.template_pack_version
    ):
        raise HTTPException(
            status_code=422,
            detail={
                "code": "lesson_ppt_template_version_conflict",
                "message": "模板版本参数不一致，请重新选择模板。",
            },
        )
    return body.template_version or body.template_pack_version


def _resolve_teacher_v6_template(
    body: TeacherLessonV6BuildRequest,
    actor: str,
) -> TemplateLayoutPackContractV1:
    version = _requested_template_version(body)
    if version is not None and not body.template_pack_id:
        raise HTTPException(
            status_code=422,
            detail={
                "code": "lesson_ppt_template_pack_missing",
                "message": "指定模板版本时必须同时指定模板。",
            },
        )
    if not body.template_pack_id:
        try:
            return compile_builtin_template_layout_contract_v1(body.theme)
        except KeyError as exc:
            raise HTTPException(
                status_code=422,
                detail={
                    "code": "lesson_ppt_template_unavailable",
                    "message": "所选内置模板不可用，请重新选择。",
                },
            ) from exc
    try:
        return ppt_template_pack_repository.resolve_v6_layout_contract(
            body.template_pack_id,
            version,
            actor,
        )
    except (FileNotFoundError, TemplatePackError, ValueError) as exc:
        raise HTTPException(
            status_code=404,
            detail={
                "code": "lesson_ppt_template_unavailable",
                "message": "所选个人模板版本不存在或尚未完成构造确认。",
            },
        ) from exc


def _resolve_locked_teacher_v6_template(
    state: dict[str, Any],
    actor: str,
) -> TemplateLayoutPackContractV1:
    pack_id = str(state.get("template_pack_id") or "")
    if pack_id:
        try:
            template = ppt_template_pack_repository.resolve_v6_layout_contract(
                pack_id,
                str(state.get("template_version") or ""),
                actor,
            )
        except (FileNotFoundError, TemplatePackError, ValueError) as exc:
            raise HTTPException(
                status_code=409,
                detail={
                    "code": "lesson_ppt_template_lock_unavailable",
                    "message": "页面内容稿锁定的模板版本不可用，请重新生成页面内容稿。",
                },
            ) from exc
    else:
        try:
            template = compile_builtin_template_layout_contract_v1(
                str(state.get("theme") or "academic-editorial")
            )
        except KeyError as exc:
            raise HTTPException(
                status_code=409,
                detail={
                    "code": "lesson_ppt_template_lock_unavailable",
                    "message": "页面内容稿锁定的内置模板不可用，请重新生成页面内容稿。",
                },
            ) from exc
    expected = (
        str(state.get("template_id") or ""),
        str(state.get("template_version") or ""),
        str(state.get("template_digest") or ""),
    )
    actual = (
        template.template_id,
        template.template_version,
        template.template_digest,
    )
    if all(expected) and expected != actual:
        raise HTTPException(
            status_code=409,
            detail={
                "code": "lesson_ppt_template_lock_drifted",
                "message": "模板合同与已确认页面内容稿不一致，请重新生成页面内容稿。",
            },
        )
    return template


class ConfirmTeacherLessonPptManuscriptRequest(BaseModel):
    manuscript_revision: str = Field(min_length=1, max_length=200)


class UpdateTeacherLessonPptManuscriptRequest(BaseModel):
    expected_manuscript_revision: str = Field(min_length=1, max_length=200)
    page_updates: list[dict[str, Any]] = Field(min_length=1, max_length=80)


class RegenerateTeacherLessonPptManuscriptPagesRequest(BaseModel):
    expected_manuscript_revision: str = Field(min_length=1, max_length=200)
    target_page_ids: list[str] = Field(default_factory=list, max_length=24)
    changed_source_block_ids: list[str] = Field(default_factory=list, max_length=200)


class TeacherLessonRepresentationEditRequest(BaseModel):
    unit_id: str
    field: str
    before: Any = None
    after: Any = None
    semantic_intent: bool = False


class TeacherLessonApplyRepresentationEditRequest(TeacherLessonRepresentationEditRequest):
    decision: str = "representation_only"


class CreateTeacherLessonV6CandidateRequest(BaseModel):
    page_id: str = Field(min_length=1, max_length=200)
    instruction: str = Field(min_length=1, max_length=2000)
    base_spec_id: str = Field(min_length=1, max_length=200)
    base_spec_revision: str = Field(min_length=1, max_length=200)


class ResolveTeacherLessonV6CandidateRequest(BaseModel):
    accept: bool


class CreateImportedPptReviewRequest(BaseModel):
    package_id: str = Field(min_length=1, max_length=200)
    asset_id: str = Field(min_length=1, max_length=200)


class UpdateImportedPptSlideRequest(BaseModel):
    base_revision_id: str = Field(min_length=1, max_length=200)
    blocks: list[dict[str, Any]] = Field(default_factory=list, max_length=80)


class CreateImportedPptCandidateRequest(BaseModel):
    base_revision_id: str = Field(min_length=1, max_length=200)
    slide_id: str = Field(min_length=1, max_length=240)
    instruction: str = Field(min_length=1, max_length=2000)


class ResolveImportedPptCandidateRequest(BaseModel):
    accept: bool


class ConfirmImportedPptReviewRequest(BaseModel):
    revision_id: str = Field(min_length=1, max_length=200)


def _raise(exc: TeacherLessonAuthoringError) -> None:
    status = 404 if exc.code.endswith("not_found") else 409
    raise HTTPException(
        status_code=status,
        detail={"code": exc.code, "message": str(exc), **exc.details},
    ) from exc


def _assert_ppt_manuscript_confirmable(manuscript: dict) -> None:
    if manuscript.get("quality_status") != "passed":
        issues = [
            str(item).strip()
            for item in manuscript.get("quality_issues") or []
            if str(item).strip()
        ]
        raise TeacherLessonAuthoringError(
            "lesson_ppt_manuscript_quality_blocked",
            "页面内容稿质量未通过，修改后才能确认并生成正式 PPT。",
            details={"quality_issues": issues},
        )
    if manuscript.get("teaching_content_contract_version") != "page_teaching_v1":
        raise TeacherLessonAuthoringError(
            "lesson_ppt_manuscript_teaching_contract_outdated",
            "当前页面内容稿缺少逐页教学设计，请重新生成后再确认。",
        )


def _ppt_manuscript_state_payload(
    state: dict[str, Any] | None,
    *,
    generation_branch: str,
    current_material_revision: str = "",
) -> dict[str, Any]:
    if not state:
        return {
            "generation_branch": generation_branch,
            "revision": "",
            "status": "not_generated",
            "source_state": "current",
            "confirmable": False,
            "can_generate_ppt": False,
            "manuscript": None,
        }
    source_state = str(state.get("source_state") or "current")
    if (
        current_material_revision
        and state.get("source_material_revision")
        and state.get("source_material_revision") != current_material_revision
    ):
        source_state = "stale"
    manuscript = state.get("manuscript")
    manuscript_payload = manuscript if isinstance(manuscript, dict) else None
    quality_passed = bool(
        manuscript_payload
        and manuscript_payload.get("quality_status") == "passed"
    )
    status = str(state.get("status") or "draft")
    return {
        "generation_branch": generation_branch,
        "revision": str(state.get("revision") or ""),
        "status": status,
        "source_state": source_state,
        "confirmable": bool(
            quality_passed and source_state == "current" and status == "draft"
        ),
        "can_generate_ppt": bool(
            quality_passed and source_state == "current" and status == "confirmed"
        ),
        "task_id": str(state.get("task_id") or ""),
        "mode": str(state.get("mode") or "teaching"),
        "theme": str(state.get("theme") or "academic-editorial"),
        "template_id": str(state.get("template_id") or ""),
        "template_version": str(state.get("template_version") or ""),
        "template_digest": str(state.get("template_digest") or ""),
        "template_pack_id": str(state.get("template_pack_id") or ""),
        "generated_representation_id": str(
            state.get("generated_representation_id") or ""
        ),
        "last_good_revision": str(
            (state.get("last_good_manuscript") or {}).get("manuscript_revision")
            if isinstance(state.get("last_good_manuscript"), dict) else ""
        ),
        "last_confirmed_revision": str(
            (state.get("last_confirmed_manuscript") or {}).get("manuscript_revision")
            if isinstance(state.get("last_confirmed_manuscript"), dict) else ""
        ),
        "manuscript": manuscript_payload,
    }


_V6_KEY_REGION_SLOTS = (
    "interpretation",
    "conclusion",
    "takeaway",
    "body",
    "content",
    "task",
    "steps",
    "items",
)


def _v6_page_expression(page: dict[str, Any]) -> dict[str, str]:
    regions = [item for item in page.get("regions") or [] if isinstance(item, dict)]
    subtitle_region = next(
        (item for item in regions if str(item.get("slot_id") or "") == "subtitle"),
        None,
    )
    key_region = next(
        (
            item
            for slot_id in _V6_KEY_REGION_SLOTS
            for item in regions
            if str(item.get("slot_id") or "") == slot_id
        ),
        next(
            (
                item
                for item in regions
                if str(item.get("slot_id") or "") not in {"eyebrow", "subtitle"}
                and str(item.get("content") or "").strip()
            ),
            None,
        ),
    )
    return {
        "page_id": str(page.get("page_id") or ""),
        "title": str(page.get("title") or "").strip(),
        "subtitle": str((subtitle_region or {}).get("content") or "").strip(),
        "key_message": str((key_region or {}).get("content") or "").strip(),
        "subtitle_region_id": str((subtitle_region or {}).get("region_id") or ""),
        "key_region_id": str((key_region or {}).get("region_id") or ""),
    }


def _apply_v6_page_expression(
    page: dict[str, Any],
    *,
    field: str,
    value: Any,
    target_region_id: str = "",
) -> None:
    if field == "title":
        page["title"] = str(value or "").strip()
        return
    expression = _v6_page_expression(page)
    if field == "subtitle":
        region_id = target_region_id or expression["subtitle_region_id"]
    elif field == "key_message":
        region_id = target_region_id or expression["key_region_id"]
    else:
        raise ValueError(f"v6_expression_field_unsupported:{field}")
    if not region_id:
        raise ValueError(f"v6_expression_region_missing:{field}")
    region = next(
        (
            item
            for item in page.get("regions") or []
            if isinstance(item, dict) and str(item.get("region_id") or "") == region_id
        ),
        None,
    )
    if region is None:
        raise ValueError(f"v6_expression_region_missing:{field}")
    region["content"] = str(value or "").strip()


def _refresh_v6_ppt_manuscript(
    content: dict[str, Any],
    *,
    course_view: dict[str, Any],
    source_lesson_plan_revision_id: str,
) -> dict[str, Any]:
    """页面文案变化后同步重建页面内容稿投影，避免最终 PPT 出现第二内容源。"""
    deck = SlideDeckV6.model_validate({
        key: content[key]
        for key in SlideDeckV6.model_fields
        if key in content
    })
    teacher_source = dict(course_view.get("teacher_lesson_source") or {})
    previous = (
        dict(content.get("ppt_manuscript") or {})
        if isinstance(content.get("ppt_manuscript"), dict)
        else {}
    )
    manuscript = project_ppt_manuscript_from_deck_v1(
        deck,
        source_lesson_plan_revision_id=source_lesson_plan_revision_id,
        source_script_revision_id=str(teacher_source.get("script_revision_id") or ""),
        material_bindings=list(
            teacher_source.get("material_bindings")
            or previous.get("material_bindings")
            or []
        ),
        page_material_evidence_ids={
            str(page.get("page_id") or ""): list(
                page.get("source_material_evidence_ids") or []
            )
            for page in previous.get("pages") or []
            if isinstance(page, dict) and str(page.get("page_id") or "")
        },
    )
    content["ppt_manuscript"] = manuscript.model_dump(mode="json")
    return content["ppt_manuscript"]


def _source_course(
    tm: TaskManager,
    course_id: str,
    *,
    allow_empty: bool = False,
) -> dict[str, Any]:
    raw = tm.storage.load_course(course_id) if tm.storage else None
    source = read_teacher_outline_source(raw or {"course_id": course_id}, tm)
    if not has_teaching_structure(source) and not (allow_empty and matches_course_shell(raw, course_id)):
        raise TeacherLessonAuthoringError("course_not_found", "课程不存在或没有可用大纲。")
    return source


def _canonical_outline_revision(source: dict[str, Any]) -> str:
    """Return the revision consumed by the V3 teaching-plan engine.

    ``blueprint_revision_id`` identifies a broader course snapshot, while the
    knowledge-scope revision is the exact frozen outline contract used by the
    lesson planner. Every read, generation and edit must use the latter when
    available or a freshly generated plan becomes stale on reload.
    """
    return str(
        (source.get("course_knowledge_scope_contract") or {}).get("revision_id")
        or (source.get("course_teaching_plan") or {}).get("source_outline_revision_id")
        or source.get("blueprint_revision_id")
        or ""
    )


def _course_material_evidence(
    course_id: str,
    actor: str,
    material_asset_ids: list[str],
) -> tuple[list[str], list[dict[str, Any]]]:
    """Load only evidence from explicitly selected, course-owned materials."""
    selected_ids = list(dict.fromkeys(
        str(value or "").strip()
        for value in material_asset_ids
        if str(value or "").strip()
    ))
    if not selected_ids:
        return [], []

    allowed_material_ids: set[str] = set()
    for summary in teacher_course_space_repository.list_owned(actor, course_id):
        try:
            package = teacher_course_space_repository.load_owned(
                str(summary.get("package_id") or ""), actor
            )
        except (FileNotFoundError, MaterialStorageError):
            continue
        allowed_material_ids.update(
            str(item.get("material_asset_id") or "")
            for item in package.get("assets") or []
            if str(item.get("material_asset_id") or "")
        )
    unknown_material_ids = sorted(set(selected_ids) - allowed_material_ids)
    if unknown_material_ids:
        raise TeacherLessonAuthoringError(
            "lesson_material_source_not_found",
            "部分已选资料不属于当前课程。",
            details={"material_asset_ids": unknown_material_ids},
        )

    parsing_material_ids: list[str] = []
    failed_material_ids: list[str] = []
    for material_asset_id in selected_ids:
        asset = material_repository.get_asset(material_asset_id)
        parsed = material_repository.load_parsed_document(material_asset_id)
        status = str(getattr(asset, "status", "") or "")
        parse_status = str(getattr(parsed, "parse_status", "") or "")
        if status == "failed" or parse_status == "failed":
            failed_material_ids.append(material_asset_id)
        elif status in {"uploaded", "pending", "parsing"} and parse_status not in {"parsed", "degraded"}:
            parsing_material_ids.append(material_asset_id)
    if parsing_material_ids:
        raise TeacherLessonAuthoringError(
            "lesson_material_source_processing",
            "资料正在解析，完成后即可生成。",
            details={"material_asset_ids": parsing_material_ids},
        )
    if failed_material_ids:
        raise TeacherLessonAuthoringError(
            "lesson_material_source_parse_failed",
            "部分资料解析失败，请移除或重新上传后再生成。",
            details={"material_asset_ids": failed_material_ids},
        )

    evidence: list[dict[str, Any]] = []
    for material_asset_id in selected_ids:
        for item in material_repository.load_evidence(material_asset_id):
            if not isinstance(item, dict):
                continue
            evidence.append({
                **item,
                "asset_id": material_asset_id,
                "source_kind": "course_material",
            })
    return selected_ids, evidence


def _capture_generation_source_snapshot(
    *,
    course_id: str,
    actor: str,
    target_id: str,
    target_type: str,
    target_label: str,
    target_revision: str = "",
    task_id: str = "",
) -> None:
    teacher_course_space_repository.capture_owned_generation_source_snapshot(
        actor,
        course_id,
        target_id=target_id,
        target_type=target_type,
        target_label=target_label,
        target_revision=target_revision,
        task_id=task_id,
    )


def _lesson_plan_material_scope(
    course_id: str,
    actor: str,
    lesson_unit_id: str,
) -> dict[str, Any]:
    """Resolve the exact per-lesson material scope saved by the teacher."""
    target_id = f"lesson-plan:{lesson_unit_id}"
    material_asset_ids: list[str] = []
    source_package_id = ""
    source_asset_id = ""
    for summary in teacher_course_space_repository.list_owned(actor, course_id):
        package_id = str(summary.get("package_id") or "")
        try:
            package = teacher_course_space_repository.load_owned(package_id, actor)
        except (FileNotFoundError, MaterialStorageError):
            continue
        for relationship in teacher_course_space_repository.relationships_for_target(
            package, target_id
        ):
            material_asset_id = str(relationship.get("material_asset_id") or "")
            if material_asset_id and material_asset_id not in material_asset_ids:
                material_asset_ids.append(material_asset_id)
            if relationship.get("role") == "primary" and not source_asset_id:
                source_package_id = package_id
                source_asset_id = str(relationship.get("source_asset_id") or "")
    return {
        "source_package_id": source_package_id,
        "source_asset_id": source_asset_id,
        "material_asset_ids": material_asset_ids,
    }


def _lesson_script_material_scope(
    course_id: str,
    actor: str,
    lesson_unit_id: str,
) -> dict[str, Any]:
    """Resolve the teacher-owned material scope frozen for one script."""
    target_id = f"script:{lesson_unit_id}"
    material_asset_ids: list[str] = []
    source_package_id = ""
    source_asset_id = ""
    for summary in teacher_course_space_repository.list_owned(actor, course_id):
        package_id = str(summary.get("package_id") or "")
        try:
            package = teacher_course_space_repository.load_owned(package_id, actor)
        except (FileNotFoundError, MaterialStorageError):
            continue
        for relationship in teacher_course_space_repository.relationships_for_target(
            package,
            target_id,
        ):
            material_asset_id = str(relationship.get("material_asset_id") or "")
            if material_asset_id and material_asset_id not in material_asset_ids:
                material_asset_ids.append(material_asset_id)
            if relationship.get("role") == "primary" and not source_asset_id:
                source_package_id = package_id
                source_asset_id = str(relationship.get("source_asset_id") or "")
    return {
        "source_package_id": source_package_id,
        "source_asset_id": source_asset_id,
        "material_asset_ids": material_asset_ids,
    }


def _ppt_material_bundle(
    course_id: str,
    actor: str,
    lesson_unit_id: str,
) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    """Resolve the PPT stage's frozen, teacher-owned material relationships."""
    target_id = f"ppt-v6:{lesson_unit_id}"
    bindings: list[dict[str, Any]] = []
    seen_material_ids: set[str] = set()
    for summary in teacher_course_space_repository.list_owned(actor, course_id):
        try:
            package = teacher_course_space_repository.load_owned(
                str(summary.get("package_id") or ""), actor
            )
        except (FileNotFoundError, MaterialStorageError):
            continue
        for relationship in teacher_course_space_repository.relationships_for_target(
            package, target_id
        ):
            material_asset_id = str(
                relationship.get("material_asset_id") or ""
            )
            if not material_asset_id or material_asset_id in seen_material_ids:
                continue
            seen_material_ids.add(material_asset_id)
            bindings.append({
                "material_asset_id": material_asset_id,
                "source_asset_id": str(
                    relationship.get("source_asset_id") or ""
                ),
                "source_label": str(
                    relationship.get("source_label") or material_asset_id
                ),
                "role": (
                    "primary"
                    if relationship.get("role") == "primary"
                    else "reference"
                ),
            })
    material_ids, evidence = _course_material_evidence(
        course_id,
        actor,
        [item["material_asset_id"] for item in bindings],
    )
    binding_by_material = {
        item["material_asset_id"]: item for item in bindings
    }
    normalized_evidence: list[dict[str, Any]] = []
    for index, item in enumerate(evidence):
        material_asset_id = str(item.get("asset_id") or "")
        binding = binding_by_material.get(material_asset_id) or {}
        evidence_id = str(
            item.get("evidence_id")
            or item.get("unit_id")
            or stable_hash(
                {
                    "material_asset_id": material_asset_id,
                    "index": index,
                    "text": str(
                        item.get("summary")
                        or item.get("source_text")
                        or item.get("text")
                        or item.get("content")
                        or ""
                    ),
                },
                prefix="pptev_",
            )
        )
        normalized_evidence.append({
            **item,
            "evidence_id": evidence_id,
            "asset_id": material_asset_id,
            "source_label": str(binding.get("source_label") or material_asset_id),
            "source_role": str(binding.get("role") or "reference"),
        })
    if set(material_ids) != set(binding_by_material):
        raise TeacherLessonAuthoringError(
            "lesson_material_source_not_found",
            "部分 PPT 资料来源无法读取。",
        )
    return bindings, normalized_evidence


def _ppt_reference_terms(value: str) -> set[str]:
    text = str(value or "").lower()
    terms = set(re.findall(r"[a-z][a-z0-9_+-]{1,30}", text))
    for group in re.findall(r"[\u4e00-\u9fff]{2,20}", text):
        terms.add(group)
        terms.update(
            group[index:index + width]
            for width in (2, 3, 4)
            for index in range(max(0, len(group) - width + 1))
        )
    return terms


def _attach_ppt_reference_evidence(
    document: CourseDocument,
    evidence: list[dict[str, Any]],
) -> CourseDocument:
    """Bind selected evidence to current usable script blocks without rewriting them."""
    if not evidence:
        return document
    evidence_terms = {
        str(item.get("evidence_id") or ""): _ppt_reference_terms(" ".join([
            " ".join(str(value) for value in item.get("keywords") or []),
            str(item.get("summary") or ""),
            str(item.get("source_text") or item.get("text") or item.get("content") or ""),
        ]))
        for item in evidence
        if str(item.get("evidence_id") or "")
    }
    section_titles = {
        section.section_id: section.title for section in document.sections
    }
    changed = False
    for block in document.blocks:
        block_query = " ".join([
            str((block.payload or {}).get("title") or ""),
            json.dumps(block.payload or {}, ensure_ascii=False),
        ])
        query_terms = _ppt_reference_terms(
            block_query or section_titles.get(block.section_id, "")
        )
        ranked = sorted(
            (
                (evidence_id, len(query_terms & terms))
                for evidence_id, terms in evidence_terms.items()
            ),
            key=lambda pair: (-pair[1], pair[0]),
        )
        selected = [evidence_id for evidence_id, score in ranked if score > 0][:4]
        if selected != block.evidence_refs:
            block.evidence_refs = selected
            changed = True
    return refresh_document_revision(document) if changed else document


def _prompt_material_evidence(
    evidence: list[dict[str, Any]],
    *,
    character_budget: int = 12000,
) -> list[dict[str, str]]:
    """Keep selected evidence useful while bounding one model request."""
    result: list[dict[str, str]] = []
    remaining = character_budget
    for item in evidence:
        raw_text = str(item.get("text") or item.get("content") or "").strip()
        if not raw_text or remaining <= 0:
            continue
        text = raw_text[:remaining]
        remaining -= len(text)
        result.append({
            "asset_id": str(item.get("asset_id") or ""),
            "unit_id": str(item.get("unit_id") or item.get("evidence_id") or ""),
            "text": text,
        })
    return result


def _lesson_projection(
    source: dict[str, Any],
    repository: TeacherLessonAuthoringRepository,
    authoring_state: dict[str, Any] | None = None,
) -> list[dict[str, Any]]:
    course_id = str(source.get("course_id") or "")
    assets = (authoring_state if authoring_state is not None else repository.view(course_id)).get("lessons") or {}
    nodes = [item for item in source.get("nodes") or [] if isinstance(item, dict)]
    lessons = [
        item for item in nodes
        if int(item.get("node_level") or 0) == 1
        and str(item.get("parent_node_id") or "").lower() in {"", "root"}
    ]
    result = []
    schedule_slots = (source.get("course_profile") or {}).get("schedule_slots") or []
    for index, lesson in enumerate(lessons, start=1):
        lesson_id = str(lesson.get("node_id") or "")
        sections = [
            item for item in nodes
            if str(item.get("parent_node_id") or "") == lesson_id
        ]
        asset = assets.get(lesson_id) if isinstance(assets, dict) else None
        plan_asset = deepcopy(asset) if isinstance(asset, dict) else {
            "lesson_unit_id": lesson_id,
            "arrangement": {
                "working_revision_id": "",
                "source_state": "current",
                "revisions": [],
            },
            "working_revision_id": "",
            "source_state": "current",
            "revisions": [],
            "working_script_revision_id": "",
            "script_revisions": [],
            "ppt_assets": [],
            "material_drafts": {},
            "current_material_draft_ids": {},
        }
        arrangement_state = plan_asset.get("arrangement") or {}
        arrangement_revision_id = str(arrangement_state.get("working_revision_id") or "")
        arrangement_revision = next(
            (
                item for item in arrangement_state.get("revisions") or []
                if isinstance(item, dict) and item.get("revision_id") == arrangement_revision_id
            ),
            None,
        )
        arrangement_is_current = (
            isinstance(arrangement_revision, dict)
            and str(arrangement_state.get("source_state") or "current") == "current"
        )
        arrangement = (
            deepcopy(arrangement_revision)
            if arrangement_is_current
            else recommend_lesson_arrangement(
                source,
                lesson_id,
                source_outline_revision_id=_canonical_outline_revision(source),
            )
        )
        arrangement["ready"] = bool(
            arrangement_is_current and list(arrangement.get("blocks") or [])
        )
        arrangement["source_state"] = "current" if not arrangement_is_current else str(
            arrangement_state.get("source_state") or "current"
        )
        expected_section_ids = [
            str(section.get("node_id") or "") for section in sections
        ]
        arrangement_issues = validate_lesson_arrangement(
            arrangement,
            expected_section_ids=expected_section_ids,
        )
        plan_can_generate = not arrangement_issues
        working_script_revision_id = str(plan_asset.get("working_script_revision_id") or "")
        script_revision = next(
            (
                item for item in plan_asset.get("script_revisions") or []
                if isinstance(item, dict) and item.get("revision_id") == working_script_revision_id
            ),
            None,
        )
        script_sections = deepcopy(
            script_revision.get("sections")
            if isinstance(script_revision, dict)
            else []
        )
        current_script_revision = str((script_revision or {}).get("revision_id") or "")
        script_quality = deepcopy((script_revision or {}).get("quality_report") or {})
        plan_revision_id = str(plan_asset.get("working_revision_id") or "")
        plan_readiness = teacher_lesson_plan_readiness(plan_asset)
        script_readiness = teacher_lesson_script_readiness(
            plan_asset,
            plan_readiness=plan_readiness,
        )
        plan_ready = bool(plan_readiness["ready"])
        script_ready = bool(script_readiness["ready"])
        plan_revision = next(
            (
                item for item in plan_asset.get("revisions") or []
                if isinstance(item, dict)
                and str(item.get("revision_id") or "") == plan_revision_id
            ),
            None,
        )
        script_can_generate = teacher_lesson_script_can_generate(plan_asset, expected_section_ids)
        script_source_state = (
            "stale"
            if (
                plan_readiness["unavailable_reason"] == "source_stale"
                or script_readiness["unavailable_reason"]
                in {"source_stale", "upstream_plan_mismatch"}
            )
            else "current"
        )
        for ppt_asset in plan_asset.get("ppt_assets") or []:
            if not isinstance(ppt_asset, dict):
                continue
            source_script_revision = str(ppt_asset.get("source_script_revision_id") or "")
            if (
                ppt_asset.get("engine") == "slide_deck_v6"
                and source_script_revision != current_script_revision
            ):
                ppt_asset["source_state"] = "stale"
            ppt_readiness = teacher_lesson_ppt_asset_readiness(
                plan_asset,
                ppt_asset,
                plan_readiness=plan_readiness,
                script_readiness=script_readiness,
            )
            ppt_asset.update(ppt_readiness)
        plan_projection = {
            "lesson_unit_id": lesson_id,
            "working_revision_id": plan_revision_id,
            "source_state": str(plan_asset.get("source_state") or "current"),
            "ready": bool(plan_readiness["ready"]),
            "unavailable_reason": str(plan_readiness["unavailable_reason"] or ""),
            "can_generate": plan_can_generate,
            "generation_unavailable_reason": (
                ""
                if plan_can_generate
                else str(
                    (arrangement_issues[0] if arrangement_issues else {}).get("code")
                    or "lesson_arrangement_unavailable"
                )
            ),
            "current_revision": deepcopy(plan_revision) if isinstance(plan_revision, dict) else None,
            "ai_candidate": next(
                (
                    deepcopy(candidate)
                    for candidate in reversed(plan_asset.get("ai_candidates") or [])
                    if isinstance(candidate, dict)
                    and candidate.get("status") == "pending"
                    and candidate.get("base_revision_id") == plan_revision_id
                ),
                None,
            ),
            "ppt_assets": deepcopy(plan_asset.get("ppt_assets") or []),
        }
        result.append({
            "lesson_unit_id": lesson_id,
            "number": index,
            "title": str(lesson.get("node_name") or f"第{index}讲"),
            "duration_minutes": int(
                lesson.get("duration_minutes")
                or lecture_duration_minutes(schedule_slots, index - 1)
            ),
            "sections": [
                {
                    "section_node_id": str(section.get("node_id") or ""),
                    "title": str(section.get("node_name") or ""),
                }
                for section in sections
            ],
            "arrangement": arrangement,
            "script": {
                "current_revision_id": current_script_revision,
                "source_lesson_plan_revision_id": str(
                    (script_revision or {}).get("source_lesson_plan_revision_id")
                    or ""
                ),
                "source_state": script_source_state,
                "ready": script_ready,
                "unavailable_reason": script_readiness["unavailable_reason"],
                "can_generate": script_can_generate,
                "generation_unavailable_reason": (
                    ""
                    if script_can_generate
                    else "lesson_plan_scope_stale"
                    if plan_ready
                    else str(
                        plan_readiness["unavailable_reason"]
                        or "lesson_plan_not_ready"
                    )
                ),
                "publication_eligible": script_ready,
                "generation_source": str(
                    (script_revision or {}).get("generation_source") or ""
                ),
                "quality_contract_version": str(
                    (script_revision or {}).get("quality_contract_version") or ""
                ),
                "quality_report": script_quality,
                "sections": script_sections,
                "actor": str((script_revision or {}).get("actor") or ""),
                "updated_at": str(
                    (script_revision or {}).get("updated_at")
                    or (script_revision or {}).get("created_at")
                    or ""
                ),
                "ai_candidate": next(
                    (
                        deepcopy(candidate)
                        for candidate in reversed(plan_asset.get("script_ai_candidates") or [])
                        if isinstance(candidate, dict)
                        and candidate.get("status") == "pending"
                        and candidate.get("base_revision_id") == current_script_revision
                    ),
                    None,
                ),
            },
            "plan": plan_projection,
            "material_drafts": {
                str(target_type): deepcopy(next(
                    (
                        item for item in reversed((plan_asset.get("material_drafts") or {}).get(target_type) or [])
                        if isinstance(item, dict)
                        and item.get("revision_id") == revision_id
                    ),
                    {},
                ))
                for target_type, revision_id in (plan_asset.get("current_material_draft_ids") or {}).items()
            },
        })
    return result


def _plan_revision(
    repository: TeacherLessonAuthoringRepository,
    course_id: str,
    lesson_unit_id: str,
    revision_id: str,
) -> dict[str, Any]:
    lesson = repository.lesson(course_id, lesson_unit_id)
    revision = next(
        (
            item for item in lesson.get("revisions") or []
            if isinstance(item, dict) and item.get("revision_id") == revision_id
        ),
        None,
    )
    if not isinstance(revision, dict):
        raise TeacherLessonAuthoringError("lesson_plan_revision_not_found", "教案修订不存在。")
    return revision


def _plan_revision_has_content(revision: dict[str, Any]) -> bool:
    return teacher_lesson_plan_revision_has_content(revision)


def _current_plan_revision(
    repository: TeacherLessonAuthoringRepository,
    course_id: str,
    lesson_unit_id: str,
) -> tuple[dict[str, Any], dict[str, Any]]:
    lesson = repository.lesson(course_id, lesson_unit_id)
    revision_id = str(lesson.get("working_revision_id") or "")
    if not revision_id or str(lesson.get("source_state") or "current") != "current":
        raise TeacherLessonAuthoringError(
            "lesson_plan_not_ready",
            "请先生成当前大纲对应的本讲教案。",
        )
    revision = _plan_revision(repository, course_id, lesson_unit_id, revision_id)
    if not _plan_revision_has_content(revision):
        raise TeacherLessonAuthoringError(
            "lesson_plan_incomplete",
            "当前教案结构不完整，请重新生成或编辑保存。",
        )
    return lesson, revision


def _script_revision(
    repository: TeacherLessonAuthoringRepository,
    course_id: str,
    lesson_unit_id: str,
    revision_id: str,
) -> dict[str, Any]:
    lesson = repository.lesson(course_id, lesson_unit_id)
    revision = next(
        (
            item for item in lesson.get("script_revisions") or []
            if isinstance(item, dict) and item.get("revision_id") == revision_id
        ),
        None,
    )
    if not isinstance(revision, dict):
        raise TeacherLessonAuthoringError(
            "lesson_script_revision_not_found",
            "讲义修订不存在。",
        )
    return revision


def _script_revision_has_content(revision: dict[str, Any]) -> bool:
    return teacher_lesson_script_revision_has_content(revision)


def _current_script_revision(
    repository: TeacherLessonAuthoringRepository,
    course_id: str,
    lesson_unit_id: str,
    *,
    source_plan_revision_id: str,
) -> tuple[dict[str, Any], dict[str, Any]]:
    lesson = repository.lesson(course_id, lesson_unit_id)
    revision_id = str(lesson.get("working_script_revision_id") or "")
    if not revision_id:
        raise TeacherLessonAuthoringError(
            "lesson_script_not_ready",
            "请先生成本讲讲义，再进入 PPT 工作台。",
        )
    revision = _script_revision(repository, course_id, lesson_unit_id, revision_id)
    if (
        str(revision.get("source_lesson_plan_revision_id") or "")
        != source_plan_revision_id
        or not _script_revision_has_content(revision)
    ):
        raise TeacherLessonAuthoringError(
            "lesson_script_source_stale",
            "当前讲义不完整或对应的教案已变化，请重新生成或编辑保存。",
        )
    return lesson, revision


def _current_script_visual_context(
    tm: TaskManager,
    repository: TeacherLessonAuthoringRepository,
    course_id: str,
    lesson_unit_id: str,
) -> tuple[dict[str, Any], list[dict[str, Any]]]:
    source = _source_course(tm, course_id)
    lesson_scope(source, lesson_unit_id)
    lesson = repository.lesson(course_id, lesson_unit_id)
    revision_id = str(lesson.get("working_script_revision_id") or "")
    if not revision_id:
        raise TeacherLessonAuthoringError(
            "lesson_script_not_ready",
            "请先生成本讲讲义，再添加视觉表达。",
        )
    revision = _script_revision(repository, course_id, lesson_unit_id, revision_id)
    if not _script_revision_has_content(revision):
        raise TeacherLessonAuthoringError(
            "lesson_script_source_incomplete",
            "当前讲义内容不完整，暂时不能添加视觉表达。",
        )
    blocks: list[dict[str, Any]] = []
    for section in revision.get("sections") or []:
        if not isinstance(section, dict):
            continue
        normalized = normalize_teacher_script_section(section)
        section_node_id = str(normalized.get("section_node_id") or "")
        for block in normalized.get("blocks") or []:
            if not isinstance(block, dict):
                continue
            blocks.append({**deepcopy(block), "section_node_id": section_node_id})
    if not blocks:
        raise TeacherLessonAuthoringError(
            "lesson_script_blocks_empty",
            "当前讲义没有可绑定视觉表达的教学块。",
        )
    return revision, blocks


def _imported_ppt_review_context(
    source: dict[str, Any],
    repository: TeacherLessonAuthoringRepository,
    course_id: str,
    lesson_unit_id: str,
    *,
    actor: str = "",
) -> tuple[list[dict[str, Any]], list[dict[str, Any]], dict[str, str]]:
    """Compile the exact upstream revisions used by one imported-deck review."""
    scoped = lesson_scope(source, lesson_unit_id)
    lesson = repository.lesson(course_id, lesson_unit_id)
    sources: list[dict[str, Any]] = []
    units: list[dict[str, Any]] = []
    revisions = {
        "outline": _canonical_outline_revision(source),
        "plan": str(lesson.get("working_revision_id") or ""),
        "script": "",
        "question_bank": "",
    }

    if revisions["outline"]:
        sources.append({
            "kind": "outline",
            "label": "课程大纲",
            "revision_id": revisions["outline"],
            "status": "current",
        })
    for section in scoped["sections"]:
        units.append({
            "kind": "outline",
            "label": str(section.get("node_name") or "未命名小节"),
            "revision_id": revisions["outline"],
            "text": "\n".join(filter(None, [
                str(section.get("node_name") or ""),
                str(section.get("learning_objective") or ""),
                str(section.get("node_content") or "")[:1600],
            ])),
        })

    plan_revision: dict[str, Any] = {}
    if revisions["plan"]:
        candidate = _plan_revision(
            repository, course_id, lesson_unit_id, revisions["plan"]
        )
        if (
            str(lesson.get("source_state") or "current") == "current"
            and _plan_revision_has_content(candidate)
        ):
            plan_revision = candidate
        else:
            revisions["plan"] = ""
    if revisions["plan"]:
        sources.append({
            "kind": "lesson_plan",
            "label": "当前教案",
            "revision_id": revisions["plan"],
            "status": "current",
        })
        for section in (plan_revision.get("plan") or {}).get("sections") or []:
            if not isinstance(section, dict):
                continue
            units.append({
                "kind": "lesson_plan",
                "label": str(section.get("title") or section.get("node_name") or section.get("node_id") or "教案小节"),
                "revision_id": revisions["plan"],
                "text": json.dumps(section, ensure_ascii=False)[:4000],
            })

    current_script = str(lesson.get("working_script_revision_id") or "")
    script_revision: dict[str, Any] = {}
    if current_script:
        candidate = _script_revision(repository, course_id, lesson_unit_id, current_script)
        if (
            candidate.get("source_lesson_plan_revision_id") == revisions["plan"]
            and _script_revision_has_content(candidate)
        ):
            script_revision = candidate
        else:
            current_script = ""
    if current_script:
        revisions["script"] = current_script
        sources.append({
            "kind": "script",
            "label": "当前讲义",
            "revision_id": current_script,
            "status": "current",
        })
        for section in script_revision.get("sections") or []:
            if not isinstance(section, dict):
                continue
            units.append({
                "kind": "script",
                "label": str(section.get("title") or section.get("section_node_id") or "讲义小节"),
                "revision_id": current_script,
                "text": str(section.get("content") or "") or json.dumps(section.get("blocks") or [], ensure_ascii=False),
            })

    bundle = question_bank_repository.load_bundle(course_id)
    if isinstance(bundle, dict):
        section_ids = {str(item.get("node_id") or "") for item in scoped["sections"]}
        items = [
            item for item in bundle.get("items") or []
            if isinstance(item, dict)
            and (
                str(item.get("node_id") or "") in section_ids
                or section_ids.intersection(str(value or "") for value in item.get("node_ids") or [])
            )
        ]
        if items:
            revisions["question_bank"] = str(bundle.get("bundle_revision_id") or "")
            sources.append({
                "kind": "question_bank",
                "label": f"题库（{len(items)} 题）",
                "revision_id": revisions["question_bank"],
                "status": "current",
            })
            units.append({
                "kind": "question_bank",
                "label": "题库考查内容",
                "revision_id": revisions["question_bank"],
                "text": "\n".join(
                    str(item.get("stem") or item.get("prompt") or item.get("question") or "")
                    for item in items[:30]
                ),
            })
    if actor:
        material_bindings, material_evidence = _ppt_material_bundle(
            course_id, actor, lesson_unit_id
        )
        for binding in material_bindings:
            role = str(binding.get("role") or "reference")
            sources.append({
                "kind": "primary_material" if role == "primary" else "reference_material",
                "label": (
                    f"主参考：{binding['source_label']}"
                    if role == "primary"
                    else f"参考：{binding['source_label']}"
                ),
                "revision_id": stable_hash(binding, prefix="pptref_"),
                "status": "current",
            })
        for item in material_evidence[:48]:
            text = str(
                item.get("summary")
                or item.get("source_text")
                or item.get("text")
                or item.get("content")
                or ""
            ).strip()
            if not text:
                continue
            units.append({
                "kind": "reference_material",
                "label": str(item.get("source_label") or "PPT 参考资料"),
                "revision_id": str(item.get("evidence_id") or ""),
                "text": text[:2400],
            })
    return sources, units, revisions


def _updated_imported_ppt_slides(
    review: dict[str, Any],
    slide_id: str,
    blocks: list[dict[str, Any]],
) -> list[dict[str, Any]]:
    slides = deepcopy(review.get("slides") or [])
    slide = next((item for item in slides if item.get("slide_id") == slide_id), None)
    if not isinstance(slide, dict):
        raise TeacherLessonAuthoringError("uploaded_ppt_slide_not_found", "PPT 页面不存在。")
    existing = {
        str(item.get("block_id") or ""): item
        for item in slide.get("blocks") or []
        if isinstance(item, dict)
    }
    for patch in blocks:
        block_id = str(patch.get("block_id") or "")
        block = existing.get(block_id)
        if not isinstance(block, dict) or not block.get("editable"):
            raise TeacherLessonAuthoringError("uploaded_ppt_block_not_editable", "该文字块不支持在线编辑。")
        text = str(patch.get("text") or "").strip()
        if len(text) > 6000:
            raise TeacherLessonAuthoringError("uploaded_ppt_block_too_long", "单个 PPT 文字块不能超过 6000 字符。")
        block["text"] = text
    title_block = next((item for item in existing.values() if item.get("kind") == "title"), None)
    slide["title"] = str((title_block or {}).get("text") or "")
    slide["content_hash"] = stable_hash(slide.get("blocks") or [])[:24]
    return slides


def _ppt_asset_revision(
    repository: TeacherLessonAuthoringRepository,
    course_id: str,
    lesson_unit_id: str,
    asset_id: str,
    revision_id: str,
) -> tuple[dict[str, Any], dict[str, Any]]:
    lesson = repository.lesson(course_id, lesson_unit_id)
    asset = next(
        (
            item for item in lesson.get("ppt_assets") or []
            if isinstance(item, dict) and item.get("asset_id") == asset_id
        ),
        None,
    )
    if not isinstance(asset, dict):
        raise TeacherLessonAuthoringError("lesson_ppt_not_found", "本讲 PPT 不存在。")
    revision = next(
        (
            item for item in asset.get("revisions") or []
            if isinstance(item, dict) and item.get("revision_id") == revision_id
        ),
        None,
    )
    if not isinstance(revision, dict):
        raise TeacherLessonAuthoringError("lesson_ppt_revision_not_found", "PPT 修订不存在。")
    return asset, revision


def _teacher_v6_source(
    tm: TaskManager,
    repository: TeacherLessonAuthoringRepository,
    course_id: str,
    lesson_unit_id: str,
):
    source = _source_course(tm, course_id)
    lesson, revision = _current_plan_revision(
        repository,
        course_id,
        lesson_unit_id,
    )
    revision_id = str(revision.get("revision_id") or "")
    expected_section_ids = [
        str(item.get("node_id") or "")
        for item in lesson_scope(source, lesson_unit_id)["sections"]
    ]
    actual_section_ids = [
        str(item.get("node_id") or "")
        for item in (revision.get("plan") or {}).get("sections") or []
        if isinstance(item, dict)
    ]
    if actual_section_ids != expected_section_ids:
        raise TeacherLessonAuthoringError(
            "lesson_plan_scope_stale",
            "当前教案没有完整对应本讲大纲，请重新生成或编辑保存。",
        )
    lesson, script_revision = _current_script_revision(
        repository,
        course_id,
        lesson_unit_id,
        source_plan_revision_id=revision_id,
    )
    document, course_view, synthetic_id = teacher_lesson_v6_source(
        source,
        lesson_unit_id=lesson_unit_id,
        plan_revision=revision,
        script_revision=script_revision,
    )
    return document, course_view, synthetic_id, lesson, revision


def _teacher_v6_registry_payload(synthetic_id: str) -> dict[str, Any]:
    registry = teaching_representation_repository.load(synthetic_id)
    payload = registry.model_dump(mode="json")
    payload["slide_deck_target_schema"] = "slide_deck_v6"
    payload["slide_deck_v6_eligible"] = True
    slide_representations = [
        item for item in registry.representations
        if item.representation_type == "slide_deck" and item.status != "archived"
    ]
    selected = slide_representations[0] if slide_representations else None
    spec = next(
        (item for item in registry.specs if selected and item.spec_id == selected.spec_id),
        None,
    )
    content = (spec.payload.get("content") if spec else None) or {}
    payload.update({
        "slide_deck_target_schema": "slide_deck_v6",
        "slide_deck_candidate_schema": str(content.get("schema_version") or ""),
        "slide_deck_published_schema": str(content.get("schema_version") or ""),
        "slide_deck_candidate_status": str(content.get("candidate_status") or content.get("status") or ""),
    })
    return payload


@router.get("/courses/{course_id}/lesson-authoring")
async def get_lesson_authoring_view(
    course_id: str,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        source = _source_course(tm, course_id, allow_empty=True)
        outline_revision = _canonical_outline_revision(source)
        authoring_state = repository.expire_stale_jobs(course_id)
        if outline_revision and str(authoring_state.get("outline_revision_id") or "") != outline_revision:
            authoring_state = repository.set_outline(course_id, outline_revision)
        outline_draft_id = str(authoring_state.get("current_outline_material_draft_id") or "")
        outline_material_draft = next(
            (
                deepcopy(item)
                for item in reversed(authoring_state.get("outline_material_drafts") or [])
                if isinstance(item, dict) and str(item.get("revision_id") or "") == outline_draft_id
            ),
            None,
        )
        return {
            "schema_version": "teacher_lesson_authoring_view_v1",
            "pipeline_version": LESSON_PLAN_PIPELINE_VERSION,
            "plan_schema_version": "course_teaching_plan_v3",
            "course_id": course_id,
            "outline_revision_id": outline_revision,
            "outline_material_draft": outline_material_draft,
            "lessons": _lesson_projection(source, repository, authoring_state),
            "jobs": list((authoring_state.get("jobs") or {}).values()),
            "course_production_state": read_course_production_state(
                source,
                repository,
                tm,
                authoring_state=authoring_state,
            ),
        }
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.put("/courses/{course_id}/lessons/{lesson_unit_id}/arrangement/type")
async def update_lesson_type(
    course_id: str,
    lesson_unit_id: str,
    body: UpdateLessonTypeRequest,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    """Save the lecture type and keep its current teaching blocks editable."""
    try:
        if body.lesson_type not in LESSON_TYPES:
            raise TeacherLessonAuthoringError(
                "lesson_type_invalid",
                "请选择有效的本讲课型。",
            )
        source = _source_course(tm, course_id)
        outline_revision = _canonical_outline_revision(source)
        repository.set_outline(course_id, outline_revision)
        projected = next(
            item for item in _lesson_projection(source, repository)
            if item["lesson_unit_id"] == lesson_unit_id
        )
        current = projected.get("arrangement") or {}
        arrangement = normalize_lesson_arrangement(
            {
                "lesson_type": body.lesson_type,
                "blocks": current.get("blocks") or [],
            },
            lesson_unit_id=lesson_unit_id,
            source_outline_revision_id=outline_revision,
        )
        # 大纲阶段只决定这一讲采用什么课型。教学块可以尚未生成；
        # 生成教案时会补齐当前结构，并只对结构完整性做准入校验。
        repository.save_arrangement_revision(
            course_id,
            lesson_unit_id,
            arrangement,
            source_outline_revision_id=outline_revision,
            actor=resolve_user_id(request.headers.get("X-User-Id")),
        )
        lesson = next(
            item for item in _lesson_projection(source, repository)
            if item["lesson_unit_id"] == lesson_unit_id
        )
        return {"lesson": lesson, "lesson_types": LESSON_TYPES}
    except StopIteration as exc:
        raise HTTPException(status_code=404, detail="未找到对应讲次。") from exc
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.post("/courses/{course_id}/lessons/{lesson_unit_id}/ppt-import/reviews")
async def create_imported_ppt_review(
    course_id: str,
    lesson_unit_id: str,
    body: CreateImportedPptReviewRequest,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(get_teacher_lesson_authoring_repository),
):
    try:
        actor = resolve_user_id(request.headers.get("X-User-Id"))
        package = teacher_course_space_repository.load_owned(body.package_id, actor)
        if str(package.get("course_id") or "") != course_id:
            raise TeacherLessonAuthoringError("uploaded_ppt_course_mismatch", "上传的 PPT 不属于当前课程。")
        asset, path = teacher_course_space_repository.source_file(package, body.asset_id)
        parsed = await run_in_threadpool(
            extract_uploaded_pptx_review,
            path,
            asset_id=body.asset_id,
            filename=str(asset.get("filename") or path.name),
        )
        source = _source_course(tm, course_id)
        sources, units, revisions = _imported_ppt_review_context(
            source,
            repository,
            course_id,
            lesson_unit_id,
            actor=actor,
        )
        report = build_uploaded_ppt_review_report(
            parsed["slides"], sources=sources, reference_units=units
        )
        review = repository.save_imported_ppt_review(
            course_id,
            lesson_unit_id,
            package_id=body.package_id,
            source_asset_id=body.asset_id,
            source_filename=parsed["source_filename"],
            slides=parsed["slides"],
            report=report,
            source_outline_revision_id=revisions["outline"],
            source_lesson_plan_revision_id=revisions["plan"],
            source_script_revision_id=revisions["script"],
            actor=actor,
        )
        return {"review": review}
    except (FileNotFoundError, MaterialStorageError) as exc:
        _raise(TeacherLessonAuthoringError("uploaded_ppt_asset_not_found", "上传的 PPT 原文件不存在。"))
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.get("/courses/{course_id}/lessons/{lesson_unit_id}/ppt-import/reviews/current")
async def get_current_imported_ppt_review(
    course_id: str,
    lesson_unit_id: str,
    repository: TeacherLessonAuthoringRepository = Depends(get_teacher_lesson_authoring_repository),
):
    return {"review": repository.current_imported_ppt_review(course_id, lesson_unit_id)}


@router.patch("/courses/{course_id}/lessons/{lesson_unit_id}/ppt-import/reviews/{review_id}/slides/{slide_id}")
async def update_imported_ppt_slide(
    course_id: str,
    lesson_unit_id: str,
    review_id: str,
    slide_id: str,
    body: UpdateImportedPptSlideRequest,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(get_teacher_lesson_authoring_repository),
):
    try:
        review = repository.current_imported_ppt_review(course_id, lesson_unit_id)
        if not isinstance(review, dict) or review.get("review_id") != review_id:
            raise TeacherLessonAuthoringError("uploaded_ppt_review_not_found", "PPT 审阅记录不存在。")
        slides = _updated_imported_ppt_slides(review, slide_id, body.blocks)
        source = _source_course(tm, course_id)
        sources, units, _revisions = _imported_ppt_review_context(
            source,
            repository,
            course_id,
            lesson_unit_id,
            actor=resolve_user_id(request.headers.get("X-User-Id")),
        )
        report = build_uploaded_ppt_review_report(slides, sources=sources, reference_units=units)
        updated = repository.replace_imported_ppt_review(
            course_id,
            lesson_unit_id,
            review_id=review_id,
            base_revision_id=body.base_revision_id,
            slides=slides,
            report=report,
            actor=resolve_user_id(request.headers.get("X-User-Id")),
        )
        return {"review": updated}
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.post("/courses/{course_id}/lessons/{lesson_unit_id}/ppt-import/reviews/{review_id}/ai-candidates")
@structured_generation_stream(
    stage="imported_ppt_candidate",
    started_message="已收到当前页的修改要求。",
    waiting_message="AI 正在生成 PPT 修改候选。",
)
async def create_imported_ppt_ai_candidate(
    course_id: str,
    lesson_unit_id: str,
    review_id: str,
    body: CreateImportedPptCandidateRequest,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(get_teacher_lesson_authoring_repository),
):
    try:
        review = repository.current_imported_ppt_review(course_id, lesson_unit_id)
        if not isinstance(review, dict) or review.get("review_id") != review_id:
            raise TeacherLessonAuthoringError("uploaded_ppt_review_not_found", "PPT 审阅记录不存在。")
        if review.get("revision_id") != body.base_revision_id:
            raise TeacherLessonAuthoringError("uploaded_ppt_revision_conflict", "PPT 工作稿已更新，请重新生成 AI 候选。")
        slide = next((item for item in review.get("slides") or [] if item.get("slide_id") == body.slide_id), None)
        if not isinstance(slide, dict):
            raise TeacherLessonAuthoringError("uploaded_ppt_slide_not_found", "PPT 页面不存在。")
        blocks = [item for item in slide.get("blocks") or [] if isinstance(item, dict)]
        title_block = next((item for item in blocks if item.get("kind") == "title" and item.get("editable")), None)
        body_block = next((item for item in blocks if item.get("kind") != "title" and item.get("editable")), None)
        page = {
            "page_id": body.slide_id,
            "title": str((title_block or {}).get("text") or slide.get("title") or "未命名页面"),
            "regions": ([{
                "region_id": str((body_block or {}).get("block_id") or "body"),
                "slot_id": "body",
                "content": str((body_block or {}).get("text") or ""),
            }] if body_block else []),
            "speaker_notes": "",
            "source_block_ids": [str(item.get("block_id") or "") for item in blocks],
        }
        optimized = await tm.course_service.optimize_teacher_lesson_v6_page(
            page=page,
            instruction=body.instruction,
        )
        proposed_blocks = deepcopy(blocks)
        proposed_by_id = {str(item.get("block_id") or ""): item for item in proposed_blocks}
        if title_block:
            proposed_by_id[str(title_block.get("block_id") or "")]["text"] = str(optimized["page"].get("title") or title_block.get("text") or "")
        if body_block and optimized["page"].get("key_message"):
            proposed_by_id[str(body_block.get("block_id") or "")]["text"] = str(optimized["page"]["key_message"])
        candidate = repository.save_imported_ppt_ai_candidate(
            course_id,
            lesson_unit_id,
            review_id=review_id,
            base_revision_id=body.base_revision_id,
            slide_id=body.slide_id,
            instruction=body.instruction.strip(),
            proposed_blocks=proposed_blocks,
        )
        return {"candidate": candidate}
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.post("/courses/{course_id}/lessons/{lesson_unit_id}/ppt-import/reviews/{review_id}/ai-candidates/{candidate_id}/resolve")
async def resolve_imported_ppt_ai_candidate(
    course_id: str,
    lesson_unit_id: str,
    review_id: str,
    candidate_id: str,
    body: ResolveImportedPptCandidateRequest,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(get_teacher_lesson_authoring_repository),
):
    try:
        review = repository.current_imported_ppt_review(course_id, lesson_unit_id)
        candidate = next((item for item in (review or {}).get("ai_candidates") or [] if isinstance(item, dict) and item.get("candidate_id") == candidate_id), None)
        if not isinstance(review, dict) or review.get("review_id") != review_id or not isinstance(candidate, dict):
            raise TeacherLessonAuthoringError("uploaded_ppt_candidate_not_found", "AI PPT 修改候选不存在。")
        if candidate.get("status") != "pending":
            return {"review": review}
        if body.accept:
            editable = [item for item in candidate.get("proposed_blocks") or [] if isinstance(item, dict) and item.get("editable")]
            slides = _updated_imported_ppt_slides(review, str(candidate.get("slide_id") or ""), editable)
            source = _source_course(tm, course_id)
            sources, units, _revisions = _imported_ppt_review_context(
                source,
                repository,
                course_id,
                lesson_unit_id,
                actor=resolve_user_id(request.headers.get("X-User-Id")),
            )
            report = build_uploaded_ppt_review_report(slides, sources=sources, reference_units=units)
            repository.replace_imported_ppt_review(
                course_id,
                lesson_unit_id,
                review_id=review_id,
                base_revision_id=str(candidate.get("base_revision_id") or ""),
                slides=slides,
                report=report,
                actor=resolve_user_id(request.headers.get("X-User-Id")),
            )
        repository.mark_imported_ppt_ai_candidate(
            course_id,
            lesson_unit_id,
            review_id=review_id,
            candidate_id=candidate_id,
            status="accepted" if body.accept else "rejected",
        )
        return {"review": repository.current_imported_ppt_review(course_id, lesson_unit_id)}
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.post("/courses/{course_id}/lessons/{lesson_unit_id}/ppt-import/reviews/{review_id}/confirm")
async def confirm_imported_ppt_review(
    course_id: str,
    lesson_unit_id: str,
    review_id: str,
    body: ConfirmImportedPptReviewRequest,
    repository: TeacherLessonAuthoringRepository = Depends(get_teacher_lesson_authoring_repository),
):
    try:
        return {"review": repository.confirm_imported_ppt_review(
            course_id,
            lesson_unit_id,
            review_id=review_id,
            revision_id=body.revision_id,
        )}
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.get("/courses/{course_id}/lessons/{lesson_unit_id}/ppt-import/reviews/{review_id}/export.pptx")
async def export_imported_ppt_review(
    course_id: str,
    lesson_unit_id: str,
    review_id: str,
    request: Request,
    repository: TeacherLessonAuthoringRepository = Depends(get_teacher_lesson_authoring_repository),
):
    try:
        review = repository.current_imported_ppt_review(course_id, lesson_unit_id)
        if not isinstance(review, dict) or review.get("review_id") != review_id:
            raise TeacherLessonAuthoringError("uploaded_ppt_review_not_found", "PPT 审阅记录不存在。")
        actor = resolve_user_id(request.headers.get("X-User-Id"))
        package = teacher_course_space_repository.load_owned(str(review.get("package_id") or ""), actor)
        asset, source_path = teacher_course_space_repository.source_file(package, str(review.get("source_asset_id") or ""))
        export_dir = repository.root / "exports"
        export_dir.mkdir(parents=True, exist_ok=True)
        output = export_dir / f"imported-{uuid.uuid4().hex}.pptx"

        def render() -> None:
            from pptx import Presentation

            presentation = Presentation(source_path)
            for slide_state in review.get("slides") or []:
                slide_index = int(slide_state.get("slide_number") or 0) - 1
                if slide_index < 0 or slide_index >= len(presentation.slides):
                    continue
                slide = presentation.slides[slide_index]
                for block in slide_state.get("blocks") or []:
                    shape_index = int(block.get("shape_index") or 0)
                    if not block.get("editable") or shape_index >= len(slide.shapes):
                        continue
                    shape = slide.shapes[shape_index]
                    if getattr(shape, "has_text_frame", False):
                        shape.text = str(block.get("text") or "")
            presentation.save(output)

        await run_in_threadpool(render)
        filename = f"{str(asset.get('filename') or 'PPT').rsplit('.', 1)[0]}-已审阅.pptx"
        return FileResponse(
            output,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=filename,
        )
    except (FileNotFoundError, MaterialStorageError) as exc:
        _raise(TeacherLessonAuthoringError("uploaded_ppt_asset_not_found", "上传的 PPT 原文件不存在。"))
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.get("/courses/{course_id}/lessons/{lesson_unit_id}/ppt-v6/source")
async def get_teacher_lesson_v6_source(
    course_id: str,
    lesson_unit_id: str,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        document, _course_view, _synthetic_id, _lesson, _revision = _teacher_v6_source(
            tm, repository, course_id, lesson_unit_id
        )
        return {
            "schema_version": "course_document_envelope_v1",
            "course_id": course_id,
            "course_name": document.title,
            "source_format": "canonical",
            "document": document.model_dump(mode="json"),
            "migration": {"required": False},
        }
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.get("/courses/{course_id}/lessons/{lesson_unit_id}/ppt-v6")
async def get_teacher_lesson_v6_registry(
    course_id: str,
    lesson_unit_id: str,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        _document, _course_view, synthetic_id, lesson, _revision = _teacher_v6_source(
            tm, repository, course_id, lesson_unit_id
        )
        return {"registry": _teacher_v6_registry_payload(synthetic_id)}
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.get("/courses/{course_id}/lessons/{lesson_unit_id}/ppt-v6/manuscript")
async def get_teacher_lesson_v6_manuscript(
    course_id: str,
    lesson_unit_id: str,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        document, _course_view, _synthetic_id, _lesson, _revision = (
            _teacher_v6_source(tm, repository, course_id, lesson_unit_id)
        )
        if repository.current_imported_ppt_review(course_id, lesson_unit_id):
            return {
                "ppt_manuscript_state": _ppt_manuscript_state_payload(
                    None, generation_branch="original_ppt_review"
                )
            }
        actor = resolve_user_id(request.headers.get("X-User-Id"))
        material_bindings, _material_evidence = _ppt_material_bundle(
            course_id, actor, lesson_unit_id
        )
        material_revision = stable_hash(material_bindings, prefix="pptrefs_")
        state = repository.current_v6_ppt_manuscript(
            course_id, lesson_unit_id
        )
        return {
            "ppt_manuscript_state": _ppt_manuscript_state_payload(
                state,
                generation_branch="manuscript_first",
                current_material_revision=material_revision,
            )
        }
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.patch("/courses/{course_id}/lessons/{lesson_unit_id}/ppt-v6/manuscript")
async def update_teacher_lesson_v6_manuscript_draft(
    course_id: str,
    lesson_unit_id: str,
    body: UpdateTeacherLessonPptManuscriptRequest,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        document, _course_view, _synthetic_id, _lesson, _revision = (
            _teacher_v6_source(tm, repository, course_id, lesson_unit_id)
        )
        if repository.current_imported_ppt_review(course_id, lesson_unit_id):
            raise TeacherLessonAuthoringError(
                "lesson_ppt_original_branch_active",
                "本讲已有原版 PPT，请在原版 PPT 审阅流程中处理。",
            )
        actor = resolve_user_id(request.headers.get("X-User-Id"))
        material_bindings, _material_evidence = _ppt_material_bundle(
            course_id, actor, lesson_unit_id
        )
        material_revision = stable_hash(material_bindings, prefix="pptrefs_")
        current = repository.current_v6_ppt_manuscript(
            course_id, lesson_unit_id
        )
        state_payload = _ppt_manuscript_state_payload(
            current,
            generation_branch="manuscript_first",
            current_material_revision=material_revision,
        )
        if state_payload.get("source_state") != "current":
            raise TeacherLessonAuthoringError(
                "lesson_ppt_source_stale",
                "上游教学内容或资料已经变化，请重新生成页面内容稿。",
            )
        manuscript_payload = state_payload.get("manuscript")
        if not isinstance(manuscript_payload, dict):
            raise TeacherLessonAuthoringError(
                "lesson_ppt_manuscript_not_found", "请先生成页面内容稿。"
            )
        manuscript = PptManuscriptV1.model_validate(manuscript_payload)
        revised = revise_ppt_manuscript_v1(manuscript, body.page_updates)
        template = _resolve_locked_teacher_v6_template(current or {}, actor)
        graph = compile_course_presentation_graph(document, teaching_plan={})
        compile_slide_deck_v6_from_manuscript(
            document,
            graph,
            revised,
            template,
        )
        saved = repository.update_v6_ppt_manuscript_draft(
            course_id,
            lesson_unit_id,
            expected_manuscript_revision=body.expected_manuscript_revision,
            manuscript=revised.model_dump(mode="json"),
        )
        return {
            "ppt_manuscript_state": _ppt_manuscript_state_payload(
                saved,
                generation_branch="manuscript_first",
                current_material_revision=material_revision,
            )
        }
    except TeacherLessonAuthoringError as exc:
        _raise(exc)
    except V6BuildError as exc:
        raise HTTPException(
            status_code=422,
            detail=exc.failure.model_dump(mode="json"),
        ) from exc


@router.post(
    "/courses/{course_id}/lessons/{lesson_unit_id}/ppt-v6/manuscript/regenerate-pages"
)
async def regenerate_teacher_lesson_v6_manuscript_pages(
    course_id: str,
    lesson_unit_id: str,
    body: RegenerateTeacherLessonPptManuscriptPagesRequest,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
    visual_service: TeacherScriptVisualService = Depends(
        get_teacher_script_visual_service
    ),
):
    try:
        document, _course_view, _synthetic_id, lesson, plan_revision = (
            _teacher_v6_source(tm, repository, course_id, lesson_unit_id)
        )
        if repository.current_imported_ppt_review(course_id, lesson_unit_id):
            raise TeacherLessonAuthoringError(
                "lesson_ppt_original_branch_active",
                "本讲已有原版 PPT，请在原版 PPT 审阅流程中处理。",
            )
        actor = resolve_user_id(request.headers.get("X-User-Id"))
        material_bindings, _material_evidence = _ppt_material_bundle(
            course_id, actor, lesson_unit_id
        )
        material_revision = stable_hash(material_bindings, prefix="pptrefs_")
        current = repository.current_v6_ppt_manuscript(
            course_id, lesson_unit_id
        )
        if not isinstance(current, dict) or not current:
            raise TeacherLessonAuthoringError(
                "lesson_ppt_manuscript_not_found", "请先生成页面内容稿。"
            )
        current_revision = str(current.get("revision") or "")
        if current_revision != body.expected_manuscript_revision:
            raise TeacherLessonAuthoringError(
                "lesson_ppt_manuscript_revision_conflict",
                "页面内容稿已在其他页面修改，请重新载入后再生成。",
                details={"current_revision": current_revision},
            )
        state_payload = _ppt_manuscript_state_payload(
            current,
            generation_branch="manuscript_first",
            current_material_revision=material_revision,
        )
        manuscript_payload = state_payload.get("manuscript")
        if not isinstance(manuscript_payload, dict):
            raise TeacherLessonAuthoringError(
                "lesson_ppt_manuscript_not_found", "请先生成页面内容稿。"
            )
        manuscript = PptManuscriptV1.model_validate(manuscript_payload)
        current_plan_revision_id = str(plan_revision.get("revision_id") or "")
        current_script_revision_id = str(
            lesson.get("working_script_revision_id") or ""
        )
        plan_is_current = bool(
            current_plan_revision_id
            and current.get("source_lesson_plan_revision_id")
            == current_plan_revision_id
        )
        material_is_current = bool(
            current.get("source_material_revision") == material_revision
        )
        script_changed = bool(
            current_script_revision_id
            and current.get("source_script_revision_id")
            != current_script_revision_id
        )
        source_rebase = False
        working_manuscript = manuscript
        affected_ids: list[str] = []
        locked_conflicts: list[str] = []
        if state_payload.get("source_state") != "current" or script_changed:
            if not (plan_is_current and material_is_current and script_changed):
                raise TeacherLessonAuthoringError(
                    "lesson_ppt_source_stale",
                    "教案或资料发生变化，当前页面结构不能局部重用，请重新生成整份页面内容稿。",
                )
            (
                working_manuscript,
                affected_ids,
                locked_conflicts,
            ) = rebase_ppt_manuscript_source_blocks_v1(
                manuscript,
                document,
                source_script_revision_id=current_script_revision_id,
            )
            source_rebase = True
        target_page_ids = list(dict.fromkeys([
            *body.target_page_ids,
            *affected_ids,
        ]))
        if locked_conflicts:
            raise TeacherLessonAuthoringError(
                "lesson_ppt_manuscript_locked_source_conflict",
                "已锁定页面的讲义来源已变化，请解锁后重新生成或人工处理。",
                details={"page_ids": locked_conflicts},
            )
        affected_set = set(affected_ids)
        pages_by_id = {
            page.page_id: page for page in working_manuscript.pages
        }
        ai_target_page_ids = [
            page_id for page_id in target_page_ids
            if not (
                page_id in affected_set
                and page_id in pages_by_id
                and (
                    pages_by_id[page_id].page_type in {
                        "cover", "agenda", "summary"
                    }
                    or pages_by_id[page_id].continuation_of_page_id
                )
            )
        ]

        revised = working_manuscript
        if ai_target_page_ids:
            candidate_repository = SlideDeckV6CandidateRepository(
                repository.root / "v6_candidates"
            )
            source_task_id = str(current.get("task_id") or "")
            try:
                candidate_repository.load(source_task_id)
                candidate_repository.load_checkpoint(source_task_id)
            except (FileNotFoundError, ValueError) as exc:
                raise TeacherLessonAuthoringError(
                    "lesson_ppt_manuscript_checkpoint_missing",
                    "页面内容稿的冻结规划依据不可用，请重新生成整份内容稿。",
                ) from exc

            question_bundle = question_bank_repository.load_bundle(course_id) or {}
            accepted_questions = [
                item for item in approved_formal_tasks(question_bundle)
                if lesson_unit_id in {
                    str(item.get("node_id") or ""),
                    *(str(value) for value in item.get("node_ids") or []),
                }
                and not item.get("stale_reasons")
            ]
            script_revision = next(
                (
                    item for item in lesson.get("script_revisions") or []
                    if isinstance(item, dict)
                    and str(item.get("revision_id") or "")
                    == current_script_revision_id
                ),
                {},
            )
            script_blocks = [
                block
                for section in script_revision.get("sections") or []
                if isinstance(section, dict)
                for block in section.get("blocks") or []
                if isinstance(block, dict)
            ]
            visual_view = visual_service.list_for_lesson(
                course_id=course_id,
                lesson_unit_id=lesson_unit_id,
                script_revision_id=current_script_revision_id,
                blocks=script_blocks,
            )
            target_source_ids = {
                block_id
                for page in working_manuscript.pages
                if page.page_id in set(ai_target_page_ids)
                for block_id in page.source_script_block_ids
            }
            accepted_visuals = [
                item for item in visual_view.get("items") or []
                if isinstance(item, dict)
                and item.get("status") == "accepted"
                and not item.get("stale_reasons")
                and target_source_ids.intersection({
                    str((item.get("source") or {}).get("source_block_id") or ""),
                    *(
                        str(value)
                        for value in (item.get("source") or {}).get(
                            "source_block_ids", []
                        )
                    ),
                })
            ]
            revised = await regenerate_ppt_manuscript_pages_v1(
                working_manuscript,
                target_page_ids=ai_target_page_ids,
                ai_planner=build_ai_base_story_planner_v6(),
                accepted_question_bank_items=accepted_questions,
                accepted_visual_expressions=accepted_visuals,
            )
        latest_document, _latest_view, _latest_synthetic_id, latest_lesson, latest_plan = (
            _teacher_v6_source(tm, repository, course_id, lesson_unit_id)
        )
        latest_material_bindings, _latest_material_evidence = _ppt_material_bundle(
            course_id, actor, lesson_unit_id
        )
        latest_material_revision = stable_hash(
            latest_material_bindings, prefix="pptrefs_"
        )
        if (
            latest_document.document_revision != document.document_revision
            or str(latest_plan.get("revision_id") or "")
            != current_plan_revision_id
            or str(latest_lesson.get("working_script_revision_id") or "")
            != current_script_revision_id
            or latest_material_revision != material_revision
        ):
            raise TeacherLessonAuthoringError(
                "lesson_ppt_source_stale",
                "页面生成期间教案、讲义或资料又发生了变化，已保留原内容，请基于最新内容重试。",
            )
        template = _resolve_locked_teacher_v6_template(current, actor)
        graph = compile_course_presentation_graph(document, teaching_plan={})
        compile_slide_deck_v6_from_manuscript(
            document,
            graph,
            revised,
            template,
        )
        saved = repository.update_v6_ppt_manuscript_draft(
            course_id,
            lesson_unit_id,
            expected_manuscript_revision=body.expected_manuscript_revision,
            manuscript=revised.model_dump(mode="json"),
            source_rebase=source_rebase,
            source_lesson_plan_revision_id=current_plan_revision_id,
            source_script_revision_id=current_script_revision_id,
            source_material_revision=material_revision,
        )
        return {
            "affected_page_ids": affected_ids,
            "regenerated_page_ids": ai_target_page_ids,
            "ppt_manuscript_state": _ppt_manuscript_state_payload(
                saved,
                generation_branch="manuscript_first",
                current_material_revision=material_revision,
            ),
        }
    except TeacherLessonAuthoringError as exc:
        _raise(exc)
    except V6BuildError as exc:
        raise HTTPException(
            status_code=422,
            detail=exc.failure.model_dump(mode="json"),
        ) from exc


@router.post("/courses/{course_id}/lessons/{lesson_unit_id}/ppt-v6/manuscript/confirm")
async def confirm_teacher_lesson_v6_manuscript_draft(
    course_id: str,
    lesson_unit_id: str,
    body: ConfirmTeacherLessonPptManuscriptRequest,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        _teacher_v6_source(tm, repository, course_id, lesson_unit_id)
        if repository.current_imported_ppt_review(course_id, lesson_unit_id):
            raise TeacherLessonAuthoringError(
                "lesson_ppt_original_branch_active",
                "本讲已有原版 PPT，请在原版 PPT 审阅流程中处理。",
            )
        actor = resolve_user_id(request.headers.get("X-User-Id"))
        material_bindings, _material_evidence = _ppt_material_bundle(
            course_id, actor, lesson_unit_id
        )
        material_revision = stable_hash(material_bindings, prefix="pptrefs_")
        current = repository.current_v6_ppt_manuscript(
            course_id, lesson_unit_id
        )
        state_payload = _ppt_manuscript_state_payload(
            current,
            generation_branch="manuscript_first",
            current_material_revision=material_revision,
        )
        manuscript = state_payload.get("manuscript")
        if not isinstance(manuscript, dict):
            raise TeacherLessonAuthoringError(
                "lesson_ppt_manuscript_not_found", "请先生成 页面内容稿。"
            )
        _assert_ppt_manuscript_confirmable(manuscript)
        if state_payload.get("source_state") != "current":
            raise TeacherLessonAuthoringError(
                "lesson_ppt_source_stale",
                "上游教学内容或资料已经变化，请重新生成 页面内容稿。",
            )
        confirmed = repository.confirm_v6_ppt_manuscript_draft(
            course_id,
            lesson_unit_id,
            manuscript_revision=body.manuscript_revision,
        )
        return {
            "ppt_manuscript_state": _ppt_manuscript_state_payload(
                confirmed,
                generation_branch="manuscript_first",
                current_material_revision=material_revision,
            )
        }
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.get("/courses/{course_id}/lessons/{lesson_unit_id}/ppt-v6/{representation_id}/spec")
async def get_teacher_lesson_v6_spec(
    course_id: str,
    lesson_unit_id: str,
    representation_id: str,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        _document, _course_view, synthetic_id, lesson, _revision = _teacher_v6_source(
            tm, repository, course_id, lesson_unit_id
        )
        registry = teaching_representation_repository.load(synthetic_id)
        representation = next(
            (item for item in registry.representations if item.representation_id == representation_id),
            None,
        )
        if representation is None:
            raise TeacherLessonAuthoringError("lesson_ppt_not_found", "本讲 V6 PPT 不存在。")
        spec = next((item for item in registry.specs if item.spec_id == representation.spec_id), None)
        if spec is None:
            raise TeacherLessonAuthoringError("lesson_ppt_revision_not_found", "本讲 V6 PPT 规格不存在。")
        manuscript = (spec.payload.get("content") or {}).get("ppt_manuscript") or {}
        asset = next(
            (
                item
                for item in lesson.get("ppt_assets") or []
                if isinstance(item, dict)
                and item.get("working_representation_id") == representation_id
            ),
            {},
        )
        return {
            "representation": representation.model_dump(mode="json"),
            "spec": spec.model_dump(mode="json"),
            "ai_candidate": repository.pending_v6_ppt_ai_candidate(
                course_id,
                lesson_unit_id,
                representation_id=representation_id,
                spec_id=spec.spec_id,
                spec_revision=spec.revision,
            ),
            "ppt_manuscript_state": {
                "revision": str(manuscript.get("manuscript_revision") or ""),
                "status": str(asset.get("ppt_manuscript_status") or "draft"),
                "source_state": str(asset.get("source_state") or "current"),
                "confirmable": bool(
                    manuscript.get("manuscript_revision")
                    and manuscript.get("quality_status") == "passed"
                    and asset.get("source_state") == "current"
                ),
            },
        }
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.post(
    "/courses/{course_id}/lessons/{lesson_unit_id}/ppt-v6/"
    "{representation_id}/manuscript/confirm"
)
async def confirm_teacher_lesson_v6_manuscript(
    course_id: str,
    lesson_unit_id: str,
    representation_id: str,
    body: ConfirmTeacherLessonPptManuscriptRequest,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        _document, _course_view, synthetic_id, lesson, _revision = _teacher_v6_source(
            tm, repository, course_id, lesson_unit_id
        )
        registry = teaching_representation_repository.load(synthetic_id)
        representation = next(
            (
                item
                for item in registry.representations
                if item.representation_id == representation_id
            ),
            None,
        )
        spec = next(
            (
                item
                for item in registry.specs
                if representation and item.spec_id == representation.spec_id
            ),
            None,
        )
        manuscript = (
            (spec.payload.get("content") or {}).get("ppt_manuscript")
            if spec
            else None
        )
        if not isinstance(manuscript, dict):
            raise TeacherLessonAuthoringError(
                "lesson_ppt_manuscript_not_found", "当前 PPT 没有可确认的页面内容稿。"
            )
        _assert_ppt_manuscript_confirmable(manuscript)
        if manuscript.get("manuscript_revision") != body.manuscript_revision:
            raise TeacherLessonAuthoringError(
                "lesson_ppt_manuscript_revision_conflict",
                "页面内容稿已更新，请刷新后再确认。",
            )
        asset = repository.confirm_v6_ppt_manuscript(
            course_id,
            lesson_unit_id,
            representation_id=representation_id,
            manuscript_revision=body.manuscript_revision,
        )
        return {
            "ppt_manuscript_state": {
                "revision": body.manuscript_revision,
                "status": asset.get("ppt_manuscript_status"),
                "source_state": asset.get("source_state"),
                "confirmable": True,
            }
        }
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.post(
    "/courses/{course_id}/lessons/{lesson_unit_id}/ppt-v6/{representation_id}/ai-candidates"
)
@structured_generation_stream(
    stage="ppt_page_candidate",
    started_message="已收到当前页的修改要求。",
    waiting_message="AI 正在生成 PPT 页面候选。",
)
async def create_teacher_lesson_v6_ai_candidate(
    course_id: str,
    lesson_unit_id: str,
    representation_id: str,
    body: CreateTeacherLessonV6CandidateRequest,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        _document, _course_view, synthetic_id, _lesson, _revision = _teacher_v6_source(
            tm, repository, course_id, lesson_unit_id
        )
        registry = teaching_representation_repository.load(synthetic_id)
        representation = next(
            (item for item in registry.representations if item.representation_id == representation_id),
            None,
        )
        spec = next(
            (item for item in registry.specs if representation and item.spec_id == representation.spec_id),
            None,
        )
        if representation is None or spec is None:
            raise TeacherLessonAuthoringError("lesson_ppt_not_found", "本讲 V6 PPT 不存在。")
        if spec.spec_id != body.base_spec_id or spec.revision != body.base_spec_revision:
            raise TeacherLessonAuthoringError(
                "lesson_ppt_revision_conflict", "PPT 已经变化，请基于当前页面重新优化。"
            )
        content = spec.payload.get("content") or {}
        pages = content.get("pages") if isinstance(content.get("pages"), list) else []
        page = next(
            (item for item in pages if str(item.get("page_id") or "") == body.page_id),
            None,
        )
        if not isinstance(page, dict):
            raise TeacherLessonAuthoringError("lesson_ppt_page_not_found", "当前 PPT 页面不存在。")
        optimized = await tm.course_service.optimize_teacher_lesson_v6_page(
            page=page,
            instruction=body.instruction,
        )
        candidate = repository.save_v6_ppt_ai_candidate(
            course_id,
            lesson_unit_id,
            representation_id=representation_id,
            base_spec_id=spec.spec_id,
            base_spec_revision=spec.revision,
            page_id=body.page_id,
            instruction=body.instruction.strip(),
            candidate_page=optimized["page"],
            changed_fields=list(optimized.get("changed_fields") or []),
        )
        return {"candidate": candidate}
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.post(
    "/courses/{course_id}/lessons/{lesson_unit_id}/ppt-v6/{representation_id}/ai-candidates/{candidate_id}/resolve"
)
async def resolve_teacher_lesson_v6_ai_candidate(
    course_id: str,
    lesson_unit_id: str,
    representation_id: str,
    candidate_id: str,
    body: ResolveTeacherLessonV6CandidateRequest,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        _document, course_view, synthetic_id, _lesson, plan_revision = _teacher_v6_source(
            tm, repository, course_id, lesson_unit_id
        )
        registry = teaching_representation_repository.load(synthetic_id)
        representation = next(
            (item for item in registry.representations if item.representation_id == representation_id),
            None,
        )
        spec = next(
            (item for item in registry.specs if representation and item.spec_id == representation.spec_id),
            None,
        )
        candidate = repository.pending_v6_ppt_ai_candidate(
            course_id,
            lesson_unit_id,
            representation_id=representation_id,
            spec_id=str(spec.spec_id if spec else ""),
            spec_revision=str(spec.revision if spec else ""),
        )
        if not isinstance(candidate, dict) or candidate.get("candidate_id") != candidate_id:
            raise TeacherLessonAuthoringError(
                "lesson_ppt_candidate_not_found", "AI PPT 候选不存在或已过期。"
            )
        if not body.accept:
            resolved = repository.mark_v6_ppt_ai_candidate(
                course_id, lesson_unit_id, candidate_id, status="rejected"
            )
            return {"candidate": resolved, "status": "rejected"}
        if representation is None or spec is None:
            raise TeacherLessonAuthoringError("lesson_ppt_not_found", "本讲 V6 PPT 不存在。")
        payload = deepcopy(spec.payload)
        content = payload.get("content") or {}
        pages = content.get("pages") if isinstance(content.get("pages"), list) else []
        page = next(
            (item for item in pages if str(item.get("page_id") or "") == candidate.get("page_id")),
            None,
        )
        if not isinstance(page, dict):
            raise TeacherLessonAuthoringError("lesson_ppt_page_not_found", "当前 PPT 页面不存在。")
        candidate_page = candidate.get("candidate_page") or {}
        for field in candidate.get("changed_fields") or []:
            if field in {"title", "subtitle", "key_message"}:
                _apply_v6_page_expression(
                    page,
                    field=field,
                    value=deepcopy(candidate_page.get(field)),
                    target_region_id=str(candidate_page.get(f"{field}_region_id") or ""),
                )
        manuscript = _refresh_v6_ppt_manuscript(
            content,
            course_view=course_view,
            source_lesson_plan_revision_id=str(plan_revision.get("revision_id") or ""),
        )
        now = datetime.now(timezone.utc).isoformat()
        spec_revision = stable_hash(payload, prefix="tsr_")
        edited_spec = TeachingRepresentationSpec(
            spec_id=stable_hash({
                "course_id": spec.course_id,
                "representation_type": spec.representation_type,
                "source_bindings": [item.model_dump(mode="json") for item in spec.source_bindings],
                "payload": payload,
            }, prefix="trs_"),
            course_id=spec.course_id,
            representation_type=spec.representation_type,
            source_bindings=spec.source_bindings,
            unit_bindings=spec.unit_bindings,
            payload=payload,
            revision=spec_revision,
            created_at=now,
            updated_at=now,
        )
        teaching_representation_repository.register_spec(edited_spec)
        edited_representation = representation.model_copy(deep=True)
        edited_representation.spec_id = edited_spec.spec_id
        edited_representation.semantic_fingerprint = stable_hash(content, prefix="sem_")
        edited_representation.render_fingerprint = stable_hash(
            {"spec_revision": spec_revision, "renderer": "slide_deck_v6"}, prefix="rnd_"
        )
        edited_representation.revision = stable_hash({
            "spec_revision": spec_revision,
            "source_revision_vector": edited_representation.source_revision_vector,
        }, prefix="rpr_")
        edited_representation.updated_at = now
        updated_registry = teaching_representation_repository.register_representation(
            edited_representation
        )
        repository.bind_v6_ppt_revision(
            course_id,
            lesson_unit_id,
            source_lesson_plan_revision_id=str(plan_revision.get("revision_id") or ""),
            source_script_revision_id=str(
                (course_view.get("teacher_lesson_source") or {}).get("script_revision_id") or ""
            ),
            synthetic_course_id=synthetic_id,
            representation_id=edited_representation.representation_id,
            spec_id=edited_spec.spec_id,
            candidate_status=str(content.get("status") or content.get("candidate_status") or "v6_ready"),
            ppt_manuscript_revision=str(manuscript.get("manuscript_revision") or ""),
            ppt_manuscript_status="draft",
        )
        resolved = repository.mark_v6_ppt_ai_candidate(
            course_id,
            lesson_unit_id,
            candidate_id,
            status="accepted",
            result_spec_id=edited_spec.spec_id,
        )
        return {
            "candidate": resolved,
            "status": "accepted",
            "registry": updated_registry.model_dump(mode="json"),
            "spec": edited_spec.model_dump(mode="json"),
        }
    except TeacherLessonAuthoringError as exc:
        _raise(exc)
    except ValueError as exc:
        raise HTTPException(
            status_code=422,
            detail={"code": "teacher_v6_edit_quality_blocked", "message": str(exc)},
        ) from exc


@router.post(
    "/courses/{course_id}/lessons/{lesson_unit_id}/ppt-v6/manuscript/build/stream"
)
async def build_teacher_lesson_v6_manuscript(
    course_id: str,
    lesson_unit_id: str,
    body: TeacherLessonV6BuildRequest,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        document, course_view, _synthetic_id, lesson, revision = _teacher_v6_source(
            tm, repository, course_id, lesson_unit_id
        )
        if repository.current_imported_ppt_review(course_id, lesson_unit_id):
            raise TeacherLessonAuthoringError(
                "lesson_ppt_original_branch_active",
                "本讲已有原版 PPT，请在原版 PPT 审阅流程中处理。",
            )
        actor = resolve_user_id(request.headers.get("X-User-Id"))
        template = _resolve_teacher_v6_template(body, actor)
        material_bindings, material_evidence = _ppt_material_bundle(
            course_id, actor, lesson_unit_id
        )
        document = _attach_ppt_reference_evidence(document, material_evidence)
        teacher_source = dict(course_view.get("teacher_lesson_source") or {})
        teacher_source["material_bindings"] = material_bindings
        course_view["teacher_lesson_source"] = teacher_source
        course_view["evidence_catalog"] = material_evidence
    except TeacherLessonAuthoringError as exc:
        _raise(exc)
    source_plan_revision = str(
        revision.get("revision_id") or lesson.get("working_revision_id") or ""
    )
    source_script_revision = str(
        (course_view.get("teacher_lesson_source") or {}).get("script_revision_id")
        or ""
    )
    source_material_revision = stable_hash(material_bindings, prefix="pptrefs_")
    resume_from_job_id = _teacher_ppt_resume_job_id(
        repository,
        course_id,
        lesson_unit_id,
        body.resume_task_id,
        job_type="teacher_lesson_ppt_manuscript_generation",
        source_lesson_plan_revision_id=source_plan_revision,
        source_lesson_script_revision_id=source_script_revision,
        source_material_revision=source_material_revision,
    )
    request_snapshot = {
        "lesson_unit_id": lesson_unit_id,
        "mode": body.mode,
        "theme": body.theme,
        "template_pack_id": body.template_pack_id,
        "template_version": _requested_template_version(body),
        "force_rebuild": body.force_rebuild,
        "source_lesson_plan_revision_id": source_plan_revision,
        "source_script_revision_id": source_script_revision,
        "source_material_revision": source_material_revision,
    }
    job = repository.create_job(
        course_id,
        lesson_unit_id,
        job_type="teacher_lesson_ppt_manuscript_generation",
        request_id=f"teacher-ppt-manuscript-{uuid.uuid4().hex}",
        source_outline_revision_id=str(document.document_revision or ""),
    )
    task_id = str(job["id"])
    job = repository.update_job(
        course_id,
        task_id,
        request_snapshot=request_snapshot,
        input_fingerprint=stable_hash(request_snapshot, prefix="teacher-ppt-manuscript-input"),
        resume_from_job_id=resume_from_job_id,
        source_lesson_plan_revision_id=source_plan_revision,
        source_script_revision_id=source_script_revision,
        source_material_revision=source_material_revision,
    )
    _capture_generation_source_snapshot(
        course_id=course_id,
        actor=actor,
        target_id=f"ppt-v6:{lesson_unit_id}",
        target_type="ppt",
        target_label=f"{lesson.get('node_name') or lesson_unit_id} PPT",
        target_revision=source_script_revision or source_plan_revision,
        task_id=task_id,
    )
    candidate_repository = SlideDeckV6CandidateRepository(
        repository.root / "v6_candidates"
    )
    if body.resume_task_id:
        try:
            candidate_repository.clone_checkpoint(
                body.resume_task_id,
                task_id,
            )
        except (FileNotFoundError, ValueError) as exc:
            _fail_teacher_ppt_job(
                repository,
                course_id,
                task_id,
                code="lesson_ppt_manuscript_resume_checkpoint_missing",
                message="页面内容稿的恢复检查点不可用，请重新生成。",
                retryable=False,
                phase="ppt_manuscript_resume_blocked",
            )
            raise HTTPException(
                status_code=409,
                detail={
                    "code": "lesson_ppt_manuscript_resume_checkpoint_missing",
                    "message": "页面内容稿的恢复检查点不可用，请重新生成。",
                },
            ) from exc
    orchestrator = SlideDeckV6Orchestrator(
        representation_repository=teaching_representation_repository,
        candidate_repository=candidate_repository,
        progress_root=repository.root / "v6_progress",
    )
    story_planner = build_ai_base_story_planner_v6()
    visual_planner = build_ai_base_visual_planner_v2()

    async def event_stream():
        queue: asyncio.Queue[dict[str, Any] | None] = asyncio.Queue()
        sequence = 0

        async def progress(payload: dict[str, object]) -> None:
            _teacher_ppt_job_must_be_active(repository, course_id, task_id)
            repository.update_job_live(
                course_id,
                task_id,
                status="running",
                phase=str(payload.get("stage") or "building"),
                progress=int(payload.get("percent") or 0),
                message="正在生成可逐页审阅的 页面内容稿",
            )
            await queue.put({
                "event": "slide_build_progress_v2",
                "task_id": task_id,
                "progress": int(payload.get("percent") or 0),
                "stage": str(payload.get("stage") or "building"),
                "message": "正在生成可逐页审阅的 页面内容稿",
                "slide_build_progress_v2": deepcopy(payload),
                "target_schema": "ppt_manuscript_v1",
            })

        def source_revision_provider() -> str:
            _teacher_ppt_job_must_be_active(repository, course_id, task_id)
            current = repository.lesson(course_id, lesson_unit_id)
            current_script = next(
                (
                    item for item in current.get("script_revisions") or []
                    if isinstance(item, dict)
                    and item.get("revision_id") == source_script_revision
                ),
                {},
            )
            try:
                current_bindings, _current_evidence = _ppt_material_bundle(
                    course_id, actor, lesson_unit_id
                )
                materials_current = (
                    stable_hash(current_bindings, prefix="pptrefs_")
                    == source_material_revision
                )
            except TeacherLessonAuthoringError:
                materials_current = False
            return (
                str(document.document_revision or "")
                if current.get("working_revision_id") == source_plan_revision
                and current.get("source_state", "current") == "current"
                and current.get("working_script_revision_id") == source_script_revision
                and current_script.get("source_lesson_plan_revision_id")
                == source_plan_revision
                and _script_revision_has_content(current_script)
                and materials_current
                else ""
            )

        async def run() -> None:
            try:
                started = repository.update_job(
                    course_id,
                    task_id,
                    status="running",
                    phase="ppt_manuscript_building",
                    message="正在生成可逐页审阅的 页面内容稿",
                )
                if str(started.get("status") or "") != "running":
                    raise _TeacherPptV6JobStopped(started)
                result = await orchestrator.build(
                    task_id=task_id,
                    document=document,
                    course_data=course_view,
                    mode=body.mode,
                    theme=body.theme,
                    story_planner=story_planner,
                    visual_planner=visual_planner,
                    source_revision_provider=source_revision_provider,
                    template_contract=template,
                    template_digest_provider=lambda: template.template_digest,
                    publish_result=False,
                    manuscript_only=True,
                    progress_callback=progress,
                )
                _teacher_ppt_job_must_be_active(repository, course_id, task_id)
                state = repository.save_v6_ppt_manuscript(
                    course_id,
                    lesson_unit_id,
                    source_lesson_plan_revision_id=source_plan_revision,
                    source_script_revision_id=source_script_revision,
                    source_material_revision=source_material_revision,
                    task_id=task_id,
                    mode=body.mode,
                    theme=body.theme,
                    template_id=template.template_id,
                    template_version=template.template_version,
                    template_digest=template.template_digest,
                    template_pack_id=(
                        body.template_pack_id if body.template_pack_id else ""
                    ),
                    manuscript=dict(result.get("ppt_manuscript") or {}),
                )
                completed = repository.update_job(
                    course_id,
                    task_id,
                    status="completed",
                    phase="ppt_manuscript_complete",
                    progress=100,
                    message="页面内容稿已生成",
                    stream_complete=True,
                    result_revision_id=str(state.get("revision") or ""),
                )
                await queue.put({
                    "event": "build_complete",
                    "task_id": task_id,
                    "job": completed,
                    "progress": 100,
                    "stage": "manuscript_complete",
                    "target_schema": "ppt_manuscript_v1",
                    "build": result,
                    "ppt_manuscript_state": _ppt_manuscript_state_payload(
                        state,
                        generation_branch="manuscript_first",
                        current_material_revision=source_material_revision,
                    ),
                })
            except V6BuildError as exc:
                failure = exc.failure.model_dump(mode="json")
                current = repository.get_job(course_id, task_id)
                if str(current.get("status") or "") in {"paused", "cancelled"}:
                    await queue.put(_teacher_ppt_stopped_event(current))
                    return
                failed = _fail_teacher_ppt_job(
                    repository,
                    course_id,
                    task_id,
                    code=str(failure.get("code") or "teacher_lesson_v6_manuscript_failed"),
                    message=str(failure.get("message") or "页面内容稿生成失败"),
                    retryable=bool(failure.get("retryable")),
                    phase=str(failure.get("stage") or "ppt_manuscript_failed"),
                )
                await queue.put({
                    "event": "build_failed",
                    "task_id": task_id,
                    "job": failed,
                    "progress": 100,
                    "stage": failure.get("stage") or "failed",
                    **failure,
                })
            except _TeacherPptV6JobStopped as exc:
                await queue.put(_teacher_ppt_stopped_event(exc.job))
            except Exception as exc:
                current = repository.get_job(course_id, task_id)
                if str(current.get("status") or "") in {"paused", "cancelled"}:
                    await queue.put(_teacher_ppt_stopped_event(current))
                    return
                failed = _fail_teacher_ppt_job(
                    repository,
                    course_id,
                    task_id,
                    code="teacher_lesson_v6_manuscript_failed",
                    message=str(exc),
                    retryable=True,
                    phase="ppt_manuscript_failed",
                )
                await queue.put({
                    "event": "build_failed",
                    "task_id": task_id,
                    "job": failed,
                    "progress": 100,
                    "stage": "failed",
                    "code": "teacher_lesson_v6_manuscript_failed",
                    "message": str(exc),
                    "retryable": True,
                })
            finally:
                await queue.put(None)

        sequence += 1
        queued_event = _teacher_ppt_queued_event(
            job,
            target_schema="ppt_manuscript_v1",
        )
        yield (
            f"id: {sequence}\nevent: build_queued\ndata: "
            f"{json.dumps({**queued_event, 'sequence': sequence}, ensure_ascii=False)}\n\n"
        )
        worker = asyncio.create_task(run())
        repository.track_runtime_job(course_id, worker)
        while True:
            payload = await queue.get()
            if payload is None:
                break
            sequence += 1
            name = str(payload.get("event") or "message")
            yield (
                f"id: {sequence}\nevent: {name}\ndata: "
                f"{json.dumps({**payload, 'sequence': sequence}, ensure_ascii=False)}\n\n"
            )
        await worker

    return StreamingResponse(
        event_stream(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


@router.post("/courses/{course_id}/lessons/{lesson_unit_id}/ppt-v6/build/stream")
async def build_teacher_lesson_v6(
    course_id: str,
    lesson_unit_id: str,
    body: TeacherLessonV6BuildRequest,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        document, course_view, synthetic_id, lesson, revision = _teacher_v6_source(
            tm, repository, course_id, lesson_unit_id
        )
        if repository.current_imported_ppt_review(course_id, lesson_unit_id):
            raise TeacherLessonAuthoringError(
                "lesson_ppt_original_branch_active",
                "本讲已有原版 PPT，请在原版 PPT 审阅流程中处理。",
            )
        actor = resolve_user_id(request.headers.get("X-User-Id"))
        material_bindings, material_evidence = _ppt_material_bundle(
            course_id, actor, lesson_unit_id
        )
        source_material_revision = stable_hash(
            material_bindings, prefix="pptrefs_"
        )
        manuscript_state = repository.current_v6_ppt_manuscript(
            course_id, lesson_unit_id
        )
        manuscript_state_payload = _ppt_manuscript_state_payload(
            manuscript_state,
            generation_branch="manuscript_first",
            current_material_revision=source_material_revision,
        )
        if not manuscript_state_payload.get("can_generate_ppt"):
            raise TeacherLessonAuthoringError(
                "lesson_ppt_manuscript_not_confirmed",
                "请先生成并确认当前版本的 页面内容稿，再生成 PPT。",
            )
        try:
            confirmed_manuscript = PptManuscriptV1.model_validate(
                manuscript_state_payload.get("manuscript")
            )
        except ValueError as exc:
            raise TeacherLessonAuthoringError(
                "lesson_ppt_manuscript_invalid",
                "页面内容稿结构无效，请重新生成页面内容稿。",
            ) from exc
        document = _attach_ppt_reference_evidence(document, material_evidence)
        teacher_source = dict(course_view.get("teacher_lesson_source") or {})
        teacher_source["material_bindings"] = material_bindings
        course_view["teacher_lesson_source"] = teacher_source
        course_view["evidence_catalog"] = material_evidence
    except TeacherLessonAuthoringError as exc:
        _raise(exc)
    source_plan_revision = str(
        manuscript_state.get("source_lesson_plan_revision_id") or ""
    )
    source_script_revision = str(
        manuscript_state.get("source_script_revision_id") or ""
    )
    manuscript_mode = str(manuscript_state.get("mode") or "teaching")
    manuscript_theme = str(
        manuscript_state.get("theme") or "academic-editorial"
    )
    template = _resolve_locked_teacher_v6_template(manuscript_state, actor)
    resume_from_job_id = _teacher_ppt_resume_job_id(
        repository,
        course_id,
        lesson_unit_id,
        body.resume_task_id,
        job_type="teacher_lesson_ppt_generation",
        source_lesson_plan_revision_id=source_plan_revision,
        source_lesson_script_revision_id=source_script_revision,
        source_material_revision=source_material_revision,
    )
    request_snapshot = {
        "lesson_unit_id": lesson_unit_id,
        "mode": manuscript_mode,
        "theme": manuscript_theme,
        "source_lesson_plan_revision_id": source_plan_revision,
        "source_script_revision_id": source_script_revision,
        "source_material_revision": source_material_revision,
        "force_rebuild": body.force_rebuild,
        "ppt_manuscript_task_id": str(manuscript_state.get("task_id") or ""),
        "ppt_manuscript_revision": str(
            manuscript_state.get("revision")
            or confirmed_manuscript.manuscript_revision
            or ""
        ),
    }
    job = repository.create_job(
        course_id,
        lesson_unit_id,
        job_type="teacher_lesson_ppt_generation",
        request_id=f"teacher-ppt-{uuid.uuid4().hex}",
        source_outline_revision_id=str(document.document_revision or ""),
    )
    task_id = str(job["id"])
    job = repository.update_job(
        course_id,
        task_id,
        request_snapshot=request_snapshot,
        input_fingerprint=stable_hash(request_snapshot, prefix="teacher-ppt-input"),
        resume_from_job_id=resume_from_job_id,
        source_lesson_plan_revision_id=source_plan_revision,
        source_script_revision_id=source_script_revision,
        source_material_revision=source_material_revision,
    )
    _capture_generation_source_snapshot(
        course_id=course_id,
        actor=actor,
        target_id=f"ppt-v6:{lesson_unit_id}",
        target_type="ppt",
        target_label=f"{lesson.get('node_name') or lesson_unit_id} PPT",
        target_revision=source_script_revision or source_plan_revision,
        task_id=task_id,
    )
    candidate_repository = SlideDeckV6CandidateRepository(
        repository.root / "v6_candidates"
    )
    checkpoint_task_id = body.resume_task_id or str(manuscript_state.get("task_id") or "")
    try:
        candidate_repository.clone_checkpoint(
            checkpoint_task_id, task_id
        )
    except (FileNotFoundError, ValueError) as exc:
        _fail_teacher_ppt_job(
            repository,
            course_id,
            task_id,
            code="lesson_ppt_manuscript_checkpoint_missing",
            message="页面内容稿的生成检查点不可用，请重新生成页面内容稿。",
            retryable=False,
            phase="ppt_resume_blocked",
        )
        raise HTTPException(
            status_code=409,
            detail={
                "code": "lesson_ppt_manuscript_checkpoint_missing",
                "message": "页面内容稿的生成检查点不可用，请重新生成页面内容稿。",
            },
        ) from exc
    orchestrator = SlideDeckV6Orchestrator(
        representation_repository=teaching_representation_repository,
        candidate_repository=candidate_repository,
        progress_root=repository.root / "v6_progress",
    )
    story_planner = build_ai_base_story_planner_v6()
    visual_planner = build_ai_base_visual_planner_v2()

    async def event_stream():
        queue: asyncio.Queue[dict[str, Any] | None] = asyncio.Queue()
        sequence = 0

        async def progress(payload: dict[str, object]) -> None:
            _teacher_ppt_job_must_be_active(repository, course_id, task_id)
            repository.update_job_live(
                course_id,
                task_id,
                status="running",
                phase=str(payload.get("stage") or "building"),
                progress=int(payload.get("percent") or 0),
                message="正在从已确认的 页面内容稿编译可编辑页面",
            )
            await queue.put({
                "event": "slide_build_progress_v2",
                "task_id": task_id,
                "progress": int(payload.get("percent") or 0),
                "stage": str(payload.get("stage") or "building"),
                "message": "正在从已确认的 页面内容稿编译可编辑页面",
                "slide_build_progress_v2": deepcopy(payload),
                "target_schema": "slide_deck_v6",
            })

        def source_revision_provider() -> str:
            _teacher_ppt_job_must_be_active(repository, course_id, task_id)
            current = repository.lesson(course_id, lesson_unit_id)
            current_script = next(
                (
                    item for item in current.get("script_revisions") or []
                    if isinstance(item, dict)
                    and item.get("revision_id") == source_script_revision
                ),
                {},
            )
            try:
                current_bindings, _current_evidence = _ppt_material_bundle(
                    course_id, actor, lesson_unit_id
                )
                materials_current = (
                    stable_hash(current_bindings, prefix="pptrefs_")
                    == source_material_revision
                )
            except TeacherLessonAuthoringError:
                materials_current = False
            return (
                str(document.document_revision or "")
                if current.get("working_revision_id") == source_plan_revision
                and current.get("source_state", "current") == "current"
                and current.get("working_script_revision_id") == source_script_revision
                and current_script.get("source_lesson_plan_revision_id")
                == source_plan_revision
                and _script_revision_has_content(current_script)
                and materials_current
                else ""
            )

        async def run() -> None:
            try:
                started = repository.update_job(
                    course_id,
                    task_id,
                    status="running",
                    phase="ppt_building",
                    message="正在从已确认的 页面内容稿编译可编辑页面",
                )
                if str(started.get("status") or "") != "running":
                    raise _TeacherPptV6JobStopped(started)
                result = await orchestrator.build(
                    task_id=task_id,
                    document=document,
                    course_data=course_view,
                    mode=manuscript_mode,
                    theme=manuscript_theme,
                    story_planner=story_planner,
                    visual_planner=visual_planner,
                    source_revision_provider=source_revision_provider,
                    template_contract=template,
                    template_digest_provider=lambda: template.template_digest,
                    publish_result=True,
                    confirmed_manuscript=confirmed_manuscript,
                    progress_callback=progress,
                )
                _teacher_ppt_job_must_be_active(repository, course_id, task_id)
                repository.bind_v6_ppt_revision(
                    course_id,
                    lesson_unit_id,
                    source_lesson_plan_revision_id=source_plan_revision,
                    source_script_revision_id=source_script_revision,
                    synthetic_course_id=synthetic_id,
                    representation_id=str(result.get("representation_id") or ""),
                    spec_id=str(result.get("spec_id") or ""),
                    candidate_status=str(result.get("candidate_status") or result.get("status") or ""),
                    ppt_manuscript_revision=str(
                        result.get("ppt_manuscript_revision") or ""
                    ),
                    ppt_manuscript_status="confirmed",
                )
                repository.bind_v6_ppt_manuscript_result(
                    course_id,
                    lesson_unit_id,
                    manuscript_revision=confirmed_manuscript.manuscript_revision,
                    representation_id=str(result.get("representation_id") or ""),
                )
                completed = repository.update_job(
                    course_id,
                    task_id,
                    status="completed",
                    phase="ppt_complete",
                    progress=100,
                    message="PPT 已生成",
                    stream_complete=True,
                    result_revision_id=str(result.get("representation_id") or ""),
                )
                await queue.put({
                    "event": "build_complete",
                    "task_id": task_id,
                    "job": completed,
                    "progress": 100,
                    "stage": "complete",
                    "target_schema": "slide_deck_v6",
                    "quality": result.get("quality") or {},
                    "build": result,
                    "registry": _teacher_v6_registry_payload(synthetic_id),
                })
            except V6BuildError as exc:
                failure = exc.failure.model_dump(mode="json")
                current = repository.get_job(course_id, task_id)
                if str(current.get("status") or "") in {"paused", "cancelled"}:
                    await queue.put(_teacher_ppt_stopped_event(current))
                    return
                failed = _fail_teacher_ppt_job(
                    repository,
                    course_id,
                    task_id,
                    code=str(failure.get("code") or "teacher_lesson_v6_failed"),
                    message=str(failure.get("message") or "PPT 生成失败"),
                    retryable=bool(failure.get("retryable")),
                    phase=str(failure.get("stage") or "ppt_failed"),
                )
                await queue.put({
                    "event": "build_failed",
                    "task_id": task_id,
                    "job": failed,
                    "progress": 100,
                    "stage": failure.get("stage") or "failed",
                    **failure,
                })
            except _TeacherPptV6JobStopped as exc:
                await queue.put(_teacher_ppt_stopped_event(exc.job))
            except Exception as exc:
                current = repository.get_job(course_id, task_id)
                if str(current.get("status") or "") in {"paused", "cancelled"}:
                    await queue.put(_teacher_ppt_stopped_event(current))
                    return
                failed = _fail_teacher_ppt_job(
                    repository,
                    course_id,
                    task_id,
                    code="teacher_lesson_v6_failed",
                    message=str(exc),
                    retryable=True,
                    phase="ppt_failed",
                )
                await queue.put({
                    "event": "build_failed",
                    "task_id": task_id,
                    "job": failed,
                    "progress": 100,
                    "stage": "failed",
                    "code": "teacher_lesson_v6_failed",
                    "message": str(exc),
                    "retryable": True,
                })
            finally:
                await queue.put(None)

        sequence += 1
        queued_event = _teacher_ppt_queued_event(
            job,
            target_schema="slide_deck_v6",
        )
        yield (
            f"id: {sequence}\nevent: build_queued\ndata: "
            f"{json.dumps({**queued_event, 'sequence': sequence}, ensure_ascii=False)}\n\n"
        )
        worker = asyncio.create_task(run())
        repository.track_runtime_job(course_id, worker)
        while True:
            payload = await queue.get()
            if payload is None:
                break
            sequence += 1
            name = str(payload.get("event") or "message")
            yield f"id: {sequence}\nevent: {name}\ndata: {json.dumps({**payload, 'sequence': sequence}, ensure_ascii=False)}\n\n"
        await worker

    return StreamingResponse(
        event_stream(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


@router.get("/courses/{course_id}/lessons/{lesson_unit_id}/ppt-v6/{representation_id}/export.pptx")
async def export_teacher_lesson_v6(
    course_id: str,
    lesson_unit_id: str,
    representation_id: str,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        _document, _course_view, synthetic_id, lesson, _revision = _teacher_v6_source(
            tm, repository, course_id, lesson_unit_id
        )
        registry = teaching_representation_repository.load(synthetic_id)
        representation = next(
            (item for item in registry.representations if item.representation_id == representation_id),
            None,
        )
        spec = next(
            (item for item in registry.specs if representation and item.spec_id == representation.spec_id),
            None,
        )
        if representation is None or spec is None:
            raise TeacherLessonAuthoringError("lesson_ppt_not_found", "本讲 V6 PPT 不存在。")
        manuscript = (spec.payload.get("content") or {}).get("ppt_manuscript")
        if isinstance(manuscript, dict):
            asset = next(
                (
                    item
                    for item in lesson.get("ppt_assets") or []
                    if isinstance(item, dict)
                    and item.get("working_representation_id") == representation_id
                ),
                {},
            )
            if (
                asset.get("ppt_manuscript_status") != "confirmed"
                or asset.get("ppt_manuscript_revision")
                != manuscript.get("manuscript_revision")
            ):
                raise TeacherLessonAuthoringError(
                    "lesson_ppt_manuscript_not_confirmed",
                    "请先确认 页面内容稿，再导出正式 PPTX。",
                )
        output = repository.root / "exports" / f"{synthetic_id}-{representation_id}-{spec.revision}.pptx"
        output.parent.mkdir(parents=True, exist_ok=True)
        content = spec.payload.get("content") or {}
        export_theme = str(content.get("theme") or "academic-editorial")
        await run_in_threadpool(
            export_slide_deck_pptx,
            spec,
            output,
            theme=export_theme,
        )
        return FileResponse(
            output,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=f"{lesson_unit_id}-V6课堂课件.pptx",
        )
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.post("/courses/{course_id}/lessons/{lesson_unit_id}/ppt-v6/{representation_id}/edits/preview")
async def preview_teacher_lesson_v6_edit(
    course_id: str,
    lesson_unit_id: str,
    representation_id: str,
    body: TeacherLessonRepresentationEditRequest,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        _document, _course_view, synthetic_id, _lesson, _revision = _teacher_v6_source(
            tm, repository, course_id, lesson_unit_id
        )
        registry = teaching_representation_repository.load(synthetic_id)
        representation = next((item for item in registry.representations if item.representation_id == representation_id), None)
        spec = next((item for item in registry.specs if representation and item.spec_id == representation.spec_id), None)
        if representation is None or spec is None:
            raise TeacherLessonAuthoringError("lesson_ppt_not_found", "本讲 V6 PPT 不存在。")
        classification = classify_representation_edit(
            field=body.field,
            before=body.before,
            after=body.after,
            semantic_intent=body.semantic_intent,
        )
        impact = representation_edit_impact(
            registry,
            spec,
            unit_id=body.unit_id,
            field=body.field,
            semantic_change=classification.get("semantic_change"),
        )
        return {"status": "preview", **classification, "impact": impact}
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.post("/courses/{course_id}/lessons/{lesson_unit_id}/ppt-v6/{representation_id}/edits/apply")
async def apply_teacher_lesson_v6_edit(
    course_id: str,
    lesson_unit_id: str,
    representation_id: str,
    body: TeacherLessonApplyRepresentationEditRequest,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    if body.decision != "representation_only":
        raise HTTPException(
            status_code=422,
            detail={
                "code": "teacher_semantic_edit_requires_lesson_plan",
                "message": "语义修改请回到本讲教案；PPT 工作台只保存表达层修改。",
            },
        )
    try:
        _document, _course_view, synthetic_id, _lesson, _revision = _teacher_v6_source(
            tm, repository, course_id, lesson_unit_id
        )
        registry = teaching_representation_repository.load(synthetic_id)
        representation = next((item for item in registry.representations if item.representation_id == representation_id), None)
        spec = next((item for item in registry.specs if representation and item.spec_id == representation.spec_id), None)
        if representation is None or spec is None:
            raise TeacherLessonAuthoringError("lesson_ppt_not_found", "本讲 V6 PPT 不存在。")
        payload = deepcopy(spec.payload)
        content = payload.get("content") or {}
        pages = content.get("pages") if isinstance(content.get("pages"), list) else None
        if pages is None or content.get("schema_version") != "slide_deck_v6":
            raise HTTPException(
                status_code=422,
                detail={"code": "teacher_v6_edit_unsupported", "message": "当前 V6 规格不支持页面编辑。"},
            )
        page = next(
            (item for item in pages if str(item.get("page_id") or "") == body.unit_id),
            None,
        )
        if page is None:
            raise HTTPException(status_code=404, detail="V6 page not found")
        if body.field not in {"title", "subtitle", "key_message"}:
            raise HTTPException(
                status_code=422,
                detail={"code": "teacher_v6_edit_field_unsupported", "message": "当前字段请通过原 V6 专用编辑器处理。"},
            )
        _apply_v6_page_expression(
            page,
            field=body.field,
            value=deepcopy(body.after),
        )
        manuscript = _refresh_v6_ppt_manuscript(
            content,
            course_view=_course_view,
            source_lesson_plan_revision_id=str(_revision.get("revision_id") or ""),
        )
        now = datetime.now(timezone.utc).isoformat()
        spec_revision = stable_hash(payload, prefix="tsr_")
        edited_spec = TeachingRepresentationSpec(
            spec_id=stable_hash({
                "course_id": spec.course_id,
                "representation_type": spec.representation_type,
                "source_bindings": [item.model_dump(mode="json") for item in spec.source_bindings],
                "payload": payload,
            }, prefix="trs_"),
            course_id=spec.course_id,
            representation_type=spec.representation_type,
            source_bindings=spec.source_bindings,
            unit_bindings=spec.unit_bindings,
            payload=payload,
            revision=spec_revision,
            created_at=now,
            updated_at=now,
        )
        teaching_representation_repository.register_spec(edited_spec)
        edited_representation = representation.model_copy(deep=True)
        edited_representation.spec_id = edited_spec.spec_id
        edited_representation.semantic_fingerprint = stable_hash(content, prefix="sem_")
        edited_representation.render_fingerprint = stable_hash({
            "spec_revision": spec_revision,
            "renderer": "slide_deck_v6",
        }, prefix="rnd_")
        edited_representation.revision = stable_hash({
            "spec_revision": spec_revision,
            "source_revision_vector": edited_representation.source_revision_vector,
        }, prefix="rpr_")
        edited_representation.updated_at = now
        updated = teaching_representation_repository.register_representation(edited_representation)
        repository.bind_v6_ppt_revision(
            course_id,
            lesson_unit_id,
            source_lesson_plan_revision_id=str(_revision.get("revision_id") or ""),
            source_script_revision_id=str(
                (_course_view.get("teacher_lesson_source") or {}).get("script_revision_id") or ""
            ),
            synthetic_course_id=synthetic_id,
            representation_id=edited_representation.representation_id,
            spec_id=edited_spec.spec_id,
            candidate_status=str(content.get("status") or content.get("candidate_status") or "v6_ready"),
            ppt_manuscript_revision=str(manuscript.get("manuscript_revision") or ""),
            ppt_manuscript_status="draft",
        )
        return {"status": "applied_to_representation", "registry": updated.model_dump(mode="json")}
    except TeacherLessonAuthoringError as exc:
        _raise(exc)
    except ValueError as exc:
        raise HTTPException(
            status_code=422,
            detail={"code": "teacher_v6_edit_quality_blocked", "message": str(exc)},
        ) from exc


@router.get("/courses/{course_id}/knowledge-evidence")
async def get_lesson_knowledge_evidence(
    course_id: str,
    lesson_unit_id: str = "",
    tm: TaskManager = Depends(require_task_manager),
):
    try:
        source = _source_course(tm, course_id)
        nodes = [item for item in source.get("nodes") or [] if isinstance(item, dict)]
        if lesson_unit_id:
            scope = lesson_scope(source, lesson_unit_id)
            section_ids = {str(item.get("node_id") or "") for item in scope["sections"]}
            nodes = [item for item in nodes if str(item.get("node_id") or "") in section_ids]
        points: list[dict[str, Any]] = []
        for node in nodes:
            for group in node.get("knowledge_structure") or []:
                if not isinstance(group, dict):
                    continue
                for point in group.get("knowledge_points") or []:
                    if not isinstance(point, dict):
                        continue
                    sources = point.get("source_refs") or point.get("evidence_refs") or []
                    if isinstance(sources, str):
                        sources = [sources]
                    points.append({
                        "section_node_id": str(node.get("node_id") or ""),
                        "section_title": str(node.get("node_name") or ""),
                        "name": str(point.get("name") or ""),
                        "statement": str(point.get("statement") or point.get("description") or ""),
                        "sources": [str(item) for item in sources if str(item).strip()],
                        "conflict": bool(point.get("conflict") or point.get("needs_manual_review")),
                    })
        return {
            "schema_version": "teacher_lesson_knowledge_evidence_v1",
            "course_id": course_id,
            "lesson_unit_id": lesson_unit_id,
            "points": points,
            "conflict_count": sum(1 for item in points if item["conflict"]),
        }
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.post("/courses/{course_id}/lessons/{lesson_unit_id}/plan/generate", status_code=202)
async def generate_lesson_plan(
    course_id: str,
    lesson_unit_id: str,
    body: GenerateLessonPlanRequest,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        _validate_new_attempt(repository, course_id, lesson_unit_id, "teacher_lesson_plan_generation", body)
        source = _source_course(tm, course_id)
        scope = lesson_scope(source, lesson_unit_id)
        outline_revision = _canonical_outline_revision(source)
        if (
            body.batch_source_revision_id
            and body.batch_source_revision_id != outline_revision
        ):
            raise TeacherLessonAuthoringError(
                "lesson_plan_batch_source_changed",
                "批量生成期间课程大纲已变化，请重新发起生成。",
            )
        source_evidence: list[dict[str, Any]] = []
        source_filename = ""
        primary_source_kind = ""
        primary_material_asset_id = ""
        actor = resolve_user_id(request.headers.get("X-User-Id"))
        previous: dict[str, Any] | None = None
        effective_source_package_id = body.source_package_id
        effective_source_asset_id = body.source_asset_id
        effective_requirements = body.requirements.strip()
        effective_material_asset_ids = list(body.material_asset_ids)
        if body.resume_job_id:
            previous = _validated_teacher_asset_resume_job(
                repository,
                course_id=course_id,
                resume_job_id=body.resume_job_id,
                lesson_unit_id=lesson_unit_id,
                job_type="teacher_lesson_plan_generation",
                source_revision_field="source_outline_revision_id",
                source_revision_id=outline_revision,
            )
            snapshot = previous.get("request_snapshot")
            snapshot = snapshot if isinstance(snapshot, dict) else previous
            effective_source_package_id = str(
                snapshot.get("source_package_id") or ""
            )
            effective_source_asset_id = str(
                snapshot.get("source_asset_id") or ""
            )
            effective_requirements = str(
                snapshot.get("requirements") or ""
            ).strip()
            effective_material_asset_ids = [
                str(item)
                for item in snapshot.get("material_asset_ids") or []
                if str(item)
            ]
        if bool(effective_source_package_id) != bool(effective_source_asset_id):
            raise TeacherLessonAuthoringError(
                "lesson_primary_source_incomplete",
                "主来源信息不完整。",
            )
        if effective_source_package_id and effective_source_asset_id:
            try:
                package = teacher_course_space_repository.load_owned(
                    effective_source_package_id,
                    actor,
                )
                source_asset, source_path = teacher_course_space_repository.source_file(
                    package,
                    effective_source_asset_id,
                )
            except (FileNotFoundError, MaterialStorageError) as exc:
                raise TeacherLessonAuthoringError(
                    "lesson_primary_source_not_found",
                    "主来源不存在或无权访问。",
                ) from exc
            if str(package.get("course_id") or "") != course_id:
                raise TeacherLessonAuthoringError(
                    "lesson_primary_source_course_mismatch",
                    "主来源不属于当前课程。",
                )
            source_filename = str(source_asset.get("filename") or "")
            extension = str(
                source_asset.get("extension") or source_path.suffix or ""
            ).lower()
            if not extension.startswith(".") and extension:
                extension = f".{extension}"
            if extension == ".pptx":
                primary_source_kind = "uploaded_ppt"
                source_evidence = await run_in_threadpool(
                    extract_uploaded_pptx_evidence,
                    source_path,
                    asset_id=effective_source_asset_id,
                )
            elif extension in {".docx", ".pdf", ".md", ".markdown", ".txt"}:
                primary_source_kind = "uploaded_lesson_plan"
                primary_material_asset_id = str(
                    source_asset.get("material_asset_id") or ""
                )
                try:
                    material = (
                        material_repository.get_asset(primary_material_asset_id)
                        if primary_material_asset_id
                        else None
                    )
                    document = (
                        await parse_material_asset(material_repository, material)
                        if material is not None
                        else await parse_document_path(
                            source_path,
                            asset_id=effective_source_asset_id,
                            filename=source_filename or source_path.name,
                        )
                    )
                except Exception as exc:
                    raise TeacherLessonAuthoringError(
                        "lesson_primary_source_parse_failed",
                        "原教案解析失败，请检查文件后重试。",
                    ) from exc
                if document.parse_status not in {"parsed", "degraded"} or not document.blocks:
                    raise TeacherLessonAuthoringError(
                        "lesson_primary_source_parse_failed",
                        "原教案没有提取到可用于生成的内容。",
                    )
                source_evidence = compile_original_lesson_plan_evidence(
                    document,
                    asset_id=(primary_material_asset_id or effective_source_asset_id),
                    filename=source_filename or source_path.name,
                    sections=scope["sections"],
                )
            else:
                raise TeacherLessonAuthoringError(
                    "lesson_primary_source_unsupported",
                    "主来源暂时支持 DOCX、PDF、Markdown、TXT 或 PPTX。",
                )
        selected_material_ids, selected_evidence = _course_material_evidence(
            course_id, actor, effective_material_asset_ids
        )
        source_evidence.extend(selected_evidence)
        repository.set_outline(course_id, outline_revision)
        arrangement = repository.current_arrangement(course_id, lesson_unit_id)
        if arrangement is None:
            arrangement = recommend_lesson_arrangement(
                source,
                lesson_unit_id,
                source_outline_revision_id=outline_revision,
            )
            repository.save_arrangement_revision(
                course_id,
                lesson_unit_id,
                arrangement,
                source_outline_revision_id=outline_revision,
                actor=actor,
            )
            arrangement = repository.current_arrangement(course_id, lesson_unit_id)
        arrangement_issues = validate_lesson_arrangement(
            arrangement or {},
            expected_section_ids=[
                str(item.get("node_id") or "") for item in scope["sections"]
            ],
        )
        if arrangement is None or arrangement_issues:
            raise TeacherLessonAuthoringError(
                "lesson_arrangement_invalid",
                "本讲教学结构不完整，请调整后重试。",
                details={"blocking_issues": arrangement_issues},
            )
        input_fingerprint = stable_hash({
            "lesson_unit_id": lesson_unit_id,
            "source_outline_revision_id": outline_revision,
            "source_package_id": effective_source_package_id,
            "source_asset_id": effective_source_asset_id,
            "requirements": effective_requirements,
            "material_asset_ids": sorted(selected_material_ids),
            "arrangement": arrangement,
        }, prefix="teacher-lesson-plan-input")
        resume_checkpoint: dict[str, Any] = {}
        if previous:
            if previous.get("input_fingerprint") == input_fingerprint:
                resume_checkpoint = deepcopy(previous.get("checkpoint") or {})
        job = repository.create_job(
            course_id,
            lesson_unit_id,
            request_id=body.request_id,
            source_outline_revision_id=outline_revision,
        )
        job = repository.update_job(
            course_id,
            str(job["id"]),
            input_fingerprint=input_fingerprint,
            retry_of_job_id=body.retry_of_job_id,
            attempt_mode="revised_inputs" if body.retry_of_job_id else "resume_original" if body.resume_job_id else "initial",
            resume_from_job_id=(body.resume_job_id if resume_checkpoint else ""),
            requirements=effective_requirements,
            material_asset_ids=selected_material_ids,
            request_snapshot={
                "source_outline_revision_id": outline_revision,
                "source_package_id": effective_source_package_id,
                "source_asset_id": effective_source_asset_id,
                "requirements": effective_requirements,
                "material_asset_ids": selected_material_ids,
            },
            **({
                "parent_job_id": body.batch_parent_job_id,
                "batch_position": body.batch_position,
                "batch_size": body.batch_size,
                "phase": "queued",
                "message": f"第 {body.batch_position} 讲教案已入队",
            } if body.batch_parent_job_id else {}),
        )
        if source_evidence:
            job_source_kind = (
                "mixed_course_sources"
                if effective_source_asset_id and selected_material_ids
                else primary_source_kind
                if effective_source_asset_id
                else "course_materials"
            )
            job = repository.update_job(
                course_id,
                str(job["id"]),
                source_asset_id=(effective_source_asset_id or selected_material_ids[0]),
                source_package_id=effective_source_package_id,
                source_filename=(
                    source_filename
                    or f"{len(selected_material_ids)} 份课程资料"
                ),
                source_kind=job_source_kind,
                source_material_asset_id=primary_material_asset_id,
            )
        if job.get("status") in {"running", "completed", "completed_with_warnings"}:
            return {"job": job}

        _capture_generation_source_snapshot(
            course_id=course_id,
            actor=actor,
            target_id=f"lesson-plan:{lesson_unit_id}",
            target_type="lesson_plan",
            target_label=f"{scope['lesson'].get('node_name') or lesson_unit_id} 教案",
            target_revision=outline_revision,
            task_id=str(job.get("id") or ""),
        )

        service = TeacherLessonAuthoringService(repository)

        async def planner(
            course: dict[str, Any],
            lesson_id: str,
            on_progress,
        ) -> dict[str, Any]:
            scoped_course = deepcopy(course)
            normalized_requirements = effective_requirements
            if normalized_requirements:
                scoped_course["requirements"] = normalized_requirements
                scoped_course.setdefault("metadata", {}).setdefault(
                    "teacher_lesson_requirements", {}
                )[lesson_id] = normalized_requirements
                for plan_key in ("course_plan", "course_outline"):
                    scoped_plan = scoped_course.get(plan_key)
                    if not isinstance(scoped_plan, dict):
                        continue
                    for chapter in scoped_plan.get("chapters") or []:
                        if not isinstance(chapter, dict):
                            continue
                        if not chapter_matches_lesson(
                            scoped_plan,
                            chapter,
                            lesson_id,
                        ):
                            continue
                        chapter["teacher_requirements"] = normalized_requirements
                        for section in chapter.get("sections") or []:
                            if isinstance(section, dict):
                                section["teacher_requirements"] = normalized_requirements
            async def persist_checkpoint(checkpoint: dict[str, Any]) -> None:
                await asyncio.to_thread(
                    repository.update_job,
                    course_id,
                    str(job["id"]),
                    checkpoint=checkpoint,
                )

            return await tm.course_service.prepare_teacher_lesson_plan(
                course_data=scoped_course,
                lesson_unit_id=lesson_id,
                on_phase=on_progress,
                source_evidence=source_evidence,
                lesson_arrangement=arrangement,
                resume_checkpoint=resume_checkpoint,
                on_checkpoint=persist_checkpoint,
            )

        async def run() -> None:
            async def run_current_lesson() -> None:
                async def repair_generated_plan(*, plan, issues):
                    return await tm.course_service.optimize_teacher_lesson_plan(
                        plan=plan,
                        instruction=(
                            "这是尚未交付的模型教案。请修复以下质量问题后返回完整候选，"
                            "只改必要的教学表达与活动，不改小节、教学块身份、顺序、时间、知识事实和资料来源。"
                            + json.dumps(issues, ensure_ascii=False)
                        ),
                        lesson_context={"lesson_unit_id": lesson_unit_id, "requirements": effective_requirements},
                        material_evidence=source_evidence,
                    )

                await service.run_plan_job(
                    course_id=course_id,
                    lesson_unit_id=lesson_unit_id,
                    job_id=str(job["id"]),
                    course_data=source,
                    planner=planner,
                    repairer=repair_generated_plan,
                )

            await _run_lesson_plan_job(
                course_id=course_id,
                job_id=str(job["id"]),
                repository=repository,
                run=run_current_lesson,
            )

        task = asyncio.create_task(run())
        repository.track_runtime_job(course_id, task)
        _background_jobs.add(task)
        task.add_done_callback(_background_jobs.discard)
        return {"job": {**job, "actor": actor}}
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.post("/courses/{course_id}/lesson-plans/generate-all", status_code=202)
async def generate_all_lesson_plans(
    course_id: str,
    body: GenerateAllLessonPlansRequest,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    """Queue every lecture at once; each child reports its own real status."""
    try:
        source = _source_course(tm, course_id)
        outline_revision = _canonical_outline_revision(source)
        lessons = _lesson_projection(source, repository)
        if not lessons:
            raise TeacherLessonAuthoringError(
                "teacher_lesson_batch_empty",
                "课程大纲中还没有可生成教案的讲次。",
            )
        actor = resolve_user_id(request.headers.get("X-User-Id"))
        skipped_lessons: list[dict[str, str]] = []
        target_lessons: list[dict[str, Any]] = []
        lesson_ids = {
            str(lesson.get("lesson_unit_id") or "")
            for lesson in lessons
            if str(lesson.get("lesson_unit_id") or "")
        }
        prior_jobs = list((repository.view(course_id).get("jobs") or {}).values())
        resume_jobs_by_lesson: dict[str, dict[str, Any]] = {}
        for resume_job_id in dict.fromkeys(body.resume_job_ids):
            try:
                candidate = repository.get_job(course_id, resume_job_id)
            except TeacherLessonAuthoringError as exc:
                raise TeacherLessonAuthoringError(
                    "teacher_asset_resume_conflict",
                    "要恢复的教案任务不存在，请刷新状态后重试。",
                    details={"resume_job_id": resume_job_id, "reason": "job_not_found"},
                ) from exc
            lesson_unit_id = str(candidate.get("lesson_unit_id") or "")
            if lesson_unit_id not in lesson_ids:
                raise TeacherLessonAuthoringError(
                    "teacher_asset_resume_conflict",
                    "要恢复的教案任务已不属于当前课程讲次，请刷新状态后重试。",
                    details={"resume_job_id": resume_job_id, "reason": "lesson_mismatch"},
                )
            validated = _validated_teacher_asset_resume_job(
                repository,
                course_id=course_id,
                resume_job_id=resume_job_id,
                lesson_unit_id=lesson_unit_id,
                job_type="teacher_lesson_plan_generation",
                source_revision_field="source_outline_revision_id",
                source_revision_id=outline_revision,
            )
            if lesson_unit_id in resume_jobs_by_lesson:
                raise TeacherLessonAuthoringError(
                    "teacher_asset_resume_conflict",
                    "同一讲次不能同时恢复多个教案任务。",
                    details={"resume_job_id": resume_job_id, "reason": "duplicate_lesson"},
                )
            resume_jobs_by_lesson[lesson_unit_id] = validated
        for lesson in lessons:
            lesson_unit_id = str(lesson.get("lesson_unit_id") or "")
            plan = lesson.get("plan") or {}
            latest_plan_job = _latest_teacher_asset_job(
                prior_jobs,
                lesson_unit_id,
                "teacher_lesson_plan_generation",
            )
            latest_plan_status = str((latest_plan_job or {}).get("status") or "")
            if resume_jobs_by_lesson:
                if lesson_unit_id in resume_jobs_by_lesson:
                    if not bool(plan.get("can_generate")):
                        raise TeacherLessonAuthoringError(
                            "teacher_asset_resume_conflict",
                            "要恢复的教案任务当前不再满足生成条件，请刷新状态后重试。",
                            details={
                                "resume_job_id": resume_jobs_by_lesson[lesson_unit_id]["id"],
                                "reason": "generation_unavailable",
                            },
                        )
                    target_lessons.append(lesson)
                else:
                    skipped_lessons.append({
                        "lesson_unit_id": lesson_unit_id,
                        "reason": "not_selected_for_resume",
                    })
            elif bool(plan.get("ready")):
                skipped_lessons.append({
                    "lesson_unit_id": lesson_unit_id,
                    "reason": "already_ready",
                })
            elif latest_plan_status not in {"", "cancelled", "canceled"}:
                skipped_lessons.append({
                    "lesson_unit_id": lesson_unit_id,
                    "reason": "existing_job_requires_explicit_action",
                })
            elif bool(plan.get("can_generate")):
                target_lessons.append(lesson)
            else:
                skipped_lessons.append({
                    "lesson_unit_id": lesson_unit_id,
                    "reason": str(
                        plan.get("generation_unavailable_reason")
                        or "lesson_arrangement_unavailable"
                    ),
                })
        skipped_lesson_ids = [
            item["lesson_unit_id"] for item in skipped_lessons
        ]
        material_scopes: dict[str, dict[str, Any]] = {}
        arrangements: dict[str, dict[str, Any]] = {}
        # Freeze and validate the whole launch set before any child job exists.
        for lesson in target_lessons:
            lesson_unit_id = str(lesson.get("lesson_unit_id") or "")
            material_scopes[lesson_unit_id] = _lesson_plan_material_scope(
                course_id,
                actor,
                lesson_unit_id,
            )
            arrangement = lesson.get("arrangement")
            if not isinstance(arrangement, dict) or not arrangement.get("blocks"):
                raise TeacherLessonAuthoringError(
                    "lesson_arrangement_unavailable",
                    f"{lesson.get('title') or lesson_unit_id} 暂时无法形成教学结构。",
                )
            arrangement_issues = validate_lesson_arrangement(
                arrangement,
                expected_section_ids=[
                    str(item.get("section_node_id") or "")
                    for item in lesson.get("sections") or []
                ],
            )
            if arrangement_issues:
                raise TeacherLessonAuthoringError(
                    "lesson_arrangement_invalid",
                    f"{lesson.get('title') or lesson_unit_id} 的教学结构不完整。",
                    details={"blocking_issues": arrangement_issues},
                )
            if repository.current_arrangement(course_id, lesson_unit_id) is None:
                await run_in_threadpool(
                    repository.save_arrangement_revision,
                    course_id,
                    lesson_unit_id,
                    arrangement,
                    source_outline_revision_id=outline_revision,
                    actor=actor,
                )
            arrangements[lesson_unit_id] = deepcopy(arrangement)
        parent_job_id = f"tlj-batch-{uuid.uuid4().hex}"
        request_prefix = body.request_id.strip() or parent_job_id
        jobs: list[dict[str, Any]] = []
        for batch_position, lesson in enumerate(target_lessons, start=1):
            lesson_unit_id = str(lesson.get("lesson_unit_id") or "")
            material_scope = material_scopes[lesson_unit_id]
            arrangement = arrangements[lesson_unit_id]
            child_body = GenerateLessonPlanRequest(
                request_id=f"{request_prefix}-{lesson_unit_id}",
                resume_job_id=str(
                    (resume_jobs_by_lesson.get(lesson_unit_id) or {}).get("id")
                    or ""
                ),
                source_package_id=str(material_scope["source_package_id"]),
                source_asset_id=str(material_scope["source_asset_id"]),
                requirements=body.requirements,
                material_asset_ids=list(material_scope["material_asset_ids"]),
                batch_parent_job_id=parent_job_id,
                batch_position=batch_position,
                batch_size=len(target_lessons),
                batch_source_revision_id=outline_revision,
            )
            result = await generate_lesson_plan(
                course_id,
                lesson_unit_id,
                child_body,
                request,
                tm,
                repository,
            )
            job = deepcopy(result.get("job") or {})
            if job.get("id"):
                job = await run_in_threadpool(
                    repository.update_job,
                    course_id,
                    str(job["id"]),
                    parent_job_id=parent_job_id,
                    batch_position=batch_position,
                    batch_size=len(target_lessons),
                )
                jobs.append(job)
        parent_job = {
            "id": parent_job_id,
            "course_id": course_id,
            "type": "teacher_lesson_plan_batch",
            "status": "running" if jobs else "completed",
            "child_job_ids": [str(item.get("id") or "") for item in jobs],
            "skipped_lesson_ids": skipped_lesson_ids,
            "skipped_lessons": skipped_lessons,
            "total": len(lessons),
            "started": len(jobs),
            "lesson_statuses": [
                {
                    "lesson_id": str(item.get("lesson_id") or item.get("lesson_unit_id") or ""),
                    "status": str(item.get("status") or "pending"),
                    "stage": str(item.get("stage") or item.get("phase") or "queued"),
                    "message": str(item.get("message") or ""),
                    "job_id": str(item.get("id") or ""),
                }
                for item in jobs
            ],
            "created_at": datetime.now(timezone.utc).isoformat(),
        }
        return {
            "parent_job": parent_job,
            "jobs": jobs,
            "skipped_lesson_ids": skipped_lesson_ids,
            "skipped_lessons": skipped_lessons,
        }
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.get("/courses/{course_id}/lesson-jobs/{job_id}")
async def get_lesson_job(
    course_id: str,
    job_id: str,
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        job = await run_in_threadpool(repository.get_job, course_id, job_id)
        if str(job.get("status") or "") in {"pending", "running"}:
            job = await run_in_threadpool(
                repository.expire_stale_job,
                course_id,
                job_id,
            )
        return {"job": job}
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.delete("/courses/{course_id}/lesson-jobs/{job_id}")
async def cancel_lesson_job(
    course_id: str,
    job_id: str,
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        job = await run_in_threadpool(repository.cancel_job, course_id, job_id)
        return {"job": job}
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.post("/courses/{course_id}/lesson-jobs/{job_id}/pause")
async def pause_lesson_job(
    course_id: str,
    job_id: str,
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        job = await run_in_threadpool(repository.pause_job, course_id, job_id)
        return {"job": job}
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.get("/courses/{course_id}/lesson-jobs/{job_id}/stream")
async def stream_lesson_job(
    course_id: str,
    job_id: str,
    request: Request,
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    """Stream the in-memory working copy while durable saves stay semantic."""
    try:
        await run_in_threadpool(repository.get_job, course_id, job_id)
    except TeacherLessonAuthoringError as exc:
        _raise(exc)

    async def event_stream():
        last_sequence = -1
        last_updated_at = ""
        while True:
            if await request.is_disconnected():
                return
            try:
                job = await run_in_threadpool(
                    repository.expire_stale_job,
                    course_id,
                    job_id,
                )
            except TeacherLessonAuthoringError:
                payload = {
                    "event": "error",
                    "job_id": job_id,
                    "message": "本讲生成任务不存在或已被清理。",
                }
                yield f"event: error\ndata: {json.dumps(payload, ensure_ascii=False)}\n\n"
                return
            sequence = int(job.get("stream_sequence") or 0)
            updated_at = str(job.get("updated_at") or "")
            status = str(job.get("status") or "")
            terminal = status in {
                "completed",
                "completed_with_warnings",
                "failed",
                "cancelled",
                "paused",
            }
            if sequence > last_sequence or updated_at != last_updated_at:
                last_sequence = sequence
                last_updated_at = updated_at
                script_job = str(job.get("type") or "") == "teacher_lesson_script_generation"
                event = (
                    "lesson_script_complete"
                    if script_job and status in {"completed", "completed_with_warnings"}
                    else "lesson_script_cancelled"
                    if script_job and status == "cancelled"
                    else "lesson_script_paused"
                    if script_job and status == "paused"
                    else "lesson_script_failed"
                    if script_job and status == "failed"
                    else "lesson_script_stream"
                    if script_job
                    else "lesson_plan_complete"
                    if status in {"completed", "completed_with_warnings"}
                    else "lesson_plan_cancelled"
                    if status == "cancelled"
                    else "lesson_plan_paused"
                    if status == "paused"
                    else "lesson_plan_failed"
                    if status == "failed"
                    else "lesson_plan_stream"
                )
                payload = {
                    "event": event,
                    "job": {
                        "id": job_id,
                        "schema_version": str(job.get("schema_version") or ""),
                        "course_id": course_id,
                        "lesson_unit_id": str(job.get("lesson_unit_id") or ""),
                        "type": str(job.get("type") or ""),
                        "status": status,
                        "phase": str(job.get("phase") or ""),
                        "progress": int(job.get("progress") or 0),
                        "message": str(job.get("message") or ""),
                        "warnings": deepcopy(job.get("warnings") or []),
                        "error": deepcopy(job.get("error")),
                        "result_revision_id": str(job.get("result_revision_id") or ""),
                        "stream_sequence": sequence,
                        "stream_batches": deepcopy(job.get("stream_batches") or {}),
                        "stream_mode": str(job.get("stream_mode") or ""),
                        "stream_events": deepcopy(job.get("stream_events") or []),
                        "last_stream_event": deepcopy(job.get("last_stream_event") or {}),
                        "stream_complete": bool(job.get("stream_complete")),
                        "checkpoint": deepcopy(job.get("checkpoint") or {}),
                        "cancel_requested": bool(job.get("cancel_requested")),
                        "pause_requested": bool(job.get("pause_requested")),
                        "parent_job_id": str(job.get("parent_job_id") or ""),
                        "batch_position": int(job.get("batch_position") or 0),
                        "batch_size": int(job.get("batch_size") or 0),
                        "retryable": bool(job.get("retryable")),
                        "heartbeat_at": str(job.get("heartbeat_at") or ""),
                        "requirements": str(job.get("requirements") or ""),
                        "total_blocks": int(job.get("total_blocks") or 0),
                        "completed_blocks": int(job.get("completed_blocks") or 0),
                        "current_block_id": str(job.get("current_block_id") or ""),
                        "current_block_title": str(job.get("current_block_title") or ""),
                        "block_states": deepcopy(job.get("block_states") or {}),
                        "result_sections": deepcopy(job.get("result_sections") or []),
                        "updated_at": updated_at,
                    },
                }
                if script_job:
                    last_stream_event = job.get("last_stream_event") or {}
                    if isinstance(last_stream_event, dict):
                        payload.update({
                            "lesson_unit_id": str(
                                last_stream_event.get("lesson_unit_id")
                                or job.get("lesson_unit_id")
                                or ""
                            ),
                            "block_id": str(
                                last_stream_event.get("block_id") or ""
                            ),
                            "shard_id": str(
                                last_stream_event.get("shard_id") or ""
                            ),
                            "sequence": int(
                                last_stream_event.get("sequence") or sequence
                            ),
                            "delta": str(last_stream_event.get("delta") or ""),
                        })
                yield (
                    f"id: {sequence}\n"
                    f"event: {event}\n"
                    f"data: {json.dumps(payload, ensure_ascii=False)}\n\n"
                )
            if terminal:
                return
            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_stream(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


@router.patch("/courses/{course_id}/lessons/{lesson_unit_id}/plan/draft")
async def save_lesson_plan_draft(
    course_id: str,
    lesson_unit_id: str,
    body: SaveLessonPlanDraftRequest,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        source = _source_course(tm, course_id)
        canonical_outline_revision = _canonical_outline_revision(source)
        if canonical_outline_revision:
            repository.set_outline(course_id, canonical_outline_revision)
        TeacherLessonAuthoringService(repository).save_plan_draft(
            course_id=course_id,
            lesson_unit_id=lesson_unit_id,
            course_data=source,
            plan=body.plan,
            source_outline_revision_id=body.source_outline_revision_id,
            actor=resolve_user_id(request.headers.get("X-User-Id")),
            expected_current_revision_id=body.expected_current_revision_id,
        )
        projected = next(
            item for item in _lesson_projection(source, repository)
            if item["lesson_unit_id"] == lesson_unit_id
        )
        return {"lesson": projected}
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.post(
    "/courses/{course_id}/lessons/{lesson_unit_id}/script/generate",
    status_code=202,
)
async def generate_lesson_script(
    course_id: str,
    lesson_unit_id: str,
    body: GenerateLessonScriptRequest,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        _validate_new_attempt(repository, course_id, lesson_unit_id, "teacher_lesson_script_generation", body)
        source = _source_course(tm, course_id)
        scope = lesson_scope(source, lesson_unit_id)
        lesson, plan_revision = _current_plan_revision(
            repository,
            course_id,
            lesson_unit_id,
        )
        plan_revision_id = str(plan_revision.get("revision_id") or "")
        if (
            body.batch_source_revision_id
            and body.batch_source_revision_id != plan_revision_id
        ):
            raise TeacherLessonAuthoringError(
                "lesson_script_batch_source_changed",
                "批量生成期间本讲教案已变化，请重新发起生成。",
            )
        expected_plan_section_ids = [
            str(item.get("node_id") or "") for item in scope["sections"]
        ]
        if not _plan_revision_covers_sections(
            plan_revision,
            expected_plan_section_ids,
        ):
            raise TeacherLessonAuthoringError(
                "lesson_plan_scope_stale",
                "当前教案没有完整对应本讲大纲，请重新生成或编辑保存。",
            )
        plan_sections = {
            str(item.get("node_id") or ""): item
            for item in (plan_revision.get("plan") or {}).get("sections") or []
            if isinstance(item, dict) and item.get("node_id")
        }
        actor = resolve_user_id(request.headers.get("X-User-Id"))
        previous: dict[str, Any] | None = None
        effective_requirements = body.requirements.strip()
        effective_material_asset_ids = list(body.material_asset_ids)
        if body.resume_job_id:
            previous = _validated_teacher_asset_resume_job(
                repository,
                course_id=course_id,
                resume_job_id=body.resume_job_id,
                lesson_unit_id=lesson_unit_id,
                job_type="teacher_lesson_script_generation",
                source_revision_field="source_lesson_plan_revision_id",
                source_revision_id=plan_revision_id,
            )
            # Resume means continue the frozen task.  A temporary failure
            # while reloading the source tray must not silently change the
            # prompt, sources, fingerprint, or reusable checkpoint.
            snapshot = previous.get("request_snapshot")
            snapshot = snapshot if isinstance(snapshot, dict) else previous
            effective_requirements = str(
                snapshot.get("requirements") or ""
            ).strip()
            effective_material_asset_ids = [
                str(item)
                for item in snapshot.get("material_asset_ids") or []
                if str(item)
            ]
        selected_material_ids, source_evidence = _course_material_evidence(
            course_id, actor, effective_material_asset_ids
        )
        prompt_evidence = _prompt_material_evidence(source_evidence)
        register = getattr(tm.course_service, "register_course_generation_metadata", None)
        if callable(register):
            register(course_id, source)
        input_fingerprint = stable_hash({
            "lesson_unit_id": lesson_unit_id,
            "source_lesson_plan_revision_id": plan_revision_id,
            "requirements": effective_requirements,
            "material_asset_ids": sorted(selected_material_ids),
        }, prefix="teacher-script-input")
        seed_sections: list[dict[str, Any]] = []
        if previous:
            if (
                previous.get("input_fingerprint") == input_fingerprint
                and previous.get("source_lesson_plan_revision_id") == plan_revision_id
            ):
                seed_sections = [
                    deepcopy(item)
                    for item in previous.get("result_sections") or []
                    if isinstance(item, dict)
                ]

        job = repository.create_job(
            course_id,
            lesson_unit_id,
            job_type="teacher_lesson_script_generation",
            request_id=body.request_id,
            source_outline_revision_id=_canonical_outline_revision(source),
        )
        job = repository.update_job(
            course_id,
            str(job["id"]),
            source_lesson_plan_revision_id=plan_revision_id,
            input_fingerprint=input_fingerprint,
            retry_of_job_id=body.retry_of_job_id,
            attempt_mode="revised_inputs" if body.retry_of_job_id else "resume_original" if body.resume_job_id else "initial",
            resume_from_job_id=(body.resume_job_id if seed_sections else ""),
            requirements=effective_requirements,
            material_asset_ids=selected_material_ids,
            actor=actor,
            request_snapshot={
                "source_lesson_plan_revision_id": plan_revision_id,
                "requirements": effective_requirements,
                "material_asset_ids": selected_material_ids,
            },
            **({
                "parent_job_id": body.batch_parent_job_id,
                "batch_position": body.batch_position,
                "batch_size": body.batch_size,
                "phase": "queued",
                "message": f"第 {body.batch_position} 讲讲义已入队",
            } if body.batch_parent_job_id else {}),
        )
        if job.get("status") in {"running", "completed", "completed_with_warnings"}:
            return {"job": job}

        _capture_generation_source_snapshot(
            course_id=course_id,
            actor=actor,
            target_id=f"script:{lesson_unit_id}",
            target_type="script",
            target_label=f"{scope['lesson'].get('node_name') or lesson_unit_id} 讲义",
            target_revision=plan_revision_id,
            task_id=str(job.get("id") or ""),
        )

        lesson_title = str(scope["lesson"].get("node_name") or "")
        lesson_section_titles = [
            str(item.get("node_name") or "") for item in scope["sections"]
        ]

        async def generate_block(
            outline_section: dict[str, Any],
            current_plan: dict[str, Any],
            module: dict[str, Any],
            shard_context: dict[str, Any],
            on_content_delta=None,
            on_content_reset=None,
        ) -> str:
            module_id = str(module.get("module_id") or "")
            single_outline = deepcopy(outline_section)
            single_outline["module_plan"] = [{
                **deepcopy(module),
                "label": str(module.get("title") or module_id),
            }]
            single_plan = deepcopy(current_plan)
            single_plan["teaching_modules"] = [{
                **deepcopy(module),
                "label": str(module.get("title") or module_id),
            }]
            stream_prefix = {"resolved": False, "buffer": ""}

            async def forward_stream_reset():
                stream_prefix.update({"resolved": False, "buffer": ""})
                if on_content_reset:
                    await on_content_reset()

            async def forward_stream_delta(delta: str):
                if not on_content_delta or not str(delta or ""):
                    return
                if stream_prefix["resolved"]:
                    await on_content_delta(str(delta))
                    return
                stream_prefix["buffer"] += str(delta)
                if "\n" not in stream_prefix["buffer"]:
                    return
                first_line, remainder = stream_prefix["buffer"].split("\n", 1)
                stream_prefix["resolved"] = True
                stream_prefix["buffer"] = ""
                visible = (
                    remainder.lstrip("\n")
                    if first_line.strip().startswith("## ")
                    else f"{first_line}\n{remainder}"
                )
                if visible:
                    await on_content_delta(visible)
            try:
                generated = await tm.course_service.generate_teacher_script_section(
                    course_id=course_id,
                    outline_section=single_outline,
                    current_plan_section=single_plan,
                    lesson_context={
                        "lesson_title": lesson_title,
                        "lesson_sections": lesson_section_titles,
                        "current_block": {
                            "block_id": module.get("block_id"),
                            "module_id": module_id,
                            "title": module.get("title"),
                            "role": module.get("role"),
                        },
                        "script_shard_context": deepcopy(shard_context),
                        "material_asset_ids": selected_material_ids,
                        "selected_material_evidence": prompt_evidence,
                    },
                    requirements=effective_requirements,
                    user_id=actor,
                    on_content_delta=forward_stream_delta,
                    on_content_reset=forward_stream_reset,
                )
            except (
                asyncio.TimeoutError,
                AIProviderRequestError,
                AIProviderUnavailable,
            ) as exc:
                error_detail = str(exc).strip()
                timeout_failed = isinstance(
                    exc,
                    (TeacherScriptGenerationTimeout, asyncio.TimeoutError),
                )
                quality_failed = error_detail.startswith(
                    "讲义未通过当前教案的质量检查"
                )
                raise TeacherLessonAuthoringError(
                    (
                        "lesson_script_model_timeout"
                        if timeout_failed
                        else "lesson_script_block_quality_failed"
                        if quality_failed
                        else "lesson_script_provider_failed"
                    ),
                    (
                        f"{module.get('title') or module_id}模型调用超时，请重试。"
                        if timeout_failed
                        else f"{module.get('title') or module_id}未通过硬校验，请重试。"
                        if quality_failed
                        else f"{module.get('title') or module_id}生成失败，请重试。"
                    ),
                    details={
                        "reason": (
                            error_detail or "讲义模型调用超时"
                        )[:1000]
                    },
                ) from exc
            blocks = [
                item for item in generated.get("blocks") or [] if isinstance(item, dict)
            ]
            content = str((blocks[0] if blocks else {}).get("content") or "").strip()
            if not content:
                raise TeacherLessonAuthoringError(
                    "lesson_script_block_empty",
                    f"{module.get('title') or module_id} 没有生成有效内容，请重试。",
                )
            return content

        async def generate_script_shard(
            entries: list[dict[str, Any]],
            shard_context: dict[str, Any],
            *,
            on_block_delta,
            on_shard_reset,
        ) -> dict[str, str]:
            if not entries:
                return {}
            modules = [deepcopy(entry["module"]) for entry in entries]
            shard_id = str(shard_context.get("shard_id") or uuid.uuid4().hex)
            synthetic_node_id = f"{lesson_unit_id}:{shard_id}"
            combined_outline = deepcopy(entries[0]["outline_section"])
            combined_outline.update({
                "node_id": synthetic_node_id,
                "node_name": lesson_title or "当前讲次",
                "learning_objective": "；".join(dict.fromkeys(
                    str(entry["contract"].get("learning_objective") or "")
                    for entry in entries
                    if str(entry["contract"].get("learning_objective") or "")
                )),
                "module_plan": [
                    {**module, "label": str(module.get("title") or "教学块")}
                    for module in modules
                ],
            })
            combined_plan = deepcopy(entries[0]["plan_section"])
            combined_plan.update({
                "node_id": synthetic_node_id,
                "title": lesson_title or "当前讲次",
                "learning_objective": combined_outline["learning_objective"],
                "key_points": list(dict.fromkeys(
                    value
                    for entry in entries
                    for value in entry["contract"].get("key_points") or []
                    if value
                )),
                "key_difficulties": list(dict.fromkeys(
                    value
                    for entry in entries
                    for value in entry["contract"].get("key_difficulties") or []
                    if value
                )),
                "teaching_modules": [
                    {**module, "label": str(module.get("title") or "教学块")}
                    for module in modules
                ],
            })
            block_ids = [str(module.get("block_id") or "") for module in modules]
            stream_parser = {"buffer": "", "block_index": -1}

            async def forward_shard_reset():
                stream_parser.update({"buffer": "", "block_index": -1})
                await on_shard_reset()

            async def emit_current(value: str):
                index = int(stream_parser["block_index"])
                if value and 0 <= index < len(block_ids):
                    await on_block_delta(block_ids[index], value)

            async def forward_shard_delta(delta: str):
                stream_parser["buffer"] += str(delta or "")
                while stream_parser["buffer"]:
                    buffer = str(stream_parser["buffer"])
                    if buffer.startswith("## "):
                        newline = buffer.find("\n")
                        if newline < 0:
                            return
                        stream_parser["block_index"] = min(
                            int(stream_parser["block_index"]) + 1,
                            len(block_ids) - 1,
                        )
                        stream_parser["buffer"] = buffer[newline + 1:]
                        continue
                    if int(stream_parser["block_index"]) < 0:
                        heading = buffer.find("## ")
                        if heading < 0:
                            return
                        stream_parser["buffer"] = buffer[heading:]
                        continue
                    boundary = buffer.find("\n## ")
                    if boundary >= 0:
                        await emit_current(buffer[:boundary])
                        stream_parser["buffer"] = buffer[boundary + 1:]
                        continue
                    last_newline = buffer.rfind("\n")
                    if last_newline >= 0 and buffer[last_newline + 1:].startswith("#"):
                        await emit_current(buffer[:last_newline + 1])
                        stream_parser["buffer"] = buffer[last_newline + 1:]
                        return
                    await emit_current(buffer)
                    stream_parser["buffer"] = ""

            try:
                generated = await tm.course_service.generate_teacher_script_section(
                    course_id=course_id,
                    outline_section=combined_outline,
                    current_plan_section=combined_plan,
                    lesson_context={
                        "lesson_title": lesson_title,
                        "lesson_sections": lesson_section_titles,
                        "script_shard_context": deepcopy(shard_context),
                        "material_asset_ids": selected_material_ids,
                        "selected_material_evidence": prompt_evidence,
                    },
                    requirements=effective_requirements,
                    user_id=actor,
                    on_content_delta=forward_shard_delta,
                    on_content_reset=forward_shard_reset,
                    allow_partial_quality=True,
                )
            except (
                asyncio.TimeoutError,
                AIProviderRequestError,
                AIProviderUnavailable,
            ) as exc:
                error_detail = str(exc).strip()
                timeout_failed = isinstance(
                    exc,
                    (TeacherScriptGenerationTimeout, asyncio.TimeoutError),
                )
                raise TeacherLessonAuthoringError(
                    (
                        "lesson_script_model_timeout"
                        if timeout_failed
                        else "lesson_script_block_quality_failed"
                        if error_detail.startswith("讲义未通过当前教案的质量检查")
                        else "lesson_script_provider_failed"
                    ),
                    (
                        f"讲义分片模型调用超时：{error_detail or '请重试'}"
                        if timeout_failed
                        else f"讲义分片生成失败：{error_detail or '请重试'}"
                    ),
                    details={
                        "reason": (
                            error_detail or "讲义模型调用超时"
                        )[:1000],
                        "shard_id": shard_id,
                    },
                ) from exc
            generated_blocks = [
                item for item in generated.get("blocks") or []
                if isinstance(item, dict)
            ]
            if len(generated_blocks) != len(block_ids):
                raise TeacherLessonAuthoringError(
                    "lesson_script_shard_incomplete",
                    "讲义分片返回的教学块数量与教案不一致。",
                )
            return {
                block_id: str(generated_blocks[index].get("content") or "").strip()
                for index, block_id in enumerate(block_ids)
            }

        async def run() -> None:
            try:
                await TeacherLessonAuthoringService(repository).run_script_job(
                    course_id=course_id,
                    lesson_unit_id=lesson_unit_id,
                    job_id=str(job["id"]),
                    source_plan_revision_id=plan_revision_id,
                    outline_sections=scope["sections"],
                    plan_sections=plan_sections,
                    generator=generate_block,
                    shard_generator=generate_script_shard,
                    repair_generator=generate_block,
                    seed_sections=seed_sections,
                    requirements=effective_requirements,
                    material_asset_ids=selected_material_ids,
                    actor=actor,
                )
            except asyncio.CancelledError:
                raise
            except Exception as exc:
                current = repository.get_job(course_id, str(job["id"]))
                if current.get("status") not in {
                    "completed", "completed_with_warnings", "failed",
                }:
                    repository.update_job(
                        course_id,
                        str(job["id"]),
                        status="failed",
                        phase="lesson_script_failed",
                        message="本讲讲义生成失败",
                        stream_sequence=int(current.get("stream_sequence") or 0) + 1,
                        stream_complete=True,
                        error={
                            "code": "lesson_script_generation_failed",
                            "message": str(exc),
                            "retryable": True,
                        },
                    )

        task = asyncio.create_task(run())
        repository.track_runtime_job(course_id, task)
        _background_jobs.add(task)
        task.add_done_callback(_background_jobs.discard)
        return {"job": job}
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.post("/courses/{course_id}/lesson-scripts/generate-all", status_code=202)
async def generate_all_lesson_scripts(
    course_id: str,
    body: GenerateAllLessonScriptsRequest,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    """Queue all current lesson scripts while keeping each lesson independent."""
    try:
        source = _source_course(tm, course_id)
        lessons = _lesson_projection(source, repository)
        if not lessons:
            raise TeacherLessonAuthoringError(
                "teacher_lesson_script_batch_empty",
                "课程大纲中还没有可生成讲义的讲次。",
            )
        actor = resolve_user_id(request.headers.get("X-User-Id"))
        skipped_lessons: list[dict[str, str]] = []
        target_lessons: list[dict[str, Any]] = []
        lessons_by_id = {
            str(lesson.get("lesson_unit_id") or ""): lesson
            for lesson in lessons
            if str(lesson.get("lesson_unit_id") or "")
        }
        prior_jobs = list((repository.view(course_id).get("jobs") or {}).values())
        resume_jobs_by_lesson: dict[str, dict[str, Any]] = {}
        plan_revision_ids: dict[str, str] = {}
        for resume_job_id in dict.fromkeys(body.resume_job_ids):
            try:
                candidate = repository.get_job(course_id, resume_job_id)
            except TeacherLessonAuthoringError as exc:
                raise TeacherLessonAuthoringError(
                    "teacher_asset_resume_conflict",
                    "要恢复的讲义任务不存在，请刷新状态后重试。",
                    details={"resume_job_id": resume_job_id, "reason": "job_not_found"},
                ) from exc
            lesson_unit_id = str(candidate.get("lesson_unit_id") or "")
            if lesson_unit_id not in lessons_by_id:
                raise TeacherLessonAuthoringError(
                    "teacher_asset_resume_conflict",
                    "要恢复的讲义任务已不属于当前课程讲次，请刷新状态后重试。",
                    details={"resume_job_id": resume_job_id, "reason": "lesson_mismatch"},
                )
            _lesson, plan_revision = _current_plan_revision(
                repository,
                course_id,
                lesson_unit_id,
            )
            plan_revision_id = str(plan_revision.get("revision_id") or "")
            validated = _validated_teacher_asset_resume_job(
                repository,
                course_id=course_id,
                resume_job_id=resume_job_id,
                lesson_unit_id=lesson_unit_id,
                job_type="teacher_lesson_script_generation",
                source_revision_field="source_lesson_plan_revision_id",
                source_revision_id=plan_revision_id,
            )
            if lesson_unit_id in resume_jobs_by_lesson:
                raise TeacherLessonAuthoringError(
                    "teacher_asset_resume_conflict",
                    "同一讲次不能同时恢复多个讲义任务。",
                    details={"resume_job_id": resume_job_id, "reason": "duplicate_lesson"},
                )
            resume_jobs_by_lesson[lesson_unit_id] = validated
            plan_revision_ids[lesson_unit_id] = plan_revision_id
        for lesson in lessons:
            lesson_unit_id = str(lesson.get("lesson_unit_id") or "")
            script = lesson.get("script") or {}
            latest_script_job = _latest_teacher_asset_job(
                prior_jobs,
                lesson_unit_id,
                "teacher_lesson_script_generation",
            )
            latest_script_status = str(
                (latest_script_job or {}).get("status") or ""
            )
            if resume_jobs_by_lesson:
                if lesson_unit_id in resume_jobs_by_lesson:
                    if not bool(script.get("can_generate")):
                        raise TeacherLessonAuthoringError(
                            "teacher_asset_resume_conflict",
                            "要恢复的讲义任务当前不再满足生成条件，请刷新状态后重试。",
                            details={
                                "resume_job_id": resume_jobs_by_lesson[lesson_unit_id]["id"],
                                "reason": "generation_unavailable",
                            },
                        )
                    target_lessons.append(lesson)
                else:
                    skipped_lessons.append({
                        "lesson_unit_id": lesson_unit_id,
                        "reason": "not_selected_for_resume",
                    })
            elif bool(script.get("ready")):
                skipped_lessons.append({
                    "lesson_unit_id": lesson_unit_id,
                    "reason": "already_ready",
                })
            elif latest_script_status not in {"", "cancelled", "canceled"}:
                skipped_lessons.append({
                    "lesson_unit_id": lesson_unit_id,
                    "reason": "existing_job_requires_explicit_action",
                })
            elif bool(script.get("can_generate")):
                target_lessons.append(lesson)
            else:
                skipped_lessons.append({
                    "lesson_unit_id": lesson_unit_id,
                    "reason": str(
                        script.get("generation_unavailable_reason")
                        or "lesson_plan_not_ready"
                    ),
                })
        skipped_lesson_ids = [
            item["lesson_unit_id"] for item in skipped_lessons
        ]

        # Validate the entire launch set before creating any child.  A stale
        # plan must not leave the teacher with a half-enqueued batch.
        material_scopes: dict[str, dict[str, Any]] = {}
        for lesson in target_lessons:
            lesson_unit_id = str(lesson.get("lesson_unit_id") or "")
            if lesson_unit_id not in plan_revision_ids:
                _lesson, plan_revision = _current_plan_revision(
                    repository,
                    course_id,
                    lesson_unit_id,
                )
                plan_revision_ids[lesson_unit_id] = str(
                    plan_revision.get("revision_id") or ""
                )
            material_scopes[lesson_unit_id] = _lesson_script_material_scope(
                course_id,
                actor,
                lesson_unit_id,
            )

        parent_job_id = f"tls-batch-{uuid.uuid4().hex}"
        request_prefix = body.request_id.strip() or parent_job_id
        jobs: list[dict[str, Any]] = []
        for batch_position, lesson in enumerate(target_lessons, start=1):
            lesson_unit_id = str(lesson.get("lesson_unit_id") or "")
            material_scope = material_scopes[lesson_unit_id]
            resume_job_id = str(
                (resume_jobs_by_lesson.get(lesson_unit_id) or {}).get("id")
                or ""
            )
            child_body = GenerateLessonScriptRequest(
                request_id=f"{request_prefix}-{lesson_unit_id}",
                resume_job_id=resume_job_id,
                requirements=body.requirements,
                material_asset_ids=list(material_scope["material_asset_ids"]),
                batch_parent_job_id=parent_job_id,
                batch_position=batch_position,
                batch_size=len(target_lessons),
                batch_source_revision_id=plan_revision_ids[lesson_unit_id],
            )
            result = await generate_lesson_script(
                course_id,
                lesson_unit_id,
                child_body,
                request,
                tm,
                repository,
            )
            job = deepcopy(result.get("job") or {})
            if job.get("id"):
                jobs.append(job)

        parent_status = (
            "completed"
            if not jobs or all(
                str(job.get("status") or "")
                in {"completed", "completed_with_warnings"}
                for job in jobs
            )
            else "running"
        )
        child_job_ids = [str(item.get("id") or "") for item in jobs]
        parent_job = {
            "id": parent_job_id,
            "course_id": course_id,
            "type": "teacher_lesson_script_batch",
            "status": parent_status,
            "child_job_ids": child_job_ids,
            "skipped_lesson_ids": skipped_lesson_ids,
            "skipped_lessons": skipped_lessons,
            "total": len(lessons),
            "started": len(jobs),
            "lesson_statuses": [
                {
                    "lesson_id": str(
                        item.get("lesson_id") or item.get("lesson_unit_id") or ""
                    ),
                    "status": str(item.get("status") or "pending"),
                    "stage": str(item.get("stage") or item.get("phase") or "queued"),
                    "message": str(item.get("message") or ""),
                    "job_id": str(item.get("id") or ""),
                }
                for item in jobs
            ],
            "created_at": datetime.now(timezone.utc).isoformat(),
        }
        return {
            "parent_job": parent_job,
            "jobs": jobs,
            "child_job_ids": child_job_ids,
            "skipped_lesson_ids": skipped_lesson_ids,
            "skipped_lessons": skipped_lessons,
        }
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.get("/courses/{course_id}/lessons/{lesson_unit_id}/script/visuals")
async def get_lesson_script_visuals(
    course_id: str,
    lesson_unit_id: str,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
    visual_service: TeacherScriptVisualService = Depends(
        get_teacher_script_visual_service
    ),
):
    try:
        resolve_user_id(request.headers.get("X-User-Id"))
        revision, blocks = _current_script_visual_context(
            tm, repository, course_id, lesson_unit_id
        )
        return await run_in_threadpool(
            visual_service.list_for_lesson,
            course_id=course_id,
            lesson_unit_id=lesson_unit_id,
            script_revision_id=str(revision.get("revision_id") or ""),
            blocks=blocks,
        )
    except (TeacherLessonAuthoringError, RepresentationConflict) as exc:
        if isinstance(exc, TeacherLessonAuthoringError):
            _raise(exc)
        _raise(TeacherLessonAuthoringError(
            "lesson_script_visual_conflict", str(exc)
        ))


@router.post("/courses/{course_id}/lessons/{lesson_unit_id}/script/visuals")
async def create_lesson_script_visual(
    course_id: str,
    lesson_unit_id: str,
    body: CreateLessonScriptVisualRequest,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
    visual_service: TeacherScriptVisualService = Depends(
        get_teacher_script_visual_service
    ),
):
    try:
        resolve_user_id(request.headers.get("X-User-Id"))
        revision, blocks = _current_script_visual_context(
            tm, repository, course_id, lesson_unit_id
        )
        revision_id = str(revision.get("revision_id") or "")
        if body.script_revision_id != revision_id:
            raise TeacherLessonAuthoringError(
                "lesson_script_revision_conflict",
                "讲义工作稿已经变化，请基于当前教学块重新生成视觉表达。",
            )
        block = next(
            (
                item for item in blocks
                if str(item.get("section_node_id") or "") == body.section_node_id
                and str(item.get("block_id") or "") == body.block_id
            ),
            None,
        )
        if block is None:
            raise TeacherLessonAuthoringError(
                "lesson_script_visual_block_not_found",
                "当前讲义教学块不存在，请重新载入。",
            )
        if body.expression_type == "animation" and not script_animation_runtime_enabled():
            raise TeacherLessonAuthoringError(
                "lesson_script_animation_gray_disabled",
                "教学动画仍处于灰度储备阶段，当前课程不参与实际运行。",
            )
        if body.expression_type == "diagram":
            item = await visual_service.create_candidate_with_ai_diagram(
                provider=tm.course_service,
                course_id=course_id,
                lesson_unit_id=lesson_unit_id,
                script_revision_id=revision_id,
                section_node_id=body.section_node_id,
                block=block,
                instruction=body.instruction,
            )
        elif body.expression_type == "animation":
            item = await visual_service.create_candidate_with_ai_animation(
                provider=tm.course_service,
                course_id=course_id,
                lesson_unit_id=lesson_unit_id,
                script_revision_id=revision_id,
                section_node_id=body.section_node_id,
                block=block,
                instruction=body.instruction,
            )
        else:
            item = await run_in_threadpool(
                visual_service.create_candidate,
                course_id=course_id,
                lesson_unit_id=lesson_unit_id,
                script_revision_id=revision_id,
                section_node_id=body.section_node_id,
                block=block,
                expression_type=body.expression_type,
                instruction=body.instruction,
            )
        return {"item": item}
    except (TeacherLessonAuthoringError, RepresentationConflict, ValueError) as exc:
        if isinstance(exc, TeacherLessonAuthoringError):
            _raise(exc)
        _raise(TeacherLessonAuthoringError(
            "lesson_script_visual_generation_failed", str(exc)
        ))


@router.post(
    "/courses/{course_id}/lessons/{lesson_unit_id}/script/visuals/{representation_id}/resolve"
)
async def resolve_lesson_script_visual(
    course_id: str,
    lesson_unit_id: str,
    representation_id: str,
    body: ResolveLessonScriptVisualRequest,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
    visual_service: TeacherScriptVisualService = Depends(
        get_teacher_script_visual_service
    ),
):
    try:
        resolve_user_id(request.headers.get("X-User-Id"))
        revision, _blocks = _current_script_visual_context(
            tm, repository, course_id, lesson_unit_id
        )
        revision_id = str(revision.get("revision_id") or "")
        if body.script_revision_id != revision_id:
            raise TeacherLessonAuthoringError(
                "lesson_script_revision_conflict",
                "讲义工作稿已经变化，不能处理旧的视觉表达候选。",
            )
        item = await run_in_threadpool(
            visual_service.resolve_candidate,
            course_id=course_id,
            lesson_unit_id=lesson_unit_id,
            script_revision_id=revision_id,
            representation_id=representation_id,
            accept=body.accept,
        )
        return {"item": item}
    except (TeacherLessonAuthoringError, RepresentationConflict) as exc:
        if isinstance(exc, TeacherLessonAuthoringError):
            _raise(exc)
        _raise(TeacherLessonAuthoringError(
            "lesson_script_visual_conflict", str(exc)
        ))


@router.put("/courses/{course_id}/lessons/{lesson_unit_id}/script/draft")
async def save_lesson_script_draft(
    course_id: str,
    lesson_unit_id: str,
    body: SaveLessonScriptDraftRequest,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        source = _source_course(tm, course_id)
        scope = lesson_scope(source, lesson_unit_id)
        lesson = repository.lesson(course_id, lesson_unit_id)
        current_revision = str(lesson.get("working_script_revision_id") or "")
        if body.base_revision_id != current_revision:
            raise TeacherLessonAuthoringError(
                "lesson_script_revision_conflict",
                "讲义工作稿已经变化，请重新载入后再保存。",
            )
        # The first teacher-authored draft has no model revision to build on.
        # Treat an empty base id as the explicit empty working state so a
        # teacher can take over after generation fails (or write from scratch)
        # without creating a parallel persistence path.
        base_revision = (
            _script_revision(
                repository, course_id, lesson_unit_id, body.base_revision_id
            )
            if body.base_revision_id
            else {}
        )
        base_sections = {
            str(item.get("section_node_id") or ""): item
            for item in base_revision.get("sections") or []
            if isinstance(item, dict) and item.get("section_node_id")
        }
        expected_ids = [str(item.get("node_id") or "") for item in scope["sections"]]
        actual_ids = [str(item.get("section_node_id") or "") for item in body.sections]
        if actual_ids != expected_ids:
            raise TeacherLessonAuthoringError(
                "lesson_script_scope_conflict",
                "讲义小节与当前大纲不一致，请重新载入。",
            )
        lesson, plan_revision = _current_plan_revision(
            repository,
            course_id,
            lesson_unit_id,
        )
        plan_revision_id = str(plan_revision.get("revision_id") or "")
        plan_sections = {
            str(item.get("node_id") or ""): item
            for item in (plan_revision.get("plan") or {}).get("sections") or []
            if isinstance(item, dict) and item.get("node_id")
        }
        outline_sections = {
            str(item.get("node_id") or ""): item for item in scope["sections"]
        }
        normalized_sections = []
        for item in body.sections:
            section_id = str(item.get("section_node_id") or "")
            contract = compile_teacher_script_module_contract(
                outline_sections.get(section_id) or {},
                plan_sections.get(section_id) or {},
            )
            normalized = normalize_teacher_script_section(item, contract)
            previous_blocks = {
                str(block.get("block_id") or ""): block
                for block in (base_sections.get(section_id) or {}).get("blocks") or []
                if isinstance(block, dict) and block.get("block_id")
            }
            for block in normalized.get("blocks") or []:
                if not isinstance(block, dict):
                    continue
                previous = previous_blocks.get(str(block.get("block_id") or "")) or {}
                block["generation_source"] = (
                    str(previous.get("generation_source") or "")
                    if str(previous.get("content") or "").strip()
                    == str(block.get("content") or "").strip()
                    else "teacher_edit"
                ) or "teacher_edit"
            normalized["quality_report"] = validate_teacher_script_section(
                normalized,
                contract,
            )
            normalized_sections.append(normalized)
        repository.save_script_revision(
            course_id,
            lesson_unit_id,
            normalized_sections,
            source_lesson_plan_revision_id=plan_revision_id,
            generation_source="teacher_edit",
            requirements=str(base_revision.get("requirements") or ""),
            material_asset_ids=list(base_revision.get("material_asset_ids") or []),
            actor=resolve_user_id(request.headers.get("X-User-Id")),
        )
        projected = next(
            item for item in _lesson_projection(source, repository)
            if item["lesson_unit_id"] == lesson_unit_id
        )
        return {"lesson": projected}
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.post("/courses/{course_id}/lessons/{lesson_unit_id}/script/rewrite-candidate")
@structured_generation_stream(
    stage="lesson_script_candidate",
    started_message="已收到讲义修改要求。",
    waiting_message="AI 正在生成可审阅的讲义候选。",
)
async def rewrite_lesson_script_candidate(
    course_id: str,
    lesson_unit_id: str,
    body: RewriteLessonScriptRequest,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        source = _source_course(tm, course_id)
        scope = lesson_scope(source, lesson_unit_id)
        lesson = repository.lesson(course_id, lesson_unit_id)
        if lesson.get("working_script_revision_id") != body.base_revision_id:
            raise TeacherLessonAuthoringError(
                "lesson_script_revision_conflict",
                "讲义工作稿已经变化，请重新载入后再优化。",
            )
        revision = _script_revision(
            repository, course_id, lesson_unit_id, body.base_revision_id
        )
        section = next(
            (
                item for item in revision.get("sections") or []
                if item.get("section_node_id") == body.section_node_id
            ),
            None,
        )
        outline_section = next(
            (
                item for item in scope["sections"]
                if item.get("node_id") == body.section_node_id
            ),
            None,
        )
        if not isinstance(section, dict) or not isinstance(outline_section, dict):
            raise TeacherLessonAuthoringError(
                "lesson_script_section_not_found",
                "当前讲义小节不存在。",
            )
        selected_material_ids, source_evidence = _course_material_evidence(
            course_id,
            resolve_user_id(request.headers.get("X-User-Id")),
            body.material_asset_ids,
        )
        plan_revision = _plan_revision(
            repository,
            course_id,
            lesson_unit_id,
            str(revision.get("source_lesson_plan_revision_id") or ""),
        )
        plan_section = next(
            (
                item for item in (plan_revision.get("plan") or {}).get("sections") or []
                if isinstance(item, dict) and item.get("node_id") == body.section_node_id
            ),
            {},
        )
        script_headings = [
            str(item.get("title") or "").strip()
            for item in section.get("blocks") or []
            if isinstance(item, dict) and str(item.get("title") or "").strip()
        ]
        result = await tm.course_service.rewrite_selection(
            course_id=course_id,
            node=outline_section,
            selected_text=str(section.get("content") or ""),
            node_content=str(section.get("content") or ""),
            heading_path=[str(section.get("title") or "")],
            user_requirement="\n".join(filter(None, [
                body.instruction.strip(),
                "保持当前教案结构和事实边界；涉及高风险事实而选定资料无法支持时标注“需核验”，不得给出无依据的绝对结论。",
                (
                    "完整保留并仅使用这些二级标题，顺序和名称均不得改变："
                    + "、".join(f"## {title}" for title in script_headings)
                ) if script_headings else "",
            ])),
            action_type="rewrite",
            course_context=json.dumps({
                "lesson_sections": [
                    str(item.get("title") or "") for item in revision.get("sections") or []
                ],
                "current_lesson_plan": plan_section,
                "teacher_requirements": str(revision.get("requirements") or ""),
                "material_asset_ids": selected_material_ids,
                "selected_material_evidence": _prompt_material_evidence(source_evidence),
            }, ensure_ascii=False),
            user_id=resolve_user_id(request.headers.get("X-User-Id")),
        )
        replacement_text = str(result.get("replacement_text") or "").strip()
        if not replacement_text:
            raise TeacherLessonAuthoringError(
                "lesson_script_candidate_empty",
                "AI 没有生成可审阅的讲义修改。",
            )
        candidate = repository.save_script_ai_candidate(
            course_id,
            lesson_unit_id,
            base_revision_id=body.base_revision_id,
            section_node_id=body.section_node_id,
            instruction=body.instruction.strip(),
            replacement_text=replacement_text,
            source_lesson_plan_revision_id=str(
                revision.get("source_lesson_plan_revision_id") or ""
            ),
            material_asset_ids=selected_material_ids,
        )
        return {"candidate": candidate}
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.post(
    "/courses/{course_id}/lessons/{lesson_unit_id}/script/ai-candidates/{candidate_id}/resolve"
)
async def resolve_lesson_script_candidate(
    course_id: str,
    lesson_unit_id: str,
    candidate_id: str,
    body: ResolveLessonScriptCandidateRequest,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        source = _source_course(tm, course_id)
        scope = lesson_scope(source, lesson_unit_id)
        lesson = repository.lesson(course_id, lesson_unit_id)
        candidate = repository.script_ai_candidate(
            course_id, lesson_unit_id, candidate_id
        )
        if candidate.get("status") != "pending":
            projected = next(
                item for item in _lesson_projection(source, repository)
                if item["lesson_unit_id"] == lesson_unit_id
            )
            return {"lesson": projected, "candidate": candidate}
        base_revision_id = str(candidate.get("base_revision_id") or "")
        if lesson.get("working_script_revision_id") != base_revision_id:
            raise TeacherLessonAuthoringError(
                "lesson_script_revision_conflict",
                "讲义工作稿已经变化，不能覆盖新修改。",
            )
        if not body.accept:
            resolved = repository.mark_script_ai_candidate(
                course_id,
                lesson_unit_id,
                candidate_id,
                status="rejected",
            )
            projected = next(
                item for item in _lesson_projection(source, repository)
                if item["lesson_unit_id"] == lesson_unit_id
            )
            return {"lesson": projected, "candidate": resolved}

        base_revision = _script_revision(
            repository, course_id, lesson_unit_id, base_revision_id
        )
        plan_revision_id = str(candidate.get("source_lesson_plan_revision_id") or "")
        plan_revision = _plan_revision(
            repository, course_id, lesson_unit_id, plan_revision_id
        )
        plan_sections = {
            str(item.get("node_id") or ""): item
            for item in (plan_revision.get("plan") or {}).get("sections") or []
            if isinstance(item, dict) and item.get("node_id")
        }
        outline_sections = {
            str(item.get("node_id") or ""): item for item in scope["sections"]
        }
        target_section_id = str(candidate.get("section_node_id") or "")
        normalized_sections: list[dict[str, Any]] = []
        for section in base_revision.get("sections") or []:
            if not isinstance(section, dict):
                continue
            section_id = str(section.get("section_node_id") or "")
            candidate_section = deepcopy(section)
            if section_id == target_section_id:
                candidate_section.pop("blocks", None)
                candidate_section["content"] = str(
                    candidate.get("replacement_text") or ""
                ).strip()
            contract = compile_teacher_script_module_contract(
                outline_sections.get(section_id) or {},
                plan_sections.get(section_id) or {},
            )
            normalized = normalize_teacher_script_section(
                candidate_section,
                contract,
            )
            previous_blocks = {
                str(block.get("block_id") or ""): block
                for block in section.get("blocks") or []
                if isinstance(block, dict) and block.get("block_id")
            }
            for block in normalized.get("blocks") or []:
                if not isinstance(block, dict):
                    continue
                previous = previous_blocks.get(str(block.get("block_id") or "")) or {}
                block["generation_source"] = (
                    "ai_optimization"
                    if section_id == target_section_id
                    else str(previous.get("generation_source") or "") or "model"
                )
            normalized["quality_report"] = validate_teacher_script_section(
                normalized,
                contract,
            )
            normalized_sections.append(normalized)
        saved = repository.save_script_revision(
            course_id,
            lesson_unit_id,
            normalized_sections,
            source_lesson_plan_revision_id=plan_revision_id,
            generation_source="ai_optimization",
            requirements=str(base_revision.get("requirements") or ""),
            material_asset_ids=list(
                candidate.get("material_asset_ids")
                or base_revision.get("material_asset_ids")
                or []
            ),
            actor=resolve_user_id(request.headers.get("X-User-Id")),
            expected_working_revision_id=base_revision_id,
        )
        accepted_revision_id = str(saved.get("working_script_revision_id") or "")
        resolved = repository.mark_script_ai_candidate(
            course_id,
            lesson_unit_id,
            candidate_id,
            status="accepted",
            result_revision_id=accepted_revision_id,
        )
        projected = next(
            item for item in _lesson_projection(source, repository)
            if item["lesson_unit_id"] == lesson_unit_id
        )
        return {"lesson": projected, "candidate": resolved}
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.post("/courses/{course_id}/lessons/{lesson_unit_id}/plan/ai-candidates")
@structured_generation_stream(
    stage="lesson_plan_candidate",
    started_message="已收到教案修改要求。",
    waiting_message="AI 正在生成可审阅的教案候选。",
)
async def create_lesson_plan_candidate(
    course_id: str,
    lesson_unit_id: str,
    body: CreateLessonPlanCandidateRequest,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        actor = resolve_user_id(request.headers.get("X-User-Id"))
        selected_material_ids, source_evidence = _course_material_evidence(
            course_id,
            actor,
            body.material_asset_ids,
        )
        lesson = repository.lesson(course_id, lesson_unit_id)
        if lesson.get("working_revision_id") != body.base_revision_id:
            raise TeacherLessonAuthoringError("lesson_plan_revision_conflict", "教案草稿已经变化，请重新打开后再优化。")
        revision = next(
            (
                item for item in lesson.get("revisions") or []
                if item.get("revision_id") == body.base_revision_id
            ),
            None,
        )
        if not isinstance(revision, dict):
            raise TeacherLessonAuthoringError("lesson_plan_revision_not_found", "教案草稿不存在。")
        source = _source_course(tm, course_id)
        scope = lesson_scope(source, lesson_unit_id)
        arrangement = repository.current_arrangement(course_id, lesson_unit_id) or {}
        lesson_node = scope.get("lesson") or {}
        knowledge_context = (
            course_knowledge_base_prompt_context(
                source.get("course_knowledge_base") or {},
                body.section_node_id,
            )
            if body.section_node_id
            else ""
        )
        optimized = await tm.course_service.optimize_teacher_lesson_plan(
            plan=deepcopy(revision.get("plan") or {}),
            instruction=body.instruction,
            section_node_id=body.section_node_id,
            target_field=body.target_field,
            target_item_id=body.target_item_id,
            selected_text=body.selected_text,
            lesson_context={
                "lesson_unit_id": lesson_unit_id,
                "title": str(lesson_node.get("node_name") or lesson_node.get("title") or ""),
                "lesson_type": str(arrangement.get("lesson_type") or ""),
                "duration_minutes": int(float(
                    lesson_node.get("duration_minutes")
                    or arrangement.get("total_minutes")
                    or 0
                )),
            },
            knowledge_context=knowledge_context,
            material_evidence=_prompt_material_evidence(
                source_evidence,
                character_budget=8000,
            ),
        )
        candidate = repository.save_ai_candidate(
            course_id,
            lesson_unit_id,
            base_revision_id=body.base_revision_id,
            instruction=body.instruction,
            section_node_id=body.section_node_id,
            target_field=body.target_field,
            target_item_id=body.target_item_id,
            selected_text=body.selected_text,
            plan=optimized["plan"],
            material_asset_ids=selected_material_ids,
        )
        return {"candidate": candidate}
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.post("/courses/{course_id}/lessons/{lesson_unit_id}/plan/ai-candidates/{candidate_id}/resolve")
async def resolve_lesson_plan_candidate(
    course_id: str,
    lesson_unit_id: str,
    candidate_id: str,
    body: ResolveLessonPlanCandidateRequest,
    request: Request,
    tm: TaskManager = Depends(require_task_manager),
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        source = _source_course(tm, course_id)
        canonical_outline_revision = _canonical_outline_revision(source)
        if canonical_outline_revision:
            repository.set_outline(course_id, canonical_outline_revision)
        TeacherLessonAuthoringService(repository).resolve_ai_candidate(
            course_id=course_id,
            lesson_unit_id=lesson_unit_id,
            course_data=source,
            candidate_id=candidate_id,
            accept=body.accept,
            actor=resolve_user_id(request.headers.get("X-User-Id")),
        )
        projected = next(
            item for item in _lesson_projection(source, repository)
            if item["lesson_unit_id"] == lesson_unit_id
        )
        return {"lesson": projected}
    except TeacherLessonAuthoringError as exc:
        _raise(exc)


@router.get("/courses/{course_id}/lessons/{lesson_unit_id}/ppt/export.pptx")
async def export_lesson_ppt(
    course_id: str,
    lesson_unit_id: str,
    asset_id: str,
    revision_id: str,
    repository: TeacherLessonAuthoringRepository = Depends(
        get_teacher_lesson_authoring_repository
    ),
):
    try:
        _asset, revision = _ppt_asset_revision(
            repository,
            course_id,
            lesson_unit_id,
            asset_id,
            revision_id,
        )
        structured = teacher_lesson_deck_to_structured_slide_deck(
            deepcopy(revision.get("deck") or {}),
            source_revision_id=str(revision.get("source_lesson_plan_revision_id") or ""),
        )
        export_dir = repository.root / "exports"
        export_dir.mkdir(parents=True, exist_ok=True)
        output = export_dir / f"{course_id}-{lesson_unit_id}-{revision_id}.pptx"
        await run_in_threadpool(
            export_structured_slide_deck,
            structured,
            output,
            require_quality=False,
            theme="qingfeng-classroom",
        )
        return FileResponse(
            output,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=f"{lesson_unit_id}-课堂课件.pptx",
        )
    except TeacherLessonAuthoringError as exc:
        _raise(exc)
