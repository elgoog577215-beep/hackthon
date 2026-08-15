"""Renderer-bound adapters for the closed slide-deck V6 template contract."""

from __future__ import annotations

import json
import re
from functools import lru_cache
from pathlib import Path
from typing import Any

from slide_asset_repository import SlideAssetRepository, slide_asset_repository
from slide_deck import SlideBlockSpec, SlideSpec
from slide_deck_renderer import (
    SlideDeckQualityError,
    _render_slide,
    validate_theme,
)
from slide_deck_v6 import SlideDeckV6, SlidePageV6

_LAYOUT_ADAPTER_PATH = (
    Path(__file__).resolve().parents[1]
    / "frontend"
    / "src"
    / "data"
    / "slide-deck-v6-layout-adapters.json"
)


@lru_cache(maxsize=1)
def _layout_adapters() -> dict[str, dict[str, Any]]:
    payload = json.loads(_LAYOUT_ADAPTER_PATH.read_text(encoding="utf-8"))
    if payload.get("schema_version") != "slide_deck_v6_layout_adapters_v1":
        raise ValueError("Unsupported V6 layout adapter contract")
    layouts = payload.get("layouts") or {}
    if not isinstance(layouts, dict) or not layouts:
        raise ValueError("V6 layout adapter contract is empty")
    adapters: dict[str, dict[str, Any]] = {}
    for slug, value in layouts.items():
        if not isinstance(value, dict):
            continue
        adapter: dict[str, Any] = {
            "renderer_layout": str(value.get("renderer_layout") or ""),
            "basic_layout": str(value.get("basic_layout") or ""),
            "capacity_profile": str(value.get("capacity_profile") or ""),
        }
        policy = value.get("variant_policy")
        if isinstance(policy, dict):
            adapter["variant_policy"] = {
                str(key): str(item or "") for key, item in policy.items()
            }
        adapters[str(slug)] = adapter
    return adapters


def _layout_slug(template_layout_id: str) -> str:
    slug = str(template_layout_id or "").rsplit("/", 1)[-1]
    if slug not in _layout_adapters():
        raise ValueError(f"v6_template_layout_adapter_missing:{template_layout_id}")
    return slug


def _layout_variant(
    page: SlidePageV6,
    adapter: dict[str, Any],
) -> tuple[str, str]:
    policy = adapter.get("variant_policy")
    if not isinstance(policy, dict):
        return "", ""
    artifact_kind = str(policy.get("artifact_content_kind") or "")
    if not artifact_kind:
        return "", ""
    has_artifact = any(region.content_kind == artifact_kind for region in page.regions)
    has_support = any(region.content_kind != artifact_kind for region in page.regions)
    wide_min_columns = int(policy.get("wide_min_columns") or 0)
    artifact_region = next(
        (
            region
            for region in page.regions
            if region.content_kind == artifact_kind
        ),
        None,
    )
    if (
        artifact_kind == "table"
        and artifact_region is not None
        and len(_parse_markdown_table(artifact_region.content)[1]) == 1
        and policy.get("detail_variant")
        and (
            page.continuation_index > 1
            or (
                not has_support
                and _table_row_requires_detail(artifact_region.content)
            )
        )
    ):
        return str(policy["detail_variant"]), "full"
    if page.continuation_index > 1:
        return str(policy.get("continuation_variant") or ""), "full"
    if (
        has_artifact
        and has_support
        and artifact_kind == "table"
        and artifact_region is not None
        and wide_min_columns
        and len(_parse_markdown_table(artifact_region.content)[0]) >= wide_min_columns
    ):
        return (
            str(policy.get("wide_variant") or policy.get("split_variant") or ""),
            str(policy.get("wide_support_mode") or "band"),
        )
    if has_artifact and has_support:
        return str(policy.get("split_variant") or ""), "split"
    return str(policy.get("full_variant") or ""), "full"


def _audience_title(page: SlidePageV6) -> str:
    title = str(page.title or "").strip()
    if page.continuation_count <= 1:
        return title
    return re.sub(
        r"\s*[（(]\s*\d+\s*/\s*\d+\s*[）)]\s*$",
        "",
        title,
    ).strip()


def _table_row_requires_detail(value: str) -> bool:
    headers, rows = _parse_markdown_table(value)
    if len(rows) != 1:
        return False
    cells = rows[0]
    column_count = max(1, len(headers), len(cells))
    safe_column_chars = max(8, round(108 / column_count))
    return max((len(cell) for cell in cells), default=0) > safe_column_chars


def _speaker_notes(page: SlidePageV6) -> str:
    sections = [
        f"source_document_revision: {page.speaker_notes.source_document_revision}",
        f"teaching_unit_id: {page.speaker_notes.teaching_unit_id}",
        "source_section_ids: " + json.dumps(
            page.speaker_notes.source_section_ids,
            ensure_ascii=False,
        ),
    ]
    sections.extend(
        "\n".join([
            f"[{block.block_id} @ {block.block_revision}]",
            f"source_kind: {block.source_kind}",
            f"asset_refs: {json.dumps(block.asset_refs, ensure_ascii=False)}",
            block.full_text,
            f"source_payload: {json.dumps(block.source_payload, ensure_ascii=False, sort_keys=True)}",
        ])
        for block in page.speaker_notes.source_blocks
    )
    return "\n\n".join(sections)


def _region_block(page: SlidePageV6, region: Any) -> SlideBlockSpec:
    metadata: dict[str, Any] = {
        "v6_slot_id": region.slot_id,
        "v6_region_id": region.region_id,
        "source_block_ids": list(region.source_block_ids),
        "source_asset_refs": list(region.source_asset_refs),
    }
    block_type = "statement"
    items: list[str] = []
    if region.content_kind == "code":
        block_type = "code"
    elif region.content_kind == "items":
        block_type = "bullets"
        items = [line.strip() for line in region.content.splitlines() if line.strip()]
    elif region.content_kind == "steps":
        block_type = "process"
        items = [line.strip() for line in region.content.splitlines() if line.strip()]
    elif region.content_kind == "formula":
        metadata["formula"] = True
    elif region.content_kind == "table":
        metadata["table_source"] = True
    return SlideBlockSpec(
        block_id=region.region_id,
        type=block_type,
        # Slot identifiers describe the template contract; they are not
        # audience-facing copy and must never become visible headings.
        title="",
        content="" if items else region.content,
        items=items,
        metadata=metadata,
    )


def _parse_markdown_table(value: str) -> tuple[list[str], list[list[str]]]:
    rows = []
    for line in str(value or "").splitlines():
        stripped = line.strip()
        if not stripped.startswith("|") or not stripped.endswith("|"):
            continue
        if re.fullmatch(r"[|:\-\s]+", stripped) and "-" in stripped:
            continue
        cells = [
            cell.replace(r"\|", "|").strip()
            for cell in re.split(r"(?<!\\)\|", stripped.strip("|"))
        ]
        if cells:
            rows.append(cells)
    return (rows[0], rows[1:]) if rows else ([], [])


def _visuals(page: SlidePageV6) -> list[dict[str, Any]]:
    decision = page.visual_decision.decision
    by_kind = {region.content_kind: region for region in page.regions}
    if decision == "formula" and "formula" in by_kind:
        return [{
            "kind": "formula",
            "caption": "",
            "alt_text": _audience_title(page),
            "parameters": {"formula": by_kind["formula"].content},
        }]
    if decision in {"table", "data"} and "table" in by_kind:
        headers, rows = _parse_markdown_table(by_kind["table"].content)
        return [{
            "kind": "table",
            "caption": "",
            "alt_text": _audience_title(page),
            "parameters": {"headers": headers, "rows": rows},
        }]
    if decision == "diagram":
        payload = page.visual_decision.visual_payload
        nodes = [
            {
                "node_id": str(node.get("node_id") or ""),
                "label": str(node.get("label") or ""),
                "emphasis": str(node.get("emphasis") or ("primary" if index == 0 else "supporting")),
                "source_fragment_ids": list(node.get("source_block_ids") or []),
            }
            for index, node in enumerate(payload.get("nodes") or [])
        ]
        edges = [
            {
                "source": str(edge.get("source") or ""),
                "target": str(edge.get("target") or ""),
                "label": str(edge.get("label") or ""),
                "relation": str(edge.get("relation") or "sequence"),
            }
            for edge in payload.get("edges") or []
        ]
        if len(nodes) < 2 or not edges:
            raise ValueError(f"v6_visual_diagram_payload_missing:{page.page_id}")
        return [{
            "kind": "rule_diagram",
            "caption": page.title,
            "nodes": nodes,
            "edges": edges,
            "source_fragment_ids": list(page.source_block_ids),
            "alt_text": page.title,
            "parameters": {
                "direction": str(payload.get("direction") or "vertical"),
                "template": "process",
                "relation_evidence": list(page.source_block_ids),
            },
        }]
    if decision in {"image", "experiment"}:
        asset_ids = list(page.visual_decision.source_asset_ids)
        if not asset_ids:
            asset_ids = [
                asset_ref
                for region in page.regions
                for asset_ref in region.source_asset_refs
            ]
        if not asset_ids:
            raise ValueError(f"v6_visual_source_asset_missing:{page.page_id}")
        return [{
            "kind": "source_image",
            "caption": page.title,
            "alt_text": page.title,
            "asset_id": asset_ids[0],
            "source_fragment_ids": list(page.source_block_ids),
            "parameters": {"asset_ref": asset_ids[0]},
        }]
    return []


def adapt_v6_page_to_slide_spec(page: SlidePageV6 | dict[str, Any]) -> SlideSpec:
    resolved_page = (
        page if isinstance(page, SlidePageV6) else SlidePageV6.model_validate(page)
    )
    slug = _layout_slug(resolved_page.resolved_layout)
    adapter = _layout_adapters()[slug]
    renderer_layout = adapter["renderer_layout"]
    layout_variant, artifact_support_mode = _layout_variant(resolved_page, adapter)
    practice_artifact_kind = (
        slug.removeprefix("practice-")
        if slug in {"practice-code", "practice-formula", "practice-table"}
        else ""
    )
    subtitle = next(
        (
            region.content
            for region in resolved_page.regions
            if region.slot_id == "subtitle"
        ),
        "",
    )
    if slug == "cover-minimal":
        subtitle = ""
    eyebrow = next(
        (
            region.content
            for region in resolved_page.regions
            if region.slot_id == "eyebrow"
        ),
        "",
    )
    return SlideSpec(
        unit_id=resolved_page.page_id,
        position=resolved_page.page_ordinal,
        layout=adapter["basic_layout"] or "concept",
        slide_purpose=slug,
        eyebrow=eyebrow,
        title=_audience_title(resolved_page),
        subtitle=subtitle,
        composition="diagram-full" if slug == "evidence-diagram" else "",
        visuals=_visuals(resolved_page),
        blocks=[_region_block(resolved_page, region) for region in resolved_page.regions],
        speaker_notes=_speaker_notes(resolved_page),
        source_block_ids=list(resolved_page.source_block_ids),
        quality={
            "passed": True,
            "render_contract": "template_layout_contract_v1",
            "audience_label_policy": "source_only",
            "v6_template_layout_id": resolved_page.resolved_layout,
            "v6_layout_slug": slug,
            "v6_layout_variant": layout_variant,
            "v6_artifact_support_mode": artifact_support_mode,
            "v6_capacity_profile": adapter.get("capacity_profile", ""),
            "v6_continuation_index": resolved_page.continuation_index,
            "v6_continuation_count": resolved_page.continuation_count,
            "v6_title_max_lines": resolved_page.title_max_lines,
            "v6_practice_artifact_kind": practice_artifact_kind,
            "resolved_layout": renderer_layout,
            "task_prompt_mode": (
                "artifact-guided" if practice_artifact_kind
                else "action" if slug == "practice-prompt"
                else ""
            ),
            "prompt_label": (
                "执行并核验" if practice_artifact_kind
                else "执行步骤" if slug == "practice-prompt"
                else ""
            ),
        },
    )


def _mark_v6_title_shape(slide: Any, unit: SlideSpec) -> None:
    normalized_title = re.sub(r"\s+", "", str(unit.title or ""))
    if not normalized_title:
        return
    max_lines = max(1, int(unit.quality.get("v6_title_max_lines") or 1))
    matches = [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        and re.sub(r"\s+", "", str(shape.text or "")) == normalized_title
    ]
    if not matches:
        return
    title_shape = max(matches, key=lambda shape: int(shape.height) * int(shape.width))
    title_shape.name = f"{title_shape.name} [v6-title-max-lines={max_lines}]"


def _validate_deck_for_export(deck: SlideDeckV6) -> None:
    checks = {
        "formal_block_visible_coverage": deck.quality.formal_block_visible_coverage == 1.0,
        "full_text_note_binding": deck.quality.full_text_note_binding == 1.0,
        "source_order_preserved": deck.quality.source_order_preserved,
        "template_contract_passed": deck.quality.template_contract_passed,
        "subject_artifacts_passed": deck.quality.subject_artifacts_passed,
        "web_pptx_contract_shared": deck.quality.web_pptx_contract_shared,
    }
    blockers = [
        {
            "severity": "critical",
            "code": f"v6_{name}_failed",
            "message": f"V6 export requires {name}",
        }
        for name, passed in checks.items()
        if not passed
    ]
    blockers.extend(item.model_dump(mode="json") for item in deck.quality.blockers)
    if blockers:
        raise SlideDeckQualityError(
            {"passed": False, "score": 0, "blockers": blockers, "warnings": []}
        )


_V6_PUBLICATION_METADATA_FIELDS = frozenset({
    "ai_batch_diagnostics",
    "build_signature",
    "course_presentation_graph",
    "planning_status",
    "source_contract",
    "story_plan",
    "visual_plan",
})


def _validated_export_deck(
    content: SlideDeckV6 | dict[str, Any],
) -> SlideDeckV6:
    if isinstance(content, SlideDeckV6):
        return content
    core_fields = set(SlideDeckV6.model_fields)
    unknown_fields = (
        set(content)
        - core_fields
        - _V6_PUBLICATION_METADATA_FIELDS
    )
    if unknown_fields:
        # Preserve the strict model's field-specific validation error instead
        # of silently accepting an undeclared publication extension.
        return SlideDeckV6.model_validate(content)
    return SlideDeckV6.model_validate({
        key: value
        for key, value in content.items()
        if key in core_fields
    })


def export_slide_deck_v6_pptx(
    content: SlideDeckV6 | dict[str, Any],
    output_path: str | Path,
    *,
    asset_repository: SlideAssetRepository | None = None,
) -> Path:
    deck = _validated_export_deck(content)
    _validate_deck_for_export(deck)

    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    theme = dict(validate_theme(deck.theme))
    theme.update(deck.template_theme_overrides)
    assets = asset_repository or slide_asset_repository
    slides = [adapt_v6_page_to_slide_spec(page) for page in deck.pages]
    for index, unit in enumerate(slides):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        _render_slide(slide, unit, index + 1, len(slides), theme, assets)
        _mark_v6_title_shape(slide, unit)
        slide.notes_slide.notes_text_frame.text = unit.speaker_notes

    path = Path(output_path)
    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)
    return path


__all__ = [
    "adapt_v6_page_to_slide_spec",
    "export_slide_deck_v6_pptx",
]
