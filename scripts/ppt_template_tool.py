#!/usr/bin/env python3
"""Inspect native PPTX or certify actual filled teaching layout specimens."""
from __future__ import annotations

import argparse
import json
import os
import subprocess
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "backend"))

from pptx import Presentation
from pptx.util import Pt

from ppt_layout_execution import ASSET_ROOT, LAYOUT_VERSION, LayoutExecution, certification_version, compile_teaching_template, file_digest, tool_identity
from ppt_layout_samples import formula_comparison_sample, layout_sample, matrix_boundary_sample
from ppt_native_scene import audit_scene, clone_native_slide, inspect_native_bindings, inspect_native_connections, render_scene
from ppt_page_scene import resolve_page_scenes
from ppt_teaching_content import PageTeachingV2


from ppt_render_audit import render_evidence


def certify_export(path, count):
    from slide_deck_renderer import audit_exported_pptx
    report = audit_exported_pptx(path, expected_slide_count=count, require_pixel_audit=False)
    if not report["passed"]:
        raise ValueError(f"template_export_audit_failed:{report['blockers']}")
    return report


def native_artwork(source_path, execution, output):
    """Render only the template artwork for Web; PPTX keeps original objects."""
    import base64
    from types import SimpleNamespace
    deck = Presentation(source_path)
    count = len(deck.slides)
    source = deck.slides[execution.source_slide_number - 1]
    slide = clone_native_slide(deck, source, execution)
    for binding in execution.connections.values():
        connector = next(s for s in slide.shapes if s.shape_id == binding.shape_id)
        connector._element.getparent().remove(connector._element)
    from slide_deck_v6_personal_renderer import _remove_original_slides
    _remove_original_slides(deck, count)
    output.mkdir(parents=True, exist_ok=True)
    path = output / 'native-artwork.pptx'
    deck.save(path)
    render_evidence(path, [SimpleNamespace(width=960, height=540, objects=[], edges=[], logical_page_id='artwork', state_id='static')])
    png = output / 'native-artwork-1.png'
    return base64.b64encode(png.read_bytes()).decode(), file_digest(png)


def certify_native(source_path, specification_path, output, *, asset_repository=None):
    """Certify explicit targets against operator supplied boundary specimens."""
    specification = json.loads(specification_path.read_text())
    execution = LayoutExecution.model_validate(specification["execution"])
    if execution.mode != "native_fill" or file_digest(source_path) != execution.source_sha256:
        raise ValueError("native_template_source_changed")
    deck = Presentation(source_path)
    if deck.slide_width != Pt(960) or deck.slide_height != Pt(540):
        raise ValueError("native_source_page_geometry_unsupported")
    originals = list(deck.slides)
    if not 1 <= execution.source_slide_number <= len(originals):
        raise ValueError("native_source_slide_missing")
    source = originals[execution.source_slide_number - 1]
    execution.targets = inspect_native_bindings(source, execution.targets)
    inspect_native_connections(source, execution)
    execution.static_artwork_data, execution.static_artwork_sha256 = native_artwork(source_path, execution, output)
    template = compile_teaching_template(specification.get("theme", "academic-editorial"), certification_required=False)
    layout = template.get_layout(template.layout_id(execution.component_id)).model_copy(deep=True)
    layout.execution = execution
    samples = specification["samples"]
    if set(samples) != {"short", "normal", "long"}:
        raise ValueError("native_boundary_samples_missing")
    scenes = []
    for name in ("short", "normal", "long"):
        try:
            resolved = resolve_page_scenes(page_id=name, title=specification.get("title", "原生模板填充测试"),
                content=PageTeachingV2.model_validate(samples[name]), layout=layout, template=template, source_document_revision="native-sample")
        except ValueError as exc:
            if name == "long" and "capacity" in str(exc):
                continue
            raise
        if name == "long":
            raise ValueError("long_text_was_not_rejected")
        scenes.extend(resolved)
    for scene in scenes:
        slide = clone_native_slide(deck, source, execution)
        render_scene(slide, scene, assets=asset_repository)
    from slide_deck_v6_personal_renderer import _remove_original_slides
    _remove_original_slides(deck, len(originals))
    output.mkdir(parents=True, exist_ok=True)
    path = output / "native-filled.pptx"
    deck.save(path)
    reports = [audit_scene(s, c) for s, c in zip(Presentation(path).slides, scenes, strict=True)]
    report = {"status": "passed", "component_version": LAYOUT_VERSION, "tools": tool_identity(),
        "checks": {"short": True, "normal": True, "long": True, "relations": True, "render": True},
        "object_reports": reports, "export_review": certify_export(path, len(scenes)), "classroom_review": "pending", **render_evidence(path, scenes)}
    execution.certification = report
    (output / "execution.json").write_text(execution.model_dump_json(indent=2) + "\n")
    print(json.dumps({"status": "passed", "mode": "native_fill", "physical_pages": len(scenes)}, ensure_ascii=False))
    return execution


def inspect(path):
    deck = Presentation(path)
    def shape_item(s, group_path):
        item = {"shape_id": s.shape_id, "name": s.name, "group_path": group_path,
            "geometry_pt": [round(v / 12700, 3) for v in (s.left, s.top, s.width, s.height)],
            "text": s.text if s.has_text_frame else "", "text_fill_supported": s.has_text_frame,
            "table": {"rows": len(s.table.rows), "columns": len(s.table.columns)} if s.has_table else None}
        if hasattr(s, "shapes"):
            item["children"] = [shape_item(child, group_path + [s.shape_id]) for child in s.shapes]
        if s.has_chart:
            item["unsupported"] = "chart_requires_explicit_series_binding"
        return item
    return {"status": "draft", "source_sha256": file_digest(Path(path)),
            "slides": [{"slide_number": i + 1, "slide_part": str(slide.part.partname),
                        "objects": [shape_item(s, []) for s in slide.shapes]} for i, slide in enumerate(deck.slides)]}


def certify(theme, slugs, output, publish):
    output.mkdir(parents=True, exist_ok=True)
    template = compile_teaching_template(theme, certification_required=False)
    manifest_path = ASSET_ROOT / "teaching-layouts" / f"{theme}.json"
    manifest = json.loads(manifest_path.read_text()) if manifest_path.exists() else {"schema_version": "teaching_template_certification_v2", "layouts": {}}
    for slug in slugs:
        layout = template.get_layout(template.layout_id(slug))
        deck = Presentation()
        deck.slide_width, deck.slide_height = Pt(960), Pt(540)
        scenes = []
        for length in ["short", "normal"]:
            content = layout_sample(slug, length=length)
            scenes.extend(resolve_page_scenes(page_id=f"sample-{length}", title="教学页面实际填充测试", content=content,
                layout=layout, template=template, source_document_revision="sample-v1"))
        if slug == "compare-matrix":
            for count in (3, 4):
                scenes.extend(resolve_page_scenes(page_id=f"matrix-boundary-{count}", title="比较矩阵边界容量测试", content=matrix_boundary_sample(count),
                    layout=layout, template=template, source_document_revision="sample-v1"))
        if slug in {"compare-matrix", "compare-visual"}:
            scenes.extend(resolve_page_scenes(page_id="formula-counterexample", title="矩阵与缺项反例", content=formula_comparison_sample(),
                layout=layout, template=template, source_document_revision="sample-v1"))
        for scene in scenes:
            slide = deck.slides.add_slide(deck.slide_layouts[6])
            render_scene(slide, scene)
        path = output / f"{slug}.pptx"
        deck.save(path)
        readback = Presentation(path)
        reports = [audit_scene(slide, scene) for slide, scene in zip(readback.slides, scenes, strict=True)]
        try:
            resolve_page_scenes(page_id="long", title="超容量测试", content=layout_sample(slug, length="long"),
                layout=layout, template=template, source_document_revision="sample-v1")
        except ValueError as exc:
            if "capacity" not in str(exc):
                raise
        else:
            raise ValueError("long_text_was_not_rejected")
        report = {"status": "passed", "component_version": LAYOUT_VERSION, "tools": tool_identity(),
            "checks": {"short": True, "normal": True, "long": True, "relations": True, "render": True},
            "object_reports": reports, "export_review": certify_export(path, len(scenes)), "classroom_review": "pending", **render_evidence(path, scenes)}
        manifest["layouts"][slug] = report
        (output / f"{slug}-report.json").write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n")
        print(json.dumps({"layout": slug, "status": "passed", "physical_pages": len(scenes)}, ensure_ascii=False), flush=True)
    if publish:
        manifest_path.parent.mkdir(parents=True, exist_ok=True)
        serialized = json.dumps(manifest, ensure_ascii=False, indent=2) + "\n"
        archive = manifest_path.parent / "versions" / f"{theme}--{certification_version(manifest)}.json"
        archive.parent.mkdir(parents=True, exist_ok=True)
        if archive.exists() and archive.read_text() != serialized:
            raise ValueError("immutable_template_version_collision")
        archive.write_text(serialized)
        manifest_path.write_text(serialized)
    return manifest


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description=__doc__)
    sub = parser.add_subparsers(dest="command", required=True)
    p = sub.add_parser("inspect")
    p.add_argument("source", type=Path)
    p.add_argument("--output", type=Path, required=True)
    p = sub.add_parser("certify")
    p.add_argument("--theme", default="academic-editorial")
    p.add_argument("--layouts", default="compare-visual,compare-matrix")
    p.add_argument("--output", type=Path, required=True)
    p.add_argument("--publish", action="store_true")
    p = sub.add_parser("certify-native")
    p.add_argument("source", type=Path)
    p.add_argument("specification", type=Path)
    p.add_argument("--output", type=Path, required=True)
    p.add_argument("--asset-repository", type=Path, help="Explicit local SlideAssetRepository containing the specimen images")
    p = sub.add_parser("register-native")
    p.add_argument("pack_id")
    p.add_argument("execution", type=Path)
    p.add_argument("--owner-id", required=True)
    p.add_argument("--layout", required=True)
    p.add_argument("--repository", type=Path, required=True)
    p.add_argument("--reviewed", action="store_true", help="Confirm that the explicit object bindings and actual rendered samples have been reviewed")
    p.add_argument("--publish", action="store_true")
    args = parser.parse_args()
    if args.command == "inspect":
        args.output.parent.mkdir(parents=True, exist_ok=True)
        args.output.write_text(json.dumps(inspect(args.source), ensure_ascii=False, indent=2) + "\n")
    elif args.command == "certify-native":
        assets = None
        if args.asset_repository:
            from slide_asset_repository import SlideAssetRepository
            if not args.asset_repository.is_dir():
                raise ValueError("native_specimen_asset_repository_missing")
            assets = SlideAssetRepository(args.asset_repository)
        certify_native(args.source.resolve(), args.specification, args.output.resolve(), asset_repository=assets)
    elif args.command == "register-native":
        from ppt_template_packs import PptTemplatePackRepository
        repository = PptTemplatePackRepository(args.repository)
        result = repository.register_teaching_layout(args.pack_id, args.owner_id, args.layout,
            json.loads(args.execution.read_text()), maintainer_reviewed=args.reviewed)
        if args.publish:
            result = repository.publish(args.pack_id, args.owner_id)
        print(json.dumps({"pack_id": result['pack_id'], "status": result['status'], "version": result.get('version')}))
    else:
        certify(args.theme, args.layouts.split(","), args.output.resolve(), args.publish)
