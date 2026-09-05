"""把 slides deck（dict）渲染成 .pptx —— 用 python-pptx，移植自 video-analyze/build_pptx.mjs。

后端容器为 python:3.14-slim（无 Node.js），故用 python-pptx 实现等价版式，避免引入 Node 运行时。
版面 16:9（13.333 x 7.5 英寸），配色/字体与 build_pptx.mjs 保持一致。
"""
from __future__ import annotations

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# 配色（与 build_pptx.mjs 的 C 一致）
C = {
    "ink": "1F2733", "body": "2B3445", "accent": "4E4376", "accent2": "2B5876",
    "muted": "8A93A2", "chipBg": "EEF0FB", "chipBd": "E0E3F6", "line": "E7EAF2",
    "white": "FFFFFF", "cover_ink": "AEB8E6", "cover_sub": "DFE5FF", "cover_meta": "9FB0E8",
    "sec_num": "5C5685", "sec_slash": "6E689A", "footer_bg": "F0F3F8",
}
FNT = "Microsoft YaHei"
MONO = "Consolas"
W, H = 13.333, 7.5


def _rgb(hexstr: str) -> RGBColor:
    return RGBColor.from_string(hexstr)


def _no_line(shape) -> None:
    shape.line.fill.background()


def _rect(slide, x, y, w, h, color, *, line=None, line_w=1.0, radius=None,
          shape=MSO_SHAPE.RECTANGLE):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = _rgb(color)
    if line:
        sp.line.color.rgb = _rgb(line)
        sp.line.width = Pt(line_w)
    else:
        _no_line(sp)
    try:
        sp.shadow.inherit = False
    except Exception:  # noqa: BLE001
        pass
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:  # noqa: BLE001
            pass
    return sp


def _text(slide, x, y, w, h, text, *, size, color, bold=False, align=PP_ALIGN.LEFT,
          anchor=MSO_ANCHOR.TOP, font=FNT, spacing=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text or "")
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.name = font
    f.color.rgb = _rgb(color)
    return tb


def _shape_text(sp, text, *, size, color, bold=False, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.MIDDLE, font=FNT):
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text or "")
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.name = font
    f.color.rgb = _rgb(color)


def _notes(slide, text) -> None:
    if not text:
        return
    slide.notes_slide.notes_text_frame.text = str(text)


def _footer(slide, deck) -> None:
    """模拟 build_pptx.mjs 的 BASE 母版页脚带 + 页脚文字。"""
    _rect(slide, 0, H - 0.32, W, 0.32, C["footer_bg"])
    foot = deck.get("footer") or deck.get("course") or ""
    if foot:
        _text(slide, 0.55, H - 0.34, 9, 0.32, foot, size=9, color=C["muted"],
              anchor=MSO_ANCHOR.MIDDLE)


def _head(slide, title, time_ref=""):
    _rect(slide, 0.6, 0.55, 0.14, 0.7, C["accent"])
    _text(slide, 0.95, 0.5, 10.2 if time_ref else 11.8, 0.8, title or "",
          size=28, color=C["ink"], bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if time_ref:
        badge = _rect(slide, 11.1, 0.62, 1.6, 0.5, C["accent"],
                      radius=0.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _shape_text(badge, time_ref, size=12, color=C["white"], align=PP_ALIGN.CENTER)


def _add_bullets(slide, x, y, w, h, bullets):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    any_line = False
    for b in bullets or []:
        if isinstance(b, str):
            b = {"text": b, "children": []}
        txt = b.get("text", "")
        if not txt:
            continue
        any_line = True
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = 0
        p.space_after = Pt(8)
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = "•  " + txt
        f = run.font
        f.size = Pt(18)
        f.name = FNT
        f.color.rgb = _rgb(C["body"])
        for c in (b.get("children") or [])[:3]:
            cp = tf.add_paragraph()
            cp.level = 1
            cp.space_after = Pt(4)
            cr = cp.add_run()
            cr.text = "–  " + str(c)
            cf = cr.font
            cf.size = Pt(14)
            cf.name = FNT
            cf.color.rgb = _rgb(C["muted"])
    if not any_line:
        tf.paragraphs[0].add_run().text = " "
    return tb


def _chips(slide, terms, y):
    if not terms:
        return
    x = 0.75
    max_right = W - 0.6
    for t in list(terms)[:6]:
        t = str(t)
        w = min(3.2, 0.5 + len(t) * 0.22)
        if x + w > max_right:
            x = 0.75
            y += 0.5
        if y > 6.7:
            break
        chip = _rect(slide, x, y, w, 0.4, C["chipBg"], line=C["chipBd"], line_w=1,
                     radius=0.5, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _shape_text(chip, t, size=11, color=C["accent"], align=PP_ALIGN.CENTER)
        x += w + 0.18


def _slide_cover(prs, deck, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, C["accent2"])
    _rect(slide, 0, 0, 0.35, H, C["accent"])
    _text(slide, 1, 2.1, 11, 0.5, (deck.get("course") or "").upper(), size=16, color=C["cover_ink"])
    _text(slide, 1, 2.7, 11.3, 1.8, s.get("title") or deck.get("title") or "",
          size=44, color=C["white"], bold=True)
    if s.get("subtitle"):
        _text(slide, 1, 4.5, 11, 0.8, s.get("subtitle"), size=20, color=C["cover_sub"])
    meta = "   ·   ".join(x for x in [deck.get("course"), deck.get("duration"), deck.get("model")] if x)
    if meta:
        _text(slide, 1, 6.6, 11, 0.4, meta, size=11, color=C["cover_meta"], font=MONO)
    _notes(slide, s.get("notes"))


def _slide_section(prs, deck, s, sec_no, total_sections):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(slide, 0, 0, W, H, C["accent"])
    _text(slide, 0.8, 1.2, 6, 3, f"{sec_no:02d}", size=160, color=C["sec_num"], bold=True, font=MONO)
    _text(slide, 5.4, 3.2, 3, 1, f"/ {total_sections:02d}", size=40, color=C["sec_slash"], font=MONO)
    _text(slide, 0.85, 4.7, 11.6, 1.4, s.get("title") or "", size=40, color=C["white"], bold=True)
    if s.get("subtitle"):
        _text(slide, 0.9, 6.0, 11.5, 0.7, s.get("subtitle"), size=18, color=C["cover_sub"])
    _notes(slide, s.get("notes"))


def _slide_agenda(prs, deck, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _footer(slide, deck)
    _head(slide, s.get("title") or "本节提纲")
    tb = slide.shapes.add_textbox(Inches(0.95), Inches(1.7), Inches(11.4), Inches(5.2))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for i, b in enumerate(list(s.get("bullets") or [])[:10]):
        t = b.get("text", "") if isinstance(b, dict) else str(b)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(14)
        run = p.add_run()
        run.text = f"  {i + 1:02d}   {t}"
        f = run.font
        f.size = Pt(20)
        f.name = FNT
        f.color.rgb = _rgb(C["body"])
    _notes(slide, s.get("notes"))


def _slide_qa(prs, deck, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _footer(slide, deck)
    _head(slide, s.get("title") or "课堂思考")
    qs = list(s.get("bullets") or [])[:7]
    top, bottom, gap = 1.7, 6.4, 0.16
    n = max(1, len(qs))
    ch = max(0.5, min(0.92, (bottom - top - gap * max(0, len(qs) - 1)) / n))
    y = top
    for i, b in enumerate(qs):
        t = b.get("text", "") if isinstance(b, dict) else str(b)
        _rect(slide, 0.75, y, 11.8, ch, C["white"], line=C["line"], line_w=1,
              radius=0.06, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        _rect(slide, 0.75, y, 0.12, ch, C["accent"])
        _text(slide, 0.95, y, 0.9, ch, f"Q{i + 1}", size=18, color=C["accent"], bold=True,
              anchor=MSO_ANCHOR.MIDDLE, font=MONO)
        _text(slide, 1.85, y, 10.5, ch, t, size=16, color=C["body"], anchor=MSO_ANCHOR.MIDDLE)
        y += ch + gap
    _chips(slide, s.get("key_terms"), min(max(y + 0.05, 6.4), 6.7))
    _notes(slide, s.get("notes"))


def _slide_content(prs, deck, s):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _footer(slide, deck)
    is_summary = s.get("layout") == "summary"
    title = ("本节小结 · " if is_summary else "") + (s.get("title") or "")
    time_ref = s.get("time_ref") if s.get("layout") == "content" else ""
    _head(slide, title, time_ref or "")
    has_chips = bool(s.get("key_terms"))
    _add_bullets(slide, 0.95, 1.7, 11.5, 4.5 if has_chips else 5.0, s.get("bullets"))
    if has_chips:
        _chips(slide, s.get("key_terms"), 6.45)
    _notes(slide, s.get("notes"))


def build_pptx(deck: dict, out_path: str) -> str:
    """渲染 deck 为 pptx 文件，返回输出路径。"""
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    prs.core_properties.title = deck.get("title") or "课堂 PPT"
    prs.core_properties.author = "edu_ai_home"

    slides = deck.get("slides") or []
    total_sections = sum(1 for s in slides if s.get("layout") == "section")
    sec_no = 0
    for s in slides:
        layout = s.get("layout")
        if layout == "cover":
            _slide_cover(prs, deck, s)
        elif layout == "section":
            sec_no += 1
            _slide_section(prs, deck, s, sec_no, total_sections)
        elif layout == "agenda":
            _slide_agenda(prs, deck, s)
        elif layout == "qa":
            _slide_qa(prs, deck, s)
        else:  # content / summary
            _slide_content(prs, deck, s)

    prs.save(out_path)
    return out_path
